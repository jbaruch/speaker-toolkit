"""Producer/validator regressions at bounded PPTX evidence edges."""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches


def _round_trip(pptx_extraction, pptx_evidence, deck: Path):
    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    decoded = pptx_evidence._decode_extraction_worker_payload(
        {
            "schema_version": pptx_evidence.PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "available",
            "extraction": extraction,
        },
        expected_ocr=False,
    )
    assert decoded == extraction
    return extraction


def _rewrite_archive(path: Path, replacements: dict[str, bytes]) -> None:
    with zipfile.ZipFile(path) as archive:
        members = {name: archive.read(name) for name in archive.namelist()}
    members.update(replacements)
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, payload in members.items():
            archive.writestr(name, payload)


def _set_background_image(slide, image_path: Path) -> str:
    _, relationship_id = slide.part.get_or_add_image_part(str(image_path))
    background = slide.element.cSld.get_or_add_bgPr()
    existing_fill = background.eg_fillProperties
    if existing_fill is not None:
        background.remove(existing_fill)
    background.insert(
        0,
        parse_xml(
            f"<a:blipFill {nsdecls('a', 'r')}>"
            f'<a:blip r:embed="{relationship_id}"/>'
            "<a:stretch><a:fillRect/></a:stretch>"
            "</a:blipFill>"
        ),
    )
    return relationship_id


def test_layout_name_is_canonicalized_before_validation(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    presentation.slide_layouts[6].name = "L" * 4097
    presentation.slides.add_slide(presentation.slide_layouts[6])
    deck = tmp_path / "long-layout-name.pptx"
    presentation.save(deck)

    extraction = _round_trip(pptx_extraction, pptx_evidence, deck)
    slide = extraction["per_slide_visual"][0]
    matching_layout = next(
        layout for layout in extraction["template_layouts"] if layout["index"] == 6
    )

    assert len(slide["layout_name"]) == 4096
    assert matching_layout["name"] == slide["layout_name"]


def test_font_name_is_canonicalized_before_global_cataloging(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = "x"
    run.font.name = "F" * 4097
    deck = tmp_path / "long-font-name.pptx"
    presentation.save(deck)

    extraction = _round_trip(pptx_extraction, pptx_evidence, deck)
    font_name = extraction["per_slide_visual"][0]["shapes_summary"][0]["font_name"]
    cataloged_fonts = extraction["global_design"]["fonts_used"]

    assert len(font_name) == 4096
    assert cataloged_fonts == {font_name: 1}


def test_graphic_uri_is_canonicalized_across_all_provenance(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(2), Inches(1))
    shape.element.graphic.graphicData.uri = "u" * 4097
    deck = tmp_path / "long-graphic-uri.pptx"
    presentation.save(deck)

    extraction = _round_trip(pptx_extraction, pptx_evidence, deck)
    slide_result = extraction["per_slide_visual"][0]
    shape_uri = slide_result["shapes_summary"][0]["graphic_data_uri"]
    unsupported_uri = slide_result["unsupported_content"][0]["graphic_data_uri"]
    channel_uri = slide_result["text_channels"][0]["provenance"]["graphic_data_uri"]

    assert len(shape_uri) == 4096
    assert unsupported_uri == shape_uri
    assert channel_uri == shape_uri


def test_background_relationship_id_is_canonicalized_before_validation(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "background.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    relationship_id = _set_background_image(slide, image_path)
    deck = tmp_path / "long-background-relationship.pptx"
    presentation.save(deck)

    with zipfile.ZipFile(deck) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml")
        relationships_xml = archive.read("ppt/slides/_rels/slide1.xml.rels")
    long_relationship_id = "r" + ("x" * 4096)
    original_embed = f'r:embed="{relationship_id}"'.encode()
    replacement_embed = f'r:embed="{long_relationship_id}"'.encode()
    original_relationship = f'Id="{relationship_id}"'.encode()
    replacement_relationship = f'Id="{long_relationship_id}"'.encode()
    assert original_embed in slide_xml
    assert original_relationship in relationships_xml
    _rewrite_archive(
        deck,
        {
            "ppt/slides/slide1.xml": slide_xml.replace(
                original_embed, replacement_embed
            ),
            "ppt/slides/_rels/slide1.xml.rels": relationships_xml.replace(
                original_relationship, replacement_relationship
            ),
        },
    )

    extraction = _round_trip(pptx_extraction, pptx_evidence, deck)
    background_channel = next(
        channel
        for channel in extraction["per_slide_visual"][0]["text_channels"]
        if channel["channel"] == "background_image_ocr"
    )

    assert len(background_channel["provenance"]["relationship_id"]) == 4096


def test_group_depth_over_contract_fails_with_stable_resource_reason(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    for _index in range(64):
        group = group.shapes.add_group_shape()
    deck = tmp_path / "group-depth-65.pptx"
    presentation.save(deck)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_probe_resource_unavailable"


def test_archive_part_name_over_contract_fails_with_stable_resource_reason(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "picture.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    deck = tmp_path / "long-picture-part.pptx"
    presentation.save(deck)

    old_part_name = "ppt/media/image1.png"
    part_prefix = "ppt/media/"
    part_suffix = ".png"
    new_part_name = (
        part_prefix + ("a" * (2049 - len(part_prefix) - len(part_suffix))) + part_suffix
    )
    with zipfile.ZipFile(deck) as archive:
        members = {name: archive.read(name) for name in archive.namelist()}
    picture_bytes = members.pop(old_part_name)
    relationships_name = "ppt/slides/_rels/slide1.xml.rels"
    old_target = b"../media/image1.png"
    new_target = f"../media/{new_part_name.rsplit('/', 1)[1]}".encode()
    assert len(new_part_name) == 2049
    assert old_target in members[relationships_name]
    members[new_part_name] = picture_bytes
    members[relationships_name] = members[relationships_name].replace(
        old_target, new_target
    )
    with zipfile.ZipFile(deck, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, payload in members.items():
            archive.writestr(name, payload)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_probe_resource_unavailable"


def test_valid_encoded_space_in_part_name_round_trips(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "encoded-space.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    deck = tmp_path / "encoded-space-part.pptx"
    presentation.save(deck)

    old_part_name = "ppt/media/image1.png"
    new_part_name = "ppt/media/image%201.png"
    relationships_name = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(deck) as archive:
        members = {name: archive.read(name) for name in archive.namelist()}
    members[new_part_name] = members.pop(old_part_name)
    members[relationships_name] = members[relationships_name].replace(
        b"../media/image1.png",
        b"../media/image%201.png",
    )
    with zipfile.ZipFile(deck, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, payload in members.items():
            archive.writestr(name, payload)

    extraction = _round_trip(pptx_extraction, pptx_evidence, deck)
    picture = extraction["per_slide_visual"][0]["shapes_summary"][0]

    assert picture["picture_part_name"] == new_part_name
