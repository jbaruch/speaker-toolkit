"""Tests for delete-slides.py — slide deletion."""

from pptx import Presentation

from conftest import make_deck


def test_single_delete(delete_slides, tmp_path):
    prs = make_deck(5)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    deleted = delete_slides.delete_slides(prs2, [2])
    assert 2 in deleted
    assert len(prs2.slides) == 4


def test_multi_delete(delete_slides, tmp_path):
    prs = make_deck(5)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    deleted = delete_slides.delete_slides(prs2, [0, 2, 4])
    assert len(deleted) == 3
    assert len(prs2.slides) == 2


def test_reverse_order_preservation(delete_slides, tmp_path):
    """Deleting in reverse order should preserve slide identities."""
    prs = make_deck(5)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    # Get the rIds before deletion for comparison
    xml_slides = list(prs2.slides._sldIdLst)
    original_ids = [s.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                    for s in xml_slides]

    # Delete slides 1 and 3 — slides 0, 2, 4 should remain
    delete_slides.delete_slides(prs2, [1, 3])
    remaining_ids = [s.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                     for s in prs2.slides._sldIdLst]
    assert remaining_ids == [original_ids[0], original_ids[2], original_ids[4]]


def test_out_of_range_warning(delete_slides, tmp_path, capsys):
    prs = make_deck(3)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    deleted = delete_slides.delete_slides(prs2, [99])
    assert deleted == []
    assert len(prs2.slides) == 3
    captured = capsys.readouterr()
    assert "WARNING" in captured.err


def test_duplicate_indices(delete_slides, tmp_path):
    prs = make_deck(5)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    deleted = delete_slides.delete_slides(prs2, [2, 2, 2])
    assert len(deleted) == 1
    assert len(prs2.slides) == 4


def test_viewprops_cleaned(delete_slides, tmp_path):
    from lxml import etree

    prs = make_deck(5)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    delete_slides.delete_slides(prs2, [0, 1])
    out = str(tmp_path / "out.pptx")
    prs2.save(out)

    prs3 = Presentation(out)
    for rel in prs3.part.rels.values():
        if "viewProps" in rel.reltype:
            vp_xml = etree.fromstring(rel.target_part.blob)
            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            assert vp_xml.findall(f'.//{{{ns_p}}}sldLst') == []
            break


def test_saved_file_opens_cleanly(delete_slides, tmp_path):
    prs = make_deck(5)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    delete_slides.delete_slides(prs2, [0, 2, 4])
    out = str(tmp_path / "clean.pptx")
    prs2.save(out)

    prs3 = Presentation(out)
    assert len(prs3.slides) == 2
