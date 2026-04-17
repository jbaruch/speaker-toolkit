"""Tests for strip-template.py — template slide stripping."""

from pptx import Presentation

from conftest import make_deck


def test_all_slides_removed(strip_template, tmp_path):
    prs = make_deck(5)
    path = str(tmp_path / "tmpl.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    removed = strip_template.strip_slides(prs2)

    assert removed == 5
    assert len(prs2.slides) == 0


def test_layouts_preserved(strip_template, tmp_path):
    prs = make_deck(3)
    layout_count_before = len(prs.slide_layouts)
    path = str(tmp_path / "tmpl.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    strip_template.strip_slides(prs2)
    out = str(tmp_path / "out.pptx")
    prs2.save(out)

    prs3 = Presentation(out)
    assert len(prs3.slides) == 0
    assert len(prs3.slide_layouts) == layout_count_before


def test_empty_deck_noop(strip_template, tmp_path):
    prs = Presentation()
    path = str(tmp_path / "empty.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    removed = strip_template.strip_slides(prs2)
    assert removed == 0


def test_saved_file_opens_cleanly(strip_template, tmp_path):
    prs = make_deck(5)
    path = str(tmp_path / "tmpl.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    strip_template.strip_slides(prs2)
    out = str(tmp_path / "stripped.pptx")
    prs2.save(out)

    # If viewProps aren't cleaned, this would raise or produce a corrupt file
    prs3 = Presentation(out)
    assert len(prs3.slides) == 0


def test_viewprops_sldlst_cleaned(strip_template, pptx_repair, tmp_path):
    """After stripping, viewProps XML should have no sldLst elements."""
    from lxml import etree

    prs = make_deck(3)
    path = str(tmp_path / "tmpl.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    strip_template.strip_slides(prs2)
    out = str(tmp_path / "stripped.pptx")
    prs2.save(out)

    prs3 = Presentation(out)
    for rel in prs3.part.rels.values():
        if "viewProps" in rel.reltype:
            vp_xml = etree.fromstring(rel.target_part.blob)
            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            assert vp_xml.findall(f'.//{{{ns_p}}}sldLst') == []
            break
