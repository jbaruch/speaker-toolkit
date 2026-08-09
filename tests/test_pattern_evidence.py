"""Direct adversarial tests for the source-located pattern-evidence boundary."""

from __future__ import annotations

import copy
import hashlib
import importlib
import json
import os
import struct
import sys
import zipfile
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import pytest
from PIL import Image
from pypdf import PdfWriter
from pptx import Presentation
from pptx.util import Inches


SCRIPTS = Path(__file__).parents[1] / "skills" / "vault-ingress" / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

pattern_evidence = importlib.import_module("pattern_evidence")
pdf_evidence = importlib.import_module("pdf_evidence")
pptx_evidence = importlib.import_module("pptx_evidence")
return_validation = importlib.import_module("return_validation")
transcript_timing = importlib.import_module("transcript_timing")
SYNTHETIC_VIDEO_ID = "abcdefghijk"
SYNTHETIC_DURATION = 600.0


def _video_probe(path: Path, *, duration: float = SYNTHETIC_DURATION) -> Any:
    raw = path.read_bytes()
    return pattern_evidence.VideoArtifactProbe(
        generation=SimpleNamespace(),
        root_generation=None,
        availability=SimpleNamespace(state="local"),
        source_sha256=hashlib.sha256(raw).hexdigest(),
        source_size_bytes=len(raw),
        duration_seconds=duration,
        duration_source="format",
        container_family="iso_bmff",
        stream_count=1,
        video_stream_count=1,
        audio_stream_count=0,
        attached_picture_count=0,
        other_stream_count=0,
        parser_diagnostics=SimpleNamespace(byte_count=0),
    )


class _TestVideoAssessment:
    def __init__(
        self,
        *,
        duration: float = SYNTHETIC_DURATION,
        failure: Exception | None = None,
    ) -> None:
        self.duration = duration
        self.failure = failure
        self.calls: list[Path] = []

    def probe(self, path: str | os.PathLike[str], **_kwargs: object) -> Any:
        video_path = Path(os.fspath(path))
        self.calls.append(video_path)
        if self.failure is not None:
            raise self.failure
        return _video_probe(video_path, duration=self.duration)


def _untrusted_video_manifest(vault: Path, source_video: Path) -> dict[str, Any]:
    context_pdf = source_video.with_suffix(".context.pdf")
    _write_pdf(context_pdf, page_count=1)
    source_path = source_video.relative_to(vault).as_posix()
    return {
        "schema_version": 3,
        "source_video_id": SYNTHETIC_VIDEO_ID,
        "source_video_path": source_path,
        "unique_frame_count": 1,
        "slide_region_method": "none",
        "slide_region_applied": False,
        "slide_region_verified": False,
        "review_required": True,
        "review_reason": "context only",
        "artifacts": [
            {
                "path": context_pdf.relative_to(vault).as_posix(),
                "artifact_scope": "full_frame_context",
                "page_count": 1,
                "source_video_id": SYNTHETIC_VIDEO_ID,
                "source_video_path": source_path,
                "crop_verified": False,
                "trusted_for_authored_slide_analysis": False,
            }
        ],
    }


def _foreign_absolute_locator(name: str) -> str:
    if os.name == "nt":
        return f"/foreign/{name}"
    return rf"C:\foreign\{name}"


def _identity(name: str) -> dict[str, str]:
    return {
        "artifact_root": "vault",
        "artifact_path": f"synthetic/{name}",
        "artifact_sha256": "a" * 64,
    }


def _write_pdf(path: Path, page_count: int = 2) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    writer = PdfWriter()
    for _ in range(page_count):
        writer.add_blank_page(width=640, height=480)
    with path.open("wb") as stream:
        writer.write(stream)


def _write_pptx(path: Path, slide_count: int = 2) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    deck = Presentation()
    layout = deck.slide_layouts[6]
    for _ in range(slide_count):
        deck.slides.add_slide(layout)
    deck.save(str(path))


def _write_crc_damaged_media_pptx(path: Path) -> str:
    """Create a deck whose embedded image follows the live BadZipFile path."""
    path.parent.mkdir(parents=True, exist_ok=True)
    image_path = path.with_suffix(".png")
    Image.new("RGB", (64, 64), "navy").save(image_path)
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    deck.save(str(path))
    with zipfile.ZipFile(path) as archive:
        member = next(
            item
            for item in archive.infolist()
            if item.filename.startswith("ppt/media/") and item.file_size
        )
    package = bytearray(path.read_bytes())
    name_size, extra_size = struct.unpack_from(
        "<HH", package, member.header_offset + 26
    )
    payload_offset = member.header_offset + 30 + name_size + extra_size
    package[payload_offset + (member.compress_size // 2)] ^= 0xFF
    path.write_bytes(package)
    return member.filename


def _transcript_lines() -> list[str]:
    lines = ["A uniquely phrased synthetic opening proves the source boundary."]
    lines.extend(
        f"Synthetic transcript line {number} explains a reliable evidence "
        "boundary with enough words."
        for number in range(2, 82)
    )
    return lines


def _write_transcript(
    vault: Path,
    *,
    name: str = "talk",
    timed: bool = False,
    quality_duration: float | None = None,
) -> tuple[Path, list[str]]:
    path = vault / "transcripts" / f"{name}.txt"
    path.parent.mkdir(parents=True, exist_ok=True)
    lines = _transcript_lines()
    text = "\n".join(lines)
    policy = transcript_timing.build_quality_policy(
        400, trusted_duration_seconds=quality_duration
    )
    provenance = (
        {"kind": "fixed_default"}
        if quality_duration is None
        else {
            "kind": "youtube_duration",
            "video_id": name,
            "duration_seconds": quality_duration,
        }
    )
    if timed:
        timing_video_id = name if len(name) == 11 else SYNTHETIC_VIDEO_ID
        timing_duration = quality_duration or SYNTHETIC_DURATION
        transcript_timing.write_transcript_bundle(
            path,
            text,
            [{"text": text, "start": 2.0, "end": 6.0}],
            source="whisper",
            timing_provenance=transcript_timing.youtube_timing_provenance(
                "whisper", timing_video_id, timing_duration
            ),
            quality_policy=policy,
            quality_policy_provenance=provenance,
        )
    else:
        path.write_text(text, encoding="utf-8")
        transcript_timing.write_quality_receipt(path, text, policy, provenance)
    return path, lines


def _entry(
    pattern_id: str = "synthetic-pattern",
    *,
    absence_gate: tuple[frozenset[str], ...] | None = (frozenset({"transcript"}),),
    channels: frozenset[str] = frozenset(
        {"transcript", "timed_transcript", "talk_metadata"}
    ),
) -> SimpleNamespace:
    return SimpleNamespace(
        pattern_id=pattern_id,
        entry_type="pattern",
        observable=True,
        evaluable_from=absence_gate,
        strong_evaluable_from=absence_gate,
        absence_evaluable_from=absence_gate,
        evidence_channels=channels,
        evidence_metadata_fields=frozenset({"title"}),
        vault_dimensions=(2, 5),
        path=f"synthetic/{pattern_id}.md",
    )


def _catalog(*entries: SimpleNamespace) -> SimpleNamespace:
    return SimpleNamespace(
        entries={entry.pattern_id: entry for entry in entries},
        fingerprint="f" * 64,
    )


def _raw_transcript_return(
    line_count: int,
    *,
    pattern_id: str = "synthetic-pattern",
    timed: bool = True,
) -> dict[str, Any]:
    channel = "timed_transcript" if timed else "transcript"
    return {
        "filename": "synthetic-talk.md",
        "return_schema_version": 4,
        "status": "processed_partial",
        "slide_source": "none",
        "transcript_source": "whisper" if timed else "manual",
        "structured_data": {},
        "pattern_observations": {
            "evidence_sources": ["transcript"],
            "source_inspection": [
                {"source": "transcript", "line_ranges": [[1, line_count]]}
            ],
            "patterns_detected": [
                {
                    "pattern_id": pattern_id,
                    "confidence": "moderate",
                    "evidence_source": "transcript",
                    "evidence": "The exact source contains the synthetic opening.",
                    "evidence_citations": [
                        {
                            "source": "transcript",
                            "channel": channel,
                            "quote": _transcript_lines()[0],
                        },
                        {
                            "source": "transcript",
                            "channel": "talk_metadata",
                            "field": "title",
                        },
                    ],
                }
            ],
            "antipatterns_detected": [],
            "not_evaluable": [],
            "pattern_score": 1,
        },
    }


def _raw_static_pdf_return(page_count: int) -> dict[str, Any]:
    return {
        "filename": "synthetic-talk.md",
        "return_schema_version": 4,
        "status": "processed",
        "slide_source": "pdf",
        "structured_data": {},
        "pattern_observations": {
            "evidence_sources": ["static_slides"],
            "source_inspection": [
                {"source": "static_slides", "page_ranges": [[1, page_count]]}
            ],
            "patterns_detected": [
                {
                    "pattern_id": "slide-pattern",
                    "confidence": "moderate",
                    "evidence_source": "static_slides",
                    "evidence": "The rendered slides establish the pattern.",
                    "evidence_citations": [
                        {
                            "source": "static_slides",
                            "channel": "slides",
                            "slide_numbers": [min(2, page_count)],
                        }
                    ],
                }
            ],
            "antipatterns_detected": [],
            "not_evaluable": [],
            "pattern_score": 1,
        },
    }


def _canonical_transcript_talk(
    tmp_path: Path,
) -> tuple[Path, dict[str, Any], dict[str, Any]]:
    vault = tmp_path / "vault"
    transcript, lines = _write_transcript(vault, name=SYNTHETIC_VIDEO_ID, timed=True)
    talk: dict[str, Any] = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "whisper",
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "source_identity": {
            "schema_version": 1,
            "provider": "youtube",
            "video_id": SYNTHETIC_VIDEO_ID,
            "duration_seconds": SYNTHETIC_DURATION,
        },
        "slide_source": "none",
    }
    raw = _raw_transcript_return(len(lines))
    canonical = pattern_evidence.canonicalize_return_evidence(
        raw,
        talk,
        vault,
        _catalog(_entry()),
    )
    persisted = copy.deepcopy(talk)
    persisted.update(copy.deepcopy(canonical))
    return vault, raw, persisted


def test_artifact_root_kinds_are_canonical_and_owner_bound(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    source_root = tmp_path / "pptx-source"
    transcript, _ = _write_transcript(vault)
    rendered_pdf = vault / "slides" / "rendered.pdf"
    configured_deck = source_root / "configured.pptx"
    absolute_deck = tmp_path / "external" / "absolute.pptx"
    _write_pdf(rendered_pdf)
    _write_pptx(configured_deck)
    _write_pptx(absolute_deck)

    configured = pattern_evidence.build_evidence_context(
        vault,
        {
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": configured_deck.name,
            "slide_source": "pptx",
        },
        source_roots={"pptx_source_dir": str(source_root)},
    )
    assert configured["transcript_artifact_identity"]["artifact_root"] == "vault"
    assert (
        configured["slide_artifact_identities"]["native_deck"]["artifact_root"]
        == "pptx_source"
    )

    configured_with_pdf = pattern_evidence.build_evidence_context(
        vault,
        {
            "slides_local_path": rendered_pdf.relative_to(vault).as_posix(),
            "slide_source": "pdf",
        },
        source_roots={"pptx_source_dir": str(source_root)},
    )
    assert (
        configured_with_pdf["slide_artifact_identities"]["static_slides"][
            "artifact_root"
        ]
        == "vault"
    )

    configured_absolute = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": str(configured_deck), "slide_source": "pptx"},
        source_roots={"pptx_source_dir": str(source_root)},
    )
    assert configured_absolute["slide_counts"] == {"native_deck": 2}
    assert (
        configured_absolute["slide_artifact_identities"]["native_deck"]["artifact_path"]
        == "configured.pptx"
    )

    absolute = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": str(absolute_deck), "slide_source": "pptx"},
    )
    assert (
        absolute["slide_artifact_identities"]["native_deck"]["artifact_root"]
        == "preclaim:pptx_path"
    )


def test_artifact_identity_names_the_actual_trusted_root(tmp_path: Path) -> None:
    source_root = tmp_path / "pptx-source"
    outside = tmp_path / "outside" / "deck.pptx"

    with pytest.raises(pattern_evidence.PatternEvidenceError) as caught:
        pattern_evidence._artifact_identity(
            source_root,
            outside,
            root_kind="pptx_source",
        )

    message = str(caught.value)
    assert "'pptx_source' artifact root" in message
    assert "vault root" not in message


@pytest.mark.parametrize(
    "field",
    ["slides_local_path", "slides_pdf_path", "pdf_path"],
)
def test_return_pdf_admission_accepts_only_the_exact_local_preclaim(
    tmp_path: Path,
    field: str,
) -> None:
    vault = tmp_path / "vault"
    _write_pdf(vault / "slides" / "named.pdf")
    talk = {
        field: "slides/named.pdf",
        "google_drive_id": "different-drive-id",
        "slide_source": "pdf",
    }
    returned = {
        "slides_local_path": "slides/named.pdf",
        "slide_source": "pdf",
    }

    pattern_evidence.admit_return_artifacts(vault, talk, returned)

    redirected = copy.deepcopy(returned)
    redirected["slides_local_path"] = "slides/different-drive-id.pdf"
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match=f"exact {field} preclaim",
    ):
        pattern_evidence.admit_return_artifacts(vault, talk, redirected)


def test_return_pdf_mismatch_never_leaks_or_probes_an_invalid_preclaim(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    secret_preclaim = _foreign_absolute_locator("credential-bearing-secret.pdf")
    secret_return = "slides/other-private-artifact.pdf"
    preclaim_resolutions: list[object] = []

    def reject_preclaim_resolution(*args, **kwargs):
        preclaim_resolutions.append((args, kwargs))
        pytest.fail("mismatched return reached preclaim artifact resolution")

    monkeypatch.setattr(
        pattern_evidence,
        "_resolve_preclaim_artifact",
        reject_preclaim_resolution,
    )

    with pytest.raises(pattern_evidence.PatternEvidenceError) as caught:
        pattern_evidence.admit_return_artifacts(
            tmp_path / "vault",
            {"slides_local_path": secret_preclaim, "slide_source": "pdf"},
            {"slides_local_path": secret_return, "slide_source": "pdf"},
        )

    assert str(caught.value) == (
        "return PDF path must match the exact slides_local_path preclaim"
    )
    assert preclaim_resolutions == []
    assert secret_preclaim not in str(caught.value)
    assert secret_return not in str(caught.value)


def test_pptx_preclaim_is_lexical_until_bounded_probe(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    source_root = tmp_path / "pptx-source"
    deck = source_root / "conference" / "talk.pptx"
    captured: dict[str, object] = {}

    def bounded_probe(path: Path, *, trusted_root: Path | None = None):
        captured.update(path=path, trusted_root=trusted_root)
        return SimpleNamespace(
            slide_count=2,
            source_sha256="a" * 64,
            source_size_bytes=123,
            archive_recovery=(),
        )

    monkeypatch.setattr(pattern_evidence, "probe_pptx_artifact", bounded_probe)
    original_resolve = Path.resolve

    def guarded_resolve(path: Path, *args: Any, **kwargs: Any) -> Path:
        if path == source_root or path.suffix.casefold() == ".pptx":
            pytest.fail("parent resolved a PPTX locator")
        return original_resolve(path, *args, **kwargs)

    monkeypatch.setattr(
        Path,
        "resolve",
        guarded_resolve,
    )
    original_is_file = Path.is_file

    def guarded_is_file(path: Path) -> bool:
        if path.suffix.casefold() == ".pptx":
            pytest.fail("parent tested a PPTX locator")
        return original_is_file(path)

    monkeypatch.setattr(
        Path,
        "is_file",
        guarded_is_file,
    )

    context = pattern_evidence.build_evidence_context(
        vault,
        {
            "pptx_path": "conference/talk.pptx",
            "slide_source": "pptx",
        },
        source_roots={"pptx_source_dir": str(source_root)},
    )

    assert captured == {"path": deck, "trusted_root": source_root}
    assert context["slide_counts"] == {"native_deck": 2}
    assert context["slide_artifact_identities"]["native_deck"] == {
        "artifact_root": "pptx_source",
        "artifact_path": "conference/talk.pptx",
        "artifact_sha256": "a" * 64,
    }


def test_pptx_preclaim_rejects_raw_dot_segments_before_normalization(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    source_root = tmp_path / "pptx-source"

    def bounded_probe_forbidden(*_args: Any, **_kwargs: Any):
        pytest.fail("ambiguous PPTX locator reached the bounded probe")

    monkeypatch.setattr(
        pattern_evidence,
        "probe_pptx_artifact",
        bounded_probe_forbidden,
    )
    locators = (
        os.path.join("conference", ".", "talk.pptx"),
        os.path.join("conference", "..", "talk.pptx"),
        os.path.join(str(source_root), "conference", ".", "talk.pptx"),
        os.path.join(str(source_root), "conference", "..", "talk.pptx"),
    )

    for locator in locators:
        context = pattern_evidence.build_evidence_context(
            vault,
            {"pptx_path": locator, "slide_source": "pptx"},
            source_roots={"pptx_source_dir": str(source_root)},
        )

        assert "native_deck" not in context["verified_evidence_sources"]
        assert context["source_reasons"]["native_deck"].endswith(
            "artifact_locator_dot_segment"
        )


@pytest.mark.parametrize(
    "locator",
    [
        r"conference\.\talk.pptx",
        r"conference\..\talk.pptx",
        "conference/./talk.pptx",
        "conference/../talk.pptx",
        r"conference\track/../talk.pptx",
        r"C:.\talk.pptx",
        r"C:..\talk.pptx",
        r"C:folder/../talk.pptx",
        r"\\server\share\..\talk.pptx",
        r"\\server\..\talk.pptx",
        "//server/../talk.pptx",
    ],
)
def test_pptx_dot_segment_guard_is_platform_independent(locator: str) -> None:
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match=(
            r"artifact_locator_(dot_segment|windows_drive_relative|"
            r"ambiguous_double_slash)"
        ),
    ):
        pattern_evidence._reject_ambiguous_path_segments(
            locator,
            label="pptx_path",
        )


@pytest.mark.parametrize(
    "locator",
    [
        r"C:talk.pptx",
        r"c:conference\talk.pptx",
        r"D:conference/mixed\talk.pptx",
        "E:",
        r"\conference\talk.pptx",
        r"\conference/mixed\talk.pptx",
    ],
)
def test_pptx_windows_cwd_dependent_guard_is_platform_independent(
    locator: str,
) -> None:
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match=r"artifact_locator_windows_(drive_relative|current_drive_rooted)",
    ) as caught:
        pattern_evidence._reject_ambiguous_path_segments(
            locator,
            label="pptx_path",
        )

    assert locator not in str(caught.value)


@pytest.mark.parametrize(
    "locator",
    [
        r"\\?\C:\conference\talk.pptx",
        r"\\?\UNC\server\share\talk.pptx",
        r"\\.\pipe\talk.pptx",
        r"\??\C:\conference\talk.pptx",
        "//?/C:/conference/talk.pptx",
    ],
)
def test_pptx_windows_device_namespace_guard_is_platform_independent(
    locator: str,
) -> None:
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match=(r"artifact_locator_(windows_device_namespace|ambiguous_double_slash)"),
    ) as caught:
        pattern_evidence._reject_ambiguous_path_segments(
            locator,
            label="pptx_path",
        )

    assert locator not in str(caught.value)


@pytest.mark.parametrize(
    "locator",
    [
        "conference/track/talk.pptx",
        r"C:\conference\talk.pptx",
        "C:/conference/talk.pptx",
        r"C:\conference/mixed\talk.pptx",
        r"\\server\share\conference\talk.pptx",
        r"\\server\share/conference\talk.pptx",
    ],
)
def test_pptx_locator_guard_accepts_canonical_relative_or_absolute_flavors(
    locator: str,
) -> None:
    pattern_evidence._reject_ambiguous_path_segments(locator, label="pptx_path")


@pytest.mark.parametrize(
    "locator,reason_code",
    [
        (r"conference\track\talk.pptx", "artifact_locator_noncanonical_relative"),
        (
            r"conference\track/mixed/talk.pptx",
            "artifact_locator_noncanonical_relative",
        ),
        (
            "//server/share/conference/talk.pptx",
            "artifact_locator_ambiguous_double_slash",
        ),
        ("conference./talk.pptx", "artifact_locator_windows_trimmed_component"),
        (
            "conference/deck:stream.pptx",
            "artifact_locator_windows_reserved_character",
        ),
        ("CON/talk.pptx", "artifact_locator_windows_reserved_name"),
    ],
)
def test_pptx_locator_guard_rejects_cross_host_relative_or_dual_flavors(
    locator: str,
    reason_code: str,
) -> None:
    with pytest.raises(pattern_evidence.PatternEvidenceError, match=reason_code):
        pattern_evidence._reject_ambiguous_path_segments(
            locator,
            label="pptx_path",
        )


def test_forward_slash_root_is_classified_without_host_reinterpretation() -> None:
    pattern_evidence._reject_ambiguous_path_segments(
        "/conference/talk.pptx",
        label="pptx_path",
    )


@pytest.mark.parametrize(
    "locator",
    [
        r"C:conference\talk.pptx",
        r"\conference\talk.pptx",
    ],
)
def test_pptx_cwd_dependent_preclaim_never_reaches_bounded_probe(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    locator: str,
) -> None:
    monkeypatch.setattr(
        pattern_evidence,
        "probe_pptx_artifact",
        lambda *_args, **_kwargs: pytest.fail("ambiguous locator reached probe"),
    )
    context = pattern_evidence.build_evidence_context(
        tmp_path / "vault",
        {"pptx_path": locator, "slide_source": "pptx"},
        source_roots={"pptx_source_dir": str(tmp_path / "pptx-source")},
    )

    assert "native_deck" not in context["verified_evidence_sources"]
    reason = context["source_reasons"]["native_deck"]
    expected_code = (
        "artifact_locator_windows_drive_relative"
        if locator.startswith("C:")
        else "artifact_locator_windows_current_drive_rooted"
    )
    assert reason == f"pptx_path: {expected_code}"
    assert locator not in reason


@pytest.mark.parametrize(
    "locator,reason_code",
    [
        ("~/conference/talk.pptx", "artifact_locator_home_expansion_unsupported"),
        (
            r"conference\track\talk.pptx",
            "artifact_locator_noncanonical_relative",
        ),
        (
            "//server/share/conference/talk.pptx",
            "artifact_locator_ambiguous_double_slash",
        ),
        (
            _foreign_absolute_locator("talk.pptx"),
            "artifact_locator_foreign_absolute",
        ),
        ("conference./talk.pptx", "artifact_locator_windows_trimmed_component"),
        (
            "conference/deck:stream.pptx",
            "artifact_locator_windows_reserved_character",
        ),
        ("CON/talk.pptx", "artifact_locator_windows_reserved_name"),
    ],
)
def test_rejected_pptx_locator_never_selects_a_shadow_file_or_reaches_probe(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    locator: str,
    reason_code: str,
) -> None:
    vault = tmp_path / "vault"
    source_root = tmp_path / "pptx-source"
    if os.name != "nt" and not locator.startswith("//"):
        _write_pptx(source_root / locator)
    monkeypatch.setattr(
        pattern_evidence,
        "probe_pptx_artifact",
        lambda *_args, **_kwargs: pytest.fail("rejected locator reached probe"),
    )

    context = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": locator, "slide_source": "pptx"},
        source_roots={"pptx_source_dir": str(source_root)},
    )

    assert "native_deck" not in context["verified_evidence_sources"]
    reason = context["source_reasons"]["native_deck"]
    assert reason == f"pptx_path: {reason_code}"
    assert locator not in reason


@pytest.mark.parametrize(
    "configured_root,reason_code",
    [
        ("", "artifact_locator_empty_or_whitespace"),
        ("relative-presentations", "artifact_root_not_native_absolute"),
        ("~/presentations", "artifact_locator_home_expansion_unsupported"),
        (
            _foreign_absolute_locator("presentations"),
            "artifact_locator_foreign_absolute",
        ),
    ],
)
def test_invalid_configured_pptx_root_fails_closed_without_vault_fallback(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    configured_root: str,
    reason_code: str,
) -> None:
    vault = tmp_path / "vault"
    _write_pptx(vault / "deck.pptx")
    monkeypatch.setattr(
        pattern_evidence,
        "probe_pptx_artifact",
        lambda *_args, **_kwargs: pytest.fail(
            "invalid configured root fell back to the vault"
        ),
    )

    context = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": "deck.pptx", "slide_source": "pptx"},
        source_roots={"pptx_source_dir": configured_root},
    )

    assert "native_deck" not in context["verified_evidence_sources"]
    assert context["source_reasons"]["native_deck"] == (
        f"pptx_source_dir: {reason_code}"
    )


def test_null_configured_pptx_root_retains_vault_fallback(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    deck = vault / "deck.pptx"
    _write_pptx(deck)

    context = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": deck.name, "slide_source": "pptx"},
        source_roots={"pptx_source_dir": None},
    )

    assert context["slide_counts"] == {"native_deck": 2}
    assert context["slide_artifact_identities"]["native_deck"]["artifact_root"] == (
        "vault"
    )


def test_foreign_pdf_and_video_locators_fail_only_their_independent_lanes(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(
        vault,
        name=SYNTHETIC_VIDEO_ID,
        timed=False,
    )
    foreign_pdf = _foreign_absolute_locator("slides.pdf")
    foreign_video = _foreign_absolute_locator(f"{SYNTHETIC_VIDEO_ID}.mp4")
    monkeypatch.setattr(
        pattern_evidence,
        "probe_pdf_artifact",
        lambda *_args, **_kwargs: pytest.fail("foreign PDF reached probe"),
    )
    context = pattern_evidence.build_evidence_context(
        vault,
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "slides_local_path": foreign_pdf,
            "slide_source": "pdf",
            "video_local_path": foreign_video,
            "youtube_id": SYNTHETIC_VIDEO_ID,
        },
    )

    assert context["verified_evidence_sources"] == {"transcript"}
    assert context["source_reasons"]["static_slides"] == (
        "slides_local_path: artifact_locator_foreign_absolute"
    )
    assert context["source_reasons"]["delivery_video"] == (
        "video_local_path: artifact_locator_foreign_absolute"
    )


def test_foreign_video_manifest_pdf_never_reaches_current_context_probe(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    foreign_pdf = _foreign_absolute_locator(f"{SYNTHETIC_VIDEO_ID}.slide-region.pdf")
    monkeypatch.setattr(
        pattern_evidence,
        "probe_pdf_artifact",
        lambda *_args, **_kwargs: pytest.fail("foreign manifest PDF reached probe"),
    )
    manifest = {
        "unique_frame_count": 1,
        "artifacts": [
            {
                "artifact_scope": "slide_region",
                "source_video_id": SYNTHETIC_VIDEO_ID,
                "page_count": 1,
                "path": foreign_pdf,
            }
        ],
    }

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="artifact_locator_foreign_absolute",
    ) as caught:
        pattern_evidence._probe_video_manifest_artifacts(
            tmp_path / "vault",
            manifest,
            SYNTHETIC_VIDEO_ID,
        )

    assert foreign_pdf not in str(caught.value)


def test_schema_v3_manifest_source_takes_precedence_over_top_level_video_path(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    legacy_video = vault / "videos" / f"{SYNTHETIC_VIDEO_ID}.mp4"
    manifest_video = (
        vault / "slides-rebuild" / SYNTHETIC_VIDEO_ID / f"{SYNTHETIC_VIDEO_ID}.mp4"
    )
    legacy_video.parent.mkdir(parents=True)
    manifest_video.parent.mkdir(parents=True)
    legacy_video.write_bytes(b"legacy top-level bytes")
    manifest_video.write_bytes(b"manifest-owned bytes")
    assessment = _TestVideoAssessment()

    context = pattern_evidence.build_evidence_context(
        vault,
        {
            "youtube_id": SYNTHETIC_VIDEO_ID,
            "video_local_path": legacy_video.relative_to(vault).as_posix(),
            "structured_data": {
                "video_extraction": {
                    "schema_version": 3,
                    "source_video_id": SYNTHETIC_VIDEO_ID,
                    "source_video_path": manifest_video.relative_to(vault).as_posix(),
                }
            },
        },
        video_evidence_assessment=assessment,
    )

    assert context["local_video_path"] == manifest_video
    assert (
        context["video_artifact_identity"]["artifact_sha256"]
        == hashlib.sha256(manifest_video.read_bytes()).hexdigest()
    )
    assert assessment.calls == [manifest_video]


def test_schema_v3_video_identity_stays_fresh_through_a_symlinked_vault_root(
    tmp_path: Path,
) -> None:
    storage = tmp_path / "vault-storage"
    storage.mkdir()
    vault = tmp_path / "vault"
    try:
        vault.symlink_to(storage, target_is_directory=True)
    except OSError as exc:  # pragma: no cover - unusual restricted platform
        pytest.skip(f"symlinks unavailable: {exc}")
    source_video = (
        storage / "slides-rebuild" / SYNTHETIC_VIDEO_ID / f"{SYNTHETIC_VIDEO_ID}.mp4"
    )
    source_video.parent.mkdir(parents=True)
    source_video.write_bytes(b"canonical storage video bytes")
    manifest = {
        "schema_version": 3,
        "source_video_id": SYNTHETIC_VIDEO_ID,
        "source_video_path": os.fspath(source_video),
    }
    assessment = _TestVideoAssessment()
    talk: dict[str, Any] = {
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "structured_data": {"video_extraction": manifest},
    }

    context = pattern_evidence.build_evidence_context(
        vault,
        talk,
        video_evidence_assessment=assessment,
    )

    assert context["local_video_path"] == source_video
    assert context["video_artifact_identity"] == {
        "artifact_root": "vault",
        "artifact_path": (
            f"slides-rebuild/{SYNTHETIC_VIDEO_ID}/{SYNTHETIC_VIDEO_ID}.mp4"
        ),
        "artifact_sha256": hashlib.sha256(source_video.read_bytes()).hexdigest(),
    }
    talk["pattern_observations"] = {
        "evidence_schema_version": 1,
        "evidence_sources": ["delivery_video"],
        "source_inspection": [
            {
                "source": "delivery_video",
                **context["video_artifact_identity"],
                "time_ranges": [[0.0, SYNTHETIC_DURATION]],
                "duration_seconds": SYNTHETIC_DURATION,
                "coverage_complete": True,
            }
        ],
        "patterns_detected": [],
        "antipatterns_detected": [],
    }
    freshness_assessment = _TestVideoAssessment()

    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            talk,
            vault_root=vault,
            video_evidence_assessment=freshness_assessment,
        )
        == ()
    )
    assert assessment.calls == [source_video]
    assert freshness_assessment.calls == [source_video]

    # Older canonicalization represented a video reached through a symlinked
    # vault as a parent-rooted preclaim identity. Keep that exact persisted
    # receipt readable so installing the bounded probe alone does not require
    # reprocessing otherwise-current evidence.
    legacy_record = talk["pattern_observations"]["source_inspection"][0]
    legacy_record["artifact_root"] = "preclaim:video_local_path"
    legacy_record["artifact_path"] = source_video.name
    legacy_assessment = _TestVideoAssessment()
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            talk,
            vault_root=vault,
            video_evidence_assessment=legacy_assessment,
        )
        == ()
    )
    assert legacy_assessment.calls == [source_video]


def test_current_return_video_does_not_retroactively_authorize_transcript(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(
        vault,
        name=SYNTHETIC_VIDEO_ID,
        timed=False,
    )
    source_video = (
        vault / "slides-rebuild" / SYNTHETIC_VIDEO_ID / f"{SYNTHETIC_VIDEO_ID}.mp4"
    )
    source_video.parent.mkdir(parents=True)
    source_video.write_bytes(b"current-return video bytes")
    text = transcript.read_text(encoding="utf-8")
    transcript_timing.write_quality_receipt(
        transcript,
        text,
        transcript_timing.build_quality_policy(
            400,
            trusted_duration_seconds=SYNTHETIC_DURATION,
        ),
        {
            "kind": "local_media_duration",
            "media_sha256": hashlib.sha256(source_video.read_bytes()).hexdigest(),
            "duration_seconds": SYNTHETIC_DURATION,
        },
    )
    ret = {
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "slide_source": "video_extracted",
        "structured_data": {
            "video_extraction": _untrusted_video_manifest(vault, source_video)
        },
    }

    context = pattern_evidence.build_evidence_context(
        vault,
        {
            "youtube_id": SYNTHETIC_VIDEO_ID,
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
        },
        ret,
        video_evidence_assessment=_TestVideoAssessment(),
    )

    assert "transcript" not in context["verified_evidence_sources"]
    assert "not bound to the claimed talk" in context["transcript_reason"]
    assert "delivery_video" in context["verified_evidence_sources"]


def test_current_return_video_receipt_is_preferred_for_delivery_evidence(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    preclaim_video = vault / "videos" / f"{SYNTHETIC_VIDEO_ID}.mp4"
    returned_video = (
        vault / "slides-rebuild" / SYNTHETIC_VIDEO_ID / f"{SYNTHETIC_VIDEO_ID}.mp4"
    )
    preclaim_video.parent.mkdir(parents=True)
    returned_video.parent.mkdir(parents=True)
    preclaim_video.write_bytes(b"preclaim video bytes")
    returned_video.write_bytes(b"current-return video bytes")
    assessment = _TestVideoAssessment()
    ret = {
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "slide_source": "video_extracted",
        "structured_data": {
            "video_extraction": _untrusted_video_manifest(vault, returned_video)
        },
    }

    context = pattern_evidence.build_evidence_context(
        vault,
        {
            "youtube_id": SYNTHETIC_VIDEO_ID,
            "video_local_path": preclaim_video.relative_to(vault).as_posix(),
        },
        ret,
        video_evidence_assessment=assessment,
    )

    assert context["local_video_path"] == returned_video
    assert (
        context["video_artifact_identity"]["artifact_sha256"]
        == hashlib.sha256(returned_video.read_bytes()).hexdigest()
    )
    assert assessment.calls == [preclaim_video, returned_video]


def test_video_freshness_reuses_the_context_receipt_without_a_second_probe(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    video = vault / "videos" / f"{SYNTHETIC_VIDEO_ID}.mp4"
    video.parent.mkdir(parents=True)
    video.write_bytes(b"one exact persisted video generation")
    digest = hashlib.sha256(video.read_bytes()).hexdigest()
    assessment = _TestVideoAssessment()
    talk = {
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "video_local_path": video.relative_to(vault).as_posix(),
        "pattern_observations": {
            "evidence_schema_version": 1,
            "evidence_sources": ["delivery_video"],
            "source_inspection": [
                {
                    "source": "delivery_video",
                    "artifact_root": "vault",
                    "artifact_path": video.relative_to(vault).as_posix(),
                    "artifact_sha256": digest,
                    "time_ranges": [[0.0, SYNTHETIC_DURATION]],
                    "duration_seconds": SYNTHETIC_DURATION,
                    "coverage_complete": True,
                }
            ],
            "patterns_detected": [],
            "antipatterns_detected": [],
        },
    }

    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        talk,
        vault_root=vault,
        video_evidence_assessment=assessment,
    )

    assert reasons == ()
    assert assessment.calls == [video]


def test_traversal_and_symlinked_artifacts_never_become_sources(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    source_root = tmp_path / "source"
    outside = tmp_path / "outside.pptx"
    _write_pptx(outside)

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="transcripts/<artifact>",
    ):
        pattern_evidence.validate_transcript_path("transcripts/../escape.txt")

    traversal = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": "../outside.pptx", "slide_source": "pptx"},
        source_roots={"pptx_source_dir": str(source_root)},
    )
    assert "native_deck" not in traversal["verified_evidence_sources"]
    assert traversal["source_reasons"]["native_deck"].endswith(
        "artifact_locator_dot_segment"
    )

    source_root.mkdir(parents=True, exist_ok=True)
    real = source_root / "real.pptx"
    _write_pptx(real)
    link = source_root / "linked.pptx"
    try:
        link.symlink_to(real.name)
    except OSError as exc:  # pragma: no cover - unusual restricted platform
        pytest.skip(f"symlinks unavailable: {exc}")
    linked = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": link.name, "slide_source": "pptx"},
        source_roots={"pptx_source_dir": str(source_root)},
    )
    assert "native_deck" not in linked["verified_evidence_sources"]
    assert "symbolic link" in linked["source_reasons"]["native_deck"]

    absolute_linked = pattern_evidence.build_evidence_context(
        vault,
        {"pptx_path": str(link), "slide_source": "pptx"},
        source_roots={"pptx_source_dir": str(source_root)},
    )
    assert "native_deck" not in absolute_linked["verified_evidence_sources"]
    assert "symbolic link" in absolute_linked["source_reasons"]["native_deck"]


@pytest.mark.parametrize(
    (
        "source",
        "range_field",
        "ranges",
        "expected_receipt_complete",
        "expected_absence_complete",
    ),
    [
        ("transcript", "line_ranges", [[1, 4]], True, True),
        ("transcript", "line_ranges", [[1, 2], [3, 4]], True, True),
        ("transcript", "line_ranges", [[1, 2], [4, 4]], False, False),
        ("transcript", "line_ranges", [[2, 4]], False, False),
        ("static_slides", "page_ranges", [[1, 2], [3, 4]], True, True),
        ("static_slides", "page_ranges", [[1, 2], [4, 4]], False, False),
        ("delivery_video", "time_ranges", [[0, 2], [2, 4]], True, False),
        ("delivery_video", "time_ranges", [[0, 2], [2.5, 4]], False, False),
    ],
)
def test_source_inspection_completeness_is_derived_from_exact_ranges(
    source: str,
    range_field: str,
    ranges: list[list[int | float]],
    expected_receipt_complete: bool,
    expected_absence_complete: bool,
) -> None:
    context = {
        "verified_evidence_sources": {source},
        "transcript_line_count": 4,
        "transcript_artifact_identity": _identity("transcript.txt"),
        "slide_counts": {source: 4},
        "slide_artifact_identities": {source: _identity("slides.pdf")},
        "video_duration_seconds": 4.0,
        "video_artifact_identity": _identity("video.mp4"),
    }
    records, complete_sources, _ = pattern_evidence.canonicalize_source_inspection(
        [{"source": source, range_field: ranges}],
        context,
    )
    assert records[0]["coverage_complete"] is expected_receipt_complete
    assert (source in complete_sources) is expected_absence_complete
    assert records[0]["absence_capability_complete"] is (expected_absence_complete)
    expected_reason = (
        "incomplete_range_coverage"
        if not expected_receipt_complete
        else {
            "transcript": "authorized_transcript",
            "static_slides": "authorized_rendered_static",
            "delivery_video": "bare_delivery_video",
        }[source]
    )
    assert records[0]["absence_capability_reason"] == expected_reason


def test_video_extracted_static_pages_support_positive_inspection_not_absence() -> None:
    context = {
        "verified_evidence_sources": {"static_slides"},
        "absence_complete_evidence_sources": set(),
        "absence_capability_reasons": {
            "static_slides": "nonexhaustive_video_extraction"
        },
        "slide_counts": {"static_slides": 4},
        "slide_artifact_identities": {
            "static_slides": _identity("video-extracted.pdf")
        },
    }

    records, complete_sources, _ = pattern_evidence.canonicalize_source_inspection(
        [{"source": "static_slides", "page_ranges": [[1, 4]]}],
        context,
    )

    assert records[0]["coverage_complete"] is True
    assert records[0]["absence_capability_complete"] is False
    assert records[0]["absence_capability_reason"] == ("nonexhaustive_video_extraction")
    assert "static_slides" not in complete_sources


def test_video_extracted_pdf_path_and_digest_replace_predeclared_pdf_identity(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    declared = vault / "slides" / "declared.pdf"
    extracted = vault / "slides-rebuild" / "video.pdf"
    _write_pdf(declared, page_count=2)
    _write_pdf(extracted, page_count=2)
    extracted_bytes = extracted.read_bytes()
    assert b"pypdf" in extracted_bytes
    extracted.write_bytes(extracted_bytes.replace(b"pypdf", b"qypdf", 1))
    assert declared.read_bytes() != extracted.read_bytes()
    extracted_probe = pattern_evidence.probe_pdf_artifact(
        extracted,
        trusted_root=vault,
    )
    monkeypatch.setattr(
        pattern_evidence,
        "_trusted_video_slide_probe",
        lambda *_args, **_kwargs: (
            extracted_probe,
            "trusted extracted PDF",
            extracted,
        ),
    )
    talk = {
        "filename": "video-talk.md",
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "slides_pdf_path": declared.relative_to(vault).as_posix(),
        "slide_source": "video_extracted",
    }

    context = pattern_evidence.build_evidence_context(vault, talk)
    identity = context["slide_artifact_identities"]["static_slides"]

    assert identity["artifact_path"] == extracted.relative_to(vault).as_posix()
    assert (
        identity["artifact_sha256"]
        == hashlib.sha256(extracted.read_bytes()).hexdigest()
    )
    assert (
        identity["artifact_sha256"] != hashlib.sha256(declared.read_bytes()).hexdigest()
    )


def test_unreadable_video_extracted_pdf_does_not_hide_independent_sources(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, name=SYNTHETIC_VIDEO_ID)
    deck = vault / "decks" / "talk.pptx"
    extracted = vault / "slides-rebuild" / "video.pdf"
    _write_pptx(deck, slide_count=2)
    _write_pdf(extracted, page_count=2)

    def unavailable_video_pdf(*_args: Any, **_kwargs: Any):
        raise pdf_evidence.PdfEvidenceError(
            "synthetic video PDF generation changed",
            reason_code="pdf_artifact_changed",
        )

    monkeypatch.setattr(
        pattern_evidence,
        "_trusted_video_slide_probe",
        unavailable_video_pdf,
    )
    talk = {
        "filename": "video-talk.md",
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "pptx_path": deck.relative_to(vault).as_posix(),
        "slide_source": "video_extracted",
    }

    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        talk,
        vault_root=vault,
    )

    assert assessment["verified_capabilities"] == ("slides", "transcript")
    assert assessment["verified_evidence_sources"] == (
        "native_deck",
        "transcript",
    )
    assert "static_slides" not in assessment["verified_evidence_sources"]
    assert (
        "synthetic video PDF generation changed"
        in assessment["source_reasons"]["static_slides"]
    )
    unavailable = assessment["unavailable_evidence_sources"]["static_slides"]
    assert unavailable["reason_code"] == "pdf_artifact_changed"


def test_current_video_manifest_bounds_bad_second_context_pdf_before_persistence(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    rebuild = vault / "slides-rebuild" / SYNTHETIC_VIDEO_ID
    slide_region = rebuild / f"{SYNTHETIC_VIDEO_ID}.slide-region.pdf"
    missing_context = rebuild / f"{SYNTHETIC_VIDEO_ID}.context.pdf"
    _write_pdf(slide_region, page_count=1)
    source_video = rebuild / f"{SYNTHETIC_VIDEO_ID}.mp4"
    source_video.write_bytes(b"synthetic video")
    shared = {
        "page_count": 1,
        "source_video_id": SYNTHETIC_VIDEO_ID,
        "source_video_path": os.fspath(source_video),
    }
    manifest = {
        "schema_version": 3,
        "source_video_id": SYNTHETIC_VIDEO_ID,
        "source_video_path": os.fspath(source_video),
        "unique_frame_count": 1,
        "slide_region_method": "manual",
        "slide_region_applied": True,
        "slide_region_verified": True,
        "review_required": False,
        "review_reason": None,
        "artifacts": [
            {
                "path": os.fspath(slide_region),
                "artifact_scope": "slide_region",
                "crop_verified": True,
                "trusted_for_authored_slide_analysis": True,
                **shared,
            },
            {
                "path": os.fspath(missing_context),
                "artifact_scope": "full_frame_context",
                "crop_verified": False,
                "trusted_for_authored_slide_analysis": False,
                **shared,
            },
        ],
    }
    ret = {
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "slide_source": "video_extracted",
        "structured_data": {"video_extraction": manifest},
    }

    with pytest.raises(pdf_evidence.PdfEvidenceError) as caught:
        pattern_evidence._returned_slide_artifact(
            vault,
            {"youtube_id": SYNTHETIC_VIDEO_ID},
            ret,
            video_evidence_assessment=_TestVideoAssessment(),
        )

    assert caught.value.reason_code == "pdf_artifact_unavailable"
    assert caught.value.details["failure_kind"] == "missing"


def test_promoted_video_pdf_must_match_trusted_slide_region_digest(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    rebuild = vault / "slides-rebuild" / SYNTHETIC_VIDEO_ID
    source_video = rebuild / f"{SYNTHETIC_VIDEO_ID}.mp4"
    source_video.parent.mkdir(parents=True)
    source_video.write_bytes(b"synthetic video")
    slide_region = rebuild / f"{SYNTHETIC_VIDEO_ID}.slide-region.pdf"
    promoted = vault / "slides" / f"{SYNTHETIC_VIDEO_ID}.pdf"
    _write_pdf(slide_region, page_count=1)
    promoted.parent.mkdir(parents=True)
    writer = PdfWriter()
    writer.add_blank_page(width=640, height=480)
    writer.add_metadata({"/Title": "Different valid one-page PDF"})
    with promoted.open("wb") as stream:
        writer.write(stream)
    source_path = source_video.relative_to(vault).as_posix()
    manifest = {
        "schema_version": 3,
        "source_video_id": SYNTHETIC_VIDEO_ID,
        "source_video_path": source_path,
        "unique_frame_count": 1,
        "slide_region_method": "manual",
        "slide_region_applied": True,
        "slide_region_verified": True,
        "review_required": False,
        "review_reason": None,
        "artifacts": [
            {
                "path": slide_region.relative_to(vault).as_posix(),
                "artifact_scope": "slide_region",
                "page_count": 1,
                "source_video_id": SYNTHETIC_VIDEO_ID,
                "source_video_path": source_path,
                "crop_verified": True,
                "trusted_for_authored_slide_analysis": True,
            }
        ],
    }
    returned = {
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "slide_source": "video_extracted",
        "slides_local_path": promoted.relative_to(vault).as_posix(),
        "structured_data": {"video_extraction": manifest},
    }

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="content digest disagrees",
    ):
        pattern_evidence.admit_return_artifacts(
            vault,
            {"youtube_id": SYNTHETIC_VIDEO_ID},
            returned,
            video_evidence_assessment=_TestVideoAssessment(),
        )


def test_artifact_admission_rejects_video_manifest_outside_video_lane(
    tmp_path: Path,
) -> None:
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="requires slide_source 'video_extracted'",
    ):
        pattern_evidence.admit_return_artifacts(
            tmp_path,
            {},
            {
                "slide_source": "pptx",
                "structured_data": {
                    "video_extraction": {
                        "artifacts": [{"path": "/outside/untrusted.pdf"}]
                    }
                },
            },
        )


def test_persisted_context_manifest_reanalysis_never_touches_pdf_in_owner(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    vault = tmp_path / "vault"
    rebuild = vault / "slides-rebuild" / SYNTHETIC_VIDEO_ID
    context = rebuild / f"{SYNTHETIC_VIDEO_ID}.context.pdf"
    _write_pdf(context, page_count=1)
    source_video = rebuild / f"{SYNTHETIC_VIDEO_ID}.mp4"
    source_video.write_bytes(b"synthetic video")
    manifest = {
        "schema_version": 3,
        "source_video_id": SYNTHETIC_VIDEO_ID,
        "source_video_path": os.fspath(source_video),
        "unique_frame_count": 1,
        "slide_region_method": "none",
        "slide_region_applied": False,
        "slide_region_verified": False,
        "review_required": True,
        "review_reason": "Context only",
        "artifacts": [
            {
                "path": os.fspath(context),
                "artifact_scope": "full_frame_context",
                "page_count": 1,
                "source_video_id": SYNTHETIC_VIDEO_ID,
                "source_video_path": os.fspath(source_video),
                "crop_verified": False,
                "trusted_for_authored_slide_analysis": False,
            }
        ],
    }
    for method_name in (
        "is_file",
        "resolve",
        "stat",
        "lstat",
        "open",
        "read_bytes",
    ):
        original = getattr(Path, method_name)

        def guarded(
            path: Path,
            *args: Any,
            _original=original,
            _method_name=method_name,
            **kwargs: Any,
        ) -> Any:
            if path.suffix.casefold() == ".pdf":
                pytest.fail(f"reanalysis called {_method_name} for a PDF artifact")
            return _original(path, *args, **kwargs)

        monkeypatch.setattr(Path, method_name, guarded)
    original_hash = pattern_evidence._sha256_file

    def guarded_hash(path: Path) -> str:
        if path.suffix.casefold() == ".pdf":
            pytest.fail("reanalysis hashed a PDF artifact in the owner")
        return original_hash(path)

    monkeypatch.setattr(pattern_evidence, "_sha256_file", guarded_hash)

    probe, reason, path = pattern_evidence._trusted_video_slide_probe(
        vault,
        {"structured_data": {"video_extraction": manifest}},
        SYNTHETIC_VIDEO_ID,
    )

    assert probe is None
    assert path is None
    assert reason == "video slide-region provenance is not trusted"


def test_pptx_preserves_a_separately_declared_rendered_pdf(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    deck = vault / "decks" / "talk.pptx"
    rendered = vault / "slides" / "talk.pdf"
    _write_pptx(deck, slide_count=2)
    _write_pdf(rendered, page_count=2)
    talk = {
        "pptx_path": deck.relative_to(vault).as_posix(),
        "slides_local_path": rendered.relative_to(vault).as_posix(),
        "slide_source": "pptx",
    }

    context = pattern_evidence.build_evidence_context(vault, talk)

    assert context["verified_evidence_sources"] == {"native_deck", "static_slides"}
    assert context["absence_complete_evidence_sources"] == {"static_slides"}
    assert (
        context["slide_artifact_identities"]["native_deck"]["artifact_sha256"]
        != context["slide_artifact_identities"]["static_slides"]["artifact_sha256"]
    )


def test_multiple_exact_comparison_groups_are_preserved_independently() -> None:
    sources = {"transcript", "static_slides", "native_deck"}
    context = {
        "verified_evidence_sources": sources,
        "transcript_line_count": 4,
        "transcript_artifact_identity": _identity("transcript.txt"),
        "slide_counts": {"static_slides": 4, "native_deck": 4},
        "slide_artifact_identities": {
            "static_slides": _identity("slides.pdf"),
            "native_deck": _identity("deck.pptx"),
        },
    }
    inspection = [
        {"source": "transcript", "line_ranges": [[1, 4]]},
        {"source": "static_slides", "page_ranges": [[1, 4]]},
        {"source": "native_deck", "page_ranges": [[1, 4]]},
        {
            "source": "source_comparison",
            "evidence_sources_used": ["transcript", "static_slides"],
            "comparison_scope": "full",
        },
        {
            "source": "source_comparison",
            "evidence_sources_used": ["static_slides", "native_deck"],
            "comparison_scope": "full",
        },
    ]

    records, complete_sources, groups = pattern_evidence.canonicalize_source_inspection(
        inspection, context
    )
    comparisons = [
        record for record in records if record["source"] == "source_comparison"
    ]
    assert len(comparisons) == 2
    assert all(record["coverage_complete"] is True for record in comparisons)
    assert all(
        record["absence_capability_complete"] is False
        and record["absence_capability_reason"] == "comparison_alignment_unverified"
        for record in comparisons
    )
    assert "source_comparison" not in complete_sources
    assert groups == set()

    duplicate = copy.deepcopy(inspection)
    duplicate[-1] = {
        "source": "source_comparison",
        "evidence_sources_used": ["static_slides", "transcript"],
        "comparison_scope": "partial",
    }
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="duplicate source_comparison group",
    ):
        pattern_evidence.canonicalize_source_inspection(duplicate, context)


def test_raw_return_and_canonical_projection_have_the_same_worker_claim(
    tmp_path: Path,
) -> None:
    _, raw, persisted = _canonical_transcript_talk(tmp_path)
    canonical = {
        key: copy.deepcopy(value)
        for key, value in persisted.items()
        if key not in {"source_identity", "title", "transcript_path", "youtube_id"}
    }

    assert pattern_evidence.return_evidence_claim(canonical) == (
        pattern_evidence.return_evidence_claim(raw)
    )
    assert "evidence_schema_version" not in raw["pattern_observations"]
    raw_citation = raw["pattern_observations"]["patterns_detected"][0][
        "evidence_citations"
    ][0]
    assert "line_start" not in raw_citation
    assert "artifact_sha256" not in raw_citation

    observations = canonical["pattern_observations"]
    assert observations["evidence_schema_version"] == 1
    citation = observations["patterns_detected"][0]["evidence_citations"][0]
    assert citation["line_start"] == 1
    assert citation["start_seconds"] == 2.0
    assert citation["artifact_root"] == "vault"
    assert citation["quality_artifact_path"] == (
        f"transcripts/{SYNTHETIC_VIDEO_ID}.quality.json"
    )
    metadata = observations["patterns_detected"][0]["evidence_citations"][1]
    assert metadata["value"] == "Synthetic Talk"
    assert metadata["owner_value_after_return"] == "Synthetic Talk"


@pytest.mark.parametrize("language_owner", ["preclaim", "return"])
def test_non_english_transcript_citations_require_an_english_translation(
    tmp_path: Path,
    language_owner: str,
) -> None:
    vault = tmp_path / "vault"
    transcript, lines = _write_transcript(vault)
    talk: dict[str, Any] = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "slide_source": "none",
    }
    raw = _raw_transcript_return(len(lines), timed=False)
    if language_owner == "preclaim":
        talk["delivery_language"] = "ru"
    else:
        raw["structured_data"] = {"delivery_language": "ru"}

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="translation",
    ):
        pattern_evidence.canonicalize_return_evidence(
            raw, talk, vault, _catalog(_entry())
        )

    citation = raw["pattern_observations"]["patterns_detected"][0][
        "evidence_citations"
    ][0]
    citation["translation"] = "An English rendering of the cited source line."
    canonical = pattern_evidence.canonicalize_return_evidence(
        raw, talk, vault, _catalog(_entry())
    )
    assert (
        canonical["pattern_observations"]["patterns_detected"][0]["evidence_citations"][
            0
        ]["translation"]
        == citation["translation"]
    )


def test_relative_pptx_native_deck_evidence_is_immediately_fresh(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    vault = tmp_path / "vault"
    source_root = tmp_path / "pptx-source"
    deck = source_root / "deck.pptx"
    _write_pptx(deck, slide_count=3)
    talk: dict[str, Any] = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "pptx_path": deck.name,
        "slide_source": "pptx",
    }
    native_audit = pptx_evidence.recompute_native_deck_audit(
        deck,
        trusted_root=source_root,
    )
    raw: dict[str, Any] = {
        "filename": "synthetic-talk.md",
        "return_schema_version": 4,
        "status": "processed",
        "slide_source": "pptx",
        "structured_data": {
            "slide_count": 3,
            "native_deck_audit": native_audit,
        },
        "pattern_observations": {
            "evidence_sources": ["native_deck"],
            "source_inspection": [{"source": "native_deck", "page_ranges": [[1, 3]]}],
            "patterns_detected": [
                {
                    "pattern_id": "slide-pattern",
                    "confidence": "moderate",
                    "evidence_source": "native_deck",
                    "evidence": "The inspected slide establishes the pattern.",
                    "evidence_citations": [
                        {
                            "source": "native_deck",
                            "channel": "slides",
                            "slide_numbers": [2],
                        }
                    ],
                }
            ],
            "antipatterns_detected": [],
            "not_evaluable": [],
            "pattern_score": 1,
        },
    }
    entry = _entry(
        "slide-pattern",
        absence_gate=(frozenset({"native_deck"}),),
        channels=frozenset({"slides"}),
    )
    roots = {"pptx_source_dir": str(source_root)}
    original_hash = pattern_evidence._sha256_file

    def no_parent_pptx_hash(path: Path) -> str:
        if path.suffix.casefold() == ".pptx":
            raise AssertionError("native PPTX bytes were hashed in the parent")
        return original_hash(path)

    monkeypatch.setattr(pattern_evidence, "_sha256_file", no_parent_pptx_hash)
    canonical = pattern_evidence.canonicalize_return_evidence(
        raw, talk, vault, _catalog(entry), source_roots=roots
    )
    persisted = copy.deepcopy(talk)
    persisted.update(copy.deepcopy(canonical))
    inspection = persisted["pattern_observations"]["source_inspection"][0]
    assert inspection["artifact_root"] == "pptx_source"
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted,
            vault_root=vault,
            source_roots=roots,
        )
        == ()
    )

    foreign_owner = copy.deepcopy(persisted)
    foreign_owner["pptx_path"] = _foreign_absolute_locator("deck.pptx")
    foreign_reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        foreign_owner,
        vault_root=vault,
        source_roots=roots,
    )
    assert any(
        reason.endswith(":artifact_owner_path_drift") for reason in foreign_reasons
    )

    invalid_root_reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        persisted,
        vault_root=vault,
        source_roots={"pptx_source_dir": "relative-presentations"},
    )
    assert any(
        reason.endswith(":artifact_root_invalid:artifact_root_not_native_absolute")
        for reason in invalid_root_reasons
    )

    missing_audit = copy.deepcopy(persisted)
    missing_audit["structured_data"].pop("native_deck_audit")
    assert "native_deck_audit_missing" in (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            missing_audit,
            vault_root=vault,
            source_roots=roots,
        )
    )

    old_generation = copy.deepcopy(persisted)
    old_generation["structured_data"]["native_deck_audit"][
        "extraction_pipeline_version"
    ] = "1.4.0"
    assert "native_deck_audit_generation_drift" in (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            old_generation,
            vault_root=vault,
            source_roots=roots,
        )
    )

    wrong_lane = copy.deepcopy(persisted)
    wrong_lane["slide_source"] = "none"
    assert "native_deck_audit_unexpected" in (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            wrong_lane,
            vault_root=vault,
            source_roots=roots,
        )
    )

    _write_crc_damaged_media_pptx(deck)
    pptx_evidence.clear_pptx_artifact_probe_cache()
    degraded_reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        persisted,
        vault_root=vault,
        source_roots=roots,
    )
    assert any(
        reason.endswith(":artifact_bounded_digest_unavailable")
        for reason in degraded_reasons
    )

    deck.write_bytes(b"not a PPTX")
    pptx_evidence.clear_pptx_artifact_probe_cache()
    unavailable_reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        persisted,
        vault_root=vault,
        source_roots=roots,
    )
    assert any(
        reason.endswith(":artifact_bounded_digest_unavailable")
        for reason in unavailable_reasons
    )


@pytest.mark.parametrize("absolute_path", [False, True], ids=["relative", "absolute"])
@pytest.mark.parametrize(
    "symlink_root",
    [False, True],
    ids=["storage-root", "canonical-symlink-root"],
)
def test_static_pdf_freshness_never_touches_the_leaf_in_parent(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    absolute_path: bool,
    symlink_root: bool,
) -> None:
    storage = tmp_path / "vault-storage"
    vault = storage
    if symlink_root:
        vault = tmp_path / "canonical-vault"
        storage.mkdir()
        try:
            vault.symlink_to(storage, target_is_directory=True)
        except OSError:
            pytest.skip("directory policy does not permit symlink creation")
    rendered = storage / "slides" / "talk.pdf"
    _write_pdf(rendered, page_count=3)
    talk: dict[str, Any] = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "slides_local_path": (
            os.fspath(rendered)
            if absolute_path
            else rendered.relative_to(storage).as_posix()
        ),
        "slide_source": "pdf",
    }
    raw = _raw_static_pdf_return(3)
    entry = _entry(
        "slide-pattern",
        absence_gate=(frozenset({"static_slides"}),),
        channels=frozenset({"slides"}),
    )
    canonical = pattern_evidence.canonicalize_return_evidence(
        raw,
        talk,
        vault,
        _catalog(entry),
    )
    persisted = copy.deepcopy(talk)
    persisted.update(copy.deepcopy(canonical))
    for method_name in (
        "is_file",
        "resolve",
        "stat",
        "lstat",
        "open",
        "read_bytes",
    ):
        original = getattr(Path, method_name)

        def guarded(
            path: Path,
            *args: Any,
            _original=original,
            _method_name=method_name,
            **kwargs: Any,
        ) -> Any:
            if path.suffix.casefold() == ".pdf":
                pytest.fail(f"freshness called {_method_name} for a PDF artifact")
            return _original(path, *args, **kwargs)

        monkeypatch.setattr(Path, method_name, guarded)

    original_hash = pattern_evidence._sha256_file

    def guarded_hash(path: Path) -> str:
        if path.suffix.casefold() == ".pdf":
            pytest.fail("freshness hashed a PDF artifact in the owner process")
        return original_hash(path)

    monkeypatch.setattr(pattern_evidence, "_sha256_file", guarded_hash)

    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted,
            vault_root=vault,
        )
        == ()
    )


def test_external_static_pdf_symlink_parent_stays_fresh(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    storage = tmp_path / "external-storage"
    locator = tmp_path / "configured-slides"
    storage.mkdir()
    try:
        locator.symlink_to(storage, target_is_directory=True)
    except OSError as exc:  # pragma: no cover - unusual restricted platform
        pytest.skip(f"symlinks unavailable: {exc}")
    rendered = storage / "talk.pdf"
    _write_pdf(rendered, page_count=3)
    talk: dict[str, Any] = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "slides_local_path": os.fspath(locator / rendered.name),
        "slide_source": "pdf",
    }
    canonical = pattern_evidence.canonicalize_return_evidence(
        _raw_static_pdf_return(3),
        talk,
        vault,
        _catalog(
            _entry(
                "slide-pattern",
                absence_gate=(frozenset({"static_slides"}),),
                channels=frozenset({"slides"}),
            )
        ),
    )
    persisted = copy.deepcopy(talk)
    persisted.update(copy.deepcopy(canonical))

    inspection = persisted["pattern_observations"]["source_inspection"][0]
    assert inspection["artifact_root"] == "preclaim:slides_local_path"
    assert inspection["artifact_path"] == rendered.name
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted,
            vault_root=vault,
        )
        == ()
    )


def test_freshness_detects_artifact_digest_drift(tmp_path: Path) -> None:
    vault, _, persisted = _canonical_transcript_talk(tmp_path)
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted, vault_root=vault
        )
        == ()
    )

    transcript = vault / str(persisted["transcript_path"])
    transcript.write_text(
        transcript.read_text(encoding="utf-8").replace("uniquely", "clearly", 1),
        encoding="utf-8",
    )
    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        persisted, vault_root=vault
    )
    assert any(reason.endswith(":artifact_digest_mismatch") for reason in reasons)


def test_freshness_detects_timing_digest_drift(tmp_path: Path) -> None:
    vault, _, persisted = _canonical_transcript_talk(tmp_path)
    timing = (vault / str(persisted["transcript_path"])).with_suffix(".segments.json")
    timing.write_text(
        timing.read_text(encoding="utf-8") + " ",
        encoding="utf-8",
    )
    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        persisted, vault_root=vault
    )
    assert any(
        reason.endswith(":timing_artifact_digest_mismatch") for reason in reasons
    )


def test_freshness_detects_quality_receipt_digest_drift(tmp_path: Path) -> None:
    vault, _, persisted = _canonical_transcript_talk(tmp_path)
    receipt = (vault / str(persisted["transcript_path"])).with_suffix(".quality.json")
    receipt.write_text(
        receipt.read_text(encoding="utf-8") + " ",
        encoding="utf-8",
    )

    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        persisted, vault_root=vault
    )
    assert any(
        reason.endswith(":quality_artifact_digest_mismatch") for reason in reasons
    )


def test_current_transcript_evidence_requires_a_hash_current_quality_receipt(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, lines = _write_transcript(vault, timed=False)
    transcript_timing.quality_sidecar_path(transcript).unlink()
    talk = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "slide_source": "none",
    }

    context = pattern_evidence.build_evidence_context(vault, talk)
    assert "transcript" not in context["verified_evidence_sources"]
    assert "requires a hash-current transcript quality receipt" in str(
        context["transcript_reason"]
    )
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="claims unavailable source 'transcript'",
    ):
        pattern_evidence.canonicalize_return_evidence(
            _raw_transcript_return(len(lines), timed=False),
            talk,
            vault,
            _catalog(_entry()),
        )


def test_fixed_default_receipt_is_valid_for_an_existing_youtube_artifact(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    youtube_id = "abcdefghijk"
    transcript, _ = _write_transcript(
        vault, name=youtube_id, timed=False, quality_duration=None
    )
    talk = {
        "filename": "synthetic-talk.md",
        "youtube_id": youtube_id,
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "slide_source": "none",
    }

    context = pattern_evidence.build_evidence_context(vault, talk)

    assert "transcript" in context["verified_evidence_sources"]
    assert "verified transcript quality receipt" in str(context["transcript_reason"])


def test_local_media_quality_provenance_binds_exact_video_bytes(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    video = vault / "videos" / "talk.mp4"
    video.parent.mkdir(parents=True)
    video.write_bytes(b"source-owned synthetic video bytes")
    text = transcript.read_text(encoding="utf-8")
    duration = 600.0
    transcript_timing.write_quality_receipt(
        transcript,
        text,
        transcript_timing.build_quality_policy(400, trusted_duration_seconds=duration),
        {
            "kind": "local_media_duration",
            "media_sha256": hashlib.sha256(video.read_bytes()).hexdigest(),
            "duration_seconds": duration,
        },
    )
    talk = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "slide_source": "none",
        "video_local_path": video.relative_to(vault).as_posix(),
    }

    context = pattern_evidence.build_evidence_context(
        vault,
        talk,
        video_evidence_assessment=_TestVideoAssessment(duration=duration),
    )
    assert "transcript" in context["verified_evidence_sources"]
    assert context["quality_artifact_identity"]["quality_artifact_path"] == (
        "transcripts/talk.quality.json"
    )

    video.write_bytes(b"different video bytes")
    drifted = pattern_evidence.build_evidence_context(
        vault,
        talk,
        video_evidence_assessment=_TestVideoAssessment(duration=duration),
    )
    assert "transcript" not in drifted["verified_evidence_sources"]
    assert "digest does not match" in str(drifted["transcript_reason"])


def test_transcript_context_retries_a_concurrent_bundle_replacement(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(
        vault,
        name=SYNTHETIC_VIDEO_ID,
        timed=True,
    )
    original_text = transcript.read_text(encoding="utf-8")
    replacement_text = original_text.replace("uniquely", "distinctly", 1)
    talk = {
        "filename": "synthetic-talk.md",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "whisper",
        "youtube_id": SYNTHETIC_VIDEO_ID,
        "source_identity": {
            "schema_version": 1,
            "provider": "youtube",
            "video_id": SYNTHETIC_VIDEO_ID,
            "duration_seconds": SYNTHETIC_DURATION,
        },
        "slide_source": "none",
    }
    original_loader = pattern_evidence.load_verified_quality_receipt
    replacement_count = 0

    def replace_bundle_then_load(
        transcript_path: Path,
        text: str,
    ) -> tuple[dict[str, object] | None, str]:
        nonlocal replacement_count
        if replacement_count == 0:
            replacement_count += 1
            transcript_timing.write_transcript_bundle(
                transcript,
                replacement_text,
                [{"text": replacement_text, "start": 2.0, "end": 6.0}],
                source="whisper",
                timing_provenance=transcript_timing.youtube_timing_provenance(
                    "whisper", SYNTHETIC_VIDEO_ID, SYNTHETIC_DURATION
                ),
                quality_policy=transcript_timing.build_quality_policy(400),
                quality_policy_provenance={"kind": "fixed_default"},
                force=True,
            )
        return original_loader(transcript_path, text)

    monkeypatch.setattr(
        pattern_evidence,
        "load_verified_quality_receipt",
        replace_bundle_then_load,
    )

    context = pattern_evidence.build_evidence_context(vault, talk)

    assert replacement_count == 1
    assert context["transcript_text"] == replacement_text
    assert context["timed_segments"] == [
        {
            "text": replacement_text,
            "start_seconds": 2.0,
            "end_seconds": 6.0,
        }
    ]
    assert context["transcript_artifact_identity"]["artifact_sha256"] == (
        hashlib.sha256(transcript.read_bytes()).hexdigest()
    )
    timing_path = transcript_timing.sidecar_path(transcript)
    assert context["timing_artifact_identity"]["timing_artifact_sha256"] == (
        hashlib.sha256(timing_path.read_bytes()).hexdigest()
    )
    quality_path = transcript_timing.quality_sidecar_path(transcript)
    assert context["quality_artifact_identity"]["quality_artifact_sha256"] == (
        hashlib.sha256(quality_path.read_bytes()).hexdigest()
    )
    assert context["verified_evidence_sources"] == {"transcript"}


def test_freshness_detects_metadata_and_owner_path_drift(tmp_path: Path) -> None:
    vault, _, persisted = _canonical_transcript_talk(tmp_path)
    metadata_drift = copy.deepcopy(persisted)
    metadata_drift["title"] = "Changed Owner Title"
    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        metadata_drift, vault_root=vault
    )
    assert any(reason.endswith(":metadata_value_drift") for reason in reasons)

    owner_drift = copy.deepcopy(persisted)
    replacement, _ = _write_transcript(vault, name="replacement")
    owner_drift["transcript_path"] = replacement.relative_to(vault).as_posix()
    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        owner_drift, vault_root=vault
    )
    assert any(reason.endswith(":artifact_owner_path_drift") for reason in reasons)


def test_freshness_detects_transcript_quality_context_drift(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    youtube_id = "abcdefghijk"
    transcript, lines = _write_transcript(
        vault,
        name=youtube_id,
        timed=True,
        quality_duration=600.0,
    )
    talk: dict[str, Any] = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "youtube_id": youtube_id,
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "whisper",
        "slide_source": "none",
        "source_identity": {
            "schema_version": 1,
            "provider": "youtube",
            "video_id": youtube_id,
            "duration_seconds": 600.0,
        },
    }
    canonical = pattern_evidence.canonicalize_return_evidence(
        _raw_transcript_return(len(lines)),
        talk,
        vault,
        _catalog(_entry()),
    )
    persisted = copy.deepcopy(talk)
    persisted.update(copy.deepcopy(canonical))
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted, vault_root=vault
        )
        == ()
    )

    drifted = copy.deepcopy(persisted)
    # Two seconds is well inside the retired 60s/5% tolerance and outside the
    # acquisition contract's one-second integer/probe allowance.
    drifted["source_identity"]["duration_seconds"] = 602.0
    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        drifted, vault_root=vault
    )
    assert "source_inspection[0]:transcript_quality_context_drift" in reasons


def test_v4_is_source_located_while_legacy_is_explicitly_unverified(
    tmp_path: Path,
) -> None:
    _, _, persisted = _canonical_transcript_talk(tmp_path)
    assert persisted["pattern_observations"]["evidence_schema_version"] == 1

    legacy = copy.deepcopy(persisted)
    legacy["return_schema_version"] = 3
    legacy["pattern_observations"].pop("evidence_schema_version")
    legacy["pattern_observations"].pop("source_inspection")
    for detection in legacy["pattern_observations"]["patterns_detected"]:
        detection["evidence_citations"] = []
    assert pattern_evidence.assess_persisted_pattern_evidence_freshness(
        legacy, vault_root=tmp_path
    ) == ("evidence_schema_unverified",)


def test_transcript_resolution_uses_only_preclaim_owner_identity(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    owner_path, _ = _write_transcript(vault, name="owner")
    redirected_path, _ = _write_transcript(vault, name="redirected")
    talk = {"transcript_path": owner_path.relative_to(vault).as_posix()}

    resolved, _ = pattern_evidence.resolve_transcript_artifact(
        vault,
        talk,
        {"transcript_path": redirected_path.relative_to(vault).as_posix()},
    )
    assert resolved == owner_path.resolve()

    unregistered, _ = pattern_evidence.resolve_transcript_artifact(
        vault,
        {},
        {"transcript_path": redirected_path.relative_to(vault).as_posix()},
    )
    assert unregistered is None

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="does not match the claimed talk's youtube_id",
    ):
        pattern_evidence.resolve_transcript_artifact(
            vault,
            {
                "youtube_id": "abcdefghijk",
                "transcript_path": owner_path.relative_to(vault).as_posix(),
            },
        )


def _validation_return(
    *,
    detections: list[dict[str, Any]],
    not_evaluable: list[dict[str, Any]],
) -> dict[str, Any]:
    return {
        "filename": "synthetic-talk.md",
        "return_schema_version": 4,
        "queue_claim": {
            "run_id": "synthetic-run",
            "batch_id": "batch-1",
            "reprocess_generation": 1,
        },
        "status": "processed_partial",
        "slide_source": "none",
        "transcript_source": "manual",
        "rhetoric_notes": "Synthetic rhetoric analysis.",
        "areas_for_improvement": "Synthetic improvement analysis.",
        "adherence_assessment": "",
        "new_patterns": "",
        "summary_updates": "",
        "structured_data": {},
        "verbatim_examples": {},
        "pattern_observations": {
            "evidence_sources": ["transcript"],
            "source_inspection": [{"source": "transcript", "line_ranges": [[1, 1]]}],
            "patterns_detected": detections,
            "antipatterns_detected": [],
            "not_evaluable": not_evaluable,
            "pattern_score": {
                "patterns_used": len(detections),
                "antipatterns_detected": 0,
                "score": len(detections),
            },
        },
        "catalog_feedback": {
            "unmatched_observations": [],
            "confusable_pairs": [],
            "definition_problems": [],
            "scoring_problems": [],
            "tensions": [],
        },
    }


def _validation_catalog(*entries: SimpleNamespace):
    return return_validation.PatternCatalog(
        entries={
            entry.pattern_id: return_validation.CatalogEntry(**vars(entry))
            for entry in entries
        },
        fingerprint="f" * 64,
    )


def _transcript_detection(
    *,
    pattern_id: str = "synthetic-pattern",
    citations: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    detection: dict[str, Any] = {
        "pattern_id": pattern_id,
        "confidence": "moderate",
        "evidence_source": "transcript",
        "evidence": "The exact source contains the synthetic opening.",
    }
    if citations is not None:
        detection["evidence_citations"] = citations
    return detection


def test_v4_detection_requires_source_located_citations() -> None:
    ret = _validation_return(
        detections=[_transcript_detection()],
        not_evaluable=[],
    )

    with pytest.raises(
        return_validation.ReturnValidationError,
        match="evidence_citations must be a non-empty array",
    ):
        return_validation.validate_return(ret, _validation_catalog(_entry()))


@pytest.mark.parametrize(
    ("unknown_field", "unknown_value"),
    [
        ("model_reasoning", "trust me"),
        ("line_start", 1),
        ("artifact_sha256", "a" * 64),
        ("quality_artifact_path", "transcripts/talk.quality.json"),
    ],
)
def test_v4_raw_citation_rejects_unknown_and_engine_owned_fields(
    unknown_field: str,
    unknown_value: object,
) -> None:
    citation: dict[str, Any] = {
        "source": "transcript",
        "channel": "transcript",
        "quote": _transcript_lines()[0],
        unknown_field: unknown_value,
    }
    ret = _validation_return(
        detections=[_transcript_detection(citations=[citation])],
        not_evaluable=[],
    )

    with pytest.raises(
        return_validation.ReturnValidationError,
        match="unknown fields",
    ):
        return_validation.validate_return(ret, _validation_catalog(_entry()))


def test_timed_transcript_detection_fails_closed_without_timing_sidecar(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, lines = _write_transcript(vault, timed=False)
    talk = {
        "filename": "synthetic-talk.md",
        "title": "Synthetic Talk",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "slide_source": "none",
    }

    raw = _raw_transcript_return(len(lines), timed=True)
    raw["transcript_source"] = "manual"
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="no verified timestamp",
    ):
        pattern_evidence.canonicalize_return_evidence(
            raw,
            talk,
            vault,
            _catalog(_entry()),
        )


def test_slide_citation_requires_a_verified_predeclared_local_artifact() -> None:
    detection = {
        "pattern_id": "slide-pattern",
        "confidence": "moderate",
        "evidence_source": "static_slides",
        "evidence": "The inspected slides establish the pattern.",
        "evidence_citations": [
            {
                "source": "static_slides",
                "channel": "slides",
                "slide_numbers": [1, 2],
            }
        ],
    }

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="predeclared local artifact",
    ):
        pattern_evidence.canonicalize_detection_citations(
            detection,
            evidence_channels=frozenset({"slides"}),
            evidence_metadata_fields=frozenset(),
            context={
                "slide_counts": {},
                "source_reasons": {
                    "static_slides": "no verified local slide artifact",
                },
            },
        )


def test_v4_video_citation_rejects_nonfinite_ranges() -> None:
    video_entry = _entry(
        "video-pattern",
        absence_gate=(frozenset({"delivery_video"}),),
        channels=frozenset({"video"}),
    )
    detection = {
        "pattern_id": video_entry.pattern_id,
        "confidence": "moderate",
        "evidence_source": "delivery_video",
        "evidence": "The delivery includes one observable physical action.",
        "evidence_citations": [
            {
                "source": "delivery_video",
                "channel": "video",
                "start_seconds": 2.0,
                "end_seconds": float("inf"),
            }
        ],
    }
    ret = _validation_return(detections=[detection], not_evaluable=[])
    ret["pattern_observations"]["evidence_sources"] = ["delivery_video"]
    ret["pattern_observations"]["source_inspection"] = [
        {
            "source": "delivery_video",
            "time_ranges": [[0.0, 10.0]],
        }
    ]

    with pytest.raises(
        return_validation.ReturnValidationError,
        match="finite non-negative start_seconds/end_seconds",
    ):
        return_validation.validate_return(ret, _validation_catalog(video_entry))


def test_catalog_rejects_duplicate_evidence_channels(tmp_path: Path) -> None:
    catalog_root = tmp_path / "patterns"
    category = catalog_root / "build"
    category.mkdir(parents=True)
    (catalog_root / "_index.md").write_text("# Synthetic catalog\n", encoding="utf-8")
    (category / "duplicate.md").write_text(
        """---
id: duplicate
type: pattern
observable: true
evidence_channels: [video, video]
vault_dimensions: [1]
---
# Duplicate
""",
        encoding="utf-8",
    )

    with pytest.raises(
        return_validation.ReturnValidationError,
        match="duplicate-free evidence_channels",
    ):
        return_validation.load_catalog(catalog_root)


@pytest.mark.parametrize(
    ("field", "message"),
    [
        ("rhetoric_notes", "not immutable source metadata"),
        ("transcript_path", "not permitted for this pattern"),
    ],
)
def test_talk_metadata_rejects_generated_and_irrelevant_fields(
    field: str,
    message: str,
) -> None:
    detection = _transcript_detection(
        citations=[
            {
                "source": "transcript",
                "channel": "talk_metadata",
                "field": field,
            }
        ]
    )
    context = {
        "metadata": {
            "title": "Synthetic Talk",
            "transcript_path": "transcripts/talk.txt",
        },
        "post_return_metadata": {},
    }

    with pytest.raises(pattern_evidence.PatternEvidenceError, match=message):
        pattern_evidence.canonicalize_detection_citations(
            detection,
            evidence_channels=frozenset({"transcript", "talk_metadata"}),
            evidence_metadata_fields=frozenset({"title"}),
            context=context,
        )


def test_talk_metadata_values_are_stamped_from_owner_state() -> None:
    quote = _transcript_lines()[0]
    detection = _transcript_detection(
        citations=[
            {
                "source": "transcript",
                "channel": "transcript",
                "quote": quote,
            },
            {
                "source": "transcript",
                "channel": "talk_metadata",
                "field": "slide_count",
            },
        ]
    )
    context = {
        "transcript_text": quote,
        "timed_segments": [],
        "transcript_artifact_identity": _identity("talk.txt"),
        "quality_artifact_identity": {
            "quality_artifact_root": "vault",
            "quality_artifact_path": "synthetic/talk.quality.json",
            "quality_artifact_sha256": "b" * 64,
        },
        "metadata": {"slide_count": 62},
        "post_return_metadata": {"slide_count": 63},
        "delivery_language": "en",
    }

    canonical = pattern_evidence.canonicalize_detection_citations(
        detection,
        evidence_channels=frozenset({"transcript", "talk_metadata"}),
        evidence_metadata_fields=frozenset({"slide_count"}),
        context=context,
    )
    metadata = canonical["evidence_citations"][1]
    assert metadata["value"] == 62
    assert metadata["owner_value_after_return"] == 63


def test_pending_catalog_gate_fails_closed_and_uses_reason_code_only() -> None:
    pending = _entry("pending-pattern", absence_gate=None)
    catalog = return_validation.PatternCatalog(
        entries={pending.pattern_id: return_validation.CatalogEntry(**vars(pending))},
        fingerprint="f" * 64,
    )
    detection = {
        "pattern_id": pending.pattern_id,
        "confidence": "moderate",
        "evidence_source": "transcript",
        "evidence": "Synthetic evidence.",
        "evidence_citations": [
            {
                "source": "transcript",
                "channel": "transcript",
                "quote": "A uniquely phrased synthetic opening proves the source.",
            }
        ],
    }
    with pytest.raises(
        return_validation.ReturnValidationError,
        match="no owner-approved source gate",
    ):
        return_validation.validate_return(
            _validation_return(detections=[detection], not_evaluable=[]),
            catalog,
        )

    valid = _validation_return(
        detections=[],
        not_evaluable=[
            {
                "pattern_id": pending.pattern_id,
                "reason_code": "source_gate_pending_owner_review",
            }
        ],
    )
    return_validation.validate_return(valid, catalog)

    invalid = copy.deepcopy(valid)
    invalid["pattern_observations"]["not_evaluable"][0]["reason"] = (
        "Free prose is not authority."
    )
    with pytest.raises(
        return_validation.ReturnValidationError,
        match="free-prose waivers are not scoring authority",
    ):
        return_validation.validate_return(invalid, catalog)


def test_batch_capability_preflight_does_not_touch_unrelated_talks(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    calls: list[str] = []

    def assess(
        talk: dict[str, object],
        *,
        vault_root: str | Path,
        source_roots: dict[str, object] | None = None,
        video_evidence_assessment: object | None = None,
    ) -> dict[str, object]:
        del vault_root, source_roots, video_evidence_assessment
        filename = str(talk["filename"])
        calls.append(filename)
        if filename == "unrelated.md":
            raise AssertionError("unrelated expensive artifact was inspected")
        return {
            "verified_capabilities": (),
            "verified_evidence_sources": (),
            "acquisition_capabilities": (),
            "source_reasons": {},
        }

    monkeypatch.setattr(pattern_evidence, "assess_talk_artifact_capabilities", assess)
    result = pattern_evidence.assess_batch_artifact_capabilities(
        [
            {"filename": "batch-member.md"},
            {"filename": "unrelated.md", "video_local_path": "huge.mp4"},
        ],
        {"batch-member.md"},
        vault_root=tmp_path,
    )

    assert calls == ["batch-member.md"]
    assert set(result) == {"batch-member.md"}


def test_invalid_transcript_cannot_hide_an_independent_valid_deck(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    deck = vault / "decks" / "talk.pptx"
    _write_pptx(deck)
    talk = {
        "filename": "talk.md",
        "transcript_path": "../bad.txt",
        "pptx_path": deck.relative_to(vault).as_posix(),
        "slide_source": "pptx",
    }

    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        talk, vault_root=vault
    )

    assert assessment["verified_capabilities"] == ("slides",)
    assert assessment["verified_evidence_sources"] == ("native_deck",)
    assert assessment["repair_capabilities"] == ()
    assert "transcripts/<artifact>" in assessment["source_reasons"]["transcript"]


def test_invalid_utf8_transcript_cannot_hide_an_independent_valid_deck(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript = vault / "transcripts" / "bad.txt"
    transcript.parent.mkdir(parents=True)
    transcript.write_bytes(b"valid prefix\xffinvalid utf-8")
    deck = vault / "decks" / "talk.pptx"
    _write_pptx(deck)
    talk = {
        "filename": "talk.md",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "pptx_path": deck.relative_to(vault).as_posix(),
        "slide_source": "pptx",
    }

    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        talk, vault_root=vault
    )

    assert assessment["verified_capabilities"] == ("slides",)
    assert assessment["verified_evidence_sources"] == ("native_deck",)
    assert assessment["repair_capabilities"] == ()
    assert "not valid UTF-8" in assessment["source_reasons"]["transcript"]


def test_crc_damaged_pptx_is_a_structured_degraded_capability(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    deck = vault / "decks" / "talk.pptx"
    damaged_part = _write_crc_damaged_media_pptx(deck)

    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slide_source": "pptx",
        },
        vault_root=vault,
    )

    assert assessment["verified_capabilities"] == ()
    assert assessment["verified_evidence_sources"] == ()
    degradation = assessment["degraded_evidence_sources"]["native_deck"]
    assert degradation["schema_version"] == 1
    assert degradation["status"] == "degraded_recoverable"
    assert degradation["reason_code"] == "pptx_archive_recovery_required"
    assert degradation["archive_recovery"][0]["part_name"] == damaged_part
    assert "read degraded local PPTX" in assessment["source_reasons"]["native_deck"]
    json.dumps(assessment)


@pytest.mark.parametrize("slide_source", [None, "none"])
def test_degraded_deck_is_not_native_evidence_when_slide_source_is_absent(
    tmp_path: Path,
    slide_source: str | None,
) -> None:
    vault = tmp_path / "vault"
    deck = vault / "decks" / "talk.pptx"
    _write_crc_damaged_media_pptx(deck)
    talk: dict[str, object] = {
        "filename": "talk.md",
        "pptx_path": deck.relative_to(vault).as_posix(),
    }
    if slide_source is not None:
        talk["slide_source"] = slide_source

    context = pattern_evidence.build_evidence_context(vault, talk)
    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        talk,
        vault_root=vault,
    )

    assert "native_deck" not in context["slide_counts"]
    assert "native_deck" not in context["slide_artifact_identities"]
    assert assessment["verified_capabilities"] == ()
    assert assessment["verified_evidence_sources"] == ()
    assert "native_deck" in assessment["degraded_evidence_sources"]


def test_optional_degraded_deck_keeps_independent_static_slides(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    deck = vault / "decks" / "talk.pptx"
    slides = vault / "slides" / "talk.pdf"
    _write_crc_damaged_media_pptx(deck)
    _write_pdf(slides)

    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slides_local_path": slides.relative_to(vault).as_posix(),
            "slide_source": "pdf",
        },
        vault_root=vault,
    )

    assert assessment["verified_capabilities"] == ("slides",)
    assert assessment["verified_evidence_sources"] == ("static_slides",)
    assert "native_deck" in assessment["degraded_evidence_sources"]


def test_structural_pptx_damage_is_unavailable_without_erasing_transcript(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _lines = _write_transcript(vault, timed=False)
    deck = vault / "decks" / "talk.pptx"
    _write_pptx(deck, slide_count=1)
    with zipfile.ZipFile(deck) as archive:
        member = archive.getinfo("ppt/slides/slide1.xml")
    package = bytearray(deck.read_bytes())
    name_size, extra_size = struct.unpack_from(
        "<HH", package, member.header_offset + 26
    )
    payload_offset = member.header_offset + 30 + name_size + extra_size
    package[payload_offset + (member.compress_size // 2)] ^= 0xFF
    deck.write_bytes(package)

    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slide_source": "pptx",
        },
        vault_root=vault,
    )

    assert assessment["verified_capabilities"] == ("transcript",)
    assert assessment["verified_evidence_sources"] == ("transcript",)
    assert "structural PPTX member" in assessment["source_reasons"]["native_deck"]
    assert assessment["degraded_evidence_sources"] == {}


def test_native_deck_audit_is_recomputed_and_bound_to_canonical_render(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    deck = vault / "decks" / "talk.pptx"
    image_path = vault / "decks" / "full-bleed.png"
    deck.parent.mkdir(parents=True)
    Image.new("RGB", (640, 480), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        Inches(0),
        Inches(0),
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    presentation.save(str(deck))
    rendered = vault / "slides" / "talk.pdf"
    _write_pdf(rendered, page_count=1)
    extraction = pptx_extraction._extract_pptx_in_process(
        deck,
        ocr=False,
        rendered_pdf_path=rendered,
        inspected_page_ranges=[[1, 1]],
    )
    talk = {
        "filename": "talk.md",
        "pptx_path": deck.relative_to(vault).as_posix(),
        "slides_local_path": rendered.relative_to(vault).as_posix(),
        "slide_source": "both",
    }
    raw = {
        "filename": "talk.md",
        "return_schema_version": 4,
        "status": "processed_partial",
        "slide_source": "both",
        "transcript_source": "none",
        "structured_data": {
            "slide_count": 1,
            "native_deck_audit": extraction["native_deck_audit"],
        },
        "pattern_observations": {
            "evidence_sources": ["native_deck", "static_slides"],
            "source_inspection": [
                {"source": "native_deck", "page_ranges": [[1, 1]]},
                {"source": "static_slides", "page_ranges": [[1, 1]]},
            ],
            "patterns_detected": [],
            "antipatterns_detected": [],
            "not_evaluable": [],
            "pattern_score": 0,
        },
    }
    original = copy.deepcopy(raw)
    canonical = pattern_evidence.canonicalize_return_evidence(
        raw,
        talk,
        vault,
        _catalog(),
    )
    assert raw == original
    assert (
        canonical["structured_data"]["native_deck_audit"]
        == (extraction["native_deck_audit"])
    )
    persisted: dict[str, Any] = copy.deepcopy(talk)
    persisted.update(copy.deepcopy(canonical))
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted,
            vault_root=vault,
        )
        == ()
    )

    rendered_inspection_drift: dict[str, Any] = copy.deepcopy(persisted)
    static_inspection = next(
        record
        for record in rendered_inspection_drift["pattern_observations"][
            "source_inspection"
        ]
        if record["source"] == "static_slides"
    )
    static_inspection["page_ranges"] = []
    assert "native_deck_audit_rendered_pdf_inspection_drift" in (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            rendered_inspection_drift,
            vault_root=vault,
        )
    )

    source_drift: dict[str, Any] = copy.deepcopy(persisted)
    current_audit = source_drift["structured_data"]["native_deck_audit"]
    source_drift["structured_data"]["native_deck_audit"] = (
        pptx_evidence.build_native_deck_audit(
            source_pptx_sha256="b" * 64,
            source_pptx_size_bytes=current_audit["source_pptx_size_bytes"],
            slide_count=1,
            render_required_reasons={
                int(slide_number): reasons
                for slide_number, reasons in current_audit[
                    "render_required_reasons"
                ].items()
            },
        )
    )
    assert "native_deck_audit_source_generation_drift" in (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            source_drift,
            vault_root=vault,
        )
    )

    forged = copy.deepcopy(raw)
    forged["structured_data"]["native_deck_audit"] = (
        pptx_evidence.build_native_deck_audit(
            source_pptx_sha256=extraction["input_fingerprint"]["digest"],
            source_pptx_size_bytes=extraction["input_fingerprint"]["size_bytes"],
            slide_count=1,
            render_required_reasons={},
        )
    )
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="disagrees with a fresh extraction",
    ):
        pattern_evidence.canonicalize_return_evidence(
            forged,
            talk,
            vault,
            _catalog(),
        )

    original_render = rendered.read_bytes()
    assert b"pypdf" in original_render
    changed_render = original_render.replace(b"pypdf", b"qypdf", 1)
    assert len(changed_render) == len(original_render)
    rendered.write_bytes(changed_render)
    pdf_evidence.clear_pdf_artifact_probe_cache()
    assert "native_deck_audit_rendered_pdf_generation_drift" in (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted,
            vault_root=vault,
        )
    )
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="rendered_page_inspection.*(canonical static_slides|PDF generation)",
    ):
        pattern_evidence.canonicalize_return_evidence(
            raw,
            talk,
            vault,
            _catalog(),
        )


def test_invalid_deck_cannot_hide_an_independent_valid_transcript(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    talk = {
        "filename": "talk.md",
        "transcript_path": transcript.relative_to(vault).as_posix(),
        "transcript_source": "manual",
        "pptx_path": "../bad.pptx",
        "slide_source": "pptx",
    }

    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        talk, vault_root=vault
    )

    assert assessment["verified_capabilities"] == ("transcript",)
    assert assessment["verified_evidence_sources"] == ("transcript",)
    assert assessment["repair_capabilities"] == ()
    assert assessment["source_reasons"]["native_deck"].endswith(
        "artifact_locator_dot_segment"
    )


def test_unprobeable_video_is_not_hashed_or_allowed_to_hide_other_lanes(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    deck = vault / "decks" / "talk.pptx"
    _write_pptx(deck)
    video = vault / "videos" / "talk.mp4"
    video.parent.mkdir(parents=True)
    video.write_bytes(b"not a probeable video")
    original_hash = pattern_evidence._sha256_file

    def reject_video_hash(path: Path) -> str:
        if path == video:
            raise AssertionError("an unverified video must not be hashed")
        return original_hash(path)

    monkeypatch.setattr(pattern_evidence, "_sha256_file", reject_video_hash)
    video_assessment = _TestVideoAssessment(
        failure=pattern_evidence.VideoEvidenceError(
            "synthetic ffprobe failure",
            reason_code="video_parser_rejected",
        )
    )
    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slide_source": "pptx",
            "video_local_path": video.relative_to(vault).as_posix(),
        },
        vault_root=vault,
        video_evidence_assessment=video_assessment,
    )

    assert assessment["verified_capabilities"] == ("slides", "transcript")
    assert "delivery_video" not in assessment["verified_evidence_sources"]
    assert "synthetic ffprobe failure" in assessment["source_reasons"]["delivery_video"]
    assert assessment["unavailable_evidence_sources"]["delivery_video"] == {
        "schema_version": 1,
        "status": "unavailable",
        "reason_code": "video_parser_rejected",
        "reason": "synthetic ffprobe failure",
        "details": {},
        "artifact_path": str(video),
    }


def test_video_mutation_between_probe_and_digest_rejects_only_video_lane(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    deck = vault / "decks" / "talk.pptx"
    _write_pptx(deck)
    video = vault / "videos" / "talk.mp4"
    video.parent.mkdir(parents=True)
    video.write_bytes(b"first synthetic video generation")
    video_assessment = _TestVideoAssessment(
        failure=pattern_evidence.VideoEvidenceError(
            "synthetic video changed during supervised inspection",
            reason_code="video_artifact_changed",
        )
    )

    context = pattern_evidence.build_evidence_context(
        vault,
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slide_source": "pptx",
            "video_local_path": video.relative_to(vault).as_posix(),
        },
        video_evidence_assessment=video_assessment,
    )

    assert context["verified_evidence_sources"] == {
        "native_deck",
        "transcript",
    }
    assert context["video_artifact_identity"] == {}
    assert (
        "changed during supervised inspection"
        in context["source_reasons"]["delivery_video"]
    )


def test_unhashable_video_cannot_hide_independent_transcript_and_deck(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    deck = vault / "decks" / "talk.pptx"
    _write_pptx(deck)
    video = vault / "videos" / "talk.mp4"
    video.parent.mkdir(parents=True)
    video.write_bytes(b"synthetic video bytes")
    original_hash = pattern_evidence._sha256_file

    def reject_video_hash(path: Path) -> str:
        if path == video:
            raise AssertionError("video digest must stay inside the bounded probe")
        return original_hash(path)

    monkeypatch.setattr(pattern_evidence, "_sha256_file", reject_video_hash)
    video_assessment = _TestVideoAssessment(
        failure=pattern_evidence.VideoEvidenceError(
            "synthetic video digest worker failure",
            reason_code="video_probe_resource_unavailable",
        )
    )
    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slide_source": "pptx",
            "video_local_path": video.relative_to(vault).as_posix(),
        },
        vault_root=vault,
        video_evidence_assessment=video_assessment,
    )

    assert assessment["verified_capabilities"] == ("slides", "transcript")
    assert "delivery_video" not in assessment["verified_evidence_sources"]
    assert assessment["source_reasons"]["delivery_video"].endswith(
        "synthetic video digest worker failure"
    )


def test_missing_python_pptx_dependency_is_local_to_the_pptx_lane(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    deck = vault / "decks" / "talk.pptx"
    slides = vault / "slides" / "talk.pdf"
    _write_pptx(deck)
    _write_pdf(slides)

    def unavailable_pptx(_path: Path, **_kwargs: Any) -> Any:
        raise pattern_evidence.PptxEvidenceError(
            "PPTX evidence requires the declared python-pptx runtime dependency",
            reason_code="pptx_dependency_unavailable",
        )

    monkeypatch.setattr(
        pattern_evidence,
        "probe_pptx_artifact",
        unavailable_pptx,
    )
    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slides_local_path": slides.relative_to(vault).as_posix(),
            "slide_source": "both",
        },
        vault_root=vault,
    )

    assert assessment["verified_capabilities"] == ("slides", "transcript")
    assert set(assessment["verified_evidence_sources"]) == {
        "static_slides",
        "transcript",
    }
    assert "python-pptx" in assessment["source_reasons"]["native_deck"]


@pytest.mark.parametrize(
    ("supervisor_reason", "supervisor_details", "public_reason"),
    [
        (
            "worker_monitor_unavailable",
            {"dependency": "psutil"},
            "pdf_dependency_unavailable",
        ),
        (
            "worker_monitor_unavailable",
            {},
            "pdf_probe_monitor_unavailable",
        ),
        (
            "worker_monitor_identity_changed",
            {},
            "pdf_probe_monitor_identity_changed",
        ),
        (
            "worker_containment_unavailable",
            {},
            "pdf_probe_containment_unavailable",
        ),
        (
            "worker_memory_limit_exceeded",
            {},
            "pdf_probe_resource_unavailable",
        ),
    ],
)
def test_pdf_supervisor_failure_is_local_and_cause_correct_in_catalog_context(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    supervisor_reason: str,
    supervisor_details: dict[str, object],
    public_reason: str,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    deck = vault / "decks" / "talk.pptx"
    slides = vault / "slides" / "talk.pdf"
    _write_pptx(deck)
    _write_pdf(slides)

    def unavailable_pdf(*_args: Any, **_kwargs: Any):
        raise pdf_evidence._supervisor_failure(
            pdf_evidence.SupervisorError(
                supervisor_reason,
                supervisor_details,
            ),
            timeout_seconds=3.5,
        )

    monkeypatch.setattr(pattern_evidence, "probe_pdf_artifact", unavailable_pdf)
    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slides_local_path": slides.relative_to(vault).as_posix(),
            "slide_source": "both",
        },
        vault_root=vault,
    )

    assert assessment["verified_capabilities"] == ("slides", "transcript")
    assert set(assessment["verified_evidence_sources"]) == {"native_deck", "transcript"}
    unavailable = assessment["unavailable_evidence_sources"]["static_slides"]
    assert unavailable["reason_code"] == public_reason
    assert unavailable["details"] == {"supervisor_reason_code": supervisor_reason}


def test_missing_pypdf_dependency_is_local_to_the_pdf_lane(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    transcript, _ = _write_transcript(vault, timed=False)
    deck = vault / "decks" / "talk.pptx"
    slides = vault / "slides" / "talk.pdf"
    _write_pptx(deck)
    _write_pdf(slides)

    def unavailable_pdf(*_args: Any, **_kwargs: Any):
        raise pdf_evidence.PdfEvidenceError(
            "PDF evidence requires its declared runtime dependencies; install the "
            "speaker-toolkit project dependencies",
            reason_code="pdf_dependency_unavailable",
            details={"exception_type": "ImportError"},
        )

    monkeypatch.setattr(pattern_evidence, "probe_pdf_artifact", unavailable_pdf)
    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {
            "filename": "talk.md",
            "transcript_path": transcript.relative_to(vault).as_posix(),
            "transcript_source": "manual",
            "pptx_path": deck.relative_to(vault).as_posix(),
            "slides_local_path": slides.relative_to(vault).as_posix(),
            "slide_source": "both",
        },
        vault_root=vault,
    )

    assert set(assessment["verified_evidence_sources"]) == {"native_deck", "transcript"}
    unavailable = assessment["unavailable_evidence_sources"]["static_slides"]
    assert unavailable["reason_code"] == "pdf_dependency_unavailable"
    assert unavailable["details"] == {"exception_type": "ImportError"}


@pytest.mark.parametrize(
    "talk",
    [
        {"video_url": "not a URL", "youtube_id": "too-short"},
        {"video_url": "https://youtu.be/too-short", "youtube_id": None},
        {"slides_url": "javascript:alert(1)"},
        {"slides_url": "https://drive.google.com/file/d//view"},
        {"google_drive_id": "../"},
    ],
)
def test_malformed_remote_references_are_not_acquisition_capabilities(
    tmp_path: Path,
    talk: dict[str, object],
) -> None:
    assessment = pattern_evidence.assess_talk_artifact_capabilities(
        {"filename": "talk.md", **talk},
        vault_root=tmp_path,
    )

    assert assessment["acquisition_capabilities"] == ()


def _symlink_escaping(root: Path, name: str, outside: Path) -> Path:
    """Create root/<name> as a symlink to a real file outside root."""
    outside.parent.mkdir(parents=True, exist_ok=True)
    outside.write_text("payload\n", encoding="utf-8")
    root.mkdir(parents=True, exist_ok=True)
    link = root / name
    link.symlink_to(outside)
    return link


def test_vault_root_escape_names_the_vault_root(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    _symlink_escaping(vault, "notes.txt", tmp_path / "external" / "notes.txt")

    with pytest.raises(pattern_evidence.PatternEvidenceError) as excinfo:
        pattern_evidence._resolve_local_artifact(
            vault, "notes.txt", suffix=".txt", label="transcript_path"
        )

    assert "resolves outside the vault root" in str(excinfo.value)


def test_pptx_source_root_escape_names_the_configured_root(tmp_path: Path) -> None:
    source_root = tmp_path / "decks"
    _symlink_escaping(source_root, "notes.txt", tmp_path / "external" / "notes.txt")

    with pytest.raises(pattern_evidence.PatternEvidenceError) as excinfo:
        pattern_evidence._resolve_local_artifact(
            source_root,
            "notes.txt",
            suffix=".txt",
            label="pptx_path",
            root_kind="pptx_source",
        )

    message = str(excinfo.value)
    assert "resolves outside the configured pptx_source_dir root" in message
    assert "vault root" not in message


def test_preclaim_root_escape_names_the_field_preclaim_root(tmp_path: Path) -> None:
    preclaim_root = tmp_path / "preclaim"
    _symlink_escaping(preclaim_root, "notes.txt", tmp_path / "external" / "notes.txt")

    with pytest.raises(pattern_evidence.PatternEvidenceError) as excinfo:
        pattern_evidence._resolve_local_artifact(
            preclaim_root,
            "notes.txt",
            suffix=".txt",
            label="slides_local_path",
            root_kind="preclaim:slides_local_path",
        )

    message = str(excinfo.value)
    assert "resolves outside the slides_local_path preclaim root" in message
    assert "vault root" not in message


def test_bounded_pptx_resolver_names_the_violated_root(tmp_path: Path) -> None:
    """The bounded .pptx lane reports the field-specific root, not a generic one."""
    preclaim_root = tmp_path / "preclaim"
    preclaim_root.mkdir(parents=True, exist_ok=True)

    with pytest.raises(pattern_evidence.PatternEvidenceError) as excinfo:
        pattern_evidence._resolve_local_pptx_artifact(
            preclaim_root,
            str(tmp_path / "external" / "deck.pptx"),
            label="pptx_path",
            root_kind="preclaim:pptx_path",
        )

    message = str(excinfo.value)
    assert "resolves outside the pptx_path preclaim root" in message
    assert "the trusted root" not in message


def test_unknown_root_kind_falls_back_without_claiming_the_vault(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    _symlink_escaping(vault, "notes.txt", tmp_path / "external" / "notes.txt")

    with pytest.raises(pattern_evidence.PatternEvidenceError) as excinfo:
        pattern_evidence._resolve_local_artifact(
            vault,
            "notes.txt",
            suffix=".txt",
            label="transcript_path",
            root_kind="something_new",
        )

    assert "resolves outside the trusted root" in str(excinfo.value)
