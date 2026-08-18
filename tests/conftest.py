"""Shared fixtures and helpers for speaker-toolkit tests."""

import importlib.util
import os
import pathlib
import shutil
import subprocess
import sys

import pytest
from pptx import Presentation
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.slide import CT_Slide
from pptx.oxml.xmlchemy import BaseOxmlElement
from pptx.util import Emu


DEFAULT_PPTX_DIRECTORY_EXCLUSIONS = [
    ".venv",
    "venv",
    "node_modules",
    ".git",
    ".hg",
    ".svn",
    "__pycache__",
    ".pytest_cache",
    ".mypy_cache",
    ".ruff_cache",
    ".tox",
    ".nox",
    ".tessl",
]


def current_tracking_config(**updates: object) -> dict[str, object]:
    """Return the owner-current config generation for writer fixtures."""
    config: dict[str, object] = {
        "schema_version": 2,
        "pptx_directory_exclusions": list(DEFAULT_PPTX_DIRECTORY_EXCLUSIONS),
    }
    config.update(updates)
    return config


# ── Script import helper ──────────────────────────────────────────────

SCRIPTS_PC = os.path.join(
    os.path.dirname(__file__),
    os.pardir,
    "skills",
    "presentation-creator",
    "scripts",
)
SCRIPTS_VI = os.path.join(
    os.path.dirname(__file__),
    os.pardir,
    "skills",
    "vault-ingress",
    "scripts",
)
SCRIPTS_ILL = os.path.join(
    os.path.dirname(__file__),
    os.pardir,
    "skills",
    "illustrations",
    "scripts",
)
SCRIPTS_VP = os.path.join(
    os.path.dirname(__file__),
    os.pardir,
    "skills",
    "vault-profile",
    "scripts",
)
# Repo-root gate scripts (pre-publish checks), not owned by any one skill.
SCRIPTS_ROOT = os.path.abspath(
    os.path.join(os.path.dirname(__file__), os.pardir, "scripts")
)
if SCRIPTS_ROOT not in sys.path:
    sys.path.insert(0, SCRIPTS_ROOT)


def _import_script(path, name):
    """Import a standalone .py script as a module (no package required).

    If the module name is already in `sys.modules` (typically because a
    sibling script imported it under the same name via Python's normal
    import machinery — e.g., `extract-script.py` doing
    `import outline_schema`), reuse that cached instance instead of
    overwriting it. Replacing the cached module creates two distinct
    module objects with non-identical enums/classes, which silently
    breaks `isinstance` and identity checks across tests.
    """
    path = os.path.abspath(path)
    script_dir = os.path.dirname(path)
    if script_dir not in sys.path:
        sys.path.insert(0, script_dir)
    if name in sys.modules:
        return sys.modules[name]
    spec = importlib.util.spec_from_file_location(name, path)
    # Both are Optional for a path no loader claims. Every caller passes a real
    # .py file, so a None here is a broken fixture path, not a missing feature —
    # naming it beats an AttributeError two lines later.
    assert spec is not None and spec.loader is not None, f"no loader for {path}"
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


# ── Session-scoped script modules ────────────────────────────────────


@pytest.fixture(scope="session")
def validate_deckops():
    return _import_script(
        os.path.join(SCRIPTS_PC, "validate-deckops.py"), "validate_deckops"
    )


@pytest.fixture(scope="session")
def backgrounds_manifest_to_spec():
    return _import_script(
        os.path.join(SCRIPTS_PC, "backgrounds-manifest-to-spec.py"),
        "backgrounds_manifest_to_spec",
    )


@pytest.fixture(scope="session")
def notes_to_packed():
    return _import_script(
        os.path.join(SCRIPTS_PC, "notes-to-packed.py"), "notes_to_packed"
    )


@pytest.fixture(scope="session")
def generate_qr():
    return _import_script(os.path.join(SCRIPTS_PC, "generate-qr.py"), "generate_qr")


@pytest.fixture(scope="session")
def pptx_extraction():
    return _import_script(
        os.path.join(SCRIPTS_VI, "pptx-extraction.py"), "pptx_extraction"
    )


@pytest.fixture(scope="session")
def pptx_evidence():
    return _import_script(os.path.join(SCRIPTS_VI, "pptx_evidence.py"), "pptx_evidence")


@pytest.fixture(scope="session")
def artifact_metadata():
    return _import_script(
        os.path.join(SCRIPTS_VI, "artifact_metadata.py"), "artifact_metadata"
    )


@pytest.fixture(scope="session")
def pdf_evidence():
    return _import_script(os.path.join(SCRIPTS_VI, "pdf_evidence.py"), "pdf_evidence")


@pytest.fixture(scope="session")
def video_evidence():
    return _import_script(
        os.path.join(SCRIPTS_VI, "video_evidence.py"),
        "video_evidence",
    )


_TINY_VIDEO_BYTES: bytes | None = None


def write_tiny_video(path):
    """Materialize one valid MP4 while keeping ffmpeg local to video tests."""
    global _TINY_VIDEO_BYTES
    path = pathlib.Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    if _TINY_VIDEO_BYTES is not None:
        path.write_bytes(_TINY_VIDEO_BYTES)
        return path

    ffmpeg = shutil.which("ffmpeg")
    assert ffmpeg is not None, "source-video tests require ffmpeg"
    assert shutil.which("ffprobe") is not None, "source-video tests require ffprobe"
    created = subprocess.run(
        [
            ffmpeg,
            "-nostdin",
            "-hide_banner",
            "-loglevel",
            "error",
            "-f",
            "lavfi",
            "-i",
            "color=c=black:s=160x90:r=1",
            "-t",
            "1",
            "-an",
            "-c:v",
            "mpeg4",
            "-pix_fmt",
            "yuv420p",
            "-movflags",
            "+faststart",
            "-y",
            str(path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    assert created.returncode == 0, created.stderr
    _TINY_VIDEO_BYTES = path.read_bytes()
    return path


def synthetic_video_source_receipt(
    *,
    size_bytes: int = 4096,
    digest: str = "0" * 64,
    duration_seconds: float = 8.0,
):
    """Build a receipt for manifest fixtures that never touch a real video.

    Routed through the shipped builder so a producer-side shape change breaks
    these fixtures instead of letting them drift into a shape the readers no
    longer accept. Every field is fixed; nothing here reads the clock.
    """
    video_evidence = _import_script(
        os.path.join(SCRIPTS_VI, "video_evidence.py"),
        "video_evidence",
    )
    artifact_supervisor = _import_script(
        os.path.join(SCRIPTS_VI, "artifact_supervisor.py"),
        "artifact_supervisor",
    )
    artifact_metadata = _import_script(
        os.path.join(SCRIPTS_VI, "artifact_metadata.py"),
        "artifact_metadata",
    )
    generation = artifact_supervisor.FileGeneration(
        size=size_bytes,
        mtime_ns=1,
        ctime_ns=1,
        device=1,
        inode=1,
        mode=0o100644,
        flags=0,
        file_attributes=None,
    )
    return video_evidence.build_video_source_receipt(
        video_evidence.VideoArtifactProbe(
            generation=generation,
            root_generation=None,
            availability=artifact_metadata.ArtifactAvailability.from_generation(
                generation
            ),
            source_sha256=digest,
            source_size_bytes=size_bytes,
            duration_seconds=duration_seconds,
            duration_source="format",
            container_family="iso_bmff",
            stream_count=2,
            video_stream_count=1,
            audio_stream_count=1,
            attached_picture_count=0,
            other_stream_count=0,
            parser_diagnostics=artifact_supervisor.DiagnosticReceipt.empty(),
        )
    )


def video_source_receipt_for(path):
    """Return the receipt a real bounded probe yields for one source video.

    Manifest fixtures stamp this exactly as the extractor does, so the tests
    exercise the shipped receipt contract rather than a hand-written shape that
    could drift away from it. The probe cache keys on the file generation, so
    repeated calls for one fixture video cost a stat.
    """
    video_evidence = _import_script(
        os.path.join(SCRIPTS_VI, "video_evidence.py"),
        "video_evidence",
    )
    return video_evidence.build_video_source_receipt(
        video_evidence.VideoEvidenceAssessment().probe(str(path))
    )


@pytest.fixture(scope="session")
def vtt_cleanup():
    return _import_script(os.path.join(SCRIPTS_VI, "vtt-cleanup.py"), "vtt_cleanup")


@pytest.fixture(scope="session")
def persist_results():
    return _import_script(
        os.path.join(SCRIPTS_VI, "persist-results.py"), "persist_results"
    )


@pytest.fixture(scope="session")
def tracking_database_io():
    return _import_script(
        os.path.join(SCRIPTS_VI, "tracking_database_io.py"),
        "tracking_database_io",
    )


@pytest.fixture(scope="session")
def cooperative_lock():
    """The persistent sibling lock every writer of one owner file shares (#168)."""
    return _import_script(
        os.path.join(SCRIPTS_VI, "cooperative_lock.py"),
        "cooperative_lock",
    )


@pytest.fixture(scope="session")
def persisted_pattern_observations():
    """The persisted-observation structural classifier every reader shares (#167)."""
    return _import_script(
        os.path.join(SCRIPTS_VI, "persisted_pattern_observations.py"),
        "persisted_pattern_observations",
    )


@pytest.fixture(scope="session")
def profile_pattern_provenance():
    """The profile's pattern-provenance contract (#160)."""
    return _import_script(
        os.path.join(SCRIPTS_VP, "profile_pattern_provenance.py"),
        "profile_pattern_provenance",
    )


@pytest.fixture(scope="session")
def pptx_talk_identity():
    """The talk-identity assessment run before catalog persistence (#176)."""
    return _import_script(
        os.path.join(SCRIPTS_VI, "pptx_talk_identity.py"),
        "pptx_talk_identity",
    )


@pytest.fixture(scope="session")
def retained_stage():
    """The staged-file lifecycle both owner writers share (#243).

    Tests that inject staging faults patch this module, not an owner: the
    observation loop, the byte read, and the descriptor closes live here.
    """
    return _import_script(
        os.path.join(SCRIPTS_VI, "retained_stage.py"),
        "retained_stage",
    )


@pytest.fixture(scope="session")
def mutate_tracking_database():
    return _import_script(
        os.path.join(SCRIPTS_VI, "mutate-tracking-database.py"),
        "mutate_tracking_database",
    )


@pytest.fixture(scope="session")
def read_tracking_database():
    return _import_script(
        os.path.join(SCRIPTS_VI, "read-tracking-database.py"),
        "read_tracking_database",
    )


@pytest.fixture(scope="session")
def queue_state():
    return _import_script(
        os.path.join(SCRIPTS_VI, "queue-state.py"),
        "queue_state",
    )


@pytest.fixture(scope="session")
def scan_shownotes_module():
    return _import_script(
        os.path.join(SCRIPTS_VI, "scan-shownotes.py"),
        "scan_shownotes_module",
    )


@pytest.fixture(scope="session")
def tracking_database():
    return _import_script(
        os.path.join(SCRIPTS_VI, "tracking_database.py"), "tracking_database"
    )


@pytest.fixture(scope="session")
def migrate_tracking_database():
    return _import_script(
        os.path.join(SCRIPTS_VI, "migrate-tracking-database.py"),
        "migrate_tracking_database",
    )


@pytest.fixture(scope="session")
def return_validation():
    return _import_script(
        os.path.join(SCRIPTS_VI, "return_validation.py"), "return_validation"
    )


@pytest.fixture(scope="session")
def write_analysis():
    return _import_script(
        os.path.join(SCRIPTS_VI, "write-analysis.py"), "write_analysis"
    )


@pytest.fixture(scope="session")
def validate_returns():
    return _import_script(
        os.path.join(SCRIPTS_VI, "validate-returns.py"), "validate_returns"
    )


@pytest.fixture(scope="session")
def failure_diagnostics():
    return _import_script(
        os.path.join(SCRIPTS_VI, "failure_diagnostics.py"), "failure_diagnostics"
    )


@pytest.fixture(scope="session")
def fetch_transcript():
    return _import_script(
        os.path.join(SCRIPTS_VI, "fetch-transcript.py"), "fetch_transcript"
    )


@pytest.fixture(scope="session")
def preflight_vault():
    return _import_script(
        os.path.join(SCRIPTS_VI, "preflight-vault.py"), "preflight_vault"
    )


@pytest.fixture(scope="session")
def apply_source_repairs():
    return _import_script(
        os.path.join(SCRIPTS_VI, "apply-source-repairs.py"),
        "apply_source_repairs",
    )


@pytest.fixture(scope="session")
def aggregate_catalog_feedback():
    return _import_script(
        os.path.join(SCRIPTS_VI, "aggregate-catalog-feedback.py"),
        "aggregate_catalog_feedback",
    )


@pytest.fixture(scope="session")
def audit_source_identities():
    return _import_script(
        os.path.join(SCRIPTS_VI, "audit-source-identities.py"),
        "audit_source_identities",
    )


@pytest.fixture(scope="session")
def audit_pattern_catalog():
    return _import_script(
        os.path.join(SCRIPTS_VI, "audit-pattern-catalog.py"),
        "audit_pattern_catalog",
    )


@pytest.fixture(scope="session")
def transcript_timing():
    return _import_script(
        os.path.join(SCRIPTS_VI, "transcript_timing.py"),
        "transcript_timing",
    )


@pytest.fixture(scope="session")
def extract_resources():
    return _import_script(
        os.path.join(SCRIPTS_PC, "extract-resources.py"), "extract_resources"
    )


@pytest.fixture(scope="session")
def guardrail_check():
    return _import_script(
        os.path.join(SCRIPTS_PC, "guardrail-check.py"), "guardrail_check"
    )


@pytest.fixture(scope="session")
def outline_schema():
    return _import_script(
        os.path.join(SCRIPTS_PC, "outline_schema.py"), "outline_schema"
    )


@pytest.fixture(scope="session")
def extract_script():
    return _import_script(
        os.path.join(SCRIPTS_PC, "extract-script.py"), "extract_script"
    )


@pytest.fixture(scope="session")
def extract_slides():
    return _import_script(
        os.path.join(SCRIPTS_PC, "extract-slides.py"), "extract_slides"
    )


@pytest.fixture(scope="session")
def extract_narrative():
    return _import_script(
        os.path.join(SCRIPTS_PC, "extract-narrative.py"), "extract_narrative"
    )


@pytest.fixture(scope="session")
def check_rhetorical():
    return _import_script(
        os.path.join(SCRIPTS_PC, "check-rhetorical.py"), "check_rhetorical"
    )


@pytest.fixture(scope="session")
def classify_prose_scan():
    """The prose-scan guardrail classifier (#287)."""
    return _import_script(
        os.path.join(SCRIPTS_PC, "classify-prose-scan.py"), "classify_prose_scan"
    )


@pytest.fixture(scope="session")
def model_registry():
    return _import_script(
        os.path.join(SCRIPTS_ILL, "model_registry.py"), "model_registry"
    )


@pytest.fixture(scope="session")
def validate_profile():
    return _import_script(
        os.path.join(SCRIPTS_VP, "validate-profile.py"), "validate_profile"
    )


@pytest.fixture(scope="session")
def classify_pattern_profile():
    return _import_script(
        os.path.join(SCRIPTS_VP, "classify-pattern-profile.py"),
        "classify_pattern_profile",
    )


@pytest.fixture(scope="session")
def compute_pacing_adherence():
    return _import_script(
        os.path.join(SCRIPTS_VP, "compute-pacing-adherence.py"),
        "compute_pacing_adherence",
    )


@pytest.fixture(scope="session")
def build_expansion_manifest():
    return _import_script(
        os.path.join(SCRIPTS_ILL, "build-expansion-manifest.py"),
        "build_expansion_manifest",
    )


@pytest.fixture(scope="session")
def build_expansion_to_packed():
    return _import_script(
        os.path.join(SCRIPTS_PC, "build-expansion-to-packed.py"),
        "build_expansion_to_packed",
    )


@pytest.fixture(scope="session")
def stage_images_into_container():
    return _import_script(
        os.path.join(SCRIPTS_PC, "stage-images-into-container.py"),
        "stage_images_into_container",
    )


@pytest.fixture(scope="session")
def sync_deck_drivers():
    return _import_script(
        os.path.join(SCRIPTS_PC, "sync-deck-drivers.py"), "sync_deck_drivers"
    )


@pytest.fixture(scope="session")
def generate_illustrations():
    return _import_script(
        os.path.join(SCRIPTS_ILL, "generate-illustrations.py"), "generate_illustrations"
    )


@pytest.fixture(scope="session")
def generate_thumbnail():
    return _import_script(
        os.path.join(SCRIPTS_ILL, "generate-thumbnail.py"), "generate_thumbnail"
    )


@pytest.fixture(scope="session")
def video_slide_extraction():
    return _import_script(
        os.path.join(SCRIPTS_VI, "video-slide-extraction.py"), "video_slide_extraction"
    )


@pytest.fixture(scope="session")
def export_pdf():
    return _import_script(os.path.join(SCRIPTS_PC, "export-pdf.py"), "export_pdf")


@pytest.fixture(scope="session")
def apply_illustrations():
    return _import_script(
        os.path.join(SCRIPTS_ILL, "apply-illustrations-to-deck.py"),
        "apply_illustrations",
    )


@pytest.fixture(scope="session")
def suggest_scrim_color():
    return _import_script(
        os.path.join(SCRIPTS_ILL, "suggest-scrim-color.py"), "suggest_scrim_color"
    )


@pytest.fixture(scope="session")
def generate_talk_timings():
    return _import_script(
        os.path.join(SCRIPTS_PC, "generate-talk-timings.py"), "generate_talk_timings"
    )


# ── PPTX fixture builders ────────────────────────────────────────────

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def deck_width(presentation) -> Emu:
    """A deck's slide width, proven present.

    python-pptx types `slide_width` `Length | None` because a malformed package
    can omit `sldSz`. A `Presentation()` built in-process always carries it, so
    the invariant is asserted once here rather than ignored at every arithmetic
    site (`language-diagnostics` prefers the typed helper).
    """
    width = presentation.slide_width
    assert width is not None, "in-process decks always carry a slide width"
    return Emu(width)


def deck_height(presentation) -> Emu:
    """A deck's slide height, proven present. See `deck_width`."""
    height = presentation.slide_height
    assert height is not None, "in-process decks always carry a slide height"
    return Emu(height)


def slide_title(slide):
    """A slide's title placeholder, proven present.

    `shapes.title` is Optional because a layout may have none. Fixtures that
    set a title have already chosen a layout that has one, so the assertion
    belongs here rather than at every assignment.
    """
    title = slide.shapes.title
    assert title is not None, "this layout has no title placeholder"
    return title


def slide_element(slide) -> CT_Slide:
    """The `<p:sld>` element behind a slide, narrowed from `BaseOxmlElement`.

    python-pptx types `Slide.element` as the generic base, which carries no
    `cSld`. Fixtures that author background DrawingML reach through it, so the
    narrowing happens once here.
    """
    element = slide.element
    assert isinstance(element, CT_Slide), type(element).__name__
    return element


def background_fill_element(slide) -> BaseOxmlElement:
    """The fill element under a slide's authored `<p:bg><p:bgPr>`.

    python-pptx annotates `CT_Background.bgPr` as Optional but leaves
    `CT_BackgroundProperties.eg_fillProperties` unannotated, so a reader sees
    the `ZeroOrOneChoice` descriptor rather than the element it returns. The
    Optional chain is asserted here and the descriptor gap is suppressed once,
    instead of at each fixture that reaches through it.
    """
    background = slide_element(slide).cSld.bg
    assert background is not None, "slide has no authored background"
    properties = background.bgPr
    assert properties is not None, "slide background has no <p:bgPr>"
    fill: BaseOxmlElement | None = properties.eg_fillProperties  # pyright: ignore[reportAssignmentType]
    assert fill is not None, "slide background has no fill element"
    return fill


def background_properties(slide):
    """A slide's `<p:bgPr>`, created if absent. See `background_fill_element`."""
    return slide_element(slide).cSld.get_or_add_bgPr()


def clear_background_fill(properties) -> None:
    """Drop any existing fill under a `<p:bgPr>`. See `background_fill_element`."""
    existing: BaseOxmlElement | None = properties.eg_fillProperties  # pyright: ignore[reportAssignmentType]
    if existing is not None:
        properties.remove(existing)


def graphic_frame_element(shape) -> CT_GraphicalObjectFrame:
    """The `<p:graphicFrame>` element behind a table or chart shape.

    python-pptx types `BaseShape.element` as the union of every CT_* shape
    element, and only the graphic-frame member carries `.graphic`. `add_table`
    always returns a GraphicFrame, so narrowing once here beats an ignore at
    each fixture site that retags a table as SmartArt.
    """
    element = shape.element
    assert isinstance(element, CT_GraphicalObjectFrame), type(element).__name__
    return element


def _make_deck(slide_count, *, slide_width=None, slide_height=None):
    """Create a minimal PPTX with *slide_count* blank slides."""
    prs = Presentation()
    if slide_width:
        prs.slide_width = slide_width
    if slide_height:
        prs.slide_height = slide_height
    blank = prs.slide_layouts[6]  # Blank layout
    for _ in range(slide_count):
        prs.slides.add_slide(blank)
    return prs


@pytest.fixture
def five_slide_deck(tmp_path):
    """Return (Presentation, path) for a 5-slide deck saved to tmp_path."""
    prs = _make_deck(5)
    path = str(tmp_path / "five.pptx")
    prs.save(path)
    return prs, path


@pytest.fixture
def three_slide_deck(tmp_path):
    """Return (Presentation, path) for a 3-slide deck saved to tmp_path."""
    prs = _make_deck(3)
    path = str(tmp_path / "three.pptx")
    prs.save(path)
    return prs, path


@pytest.fixture
def deck_with_text(tmp_path):
    """Return (Presentation, path) for a 3-slide deck with text on each slide."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content
    for i in range(3):
        slide = prs.slides.add_slide(layout)
        slide_title(slide).text = f"Slide {i + 1} Title"
    path = str(tmp_path / "text_deck.pptx")
    prs.save(path)
    return prs, path


def make_deck(slide_count):
    """Public helper for tests that need a Presentation without saving."""
    return _make_deck(slide_count)


@pytest.fixture(scope="session")
def load_vault():
    return _import_script(os.path.join(SCRIPTS_VP, "load-vault.py"), "load_vault")


@pytest.fixture(scope="session")
def section15_pattern_history():
    return _import_script(
        os.path.join(SCRIPTS_VP, "section15_pattern_history.py"),
        "section15_pattern_history",
    )


@pytest.fixture(scope="session")
def classify_pptx_evidence():
    return _import_script(
        os.path.join(SCRIPTS_VI, "classify-pptx-evidence.py"),
        "classify_pptx_evidence",
    )


@pytest.fixture(scope="session")
def pptx_catalog_selection():
    return _import_script(
        os.path.join(SCRIPTS_VI, "pptx_catalog_selection.py"),
        "pptx_catalog_selection",
    )


@pytest.fixture(scope="session")
def render_vault_status():
    return _import_script(
        os.path.join(SCRIPTS_VI, "render-vault-status.py"),
        "render_vault_status",
    )
