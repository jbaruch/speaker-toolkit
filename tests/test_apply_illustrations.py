"""Tests for apply-illustrations-to-deck.py — zone parsing, scrim, title repositioning.

The parsers read outline.yaml (the single source of truth); tests build minimal
partial outlines via _write_outline.
"""

import copy
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches

from conftest import make_deck


FIXTURE = Path(__file__).parent / "fixtures" / "outline-example.yaml"
_TALK = yaml.safe_load(FIXTURE.read_text(encoding="utf-8"))["talk"]


def _write_outline(tmp_path, slides, *, composition=None, embedded_footer=None,
                   name="outline.yaml"):
    """Write a minimal partial outline.yaml for the apply parsers.

    `slides` is a list of dicts; each needs at least `n` + `format`. chapter and
    title are filled in so the partial schema validates.
    """
    norm = [{"chapter": "c", "title": "S", **s} for s in slides]
    data = {"talk": copy.deepcopy(_TALK), "slides": norm}
    if composition is not None or embedded_footer is not None:
        anchor = {"model": "imagen-4", "full": "F", "imgtxt": "I", "conventions": "C"}
        if composition is not None:
            anchor["composition"] = composition
        if embedded_footer is not None:
            anchor["embedded_footer"] = embedded_footer
        data["style_anchor"] = anchor
    p = tmp_path / name
    p.write_text(yaml.safe_dump(data, sort_keys=False), encoding="utf-8")
    return p


ZONES_SLIDES = [
    {"n": 2, "format": "FULL", "safe_zone": {"zone": "upper_third", "surface": "painted sky"}},
    {"n": 5, "format": "FULL", "safe_zone": {"zone": "lower_third", "surface": "flat gradient"}},
    {"n": 8, "format": "FULL", "safe_zone": {"zone": "left_half", "surface": "studio backdrop"}},
]

NO_ZONES_SLIDES = [{"n": 1, "format": "FULL", "image_prompt": "A title slide"}]


def test_parse_zones(apply_illustrations, tmp_path):
    outline = _write_outline(tmp_path, ZONES_SLIDES)
    zones = apply_illustrations.parse_zones(outline)
    assert zones == {2: "upper_third", 5: "lower_third", 8: "left_half"}


def test_parse_zones_empty(apply_illustrations, tmp_path):
    outline = _write_outline(tmp_path, NO_ZONES_SLIDES)
    zones = apply_illustrations.parse_zones(outline)
    assert zones == {}


def test_parse_zones_all_types(apply_illustrations, tmp_path):
    outline = _write_outline(tmp_path, [
        {"n": 1, "format": "FULL", "safe_zone": {"zone": "upper_third"}},
        {"n": 2, "format": "FULL", "safe_zone": {"zone": "middle_third"}},
        {"n": 3, "format": "FULL", "safe_zone": {"zone": "lower_third"}},
        {"n": 4, "format": "FULL", "safe_zone": {"zone": "left_half"}},
        {"n": 5, "format": "FULL", "safe_zone": {"zone": "right_half"}},
    ])
    zones = apply_illustrations.parse_zones(outline)
    assert len(zones) == 5
    assert zones[1] == "upper_third"
    assert zones[2] == "middle_third"
    assert zones[3] == "lower_third"
    assert zones[4] == "left_half"
    assert zones[5] == "right_half"


def test_zone_layout_keys(apply_illustrations):
    """All five zones have layout entries."""
    for zone in ["upper_third", "middle_third", "lower_third", "left_half", "right_half"]:
        assert zone in apply_illustrations.ZONE_LAYOUT
        layout = apply_illustrations.ZONE_LAYOUT[zone]
        assert "top_in" in layout
        assert "left_in" in layout
        assert "width_in" in layout
        assert "height_in" in layout


def test_scrim_shape_name_constant(apply_illustrations):
    """Scrim uses a named constant for detection."""
    assert apply_illustrations.SCRIM_SHAPE_NAME == "_title_scrim"


def test_ensure_scrim_adds_shape(apply_illustrations, tmp_path):
    """ensure_scrim adds a scrim rectangle to a slide."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Add a textbox so the slide has text
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))

    result = apply_illustrations.ensure_scrim(slide, "upper_third", "000000", 45000)
    assert result == 1

    # Verify a shape with the scrim name was added
    scrim_shapes = [s for s in slide.shapes if s.name == "_title_scrim"]
    assert len(scrim_shapes) == 1


def test_ensure_scrim_idempotent(apply_illustrations, tmp_path):
    """Running ensure_scrim twice doesn't add a second scrim."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))

    apply_illustrations.ensure_scrim(slide, "upper_third", "000000", 45000)
    result = apply_illustrations.ensure_scrim(slide, "upper_third", "000000", 45000)
    assert result == 0  # Already present, not added again

    scrim_shapes = [s for s in slide.shapes if s.name == "_title_scrim"]
    assert len(scrim_shapes) == 1


def test_reposition_title(apply_illustrations, tmp_path):
    """Title text shapes are repositioned to the zone."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = "Title"

    moved = apply_illustrations.reposition_title(slide, "upper_third")
    assert moved == 1


def test_reposition_title_placeholder(apply_illustrations, tmp_path):
    """Placeholder title shapes (from Title layouts) are also repositioned."""
    prs = Presentation()
    title_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Placeholder Title"

    moved = apply_illustrations.reposition_title(slide, "lower_third")
    # Title Slide has title + subtitle placeholders
    assert moved >= 1


def test_parse_zones_no_cross_slide_match(apply_illustrations, tmp_path):
    """A slide without a safe_zone stays absent; the neighbouring zone doesn't leak."""
    outline = _write_outline(tmp_path, [
        {"n": 4, "format": "FULL", "image_prompt": "A scene"},
        {"n": 5, "format": "FULL", "image_prompt": "Another scene",
         "safe_zone": {"zone": "lower_third", "surface": "gradient"}},
    ])
    zones = apply_illustrations.parse_zones(outline)
    assert 4 not in zones
    assert zones.get(5) == "lower_third"


# ── IMG+TXT layout tests ─────────────────────────────────────────────


MIXED_FORMAT_SLIDES = [
    {"n": 1, "format": "FULL", "image_prompt": "Opening scene",
     "safe_zone": {"zone": "upper_third", "surface": "sky"}},
    {"n": 2, "format": "IMG+TXT", "image_prompt": "A diagram with explanation"},
    {"n": 3, "format": "EXCEPTION", "format_justification": "real screenshot",
     "visual": "screenshot of the dashboard"},
    {"n": 4, "format": "IMG+TXT", "image_prompt": "Mixed signals",
     "safe_zone": {"zone": "lower_third", "surface": "gradient"}},
]


def test_parse_img_txt_slides(apply_illustrations, tmp_path):
    outline = _write_outline(tmp_path, MIXED_FORMAT_SLIDES)
    img_txt = apply_illustrations.parse_img_txt_slides(outline)
    # Slide 2 is IMG+TXT only — included
    assert 2 in img_txt
    # Slide 1 is FULL — excluded
    assert 1 not in img_txt
    # Slide 3 is EXCEPTION — excluded
    assert 3 not in img_txt
    # Slide 4 has both format IMG+TXT and a safe_zone — safe zone wins, excluded
    assert 4 not in img_txt


def test_parse_img_txt_slides_empty(apply_illustrations, tmp_path):
    outline = _write_outline(tmp_path, NO_ZONES_SLIDES)
    img_txt = apply_illustrations.parse_img_txt_slides(outline)
    assert img_txt == set()


def test_parse_zones_and_img_txt_disjoint(apply_illustrations, tmp_path):
    """A slide with both a safe_zone and format IMG+TXT goes to zones, not img_txt."""
    outline = _write_outline(tmp_path, MIXED_FORMAT_SLIDES)
    zones = apply_illustrations.parse_zones(outline)
    img_txt = apply_illustrations.parse_img_txt_slides(outline)
    # Slide 4 has a safe_zone — it's a FULL slide via the zone path
    assert zones.get(4) == "lower_third"
    assert 4 not in img_txt


def test_apply_img_txt_layout_repositions_picture(apply_illustrations, tmp_path):
    """Picture in an IMG+TXT slide is moved to the left column at ~60% width."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Add a picture at an arbitrary position
    img = tmp_path / "test.png"
    # Minimal 1x1 PNG
    img.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108020000"
            "00907753de0000000c4944415478da6300010000000500010d0a2db40000"
            "000049454e44ae426082"
        )
    )
    slide.shapes.add_picture(str(img), Inches(2), Inches(2), Inches(4), Inches(3))

    pic_moved, _ = apply_illustrations.apply_img_txt_layout(slide)
    assert pic_moved == 1

    pictures = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    bg = pictures[0]
    assert bg.left == Inches(apply_illustrations.IMGTXT_IMG_LEFT_IN)
    assert bg.top == Inches(apply_illustrations.IMGTXT_IMG_TOP_IN)
    assert bg.width == Inches(apply_illustrations.IMGTXT_IMG_WIDTH_IN)
    assert bg.height == Inches(apply_illustrations.IMGTXT_IMG_HEIGHT_IN)


def test_apply_img_txt_layout_repositions_title(apply_illustrations):
    """Title placeholder is repositioned to the right column."""
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Right Column Title"

    _, text_moved = apply_illustrations.apply_img_txt_layout(slide)
    assert text_moved >= 1
    assert slide.shapes.title.left == Inches(apply_illustrations.IMGTXT_TEXT_LEFT_IN)
    assert slide.shapes.title.top == Inches(apply_illustrations.IMGTXT_TITLE_TOP_IN)


def test_swap_or_insert_picture_inserts_when_missing(apply_illustrations, tmp_path):
    """When a slide has no picture shape, swap_or_insert_picture adds one."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # No picture on the slide

    img = tmp_path / "test.png"
    img.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108020000"
            "00907753de0000000c4944415478da6300010000000500010d0a2db40000"
            "000049454e44ae426082"
        )
    )

    pic = apply_illustrations.swap_or_insert_picture(slide, img)
    assert pic is not None
    pictures = [s for s in slide.shapes if s.shape_type == 13]
    assert len(pictures) == 1
    # Inserted full-bleed at slide dimensions
    assert pictures[0].left == Inches(0)
    assert pictures[0].top == Inches(0)


def test_swap_or_insert_picture_swaps_when_present(apply_illustrations, tmp_path):
    """When a slide already has a picture shape, swap_or_insert_picture swaps it."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    img1 = tmp_path / "before.png"
    img2 = tmp_path / "after.png"
    png_bytes = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108020000"
        "00907753de0000000c4944415478da6300010000000500010d0a2db40000"
        "000049454e44ae426082"
    )
    img1.write_bytes(png_bytes)
    img2.write_bytes(png_bytes)
    slide.shapes.add_picture(str(img1), Inches(2), Inches(2), Inches(4), Inches(3))
    pictures_before = [s for s in slide.shapes if s.shape_type == 13]
    assert len(pictures_before) == 1

    pic = apply_illustrations.swap_or_insert_picture(slide, img2)
    pictures_after = [s for s in slide.shapes if s.shape_type == 13]
    # Still exactly one picture — swapped, not added
    assert len(pictures_after) == 1
    # python-pptx returns fresh wrapper objects each iteration, so identity-check
    # via the underlying XML element rather than the Python object.
    assert pic._element is pictures_after[0]._element
    # Position untouched by swap (the FULL/IMG+TXT layout repositioning happens after)
    assert pic.left == Inches(2)
    assert pic.top == Inches(2)


_MIN_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108020000"
    "00907753de0000000c4944415478da6300010000000500010d0a2db40000"
    "000049454e44ae426082"
)


def test_apply_full_records_background_not_picture(apply_illustrations, tmp_path):
    """FULL slides are recorded for the VBA background pass, not given a picture shape."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = tmp_path / "deck.pptx"
    make_deck(3).save(str(deck))

    illust_dir = tmp_path / "illustrations"
    illust_dir.mkdir()
    (illust_dir / "slide-02.png").write_bytes(_MIN_PNG)

    outline = _write_outline(tmp_path, [
        {"n": 2, "format": "FULL", "safe_zone": {"zone": "upper_third", "surface": "sky"}},
    ])
    out_deck = tmp_path / "out.pptx"

    zones = apply_illustrations.parse_zones(outline)
    _, backgrounds = apply_illustrations.apply(
        deck, illust_dir, zones, set(), out_deck, "png", "000000", 45000,
    )

    # FULL slide 2 is recorded for the background pass with an absolute path…
    assert backgrounds == {2: str((illust_dir / "slide-02.png").resolve())}
    # …and NOT inserted as a picture shape (background fill is applied later via VBA)
    slide2 = Presentation(str(out_deck)).slides[1]
    pictures = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == []


def test_apply_imgtxt_keeps_picture_shape(apply_illustrations, tmp_path):
    """IMG+TXT slides still get a real picture shape (not a background fill)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = tmp_path / "deck.pptx"
    make_deck(3).save(str(deck))

    illust_dir = tmp_path / "illustrations"
    illust_dir.mkdir()
    (illust_dir / "slide-02.png").write_bytes(_MIN_PNG)

    outline = _write_outline(tmp_path, [{"n": 2, "format": "IMG+TXT"}])
    out_deck = tmp_path / "out.pptx"

    img_txt = apply_illustrations.parse_img_txt_slides(outline)
    _, backgrounds = apply_illustrations.apply(
        deck, illust_dir, {}, img_txt, out_deck, "png", "000000", 45000,
    )

    # IMG+TXT is not a background; it stays a picture shape
    assert backgrounds == {}
    slide2 = Presentation(str(out_deck)).slides[1]
    pictures = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_imgtxt_geometry_constants_consistent(apply_illustrations):
    """IMG+TXT image + text columns + margins fit the 13.333" slide."""
    img_right = (
        apply_illustrations.IMGTXT_IMG_LEFT_IN
        + apply_illustrations.IMGTXT_IMG_WIDTH_IN
    )
    text_right = (
        apply_illustrations.IMGTXT_TEXT_LEFT_IN
        + apply_illustrations.IMGTXT_TEXT_WIDTH_IN
    )
    assert img_right < apply_illustrations.IMGTXT_TEXT_LEFT_IN, (
        "Image must finish before the text column starts"
    )
    assert text_right <= apply_illustrations.SLIDE_W_IN, (
        "Text column must fit within the slide width"
    )


# ── Poster-theatrical composition ────────────────────────────────────

POSTER_SLIDES = [
    {"n": 3, "format": "FULL", "text_overlay": "One team, one bench",
     "image_prompt": "[STYLE ANCHOR] one team at a workbench"},
    {"n": 7, "format": "FULL", "text_overlay": "Many teams, many wires",
     "image_prompt": "[STYLE ANCHOR] many teams, snarl of wires"},
]


def test_parse_composition_poster(apply_illustrations, tmp_path):
    outline = _write_outline(
        tmp_path, POSTER_SLIDES,
        composition="poster-theatrical", embedded_footer="jbaruch • Devoxx 2026",
    )
    assert apply_illustrations.parse_composition(outline) == "poster-theatrical"


def test_parse_composition_absent(apply_illustrations, tmp_path):
    outline = _write_outline(tmp_path, [{"n": 1, "format": "FULL"}])
    assert apply_illustrations.parse_composition(outline) is None


def test_parse_full_slides_poster(apply_illustrations, tmp_path):
    # Poster FULL slides (no safe_zone) are collected for background-only apply.
    outline = _write_outline(
        tmp_path, POSTER_SLIDES,
        composition="poster-theatrical", embedded_footer="jbaruch • Devoxx 2026",
    )
    assert apply_illustrations.parse_full_slides(outline) == {3, 7}


def test_parse_full_slides_excludes_safe_zone(apply_illustrations, tmp_path):
    # A FULL slide with a safe_zone belongs to the zones path, not poster.
    outline = _write_outline(tmp_path, [
        {"n": 1, "format": "FULL", "image_prompt": "x",
         "safe_zone": {"zone": "upper_third"}},
        {"n": 2, "format": "FULL", "image_prompt": "y"},
    ])
    assert apply_illustrations.parse_full_slides(outline) == {2}


def test_apply_main_rejects_poster_with_safe_zone(apply_illustrations, tmp_path, monkeypatch):
    # Poster mode + a safe_zone slide is contradictory — main() must fail fast.
    import pytest
    outline = _write_outline(
        tmp_path,
        [{"n": 1, "format": "FULL", "image_prompt": "x",
          "safe_zone": {"zone": "upper_third"}}],
        composition="poster-theatrical",
    )
    deck = tmp_path / "deck.pptx"
    illust = tmp_path / "illustrations"
    monkeypatch.setattr(
        "sys.argv",
        ["apply-illustrations-to-deck.py", str(deck), str(illust), str(outline)],
    )
    with pytest.raises(SystemExit) as exc:
        apply_illustrations.main()
    assert "poster-theatrical" in str(exc.value)
