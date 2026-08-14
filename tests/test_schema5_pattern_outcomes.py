"""Focused regressions for exhaustive return-schema v5 pattern outcomes."""

from __future__ import annotations

from typing import Any

import copy
from pathlib import Path

import pytest
from pptx import Presentation


def _entry(
    return_validation,
    pattern_id: str,
    *,
    absence_gate=(frozenset({"transcript"}),),
    applicability=False,
):
    conditions = (
        (
            return_validation.NotApplicableCondition(
                "no-opportunity", "No relevant opportunity occurs."
            ),
        )
        if applicability
        else None
    )
    return return_validation.CatalogEntry(
        pattern_id=pattern_id,
        entry_type="pattern",
        observable=True,
        evaluable_from=(frozenset({"transcript"}),),
        strong_evaluable_from=(frozenset({"transcript"}),),
        absence_evaluable_from=absence_gate,
        evidence_channels=frozenset({"transcript"}),
        evidence_metadata_fields=frozenset(),
        vault_dimensions=(2,),
        path=f"synthetic/{pattern_id}.md",
        not_applicable_when=conditions,
        applicability_evaluable_from=(
            (frozenset({"transcript"}),) if applicability else None
        ),
    )


def _catalog(return_validation, *entries):
    return return_validation.PatternCatalog(
        entries={entry.pattern_id: entry for entry in entries},
        fingerprint="a" * 64,
    )


def _artifact(tmp_path: Path, transcript_timing):
    vault = tmp_path / "vault"
    transcript = vault / "transcripts" / "synthetic.txt"
    transcript.parent.mkdir(parents=True)
    first = "A singular opening proves this exact synthetic source. " + " ".join(
        ["alpha"] * 250
    )
    second = "A separate closing proves the second source line. " + " ".join(
        ["beta"] * 250
    )
    text = first + "\n" + second
    transcript_timing.write_transcript_bundle(
        transcript,
        text,
        [],
        source="captions",
        timing_provenance=None,
        quality_policy={
            "schema_version": 1,
            "min_words": 400,
            "duration_seconds": None,
        },
        quality_policy_provenance={"kind": "fixed_default"},
    )
    talk = {
        "schema_version": 5,
        "filename": "synthetic.md",
        "title": "Synthetic",
        "transcript_path": "transcripts/synthetic.txt",
        "transcript_source": "youtube_auto",
        "slide_source": "none",
    }
    return vault, talk


def _citation(*, second=False):
    return {
        "source": "transcript",
        "channel": "transcript",
        "quote": (
            "A separate closing proves the second source line."
            if second
            else "A singular opening proves this exact synthetic source."
        ),
    }


def _raw_return(
    *,
    detections=None,
    assessments=None,
    not_evaluable=None,
    line_ranges=None,
):
    detections = detections or []
    return {
        "filename": "synthetic.md",
        "return_schema_version": 5,
        "queue_claim": {
            "run_id": "schema5",
            "batch_id": "batch-1",
            "reprocess_generation": 1,
        },
        "status": "processed_partial",
        "slide_source": "none",
        "transcript_source": "youtube_auto",
        "rhetoric_notes": "Synthetic rhetoric notes.",
        "areas_for_improvement": "Synthetic improvement notes.",
        "adherence_assessment": "",
        "new_patterns": "",
        "summary_updates": "",
        "structured_data": {},
        "verbatim_examples": {},
        "pattern_observations": {
            "evidence_sources": ["transcript"],
            "source_inspection": [
                {
                    "source": "transcript",
                    "line_ranges": line_ranges or [[1, 2]],
                }
            ],
            "patterns_detected": detections,
            "antipatterns_detected": [],
            "applicability_assessments": assessments or [],
            "not_evaluable": not_evaluable or [],
            "pattern_score": {
                "patterns_used": len(detections),
                "antipatterns_detected": 0,
                "score": len(detections),
            },
        },
        "catalog_feedback": {
            "unmatched_observations": [],
            "confusable_pairs": [],
            "definition_problems": [],
            "scoring_problems": [],
            "tensions": [],
        },
    }


def _detection(pattern_id="detected", *, second=False):
    return {
        "pattern_id": pattern_id,
        "confidence": "moderate",
        "evidence_source": "transcript",
        "evidence": "The source directly establishes this pattern.",
        "evidence_citations": [_citation(second=second)],
    }


def _assessment(result="not_applicable", *, condition=True):
    value = {
        "pattern_id": "conditional",
        "result": result,
        "evidence_source": "transcript",
        "evidence": "The complete source establishes opportunity polarity.",
        "evidence_citations": [_citation()],
    }
    if condition:
        value["condition_id"] = "no-opportunity"
    return value


def _persisted_v5(return_validation, transcript_timing, tmp_path):
    import pattern_evidence

    catalog = _catalog(
        return_validation,
        _entry(return_validation, "detected"),
        _entry(return_validation, "conditional", applicability=True),
        _entry(return_validation, "undetected"),
        _entry(return_validation, "positive-only", absence_gate=None),
    )
    raw = _raw_return(
        detections=[_detection()],
        assessments=[_assessment()],
        not_evaluable=[
            {
                "pattern_id": "positive-only",
                "reason_code": "absence_not_authorized_by_catalog",
            }
        ],
    )
    return_validation.validate_return(raw, catalog)
    vault, owner = _artifact(tmp_path, transcript_timing)
    canonical: dict[str, Any] = pattern_evidence.canonicalize_return_evidence(
        raw,
        owner,
        vault,
        catalog,
        pattern_scoring_schema_version=5,
    )
    generation = return_validation.assess_scoring_generation(canonical, catalog)
    observations = return_validation.canonical_persisted_pattern_observations(
        canonical, catalog, generation
    )
    talk = copy.deepcopy(owner)
    talk.update(
        {
            "status": "processed_partial",
            "pattern_catalog_fingerprint": catalog.fingerprint,
            "pattern_scoring_schema_version": 5,
            "pattern_score": observations["pattern_score"],
            "pattern_observations": observations,
        }
    )
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            talk, vault_root=vault
        )
        == ()
    )
    return pattern_evidence, vault, talk


def _rehash_opportunity(pattern_evidence, talk):
    observations = talk["pattern_observations"]
    observations["opportunity_coverage_identity"] = (
        pattern_evidence.opportunity_coverage_identity(
            observations["pattern_outcomes"],
            pattern_catalog_fingerprint=talk["pattern_catalog_fingerprint"],
            pattern_scoring_schema_version=talk["pattern_scoring_schema_version"],
        )
    )


def test_v5_persists_one_exhaustive_outcome_and_opportunity_identity(
    return_validation, transcript_timing, tmp_path
):
    import pattern_evidence

    catalog = _catalog(
        return_validation,
        _entry(return_validation, "detected"),
        _entry(return_validation, "conditional", applicability=True),
        _entry(return_validation, "undetected"),
        _entry(return_validation, "positive-only", absence_gate=None),
    )
    raw = _raw_return(
        detections=[_detection()],
        assessments=[_assessment()],
        not_evaluable=[
            {
                "pattern_id": "positive-only",
                "reason_code": "absence_not_authorized_by_catalog",
            }
        ],
    )
    return_validation.validate_return(raw, catalog)
    vault, talk = _artifact(tmp_path, transcript_timing)

    canonical: dict[str, Any] = pattern_evidence.canonicalize_return_evidence(
        raw,
        talk,
        vault,
        catalog,
        pattern_scoring_schema_version=5,
    )
    observations: dict[str, Any] = canonical["pattern_observations"]
    transcript_inspection = observations["source_inspection"][0]
    assert transcript_inspection["coverage_complete"] is True
    assert transcript_inspection["absence_capability_complete"] is True
    assert transcript_inspection["absence_capability_reason"] == (
        "authorized_transcript"
    )

    assert observations["pattern_outcomes"] == [
        {"pattern_id": "conditional", "outcome": "not_applicable"},
        {"pattern_id": "detected", "outcome": "detected"},
        {"pattern_id": "positive-only", "outcome": "not_evaluable"},
        {"pattern_id": "undetected", "outcome": "undetected"},
    ]
    assert len(observations["opportunity_coverage_identity"]) == 64
    assert pattern_evidence.return_evidence_claim(
        canonical
    ) == pattern_evidence.return_evidence_claim(raw)
    assessment = return_validation.assess_scoring_generation(canonical, catalog)
    assert assessment.current is True
    persisted = return_validation.canonical_persisted_pattern_observations(
        canonical, catalog, assessment
    )
    assert persisted["evidence_schema_version"] == 2
    assert persisted["pattern_outcomes"] == observations["pattern_outcomes"]


@pytest.mark.parametrize(
    ("pattern_id", "replacement", "expected_reason"),
    [
        (
            "detected",
            "undetected",
            "pattern_outcomes_detected_projection_drift",
        ),
        (
            "positive-only",
            "undetected",
            "pattern_outcomes_not_evaluable_projection_drift",
        ),
        (
            "conditional",
            "undetected",
            "pattern_outcomes_not_applicable_projection_drift",
        ),
    ],
)
def test_freshness_rejects_self_consistently_rehashed_lane_contradictions(
    return_validation,
    transcript_timing,
    tmp_path,
    pattern_id,
    replacement,
    expected_reason,
):
    pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    outcome = next(
        item
        for item in talk["pattern_observations"]["pattern_outcomes"]
        if item["pattern_id"] == pattern_id
    )
    outcome["outcome"] = replacement
    _rehash_opportunity(pattern_evidence, talk)

    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        talk, vault_root=vault
    )

    assert expected_reason in reasons
    assert "opportunity_coverage_identity_drift" not in reasons


def test_freshness_rejects_unknown_and_overlapping_v5_lane_ids(
    return_validation,
    transcript_timing,
    tmp_path,
):
    pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    observations = talk["pattern_observations"]
    unknown = copy.deepcopy(observations["patterns_detected"][0])
    unknown["pattern_id"] = "not-in-outcomes"
    observations["patterns_detected"].append(unknown)
    overlap = copy.deepcopy(observations["not_evaluable"][0])
    overlap["pattern_id"] = "detected"
    observations["not_evaluable"].append(overlap)

    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        talk, vault_root=vault
    )

    assert "pattern_outcomes_unknown_lane_ids" in reasons
    assert "pattern_observation_lane_overlap" in reasons


def test_freshness_rejects_nested_and_promoted_pattern_score_tampering(
    return_validation,
    transcript_timing,
    tmp_path,
):
    pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    talk["pattern_observations"]["pattern_score"] = 99
    talk["pattern_score"] = 99

    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        talk, vault_root=vault
    )

    assert "pattern_score_projection_drift" in reasons
    assert "promoted_pattern_score_drift" in reasons


def test_freshness_rejects_absence_capability_receipt_tampering(
    return_validation,
    transcript_timing,
    tmp_path,
):
    pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    inspection = talk["pattern_observations"]["source_inspection"][0]
    inspection["absence_capability_complete"] = False
    inspection["absence_capability_reason"] = "bare_delivery_video"

    reasons = pattern_evidence.assess_persisted_pattern_evidence_freshness(
        talk, vault_root=vault
    )

    assert "source_inspection[0]:absence_capability_complete_drift" in reasons
    assert "source_inspection[0]:absence_capability_reason_drift" in reasons


def test_active_catalog_freshness_rejects_deleted_rehashed_outcome(
    return_validation,
    transcript_timing,
    tmp_path,
):
    pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    outcomes = talk["pattern_observations"]["pattern_outcomes"]
    talk["pattern_observations"]["pattern_outcomes"] = [
        item for item in outcomes if item["pattern_id"] != "undetected"
    ]
    _rehash_opportunity(pattern_evidence, talk)

    reasons = return_validation.assess_current_persisted_pattern_evidence_freshness(
        talk,
        vault_root=vault,
        catalog=_catalog(
            return_validation,
            _entry(return_validation, "detected"),
            _entry(return_validation, "conditional", applicability=True),
            _entry(return_validation, "undetected"),
            _entry(return_validation, "positive-only", absence_gate=None),
        ),
    )

    assert "pattern_outcomes_catalog_projection_drift" in reasons
    assert "opportunity_coverage_identity_drift" not in reasons


def test_active_catalog_freshness_rejects_rehashed_denominator_reclassification(
    return_validation,
    transcript_timing,
    tmp_path,
):
    pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    observations = talk["pattern_observations"]
    outcome = next(
        item
        for item in observations["pattern_outcomes"]
        if item["pattern_id"] == "undetected"
    )
    outcome["outcome"] = "not_evaluable"
    synthetic_not_evaluable = copy.deepcopy(observations["not_evaluable"][0])
    synthetic_not_evaluable["pattern_id"] = "undetected"
    observations["not_evaluable"].append(synthetic_not_evaluable)
    observations["not_evaluable_ids"].append("undetected")
    _rehash_opportunity(pattern_evidence, talk)

    reasons = return_validation.assess_current_persisted_pattern_evidence_freshness(
        talk,
        vault_root=vault,
        catalog=_catalog(
            return_validation,
            _entry(return_validation, "detected"),
            _entry(return_validation, "conditional", applicability=True),
            _entry(return_validation, "undetected"),
            _entry(return_validation, "positive-only", absence_gate=None),
        ),
    )

    assert "pattern_outcomes_catalog_projection_drift" in reasons
    assert "opportunity_coverage_identity_drift" not in reasons


def test_active_catalog_freshness_rejects_polarity_swapped_detection(
    return_validation,
    transcript_timing,
    tmp_path,
):
    pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    observations = talk["pattern_observations"]
    detection = observations["patterns_detected"].pop()
    observations["antipatterns_detected"].append(detection)
    observations["pattern_ids"] = []
    observations["antipattern_ids"] = ["detected"]
    observations["pattern_score"] = -1
    talk["pattern_score"] = -1

    reasons = return_validation.assess_current_persisted_pattern_evidence_freshness(
        talk,
        vault_root=vault,
        catalog=_catalog(
            return_validation,
            _entry(return_validation, "detected"),
            _entry(return_validation, "conditional", applicability=True),
            _entry(return_validation, "undetected"),
            _entry(return_validation, "positive-only", absence_gate=None),
        ),
    )

    assert "pattern_outcomes_catalog_projection_drift" in reasons


@pytest.mark.parametrize("bad_outcome", [[], {}])
def test_non_string_outcome_is_a_clean_pattern_evidence_error(
    bad_outcome,
):
    import pattern_evidence

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="outcome is invalid",
    ):
        pattern_evidence.opportunity_coverage_identity(
            [{"pattern_id": "synthetic", "outcome": bad_outcome}],
            pattern_catalog_fingerprint="a" * 64,
            pattern_scoring_schema_version=5,
        )


def test_complete_applicability_gate_requires_exactly_one_assessment(
    return_validation, transcript_timing, tmp_path
):
    import pattern_evidence

    catalog = _catalog(
        return_validation,
        _entry(return_validation, "conditional", applicability=True),
    )
    raw = _raw_return()
    return_validation.validate_return(raw, catalog)
    vault, talk = _artifact(tmp_path, transcript_timing)

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="requires exactly one applicability assessment",
    ):
        pattern_evidence.canonicalize_return_evidence(
            raw, talk, vault, catalog, pattern_scoring_schema_version=5
        )


def test_incomplete_applicability_gate_forbids_assessment_and_is_not_evaluable(
    return_validation, transcript_timing, tmp_path
):
    import pattern_evidence

    catalog = _catalog(
        return_validation,
        _entry(return_validation, "conditional", applicability=True),
    )
    raw = _raw_return(
        line_ranges=[[1, 1]],
        not_evaluable=[
            {
                "pattern_id": "conditional",
                "reason_code": "missing_applicability_source_coverage",
            }
        ],
    )
    return_validation.validate_return(raw, catalog)
    vault, talk = _artifact(tmp_path, transcript_timing)
    canonical: dict[str, Any] = pattern_evidence.canonicalize_return_evidence(
        raw, talk, vault, catalog, pattern_scoring_schema_version=5
    )

    assert canonical["pattern_observations"]["pattern_outcomes"] == [
        {
            "pattern_id": "conditional",
            "outcome": "not_evaluable",
        }
    ]

    with_assessment = copy.deepcopy(raw)
    with_assessment["pattern_observations"]["applicability_assessments"] = [
        _assessment()
    ]
    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="assessment without complete applicability-gate coverage",
    ):
        pattern_evidence.canonicalize_return_evidence(
            with_assessment,
            talk,
            vault,
            catalog,
            pattern_scoring_schema_version=5,
        )


def test_bare_native_deck_is_fail_closed_for_applicability_and_fresh_with_root(
    return_validation,
    tmp_path,
):
    import pattern_evidence
    import pptx_evidence

    condition = return_validation.NotApplicableCondition(
        "no-opportunity", "No relevant opportunity occurs."
    )
    native_gate = (frozenset({"native_deck"}),)
    entry = return_validation.CatalogEntry(
        pattern_id="conditional",
        entry_type="pattern",
        observable=True,
        evaluable_from=native_gate,
        strong_evaluable_from=native_gate,
        absence_evaluable_from=native_gate,
        evidence_channels=frozenset({"slides"}),
        evidence_metadata_fields=frozenset(),
        vault_dimensions=(2,),
        path="synthetic/conditional.md",
        not_applicable_when=(condition,),
        applicability_evaluable_from=native_gate,
    )
    catalog = _catalog(return_validation, entry)
    vault = tmp_path / "vault"
    source_root = tmp_path / "pptx-source"
    source_root.mkdir()
    deck_path = source_root / "deck.pptx"
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(str(deck_path))
    owner = {
        "schema_version": 5,
        "filename": "synthetic.md",
        "title": "Synthetic",
        "pptx_path": deck_path.name,
        "slide_source": "pptx",
    }
    raw = _raw_return(
        not_evaluable=[
            {
                "pattern_id": "conditional",
                "reason_code": "missing_applicability_source_coverage",
            }
        ]
    )
    raw["slide_source"] = "pptx"
    raw["structured_data"] = {
        "slide_count": 1,
        "native_deck_audit": pptx_evidence.recompute_native_deck_audit(deck_path),
    }
    raw["pattern_observations"].update(
        {
            "evidence_sources": ["native_deck"],
            "source_inspection": [
                {
                    "source": "native_deck",
                    "page_ranges": [[1, 1]],
                }
            ],
        }
    )
    roots = {"pptx_source_dir": str(source_root)}
    return_validation.validate_return(raw, catalog)
    canonical: dict[str, Any] = pattern_evidence.canonicalize_return_evidence(
        raw,
        owner,
        vault,
        catalog,
        pattern_scoring_schema_version=5,
        source_roots=roots,
    )

    inspection = canonical["pattern_observations"]["source_inspection"][0]
    assert inspection["coverage_complete"] is True
    assert inspection["absence_capability_complete"] is False
    assert inspection["absence_capability_reason"] == "bare_native_deck"
    assert canonical["pattern_observations"]["applicability_assessments"] == []
    assert canonical["pattern_observations"]["pattern_outcomes"] == [
        {
            "pattern_id": "conditional",
            "outcome": "not_evaluable",
        }
    ]

    generation = return_validation.assess_scoring_generation(canonical, catalog)
    observations = return_validation.canonical_persisted_pattern_observations(
        canonical, catalog, generation
    )
    persisted = copy.deepcopy(owner)
    persisted.update(
        {
            "status": "processed_partial",
            "structured_data": copy.deepcopy(canonical["structured_data"]),
            "pattern_catalog_fingerprint": catalog.fingerprint,
            "pattern_scoring_schema_version": 5,
            "pattern_score": observations["pattern_score"],
            "pattern_observations": observations,
        }
    )
    assert (
        pattern_evidence.assess_persisted_pattern_evidence_freshness(
            persisted,
            vault_root=vault,
            source_roots=roots,
        )
        == ()
    )


@pytest.mark.parametrize(
    ("assessment", "message"),
    [
        ({**_assessment(), "condition_id": "not-authorized"}, "must be one of"),
        (_assessment("applicable", condition=True), "condition_id is forbidden"),
    ],
)
def test_raw_applicability_polarity_is_catalog_authorized(
    return_validation, assessment, message
):
    catalog = _catalog(
        return_validation,
        _entry(return_validation, "conditional", applicability=True),
    )
    raw = _raw_return(assessments=[assessment])

    with pytest.raises(return_validation.ReturnValidationError, match=message):
        return_validation.validate_return(raw, catalog)


def test_assessment_without_catalog_gate_and_detected_overlap_are_rejected(
    return_validation,
):
    ungated = _catalog(
        return_validation,
        _entry(return_validation, "conditional"),
    )
    raw = _raw_return(assessments=[_assessment()])
    with pytest.raises(
        return_validation.ReturnValidationError,
        match="no catalog-owned applicability contract",
    ):
        return_validation.validate_return(raw, ungated)

    gated = _catalog(
        return_validation,
        _entry(return_validation, "conditional", applicability=True),
    )
    overlap = _raw_return(
        detections=[_detection("conditional")],
        assessments=[_assessment()],
    )
    with pytest.raises(
        return_validation.ReturnValidationError,
        match="both detected and applicability-assessed",
    ):
        return_validation.validate_return(overlap, gated)


def test_citations_outside_declared_inspection_are_rejected(
    return_validation, transcript_timing, tmp_path
):
    import pattern_evidence

    catalog = _catalog(
        return_validation,
        _entry(return_validation, "detected"),
    )
    raw = _raw_return(
        detections=[_detection(second=True)],
        line_ranges=[[1, 1]],
    )
    return_validation.validate_return(raw, catalog)
    vault, talk = _artifact(tmp_path, transcript_timing)

    with pytest.raises(
        pattern_evidence.PatternEvidenceError,
        match="falls outside the declared inspection ranges",
    ):
        pattern_evidence.canonicalize_return_evidence(
            raw, talk, vault, catalog, pattern_scoring_schema_version=5
        )


def test_opportunity_identity_binds_scoring_generation_independently(
    return_validation,
):
    import pattern_evidence

    outcomes = [{"pattern_id": "one", "outcome": "undetected"}]
    v5 = pattern_evidence.opportunity_coverage_identity(
        outcomes,
        pattern_catalog_fingerprint="a" * 64,
        pattern_scoring_schema_version=5,
    )
    v6 = pattern_evidence.opportunity_coverage_identity(
        outcomes,
        pattern_catalog_fingerprint="a" * 64,
        pattern_scoring_schema_version=6,
    )

    assert v5 != v6
    # The identity is bound to the generation it was built at, not to
    # whichever generation happens to be current.
    assert return_validation.RETURN_SCHEMA_VERSION == (
        return_validation.WEIGHTED_SCORE_RETURN_SCHEMA_VERSION
    )


def test_v4_source_locations_survive_root_migration_without_v5_outcomes(
    tracking_database,
):
    talk = {
        "schema_version": 4,
        "filename": "archive.md",
        "status": "processed",
        "pattern_observations": {
            "evidence_schema_version": 1,
            "patterns_detected": [
                {
                    "pattern_id": "archive",
                    "evidence_citations": [{"source": "transcript"}],
                }
            ],
            "antipatterns_detected": [],
            "source_inspection": [{"source": "transcript"}],
        },
    }
    database = {"talks": [talk]}

    migrated = tracking_database.migrate_tracking_database(database).database
    talk = migrated["talks"][0]
    assert talk["schema_version"] == 4
    assert talk["pattern_observations"]["evidence_schema_version"] == 1
    assert talk["pattern_observations"]["patterns_detected"][0]["evidence_citations"]
    assert "pattern_outcomes" not in talk["pattern_observations"]


def _scored_talk(
    filename: str,
    identity: str,
    score: int = 1,
    *,
    outcome: str = "undetected",
    scoring_schema_version: int | None = None,
):
    """Build a talk stamped at the ACTIVE scoring generation by default.

    Pinning the literal 5 here made every such talk fall out of the cohort the
    moment the generation advanced, which reads downstream as "the baseline
    population is too small" rather than "this fixture is stale".
    """
    import return_validation

    if scoring_schema_version is None:
        scoring_schema_version = return_validation.PATTERN_SCORING_SCHEMA_VERSION
    return {
        "filename": filename,
        "status": "processed",
        "pattern_score": score,
        "pattern_scoring_generation_status": "current",
        "pattern_scoring_generation_reasons": [],
        "pattern_catalog_fingerprint": "a" * 64,
        "pattern_scoring_schema_version": scoring_schema_version,
        "pattern_observations": {
            "pattern_score": score,
            "opportunity_coverage_identity": identity,
            "pattern_outcomes": [{"pattern_id": "synthetic", "outcome": outcome}],
        },
    }


def test_mixed_opportunity_baseline_preserves_eligible_cohort_but_suppresses_score(
    return_validation,
):
    import adherence_baseline

    baseline = adherence_baseline.build_current_cohort_baseline(
        [
            _scored_talk("one.md", "1" * 64, scoring_schema_version=5),
            _scored_talk("two.md", "2" * 64, scoring_schema_version=5),
        ],
        as_of="2026-07-31T12:00:00+00:00",
        pattern_catalog_fingerprint="a" * 64,
        pattern_scoring_schema_version=5,
        evidence_freshness_assessor=lambda _talk: (),
        persisted_observation_assessor=lambda _talk: (),
    )

    assert baseline["eligible_talk_count"] == 2
    assert baseline["scored_talk_count"] == 0
    assert baseline["average_pattern_score"] is None
    assert baseline["opportunity_coverage_identity"] is None
    assert baseline["raw_score_comparison_status"] == "unavailable"
    assert baseline["raw_score_comparison_reason"] == "mixed_opportunity_coverage"
    assert return_validation.PATTERN_SCORING_SCHEMA_VERSION == (
        return_validation.WEIGHTED_PATTERN_SCORING_SCHEMA_VERSION
    )


def test_all_unknown_opportunity_baseline_suppresses_zero_raw_score(
    return_validation,
):
    import adherence_baseline

    baseline = adherence_baseline.build_current_cohort_baseline(
        [
            _scored_talk(
                "one.md",
                "1" * 64,
                score=0,
                outcome="not_evaluable",
                scoring_schema_version=5,
            ),
            _scored_talk(
                "two.md",
                "1" * 64,
                score=0,
                outcome="not_evaluable",
                scoring_schema_version=5,
            ),
        ],
        as_of="2026-07-31T12:00:00+00:00",
        pattern_catalog_fingerprint="a" * 64,
        pattern_scoring_schema_version=5,
        evidence_freshness_assessor=lambda _talk: (),
        persisted_observation_assessor=lambda _talk: (),
    )

    assert baseline["eligible_talk_count"] == 2
    assert baseline["scored_talk_count"] == 0
    assert baseline["pattern_score_sum"] == 0
    assert baseline["average_pattern_score"] is None
    assert baseline["opportunity_coverage_identity"] is None
    assert baseline["raw_score_comparison_status"] == "unavailable"
    assert baseline["raw_score_comparison_reason"] == (
        adherence_baseline.NO_EVALUABLE_PATTERN_OPPORTUNITIES_REASON
    )
    assert return_validation.PATTERN_SCORING_SCHEMA_VERSION == (
        return_validation.WEIGHTED_PATTERN_SCORING_SCHEMA_VERSION
    )


def test_adherence_comparison_is_suppressed_when_opportunity_identity_differs(
    return_validation,
):
    import adherence_baseline

    baseline = adherence_baseline.build_adherence_baseline(
        [
            _scored_talk(f"talk-{index}.md", "1" * 64, scoring_schema_version=5)
            for index in range(10)
        ],
        selected_filenames=[],
        as_of="2026-07-31T12:00:00+00:00",
        pattern_catalog_fingerprint="a" * 64,
        pattern_scoring_schema_version=5,
        evidence_freshness_assessor=lambda _talk: (),
        persisted_observation_assessor=lambda _talk: (),
    )
    talk = {"_queue_claim": {"adherence_baseline": baseline}}
    canonical = {
        "pattern_observations": {
            "opportunity_coverage_identity": "2" * 64,
        }
    }
    empty = {
        "filename": "new.md",
        "return_schema_version": 5,
        "status": "processed_partial",
        "adherence_assessment": "",
    }
    return_validation.validate_v5_adherence_opportunity(talk, empty, canonical)

    misleading = {
        **empty,
        "adherence_assessment": "It rose. This is misleading.",
        "adherence_comparison": {
            "schema_version": 1,
            "baseline": baseline,
            "talk_pattern_score": 1,
        },
    }
    with pytest.raises(
        return_validation.ReturnValidationError,
        match="opportunity coverage identity mismatch",
    ):
        return_validation.validate_v5_adherence_opportunity(talk, misleading, canonical)


@pytest.mark.parametrize(
    "stamp",
    [None, True, False, "5", 5.0],
    ids=["missing", "true", "false", "string", "float"],
)
def test_an_unusable_scoring_stamp_is_not_replaced_with_the_current_one(
    return_validation,
    transcript_timing,
    tmp_path,
    stamp,
):
    """A record with no usable generation stamp has nothing to replay against.

    Substituting the current generation lets a record whose identity was computed
    under it match by coincidence, so unusable state reports as fresh — which is
    the opposite of what a freshness check is for.
    """
    _pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )
    if stamp is None:
        talk.pop("pattern_scoring_schema_version", None)
    else:
        talk["pattern_scoring_schema_version"] = stamp

    reasons = return_validation.assess_current_persisted_pattern_evidence_freshness(
        talk,
        vault_root=vault,
        catalog=_catalog(
            return_validation,
            _entry(return_validation, "detected"),
            _entry(return_validation, "conditional", applicability=True),
            _entry(return_validation, "undetected"),
            _entry(return_validation, "positive-only", absence_gate=None),
        ),
    )

    assert "pattern_scoring_schema_version_unusable" in reasons


def test_a_correctly_stamped_record_still_replays_clean(
    return_validation,
    transcript_timing,
    tmp_path,
):
    """The guard must not make every record unusable."""
    _pattern_evidence, vault, talk = _persisted_v5(
        return_validation, transcript_timing, tmp_path
    )

    reasons = return_validation.assess_current_persisted_pattern_evidence_freshness(
        talk,
        vault_root=vault,
        catalog=_catalog(
            return_validation,
            _entry(return_validation, "detected"),
            _entry(return_validation, "conditional", applicability=True),
            _entry(return_validation, "undetected"),
            _entry(return_validation, "positive-only", absence_gate=None),
        ),
    )

    assert "pattern_scoring_schema_version_unusable" not in reasons


def test_a_v6_return_canonicalizes_on_the_same_exhaustive_path_as_v5(
    return_validation, transcript_timing, tmp_path
):
    """v6 keeps every v5 semantic and adds weighting, so canonicalization must
    treat them alike.

    An `== EXHAUSTIVE_OUTCOME_RETURN_SCHEMA_VERSION` test reads a v6 return as
    pre-v5 and rejects the `applicability_assessments` its own contract
    REQUIRES — so a v6 return validated on the way in, then died on the way to
    the database. Nothing caught it because every canonicalization test built a
    v5 return.
    """
    import pattern_evidence

    catalog = _catalog(
        return_validation,
        _entry(return_validation, "detected"),
        _entry(return_validation, "conditional", applicability=True),
        _entry(return_validation, "undetected"),
        _entry(return_validation, "positive-only", absence_gate=None),
    )
    raw = _raw_return(
        detections=[_detection()],
        assessments=[_assessment()],
        not_evaluable=[
            {
                "pattern_id": "positive-only",
                "reason_code": "absence_not_authorized_by_catalog",
            }
        ],
    )
    raw["return_schema_version"] = (
        return_validation.WEIGHTED_SCORE_RETURN_SCHEMA_VERSION
    )
    block = raw["pattern_observations"]
    patterns = block["patterns_detected"]
    antipatterns = block["antipatterns_detected"]
    block["pattern_score"] = return_validation.expected_weighted_score(
        patterns, antipatterns
    )
    block["pattern_score_basis"] = return_validation.pattern_score_basis(
        patterns, antipatterns, block["not_evaluable"]
    )
    return_validation.validate_return(raw, catalog)
    vault, owner = _artifact(tmp_path, transcript_timing)

    canonical = pattern_evidence.canonicalize_return_evidence(
        raw,
        owner,
        vault,
        catalog,
        pattern_scoring_schema_version=(
            return_validation.WEIGHTED_PATTERN_SCORING_SCHEMA_VERSION
        ),
    )

    observations = canonical["pattern_observations"]
    assert isinstance(observations, dict)
    assert observations["applicability_assessments"]
    assert observations["pattern_outcomes"]
    assert observations["evidence_schema_version"] == 2
