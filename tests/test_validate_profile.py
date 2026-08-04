"""Tests for validate-profile.py — required-key + schema_version validation.

Locks in the graceful fallback for the engine-sourcing feature: a profile that
predates `presentation_engines` must still validate, since the field is
optional/additive and not part of REQUIRED_KEYS.
"""

import hashlib
import importlib
import json
import os
from datetime import datetime, timezone
from types import SimpleNamespace

import pytest
from pptx import Presentation


CATALOG = "c" * 64
OTHER_CATALOG = "d" * 64
AS_OF = "2026-07-31T13:30:45.987654-05:00"


def foreign_absolute_locator(name: str) -> str:
    if os.name == "nt":
        return f"/foreign/{name}"
    return rf"C:\foreign\{name}"


DOT_SEGMENT_VAULT_ROOT = (
    r"C:\trusted\other\..\vault"
    if os.name == "nt"
    else "/trusted/other/../vault"
)
INVALID_VAULT_ROOT_LOCATORS = (
    ("", "artifact_locator_empty_or_whitespace"),
    ("   ", "artifact_locator_empty_or_whitespace"),
    ("relative-vault", "artifact_root_not_native_absolute"),
    ("C:vault", "artifact_locator_windows_drive_relative"),
    ("~/vault", "artifact_locator_home_expansion_unsupported"),
    (foreign_absolute_locator("vault"), "artifact_locator_foreign_absolute"),
    (r"\\?\C:\vault", "artifact_locator_windows_device_namespace"),
    (r"\vault", "artifact_locator_windows_current_drive_rooted"),
    (DOT_SEGMENT_VAULT_ROOT, "artifact_locator_dot_segment"),
)


def _catalog_projection(source: str | None = "transcript"):
    """Return exhaustive outcomes for one absence-capable singleton source."""
    catalog = importlib.import_module("pattern_opportunities").load_catalog()
    complete_source = (
        frozenset({source})
        if source in {"transcript", "static_slides"}
        else None
    )
    outcomes = []
    not_evaluable = []
    assessments = []
    for pattern_id, entry in sorted(catalog.entries.items()):
        if not entry.observable:
            continue
        if entry.applicability_evaluable_from is not None:
            if (
                complete_source is None
                or complete_source not in entry.applicability_evaluable_from
            ):
                outcomes.append({"pattern_id": pattern_id, "outcome": "not_evaluable"})
                not_evaluable.append({"pattern_id": pattern_id})
                continue
            assessments.append({"pattern_id": pattern_id, "result": "applicable"})
        outcome = (
            "undetected"
            if complete_source is not None
            and entry.absence_evaluable_from is not None
            and complete_source in entry.absence_evaluable_from
            else "not_evaluable"
        )
        outcomes.append({"pattern_id": pattern_id, "outcome": outcome})
        if outcome == "not_evaluable":
            not_evaluable.append({"pattern_id": pattern_id})
    return catalog, outcomes, not_evaluable, assessments


def _set_projection(talk, source: str | None) -> None:
    _, outcomes, not_evaluable, assessments = _catalog_projection(source)
    observations = talk["pattern_observations"]
    observations.update({
        "pattern_outcomes": outcomes,
        "not_evaluable": not_evaluable,
        "applicability_assessments": assessments,
        "opportunity_coverage_identity": (
            importlib.import_module("pattern_evidence").opportunity_coverage_identity(
                outcomes,
                pattern_catalog_fingerprint=talk["pattern_catalog_fingerprint"],
                pattern_scoring_schema_version=talk[
                    "pattern_scoring_schema_version"
                ],
            )
        ),
    })


# Built programmatically per testing-standards (no fixture file). A current
# profile carries explicit zero-cohort pattern provenance rather than borrowing
# historical pattern values.
def _minimal_profile(validate_profile):
    catalog_fingerprint, scoring_schema = (
        validate_profile.active_pattern_generation_identity()
    )
    opportunities = importlib.import_module("pattern_opportunities")
    provenance = importlib.import_module("profile_pattern_provenance")
    rows = opportunities.build_pattern_opportunity_rows([])
    profile = {
        k: []
        for k in (
            "generated_date",
            "talks_analyzed",
            "speaker",
            "infrastructure",
            "presentation_modes",
            "instrument_catalog",
            "rhetoric_defaults",
            "confirmed_intents",
            "guardrail_sources",
            "pacing",
            "pattern_profile",
            "visual_style_history",
            "publishing_process",
            "design_rules",
            "badges",
        )
    } | {"schema_version": 4}
    profile["pattern_profile"] = {
        "pattern_baseline": {
            "schema_version": 2,
            "as_of": "2025-01-02T03:04:05+00:00",
            "scope": "global",
            "active_batch_excluded": False,
            "excluded_filenames": [],
            "eligible_statuses": ["processed", "processed_partial"],
            "pattern_scoring_generation_status": "current",
            "pattern_scoring_generation_reasons": [],
            "pattern_catalog_fingerprint": catalog_fingerprint,
            "pattern_scoring_schema_version": scoring_schema,
            "scored_talk_count": 0,
            "pattern_score_sum": 0,
            "average_pattern_score": None,
            "eligible_talk_count": 0,
            "opportunity_coverage_identity": None,
            "raw_score_comparison_status": "unavailable",
            "raw_score_comparison_reason": "empty_current_cohort",
        },
        "baseline_talk_filenames": [],
        "eligible_talk_count": 0,
        "talks_scored": 0,
        "average_pattern_score": None,
        "score_trend": "unavailable",
        "pattern_breadth": {
            "avg_distinct_patterns_per_talk": None,
            "trend": "unavailable",
            "note": "No current pattern cohort.",
        },
        "underused_patterns": [],
        "score_drivers": {
            "direction": "unavailable",
            "antipattern_drivers": [],
            "pattern_drivers": [],
            "note": "No current pattern cohort.",
        },
        "by_mode": [],
        "strengths": [],
        "strengths_note": "No current pattern cohort.",
        "note": "Only current observable patterns are included.",
        "pattern_usage": rows["pattern_usage"],
        "antipattern_frequency": rows["antipattern_frequency"],
        "never_used_patterns": [],
        "signature_combinations": [],
        "mastery_levels": {
            "signature": [],
            "regular": [],
            "occasional": [],
            "rare": [],
            "never_tried": [],
        },
        "classification_availability": (
            provenance.unavailable_classification_availability()
        ),
    }
    profile["rhetoric_defaults"] = {}
    profile["guardrail_sources"] = {"recurring_issues": []}
    return profile


def _run(validate_profile, profile, tmp_path, *, database=None):
    (tmp_path / "tracking-database.json").write_text(
        json.dumps(
            {"config": {}, "talks": []} if database is None else database
        ),
        encoding="utf-8",
    )
    path = tmp_path / "profile.json"
    path.write_text(json.dumps(profile))
    rc = validate_profile.main(
        ["validate-profile.py", str(path), "--vault-root", str(tmp_path)]
    )
    return rc


def test_profile_without_engines_still_validates(validate_profile, tmp_path, capsys):
    # The whole point: presentation_engines is optional/additive — a profile that
    # never heard of it is still valid.
    profile = _minimal_profile(validate_profile)
    assert "presentation_engines" not in profile
    rc = _run(validate_profile, profile, tmp_path)
    out = json.loads(capsys.readouterr().out)
    assert rc == 0
    assert out["valid"] is True
    assert out["missing_keys"] == []


def test_validate_profile_rejects_future_tracking_database(
    validate_profile, tmp_path, capsys
):
    profile = _minimal_profile(validate_profile)

    rc = _run(
        validate_profile,
        profile,
        tmp_path,
        database={"schema_version": 99, "config": {}, "talks": []},
    )

    captured = capsys.readouterr()
    out = json.loads(captured.out)
    assert rc == 1
    assert out["valid"] is False
    assert "no usable prior state" in captured.err


def test_profile_with_engines_validates(validate_profile, tmp_path, capsys):
    profile = _minimal_profile(validate_profile)
    profile["presentation_engines"] = [
        {"id": "pptx", "renderer": "pptx", "usage_count": 18, "out_of": 24}
    ]
    rc = _run(validate_profile, profile, tmp_path)
    out = json.loads(capsys.readouterr().out)
    assert rc == 0
    assert out["valid"] is True


def test_profile_missing_required_key_is_invalid(validate_profile, tmp_path, capsys):
    profile = _minimal_profile(validate_profile)
    del profile["design_rules"]
    rc = _run(validate_profile, profile, tmp_path)
    out = json.loads(capsys.readouterr().out)
    assert rc == 1
    assert out["valid"] is False
    assert "design_rules" in out["missing_keys"]


def test_profile_with_outdated_schema_version_is_invalid(
    validate_profile, tmp_path, capsys
):
    # Profiles before v4 have no exact per-pattern opportunity provenance and cannot
    # pass the owner writer's current-generation gate.
    profile = _minimal_profile(validate_profile) | {"schema_version": 1}
    rc = _run(validate_profile, profile, tmp_path)
    out = json.loads(capsys.readouterr().out)
    assert rc == 1
    assert out["valid"] is False
    assert out["schema_version"] == 1


# --- load-vault instrumentation partitioning -------------------------------


def _talks(*dates):
    return [
        {"filename": f"t{i}.md", "status": "processed", "processed_date": d}
        for i, d in enumerate(dates)
    ]


def test_partition_splits_on_the_epoch(load_vault):
    current, stale = load_vault.partition_by_instrumentation(
        _talks("2026-07-27", "2026-07-26", "2026-07-25", "2026-05-01")
    )
    assert [t["processed_date"] for t in current] == ["2026-07-27", "2026-07-26"]
    assert [t["processed_date"] for t in stale] == ["2026-07-25", "2026-05-01"]


def test_undated_talks_are_treated_as_stale(load_vault):
    """Excluding an undated talk only narrows the sample; including it would
    silently contaminate the baseline."""
    talks = _talks("2026-07-27")
    talks.append({"filename": "x.md", "status": "processed"})
    talks.append({"filename": "y.md", "status": "processed", "processed_date": None})
    current, stale = load_vault.partition_by_instrumentation(talks)
    assert len(current) == 1
    assert {t["filename"] for t in stale} == {"x.md", "y.md"}


def test_epoch_is_injectable_for_callers(load_vault):
    current, stale = load_vault.partition_by_instrumentation(
        _talks("2020-01-01"), epoch="2019-01-01"
    )
    assert len(current) == 1 and not stale


# --- load-vault exact pattern-scoring cohort -------------------------------


def _scored_talk(
    filename,
    score,
    *,
    processed_date="2026-07-27",
    status="processed",
    catalog=CATALOG,
    scoring_schema=3,
    generation_status="current",
    generation_reasons=None,
):
    if generation_reasons is None:
        generation_reasons = []
    pattern_evidence = importlib.import_module("pattern_evidence")
    _, outcome_rows, not_evaluable, assessments = _catalog_projection("transcript")
    observations = {
        "pattern_score": score,
        "patterns_detected": [],
        "antipatterns_detected": [],
        "not_evaluable": not_evaluable,
        "applicability_assessments": assessments,
        "pattern_outcomes": outcome_rows,
    }
    if scoring_schema >= 5:
        observations["opportunity_coverage_identity"] = (
            pattern_evidence.opportunity_coverage_identity(
                outcome_rows,
                pattern_catalog_fingerprint=catalog,
                pattern_scoring_schema_version=scoring_schema,
            )
        )
    return {
        "filename": filename,
        "status": status,
        "processed_date": processed_date,
        "pattern_scoring_generation_status": generation_status,
        "pattern_scoring_generation_reasons": generation_reasons,
        "pattern_catalog_fingerprint": catalog,
        "pattern_scoring_schema_version": scoring_schema,
        "pattern_score": score,
        "pattern_observations": observations,
    }


def _write_vault(vault_root, talks, *, config=None):
    transcript_timing = importlib.import_module("transcript_timing")
    transcripts = vault_root / "transcripts"
    transcripts.mkdir(exist_ok=True)
    for talk in talks:
        scoring_version = talk.get("pattern_scoring_schema_version")
        if not isinstance(scoring_version, int) or scoring_version < 4:
            continue
        observations = talk.get("pattern_observations")
        if not isinstance(observations, dict):
            continue
        if observations.get("source_inspection"):
            # A test that supplies a canonical source receipt owns its matching
            # talk metadata. Do not contaminate native-deck or configured-root
            # fixtures with an unrelated synthetic YouTube transcript owner.
            continue
        video_id = hashlib.sha256(talk["filename"].encode("utf-8")).hexdigest()[:11]
        source_duration = 60.0
        artifact = transcripts / f"{video_id}.txt"
        content = ("synthetic evidence " * 225).strip() + "\n"
        transcript_timing.write_transcript_bundle(
            artifact,
            content,
            [{"text": content, "start": 0.0, "end": 10.0}],
            source="captions",
            timing_provenance=transcript_timing.youtube_timing_provenance(
                "captions", video_id, source_duration
            ),
            quality_policy=transcript_timing.build_quality_policy(400),
            quality_policy_provenance={"kind": "fixed_default"},
        )
        quality_artifact = artifact.with_suffix(".quality.json")
        timing_artifact = artifact.with_suffix(".segments.json")
        digest = hashlib.sha256(content.encode("utf-8")).hexdigest()
        timing_digest = hashlib.sha256(timing_artifact.read_bytes()).hexdigest()
        quality_digest = hashlib.sha256(quality_artifact.read_bytes()).hexdigest()
        talk.setdefault("transcript_path", artifact.relative_to(vault_root).as_posix())
        talk["transcript_source"] = "youtube_auto"
        talk["youtube_id"] = video_id
        talk["source_identity"] = {
            "schema_version": 1,
            "provider": "youtube",
            "video_id": video_id,
            "duration_seconds": source_duration,
        }
        observations.setdefault(
            "evidence_schema_version", 2 if scoring_version >= 5 else 1
        )
        observations.setdefault("evidence_sources", ["transcript"])
        observations.setdefault(
            "source_inspection",
            [
                {
                    "source": "transcript",
                    "line_ranges": [[1, 1]],
                    "line_count": 1,
                    "coverage_complete": True,
                    "absence_capability_complete": True,
                    "absence_capability_reason": "authorized_transcript",
                    "artifact_root": "vault",
                    "artifact_path": artifact.relative_to(vault_root).as_posix(),
                    "artifact_sha256": digest,
                    "timing_artifact_root": "vault",
                    "timing_artifact_path": timing_artifact.relative_to(
                        vault_root
                    ).as_posix(),
                    "timing_artifact_sha256": timing_digest,
                    "quality_artifact_root": "vault",
                    "quality_artifact_path": quality_artifact.relative_to(
                        vault_root
                    ).as_posix(),
                    "quality_artifact_sha256": quality_digest,
                }
            ],
        )
        observations.setdefault("patterns_detected", [])
        observations.setdefault("antipatterns_detected", [])
        if scoring_version >= 5:
            observations.setdefault("applicability_assessments", [])
            inspection = observations.get("source_inspection")
            record = (
                inspection[0]
                if isinstance(inspection, list)
                and len(inspection) == 1
                and isinstance(inspection[0], dict)
                and inspection[0].get("source") == "transcript"
                else None
            )
            raw_assessments = observations.get("applicability_assessments")
            if record is not None and isinstance(raw_assessments, list):
                catalog = importlib.import_module("pattern_opportunities").load_catalog()
                located = []
                for assessment in raw_assessments:
                    if not isinstance(assessment, dict):
                        continue
                    entry = catalog.entries[assessment["pattern_id"]]
                    channel = (
                        "transcript"
                        if "transcript" in entry.evidence_channels
                        else "timed_transcript"
                    )
                    citation = {
                        "source": "transcript",
                        "channel": channel,
                        "quote": "synthetic evidence synthetic evidence",
                        "line_start": 1,
                        "line_end": 1,
                    }
                    for identity_field in (
                        "artifact_root", "artifact_path", "artifact_sha256",
                        "timing_artifact_root", "timing_artifact_path",
                        "timing_artifact_sha256", "quality_artifact_root",
                        "quality_artifact_path", "quality_artifact_sha256",
                    ):
                        if identity_field in record:
                            citation[identity_field] = record[identity_field]
                    if channel == "timed_transcript":
                        citation.update({
                            "start_seconds": 0.0,
                            "end_seconds": 10.0,
                        })
                    located.append({
                        **assessment,
                        "evidence_source": "transcript",
                        "evidence": (
                            "The complete transcript establishes applicability."
                        ),
                        "evidence_citations": [citation],
                    })
                observations["applicability_assessments"] = located
    (vault_root / "tracking-database.json").write_text(
        json.dumps(
            {
                "config": config or {"synthetic": True},
                "confirmed_intents": [],
                "talks": talks,
            }
        )
    )
    (vault_root / "rhetoric-style-summary.md").write_text("Synthetic summary\n")


def _run_load_vault(load_vault, vault_root, monkeypatch, capsys):
    active_catalog = load_vault.load_catalog()
    monkeypatch.setattr(
        load_vault,
        "load_catalog",
        lambda: SimpleNamespace(
            fingerprint=CATALOG,
            entries=active_catalog.entries,
        ),
    )
    rc = load_vault.main(["load-vault.py", str(vault_root), "--as-of", AS_OF])
    captured = capsys.readouterr()
    payload = json.loads(captured.out) if captured.out else None
    return rc, payload, captured.err


def test_load_vault_separates_pattern_generation_from_instrumentation(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
):
    exact_old = _scored_talk(
        "current-old-date.md",
        0,
        processed_date="2026-01-01",
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION,
    )
    legacy_recent = _scored_talk(
        "legacy-recent.md",
        100,
        generation_status="legacy_unbaselineable",
    )
    old_catalog_recent = _scored_talk(
        "old-catalog-recent.md",
        90,
        catalog=OTHER_CATALOG,
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION,
    )
    old_schema_recent = _scored_talk(
        "old-schema-recent.md",
        80,
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION - 1,
    )
    missing_generation = _scored_talk(
        "missing-generation.md",
        70,
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION,
    )
    missing_generation.pop("pattern_scoring_generation_status")
    pending_duplicate = _scored_talk(
        "pending-ineligible.md",
        60,
        status="pending",
        generation_status="future",
    )
    _write_vault(
        tmp_path,
        [
            exact_old,
            legacy_recent,
            old_catalog_recent,
            old_schema_recent,
            missing_generation,
            pending_duplicate,
        ],
    )
    rc, payload, error = _run_load_vault(load_vault, tmp_path, monkeypatch, capsys)

    assert rc == 0
    assert error == ""
    assert [talk["filename"] for talk in payload["baseline_talks"]] == [
        "current-old-date.md"
    ], json.dumps(payload["pattern_scoring_exclusions"], indent=2)
    assert [talk["filename"] for talk in payload["excluded_pattern_scoring_talks"]] == [
        "legacy-recent.md",
        "old-catalog-recent.md",
        "old-schema-recent.md",
        "missing-generation.md",
    ]
    assert {
        detail["filename"]: detail["reason_codes"]
        for detail in payload["pattern_scoring_exclusions"]
    } == {
        "legacy-recent.md": ["legacy_generation"],
        "old-catalog-recent.md": ["catalog_fingerprint_mismatch"],
        "old-schema-recent.md": ["scoring_schema_version_mismatch"],
        "missing-generation.md": ["missing_generation_status"],
    }
    assert [talk["filename"] for talk in payload["current_instrumentation_talks"]] == [
        "legacy-recent.md",
        "old-catalog-recent.md",
        "old-schema-recent.md",
        "missing-generation.md",
    ]
    assert [talk["filename"] for talk in payload["stale_instrumentation_talks"]] == [
        "current-old-date.md"
    ]
    assert payload["pattern_baseline"] == {
        "schema_version": 2,
        "as_of": "2026-07-31T18:30:45+00:00",
        "scope": "global",
        "active_batch_excluded": False,
        "excluded_filenames": [],
        "eligible_statuses": ["processed", "processed_partial"],
        "pattern_scoring_generation_status": "current",
        "pattern_scoring_generation_reasons": [],
        "pattern_catalog_fingerprint": CATALOG,
        "pattern_scoring_schema_version": (load_vault.PATTERN_SCORING_SCHEMA_VERSION),
        "scored_talk_count": 1,
        "pattern_score_sum": 0,
        "average_pattern_score": 0.0,
        "eligible_talk_count": 1,
        "opportunity_coverage_identity": exact_old["pattern_observations"][
            "opportunity_coverage_identity"
        ],
        "raw_score_comparison_status": "available",
        "raw_score_comparison_reason": None,
    }
    assert "snapshot observation time" in payload["baseline_note"]
    assert "processed_date" in payload["baseline_note"]
    assert "does not confer pattern-scoring" in payload["instrumentation_note"]


def test_load_vault_rejects_future_tracking_database(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
):
    _write_vault(tmp_path, [])
    database_path = tmp_path / "tracking-database.json"
    database = json.loads(database_path.read_text(encoding="utf-8"))
    database["schema_version"] = 99
    database_path.write_text(json.dumps(database), encoding="utf-8")
    before = database_path.read_bytes()

    rc, payload, error = _run_load_vault(load_vault, tmp_path, monkeypatch, capsys)

    assert rc == 1
    assert payload is None
    assert "not usable by this reader" in error
    assert database_path.read_bytes() == before


def test_load_vault_projects_confirmed_intents_without_storage_metadata(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
):
    _write_vault(tmp_path, [])
    database_path = tmp_path / "tracking-database.json"
    database = {
        "schema_version": 1,
        "config": {"schema_version": 1, "synthetic": True},
        "talks": [],
        "pptx_catalog": [],
        "qr_codes": [],
        "resources": [],
        "thumbnails": [],
        "confirmed_intents": [{
            "schema_version": 1,
            "pattern": "delayed_self_introduction",
            "intent": "deliberate",
            "rule": "Use the two-phase introduction",
            "note": "Speaker-confirmed",
            "confirmed_date": "2026-08-01",
            "source_talk": "example.md",
        }],
        "improvement_goals": [],
    }
    raw = (json.dumps(database, indent=2) + "\n").encode()
    database_path.write_bytes(raw)

    rc, payload, error = _run_load_vault(
        load_vault, tmp_path, monkeypatch, capsys
    )

    assert rc == 0
    assert error == ""
    assert payload["confirmed_intents"] == [{
        "pattern": "delayed_self_introduction",
        "intent": "deliberate",
        "rule": "Use the two-phase introduction",
        "note": "Speaker-confirmed",
    }]
    assert list(payload["confirmed_intents"][0]) == [
        "pattern", "intent", "rule", "note",
    ]
    assert database_path.read_bytes() == raw


@pytest.mark.parametrize(
    ("case", "message"),
    [
        ("unknown_status", "pattern_scoring_generation_status must be one of"),
        ("current_reasons", "pattern_scoring_generation_reasons must be exactly"),
        ("missing_fingerprint", "missing required identity fields"),
        ("malformed_fingerprint", "must be a lowercase 64-character"),
        ("divergent_score", "promoted pattern_score 5 diverges"),
    ],
)
def test_load_vault_rejects_malformed_current_generation_records(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
    case,
    message,
):
    talk = _scored_talk(
        "malformed.md",
        5,
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION,
    )
    if case == "unknown_status":
        talk["pattern_scoring_generation_status"] = "future"
    elif case == "current_reasons":
        talk["pattern_scoring_generation_reasons"] = ["contradiction"]
    elif case == "missing_fingerprint":
        talk.pop("pattern_catalog_fingerprint")
    elif case == "malformed_fingerprint":
        talk["pattern_catalog_fingerprint"] = "not-a-sha"
    elif case == "divergent_score":
        talk["pattern_observations"]["pattern_score"] = 4
    _write_vault(tmp_path, [talk])

    rc, payload, error = _run_load_vault(load_vault, tmp_path, monkeypatch, capsys)

    assert rc == 1
    assert payload is None
    assert message in error


def test_load_vault_does_not_fallback_when_current_pattern_cohort_is_empty(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
):
    legacy = _scored_talk(
        "legacy.md",
        100,
        generation_status="legacy_unbaselineable",
    )
    _write_vault(tmp_path, [legacy])

    rc, payload, error = _run_load_vault(load_vault, tmp_path, monkeypatch, capsys)

    assert rc == 0
    assert error == ""
    assert payload["baseline_talks"] == []
    assert [talk["filename"] for talk in payload["excluded_pattern_scoring_talks"]] == [
        "legacy.md"
    ]
    assert payload["pattern_baseline"]["scored_talk_count"] == 0
    assert payload["pattern_baseline"]["pattern_score_sum"] == 0
    assert payload["pattern_baseline"]["average_pattern_score"] is None


@pytest.mark.parametrize("drift", ["missing", "digest_mismatch"])
def test_load_vault_excludes_current_generation_with_stale_evidence(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
    drift,
):
    talk = _scored_talk(
        "artifact-drift.md",
        3,
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION,
    )
    _write_vault(tmp_path, [talk])
    relative = talk["pattern_observations"]["source_inspection"][0]["artifact_path"]
    artifact = tmp_path / relative
    if drift == "missing":
        artifact.unlink()
    else:
        artifact.write_text("Replacement with a different digest.\n", encoding="utf-8")

    rc, payload, error = _run_load_vault(load_vault, tmp_path, monkeypatch, capsys)

    assert rc == 0
    assert error == ""
    assert payload["baseline_talks"] == []
    assert [talk["filename"] for talk in payload["excluded_pattern_scoring_talks"]] == [
        "artifact-drift.md"
    ]
    detail = payload["pattern_scoring_exclusions"][0]
    assert detail["reason_codes"] == ["persisted_evidence_stale"]
    assert any(drift in reason for reason in detail["evidence_freshness_details"])


def test_load_vault_passes_configured_source_roots_to_freshness_assessor(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
):
    source_root = tmp_path / "configured-pptx"
    source_root.mkdir()
    deck = source_root / "synthetic-deck.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(deck)
    digest = hashlib.sha256(deck.read_bytes()).hexdigest()
    talk = _scored_talk(
        "configured-root.md",
        0,
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION,
    )
    talk["pptx_path"] = deck.name
    talk["slide_source"] = "pptx"
    talk["structured_data"] = {
        "slide_count": 1,
        "native_deck_audit": importlib.import_module(
            "pptx_evidence"
        ).recompute_native_deck_audit(
            deck,
            trusted_root=source_root,
        ),
    }
    talk["pattern_observations"].update(
        {
            "evidence_schema_version": 2,
            "evidence_sources": ["native_deck"],
            "source_inspection": [
                {
                    "source": "native_deck",
                    "page_ranges": [[1, 1]],
                    "page_count": 1,
                    "coverage_complete": True,
                    "absence_capability_complete": False,
                    "absence_capability_reason": "bare_native_deck",
                    "artifact_root": "pptx_source",
                    "artifact_path": deck.name,
                    "artifact_sha256": digest,
                }
            ],
            "patterns_detected": [],
            "antipatterns_detected": [],
        }
    )
    _set_projection(talk, None)
    _write_vault(
        tmp_path,
        [talk],
        config={"pptx_source_dir": str(source_root)},
    )

    rc, payload, error = _run_load_vault(load_vault, tmp_path, monkeypatch, capsys)

    assert rc == 0
    assert error == ""
    assert [talk["filename"] for talk in payload["baseline_talks"]] == [
        "configured-root.md"
    ], json.dumps(payload["pattern_scoring_exclusions"], indent=2)
    assert payload["pattern_scoring_exclusions"] == []


def test_load_vault_rejects_mismatched_vault_storage_root_before_freshness(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
):
    database_root = tmp_path / "database-root"
    database_root.mkdir()
    evidence_root = tmp_path / "evidence-root"
    transcript = evidence_root / "transcripts" / "configured.txt"
    transcript.parent.mkdir(parents=True)
    content = (
        ("synthetic evidence " * 225).strip()
        + "\n"
        + ("additional uninspected evidence " * 20).strip()
        + "\n"
    )
    transcript.write_text(content, encoding="utf-8")
    transcript_timing = importlib.import_module("transcript_timing")
    transcript_timing.write_quality_receipt(
        transcript,
        content,
        transcript_timing.build_quality_policy(400),
        {"kind": "fixed_default"},
    )
    quality_artifact = transcript.with_suffix(".quality.json")
    talk = _scored_talk(
        "configured-vault.md",
        0,
        scoring_schema=load_vault.PATTERN_SCORING_SCHEMA_VERSION,
    )
    talk["transcript_path"] = "transcripts/configured.txt"
    talk["transcript_source"] = "manual"
    talk["pattern_observations"].update(
        {
            "evidence_schema_version": 2,
            "evidence_sources": ["transcript"],
            "source_inspection": [
                {
                    "source": "transcript",
                    "line_ranges": [[1, 1]],
                    "line_count": 2,
                    "coverage_complete": False,
                    "absence_capability_complete": False,
                    "absence_capability_reason": "incomplete_range_coverage",
                    "artifact_root": "vault",
                    "artifact_path": "transcripts/configured.txt",
                    "artifact_sha256": hashlib.sha256(
                        content.encode("utf-8")
                    ).hexdigest(),
                    "quality_artifact_root": "vault",
                    "quality_artifact_path": "transcripts/configured.quality.json",
                    "quality_artifact_sha256": hashlib.sha256(
                        quality_artifact.read_bytes()
                    ).hexdigest(),
                }
            ],
            "patterns_detected": [],
            "antipatterns_detected": [],
        }
    )
    _set_projection(talk, None)
    _write_vault(
        database_root,
        [talk],
        config={"vault_storage_path": str(evidence_root)},
    )
    cohort_snapshot = importlib.import_module("pattern_cohort_snapshot")
    monkeypatch.setattr(
        cohort_snapshot,
        "assess_current_persisted_pattern_evidence_freshness",
        lambda *_args, **_kwargs: pytest.fail(
            "mismatched root reached the evidence freshness assessor"
        ),
    )

    rc, payload, error = _run_load_vault(load_vault, database_root, monkeypatch, capsys)

    assert rc == 1
    assert payload is None
    assert "vault_root_authority_mismatch:database_path:config_root" in error
    assert str(evidence_root) not in error


def test_load_vault_default_is_an_already_native_absolute_path(load_vault):
    vault_root, as_of = load_vault._parse_args(["load-vault.py"])

    assert as_of is None
    assert vault_root == load_vault.DEFAULT_VAULT
    assert vault_root.is_absolute()
    assert "~" not in str(vault_root)


@pytest.mark.parametrize(
    ("cli_root", "locator_reason"),
    INVALID_VAULT_ROOT_LOCATORS,
)
def test_load_vault_rejects_invalid_explicit_root_before_any_vault_io(
    load_vault,
    monkeypatch,
    capsys,
    cli_root,
    locator_reason,
):
    io_calls = []

    def forbidden_io(*_args, **_kwargs):
        io_calls.append("vault_io")
        pytest.fail("invalid CLI root reached vault I/O")

    monkeypatch.setattr(load_vault.pathlib.Path, "exists", forbidden_io)
    monkeypatch.setattr(load_vault, "snapshot_tracking_database", forbidden_io)

    rc = load_vault.main(["load-vault.py", cli_root])
    captured = capsys.readouterr()

    assert rc == 1
    assert captured.out == ""
    assert (
        f"vault_root_cli_invalid:{locator_reason}" in captured.err
    )
    assert io_calls == []
    if cli_root.strip():
        assert cli_root not in captured.err


@pytest.mark.parametrize(
    ("configured_root", "locator_reason"),
    INVALID_VAULT_ROOT_LOCATORS,
)
def test_load_vault_rejects_invalid_config_before_summary_or_freshness_io(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
    configured_root,
    locator_reason,
):
    (tmp_path / "tracking-database.json").write_text(
        json.dumps(
            {
                "config": {"vault_storage_path": configured_root},
                "talks": [],
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        load_vault,
        "configured_evidence_freshness_assessor",
        lambda *_args, **_kwargs: pytest.fail(
            "invalid configured root reached freshness assessment"
        ),
    )

    rc = load_vault.main(["load-vault.py", str(tmp_path), "--as-of", AS_OF])
    captured = capsys.readouterr()

    assert rc == 1
    assert captured.out == ""
    assert not (tmp_path / "rhetoric-style-summary.md").exists()
    assert (
        f"vault_root_config_invalid:{locator_reason}" in captured.err
    )
    if configured_root.strip():
        assert configured_root not in captured.err


def _directory_symlink(link, target):
    try:
        link.symlink_to(target, target_is_directory=True)
    except OSError:
        pytest.skip("directory policy does not permit symlink creation")


def test_load_vault_accepts_one_matching_lexical_symlink_authority(
    load_vault,
    tmp_path,
    monkeypatch,
    capsys,
):
    storage = tmp_path / "storage"
    storage.mkdir()
    locator = tmp_path / "vault-alias"
    _directory_symlink(locator, storage)
    _write_vault(storage, [], config={"vault_storage_path": str(locator)})

    rc, payload, error = _run_load_vault(
        load_vault,
        locator,
        monkeypatch,
        capsys,
    )

    assert rc == 0, error
    assert payload["vault_root"] == str(locator)


def test_load_vault_rejects_symlink_target_and_locator_authority_mismatch(
    load_vault,
    tmp_path,
    capsys,
):
    storage = tmp_path / "credential-bearing-storage"
    storage.mkdir()
    locator = tmp_path / "credential-bearing-alias"
    _directory_symlink(locator, storage)
    (storage / "tracking-database.json").write_text(
        json.dumps(
            {
                "config": {"vault_storage_path": str(storage)},
                "talks": [],
            }
        ),
        encoding="utf-8",
    )

    rc = load_vault.main(["load-vault.py", str(locator), "--as-of", AS_OF])
    captured = capsys.readouterr()

    assert rc == 1
    assert captured.out == ""
    assert (
        "vault_root_authority_mismatch:database_path:config_root"
        in captured.err
    )
    assert str(storage) not in captured.err
    assert str(locator) not in captured.err


@pytest.mark.parametrize(
    ("configured_root", "locator_reason"),
    INVALID_VAULT_ROOT_LOCATORS,
)
def test_profile_cohort_rejects_invalid_configured_vault_before_freshness(
    validate_profile,
    tmp_path,
    monkeypatch,
    configured_root,
    locator_reason,
):
    del validate_profile
    cohort_snapshot = importlib.import_module("pattern_cohort_snapshot")
    monkeypatch.setattr(
        cohort_snapshot,
        "assess_current_persisted_pattern_evidence_freshness",
        lambda *_args, **_kwargs: pytest.fail(
            "invalid root reached the evidence freshness assessor"
        ),
    )
    monkeypatch.setattr(
        cohort_snapshot,
        "VideoEvidenceAssessment",
        lambda: pytest.fail("invalid root created a video evidence assessment"),
    )

    with pytest.raises(cohort_snapshot.PatternCohortSnapshotError) as caught:
        cohort_snapshot.configured_evidence_freshness_assessor(
            tmp_path,
            {"vault_storage_path": configured_root},
        )

    assert str(caught.value) == f"vault_root_config_invalid:{locator_reason}"
    if configured_root.strip():
        assert configured_root not in str(caught.value)


@pytest.mark.parametrize(
    "config_factory",
    [
        lambda _root: {},
        lambda _root: {"vault_storage_path": None},
        lambda root: {"vault_storage_path": str(root)},
    ],
    ids=["absent", "null", "matching"],
)
def test_profile_cohort_uses_database_root_for_every_valid_authority_form(
    validate_profile,
    tmp_path,
    monkeypatch,
    config_factory,
):
    del validate_profile
    cohort_snapshot = importlib.import_module("pattern_cohort_snapshot")
    observed = []
    assessment = object()
    created = []

    def create_assessment():
        created.append(assessment)
        return assessment

    def assess(_talk, **kwargs):
        observed.append(kwargs)
        return ()

    monkeypatch.setattr(
        cohort_snapshot,
        "assess_current_persisted_pattern_evidence_freshness",
        assess,
    )
    monkeypatch.setattr(
        cohort_snapshot,
        "VideoEvidenceAssessment",
        create_assessment,
    )
    config = config_factory(tmp_path)

    assessor = cohort_snapshot.configured_evidence_freshness_assessor(
        tmp_path,
        config,
    )

    first_talk = {}
    second_talk = {"filename": "second.md"}
    assert assessor(first_talk) == ()
    assert assessor(second_talk) == ()
    assert created == [assessment]
    assert observed == [
        {
            "vault_root": tmp_path,
            "source_roots": config,
            "catalog": None,
            "video_evidence_assessment": assessment,
        },
        {
            "vault_root": tmp_path,
            "source_roots": config,
            "catalog": None,
            "video_evidence_assessment": assessment,
        },
    ]


@pytest.mark.parametrize(
    ("cli_root", "locator_reason"),
    INVALID_VAULT_ROOT_LOCATORS,
)
def test_validate_profile_rejects_invalid_cli_root_before_profile_or_database_io(
    validate_profile,
    monkeypatch,
    capsys,
    cli_root,
    locator_reason,
):
    io_calls = []

    def forbidden_io(*_args, **_kwargs):
        io_calls.append("profile_or_database_io")
        pytest.fail("invalid CLI root reached profile or database I/O")

    monkeypatch.setattr(validate_profile, "_load_input", forbidden_io)
    monkeypatch.setattr(validate_profile, "snapshot_tracking_database", forbidden_io)

    rc = validate_profile.main(
        ["validate-profile.py", "--vault-root", cli_root]
    )
    captured = capsys.readouterr()

    assert rc == 1
    assert io_calls == []
    assert f"vault_root_cli_invalid:{locator_reason}" in captured.err
    if cli_root.strip():
        assert cli_root not in captured.err


@pytest.mark.parametrize(
    ("configured_root", "locator_reason"),
    INVALID_VAULT_ROOT_LOCATORS,
)
def test_validate_profile_rejects_invalid_config_before_freshness(
    validate_profile,
    tmp_path,
    monkeypatch,
    capsys,
    configured_root,
    locator_reason,
):
    profile = _minimal_profile(validate_profile)
    profile_path = tmp_path / "profile.json"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    (tmp_path / "tracking-database.json").write_text(
        json.dumps(
            {
                "config": {"vault_storage_path": configured_root},
                "talks": [],
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        validate_profile,
        "configured_evidence_freshness_assessor",
        lambda *_args, **_kwargs: pytest.fail(
            "invalid configured root reached freshness assessment"
        ),
    )

    rc = validate_profile.main(
        [
            "validate-profile.py",
            str(profile_path),
            "--vault-root",
            str(tmp_path),
        ]
    )
    captured = capsys.readouterr()

    assert rc == 1
    assert f"vault_root_config_invalid:{locator_reason}" in captured.err
    if configured_root.strip():
        assert configured_root not in captured.err


def test_validate_profile_accepts_matching_symlink_lexical_authority(
    validate_profile,
    tmp_path,
    capsys,
):
    storage = tmp_path / "storage"
    storage.mkdir()
    locator = tmp_path / "vault-alias"
    _directory_symlink(locator, storage)
    profile = _minimal_profile(validate_profile)
    profile_path = storage / "profile.json"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    (storage / "tracking-database.json").write_text(
        json.dumps(
            {
                "config": {"vault_storage_path": str(locator)},
                "talks": [],
            }
        ),
        encoding="utf-8",
    )

    rc = validate_profile.main(
        [
            "validate-profile.py",
            str(profile_path),
            "--vault-root",
            str(locator),
        ]
    )
    captured = capsys.readouterr()

    assert rc == 0, captured.err
    assert json.loads(captured.out)["valid"] is True


def test_validate_profile_rejects_symlink_target_locator_mismatch_without_paths(
    validate_profile,
    tmp_path,
    capsys,
):
    storage = tmp_path / "credential-bearing-storage"
    storage.mkdir()
    locator = tmp_path / "credential-bearing-alias"
    _directory_symlink(locator, storage)
    profile_path = storage / "profile.json"
    profile_path.write_text(
        json.dumps(_minimal_profile(validate_profile)),
        encoding="utf-8",
    )
    (storage / "tracking-database.json").write_text(
        json.dumps(
            {
                "config": {"vault_storage_path": str(storage)},
                "talks": [],
            }
        ),
        encoding="utf-8",
    )

    rc = validate_profile.main(
        [
            "validate-profile.py",
            str(profile_path),
            "--vault-root",
            str(locator),
        ]
    )
    captured = capsys.readouterr()

    assert rc == 1
    assert (
        "vault_root_authority_mismatch:database_path:config_root"
        in captured.err
    )
    assert str(storage) not in captured.err
    assert str(locator) not in captured.err


def test_validate_profile_catches_cohort_authority_error_without_raw_root_prefix(
    validate_profile,
    tmp_path,
    monkeypatch,
    capsys,
):
    profile_path = tmp_path / "profile.json"
    profile_path.write_text(
        json.dumps(_minimal_profile(validate_profile)),
        encoding="utf-8",
    )
    (tmp_path / "tracking-database.json").write_text(
        json.dumps({"config": {}, "talks": []}),
        encoding="utf-8",
    )

    def fail_cohort(*_args, **_kwargs):
        raise validate_profile.PatternCohortSnapshotError(
            "vault_root_config_invalid:artifact_locator_dot_segment"
        )

    monkeypatch.setattr(
        validate_profile,
        "configured_evidence_freshness_assessor",
        fail_cohort,
    )

    rc = validate_profile.main(
        [
            "validate-profile.py",
            str(profile_path),
            "--vault-root",
            str(tmp_path),
        ]
    )
    captured = capsys.readouterr()

    assert rc == 1
    assert "vault_root_config_invalid:artifact_locator_dot_segment" in captured.err
    assert str(tmp_path) not in captured.err
    assert "Traceback" not in captured.err


def test_default_as_of_uses_injected_time_and_canonical_whole_seconds(load_vault):
    now = datetime(2026, 7, 31, 18, 30, 45, 987654, tzinfo=timezone.utc)

    assert load_vault.default_as_of(now) == "2026-07-31T18:30:45+00:00"
