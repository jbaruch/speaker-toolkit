"""Focused regressions for PPTX recovery and exact audit receipts."""

from __future__ import annotations

import ast
import json
import struct
import subprocess
import sys
import zipfile
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches
from pypdf import PdfWriter


def _write_deck(path: Path, *, with_image: bool) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    if with_image:
        image_path = path.with_suffix(".png")
        Image.new("RGB", (64, 64), "navy").save(image_path)
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    deck.save(path)


def _damage_member(path: Path, predicate) -> str:
    with zipfile.ZipFile(path) as archive:
        member = next(item for item in archive.infolist() if predicate(item.filename))
    package = bytearray(path.read_bytes())
    name_size, extra_size = struct.unpack_from(
        "<HH", package, member.header_offset + 26
    )
    payload_offset = member.header_offset + 30 + name_size + extra_size
    package[payload_offset + max(member.compress_size // 2, 0)] ^= 0xFF
    path.write_bytes(package)
    return member.filename


def _damage_member_crc(path: Path, predicate) -> str:
    """Alter only the central-directory CRC so reading raises BadZipFile."""
    with zipfile.ZipFile(path) as archive:
        member = next(item for item in archive.infolist() if predicate(item.filename))
    package = bytearray(path.read_bytes())
    cursor = 0
    while True:
        header = package.find(b"PK\x01\x02", cursor)
        if header < 0:
            raise AssertionError(
                f"central directory entry not found: {member.filename}"
            )
        name_size, extra_size, comment_size = struct.unpack_from(
            "<HHH", package, header + 28
        )
        name_start = header + 46
        name_end = name_start + name_size
        if package[name_start:name_end].decode("utf-8") == member.filename:
            recorded_crc = struct.unpack_from("<I", package, header + 16)[0]
            struct.pack_into("<I", package, header + 16, recorded_crc ^ 0xFFFFFFFF)
            path.write_bytes(package)
            return member.filename
        cursor = name_end + extra_size + comment_size


def _write_pdf(path: Path, *, page_count: int) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    writer = PdfWriter()
    for _index in range(page_count):
        writer.add_blank_page(width=640, height=480)
    with path.open("wb") as stream:
        writer.write(stream)


def _set_background_image(slide, image_path: Path) -> None:
    _, relationship_id = slide.part.get_or_add_image_part(str(image_path))
    background = slide.element.cSld.get_or_add_bgPr()
    existing_fill = background.eg_fillProperties
    if existing_fill is not None:
        background.remove(existing_fill)
    background.insert(
        0,
        parse_xml(
            f"<a:blipFill {nsdecls('a', 'r')}>"
            f'<a:blip r:embed="{relationship_id}"/>'
            "<a:stretch><a:fillRect/></a:stretch>"
            "</a:blipFill>"
        ),
    )


def _worker_result(payload: dict[str, object]) -> SimpleNamespace:
    return SimpleNamespace(
        payload=payload,
        observed_generations={},
        diagnostics=None,
    )


@pytest.fixture(autouse=True)
def _bounded_metadata_stub(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Keep parser tests focused while preserving the signed metadata shape."""

    def invoke(_command, payload, _sensitive_values, _limits):
        path = Path(payload["pptx_path"])
        root_value = payload["trusted_root"]
        root = Path(root_value) if isinstance(root_value, str) else None
        try:
            generation, root_generation, reparse_tag = (
                pptx_evidence._metadata_generation_in_worker(
                    path,
                    trusted_root=root,
                )
            )
        except pptx_evidence.PptxEvidenceError as exc:
            response = {
                "schema_version": 1,
                "status": "unavailable",
                "reason_code": exc.reason_code,
                "details": dict(exc.details),
            }
        else:
            response = {
                "schema_version": 1,
                "status": "available",
                "generation": generation.to_dict(),
                "root_generation": (
                    root_generation.to_dict() if root_generation is not None else None
                ),
                "reparse_tag": reparse_tag,
            }
        return _worker_result(response)

    monkeypatch.setattr(pptx_evidence, "_invoke_metadata_worker", invoke)


def test_metadata_generation_never_calls_owner_lstat(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "metadata-only-child.pptx"
    _write_deck(deck, with_image=False)
    generation = pptx_evidence.FileGeneration.from_stat(deck.stat())
    captured: dict[str, object] = {}

    def metadata_child(command, payload, sensitive_values, _limits):
        captured.update(
            command=command,
            payload=payload,
            sensitive_values=sensitive_values,
        )
        return _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "generation": generation.to_dict(),
                "root_generation": None,
                "reparse_tag": None,
            }
        )

    monkeypatch.setattr(pptx_evidence, "_invoke_metadata_worker", metadata_child)
    monkeypatch.setattr(
        Path,
        "lstat",
        lambda _path: pytest.fail("owner process called lstat"),
    )

    observed = pptx_evidence._supervised_file_generation(
        deck,
        label="PPTX artifact",
    )

    assert observed == generation
    assert captured["payload"] == {
        "pptx_path": str(deck),
        "trusted_root": None,
    }
    command_text = "\n".join(str(part) for part in captured["command"])
    assert str(deck) not in command_text
    assert captured["sensitive_values"] == (deck,)


def test_unhashable_worker_reason_enums_fail_as_malformed_results(
    pptx_evidence,
) -> None:
    metadata = {
        "schema_version": 1,
        "status": "unavailable",
        "reason_code": "pptx_artifact_unavailable",
        "details": {"failure_kind": []},
    }
    unavailable = {
        "schema_version": 1,
        "status": "unavailable",
        "reason_code": [],
        "details": {},
    }

    calls = (
        lambda: pptx_evidence._decode_metadata_payload(metadata),
        lambda: pptx_evidence._decode_pptx_probe_result(
            (json.dumps(unavailable) + "\n").encode("utf-8")
        ),
        lambda: pptx_evidence._decode_native_audit_result(
            (json.dumps(unavailable) + "\n").encode("utf-8")
        ),
        lambda: pptx_evidence._decode_extraction_worker_payload(unavailable),
    )
    for call in calls:
        with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
            call()
        assert caught.value.reason_code == "pptx_probe_malformed_result"


@pytest.mark.parametrize(
    ("supervisor_reason", "evidence_reason"),
    [
        ("worker_exit_before_barrier", "pptx_probe_start_failure"),
        ("worker_request_write_failed", "pptx_probe_start_failure"),
        ("worker_output_read_failed", "pptx_probe_crash"),
        ("worker_diagnostic_read_failed", "pptx_probe_crash"),
        ("worker_generation_binding_mismatch", "pptx_artifact_changed"),
    ],
)
def test_supervisor_failures_keep_stable_evidence_distinctions(
    pptx_evidence,
    supervisor_reason: str,
    evidence_reason: str,
) -> None:
    error = pptx_evidence._supervisor_probe_failure(
        pptx_evidence.SupervisorError(supervisor_reason),
        timeout_seconds=1,
    )

    assert error.reason_code == evidence_reason


def test_blocked_metadata_child_timeout_is_bounded_and_path_free(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "blocked-metadata.pptx"

    def timeout(*_args, **_kwargs):
        raise pptx_evidence.SupervisorError("worker_timeout")

    monkeypatch.setattr(pptx_evidence, "_invoke_metadata_worker", timeout)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._supervised_file_generation(deck, label="PPTX artifact")

    assert caught.value.reason_code == "pptx_probe_timeout"
    assert caught.value.details == {
        "timeout_seconds": pptx_evidence.PPTX_METADATA_LIMITS.wall_seconds
    }
    assert str(deck) not in str(caught.value)


def test_metadata_payload_accepts_hydrated_cloud_tag(
    pptx_evidence,
) -> None:
    generation = pptx_evidence.FileGeneration(
        size=123,
        mtime_ns=2,
        ctime_ns=3,
        device=4,
        inode=5,
        mode=0o100644,
        flags=0,
        file_attributes=pptx_evidence.PPTX_WINDOWS_REPARSE_POINT_ATTRIBUTE,
    )
    tag = min(pptx_evidence.PPTX_WINDOWS_CLOUD_REPARSE_TAGS)

    receipt = pptx_evidence._decode_metadata_payload(
        {
            "schema_version": 1,
            "status": "available",
            "generation": generation.to_dict(),
            "root_generation": None,
            "reparse_tag": tag,
        }
    )

    assert receipt.generation == generation
    assert receipt.reparse_tag == tag


def test_metadata_payload_preserves_offline_cloud_bits_for_classification(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    offline = 0x001000
    attributes = pptx_evidence.PPTX_WINDOWS_REPARSE_POINT_ATTRIBUTE | offline
    generation = pptx_evidence.FileGeneration(
        size=123,
        mtime_ns=2,
        ctime_ns=3,
        device=4,
        inode=5,
        mode=0o100644,
        flags=0,
        file_attributes=attributes,
    )
    tag = min(pptx_evidence.PPTX_WINDOWS_CLOUD_REPARSE_TAGS)
    monkeypatch.setattr(
        pptx_evidence,
        "PPTX_WINDOWS_CLOUD_FILE_ATTRIBUTES",
        offline,
    )

    receipt = pptx_evidence._decode_metadata_payload(
        {
            "schema_version": 1,
            "status": "available",
            "generation": generation.to_dict(),
            "root_generation": None,
            "reparse_tag": tag,
        }
    )

    assert receipt.generation == generation
    assert pptx_evidence._generation_cloud_placeholder_details(receipt.generation) == {
        "file_attributes": attributes
    }


def test_metadata_payload_rejects_unknown_redirecting_reparse_tag(
    pptx_evidence,
) -> None:
    generation = pptx_evidence.FileGeneration(
        size=123,
        mtime_ns=2,
        ctime_ns=3,
        device=4,
        inode=5,
        mode=0o100644,
        flags=0,
        file_attributes=pptx_evidence.PPTX_WINDOWS_REPARSE_POINT_ATTRIBUTE,
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_metadata_payload(
            {
                "schema_version": 1,
                "status": "available",
                "generation": generation.to_dict(),
                "root_generation": None,
                "reparse_tag": 0xA000000C,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_parser_worker_rejects_trusted_root_generation_swap(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    root = tmp_path / "root"
    deck = root / "deck.pptx"
    _write_deck(deck, with_image=False)
    generation = pptx_evidence.FileGeneration.from_stat(deck.stat())
    root_generation = pptx_evidence.FileGeneration.from_stat(root.stat())
    replaced_root = pptx_evidence.replace(
        root_generation,
        ctime_ns=root_generation.ctime_ns + 1,
    )
    root_reads = iter((root_generation, replaced_root))
    monkeypatch.setattr(
        pptx_evidence,
        "_worker_generation",
        lambda _path, *, trusted_root=None: generation,
    )
    monkeypatch.setattr(
        pptx_evidence,
        "_worker_root_generation",
        lambda _path: next(root_reads),
    )
    monkeypatch.setattr(
        pptx_evidence,
        "_pptx_probe_child",
        lambda _path: {
            "schema_version": 1,
            "status": "available",
            "slide_count": 1,
            "source_sha256": "a" * 64,
            "source_size_bytes": generation.size,
            "archive_recovery": [],
        },
    )
    request = SimpleNamespace(
        operation=pptx_evidence.PPTX_PROBE_OPERATION,
        limit_profile_id=pptx_evidence.PPTX_PROBE_LIMITS.profile_id,
        schema_generation=pptx_evidence.PPTX_EXTRACTION_SCHEMA_VERSION,
        pipeline_generation=pptx_evidence.PPTX_EXTRACTION_PIPELINE_VERSION,
        expected_generations={
            "pptx": generation,
            "pptx_root": root_generation,
        },
        payload={"pptx_path": str(deck), "trusted_root": str(root)},
    )

    with pytest.raises(pptx_evidence.SupervisorError) as caught:
        pptx_evidence._dispatch_supervised_worker(request)

    assert caught.value.reason_code == "worker_generation_changed"


def test_extraction_worker_rejects_malformed_request_before_generation_io(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        pptx_evidence,
        "_worker_generation",
        lambda *_args, **_kwargs: pytest.fail("malformed request touched artifact"),
    )
    monkeypatch.setattr(
        pptx_evidence,
        "_worker_root_generation",
        lambda *_args, **_kwargs: pytest.fail("malformed request touched root"),
    )
    request = SimpleNamespace(
        operation=pptx_evidence.PPTX_EXTRACT_OPERATION,
        limit_profile_id=pptx_evidence.PPTX_EXTRACT_NO_OCR_LIMITS.profile_id,
        schema_generation=pptx_evidence.PPTX_EXTRACTION_SCHEMA_VERSION,
        pipeline_generation=pptx_evidence.PPTX_EXTRACTION_PIPELINE_VERSION,
        expected_generations={},
        payload={
            "pptx_path": "/untrusted/deck.pptx",
            "trusted_root": None,
            "ocr": False,
            "rendered_pdf_path": None,
            "inspected_page_ranges": [],
            "extraction_schema_version": (pptx_evidence.PPTX_EXTRACTION_SCHEMA_VERSION),
            "extraction_pipeline_version": (
                pptx_evidence.PPTX_EXTRACTION_PIPELINE_VERSION
            ),
            "unexpected": True,
        },
    )

    with pytest.raises(pptx_evidence.SupervisorError) as caught:
        pptx_evidence._dispatch_supervised_worker(request)

    assert caught.value.reason_code == "invalid_worker_request"


@pytest.mark.skipif(
    sys.platform == "win32",
    reason="this deterministic swap uses POSIX directory symlinks",
)
def test_metadata_worker_rejects_lexical_trusted_root_swapped_to_symlink(
    pptx_evidence,
    tmp_path: Path,
) -> None:
    root = tmp_path / "root"
    root.mkdir()
    _write_deck(root / "deck.pptx", with_image=False)
    outside = tmp_path / "outside"
    outside.mkdir()
    _write_deck(outside / "deck.pptx", with_image=False)
    parked = tmp_path / "original-root"
    root.rename(parked)
    root.symlink_to(outside, target_is_directory=True)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._metadata_generation_in_worker(
            root / "deck.pptx",
            trusted_root=root,
        )

    assert caught.value.reason_code == "pptx_artifact_unavailable"
    assert caught.value.details["failure_kind"] == "root_escape"


def test_live_badzipfile_media_path_recovers_with_structured_loss(
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "damaged-media.pptx"
    _write_deck(deck, with_image=True)
    damaged_part = _damage_member_crc(deck, lambda name: name.startswith("ppt/media/"))

    with pytest.raises(zipfile.BadZipFile):
        Presentation(deck)

    probe = pptx_evidence.probe_pptx_artifact(deck)

    assert probe.slide_count == 1
    assert len(probe.source_sha256) == 64
    assert probe.archive_recovery == (
        {
            "schema_version": 1,
            "part_name": damaged_part,
            "member_kind": "embedded_media",
            "error_type": "crc_mismatch",
            "status": "recovered_with_placeholder_asset",
            "content_replaced": True,
            "replacement_sha256": (
                "431ced6916a2a21a156e38701afe55bbd7f88969fbbfc56d7fe099d47f265460"
            ),
        },
    )


def test_bounded_probe_timeout_is_structured_and_never_cached(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "slow.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()
    calls = 0

    def timeout_runner(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        raise pptx_evidence.SupervisorError("worker_timeout")

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", timeout_runner)
    for _index in range(2):
        with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
            pptx_evidence.probe_pptx_artifact(deck)
        assert caught.value.reason_code == "pptx_probe_timeout"
    assert calls == 2


def test_bounded_probe_resource_failure_is_structured(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "memory-heavy.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()

    def resource_runner(*_args, **_kwargs):
        return _worker_result(
            {
                "schema_version": 1,
                "status": "unavailable",
                "reason_code": "pptx_probe_resource_unavailable",
                "details": {},
            }
        )

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", resource_runner)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.probe_pptx_artifact(deck)
    assert caught.value.reason_code == "pptx_probe_resource_unavailable"


def test_successful_probe_is_cached_by_exact_file_generation(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "cached.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()
    calls = 0

    def successful_runner(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        return _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "slide_count": 1,
                "source_sha256": "a" * 64,
                "source_size_bytes": deck.stat().st_size,
                "archive_recovery": [],
            }
        )

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", successful_runner)
    assert pptx_evidence.probe_pptx_artifact(deck).slide_count == 1
    assert pptx_evidence.probe_pptx_artifact(deck).source_sha256 == "a" * 64
    assert calls == 1


def test_single_recovery_read_is_transient_and_not_cached(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "materializing.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()
    calls = 0
    recovery = [
        {
            "schema_version": 1,
            "part_name": "ppt/media/image1.png",
            "member_kind": "embedded_media",
            "error_type": "crc_mismatch",
            "status": "recovered_with_placeholder_asset",
            "content_replaced": True,
            "replacement_sha256": (
                "431ced6916a2a21a156e38701afe55bbd7f88969fbbfc56d7fe099d47f265460"
            ),
        }
    ]

    def materializing_runner(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        return _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "slide_count": 1,
                "source_sha256": "a" * 64,
                "source_size_bytes": deck.stat().st_size,
                "archive_recovery": recovery if calls == 1 else [],
            }
        )

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", materializing_runner)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.probe_pptx_artifact(deck)
    assert caught.value.reason_code == "pptx_probe_materialization_changed"
    assert calls == 2

    assert pptx_evidence.probe_pptx_artifact(deck).archive_recovery == ()
    assert calls == 3


def test_recovery_confirmation_shares_one_enclosing_deadline(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "deadline-recovery.pptx"
    _write_deck(deck, with_image=False)
    key = (
        str(deck),
        1,
        2,
        deck.stat().st_size,
        3,
        4,
        0,
        0,
        pptx_evidence.PPTX_EXTRACTION_SCHEMA_VERSION,
        pptx_evidence.PPTX_EXTRACTION_PIPELINE_VERSION,
    )
    deadlines: list[float | None] = []
    calls = 0
    recovery = (
        {
            "schema_version": 1,
            "part_name": "ppt/media/image1.png",
            "member_kind": "embedded_media",
            "error_type": "crc_mismatch",
            "status": "recovered_with_placeholder_asset",
            "content_replaced": True,
            "replacement_sha256": (
                "431ced6916a2a21a156e38701afe55bbd7f88969fbbfc56d7fe099d47f265460"
            ),
        },
    )

    def identity(path, *, trusted_root=None, deadline_monotonic=None):
        del path, trusted_root
        deadlines.append(deadline_monotonic)
        return deck, key

    def probe(path, *, trusted_root=None, deadline_monotonic=None):
        nonlocal calls
        del path, trusted_root
        calls += 1
        deadlines.append(deadline_monotonic)
        if calls == 2:
            raise pptx_evidence.PptxEvidenceError(
                "deadline expired",
                reason_code="pptx_batch_wall_limit",
            )
        return pptx_evidence.PptxArtifactProbe(
            slide_count=1,
            source_sha256="a" * 64,
            source_size_bytes=key[3],
            archive_recovery=recovery,
        )

    pptx_evidence.clear_pptx_artifact_probe_cache()
    monkeypatch.setattr(pptx_evidence, "_probe_file_identity", identity)
    monkeypatch.setattr(pptx_evidence, "_run_bounded_pptx_probe", probe)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.probe_pptx_artifact(deck, deadline_monotonic=123.0)

    assert caught.value.reason_code == "pptx_batch_wall_limit"
    assert calls == 2
    assert deadlines and set(deadlines) == {123.0}


def test_single_structural_crc_failure_is_not_cached_when_confirmation_is_clean(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "materializing-structure.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()
    calls = 0

    def materializing_runner(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        payload = (
            {
                "schema_version": 1,
                "status": "unavailable",
                "reason_code": "pptx_structural_damage",
                "details": {"part_names": ["ppt/slides/slide1.xml"]},
            }
            if calls == 1
            else {
                "schema_version": 1,
                "status": "available",
                "slide_count": 1,
                "source_sha256": "a" * 64,
                "source_size_bytes": deck.stat().st_size,
                "archive_recovery": [],
            }
        )
        return _worker_result(payload)

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", materializing_runner)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.probe_pptx_artifact(deck)
    assert caught.value.reason_code == "pptx_probe_materialization_changed"
    assert calls == 2

    assert pptx_evidence.probe_pptx_artifact(deck).slide_count == 1
    assert calls == 3


def test_single_child_io_failure_is_not_cached_when_confirmation_is_clean(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "hydrating.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()
    calls = 0

    def hydrating_runner(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        payload = (
            {
                "schema_version": 1,
                "status": "unavailable",
                "reason_code": "pptx_artifact_unavailable",
                "details": {"exception_type": "OSError"},
            }
            if calls == 1
            else {
                "schema_version": 1,
                "status": "available",
                "slide_count": 1,
                "source_sha256": "a" * 64,
                "source_size_bytes": deck.stat().st_size,
                "archive_recovery": [],
            }
        )
        return _worker_result(payload)

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", hydrating_runner)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.probe_pptx_artifact(deck)
    assert caught.value.reason_code == "pptx_probe_materialization_changed"
    assert calls == 2

    assert pptx_evidence.probe_pptx_artifact(deck).slide_count == 1
    assert calls == 3


def test_snapshot_open_io_failure_has_portable_unavailable_code(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "open-failure.pptx"
    _write_deck(deck, with_image=False)
    original_open = Path.open

    def failing_open(path: Path, *args, **kwargs):
        if path == deck:
            raise OSError("synthetic hydration failure")
        return original_open(path, *args, **kwargs)

    monkeypatch.setattr(Path, "open", failing_open)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._probe_pptx_artifact_in_process(deck)
    assert caught.value.reason_code == "pptx_artifact_unavailable"
    assert caught.value.details == {"exception_type": "OSError"}


def test_snapshot_generation_change_has_nonsticky_changed_code(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "generation-change.pptx"
    _write_deck(deck, with_image=False)
    actual = deck.lstat()
    original_lstat = Path.lstat
    calls = 0

    def changing_lstat(path: Path):
        nonlocal calls
        if path != deck:
            return original_lstat(path)
        calls += 1
        if calls == 1:
            return actual
        return SimpleNamespace(
            st_mode=actual.st_mode,
            st_dev=actual.st_dev,
            st_ino=actual.st_ino,
            st_size=actual.st_size,
            st_mtime_ns=actual.st_mtime_ns + 1,
            st_ctime_ns=actual.st_ctime_ns,
            st_flags=getattr(actual, "st_flags", 0),
        )

    monkeypatch.setattr(Path, "lstat", changing_lstat)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.snapshot_regular_file(deck, label="PPTX artifact")
    assert caught.value.reason_code == "pptx_artifact_changed"


def test_rendered_pdf_same_size_generation_replacement_fails_closed(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    rendered = tmp_path / "rendered.pdf"
    _write_pdf(rendered, page_count=1)
    actual = rendered.lstat()
    original_lstat = Path.lstat
    calls = 0

    def replaced_lstat(path: Path):
        nonlocal calls
        if path != rendered:
            return original_lstat(path)
        calls += 1
        if calls == 1:
            return actual
        return SimpleNamespace(
            st_mode=actual.st_mode,
            st_dev=actual.st_dev,
            st_ino=actual.st_ino,
            st_size=actual.st_size,
            st_mtime_ns=actual.st_mtime_ns + 1,
            st_ctime_ns=actual.st_ctime_ns,
            st_flags=getattr(actual, "st_flags", 0),
        )

    monkeypatch.setattr(Path, "lstat", replaced_lstat)
    with pytest.raises(pptx_evidence.PptxEvidenceError, match="changed"):
        pptx_evidence.snapshot_rendered_pdf(rendered)


def test_native_audit_recomputation_timeout_is_structured(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "slow-audit.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()

    def timeout_runner(*_args, **_kwargs):
        raise pptx_evidence.SupervisorError("worker_timeout")

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", timeout_runner)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.recompute_native_deck_audit(deck)
    assert caught.value.reason_code == "pptx_probe_timeout"


def test_native_audit_recovery_is_not_sticky_when_confirmation_is_clean(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "materializing-audit.pptx"
    _write_deck(deck, with_image=False)
    pptx_evidence.clear_pptx_artifact_probe_cache()
    audit = pptx_evidence.build_native_deck_audit(
        source_pptx_sha256="a" * 64,
        source_pptx_size_bytes=deck.stat().st_size,
        slide_count=1,
        render_required_reasons={},
    )
    calls = 0

    def materializing_runner(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        payload = (
            {
                "schema_version": 1,
                "status": "unavailable",
                "reason_code": "pptx_archive_recovery_required",
                "details": {"part_names": ["ppt/media/image1.png"]},
            }
            if calls == 1
            else {
                "schema_version": 1,
                "status": "available",
                "native_deck_audit": audit,
            }
        )
        return _worker_result(payload)

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", materializing_runner)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.recompute_native_deck_audit(deck)
    assert caught.value.reason_code == "pptx_probe_materialization_changed"
    assert calls == 2
    assert pptx_evidence.recompute_native_deck_audit(deck) == audit
    assert calls == 3


def test_dataless_cloud_placeholder_never_starts_worker_io(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "cloud-placeholder.pptx"
    _write_deck(deck, with_image=False)
    actual = deck.lstat()
    dataless_flag = pptx_evidence.PPTX_MACOS_DATALESS_FLAG or 0x40000000
    original_lstat = Path.lstat

    def dataless_lstat(path: Path):
        if path == deck:
            return SimpleNamespace(
                st_mode=actual.st_mode,
                st_dev=actual.st_dev,
                st_ino=actual.st_ino,
                st_size=actual.st_size,
                st_mtime_ns=actual.st_mtime_ns,
                st_ctime_ns=actual.st_ctime_ns,
                st_flags=dataless_flag,
            )
        return original_lstat(path)

    if pptx_evidence.PPTX_MACOS_DATALESS_FLAG == 0:
        monkeypatch.setattr(
            pptx_evidence,
            "PPTX_MACOS_DATALESS_FLAG",
            dataless_flag,
        )
    monkeypatch.setattr(Path, "lstat", dataless_lstat)
    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("dataless deck started worker I/O"),
    )
    pptx_evidence.clear_pptx_artifact_probe_cache()

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.probe_pptx_artifact(deck)
    assert caught.value.reason_code == "pptx_cloud_placeholder_unavailable"


def test_dataless_flag_uses_explicit_darwin_fallback(pptx_evidence) -> None:
    expected = int(
        getattr(
            pptx_evidence.stat_module,
            "SF_DATALESS",
            0x40000000 if sys.platform == "darwin" else 0,
        )
    )
    assert pptx_evidence.PPTX_MACOS_DATALESS_FLAG == expected
    if sys.platform == "darwin":
        assert expected == 0x40000000


def test_windows_offline_attribute_never_starts_worker_io(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "windows-placeholder.pptx"
    _write_deck(deck, with_image=False)
    actual = deck.lstat()
    offline = 0x001000
    original_lstat = Path.lstat

    def offline_lstat(path: Path):
        if path == deck:
            return SimpleNamespace(
                st_mode=actual.st_mode,
                st_dev=actual.st_dev,
                st_ino=actual.st_ino,
                st_size=actual.st_size,
                st_mtime_ns=actual.st_mtime_ns,
                st_ctime_ns=actual.st_ctime_ns,
                st_flags=0,
                st_file_attributes=offline,
            )
        return original_lstat(path)

    monkeypatch.setattr(
        pptx_evidence,
        "PPTX_WINDOWS_CLOUD_FILE_ATTRIBUTES",
        offline,
    )
    monkeypatch.setattr(Path, "lstat", offline_lstat)
    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("offline deck started worker I/O"),
    )
    pptx_evidence.clear_pptx_artifact_probe_cache()

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.probe_pptx_artifact(deck)
    assert caught.value.reason_code == "pptx_cloud_placeholder_unavailable"


def test_hydrated_gigabyte_deck_fits_explicit_artifact_cap(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    hydrated_size = 1_100_000_000
    generation = pptx_evidence.FileGeneration(
        size=hydrated_size,
        mtime_ns=1,
        ctime_ns=1,
        device=1,
        inode=1,
        mode=0o100644,
        flags=0,
        file_attributes=0,
    )
    monkeypatch.setattr(
        pptx_evidence,
        "_supervised_file_generation",
        lambda _path, **_kwargs: generation,
    )

    artifact, admitted, root_generation = pptx_evidence._admit_supervised_input(
        tmp_path / "hydrated-large.pptx",
        label="PPTX artifact",
    )

    assert artifact.is_absolute()
    assert admitted == generation
    assert root_generation is None
    assert hydrated_size < pptx_evidence.PPTX_MAX_INPUT_BYTES
    assert pptx_evidence.PPTX_MAX_INPUT_BYTES == 2 * 1024 * 1024 * 1024
    assert all(
        limits.max_memory_bytes >= 4 * 1024 * 1024 * 1024
        for limits in (
            pptx_evidence.PPTX_PROBE_LIMITS,
            pptx_evidence.PPTX_NATIVE_AUDIT_LIMITS,
            pptx_evidence.PPTX_EXTRACT_NO_OCR_LIMITS,
            pptx_evidence.PPTX_EXTRACT_OCR_LIMITS,
        )
    )


def test_artifact_larger_than_explicit_cap_fails_before_worker(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    generation = pptx_evidence.FileGeneration(
        size=pptx_evidence.PPTX_MAX_INPUT_BYTES + 1,
        mtime_ns=1,
        ctime_ns=1,
        device=1,
        inode=1,
        mode=0o100644,
        flags=0,
        file_attributes=0,
    )
    monkeypatch.setattr(
        pptx_evidence,
        "_supervised_file_generation",
        lambda _path, **_kwargs: generation,
    )
    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("oversized artifact started worker"),
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.run_supervised_pptx_extraction(
            tmp_path / "too-large.pptx",
            ocr=False,
        )

    assert caught.value.reason_code == "pptx_probe_resource_unavailable"


def test_structural_crc_damage_is_unavailable_not_placeholder_recovered(
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "damaged-slide.pptx"
    _write_deck(deck, with_image=False)
    damaged_part = _damage_member(deck, lambda name: name == "ppt/slides/slide1.xml")

    with pytest.raises(
        pptx_evidence.PptxEvidenceError,
        match=f"structural PPTX member.*{damaged_part}",
    ):
        pptx_evidence.probe_pptx_artifact(deck)


def test_corrupt_member_scan_enforces_recovery_record_cap(
    pptx_evidence,
    tmp_path: Path,
) -> None:
    package_path = tmp_path / "many-corrupt.zip"
    with zipfile.ZipFile(package_path, "w", compression=zipfile.ZIP_STORED) as archive:
        for index in range(pptx_evidence.PPTX_ARTIFACT_PROBE_MAX_RECOVERY_RECORDS + 1):
            archive.writestr(f"ppt/media/image{index}.bin", b"asset")
    package = bytearray(package_path.read_bytes())
    cursor = 0
    while True:
        header = package.find(b"PK\x01\x02", cursor)
        if header < 0:
            break
        recorded_crc = struct.unpack_from("<I", package, header + 16)[0]
        struct.pack_into("<I", package, header + 16, recorded_crc ^ 0xFFFFFFFF)
        name_size, extra_size, comment_size = struct.unpack_from(
            "<HHH", package, header + 28
        )
        cursor = header + 46 + name_size + extra_size + comment_size

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._corrupt_zip_members(bytes(package))
    assert caught.value.reason_code == "pptx_probe_result_oversized"


@pytest.mark.parametrize(
    ("limit_name", "limit_value"),
    [
        ("PPTX_ARCHIVE_MAX_MEMBERS", 0),
        ("PPTX_ARCHIVE_MAX_EXPANDED_BYTES", 1),
        ("PPTX_ARCHIVE_MAX_MEMBER_BYTES", 1),
    ],
)
def test_archive_metadata_budgets_fail_before_unbounded_expansion(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    limit_name: str,
    limit_value: int,
) -> None:
    deck = tmp_path / "bounded-archive.pptx"
    _write_deck(deck, with_image=False)
    monkeypatch.setattr(pptx_evidence, limit_name, limit_value)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._corrupt_zip_members(deck.read_bytes())

    assert caught.value.reason_code == "pptx_probe_resource_unavailable"


@pytest.mark.parametrize(
    "duplicate_name",
    ["ppt/slides/slide1.xml", "ppt/slides/SLIDE1.XML"],
)
def test_extraction_rejects_conflicting_equivalent_package_part_names(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    duplicate_name: str,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    textbox.text_frame.text = "REAL"
    deck = tmp_path / "duplicate-slide-part.pptx"
    presentation.save(deck)

    with zipfile.ZipFile(deck, "a") as archive:
        original = archive.read("ppt/slides/slide1.xml")
        conflicting = original.replace(b"REAL", b"EVIL")
        if duplicate_name == "ppt/slides/slide1.xml":
            with pytest.warns(UserWarning, match="Duplicate name"):
                archive.writestr(duplicate_name, conflicting)
        else:
            archive.writestr(duplicate_name, conflicting)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_invalid_container"


def test_extraction_rejects_package_part_segment_prefix_collision(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "segment-prefix-part.pptx"
    _write_deck(deck, with_image=True)
    with zipfile.ZipFile(deck, "a") as archive:
        archive.writestr("ppt/media", b"ambiguous part")

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_invalid_container"


@pytest.mark.parametrize(
    "forbidden_name",
    [
        "ppt/media/%69mage1.png",
        "ppt/media/%2Falias.png",
        "ppt/media/%5calias.png",
        "/ppt/media/alias.png",
        "ppt\\media\\alias.png",
        "ppt//media/alias.png",
        "ppt/./alias.png",
        "ppt/media./alias.png",
        "ppt/media/bad%ZZ.png",
        "ppt/media/alias.png?version=2",
    ],
)
def test_extraction_rejects_noncanonical_opc_member_names(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    forbidden_name: str,
) -> None:
    deck = tmp_path / "noncanonical-part.pptx"
    _write_deck(deck, with_image=True)
    with zipfile.ZipFile(deck, "a") as archive:
        archive.writestr(forbidden_name, b"ambiguous part")

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_invalid_container"


def test_extraction_rejects_duplicate_relationship_ids_before_asset_cataloging(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    first_image = tmp_path / "first.png"
    second_image = tmp_path / "second.png"
    Image.new("RGB", (64, 64), "navy").save(first_image)
    Image.new("RGB", (64, 64), "orange").save(second_image)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(first_image), Inches(1), Inches(1))
    slide.shapes.add_picture(str(second_image), Inches(3), Inches(1))
    deck = tmp_path / "duplicate-relationship-id.pptx"
    presentation.save(deck)

    relationships_name = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(deck) as archive:
        members = [(member, archive.read(member)) for member in archive.infolist()]
    duplicate = (
        b'<Relationship Id="rId2" '
        b'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
        b'relationships/image" Target="../media/image2.png"/>'
    )
    closing_tag = b"</Relationships>"
    with zipfile.ZipFile(deck, "w") as archive:
        for member, payload in members:
            if member.filename == relationships_name:
                assert closing_tag in payload
                payload = payload.replace(closing_tag, duplicate + closing_tag)
            archive.writestr(member, payload)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_invalid_container"


@pytest.mark.parametrize(
    "duplicate_entry",
    [
        pytest.param(
            b'<Override PartName="/ppt/slides/slide1.xml" '
            b'ContentType="application/octet-stream"/>',
            id="override-exact",
        ),
        pytest.param(
            b'<Override PartName="/ppt/slides/SLIDE1.XML" '
            b'ContentType="application/octet-stream"/>',
            id="override-case-equivalent",
        ),
        pytest.param(
            b'<Default Extension="PNG" ContentType="application/octet-stream"/>',
            id="default-case-equivalent",
        ),
    ],
)
def test_extraction_rejects_ambiguous_content_type_entries(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    duplicate_entry: bytes,
) -> None:
    deck = tmp_path / "ambiguous-content-types.pptx"
    _write_deck(deck, with_image=True)
    manifest_name = "[Content_Types].xml"
    closing_tag = b"</Types>"
    with zipfile.ZipFile(deck) as archive:
        members = [(member, archive.read(member)) for member in archive.infolist()]
    with zipfile.ZipFile(deck, "w") as archive:
        for member, payload in members:
            if member.filename == manifest_name:
                assert closing_tag in payload
                payload = payload.replace(closing_tag, duplicate_entry + closing_tag)
            archive.writestr(member, payload)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_invalid_container"


def test_extraction_rejects_duplicate_presentation_slide_identities(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "duplicate-slide-identity.pptx"
    _write_deck(deck, with_image=False)
    presentation_name = "ppt/presentation.xml"
    with zipfile.ZipFile(deck) as archive:
        members = [(member, archive.read(member)) for member in archive.infolist()]
    with zipfile.ZipFile(deck, "w") as archive:
        for member, payload in members:
            if member.filename == presentation_name:
                entry_start = payload.index(b"<p:sldId ")
                entry_end = payload.index(b"/>", entry_start) + 2
                duplicate_entry = payload[entry_start:entry_end]
                closing_tag = b"</p:sldIdLst>"
                assert closing_tag in payload
                payload = payload.replace(
                    closing_tag,
                    duplicate_entry + closing_tag,
                )
            archive.writestr(member, payload)

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert caught.value.reason_code == "pptx_invalid_container"


def test_malformed_container_is_structured_unavailable(
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "not-a-deck.pptx"
    deck.write_bytes(b"not an OOXML ZIP package")

    with pytest.raises(
        pptx_evidence.PptxEvidenceError,
        match="invalid PPTX ZIP container",
    ):
        pptx_evidence.probe_pptx_artifact(deck)


def test_extractor_reports_recovered_media_and_keeps_schema_versions_distinct(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "damaged-media.pptx"
    _write_deck(deck, with_image=True)
    damaged_part = _damage_member(deck, lambda name: name.startswith("ppt/media/"))

    result = pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    assert result["schema_version"] == 4
    assert result["pipeline_version"] == "1.4.0"
    assert result["archive_recovery"][0]["part_name"] == damaged_part
    assert result["corrupt_assets"] == [
        {
            "part_name": damaged_part,
            "error_type": "crc_mismatch",
            "status": "recovered_with_placeholder",
        }
    ]
    assert result["native_deck_audit"]["schema_version"] == 1
    assert result["native_deck_audit"]["extraction_schema_version"] == 4
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": result,
            }
        )
        == result
    )


def test_render_receipt_binds_exact_source_render_and_ranges(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "render-required.pptx"
    deck.parent.mkdir(parents=True, exist_ok=True)
    image_path = tmp_path / "full-bleed.png"
    Image.new("RGB", (640, 480), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    presentation.save(deck)
    rendered = tmp_path / "rendered.pdf"
    _write_pdf(rendered, page_count=1)

    result = pptx_extraction._extract_pptx_in_process(
        deck,
        ocr=False,
        rendered_pdf_path=rendered,
        inspected_page_ranges=[[1, 1]],
    )
    audit = pptx_evidence.validate_native_deck_audit(
        result["native_deck_audit"], slide_count=1
    )
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": result,
            }
        )
        == result
    )

    assert audit["render_required_slide_numbers"] == [1]
    receipt = audit["rendered_page_inspection"]
    assert receipt["inspected_page_ranges"] == [[1, 1]]
    assert receipt["complete"] is True
    tampered = json.loads(json.dumps(audit))
    tampered["rendered_page_inspection"]["inspected_page_ranges"] = []
    with pytest.raises(
        pptx_evidence.PptxEvidenceError,
        match="inspected_required_slide_numbers",
    ):
        pptx_evidence.validate_native_deck_audit(tampered, slide_count=1)


def test_full_extraction_render_binding_rejects_unrequested_or_broader_receipt(
    pptx_extraction,
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "render-binding.pptx"
    image_path = tmp_path / "full-bleed-binding.png"
    Image.new("RGB", (640, 480), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    presentation.save(deck)
    rendered = tmp_path / "render-binding.pdf"
    _write_pdf(rendered, page_count=1)
    extraction = pptx_extraction._extract_pptx_in_process(
        deck,
        ocr=False,
        rendered_pdf_path=rendered,
        inspected_page_ranges=[[1, 1]],
    )
    rendered_generation = pptx_evidence.FileGeneration.from_stat(rendered.stat())

    with pytest.raises(pptx_evidence.PptxEvidenceError) as broader:
        pptx_evidence._validate_extraction_render_binding(
            extraction,
            rendered_generation=rendered_generation,
            requested_ranges=[],
        )
    assert broader.value.reason_code == "pptx_probe_malformed_result"

    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        ),
    )
    with pytest.raises(pptx_evidence.PptxEvidenceError) as unrequested:
        pptx_evidence.run_supervised_pptx_extraction(deck, ocr=False)
    assert unrequested.value.reason_code == "pptx_probe_malformed_result"


def test_full_extraction_request_binds_options_generations_and_limit_profile(
    pptx_extraction,
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "supervised.pptx"
    _write_deck(deck, with_image=False)
    rendered = tmp_path / "supervised.pdf"
    _write_pdf(rendered, page_count=1)
    in_process = pptx_extraction._extract_pptx_in_process(
        deck,
        ocr=False,
        rendered_pdf_path=rendered,
        inspected_page_ranges=[[1, 1]],
    )
    captured = {}

    def worker(command, operation, generations, payload, limits, **kwargs):
        captured.update(
            {
                "command": command,
                "operation": operation,
                "generations": generations,
                "payload": payload,
                "profile_id": limits.profile_id,
                "kwargs": kwargs,
            }
        )
        return _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": in_process,
            }
        )

    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", worker)
    result = pptx_evidence.run_supervised_pptx_extraction(
        deck,
        ocr=False,
        rendered_pdf_path=rendered,
        inspected_page_ranges=[[1, 1]],
    )

    assert result["slide_count"] == 1
    assert captured["operation"] == pptx_evidence.PPTX_EXTRACT_OPERATION
    assert captured["profile_id"] == "pptx-extract-no-ocr-v1"
    assert set(captured["generations"]) == {"pptx", "rendered_pdf"}
    assert captured["payload"] == {
        "pptx_path": str(deck),
        "trusted_root": None,
        "ocr": False,
        "rendered_pdf_path": str(rendered),
        "inspected_page_ranges": [[1, 1]],
        "extraction_schema_version": 4,
        "extraction_pipeline_version": "1.4.0",
    }
    assert captured["kwargs"]["schema_generation"] == 4
    assert captured["kwargs"]["pipeline_generation"] == "1.4.0"
    command_text = "\n".join(str(item) for item in captured["command"])
    assert str(deck) not in command_text
    assert str(rendered) not in command_text


def test_full_extraction_enforces_remaining_source_bytes_before_worker(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "grown-after-discovery.pptx"
    _write_deck(deck, with_image=False)
    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail(
            "generation outside remaining aggregate bytes started worker"
        ),
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.run_supervised_pptx_extraction(
            deck,
            ocr=False,
            source_size_limit_bytes=deck.stat().st_size - 1,
        )

    assert caught.value.reason_code == "pptx_batch_input_limit"


def test_full_extraction_clamps_worker_to_enclosing_batch_deadline(
    pptx_extraction,
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "deadline-bound.pptx"
    _write_deck(deck, with_image=False)
    in_process = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    captured = {}

    def worker(_command, _operation, _generations, _payload, limits, **_kwargs):
        captured["wall_seconds"] = limits.wall_seconds
        return _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": in_process,
            }
        )

    monkeypatch.setattr(pptx_evidence.time, "monotonic", lambda: 100.0)
    monkeypatch.setattr(pptx_evidence, "run_authenticated_worker", worker)

    result = pptx_evidence.run_supervised_pptx_extraction(
        deck,
        ocr=False,
        deadline_monotonic=112.0,
    )

    assert result["slide_count"] == 1
    # Two seconds remain reserved for process-tree cleanup.
    assert captured["wall_seconds"] == 10.0


def test_metadata_admission_timeout_reports_enclosing_batch_wall_limit(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "metadata-deadline.pptx"
    _write_deck(deck, with_image=False)
    captured: dict[str, object] = {}

    def blocked_metadata(_command, _payload, _sensitive, limits):
        captured["wall_seconds"] = limits.wall_seconds
        raise pptx_evidence.SupervisorError("worker_timeout")

    monkeypatch.setattr(pptx_evidence.time, "monotonic", lambda: 100.0)
    monkeypatch.setattr(
        pptx_evidence,
        "_invoke_metadata_worker",
        blocked_metadata,
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.run_supervised_pptx_extraction(
            deck,
            ocr=False,
            deadline_monotonic=112.0,
        )

    assert caught.value.reason_code == "pptx_batch_wall_limit"
    assert captured["wall_seconds"] == 10.0


def test_post_metadata_check_cannot_extend_enclosing_batch_deadline(
    pptx_extraction,
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "post-metadata-deadline.pptx"
    _write_deck(deck, with_image=False)
    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    generation = pptx_evidence.FileGeneration.from_stat(deck.stat())
    metadata_calls = 0

    def metadata(_command, _payload, _sensitive, _limits):
        nonlocal metadata_calls
        metadata_calls += 1
        return _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "generation": generation.to_dict(),
                "root_generation": None,
                "reparse_tag": None,
            }
        )

    times = iter((100.0, 100.0, 111.0))
    monkeypatch.setattr(pptx_evidence.time, "monotonic", lambda: next(times))
    monkeypatch.setattr(pptx_evidence, "_invoke_metadata_worker", metadata)
    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        ),
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.run_supervised_pptx_extraction(
            deck,
            ocr=False,
            deadline_monotonic=112.0,
        )

    assert caught.value.reason_code == "pptx_batch_wall_limit"
    assert metadata_calls == 1


def test_exhausted_batch_deadline_fails_before_worker_start(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "deadline-exhausted.pptx"
    _write_deck(deck, with_image=False)
    monkeypatch.setattr(pptx_evidence.time, "monotonic", lambda: 100.0)
    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("expired deadline started worker"),
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.run_supervised_pptx_extraction(
            deck,
            ocr=False,
            deadline_monotonic=101.0,
        )

    assert caught.value.reason_code == "pptx_batch_wall_limit"


def test_deadline_clamped_worker_timeout_reports_batch_wall_limit(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "deadline-timeout.pptx"
    _write_deck(deck, with_image=False)
    monkeypatch.setattr(pptx_evidence.time, "monotonic", lambda: 100.0)

    def timeout_worker(*_args, **_kwargs):
        raise pptx_evidence.SupervisorError("worker_timeout")

    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        timeout_worker,
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.run_supervised_pptx_extraction(
            deck,
            ocr=False,
            deadline_monotonic=112.0,
        )

    assert caught.value.reason_code == "pptx_batch_wall_limit"
    assert caught.value.details["admitted_source_size_bytes"] == deck.stat().st_size


@pytest.mark.parametrize(
    ("field_path", "invalid_value"),
    [
        (("slide_width_inches",), None),
        (("slide_width_inches",), 10**1000),
        (("aspect_ratio",), "arbitrary"),
        (("aspect_ratio",), "1:1"),
        (("aspect_ratio",), f"{'1' * 5000}:1"),
        (("corrupt_assets",), [{"unexpected": True}]),
        (("template_layouts",), [{"index": 0}]),
        (("global_design",), {"fonts_used": {}}),
        (("native_timing_summary",), {}),
        (("per_slide_visual", 0), {"slide_number": 1}),
        (("per_slide_visual", 0, "native_timing"), {}),
        (("per_slide_visual", 0, "background_type"), []),
        (("per_slide_visual", 0, "text_extraction_confidence"), []),
        (("per_slide_visual", 0, "text_extraction_method"), []),
    ],
)
def test_full_extraction_rejects_malformed_nested_worker_payloads(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    field_path: tuple[object, ...],
    invalid_value: object,
) -> None:
    deck = tmp_path / "nested-contract.pptx"
    _write_deck(deck, with_image=False)
    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    payload = {
        "schema_version": 1,
        "status": "available",
        "extraction": json.loads(json.dumps(extraction)),
    }
    target = payload["extraction"]
    for component in field_path[:-1]:
        target = target[component]
    target[field_path[-1]] = invalid_value

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(payload)

    assert caught.value.reason_code == "pptx_probe_malformed_result"


@pytest.mark.parametrize("with_image", [False, True])
def test_full_extraction_accepts_complete_nested_worker_payload(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    with_image: bool,
) -> None:
    deck = tmp_path / "nested-contract-valid.pptx"
    _write_deck(deck, with_image=with_image)
    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    decoded = pptx_evidence._decode_extraction_worker_payload(
        {
            "schema_version": 1,
            "status": "available",
            "extraction": extraction,
        }
    )

    assert decoded == extraction


def test_full_extraction_accepts_derived_render_and_ocr_contracts(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    background_image = tmp_path / "background.png"
    Image.new("RGB", (64, 64), "navy").save(background_image)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.shapes.add_textbox(
        Inches(1), Inches(1), Inches(2), Inches(1)
    ).text_frame.text = "group text"
    slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(2), Inches(1)).table.cell(
        0, 0
    ).text = "table text"
    smartart = slide.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(2), Inches(1))
    smartart.element.graphic.graphicData.uri = (
        "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    )
    _set_background_image(slide, background_image)
    deck = tmp_path / "derived-contracts.pptx"
    presentation.save(str(deck))

    extraction = pptx_extraction._extract_pptx_in_process(
        deck,
        ocr=True,
        ocr_fn=lambda _blobs: {
            "attempted": True,
            "engine": "tesseract",
            "engine_version": "5.3.4",
            "result_status": "text_recovered",
            "result_confidence": 93.0,
            "error": None,
            "recovered_text": "background label",
            "trustworthy_text": True,
        },
    )
    decoded = pptx_evidence._decode_extraction_worker_payload(
        {
            "schema_version": 1,
            "status": "available",
            "extraction": extraction,
        }
    )

    assert decoded == extraction
    assert decoded["per_slide_visual"][0]["render_required_reasons"] == [
        "background_image",
        "grouped_shapes",
        "smartart",
        "table",
    ]
    assert decoded["per_slide_visual"][0]["ocr_text"] == "background label"


@pytest.mark.parametrize(
    "binding",
    [
        "shape_name",
        "layout_catalog",
        "timing_part",
        "slide_part_ordinal",
        "unsupported_uri",
        "unsupported_reason",
        "table_cells",
        "background_provenance",
    ],
)
def test_catalog_bindings_reject_detached_worker_evidence(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    binding: str,
) -> None:
    background_image = tmp_path / f"catalog-binding-{binding}.png"
    Image.new("RGB", (64, 64), "navy").save(background_image)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(2), Inches(0.5))
    textbox.text_frame.text = "shape text"
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1.5), Inches(2), Inches(1))
    table.table.cell(0, 0).text = "table text"
    smartart = slide.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(2), Inches(1))
    smartart.element.graphic.graphicData.uri = (
        "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    )
    _set_background_image(slide, background_image)
    deck = tmp_path / f"catalog-binding-{binding}.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]

    if binding == "shape_name":
        shape = next(
            item for item in slide_result["shapes_summary"] if item["has_text_frame"]
        )
        shape["name"] = "detached name"
    elif binding == "layout_catalog":
        slide_result["layout_name"] = "not present in template_layouts"
    elif binding == "timing_part":
        slide_result["native_timing"]["provenance"]["part_name"] = (
            "ppt/slides/slide999.xml"
        )
    elif binding == "slide_part_ordinal":
        slide_result["slide_part_name"] = "ppt/slides/slide999.xml"
        slide_result["native_timing"]["provenance"]["part_name"] = (
            "ppt/slides/slide999.xml"
        )
    elif binding == "unsupported_uri":
        unsupported = next(
            item
            for item in slide_result["unsupported_content"]
            if item["content_type"] == "smartart"
        )
        unsupported["graphic_data_uri"] = "urn:detached:diagram"
        channel = next(
            item
            for item in slide_result["text_channels"]
            if item["channel"] == "smartart_text"
        )
        channel["provenance"]["graphic_data_uri"] = "urn:detached:diagram"
    elif binding == "unsupported_reason":
        slide_result["unsupported_content"][0]["reason"] = "unverified"
    elif binding == "table_cells":
        channel = next(
            item
            for item in slide_result["text_channels"]
            if item["channel"] == "table_cell_text"
        )
        channel["provenance"]["cells"] = []
    else:
        channel = next(
            item
            for item in slide_result["text_channels"]
            if item["channel"] == "background_image_ocr"
        )
        channel["provenance"].pop("relationship_id")

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            },
            expected_ocr=False,
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_template_master_indices_allow_empty_master_gaps_but_remain_grouped(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "template-master-gaps.pptx"
    _write_deck(deck, with_image=False)
    layouts = json.loads(
        json.dumps(
            pptx_extraction._extract_pptx_in_process(deck, ocr=False)[
                "template_layouts"
            ]
        )
    )
    assert len(layouts) >= 2

    for layout in layouts:
        layout["master_index"] = 2
    assert pptx_evidence._valid_template_layouts(layouts) is True

    layouts[0]["master_index"] = 3
    assert pptx_evidence._valid_template_layouts(layouts) is False


def test_near_threshold_picture_uses_reported_ratio_for_render_decision(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "threshold.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=int(presentation.slide_width * 4996 / 10_000),
        height=presentation.slide_height,
    )
    deck = tmp_path / "threshold.pptx"
    presentation.save(deck)

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    slide_result = extraction["per_slide_visual"][0]

    assert (
        pptx_extraction.PPTX_TEXT_BEARING_IMAGE_AREA_RATIO
        == pptx_evidence.PPTX_TEXT_BEARING_IMAGE_AREA_RATIO
    )
    assert slide_result["image_area_ratio"] == 0.5
    assert "large_picture" in slide_result["render_required_reasons"]
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )
        == extraction
    )


def test_extraction_payload_is_bound_to_requested_ocr_mode(
    pptx_extraction,
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "ocr-mode.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    deck = tmp_path / "ocr-mode.pptx"
    presentation.save(deck)

    disabled = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as disabled_mismatch:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": disabled,
            },
            expected_ocr=True,
        )
    assert disabled_mismatch.value.reason_code == "pptx_probe_malformed_result"

    monkeypatch.setattr(
        pptx_extraction,
        "_ocr_image_result",
        lambda _blob: {
            "attempted": True,
            "engine": "tesseract",
            "engine_version": "5.3.4",
            "result_status": "text_recovered",
            "result_confidence": 99.0,
            "error": None,
            "recovered_text": "BOUND OCR",
            "trustworthy_text": True,
        },
    )
    enabled = pptx_extraction._extract_pptx_in_process(deck, ocr=True)
    with pytest.raises(pptx_evidence.PptxEvidenceError) as enabled_mismatch:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": enabled,
            },
            expected_ocr=False,
        )
    assert enabled_mismatch.value.reason_code == "pptx_probe_malformed_result"


@pytest.mark.parametrize(
    ("top_inches", "reported_top", "expected_footer"),
    [
        (6.374, 6.37, ""),
        (6.375, 6.38, "boundary footer"),
    ],
)
def test_footer_threshold_uses_reported_geometry_for_round_trip_validation(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    top_inches: float,
    reported_top: float,
    expected_footer: str,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(top_inches),
        Inches(2),
        Inches(0.5),
    )
    textbox.text_frame.text = "boundary footer"
    deck = tmp_path / f"footer-{top_inches}.pptx"
    presentation.save(deck)

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    slide_result = extraction["per_slide_visual"][0]

    assert slide_result["shapes_summary"][0]["top"] == reported_top
    assert slide_result["footer_text"] == expected_footer
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )
        == extraction
    )


def test_full_slide_picture_cannot_hide_its_geometry_and_render_requirement(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "full-slide.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    deck = tmp_path / "full-slide.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]
    slide_result["image_area_ratio"] = 0.0
    slide_result["render_required"] = False
    slide_result["render_required_reasons"] = []
    slide_result["text_extraction_confidence"] = "high"
    extraction["native_deck_audit"]["render_required_slide_numbers"] = []
    extraction["native_deck_audit"]["render_required_reasons"] = {}

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_full_extraction_rejects_unsupported_content_without_render_requirement(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "hidden-unsupported.pptx"
    _write_deck(deck, with_image=False)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide = extraction["per_slide_visual"][0]
    slide["unsupported_content"] = [
        {
            "content_type": "chart",
            "shape_name": None,
            "shape_path": [],
            "reason": "visible labels require rendering",
            "render_required": True,
        }
    ]
    slide["has_unsupported_content"] = True

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_smartart_shape_cannot_omit_unsupported_evidence(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    smartart = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(2), Inches(1))
    smartart.element.graphic.graphicData.uri = (
        "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    )
    deck = tmp_path / "hidden-smartart.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]
    slide_result["unsupported_content"] = []
    slide_result["has_unsupported_content"] = False
    slide_result["text_channels"] = [
        channel
        for channel in slide_result["text_channels"]
        if channel["channel"] != "smartart_text"
    ]
    slide_result["render_required"] = False
    slide_result["render_required_reasons"] = []
    slide_result["text_extraction_confidence"] = "high"
    extraction["native_deck_audit"]["render_required_slide_numbers"] = []
    extraction["native_deck_audit"]["render_required_reasons"] = {}

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_graphic_frame_without_uri_round_trips_as_unsupported(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    graphic_frame = slide.shapes.add_table(
        1,
        1,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    graphic_frame.element.graphic.graphicData.uri = ""
    deck = tmp_path / "missing-graphic-uri.pptx"
    presentation.save(deck)

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    slide_result = extraction["per_slide_visual"][0]
    shape = slide_result["shapes_summary"][0]

    assert shape["is_graphic_frame"] is True
    assert shape["graphic_frame_type"] == "graphic_frame"
    assert shape["graphic_data_uri"] is None
    assert slide_result["render_required_reasons"] == ["graphic_frame"]
    assert slide_result["unsupported_content"][0]["content_type"] == "graphic_frame"
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            },
            expected_ocr=False,
        )
        == extraction
    )


def test_full_extraction_rejects_receiptless_invented_ocr_text(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "invented-ocr.pptx"
    _write_deck(deck, with_image=False)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide = extraction["per_slide_visual"][0]
    slide["text_channels"].append(
        {
            "channel": "picture_ocr",
            "text": "INVENTED",
            "confidence": "low",
            "result_confidence": 99.0,
            "status": "extracted",
            "attempted": True,
            "engine": "tesseract",
            "engine_version": "5",
            "reason": None,
            "provenance": {
                "source": "embedded_picture_blobs",
                "shape_paths": [],
            },
            "ocr_receipts": [],
        }
    )
    slide["has_extracted_text"] = True

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


@pytest.mark.parametrize("mutation", ["preview", "footer", "shape_channel"])
def test_native_text_cannot_be_invented_without_bound_shape_evidence(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    mutation: str,
) -> None:
    deck = tmp_path / f"invented-native-{mutation}.pptx"
    _write_deck(deck, with_image=False)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide = extraction["per_slide_visual"][0]
    if mutation == "preview":
        slide["text_content_preview"] = "INVENTED"
    elif mutation == "footer":
        slide["footer_text"] = "INVENTED"
    else:
        slide["text_channels"].append(
            {
                "channel": "shape_text",
                "text": "INVENTED",
                "confidence": "high",
                "status": "extracted",
                "provenance": {
                    "source": "pptx_shape_text_frame",
                    "shape_path": [],
                },
            }
        )
        slide["has_extracted_text"] = True
        slide["text_content_preview"] = "INVENTED"

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_native_shape_text_channel_cannot_be_omitted(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(2), Inches(1)
    ).text_frame.text = "bound text"
    deck = tmp_path / "omitted-shape-text.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]
    slide_result["text_channels"] = []
    slide_result["has_extracted_text"] = False
    slide_result["text_content_preview"] = ""

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_textbox_capability_cannot_be_erased_with_its_channel(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(2), Inches(1)
    ).text_frame.text = "SECRET"
    deck = tmp_path / "erased-textbox.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]
    shape = slide_result["shapes_summary"][0]
    shape["has_text_frame"] = False
    shape.pop("text_preview")
    slide_result["text_channels"] = []
    slide_result["has_text_frame_shapes"] = False
    slide_result["has_extracted_text"] = False
    slide_result["text_content_preview"] = ""

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_picture_cannot_forge_native_text_frame_capability(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "picture-native-text.pptx"
    _write_deck(deck, with_image=True)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]
    shape = slide_result["shapes_summary"][0]
    shape["has_text_frame"] = True
    shape["text_preview"] = "INVENTED"
    slide_result["text_channels"].append(
        {
            "channel": "shape_text",
            "text": "INVENTED",
            "confidence": "high",
            "status": "extracted",
            "provenance": {
                "source": "pptx_shape_text_frame",
                "shape_path": shape["shape_path"],
            },
        }
    )
    slide_result["has_text_frame_shapes"] = True
    slide_result["has_extracted_text"] = True
    slide_result["text_content_preview"] = "INVENTED"

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_inserted_picture_placeholder_round_trips_as_picture_evidence(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "placeholder-picture.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    picture_placeholder = next(
        placeholder
        for placeholder in slide.placeholders
        if hasattr(placeholder, "insert_picture")
    )
    picture_placeholder.insert_picture(str(image_path))
    deck = tmp_path / "placeholder-picture.pptx"
    presentation.save(deck)

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    slide_result = extraction["per_slide_visual"][0]
    picture = next(
        shape for shape in slide_result["shapes_summary"] if shape["is_picture"]
    )

    assert picture["shape_type"] == "PLACEHOLDER (14)"
    assert picture["has_text_frame"] is False
    assert picture["picture_asset_status"] == "available"
    assert slide_result["has_image"] is True
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )
        == extraction
    )


def test_linked_picture_round_trips_as_unavailable_picture_evidence(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "linked-picture.pptx"
    _write_deck(deck, with_image=True)
    slide_name = "ppt/slides/slide1.xml"
    relationships_name = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(deck) as archive:
        members = [(member, archive.read(member)) for member in archive.infolist()]
    with zipfile.ZipFile(deck, "w") as archive:
        for member, payload in members:
            if member.filename == slide_name:
                assert b'r:embed="rId2"' in payload
                payload = payload.replace(
                    b'r:embed="rId2"',
                    b'r:link="rId2"',
                )
            elif member.filename == relationships_name:
                assert b'Target="../media/image1.png"' in payload
                payload = payload.replace(
                    b'Target="../media/image1.png"',
                    b'Target="https://example.invalid/image.png" TargetMode="External"',
                )
            archive.writestr(member, payload)

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    slide_result = extraction["per_slide_visual"][0]
    picture = slide_result["shapes_summary"][0]

    assert picture["shape_type"] == "PICTURE (13)"
    assert picture["is_picture"] is True
    assert picture["picture_asset_status"] == "unavailable"
    assert picture["picture_part_name"] is None
    assert slide_result["has_image"] is True
    assert slide_result["render_required_reasons"] == ["unreadable_picture"]
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            },
            expected_ocr=False,
        )
        == extraction
    )


@pytest.mark.parametrize("mutation", ["textbox_to_picture", "table_to_chart"])
def test_known_shape_capabilities_cannot_be_coherently_relabelled(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    mutation: str,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    textbox.text_frame.text = "bound text"
    slide.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(2), Inches(1))
    deck = tmp_path / f"shape-relabel-{mutation}.pptx"
    presentation.save(deck)
    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    shapes = extraction["per_slide_visual"][0]["shapes_summary"]

    if mutation == "textbox_to_picture":
        shape = next(
            item for item in shapes if item["shape_type"].startswith("TEXT_BOX")
        )
        shape["has_text_frame"] = False
        shape["is_picture"] = True
        shape.pop("text_preview")
        for field in ("font_name", "font_size", "font_color", "bold", "italic"):
            shape.pop(field, None)
        shape.update(
            picture_asset_status="unavailable",
            picture_part_name=None,
            picture_asset_sha256=None,
        )
    else:
        shape = next(item for item in shapes if item["shape_type"].startswith("TABLE"))
        shape["graphic_frame_type"] = "chart"
        shape["graphic_data_uri"] = pptx_evidence._GRAPHIC_DATA_URI_CHART
        for field in (
            "table_rows",
            "table_columns",
            "table_text_preview",
            "table_fonts",
        ):
            shape.pop(field)

    assert pptx_evidence._valid_shape_summary(shape) is False


@pytest.mark.parametrize(
    "mutation", ["partial_text_format", "picture_font", "picture_auto_shape"]
)
def test_shape_catalog_fields_follow_their_native_capabilities(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    mutation: str,
) -> None:
    deck = tmp_path / f"shape-catalog-{mutation}.pptx"
    image_path = tmp_path / f"shape-catalog-{mutation}.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    textbox.text_frame.text = "formatted"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(3))
    presentation.save(deck)
    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    shapes = extraction["per_slide_visual"][0]["shapes_summary"]

    if mutation == "partial_text_format":
        shape = next(item for item in shapes if item["has_text_frame"])
        assert {"font_name", "font_size", "font_color", "bold", "italic"} <= set(shape)
        shape.pop("italic")
    else:
        shape = next(item for item in shapes if item["is_picture"])
        if mutation == "picture_font":
            shape.update(
                font_name=None,
                font_size=None,
                font_color=None,
                bold=None,
                italic=None,
            )
        else:
            shape["auto_shape_type"] = "RECTANGLE (1)"

    assert pptx_evidence._valid_shape_summary(shape) is False


def test_table_capability_cannot_be_erased_with_channel_and_render_flags(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        1, 1, Inches(1), Inches(1), Inches(3), Inches(1)
    )
    table_shape.table.cell(0, 0).text = "SECRET TABLE"
    deck = tmp_path / "erased-table.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]
    shape = slide_result["shapes_summary"][0]
    shape["is_graphic_frame"] = False
    shape["graphic_frame_type"] = None
    shape["graphic_data_uri"] = None
    for field in ("table_rows", "table_columns", "table_text_preview", "table_fonts"):
        shape.pop(field)
    slide_result["text_channels"] = []
    slide_result["has_extracted_text"] = False
    slide_result["text_content_preview"] = ""
    slide_result["render_required"] = False
    slide_result["render_required_reasons"] = []
    slide_result["text_extraction_confidence"] = "high"
    fingerprint = extraction["input_fingerprint"]
    extraction["native_deck_audit"] = pptx_evidence.build_native_deck_audit(
        source_pptx_sha256=fingerprint["digest"],
        source_pptx_size_bytes=fingerprint["size_bytes"],
        slide_count=1,
        render_required_reasons={},
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_media_shape_cannot_forge_a_native_text_frame(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    movie_path = tmp_path / "movie.mp4"
    movie_path.write_bytes(b"not-a-real-movie")
    poster_path = tmp_path / "movie-poster.png"
    Image.new("RGB", (64, 64), "navy").save(poster_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_movie(
        str(movie_path),
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        str(poster_path),
        mime_type="video/mp4",
    )
    deck = tmp_path / "media-native-text.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    slide_result = extraction["per_slide_visual"][0]
    shape = next(
        item
        for item in slide_result["shapes_summary"]
        if item["shape_type"].startswith("MEDIA")
    )
    shape["has_text_frame"] = True
    shape["text_preview"] = "INVENTED"
    slide_result["text_channels"].append(
        {
            "channel": "shape_text",
            "text": "INVENTED",
            "confidence": "high",
            "status": "extracted",
            "provenance": {
                "source": "pptx_shape_text_frame",
                "shape_path": shape["shape_path"],
            },
        }
    )
    slide_result["has_text_frame_shapes"] = True
    slide_result["has_extracted_text"] = True
    slide_result["text_content_preview"] = "INVENTED"

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_duplicate_and_empty_shape_names_round_trip_with_canonical_paths(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    second = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    unnamed = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    first.name = "Same"
    second.name = "Same"
    unnamed.name = ""
    first.text_frame.text = "first"
    second.text_frame.text = "second"
    unnamed.text_frame.text = "third"
    deck = tmp_path / "duplicate-empty-names.pptx"
    presentation.save(deck)

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    paths = [
        shape["shape_path"]
        for shape in extraction["per_slide_visual"][0]["shapes_summary"]
    ]

    assert paths[:2] == [["Same"], ["Same"]]
    assert paths[2] == ["shape_3"]
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )
        == extraction
    )


def test_duplicate_and_empty_unsupported_names_preserve_multiplicity(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shapes = [
        slide.shapes.add_table(
            1,
            1,
            Inches(1),
            Inches(1 + index),
            Inches(2),
            Inches(0.5),
        )
        for index in range(3)
    ]
    for shape, name in zip(shapes, ("Same", "Same", "")):
        shape.name = name
        shape.element.graphic.graphicData.uri = (
            "http://schemas.openxmlformats.org/drawingml/2006/diagram"
        )
    deck = tmp_path / "duplicate-empty-unsupported.pptx"
    presentation.save(deck)

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    unsupported_names = [
        item["shape_name"]
        for item in extraction["per_slide_visual"][0]["unsupported_content"]
    ]

    assert unsupported_names == ["Same", "Same", "shape_3"]
    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )
        == extraction
    )


def test_picture_ocr_cannot_forge_or_duplicate_shape_paths(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "duplicate-picture-receipt.pptx"
    image_path = tmp_path / "duplicate-picture-receipt.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    channel = next(
        item
        for item in extraction["per_slide_visual"][0]["text_channels"]
        if item["channel"] == "picture_ocr"
    )
    channel["ocr_receipts"].append(dict(channel["ocr_receipts"][0]))
    channel["provenance"]["shape_paths"].append(
        list(channel["ocr_receipts"][0]["shape_path"])
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


@pytest.mark.parametrize("mutation", ["drop", "forge_sha"])
def test_picture_ocr_receipts_are_exactly_bound_to_readable_assets(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    mutation: str,
) -> None:
    deck = tmp_path / f"picture-binding-{mutation}.pptx"
    image_path = tmp_path / f"picture-binding-{mutation}.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    channel = next(
        item
        for item in extraction["per_slide_visual"][0]["text_channels"]
        if item["channel"] == "picture_ocr"
    )
    if mutation == "drop":
        channel["ocr_receipts"] = []
        channel["provenance"]["shape_paths"] = []
        channel["status"] = "unavailable"
        channel["reason"] = "no_readable_asset"
    else:
        channel["ocr_receipts"][0]["asset_sha256"] = "f" * 64

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_picture_asset_status_is_cross_bound_to_archive_recovery(
    pptx_evidence,
) -> None:
    part_name = "ppt/media/image1.png"
    replacement_sha = "a" * 64
    recovery = (
        {
            "part_name": part_name,
            "replacement_sha256": replacement_sha,
        },
    )

    def slides(status: str) -> list[dict[str, object]]:
        return [
            {
                "background_type": "solid",
                "shapes_summary": [
                    {
                        "shape_type": "PICTURE (13)",
                        "is_picture": True,
                        "picture_asset_status": status,
                        "picture_part_name": part_name,
                        "picture_asset_sha256": replacement_sha,
                    }
                ],
            }
        ]

    assert pptx_evidence._valid_recovery_asset_bindings(slides("corrupt"), recovery)
    assert not pptx_evidence._valid_recovery_asset_bindings(slides("corrupt"), ())
    assert not pptx_evidence._valid_recovery_asset_bindings(
        slides("available"), recovery
    )
    case_variant_recovery = (
        {
            "part_name": "ppt/media/Image1.png",
            "replacement_sha256": replacement_sha,
        },
    )
    assert not pptx_evidence._valid_recovery_asset_bindings(
        slides("available"),
        case_variant_recovery,
    )


def test_one_picture_part_cannot_claim_conflicting_asset_digests(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "deduplicated-picture.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    slide.shapes.add_picture(str(image_path), Inches(3), Inches(1))
    deck = tmp_path / "conflicting-picture-digests.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(pptx_extraction._extract_pptx_in_process(deck, ocr=False))
    )
    pictures = extraction["per_slide_visual"][0]["shapes_summary"]

    assert pictures[0]["picture_part_name"] == pictures[1]["picture_part_name"]
    assert pictures[0]["picture_asset_sha256"] == pictures[1]["picture_asset_sha256"]
    pictures[1]["picture_asset_sha256"] = "f" * 64

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            },
            expected_ocr=False,
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


@pytest.mark.parametrize("asset_kind", ["picture", "background"])
def test_empty_image_parts_round_trip_as_unavailable_assets(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    asset_kind: str,
) -> None:
    image_path = tmp_path / f"empty-{asset_kind}.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    if asset_kind == "picture":
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    else:
        _set_background_image(slide, image_path)
    deck = tmp_path / f"empty-{asset_kind}.pptx"
    presentation.save(deck)

    with zipfile.ZipFile(deck) as archive:
        members = [(member, archive.read(member)) for member in archive.infolist()]
    with zipfile.ZipFile(deck, "w") as archive:
        for member, payload in members:
            archive.writestr(
                member,
                b"" if member.filename == "ppt/media/image1.png" else payload,
            )

    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    slide_result = extraction["per_slide_visual"][0]
    if asset_kind == "picture":
        picture = slide_result["shapes_summary"][0]
        assert picture["picture_asset_status"] == "unavailable"
        assert picture["picture_part_name"] is None
        assert picture["picture_asset_sha256"] is None
        assert slide_result["render_required_reasons"] == ["unreadable_picture"]
    else:
        assert slide_result["background_asset_status"] == "unavailable"
        assert slide_result["background_part_name"] is None
        assert slide_result["background_asset_sha256"] is None
        assert slide_result["render_required_reasons"] == ["background_image"]

    assert (
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            },
            expected_ocr=False,
        )
        == extraction
    )


def test_background_ocr_cannot_duplicate_its_single_bound_asset(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "background-duplicate.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background_image(slide, image_path)
    deck = tmp_path / "background-duplicate.pptx"
    presentation.save(deck)
    extraction = json.loads(
        json.dumps(
            pptx_extraction._extract_pptx_in_process(
                deck,
                ocr=True,
                ocr_fn=lambda _blobs: {
                    "attempted": True,
                    "engine": "tesseract",
                    "engine_version": "5.3.4",
                    "result_status": "text_recovered",
                    "result_confidence": 93.0,
                    "error": None,
                    "recovered_text": "background label",
                    "trustworthy_text": True,
                },
            )
        )
    )
    slide_result = extraction["per_slide_visual"][0]
    channel = next(
        item
        for item in slide_result["text_channels"]
        if item["channel"] == "background_image_ocr"
    )
    channel["ocr_receipts"].append(dict(channel["ocr_receipts"][0]))
    channel["text"] = "background label | background label"
    slide_result["ocr_text"] = channel["text"]

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )

    assert caught.value.reason_code == "pptx_probe_malformed_result"


def test_unhashable_channel_and_ocr_enums_fail_as_malformed_results(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "unhashable-ocr.pptx"
    image_path = tmp_path / "unhashable-ocr.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    presentation.save(deck)
    baseline = pptx_extraction._extract_pptx_in_process(deck, ocr=False)

    for field, invalid in (("status", []), ("confidence", [])):
        extraction = json.loads(json.dumps(baseline))
        extraction["per_slide_visual"][0]["text_channels"][0][field] = invalid
        with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
            pptx_evidence._decode_extraction_worker_payload(
                {
                    "schema_version": 1,
                    "status": "available",
                    "extraction": extraction,
                }
            )
        assert caught.value.reason_code == "pptx_probe_malformed_result"

    extraction = json.loads(json.dumps(baseline))
    extraction["per_slide_visual"][0]["text_channels"][0]["ocr_receipts"][0][
        "result_status"
    ] = []
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence._decode_extraction_worker_payload(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": extraction,
            }
        )
    assert caught.value.reason_code == "pptx_probe_malformed_result"


@pytest.mark.parametrize(
    "changes",
    [
        {"result_confidence": None},
        {"result_confidence": 49.9},
        {
            "result_status": "low_confidence_text",
            "result_confidence": 90.0,
            "trustworthy_text": False,
        },
        {"engine": "injected"},
        {"engine_version": None},
        {"error": "unexpected"},
    ],
)
def test_ocr_receipt_status_is_bound_to_production_semantics(
    pptx_evidence,
    changes,
) -> None:
    receipt = {
        "attempted": True,
        "engine": "tesseract",
        "engine_version": "5.3.4",
        "result_status": "text_recovered",
        "result_confidence": 93.0,
        "error": None,
        "part_name": "ppt/media/image1.png",
        "asset_sha256": "a" * 64,
        "shape_path": ["Picture 1"],
        "recovered_text": "label",
        "trustworthy_text": True,
    }
    receipt.update(changes)

    assert pptx_evidence._valid_ocr_receipt(receipt) is False


@pytest.mark.parametrize("lane", ["animation", "build"])
def test_native_timing_rejects_counts_outside_required_container(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
    lane: str,
) -> None:
    deck = tmp_path / f"timing-{lane}.pptx"
    _write_deck(deck, with_image=False)
    extraction = pptx_extraction._extract_pptx_in_process(deck, ocr=False)
    timing = json.loads(json.dumps(extraction["per_slide_visual"][0]["native_timing"]))
    if lane == "animation":
        timing["animation_behavior_counts"]["general"] = 1
        timing["animation_behavior_counts"]["total"] = 1
        timing["has_animation_behaviors"] = True
    else:
        timing["build_entry_counts"]["paragraph"] = 1
        timing["build_entry_counts"]["total"] = 1
        timing["has_build_entries"] = True

    assert pptx_evidence._valid_native_timing(timing) is False


def test_native_audit_rejects_unbounded_reason_key_as_structured_error(
    pptx_extraction,
    pptx_evidence,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "unbounded-reason-key.pptx"
    _write_deck(deck, with_image=False)
    audit = json.loads(
        json.dumps(
            pptx_extraction._extract_pptx_in_process(deck, ocr=False)[
                "native_deck_audit"
            ]
        )
    )
    audit["render_required_slide_numbers"] = [1]
    audit["render_required_reasons"] = {"1" * 5000: ["chart"]}

    with pytest.raises(pptx_evidence.PptxEvidenceError, match="bounded canonical"):
        pptx_evidence.validate_native_deck_audit(audit, slide_count=1)


def test_public_extractor_runs_healthy_deck_through_real_worker(
    pptx_extraction,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "healthy-worker.pptx"
    _write_deck(deck, with_image=False)

    result = pptx_extraction.extract_pptx(deck, ocr=False)

    assert result["schema_version"] == 4
    assert result["pipeline_version"] == "1.4.0"
    assert result["slide_count"] == 1
    assert result["pptx_path"] == str(deck)


def test_full_extraction_rejects_render_generation_change(
    pptx_extraction,
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "source.pptx"
    _write_deck(deck, with_image=False)
    rendered = tmp_path / "rendered.pdf"
    _write_pdf(rendered, page_count=1)
    in_process = pptx_extraction._extract_pptx_in_process(
        deck,
        ocr=False,
        rendered_pdf_path=rendered,
        inspected_page_ranges=[[1, 1]],
    )

    def replacing_worker(*_args, **_kwargs):
        original = bytearray(rendered.read_bytes())
        original[-1] ^= 1
        rendered.write_bytes(original)
        return _worker_result(
            {
                "schema_version": 1,
                "status": "available",
                "extraction": in_process,
            }
        )

    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        replacing_worker,
    )
    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_evidence.run_supervised_pptx_extraction(
            deck,
            ocr=False,
            rendered_pdf_path=rendered,
            inspected_page_ranges=[[1, 1]],
        )

    assert caught.value.reason_code == "pptx_artifact_changed"


@pytest.mark.parametrize(
    ("size_bytes", "slide_count"),
    [(True, 1), (1, True), (0, 1), (1, 0)],
)
def test_native_audit_builder_rejects_nonpositive_and_boolean_integers(
    pptx_evidence,
    size_bytes,
    slide_count,
) -> None:
    with pytest.raises(pptx_evidence.PptxEvidenceError):
        pptx_evidence.build_native_deck_audit(
            source_pptx_sha256="a" * 64,
            source_pptx_size_bytes=size_bytes,
            slide_count=slide_count,
            render_required_reasons={},
        )


def test_native_audit_rejects_slide_count_beyond_archive_member_bound(
    pptx_evidence,
) -> None:
    with pytest.raises(pptx_evidence.PptxEvidenceError):
        pptx_evidence.build_native_deck_audit(
            source_pptx_sha256="a" * 64,
            source_pptx_size_bytes=1,
            slide_count=pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS + 1,
            render_required_reasons={},
        )

    audit = pptx_evidence.build_native_deck_audit(
        source_pptx_sha256="a" * 64,
        source_pptx_size_bytes=1,
        slide_count=1,
        render_required_reasons={},
    )
    audit["slide_count"] = 10**31
    with pytest.raises(pptx_evidence.PptxEvidenceError):
        pptx_evidence.validate_native_deck_audit(audit)


def test_render_receipt_rejects_huge_interval_before_coverage_expansion(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        pptx_evidence,
        "_pages_covered_by_ranges",
        lambda *_args, **_kwargs: pytest.fail("huge interval reached coverage scan"),
    )
    huge = 10**31
    receipt = {
        "schema_version": pptx_evidence.RENDER_INSPECTION_SCHEMA_VERSION,
        "source_pptx_sha256": "a" * 64,
        "rendered_pdf_sha256": "b" * 64,
        "rendered_pdf_size_bytes": 1,
        "rendered_page_count": huge,
        "inspected_page_ranges": [[1, huge]],
        "inspected_required_slide_numbers": [1],
        "complete": True,
        "binding_sha256": "c" * 64,
    }

    with pytest.raises(pptx_evidence.PptxEvidenceError):
        pptx_evidence.validate_rendered_page_inspection(
            receipt,
            required_slide_numbers=[1],
            slide_count=huge,
        )


def test_render_builder_rejects_boolean_slide_count_before_file_io(
    pptx_evidence,
    tmp_path: Path,
) -> None:
    with pytest.raises(pptx_evidence.PptxEvidenceError, match="positive integer"):
        pptx_evidence.build_rendered_page_inspection(
            source_pptx_sha256="a" * 64,
            rendered_pdf_path=tmp_path / "missing.pdf",
            inspected_page_ranges=[],
            required_slide_numbers=[],
            slide_count=True,
        )


def test_native_audit_rejects_future_or_extra_fields(
    pptx_evidence,
) -> None:
    audit = pptx_evidence.build_native_deck_audit(
        source_pptx_sha256="a" * 64,
        source_pptx_size_bytes=1,
        slide_count=1,
        render_required_reasons={},
    )
    future = dict(audit, schema_version=2)
    with pytest.raises(pptx_evidence.PptxEvidenceError, match="schema_version"):
        pptx_evidence.validate_native_deck_audit(future)
    extra = dict(audit, unexpected=True)
    with pytest.raises(pptx_evidence.PptxEvidenceError, match="unknown"):
        pptx_evidence.validate_native_deck_audit(extra)


@pytest.mark.parametrize("value", [float("nan"), float("inf"), -1.0, 101.0, True])
def test_finite_confidence_rejects_nonfinite_or_out_of_domain(
    pptx_evidence,
    value,
) -> None:
    assert pptx_evidence.finite_confidence(value) is None


def test_supervised_worker_main_reports_closed_supervisor_failure(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    def fail() -> int:
        raise pptx_evidence.SupervisorError("worker_output_limit_exceeded")

    monkeypatch.setattr(pptx_evidence, "_run_supervised_worker_child", fail)
    monkeypatch.setattr(
        pptx_evidence.sys,
        "argv",
        [pptx_evidence.__file__, pptx_evidence.PPTX_SUPERVISED_WORKER_FLAG],
    )

    assert pptx_evidence._main() == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert (
        captured.err == "pptx supervised worker failed: worker_output_limit_exceeded\n"
    )
    assert "Traceback" not in captured.err


def test_supervised_worker_main_closes_unexpected_failure_diagnostic(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    leaked_path = "/private/vault/source.pptx"

    def fail() -> int:
        raise RuntimeError(f"failure at {leaked_path}")

    monkeypatch.setattr(pptx_evidence, "_run_supervised_worker_child", fail)
    monkeypatch.setattr(
        pptx_evidence.sys,
        "argv",
        [pptx_evidence.__file__, pptx_evidence.PPTX_SUPERVISED_WORKER_FLAG],
    )

    assert pptx_evidence._main() == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == "pptx supervised worker failed: unexpected_error\n"
    assert leaked_path not in captured.err
    assert "Traceback" not in captured.err


def test_supervised_worker_main_preserves_success_output_contract(
    pptx_evidence,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    monkeypatch.setattr(pptx_evidence, "_run_supervised_worker_child", lambda: 0)
    monkeypatch.setattr(
        pptx_evidence.sys,
        "argv",
        [pptx_evidence.__file__, pptx_evidence.PPTX_SUPERVISED_WORKER_FLAG],
    )

    assert pptx_evidence._main() == 0
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == ""


def test_extractor_confines_broad_handlers_to_process_boundaries(
    pptx_extraction,
) -> None:
    tree = ast.parse(Path(pptx_extraction.__file__).read_text(encoding="utf-8"))
    broad_handlers: list[tuple[str, int]] = []
    for function in (
        node for node in ast.walk(tree) if isinstance(node, ast.FunctionDef)
    ):
        for node in ast.walk(function):
            if not isinstance(node, ast.ExceptHandler):
                continue
            if isinstance(node.type, ast.Name) and node.type.id == "Exception":
                broad_handlers.append((function.name, node.lineno))

    assert sorted(name for name, _line in broad_handlers) == ["_main", "main"]


def test_single_file_cli_reports_parser_failure_without_traceback(
    pptx_extraction,
    tmp_path: Path,
) -> None:
    deck = tmp_path / "broken.pptx"
    deck.write_bytes(b"not an OOXML ZIP package")

    completed = subprocess.run(
        [sys.executable, pptx_extraction.__file__, str(deck), "--no-ocr"],
        capture_output=True,
        text=True,
        check=False,
        timeout=30,
    )

    assert completed.returncode == 1
    assert completed.stdout == ""
    assert completed.stderr == "ERROR: pptx_invalid_container\n"
    assert str(deck) not in completed.stderr
    assert "Traceback" not in completed.stderr
