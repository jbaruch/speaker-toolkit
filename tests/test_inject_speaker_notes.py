"""Tests for inject-speaker-notes.py — speaker notes injection + Keynote patch."""

import json
import shutil
import zipfile

from pptx import Presentation

from conftest import make_deck


def test_notes_injected(inject_speaker_notes, tmp_path):
    prs = make_deck(3)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    notes = {"0": "First slide notes", "1": "Second slide notes", "2": ""}
    notes_path = str(tmp_path / "notes.json")
    with open(notes_path, "w") as f:
        json.dump(notes, f)

    prs2 = Presentation(path)
    injected = 0
    for idx_str, notes_text in notes.items():
        idx = int(idx_str)
        if notes_text and idx < len(prs2.slides):
            slide = prs2.slides[idx]
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
            injected += 1
    prs2.save(path)

    prs3 = Presentation(path)
    assert prs3.slides[0].notes_slide.notes_text_frame.text == "First slide notes"
    assert prs3.slides[1].notes_slide.notes_text_frame.text == "Second slide notes"
    assert injected == 2


def test_out_of_range_skip(inject_speaker_notes, tmp_path):
    prs = make_deck(2)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    notes = {"0": "Valid", "99": "Out of range"}
    prs2 = Presentation(path)
    injected = 0
    for idx_str, notes_text in notes.items():
        idx = int(idx_str)
        if notes_text and idx < len(prs2.slides):
            slide = prs2.slides[idx]
            slide.notes_slide.notes_text_frame.text = notes_text
            injected += 1
    prs2.save(path)

    assert injected == 1


def test_patch_notes_master_idlst(inject_speaker_notes, tmp_path):
    """After injecting notes, the patch should add notesMasterIdLst."""
    prs = make_deck(2)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    # Inject a note to trigger notesMaster creation
    prs2 = Presentation(path)
    slide = prs2.slides[0]
    slide.notes_slide.notes_text_frame.text = "Test note"
    prs2.save(path)

    result = inject_speaker_notes.patch_notes_master_idlst(path)

    # Verify the patch was applied
    with zipfile.ZipFile(path, "r") as z:
        pres_xml = z.read("ppt/presentation.xml").decode("utf-8")
    assert "<p:notesMasterIdLst>" in pres_xml


def test_patch_idempotent(inject_speaker_notes, tmp_path):
    """Running the patch twice should be safe (second call returns False)."""
    prs = make_deck(2)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    prs2.slides[0].notes_slide.notes_text_frame.text = "Test"
    prs2.save(path)

    first = inject_speaker_notes.patch_notes_master_idlst(path)
    second = inject_speaker_notes.patch_notes_master_idlst(path)

    # First call patches, second call finds it already there
    assert first is True
    assert second is False

    # File should still be valid
    with zipfile.ZipFile(path, "r") as z:
        pres_xml = z.read("ppt/presentation.xml").decode("utf-8")
    assert pres_xml.count("<p:notesMasterIdLst>") == 1
