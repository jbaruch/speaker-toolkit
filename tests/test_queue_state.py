"""Queue-state contract tests for vault-ingress.

All timestamps are injected and every fixture is local. The CLI never reaches a
network or reads subagent returns.
"""

import copy
import hashlib
import json
import os
import shutil
import struct
import subprocess
import sys
from pathlib import Path
import zipfile

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter


SCRIPT = (
    Path(__file__).parents[1]
    / "skills"
    / "vault-ingress"
    / "scripts"
    / "queue-state.py"
)
NOW = "2026-07-31T18:00:00+00:00"


def _talk(
    video_id,
    *,
    status="pending",
    filename=None,
    video=True,
    youtube_identity=True,
) -> dict[str, object]:
    filename = filename or f"playlist-{video_id}.md"
    return {
        "schema_version": 5,
        "filename": filename,
        "title": filename,
        "status": status,
        "video_url": (
            f"https://www.youtube.com/watch?v={video_id}" if video else ""
        ),
        "youtube_id": video_id if youtube_identity else None,
    }


def _scored_talk(
    video_id,
    *,
    filename,
    fingerprint,
    scoring_schema,
    status="processed",
    score=2,
    generation_status="current",
    generation_reasons=None,
):
    if generation_reasons is None:
        generation_reasons = []
    talk = _talk(video_id, status=status, filename=filename)
    talk.update({
        "processed_date": "2001-01-01T00:00:00+00:00",
        "pattern_scoring_generation_status": generation_status,
        "pattern_scoring_generation_reasons": generation_reasons,
        "pattern_catalog_fingerprint": fingerprint,
        "pattern_scoring_schema_version": scoring_schema,
        "pattern_score": score,
        "pattern_observations": {"pattern_score": score},
    })
    return talk


def _write_db(tmp_path, talks, *, config=None, current=True):
    transcripts = tmp_path / "transcripts"
    transcripts.mkdir(exist_ok=True)
    for talk in talks:
        if talk.get("pattern_scoring_schema_version") != 5:
            continue
        observations = talk.get("pattern_observations")
        if not isinstance(observations, dict):
            continue
        artifact = transcripts / f"{Path(str(talk['filename'])).stem}.txt"
        content = ("synthetic evidence " * 225).strip() + "\n"
        artifact.write_text(content, encoding="utf-8")
        digest = hashlib.sha256(content.encode("utf-8")).hexdigest()
        relative = artifact.relative_to(tmp_path).as_posix()
        quality = artifact.with_suffix(".quality.json")
        quality.write_text(
            json.dumps({
                "schema_version": 1,
                "transcript_sha256": digest,
                "policy": {
                    "schema_version": 1,
                    "min_words": 400,
                    "duration_seconds": None,
                },
                "provenance": {"kind": "fixed_default"},
            }),
            encoding="utf-8",
        )
        quality_digest = hashlib.sha256(quality.read_bytes()).hexdigest()
        quality_relative = quality.relative_to(tmp_path).as_posix()
        talk.setdefault("transcript_path", relative)
        observations.setdefault("evidence_schema_version", 2)
        observations.setdefault("evidence_sources", ["transcript"])
        observations.setdefault(
            "source_inspection",
            [
                {
                    "source": "transcript",
                    "line_ranges": [[1, 1]],
                    "line_count": 1,
                    "coverage_complete": True,
                    "artifact_root": "vault",
                    "artifact_path": relative,
                    "artifact_sha256": digest,
                    "quality_artifact_root": "vault",
                    "quality_artifact_path": quality_relative,
                    "quality_artifact_sha256": quality_digest,
                }
            ],
        )
        score = observations.get("pattern_score", 0)
        patterns = [
            {
                "pattern_id": f"fixture-pattern-{index}",
                "confidence": "moderate",
                "evidence_source": "transcript",
                "evidence": "Synthetic source-located fixture evidence.",
                "evidence_citations": [{
                    "source": "transcript",
                    "channel": "transcript",
                    "quote": "synthetic evidence synthetic evidence",
                    "line_start": 1,
                    "line_end": 1,
                    "artifact_root": "vault",
                    "artifact_path": relative,
                    "artifact_sha256": digest,
                    "quality_artifact_root": "vault",
                    "quality_artifact_path": quality_relative,
                    "quality_artifact_sha256": quality_digest,
                }],
            }
            for index in range(max(int(score), 0))
        ]
        antipatterns = [
            {
                "pattern_id": f"fixture-antipattern-{index}",
                "confidence": "moderate",
                "evidence_source": "transcript",
                "evidence": "Synthetic source-located fixture evidence.",
                "evidence_citations": [{
                    "source": "transcript",
                    "channel": "transcript",
                    "quote": "synthetic evidence synthetic evidence",
                    "line_start": 1,
                    "line_end": 1,
                    "artifact_root": "vault",
                    "artifact_path": relative,
                    "artifact_sha256": digest,
                    "quality_artifact_root": "vault",
                    "quality_artifact_path": quality_relative,
                    "quality_artifact_sha256": quality_digest,
                }],
            }
            for index in range(max(-int(score), 0))
        ]
        observations.setdefault("patterns_detected", patterns)
        observations.setdefault("antipatterns_detected", antipatterns)
        observations.setdefault("applicability_assessments", [])
        observations.setdefault("not_evaluable", [])
        outcomes = [
            {"pattern_id": item["pattern_id"], "outcome": "detected"}
            for item in patterns + antipatterns
        ]
        outcomes.sort(key=lambda item: item["pattern_id"])
        observations.setdefault("pattern_outcomes", outcomes)
        fingerprint = talk.get("pattern_catalog_fingerprint")
        if fingerprint is not None:
            payload = {
                "pattern_scoring_schema_version": 5,
                "pattern_catalog_fingerprint": fingerprint,
                "opportunity_states": [
                    {
                        "pattern_id": item["pattern_id"],
                        "opportunity_state": "evaluable",
                    }
                    for item in outcomes
                ],
            }
            observations.setdefault(
                "opportunity_coverage_identity",
                hashlib.sha256(json.dumps(
                    payload,
                    ensure_ascii=False,
                    separators=(",", ":"),
                    sort_keys=True,
                ).encode("utf-8")).hexdigest(),
            )
    path = tmp_path / "tracking-database.json"
    database = {
        "config": {**(config or {})},
        "talks": talks,
        "pptx_catalog": [],
        "qr_codes": [],
        "resources": [],
        "thumbnails": [],
        "confirmed_intents": [],
        "improvement_goals": [],
    }
    if current:
        database["schema_version"] = 1
        database["config"]["schema_version"] = 1
    else:
        for talk in talks:
            talk.pop("schema_version", None)
    path.write_text(json.dumps(database, indent=2), encoding="utf-8")
    return path


def _write_verified_transcript(tmp_path, name="talk"):
    transcript = tmp_path / "transcripts" / f"{name}.txt"
    transcript.parent.mkdir(exist_ok=True)
    text = " ".join(["substantive transcript evidence"] * 200)
    transcript.write_text(text, encoding="utf-8")
    transcript.with_suffix(".quality.json").write_text(
        json.dumps({
            "schema_version": 1,
            "transcript_sha256": hashlib.sha256(
                text.encode("utf-8")).hexdigest(),
            "policy": {
                "schema_version": 1,
                "min_words": 400,
                "duration_seconds": None,
            },
            "provenance": {"kind": "fixed_default"},
        }),
        encoding="utf-8",
    )
    return transcript


def _write_crc_damaged_media_pptx(path):
    """Create a deck whose media needs loss-reporting placeholder recovery."""
    path.parent.mkdir(parents=True, exist_ok=True)
    image_path = path.with_suffix(".png")
    Image.new("RGB", (64, 64), "navy").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    presentation.save(path)
    with zipfile.ZipFile(path) as archive:
        member = next(
            item
            for item in archive.infolist()
            if item.filename.startswith("ppt/media/") and item.file_size
        )
    package = bytearray(path.read_bytes())
    name_size, extra_size = struct.unpack_from(
        "<HH", package, member.header_offset + 26
    )
    payload_offset = member.header_offset + 30 + name_size + extra_size
    package[payload_offset + (member.compress_size // 2)] ^= 0xFF
    path.write_bytes(package)
    return member.filename


def _read_db(path):
    return json.loads(path.read_text(encoding="utf-8"))


def _run(path, *arguments):
    environment = os.environ.copy()
    environment["PYTHONDONTWRITEBYTECODE"] = "1"
    return subprocess.run(
        [sys.executable, str(SCRIPT), str(path), *arguments],
        capture_output=True,
        text=True,
        env=environment,
    )


def _claim(path, *, run_id="run-1", batch_id="batch-1", limit=5, filenames=()):
    arguments = [
        "claim",
        "--run-id", run_id,
        "--batch-id", batch_id,
        "--now", NOW,
        "--limit", str(limit),
    ]
    for filename in filenames:
        arguments.extend(("--filename", filename))
    return _run(path, *arguments)


def test_claim_recovers_the_two_stranded_transcript_statuses(tmp_path):
    """The two real vault rows with videos must re-enter the processable queue."""
    talks = [
        _talk("eixm_f7Jpdc", status="skipped_no_transcript"),
        _talk("QS-_4k7o7A4", status="skipped_no_transcript"),
        _talk(
            "abcdefghijk", status="skipped_no_video", video=False,
            youtube_identity=False, filename="catalog-no-source.md",
        ),
        _talk("lmnopqrstuv", status="skipped_duplicate"),
        _talk(
            "wxyzABCDEF0", status="pending", video=False,
            youtube_identity=False, filename="catalog-pending.md",
        ),
    ]
    path = _write_db(tmp_path, talks)

    result = _claim(path, limit=10)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert [item["filename"] for item in payload["claimed"]] == [
        "playlist-QS-_4k7o7A4.md",
        "playlist-eixm_f7Jpdc.md",
    ]
    assert {item["previous_status"] for item in payload["claimed"]} == {"pending"}
    assert all(item["reprocess_generation"] == 1 for item in payload["claimed"])
    records = {talk["filename"]: talk for talk in _read_db(path)["talks"]}
    assert records["playlist-eixm_f7Jpdc.md"]["status"] == "reprocessing-inflight"
    assert records["playlist-QS-_4k7o7A4.md"]["status"] == "reprocessing-inflight"
    assert records["catalog-no-source.md"]["status"] == "skipped_no_sources"
    assert records["playlist-lmnopqrstuv.md"]["status"] == "skipped_duplicate"
    assert records["catalog-pending.md"]["status"] == "pending"
    assert not [item for item in tmp_path.iterdir() if item.name.endswith(".partial")]


def test_legacy_no_video_status_with_video_is_normalized_and_claimed(tmp_path):
    path = _write_db(
        tmp_path,
        [_talk("eg6gqvUFh6Q", status="skipped_no_video")],
    )

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"] == [{
        "filename": "playlist-eg6gqvUFh6Q.md",
        "previous_status": "skipped_no_video",
        "status": "pending",
        "video_present": True,
        "source_capabilities": ["video", "transcript"],
    }]
    assert payload["claimed"][0]["previous_status"] == "pending"


@pytest.mark.parametrize("source_kind,expected_capability", [
    ("remote_slides", "slides"),
    ("pptx", "slides"),
    ("pdf", "slides"),
    ("transcript", "transcript"),
])
def test_legacy_no_video_talk_with_nonvideo_source_is_claimable(
        tmp_path, source_kind, expected_capability):
    talk = _talk(
        "abcdefghijk", status="skipped_no_video", video=False,
        youtube_identity=False,
        filename="source-only.md",
    )
    if source_kind == "remote_slides":
        talk["slides_url"] = "https://drive.google.com/open?id=deck"
    elif source_kind == "pptx":
        deck = tmp_path / "decks" / "talk.pptx"
        deck.parent.mkdir()
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(deck)
        talk["pptx_path"] = "decks/talk.pptx"
    elif source_kind == "pdf":
        pdf = tmp_path / "slides" / "talk.pdf"
        pdf.parent.mkdir()
        writer = PdfWriter()
        writer.add_blank_page(width=640, height=480)
        with pdf.open("wb") as stream:
            writer.write(stream)
        talk["slides_local_path"] = "slides/talk.pdf"
    else:
        transcript = tmp_path / "transcripts" / "abcdefghijk.txt"
        transcript.parent.mkdir()
        transcript.write_text(" ".join(["synthetic"] * 450), encoding="utf-8")
        talk["transcript_path"] = "transcripts/abcdefghijk.txt"
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"] == [{
        "filename": "source-only.md",
        "previous_status": "skipped_no_video",
        "status": "pending",
        "video_present": False,
        "source_capabilities": [expected_capability],
    }]
    assert payload["claimed"][0]["filename"] == "source-only.md"
    assert payload["claimed"][0]["previous_status"] == "pending"


@pytest.mark.parametrize(
    "broken_field",
    ["pptx_path", "slides_local_path", "transcript_path"],
)
def test_broken_local_path_is_not_a_processable_capability(
    tmp_path,
    broken_field,
):
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="broken-local.md",
    )
    talk[broken_field] = {
        "pptx_path": "missing/source.pptx",
        "slides_local_path": "missing/source.pdf",
        "transcript_path": "transcripts/source.txt",
    }[broken_field]
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"] == [
        {
            "filename": "broken-local.md",
            "previous_status": "skipped_no_video",
            "status": "skipped_no_sources",
            "video_present": False,
            "source_capabilities": [],
        }
    ]
    assert payload["claimed"] == []


def test_low_quality_transcript_is_not_a_processable_capability(tmp_path):
    transcript = tmp_path / "transcripts" / "abcdefghijk.txt"
    transcript.parent.mkdir()
    transcript.write_text("too short to be a talk", encoding="utf-8")
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="short-transcript.md",
    )
    talk["transcript_path"] = "transcripts/abcdefghijk.txt"
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"][0]["status"] == "skipped_no_sources"
    assert payload["normalizations"][0]["source_capabilities"] == []
    assert payload["claimed"] == []


def test_symlinked_transcript_is_not_a_processable_capability(tmp_path):
    source = tmp_path / "external-transcript.txt"
    source.write_text(" ".join(["synthetic"] * 450), encoding="utf-8")
    transcript = tmp_path / "transcripts" / "abcdefghijk.txt"
    transcript.parent.mkdir()
    transcript.symlink_to(source)
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="symlinked-transcript.md",
    )
    talk["transcript_path"] = "transcripts/abcdefghijk.txt"
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"][0]["status"] == "skipped_no_sources"
    assert payload["normalizations"][0]["source_capabilities"] == []
    assert payload["claimed"] == []


def test_local_artifact_path_escape_is_not_a_processable_capability(tmp_path):
    outside = tmp_path.parent / f"{tmp_path.name}-outside.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(outside)
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="escaped-deck.md",
    )
    talk["pptx_path"] = f"../{outside.name}"
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"][0]["status"] == "skipped_no_sources"
    assert payload["normalizations"][0]["source_capabilities"] == []
    assert payload["claimed"] == []


def test_invalid_transcript_does_not_hide_valid_deck_from_queue(tmp_path):
    deck = tmp_path / "decks" / "talk.pptx"
    deck.parent.mkdir()
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(deck)
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="deck-survives.md",
    )
    talk.update({
        "transcript_path": "../bad.txt",
        "pptx_path": "decks/talk.pptx",
    })
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"][0]["source_capabilities"] == ["slides"]
    assert payload["claimed"][0]["filename"] == "deck-survives.md"


def test_invalid_deck_does_not_hide_valid_transcript_from_queue(tmp_path):
    transcript = _write_verified_transcript(tmp_path)
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="transcript-survives.md",
    )
    talk.update({
        "transcript_path": transcript.relative_to(tmp_path).as_posix(),
        "transcript_source": "manual",
        "pptx_path": "../bad.pptx",
    })
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"][0]["source_capabilities"] == [
        "transcript"]
    assert payload["claimed"][0]["filename"] == "transcript-survives.md"


def test_configured_pptx_root_is_used_for_queue_eligibility(tmp_path):
    source_root = tmp_path / "configured-pptx"
    source_root.mkdir()
    deck = source_root / "talk.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(deck)
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="configured-pptx.md",
    )
    talk["pptx_path"] = deck.name
    path = _write_db(
        tmp_path,
        [talk],
        config={"pptx_source_dir": str(source_root)},
    )

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"][0]["source_capabilities"] == ["slides"]
    assert payload["claimed"][0]["filename"] == "configured-pptx.md"


@pytest.mark.skipif(
    shutil.which("ffmpeg") is None or shutil.which("ffprobe") is None,
    reason="local video capability requires ffmpeg and ffprobe",
)
def test_readable_identity_bound_local_video_is_claimable(tmp_path):
    video = tmp_path / "videos" / "abcdefghijk.mp4"
    video.parent.mkdir()
    created = subprocess.run(
        [
            "ffmpeg",
            "-loglevel",
            "error",
            "-f",
            "lavfi",
            "-i",
            "color=c=black:s=160x90:r=1",
            "-t",
            "1",
            "-pix_fmt",
            "yuv420p",
            "-y",
            str(video),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    assert created.returncode == 0, created.stderr
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="local-video.md",
    )
    talk["video_local_path"] = "videos/abcdefghijk.mp4"
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"][0]["source_capabilities"] == ["video"]
    assert payload["claimed"][0]["filename"] == "local-video.md"


def test_retry_status_with_only_a_missing_local_artifact_is_not_claimed(tmp_path):
    talk = _talk(
        "abcdefghijk",
        status="skipped_download_failed",
        video=False,
        youtube_identity=False,
        filename="retry-missing.md",
    )
    talk["transcript_path"] = "transcripts/abcdefghijk.txt"
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["claimed"] == []
    assert payload["remaining_eligible"] == 0
    assert _read_db(path)["talks"][0]["status"] == "skipped_download_failed"


def test_required_degraded_pptx_cannot_create_a_fresh_current_claim(tmp_path):
    deck = tmp_path / "decks" / "damaged.pptx"
    damaged_part = _write_crc_damaged_media_pptx(deck)
    talk = _talk(
        "abcdefghijk",
        filename="required-damaged-deck.md",
    )
    talk.update({
        "pptx_path": "decks/damaged.pptx",
        "slide_source": "pptx",
    })
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    result = _claim(path, filenames=(str(talk["filename"]),))

    assert result.returncode == 2
    payload = json.loads(result.stdout)
    assert "placeholder archive recovery" in payload["error"]
    assert damaged_part in payload["error"]
    assert "restore or re-export" in payload["error"]
    assert path.read_bytes() == before
    persisted = _read_db(path)["talks"][0]
    assert persisted["status"] == "pending"
    assert "_queue_claim" not in persisted


def test_claim_reassesses_after_eligibility_and_refuses_mutated_deck(
    queue_state,
):
    talk = _talk("abcdefghijk", filename="changed-after-selection.md")
    talk.update({"pptx_path": "decks/talk.pptx", "slide_source": "pptx"})
    clean = {
        "verified_capabilities": ("slides",),
        "verified_evidence_sources": ("native_deck",),
        "acquisition_capabilities": (),
        "repair_capabilities": (),
        "degraded_evidence_sources": {},
    }
    degraded = {
        "verified_capabilities": (),
        "verified_evidence_sources": (),
        "acquisition_capabilities": (),
        "repair_capabilities": (),
        "degraded_evidence_sources": {
            "native_deck": {
                "reason_code": "pptx_archive_recovery_required",
                "archive_recovery": [
                    {
                        "part_name": "ppt/media/image1.png",
                        "status": "recovered_with_placeholder_asset",
                    }
                ],
            }
        },
    }
    assessments = iter([clean, degraded])

    def assessor(_talk):
        return next(assessments)

    assert queue_state.has_claimable_source(
        talk, capability_assessor=assessor
    ) is True
    before = copy.deepcopy(talk)
    with pytest.raises(queue_state.QueueStateError, match="cannot claim"):
        queue_state.claim_talk(
            talk,
            "run",
            "batch",
            NOW,
            {},
            capability_assessor=assessor,
        )
    assert talk == before


def test_unused_optional_degraded_pptx_does_not_block_independent_source(
    tmp_path,
):
    deck = tmp_path / "decks" / "optional-damaged.pptx"
    _write_crc_damaged_media_pptx(deck)
    talk = _talk(
        "abcdefghijk",
        filename="optional-damaged-deck.md",
    )
    talk.update({
        "pptx_path": "decks/optional-damaged.pptx",
        "slide_source": "pdf",
    })
    path = _write_db(tmp_path, [talk])

    result = _claim(path, filenames=(str(talk["filename"]),))

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert [item["filename"] for item in payload["claimed"]] == [
        "optional-damaged-deck.md"
    ]


@pytest.mark.parametrize("slide_source", [None, "none"])
def test_degraded_deck_without_declared_or_independent_lane_is_not_claimed(
    tmp_path,
    slide_source,
):
    deck = tmp_path / "decks" / "undeclared-damaged.pptx"
    _write_crc_damaged_media_pptx(deck)
    talk = _talk(
        "abcdefghijk",
        filename="undeclared-damaged-deck.md",
        video=False,
        youtube_identity=False,
    )
    talk["pptx_path"] = "decks/undeclared-damaged.pptx"
    if slide_source is not None:
        talk["slide_source"] = slide_source
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    assert json.loads(result.stdout)["claimed"] == []
    assert path.read_bytes() == before


def test_optional_degraded_deck_allows_an_independent_healthy_transcript(
    tmp_path,
):
    deck = tmp_path / "decks" / "optional-damaged.pptx"
    _write_crc_damaged_media_pptx(deck)
    transcript = _write_verified_transcript(tmp_path, "independent")
    talk = _talk(
        "abcdefghijk",
        filename="independent-transcript.md",
        video=False,
        youtube_identity=False,
    )
    talk.update({
        "pptx_path": "decks/optional-damaged.pptx",
        "slide_source": "none",
        "transcript_path": transcript.relative_to(tmp_path).as_posix(),
        "transcript_source": "manual",
    })
    path = _write_db(tmp_path, [talk])

    result = _claim(path, filenames=(str(talk["filename"]),))

    assert result.returncode == 0, result.stderr
    assert [item["filename"] for item in json.loads(result.stdout)["claimed"]] == [
        "independent-transcript.md"
    ]


def test_manual_provenance_label_without_artifact_is_not_a_capability(tmp_path):
    talk = _talk(
        "abcdefghijk",
        status="skipped_no_video",
        video=False,
        youtube_identity=False,
        filename="label-only.md",
    )
    talk["transcript_source"] = "manual"
    path = _write_db(tmp_path, [talk])

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["normalizations"] == [{
        "filename": "label-only.md",
        "previous_status": "skipped_no_video",
        "status": "skipped_no_sources",
        "video_present": False,
        "source_capabilities": [],
    }]
    assert payload["claimed"] == []


def test_normalize_requeues_every_noncurrent_or_stale_generation_atomically(
        tmp_path, return_validation):
    fingerprint = return_validation.load_catalog().fingerprint
    scoring_schema = return_validation.PATTERN_SCORING_SCHEMA_VERSION
    other_fingerprint = "0" * 64 if fingerprint != "0" * 64 else "1" * 64
    current = _scored_talk(
        "AAAAAAAAAAA", filename="current-old-date.md",
        fingerprint=fingerprint, scoring_schema=scoring_schema,
    )
    legacy = _scored_talk(
        "BBBBBBBBBBB", filename="legacy-recent.md",
        fingerprint=fingerprint, scoring_schema=scoring_schema, score=100,
        generation_status="legacy_unbaselineable",
    )
    legacy.pop("pattern_catalog_fingerprint")
    legacy.pop("pattern_scoring_schema_version")
    old_catalog = _scored_talk(
        "CCCCCCCCCCC", filename="old-catalog-recent.md",
        fingerprint=other_fingerprint, scoring_schema=scoring_schema,
    )
    old_schema = _scored_talk(
        "DDDDDDDDDDD", filename="old-schema-recent.md",
        fingerprint=fingerprint, scoring_schema=scoring_schema - 1,
        status="processed_partial",
    )
    missing = _scored_talk(
        "EEEEEEEEEEE", filename="missing-generation.md",
        fingerprint=fingerprint, scoring_schema=scoring_schema,
    )
    missing.pop("pattern_scoring_generation_status")
    legacy_source = _talk(
        "FFFFFFFFFFF", filename="legacy-source-status.md",
        status="skipped_no_video",
    )
    path = _write_db(
        tmp_path,
        [current, legacy, old_catalog, old_schema, missing, legacy_source],
    )

    result = _run(path, "normalize")

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["changed"] == 6
    generation_changes = {
        item["filename"]: item
        for item in payload["normalizations"]
        if "reason_codes" in item
    }
    assert {
        filename: item["reason_codes"]
        for filename, item in generation_changes.items()
    } == {
        "current-old-date.md": ["persisted_evidence_stale"],
        "legacy-recent.md": ["legacy_generation"],
        "old-catalog-recent.md": ["catalog_fingerprint_mismatch"],
        "old-schema-recent.md": ["scoring_schema_version_mismatch"],
        "missing-generation.md": ["missing_generation_status"],
    }
    assert {
        filename: item["reprocess_reason"]
        for filename, item in generation_changes.items()
    } == {
        filename: "pattern_scoring_generation:" + "+".join(item["reason_codes"])
        for filename, item in generation_changes.items()
    }
    records = {talk["filename"]: talk for talk in _read_db(path)["talks"]}
    for filename in generation_changes:
        assert records[filename]["status"] == "needs-reprocessing"
        assert records[filename]["reprocess_reason"] == \
            generation_changes[filename]["reprocess_reason"]
    assert records["legacy-source-status.md"]["status"] == "pending"

    first_bytes = path.read_bytes()
    repeated = _run(path, "normalize")

    assert repeated.returncode == 0, repeated.stderr
    assert json.loads(repeated.stdout)["normalizations"] == []
    assert path.read_bytes() == first_bytes


def test_normalize_preserves_ordered_generation_reasons(tmp_path, return_validation):
    fingerprint = return_validation.load_catalog().fingerprint
    scoring_schema = return_validation.PATTERN_SCORING_SCHEMA_VERSION
    other_fingerprint = "0" * 64 if fingerprint != "0" * 64 else "1" * 64
    talk = _scored_talk(
        "GGGGGGGGGGG", filename="both-generations-stale.md",
        fingerprint=other_fingerprint, scoring_schema=scoring_schema - 1,
    )
    path = _write_db(tmp_path, [talk])

    result = _run(path, "normalize")

    assert result.returncode == 0, result.stderr
    change = json.loads(result.stdout)["normalizations"][0]
    assert change["reason_codes"] == [
        "catalog_fingerprint_mismatch",
        "scoring_schema_version_mismatch",
    ]
    assert change["reprocess_reason"] == (
        "pattern_scoring_generation:catalog_fingerprint_mismatch+"
        "scoring_schema_version_mismatch"
    )


@pytest.mark.parametrize("drift", ["missing", "digest_mismatch"])
def test_normalize_requeues_current_generation_when_evidence_artifact_drifts(
    tmp_path,
    return_validation,
    drift,
):
    fingerprint = return_validation.load_catalog().fingerprint
    talk = _scored_talk(
        "GGGGGGGGGGG",
        filename="artifact-drift.md",
        fingerprint=fingerprint,
        scoring_schema=return_validation.PATTERN_SCORING_SCHEMA_VERSION,
    )
    path = _write_db(tmp_path, [talk])
    stored = _read_db(path)["talks"][0]
    relative = stored["pattern_observations"]["source_inspection"][0][
        "artifact_path"
    ]
    artifact = tmp_path / relative
    if drift == "missing":
        artifact.unlink()
    else:
        artifact.write_text("Replacement with a different digest.\n", encoding="utf-8")

    result = _run(path, "normalize")

    assert result.returncode == 0, result.stderr
    change = json.loads(result.stdout)["normalizations"][0]
    assert change["filename"] == "artifact-drift.md"
    assert change["reason_codes"] == ["persisted_evidence_stale"]
    assert any(drift in detail for detail in change["evidence_freshness_details"])
    assert change["reprocess_reason"] == (
        "pattern_scoring_generation:persisted_evidence_stale"
    )
    repaired = _read_db(path)["talks"][0]
    assert repaired["status"] == "needs-reprocessing"


@pytest.mark.parametrize(
    "case,message",
    [
        ("unknown_status", "pattern_scoring_generation_status must be one of"),
        ("malformed_status", "pattern_scoring_generation_status must be one of"),
        ("current_reasons", "pattern_scoring_generation_reasons must be exactly"),
        ("incomplete_identity", "missing required identity fields"),
        ("malformed_fingerprint", "must be a lowercase 64-character"),
        ("malformed_schema", "pattern_scoring_schema_version must be an integer"),
        ("divergent_score", "promoted pattern_score 2 diverges"),
    ],
)
def test_normalize_rejects_malformed_generation_without_any_write(
        tmp_path, return_validation, case, message):
    fingerprint = return_validation.load_catalog().fingerprint
    scoring_schema = return_validation.PATTERN_SCORING_SCHEMA_VERSION
    malformed = _scored_talk(
        "HHHHHHHHHHH", filename="malformed-generation.md",
        fingerprint=fingerprint, scoring_schema=scoring_schema,
    )
    if case == "unknown_status":
        malformed["pattern_scoring_generation_status"] = "future"
    elif case == "malformed_status":
        malformed["pattern_scoring_generation_status"] = True
    elif case == "current_reasons":
        malformed["pattern_scoring_generation_reasons"] = ["contradiction"]
    elif case == "incomplete_identity":
        malformed.pop("pattern_catalog_fingerprint")
    elif case == "malformed_fingerprint":
        malformed["pattern_catalog_fingerprint"] = "not-a-sha"
    elif case == "malformed_schema":
        malformed["pattern_scoring_schema_version"] = True
    elif case == "divergent_score":
        malformed["pattern_observations"]["pattern_score"] = 1
    legacy_source = _talk(
        "IIIIIIIIIII", filename="would-normalize.md",
        status="skipped_no_video",
    )
    path = _write_db(tmp_path, [legacy_source, malformed])
    before = path.read_bytes()

    result = _run(path, "normalize")

    assert result.returncode == 2
    assert message in json.loads(result.stdout)["error"]
    assert path.read_bytes() == before
    assert not [item for item in tmp_path.iterdir() if item.name.endswith(".partial")]


def test_normalize_does_not_inspect_generation_for_ineligible_statuses(tmp_path):
    inflight = _talk(
        "JJJJJJJJJJJ", filename="inflight.md",
        status="reprocessing-inflight",
    )
    inflight.update({
        "reprocess_generation": 1,
        "_queue_claim": {
            "schema_version": 1,
            "run_id": "active-run",
            "batch_id": "active-batch",
            "claimed_at": NOW,
            "previous_status": "pending",
            "reprocess_generation": 1,
            "state": "claimed",
        },
        "pattern_scoring_generation_status": "future",
        "pattern_score": 4,
        "pattern_observations": {"pattern_score": 5},
    })
    pending = _talk("KKKKKKKKKKK", filename="pending.md")
    pending["pattern_scoring_generation_status"] = "future"
    skipped = _talk(
        "LLLLLLLLLLL", filename="skipped.md", status="skipped_duplicate",
    )
    skipped["pattern_scoring_generation_status"] = True
    path = _write_db(tmp_path, [inflight, pending, skipped])
    before = path.read_bytes()

    result = _run(path, "normalize")

    assert result.returncode == 0, result.stderr
    assert json.loads(result.stdout)["normalizations"] == []
    assert path.read_bytes() == before


def test_generation_requeue_preserves_completed_claim_until_next_claim(
        tmp_path, return_validation):
    fingerprint = return_validation.load_catalog().fingerprint
    scoring_schema = return_validation.PATTERN_SCORING_SCHEMA_VERSION
    talk = _scored_talk(
        "MMMMMMMMMMM", filename="completed-legacy.md",
        fingerprint=fingerprint, scoring_schema=scoring_schema,
        generation_status="legacy_unbaselineable",
    )
    talk.pop("pattern_catalog_fingerprint")
    talk.pop("pattern_scoring_schema_version")
    talk["reprocess_generation"] = 1
    completed_claim = {
        "schema_version": 2,
        "run_id": "completed-run",
        "batch_id": "completed-batch",
        "claimed_at": "2026-07-31T17:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "completed",
        "released_at": "2026-07-31T17:30:00+00:00",
        "release_reason": "return_persisted",
        "result_status": "processed",
        "result_payload_sha256": "0" * 64,
    }
    talk["_queue_claim"] = copy.deepcopy(completed_claim)
    path = _write_db(tmp_path, [talk])

    normalized = _run(path, "normalize")

    assert normalized.returncode == 0, normalized.stderr
    requeued = _read_db(path)["talks"][0]
    assert requeued["status"] == "needs-reprocessing"
    assert requeued["_queue_claim"] == completed_claim
    assert requeued.get("_queue_claim_history", []) == []

    claimed = _claim(
        path, run_id="replacement-run", batch_id="replacement-batch",
        filenames=(talk["filename"],),
    )

    assert claimed.returncode == 0, claimed.stderr
    replacement = _read_db(path)["talks"][0]
    assert replacement["status"] == "reprocessing-inflight"
    assert replacement["_queue_claim"]["reprocess_generation"] == 2
    assert replacement["_queue_claim_history"] == [completed_claim]


@pytest.mark.parametrize("reason", ["please_run_this_again", ["not", "a", "reason"]])
def test_completed_claim_status_drift_rejects_unowned_reprocess_reason(
        tmp_path, reason):
    talk = _talk(
        "NNNNNNNNNNN", filename="unowned-requeue.md",
        status="needs-reprocessing",
    )
    talk.update({
        "reprocess_reason": reason,
        "reprocess_generation": 1,
        "_queue_claim": {
            "schema_version": 2,
            "run_id": "completed-run",
            "batch_id": "completed-batch",
            "claimed_at": "2026-07-31T17:00:00+00:00",
            "previous_status": "pending",
            "reprocess_generation": 1,
            "state": "completed",
            "released_at": "2026-07-31T17:30:00+00:00",
            "release_reason": "return_persisted",
            "result_status": "processed",
            "result_payload_sha256": "0" * 64,
        },
    })
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    result = _run(path, "normalize")

    assert result.returncode == 2
    assert "completed claim result_status" in result.stderr
    assert path.read_bytes() == before


def test_legacy_true_no_source_talk_is_the_only_one_skipped(tmp_path):
    no_source = _talk(
        "abcdefghijk", status="skipped_no_transcript", video=False,
        youtube_identity=False,
        filename="no-source.md",
    )
    slides_only = _talk(
        "lmnopqrstuv", status="skipped_no_transcript", video=False,
        youtube_identity=False,
        filename="slides-only.md",
    )
    slides_only["google_drive_id"] = "drive-artifact"
    path = _write_db(tmp_path, [no_source, slides_only])

    result = _claim(path, limit=10)

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert [item["filename"] for item in payload["claimed"]] == ["slides-only.md"]
    records = {talk["filename"]: talk for talk in _read_db(path)["talks"]}
    assert records["no-source.md"]["status"] == "skipped_no_sources"
    assert records["slides-only.md"]["status"] == "reprocessing-inflight"


def test_video_bearing_download_failure_is_retryable(tmp_path):
    path = _write_db(
        tmp_path,
        [_talk("iPYc7LCH608", status="skipped_download_failed")],
    )

    result = _claim(path)

    assert result.returncode == 0, result.stderr
    claim = json.loads(result.stdout)["claimed"][0]
    assert claim["previous_status"] == "skipped_download_failed"


def test_claim_is_idempotent_for_an_existing_run_and_batch(tmp_path):
    path = _write_db(tmp_path, [_talk("eg6gqvUFh6Q")])
    first = _claim(path)
    first_bytes = path.read_bytes()

    second = _claim(path)

    assert first.returncode == second.returncode == 0
    payload = json.loads(second.stdout)
    assert payload["idempotent_replay"] is True
    assert payload["claimed"][0]["reprocess_generation"] == 1
    assert path.read_bytes() == first_bytes


def test_new_claim_is_v5_with_one_immutable_batch_baseline(tmp_path):
    talks = [_talk("eg6gqvUFh6Q"), _talk("iPYc7LCH608")]
    path = _write_db(tmp_path, talks)

    result = _claim(path, limit=2)

    assert result.returncode == 0, result.stderr
    claims = json.loads(result.stdout)["claimed"]
    assert {claim["schema_version"] for claim in claims} == {5}
    assert {claim["required_return_schema_version"] for claim in claims} == {5}
    assert claims[0]["adherence_baseline"] == claims[1]["adherence_baseline"]
    baseline = claims[0]["adherence_baseline"]
    assert baseline["as_of"] == NOW
    assert baseline["excluded_filenames"] == sorted(
        talk["filename"] for talk in talks)


def test_inspect_dual_reads_a_schema_v3_adherence_claim(tmp_path):
    path = _write_db(tmp_path, [_talk("eg6gqvUFh6Q")])
    claimed = _claim(path)
    assert claimed.returncode == 0, claimed.stderr
    database = _read_db(path)
    claim = database["talks"][0]["_queue_claim"]
    claim["schema_version"] = 3
    claim["required_return_schema_version"] = 3
    baseline = claim["adherence_baseline"]
    baseline["schema_version"] = 1
    for field in (
        "eligible_talk_count",
        "opportunity_coverage_identity",
        "raw_score_comparison_status",
        "raw_score_comparison_reason",
    ):
        baseline.pop(field, None)
    path.write_text(json.dumps(database))
    before = path.read_bytes()

    inspected = _run(path, "inspect", "--run-id", "run-1")

    assert inspected.returncode == 0, inspected.stderr
    assert json.loads(inspected.stdout)["claims"][0]["schema_version"] == 3
    assert path.read_bytes() == before


def test_schema_v5_claim_cannot_request_a_schema_v3_return(tmp_path):
    path = _write_db(tmp_path, [_talk("eg6gqvUFh6Q")])
    claimed = _claim(path)
    assert claimed.returncode == 0, claimed.stderr
    database = _read_db(path)
    database["talks"][0]["_queue_claim"][
        "required_return_schema_version"
    ] = 3
    path.write_text(json.dumps(database))
    before = path.read_bytes()

    inspected = _run(path, "inspect", "--run-id", "run-1")

    assert inspected.returncode == 2
    assert "must equal its claim schema version 5" in inspected.stderr
    assert path.read_bytes() == before


def test_claim_baseline_failure_is_copy_on_write(tmp_path):
    malformed = _talk("eg6gqvUFh6Q", status="processed")
    malformed.update({
        "pattern_scoring_generation_status": "current",
        "pattern_scoring_generation_reasons": [],
        "pattern_score": 1,
        "pattern_observations": {"pattern_score": 1},
    })
    path = _write_db(tmp_path, [malformed, _talk("iPYc7LCH608")])
    before = path.read_bytes()

    result = _claim(
        path,
        filenames=("playlist-iPYc7LCH608.md",),
    )

    assert result.returncode == 2
    assert "missing required identity fields" in result.stderr
    assert path.read_bytes() == before


def test_same_run_and_batch_reclaims_a_stale_recovered_generation(tmp_path):
    path = _write_db(
        tmp_path,
        [_talk("eg6gqvUFh6Q", status="needs-reprocessing")],
    )
    first = _run(
        path,
        "claim",
        "--run-id", "reparse",
        "--batch-id", "25",
        "--now", "2026-07-31T17:00:00+00:00",
    )
    assert first.returncode == 0, first.stderr
    first_claim = json.loads(first.stdout)["claimed"][0]
    recovered = _run(
        path,
        "recover",
        "--now", "2026-07-31T18:00:00+00:00",
        "--stale-after-seconds", "3600",
    )
    assert recovered.returncode == 0, recovered.stderr

    retried = _run(
        path,
        "claim",
        "--run-id", "reparse",
        "--batch-id", "25",
        "--now", "2026-07-31T18:01:00+00:00",
    )

    assert retried.returncode == 0, retried.stderr
    payload = json.loads(retried.stdout)
    assert payload["idempotent_replay"] is False
    assert payload["claimed"][0]["reprocess_generation"] == 2
    assert payload["claimed"][0]["state"] == "claimed"
    talk = _read_db(path)["talks"][0]
    assert talk["status"] == "reprocessing-inflight"
    assert talk["reprocess_generation"] == 2
    assert talk["_queue_claim"]["reprocess_generation"] == 2
    archived = talk["_queue_claim_history"][0]
    assert archived["schema_version"] == 5
    assert archived["adherence_baseline"] == first_claim["adherence_baseline"]
    assert archived["state"] == "stale_recovered"
    assert archived["released_at"] == "2026-07-31T18:00:00+00:00"
    assert archived["release_reason"] == "lease_expired"
    assert talk["_queue_claim"]["adherence_baseline"]["as_of"] == \
        "2026-07-31T18:01:00+00:00"
    assert (talk["_queue_claim"]["adherence_baseline"] !=
            archived["adherence_baseline"])

    retried_bytes = path.read_bytes()
    replay = _run(
        path,
        "claim",
        "--run-id", "reparse",
        "--batch-id", "25",
        "--now", "2026-07-31T18:02:00+00:00",
    )
    assert replay.returncode == 0, replay.stderr
    replay_payload = json.loads(replay.stdout)
    assert replay_payload["idempotent_replay"] is True
    assert len(replay_payload["claimed"]) == 1
    assert replay_payload["claimed"][0]["reprocess_generation"] == 2
    assert path.read_bytes() == retried_bytes


def test_v4_batch_epoch_can_span_current_and_history_after_member_reclaim(
        tmp_path):
    path = _write_db(
        tmp_path,
        [_talk("eg6gqvUFh6Q"), _talk("iPYc7LCH608")],
    )
    claimed = _claim(path, run_id="old-run", batch_id="old-batch", limit=2)
    assert claimed.returncode == 0, claimed.stderr
    database = _read_db(path)
    for talk in database["talks"]:
        talk["status"] = "processed"
        talk["_queue_claim"].update({
            "state": "completed",
            "released_at": "2026-07-31T18:05:00+00:00",
            "release_reason": "return_persisted",
            "result_status": "processed",
            "result_payload_sha256": "0" * 64,
        })

    reclaimed = database["talks"][0]
    old_claim = copy.deepcopy(reclaimed["_queue_claim"])
    reclaimed["_queue_claim_history"] = [old_claim]
    new_claim = copy.deepcopy(old_claim)
    for field in (
            "released_at", "release_reason", "result_status",
            "result_payload_sha256"):
        new_claim.pop(field)
    new_claim.update({
        "run_id": "new-run",
        "batch_id": "new-batch",
        "claimed_at": "2026-07-31T19:00:00+00:00",
        "reprocess_generation": 2,
        "state": "claimed",
    })
    new_claim["adherence_baseline"]["as_of"] = new_claim["claimed_at"]
    new_claim["adherence_baseline"]["excluded_filenames"] = [
        reclaimed["filename"]]
    reclaimed["_queue_claim"] = new_claim
    reclaimed["reprocess_generation"] = 2
    reclaimed["status"] = "reprocessing-inflight"

    path.write_text(json.dumps(database))
    result = _run(path, "inspect", "--run-id", "new-run")

    assert result.returncode == 0, result.stderr


def test_v4_current_batch_cannot_split_claimed_at(tmp_path):
    path = _write_db(
        tmp_path,
        [_talk("eg6gqvUFh6Q"), _talk("iPYc7LCH608")],
    )
    claimed = _claim(path, run_id="run", batch_id="batch", limit=2)
    assert claimed.returncode == 0, claimed.stderr
    database = _read_db(path)
    second_claim = database["talks"][1]["_queue_claim"]
    second_claim["claimed_at"] = "2026-07-31T18:00:01+00:00"
    second_claim["adherence_baseline"]["as_of"] = second_claim["claimed_at"]

    path.write_text(json.dumps(database))
    before = path.read_bytes()

    result = _run(path, "inspect", "--run-id", "run")

    assert result.returncode == 2
    assert "do not share one claimed_at" in result.stderr
    assert path.read_bytes() == before


def test_claim_is_idempotent_for_a_completed_same_run_and_batch(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="processed")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 2,
        "run_id": "reparse",
        "batch_id": "25",
        "claimed_at": "2026-07-31T17:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "completed",
        "released_at": "2026-07-31T17:30:00+00:00",
        "release_reason": "return_persisted",
        "result_status": "processed",
        "result_payload_sha256": "0" * 64,
    }
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    replay = _claim(
        path,
        run_id="reparse",
        batch_id="25",
        filenames=(talk["filename"],),
    )

    assert replay.returncode == 0, replay.stderr
    payload = json.loads(replay.stdout)
    assert payload["idempotent_replay"] is True
    assert payload["claimed"][0]["state"] == "completed"
    assert payload["claimed"][0]["result_status"] == "processed"
    assert path.read_bytes() == before


def test_completed_v1_replay_keeps_stdout_and_disk_at_v1(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="processed")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 1,
        "run_id": "legacy-run",
        "batch_id": "legacy-batch",
        "claimed_at": "2026-07-31T17:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "completed",
        "released_at": "2026-07-31T17:30:00+00:00",
        "release_reason": "return_persisted",
        "result_status": "processed",
    }
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    result = _claim(
        path,
        run_id="legacy-run",
        batch_id="legacy-batch",
        filenames=(talk["filename"],),
    )

    assert result.returncode == 0, result.stderr
    claim = json.loads(result.stdout)["claimed"][0]
    assert claim["schema_version"] == 1
    assert "result_payload_sha256" not in claim
    assert path.read_bytes() == before


def test_recover_uses_injected_time_and_exact_stale_threshold(tmp_path):
    path = _write_db(
        tmp_path,
        [_talk("eg6gqvUFh6Q", status="needs-reprocessing")],
    )
    claimed = _run(
        path,
        "claim",
        "--run-id", "reparse",
        "--batch-id", "25",
        "--now", "2026-07-31T17:00:00+00:00",
    )
    assert claimed.returncode == 0, claimed.stderr

    fresh = _run(
        path,
        "recover",
        "--now", "2026-07-31T17:59:59+00:00",
        "--stale-after-seconds", "3600",
    )
    assert fresh.returncode == 0
    assert json.loads(fresh.stdout)["recovered"] == []

    stale = _run(
        path,
        "recover",
        "--now", "2026-07-31T18:00:00+00:00",
        "--stale-after-seconds", "3600",
    )
    assert stale.returncode == 0, stale.stderr
    assert json.loads(stale.stdout)["recovered"] == [{
        "filename": "playlist-eg6gqvUFh6Q.md",
        "run_id": "reparse",
        "batch_id": "25",
        "reprocess_generation": 1,
        "status": "needs-reprocessing",
        "age_seconds": 3600,
    }]
    talk = _read_db(path)["talks"][0]
    assert talk["status"] == "needs-reprocessing"
    assert talk["_queue_claim"]["state"] == "stale_recovered"
    assert talk["reprocess_generation"] == 1


def test_inspect_reconstructs_claims_after_stale_recovery(tmp_path):
    path = _write_db(tmp_path, [_talk("eg6gqvUFh6Q")])
    assert _claim(path, run_id="reparse", batch_id="25").returncode == 0
    assert _run(
        path,
        "recover",
        "--now", "2026-07-31T19:00:00+00:00",
        "--stale-after-seconds", "1",
    ).returncode == 0

    result = _run(path, "inspect", "--run-id", "reparse")

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["batches"] == [{
        "batch_id": "25",
        "filenames": ["playlist-eg6gqvUFh6Q.md"],
    }]
    assert payload["claims"][0]["state"] == "stale_recovered"


def test_legacy_database_inspect_is_read_only(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="processed")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 1,
        "run_id": "legacy-run",
        "batch_id": "legacy-batch",
        "claimed_at": "2026-07-31T17:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "completed",
        "released_at": "2026-07-31T17:30:00+00:00",
        "release_reason": "return_persisted",
        "result_status": "processed",
    }
    path = _write_db(tmp_path, [talk], current=False)
    before = path.read_bytes()

    result = _run(path, "inspect", "--run-id", "legacy-run")

    assert result.returncode == 0, result.stderr
    assert json.loads(result.stdout)["claims"][0]["state"] == "completed"
    assert path.read_bytes() == before


def test_legacy_database_recover_closes_lease_without_migrating(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="reprocessing-inflight")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 1,
        "run_id": "legacy-run",
        "batch_id": "legacy-batch",
        "claimed_at": "2026-07-31T17:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "claimed",
    }
    path = _write_db(tmp_path, [talk], current=False)

    result = _run(
        path,
        "recover",
        "--now",
        "2026-07-31T18:00:00+00:00",
        "--stale-after-seconds",
        "3600",
    )

    assert result.returncode == 0, result.stderr
    database = _read_db(path)
    assert "schema_version" not in database
    assert "schema_version" not in database["config"]
    assert "schema_version" not in database["talks"][0]
    assert database["talks"][0]["status"] == "needs-reprocessing"
    assert database["talks"][0]["_queue_claim"]["schema_version"] == 2
    assert database["talks"][0]["_queue_claim"]["state"] == "stale_recovered"


@pytest.mark.parametrize("action", ["normalize", "claim"])
def test_legacy_database_rejects_new_queue_mutations(tmp_path, action):
    path = _write_db(tmp_path, [_talk("eg6gqvUFh6Q")], current=False)
    before = path.read_bytes()
    arguments = (
        ("normalize",)
        if action == "normalize"
        else (
            "claim",
            "--run-id",
            "new-run",
            "--batch-id",
            "1",
            "--now",
            NOW,
        )
    )

    result = _run(path, *arguments)

    assert result.returncode == 2
    assert "migrate-tracking-database.py" in json.loads(result.stdout)["error"]
    assert path.read_bytes() == before


def test_inspect_accepts_a_completed_persistence_claim(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="processed")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 2,
        "run_id": "reparse",
        "batch_id": "25",
        "claimed_at": "2026-07-31T17:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "completed",
        "released_at": "2026-07-31T17:30:00+00:00",
        "release_reason": "return_persisted",
        "result_status": "processed",
        "result_payload_sha256": "0" * 64,
    }
    path = _write_db(tmp_path, [talk])

    result = _run(path, "inspect", "--run-id", "reparse")

    assert result.returncode == 0, result.stderr
    claim = json.loads(result.stdout)["claims"][0]
    assert claim["state"] == "completed"
    assert claim["result_status"] == "processed"


def test_inspect_dual_reads_legacy_completed_claim_without_mutating(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="processed")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 1,
        "run_id": "legacy-run",
        "batch_id": "legacy-batch",
        "claimed_at": "2026-07-31T17:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "completed",
        "released_at": "2026-07-31T17:30:00+00:00",
        "release_reason": "return_persisted",
        "result_status": "processed",
    }
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    result = _run(path, "inspect", "--run-id", "legacy-run")

    assert result.returncode == 0, result.stderr
    claim = json.loads(result.stdout)["claims"][0]
    assert claim["schema_version"] == 1
    assert "result_payload_sha256" not in claim
    assert path.read_bytes() == before


def test_unknown_future_claim_schema_fails_without_rewriting(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="reprocessing-inflight")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 99,
        "run_id": "future-run",
        "batch_id": "future-batch",
        "claimed_at": NOW,
        "previous_status": "pending",
        "reprocess_generation": 1,
        "state": "claimed",
    }
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    result = _run(path, "inspect", "--run-id", "future-run")

    assert result.returncode == 2
    assert (
        "queue_claim_schema_version_unsupported"
        in json.loads(result.stdout)["error"]
    )
    assert path.read_bytes() == before


def test_active_claim_with_terminal_status_rejects_every_command_but_recover(
        tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="processed")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 1,
        "run_id": "reparse",
        "batch_id": "25",
        "claimed_at": NOW,
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "claimed",
    }
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    rejected = _run(path, "inspect", "--run-id", "reparse")

    assert rejected.returncode == 2
    assert "stranded lease" in json.loads(rejected.stdout)["error"]
    assert path.read_bytes() == before

    recovered = _run(
        path,
        "recover",
        "--now", NOW,
        "--stale-after-seconds", "999",
    )

    assert recovered.returncode == 0, recovered.stderr
    assert json.loads(recovered.stdout)["recovered"] == [{
        "filename": "playlist-eg6gqvUFh6Q.md",
        "run_id": "reparse",
        "batch_id": "25",
        "reprocess_generation": 1,
        "status": "needs-reprocessing",
        "age_seconds": 0,
        "status_before": "processed",
        "release_reason": "state_status_drift",
    }]
    repaired = _read_db(path)["talks"][0]
    assert repaired["status"] == "needs-reprocessing"
    assert repaired["_queue_claim"]["state"] == "stale_recovered"
    assert repaired["_queue_claim"]["release_reason"] == "state_status_drift"


def test_duplicate_filenames_reject_without_rewriting(tmp_path):
    talk = _talk("eg6gqvUFh6Q")
    path = _write_db(tmp_path, [talk, dict(talk)])
    before = path.read_bytes()

    result = _run(path, "normalize")

    assert result.returncode == 2
    assert "duplicate talk filename" in json.loads(result.stdout)["error"]
    assert path.read_bytes() == before


def test_tracking_database_symlink_is_rejected_before_read_or_write(tmp_path):
    target = _write_db(tmp_path, [_talk("eg6gqvUFh6Q")])
    before = target.read_bytes()
    link = tmp_path / "tracking-link.json"
    link.symlink_to(target.name)

    result = _claim(link)

    assert result.returncode == 2
    assert "symbolic link" in json.loads(result.stdout)["error"]
    assert link.is_symlink()
    assert target.read_bytes() == before


@pytest.mark.parametrize(
    "talk,error",
    [
        (
            _talk(
                "eg6gqvUFh6Q",
                filename="playlist-iPYc7LCH608.md",
            ),
            "filename id",
        ),
        (
            {
                **_talk("eg6gqvUFh6Q", filename="talk.md"),
                "video_url": "https://youtu.be/iPYc7LCH608",
            },
            "disagrees with video_url id",
        ),
    ],
)
def test_locally_decidable_video_identity_mismatches_reject(tmp_path, talk, error):
    path = _write_db(tmp_path, [talk])

    result = _run(path, "normalize")

    assert result.returncode == 2
    assert error in json.loads(result.stdout)["error"]


def test_malformed_claim_timestamp_rejects(tmp_path):
    talk = _talk("eg6gqvUFh6Q", status="reprocessing-inflight")
    talk["reprocess_generation"] = 1
    talk["_queue_claim"] = {
        "schema_version": 1,
        "run_id": "reparse",
        "batch_id": "25",
        "claimed_at": "yesterday",
        "previous_status": "pending",
        "reprocess_generation": 1,
        "state": "claimed",
    }
    path = _write_db(tmp_path, [talk])

    result = _run(path, "inspect", "--run-id", "reparse")

    assert result.returncode == 2
    assert "claim.claimed_at" in json.loads(result.stdout)["error"]


@pytest.mark.parametrize(
    "talk",
    [
        _talk("eg6gqvUFh6Q", status="skipped_duplicate"),
        _talk(
            "iPYc7LCH608", status="pending", video=False,
            youtube_identity=False, filename="catalog-pending.md",
        ),
    ],
)
def test_explicit_claim_rejects_invalid_transitions(tmp_path, talk):
    path = _write_db(tmp_path, [talk])
    before = path.read_bytes()

    result = _claim(path, filenames=(talk["filename"],))

    assert result.returncode == 2
    assert "cannot claim" in json.loads(result.stdout)["error"]
    assert path.read_bytes() == before


def test_cli_requires_timezone_aware_now(tmp_path):
    path = _write_db(tmp_path, [_talk("eg6gqvUFh6Q")])

    result = _run(
        path,
        "claim",
        "--run-id", "reparse",
        "--batch-id", "25",
        "--now", "2026-07-31T18:00:00",
    )

    assert result.returncode == 2
    payload = json.loads(result.stdout)
    assert payload["ok"] is False
    assert "has no timezone" in payload["error"]
