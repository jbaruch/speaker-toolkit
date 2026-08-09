"""Tests for pptx-extraction.py — PPTX visual data extraction."""

import io
import importlib
import json
import os
import shlex
import shutil
import struct
import subprocess
import sys
import zipfile
from dataclasses import replace
from pathlib import Path
from types import SimpleNamespace

import pytest
from pypdf import PdfWriter
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Emu, Inches, Pt

from conftest import (
    deck_height,
    deck_width,
    graphic_frame_element,
    make_deck,
    slide_title,
)


def _use_in_process_directory_discovery(pptx_extraction, monkeypatch):
    """Keep batch unit tests fast while production discovery stays supervised."""

    def discover(directory, patterns, directory_exclusions, *, deadline):
        assert deadline > 0
        files, skipped, _started = pptx_extraction._discover_pptx_files(
            directory,
            patterns,
            directory_exclusions,
        )
        reasons = pptx_extraction.directory_incomplete_reason_codes(skipped)
        return [relative for _path, relative in files], skipped, not reasons, reasons

    monkeypatch.setattr(
        pptx_extraction,
        "_run_supervised_directory_discovery",
        discover,
    )


def _directory_manifest(*, files=None, skipped=None, **updates):
    skipped = [] if skipped is None else skipped
    policy_skips = {
        "pptx_batch_conflict_copy",
        "pptx_batch_directory_excluded",
        "pptx_batch_office_lock_file",
        "pptx_batch_reparse_point_rejected",
        "pptx_batch_skip_pattern",
        "pptx_batch_static_export",
        "pptx_batch_symlink_rejected",
    }
    reasons = sorted(
        {
            item["reason"]
            for item in skipped
            if isinstance(item.get("reason"), str)
            and item["reason"] not in policy_skips
        }
    )
    manifest = {
        "schema_version": 2,
        "kind": "directory",
        "complete": not reasons,
        "directory_exclusions": [],
        "incomplete_reason_codes": reasons,
        "files": [] if files is None else files,
        "skipped": skipped,
    }
    manifest.update(updates)
    return manifest


def test_slide_count(pptx_extraction, tmp_path):
    prs = make_deck(5)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    result = pptx_extraction._extract_pptx_in_process(path)
    assert result["slide_count"] == 5


def test_public_extractor_routes_only_through_supervisor(
    pptx_extraction,
    monkeypatch,
):
    calls = []

    def supervised(path, **options):
        calls.append((path, options))
        return {"slide_count": 7}

    monkeypatch.setattr(
        pptx_extraction,
        "run_supervised_pptx_extraction",
        supervised,
    )
    result = pptx_extraction.extract_pptx(
        "deck.pptx",
        ocr=False,
        rendered_pdf_path="deck.pdf",
        inspected_page_ranges=[[1, 2]],
    )

    assert result == {"slide_count": 7}
    assert calls == [
        (
            "deck.pptx",
            {
                "ocr": False,
                "rendered_pdf_path": "deck.pdf",
                "inspected_page_ranges": [[1, 2]],
            },
        )
    ]


def test_public_extractor_rejects_in_process_callable(
    pptx_extraction,
):
    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction.extract_pptx(
            "deck.pptx",
            ocr_fn=lambda _blobs: "not serializable",
        )

    assert caught.value.reason_code == "pptx_evidence_invalid"


def test_contained_extraction_never_nests_the_pdf_supervisor(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    deck_path = tmp_path / "deck.pptx"
    make_deck(1).save(str(deck_path))
    rendered_path = tmp_path / "deck.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=640, height=480)
    with rendered_path.open("wb") as stream:
        writer.write(stream)
    pdf_evidence = importlib.import_module("pdf_evidence")
    monkeypatch.setattr(
        pdf_evidence,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("nested PDF supervisor started"),
    )

    result = pptx_extraction._extract_pptx_in_process(
        deck_path,
        ocr=False,
        rendered_pdf_path=rendered_path,
        inspected_page_ranges=[[1, 1]],
    )

    receipt = result["native_deck_audit"]["rendered_page_inspection"]
    assert receipt["rendered_page_count"] == 1
    assert receipt["inspected_page_ranges"] == [[1, 1]]


def test_contained_extraction_rejects_rendered_pdf_repair_diagnostics(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    deck_path = tmp_path / "deck.pptx"
    make_deck(1).save(str(deck_path))
    rendered_path = tmp_path / "deck.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=640, height=480)
    with rendered_path.open("wb") as stream:
        writer.write(stream)
    pdf_evidence = importlib.import_module("pdf_evidence")
    pptx_evidence = importlib.import_module("pptx_evidence")

    def repair_required(*_args, **_kwargs):
        raise pdf_evidence.PdfEvidenceError(
            "synthetic repair diagnostics",
            reason_code="pdf_parser_repair_required",
            details={
                "diagnostic_receipt": {
                    "byte_count": 1,
                    "sha256": "f" * 64,
                    "truncated": False,
                }
            },
        )

    monkeypatch.setattr(
        pptx_evidence,
        "_inspect_pdf_in_contained_worker",
        repair_required,
    )

    with pytest.raises(pptx_evidence.PptxEvidenceError) as caught:
        pptx_extraction._extract_pptx_in_process(
            deck_path,
            ocr=False,
            rendered_pdf_path=rendered_path,
            inspected_page_ranges=[[1, 1]],
        )

    assert caught.value.reason_code == "pdf_parser_repair_required"
    assert caught.value.details["diagnostic_receipt"]["byte_count"] == 1


def test_slide_dimensions(pptx_extraction, tmp_path):
    prs = make_deck(1)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    result = pptx_extraction._extract_pptx_in_process(path)
    # Default slide dimensions should be reasonable
    assert result["slide_width_inches"] > 0
    assert result["slide_height_inches"] > 0


def test_shape_text_extraction(pptx_extraction, tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide_title(slide).text = "Hello World"
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    result = pptx_extraction._extract_pptx_in_process(path)
    slide_data = result["per_slide_visual"][0]
    assert "Hello World" in slide_data["text_content_preview"]


def test_font_tracking(pptx_extraction, tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Test text"
    run.font.name = "Arial"
    run.font.size = Pt(24)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    result = pptx_extraction._extract_pptx_in_process(path)
    assert "Arial" in result["global_design"]["fonts_used"]


def test_skip_static(pptx_extraction):
    skip, reason = pptx_extraction.should_skip("presentation-static.pptx", [])
    assert skip is True
    assert "static" in reason


def test_skip_conflict_copy(pptx_extraction):
    skip, reason = pptx_extraction.should_skip("deck (1).pptx", [])
    assert skip is True
    assert "conflict" in reason


@pytest.mark.parametrize(
    "reason_code",
    [
        "pptx_artifact_changed",
        "pptx_dependency_unavailable",
        "pptx_probe_containment_unavailable",
        "pptx_probe_crash",
        "pptx_probe_malformed_result",
        "pptx_probe_monitor_identity_changed",
        "pptx_probe_monitor_unavailable",
        "pptx_probe_request_oversized",
        "pptx_probe_resource_unavailable",
        "pptx_probe_result_oversized",
        "pptx_probe_start_failure",
        "pptx_probe_timeout",
    ],
)
def test_batch_preserves_current_and_legacy_supervisor_reasons(
    pptx_extraction,
    reason_code: str,
) -> None:
    assert (
        pptx_extraction._batch_error_reason(SimpleNamespace(reason_code=reason_code))
        == reason_code
    )


def test_skip_custom_pattern(pptx_extraction):
    skip, reason = pptx_extraction.should_skip("my-template.pptx", ["template"])
    assert skip is True


def test_no_skip_normal_file(pptx_extraction):
    skip, _ = pptx_extraction.should_skip("great-talk.pptx", [])
    assert skip is False


def test_skip_office_lock_file(pptx_extraction):
    skip, reason = pptx_extraction.should_skip("~$great-talk.pptx", [])
    assert skip is True
    assert reason == "Office lock file"


def test_bounded_discovery_rejects_symlinks_and_is_deterministic(
    pptx_extraction,
    tmp_path,
):
    root = tmp_path / "decks"
    nested = root / "nested"
    nested.mkdir(parents=True)
    (root / "z.pptx").write_bytes(b"z")
    (nested / "a.pptx").write_bytes(b"a")
    outside = tmp_path / "outside.pptx"
    outside.write_bytes(b"outside")
    (root / "linked.pptx").symlink_to(outside)
    (root / "linked-directory").symlink_to(tmp_path, target_is_directory=True)

    first, first_skipped, _started = pptx_extraction._discover_pptx_files(root, [])
    second, second_skipped, _started = pptx_extraction._discover_pptx_files(root, [])

    assert [item[1] for item in first] == ["z.pptx", "nested/a.pptx"]
    assert [item[1] for item in second] == ["z.pptx", "nested/a.pptx"]
    assert (
        first_skipped
        == second_skipped
        == [
            {
                "path": "linked-directory",
                "reason": "pptx_batch_symlink_rejected",
            },
            {
                "path": "linked.pptx",
                "reason": "pptx_batch_symlink_rejected",
            },
        ]
    )


def test_directory_exclusion_prunes_before_scandir_and_keeps_authored_default_deck(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    dependency = root / ".venv" / "lib" / "site-packages" / "pptx" / "templates"
    authored = root / "templates"
    dependency.mkdir(parents=True)
    authored.mkdir(parents=True)
    (dependency / "default.pptx").write_bytes(b"dependency")
    (authored / "default.pptx").write_bytes(b"authored")
    original_scandir = pptx_extraction.os.scandir

    def guarded_scandir(path):
        assert Path(path) != root / ".venv", "excluded directory was inspected"
        return original_scandir(path)

    monkeypatch.setattr(pptx_extraction.os, "scandir", guarded_scandir)

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(
        root,
        [],
        [".VENV"],
    )

    assert [relative for _path, relative in discovered] == ["templates/default.pptx"]
    assert skipped == [{"path": ".venv", "reason": "pptx_batch_directory_excluded"}]
    assert pptx_extraction.directory_incomplete_reason_codes(skipped) == []


def test_excluded_directory_cannot_starve_authored_sibling_at_entry_cap(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    excluded = root / ".venv"
    excluded.mkdir(parents=True)
    (excluded / "dependency.pptx").write_bytes(b"dependency")
    (root / "talk.pptx").write_bytes(b"authored")
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_ENTRIES", 1)

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(
        root,
        [],
        [".venv"],
    )

    assert [relative for _path, relative in discovered] == ["talk.pptx"]
    assert skipped == [{"path": ".venv", "reason": "pptx_batch_directory_excluded"}]


def test_policy_exclusion_enumeration_has_its_own_closed_cap(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    (root / ".venv").mkdir(parents=True)
    (root / "talk.pptx").write_bytes(b"authored")
    monkeypatch.setattr(
        pptx_extraction,
        "_BATCH_MAX_POLICY_EXCLUDED_ENTRIES",
        0,
    )

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(
        root,
        [],
        [".venv"],
    )

    assert discovered == []
    assert skipped == [{"path": ".", "reason": "pptx_batch_entry_limit"}]


@pytest.mark.skipif(
    sys.platform == "win32",
    reason="POSIX symlink precedence is tested separately from Windows reparse points",
)
def test_named_exclusion_keeps_symlink_rejection_precedence(
    pptx_extraction,
    tmp_path,
):
    root = tmp_path / "decks"
    outside = tmp_path / "outside"
    root.mkdir()
    outside.mkdir()
    (root / ".venv").symlink_to(outside, target_is_directory=True)

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(
        root,
        [],
        [".venv"],
    )

    assert discovered == []
    assert skipped == [{"path": ".venv", "reason": "pptx_batch_symlink_rejected"}]


def test_named_exclusion_keeps_reparse_rejection_precedence(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    excluded = root / ".venv"
    excluded.mkdir(parents=True)
    excluded_inode = excluded.lstat().st_ino
    monkeypatch.setattr(
        pptx_extraction,
        "_is_windows_reparse_point",
        lambda value: value.st_ino == excluded_inode,
    )

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(
        root,
        [],
        [".venv"],
    )

    assert discovered == []
    assert skipped == [{"path": ".venv", "reason": "pptx_batch_reparse_point_rejected"}]


def test_excluded_directories_do_not_consume_directory_budget(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    included = root / "included"
    excluded = root / ".venv"
    included.mkdir(parents=True)
    excluded.mkdir()
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_DIRECTORIES", 2)

    _files, exact_skipped, _started = pptx_extraction._discover_pptx_files(root, [])
    assert exact_skipped == [{"path": ".", "reason": "pptx_batch_directory_limit"}]

    _files, excluded_skipped, _started = pptx_extraction._discover_pptx_files(
        root,
        [],
        [".venv"],
    )
    assert excluded_skipped == [
        {"path": ".venv", "reason": "pptx_batch_directory_excluded"}
    ]
    assert pptx_extraction.directory_incomplete_reason_codes(excluded_skipped) == []


def test_exact_directory_budget_is_complete_but_one_more_directory_is_partial(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    first = root / "first"
    first.mkdir(parents=True)
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_DIRECTORIES", 2)

    _files, skipped, _started = pptx_extraction._discover_pptx_files(root, [])
    assert skipped == []

    (root / "second").mkdir()
    _files, skipped, _started = pptx_extraction._discover_pptx_files(root, [])
    assert skipped == [{"path": ".", "reason": "pptx_batch_directory_limit"}]


def test_depth_budget_marks_only_unvisited_descendants_partial(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    child = root / "child"
    child.mkdir(parents=True)
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_DEPTH", 1)

    _files, skipped, _started = pptx_extraction._discover_pptx_files(root, [])
    assert skipped == []

    (child / "grandchild").mkdir()
    _files, skipped, _started = pptx_extraction._discover_pptx_files(root, [])
    assert skipped == [{"path": "child/grandchild", "reason": "pptx_batch_depth_limit"}]


def test_bounded_discovery_rejects_symlink_root(pptx_extraction, tmp_path):
    root = tmp_path / "decks"
    root.mkdir()
    linked_root = tmp_path / "linked-root"
    linked_root.symlink_to(root, target_is_directory=True)

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction._discover_pptx_files(linked_root, [])

    assert caught.value.reason_code == "pptx_batch_root_invalid"


def test_bounded_discovery_rejects_windows_reparse_root(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    root_inode = root.lstat().st_ino
    monkeypatch.setattr(
        pptx_extraction,
        "_is_windows_reparse_point",
        lambda value: value.st_ino == root_inode,
    )

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction._discover_pptx_files(root, [])

    assert caught.value.reason_code == "pptx_batch_root_invalid"


def test_bounded_discovery_rejects_reparse_directory_and_leaf(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    junction = root / "junction"
    junction.mkdir(parents=True)
    leaf = root / "cloud.pptx"
    leaf.write_bytes(b"cloud")
    rejected_inodes = {junction.lstat().st_ino, leaf.lstat().st_ino}
    monkeypatch.setattr(
        pptx_extraction,
        "_is_windows_reparse_point",
        lambda value: value.st_ino in rejected_inodes,
    )

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(root, [])

    assert discovered == []
    assert skipped == [
        {"path": "cloud.pptx", "reason": "pptx_batch_reparse_point_rejected"},
        {"path": "junction", "reason": "pptx_batch_reparse_point_rejected"},
    ]


def test_windows_leaf_policy_accepts_only_hydrated_supported_cloud_tags(
    pptx_extraction,
):
    hydrated = SimpleNamespace(
        st_file_attributes=pptx_extraction._WINDOWS_REPARSE_POINT_ATTRIBUTE,
        st_reparse_tag=next(iter(pptx_extraction._WINDOWS_CLOUD_REPARSE_TAGS)),
    )
    offline = SimpleNamespace(
        st_file_attributes=(
            pptx_extraction._WINDOWS_REPARSE_POINT_ATTRIBUTE
            | pptx_extraction._WINDOWS_OFFLINE_ATTRIBUTE
        ),
        st_reparse_tag=hydrated.st_reparse_tag,
    )
    recalling = SimpleNamespace(
        st_file_attributes=(
            pptx_extraction._WINDOWS_REPARSE_POINT_ATTRIBUTE
            | pptx_extraction._WINDOWS_RECALL_ON_DATA_ACCESS_ATTRIBUTE
        ),
        st_reparse_tag=hydrated.st_reparse_tag,
    )
    unknown_redirect = SimpleNamespace(
        st_file_attributes=pptx_extraction._WINDOWS_REPARSE_POINT_ATTRIBUTE,
        st_reparse_tag=0xA0000003,
    )

    assert pptx_extraction._windows_leaf_rejection_reason(hydrated) is None
    assert (
        pptx_extraction._windows_leaf_rejection_reason(offline)
        == "pptx_batch_cloud_placeholder_unavailable"
    )
    assert (
        pptx_extraction._windows_leaf_rejection_reason(recalling)
        == "pptx_batch_cloud_placeholder_unavailable"
    )
    assert (
        pptx_extraction._windows_leaf_rejection_reason(unknown_redirect)
        == "pptx_batch_reparse_point_rejected"
    )


@pytest.mark.parametrize(
    "manifest",
    [
        _directory_manifest(unexpected=True),
        _directory_manifest(schema_version=True),
        _directory_manifest(
            files=[{"path": "deck.pptx"}, {"path": "deck.pptx"}],
        ),
        _directory_manifest(files=[{"path": "../deck.pptx"}]),
        _directory_manifest(files=[{"path": "/outside/deck.pptx"}]),
        _directory_manifest(files=[{"path": "not-a-deck.txt"}]),
        _directory_manifest(files=[{"path": "~$locked.pptx"}]),
        _directory_manifest(
            skipped=[{"path": ".", "reason": "pptx_batch_unknown_typo"}],
        ),
        _directory_manifest(skipped=[{"path": ".", "reason": []}]),
        _directory_manifest(
            files=[{"path": "deck.pptx"}],
            skipped=[{"path": "deck.pptx", "reason": "pptx_batch_skip_pattern"}],
        ),
        _directory_manifest(
            skipped=[
                {"path": ".", "reason": "pptx_batch_wall_limit"},
                {"path": ".", "reason": "pptx_batch_wall_limit"},
            ],
        ),
        _directory_manifest(complete=False),
        _directory_manifest(incomplete_reason_codes=["pptx_batch_wall_limit"]),
    ],
)
def test_directory_manifest_rejects_malformed_duplicate_or_noncanonical_data(
    pptx_extraction,
    manifest,
):
    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction._decode_directory_manifest(manifest)

    assert caught.value.reason_code == "pptx_batch_manifest_invalid"


def test_directory_manifest_allows_distinct_root_findings(pptx_extraction):
    manifest = _directory_manifest(
        skipped=[
            {"path": ".", "reason": "pptx_batch_wall_limit"},
            {"path": ".", "reason": "pptx_batch_entry_limit"},
        ],
    )

    assert pptx_extraction._decode_directory_manifest(manifest) == (
        [],
        manifest["skipped"],
        False,
        ["pptx_batch_entry_limit", "pptx_batch_wall_limit"],
    )


@pytest.mark.parametrize(
    ("manifest", "expected_exclusions"),
    [
        (
            _directory_manifest(directory_exclusions=["templates"]),
            [".venv"],
        ),
        (
            _directory_manifest(
                directory_exclusions=[".venv"],
                skipped=[
                    {
                        "path": "templates",
                        "reason": "pptx_batch_directory_excluded",
                    }
                ],
            ),
            [".venv"],
        ),
        (
            _directory_manifest(
                directory_exclusions=[".venv"],
                files=[{"path": ".venv/lib/default.pptx"}],
            ),
            [".venv"],
        ),
        (
            _directory_manifest(
                directory_exclusions=[".venv"],
                skipped=[
                    {
                        "path": ".venv",
                        "reason": "pptx_batch_directory_excluded",
                    },
                    {
                        "path": ".venv/lib",
                        "reason": "pptx_batch_directory_unavailable",
                    },
                ],
            ),
            [".venv"],
        ),
        (
            _directory_manifest(
                skipped=[
                    {
                        "path": "unavailable",
                        "reason": "pptx_batch_directory_unavailable",
                    },
                    {
                        "path": "unavailable/deck.pptx",
                        "reason": "pptx_batch_entry_unavailable",
                    },
                ],
            ),
            [],
        ),
    ],
)
def test_directory_manifest_binds_exclusions_and_rejects_descendant_claims(
    pptx_extraction,
    manifest,
    expected_exclusions,
):
    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction._decode_directory_manifest(
            manifest,
            expected_directory_exclusions=expected_exclusions,
        )

    assert caught.value.reason_code == "pptx_batch_manifest_invalid"


@pytest.mark.skipif(
    sys.platform == "win32",
    reason="backslash is a path separator rather than a legal filename",
)
def test_directory_discovery_deduplicates_collapsed_invalid_path_receipts(
    pptx_extraction,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    (root / "one\\bad.pptx").write_bytes(b"one")
    (root / "two\\bad.pptx").write_bytes(b"two")

    files, skipped, _started = pptx_extraction._discover_pptx_files(root, [])
    manifest = _directory_manifest(
        files=[{"path": relative} for _path, relative in files],
        skipped=skipped,
    )

    assert files == []
    assert skipped == [{"path": ".", "reason": "pptx_batch_path_invalid"}]
    assert pptx_extraction._decode_directory_manifest(manifest) == (
        [],
        skipped,
        False,
        ["pptx_batch_path_invalid"],
    )


def test_directory_owner_never_touches_root_before_authenticated_manifest(
    pptx_extraction,
    monkeypatch,
    capsys,
):
    root = "/untrusted/root-that-need-not-exist"
    calls = []

    def forbidden(*_args, **_kwargs):
        pytest.fail("directory owner touched the filesystem")

    for name in ("lstat", "stat", "is_dir", "is_file", "resolve"):
        monkeypatch.setattr(pptx_extraction.Path, name, forbidden)
    monkeypatch.setattr(pptx_extraction.os, "scandir", forbidden)
    monkeypatch.setattr(pptx_extraction.os, "lstat", forbidden)

    def authenticated_worker(
        command, operation, generations, payload, limits, **kwargs
    ):
        calls.append((command, operation, generations, payload, limits, kwargs))
        return SimpleNamespace(
            payload=_directory_manifest(files=[{"path": "nested/deck.pptx"}])
        )

    monkeypatch.setattr(
        pptx_extraction, "run_authenticated_worker", authenticated_worker
    )
    monkeypatch.setattr(
        pptx_extraction,
        "extract_pptx",
        lambda path, **_options: {
            "pptx_path": str(path),
            "slide_count": 1,
            "input_fingerprint": {"size_bytes": 1},
        },
    )

    assert pptx_extraction.main(["--directory", root, "--no-ocr"]) == 0

    captured = capsys.readouterr()
    assert captured.err == ""
    assert json.loads(captured.out)["results"][0]["pptx_path"] == "nested/deck.pptx"
    command, operation, generations, payload, limits, kwargs = calls[0]
    assert command == [
        sys.executable,
        pptx_extraction.os.path.abspath(pptx_extraction.__file__),
        pptx_extraction._DIRECTORY_WORKER_FLAG,
    ]
    assert operation == "pptx_resolve_input"
    assert generations == {}
    assert payload == {
        "root_path": root,
        "skip_patterns": [],
        "directory_exclusions": [],
    }
    assert limits.profile_id == "pptx-directory-discovery-v2"
    assert kwargs["schema_generation"] == pptx_extraction.SCHEMA_VERSION
    assert kwargs["pipeline_generation"] == pptx_extraction.PIPELINE_VERSION
    assert kwargs["immutable_process_identity"] == command[:2]


@pytest.mark.skipif(
    os.name != "posix",
    reason="real nested-runtime fixture uses an executable POSIX shim",
)
def test_nested_runtime_reaches_every_fixed_pptx_worker(
    pptx_extraction,
    monkeypatch,
    request,
    tmp_path,
):
    root = tmp_path / "Presentations"
    source_scripts = Path(pptx_extraction.__file__).parent
    worker_scripts = root / "toolkit" / "vault-ingress" / "scripts"
    shutil.copytree(
        source_scripts,
        worker_scripts,
        ignore=shutil.ignore_patterns("__pycache__", "*.pyc"),
    )
    runtime = root / ".venv" / "bin" / "python3"
    runtime.parent.mkdir(parents=True)
    runtime.write_text(
        f'#!/bin/sh\nexec {shlex.quote(sys.executable)} "$@"\n',
        encoding="utf-8",
    )
    runtime.chmod(0o700)
    deck = root / "Conference" / "deck.pptx"
    deck.parent.mkdir()
    make_deck(1).save(deck)

    pptx_evidence = importlib.import_module("pptx_evidence")
    pptx_evidence.clear_pptx_artifact_probe_cache()
    request.addfinalizer(pptx_evidence.clear_pptx_artifact_probe_cache)
    operations = []
    authenticated_worker = pptx_evidence.run_authenticated_worker

    def record_operation(command, operation, *args, **kwargs):
        operations.append(operation)
        return authenticated_worker(command, operation, *args, **kwargs)

    monkeypatch.setattr(
        pptx_evidence,
        "run_authenticated_worker",
        record_operation,
    )
    monkeypatch.setattr(sys, "executable", str(runtime))
    monkeypatch.setattr(
        pptx_extraction,
        "__file__",
        str(worker_scripts / "pptx-extraction.py"),
    )
    monkeypatch.setattr(
        pptx_evidence,
        "__file__",
        str(worker_scripts / "pptx_evidence.py"),
    )

    batch = pptx_extraction.batch_extract(root, [], ocr=False)
    results = batch["results"]
    skipped = batch["skipped"]
    batch_operations = set(operations)
    operation_count = len(operations)
    probe = pptx_evidence.probe_pptx_artifact(deck, trusted_root=root)
    probe_operations = set(operations[operation_count:])
    operation_count = len(operations)
    audit = pptx_evidence.recompute_native_deck_audit(deck, trusted_root=root)
    audit_operations = set(operations[operation_count:])

    assert batch["schema_version"] == 1
    assert batch["kind"] == "pptx_directory_batch"
    assert batch["complete"] is True
    assert batch["incomplete_reason_codes"] == []
    assert skipped == []
    assert len(results) == 1
    assert results[0]["pptx_path"] == "Conference/deck.pptx"
    assert results[0]["slide_count"] == 1
    assert probe.source_sha256 == results[0]["input_fingerprint"]["digest"]
    assert audit["source_pptx_sha256"] == probe.source_sha256
    assert batch_operations >= {
        pptx_evidence.PPTX_METADATA_OPERATION,
        pptx_evidence.PPTX_EXTRACT_OPERATION,
    }
    assert probe_operations >= {
        pptx_evidence.PPTX_METADATA_OPERATION,
        pptx_evidence.PPTX_PROBE_OPERATION,
    }
    assert audit_operations >= {
        pptx_evidence.PPTX_METADATA_OPERATION,
        pptx_evidence.PPTX_NATIVE_AUDIT_OPERATION,
    }


def test_directory_cli_passes_only_the_exact_configured_skip_patterns(
    pptx_extraction,
    monkeypatch,
    capsys,
):
    root = "/untrusted/root-that-need-not-exist"
    captured_payload = {}

    def authenticated_worker(
        _command, _operation, _generations, payload, _limits, **_kwargs
    ):
        captured_payload.update(payload)
        return SimpleNamespace(
            payload=_directory_manifest(directory_exclusions=[".venv"])
        )

    monkeypatch.setattr(
        pptx_extraction, "run_authenticated_worker", authenticated_worker
    )

    assert (
        pptx_extraction.main(
            [
                "--directory",
                root,
                "--skip=draft",
                "--skip=-speaker-master",
                "--exclude-directory=.venv",
                "--no-ocr",
            ]
        )
        == 0
    )

    captured = capsys.readouterr()
    assert captured.err == ""
    assert json.loads(captured.out) == {
        "schema_version": 1,
        "kind": "pptx_directory_batch",
        "complete": True,
        "incomplete_reason_codes": [],
        "results": [],
        "skipped": [],
    }
    assert captured_payload == {
        "root_path": root,
        "skip_patterns": ["draft", "-speaker-master"],
        "directory_exclusions": [".venv"],
    }


def test_directory_relative_root_is_rejected_before_supervision(
    pptx_extraction,
    monkeypatch,
):
    monkeypatch.setattr(
        pptx_extraction,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("relative root started worker"),
    )

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction._run_supervised_directory_discovery(
            ".",
            [],
            deadline=pptx_extraction.time.monotonic() + 30,
        )

    assert caught.value.reason_code == "pptx_batch_root_invalid"
    assert caught.value.details == {"locator_failure": "artifact_locator_dot_segment"}


def test_directory_foreign_home_and_noncanonical_roots_never_launch_worker(
    pptx_extraction,
    monkeypatch,
):
    foreign_roots = (
        ["/foreign/decks"]
        if pptx_extraction.os.name == "nt"
        else [r"C:\conference\decks", r"\\server\share\decks"]
    )
    cases = [
        ("decks", "artifact_root_not_native_absolute"),
        ("~/decks", "artifact_locator_home_expansion_unsupported"),
        (r"conference\decks", "artifact_locator_noncanonical_relative"),
        *[(foreign, "artifact_locator_foreign_absolute") for foreign in foreign_roots],
    ]
    monkeypatch.setattr(
        pptx_extraction,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("invalid root started worker"),
    )

    for locator, expected_failure in cases:
        with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
            pptx_extraction._run_supervised_directory_discovery(
                locator,
                [],
                deadline=pptx_extraction.time.monotonic() + 30,
            )
        assert caught.value.reason_code == "pptx_batch_root_invalid"
        assert caught.value.details == {"locator_failure": expected_failure}
        assert locator not in str(caught.value)
        assert locator not in repr(caught.value.details)


def test_directory_worker_revalidates_root_before_discovery(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    request = pptx_extraction.WorkerRequest(
        request_id="a" * 64,
        operation=pptx_extraction._DIRECTORY_OPERATION,
        request_sha256="b" * 64,
        limit_profile_id=pptx_extraction._DIRECTORY_LIMITS.profile_id,
        schema_generation=pptx_extraction.SCHEMA_VERSION,
        pipeline_generation=pptx_extraction.PIPELINE_VERSION,
        expected_generations={},
        payload={
            "root_path": "relative",
            "skip_patterns": [],
            "directory_exclusions": [],
        },
        key=b"k" * 32,
    )
    monkeypatch.setattr(
        pptx_extraction,
        "_discover_pptx_files",
        lambda *_args, **_kwargs: pytest.fail("invalid root reached discovery"),
    )

    with pytest.raises(pptx_extraction.SupervisorError) as caught:
        pptx_extraction._dispatch_directory_worker(request)

    assert caught.value.reason_code == "pptx_batch_root_invalid"
    assert caught.value.details == {
        "locator_failure": "artifact_root_not_native_absolute"
    }

    observed = {}

    def discover(root, patterns, exclusions):
        observed["root"] = root
        observed["patterns"] = patterns
        observed["exclusions"] = exclusions
        return [], [], 0.0

    monkeypatch.setattr(pptx_extraction, "_discover_pptx_files", discover)
    payload = pptx_extraction._dispatch_directory_worker(
        replace(
            request,
            payload={
                "root_path": str(tmp_path),
                "skip_patterns": [],
                "directory_exclusions": [".venv"],
            },
        )
    )

    assert observed == {
        "root": str(tmp_path),
        "patterns": [],
        "exclusions": [".venv"],
    }
    assert payload == _directory_manifest(directory_exclusions=[".venv"])


def test_batch_reuses_native_absolute_root_if_working_directory_changes(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    original = tmp_path / "original"
    root = original / "decks"
    other = tmp_path / "other"
    root.mkdir(parents=True)
    other.mkdir()
    monkeypatch.chdir(original)

    def discover(directory, _patterns, _exclusions, *, deadline):
        assert directory == root
        assert deadline > 0
        monkeypatch.chdir(other)
        return ["deck.pptx"], [], True, []

    observed = []

    def extract(path, **_options):
        observed.append(path)
        return {
            "pptx_path": str(path),
            "slide_count": 1,
            "input_fingerprint": {"size_bytes": 1},
        }

    monkeypatch.setattr(
        pptx_extraction,
        "_run_supervised_directory_discovery",
        discover,
    )
    monkeypatch.setattr(pptx_extraction, "extract_pptx", extract)

    batch = pptx_extraction.batch_extract(root, [], ocr=False)

    assert batch["complete"] is True
    assert batch["skipped"] == []
    assert observed == [root / "deck.pptx"]
    assert batch["results"][0]["pptx_path"] == "deck.pptx"


def test_directory_skip_pattern_iterators_are_rejected_without_consumption(
    pptx_extraction,
    monkeypatch,
):
    class ExplodingIterator:
        def __iter__(self):
            raise AssertionError("owner consumed an unbounded iterator")

    monkeypatch.setattr(
        pptx_extraction,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("invalid request started worker"),
    )

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction._run_supervised_directory_discovery(
            "/lexical/root",
            ExplodingIterator(),
            deadline=pptx_extraction.time.monotonic() + 30,
        )

    assert caught.value.reason_code == "pptx_batch_request_invalid"


@pytest.mark.parametrize(
    "exclusions",
    [
        ["nested/.venv"],
        ["venv*"],
        ["venv", "VENV"],
        (item for item in [".venv"]),
    ],
)
def test_directory_exclusions_are_rejected_before_worker_launch(
    pptx_extraction,
    monkeypatch,
    exclusions,
):
    monkeypatch.setattr(
        pptx_extraction,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: pytest.fail("invalid request started worker"),
    )

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction._run_supervised_directory_discovery(
            "/lexical/root",
            [],
            exclusions,
            deadline=pptx_extraction.time.monotonic() + 30,
        )

    assert caught.value.reason_code == "pptx_batch_request_invalid"


def test_directory_discovery_timeout_blocks_every_extraction(
    pptx_extraction,
    monkeypatch,
):
    monkeypatch.setattr(
        pptx_extraction,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(
            pptx_extraction.SupervisorError("worker_timeout")
        ),
    )
    monkeypatch.setattr(
        pptx_extraction,
        "extract_pptx",
        lambda *_args, **_kwargs: pytest.fail("timeout launched an extractor"),
    )

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction.batch_extract("/blocked/discovery", [], ocr=False)

    assert caught.value.reason_code == "pptx_batch_discovery_timeout"
    assert caught.value.details == {"supervisor_reason_code": "worker_timeout"}


@pytest.mark.parametrize(
    ("supervisor_reason", "expected_reason"),
    [
        ("worker_memory_limit_exceeded", "pptx_batch_discovery_resource_unavailable"),
        ("unsafe_worker_process_metadata", "pptx_batch_discovery_start_failure"),
        ("worker_start_failed", "pptx_batch_discovery_start_failure"),
        ("worker_exit", "pptx_batch_discovery_worker_failure"),
        (
            "worker_response_authentication_failed",
            "pptx_batch_discovery_protocol_invalid",
        ),
    ],
)
def test_directory_supervisor_failures_keep_closed_reason_distinctions(
    pptx_extraction,
    monkeypatch,
    supervisor_reason,
    expected_reason,
):
    monkeypatch.setattr(
        pptx_extraction,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(
            pptx_extraction.SupervisorError(supervisor_reason)
        ),
    )

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction.batch_extract(
            "/lexical/root",
            [],
            ocr=False,
        )

    assert caught.value.reason_code == expected_reason
    assert caught.value.details == {"supervisor_reason_code": supervisor_reason}


def test_directory_cli_returns_nonzero_structured_discovery_failure(
    pptx_extraction,
    monkeypatch,
    capsys,
):
    monkeypatch.setattr(
        pptx_extraction,
        "run_authenticated_worker",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(
            pptx_extraction.SupervisorError("unsafe_worker_process_metadata")
        ),
    )

    assert pptx_extraction.main(["--directory", "/lexical/root", "--no-ocr"]) == 1

    captured = capsys.readouterr()
    assert json.loads(captured.out) == {
        "schema_version": 1,
        "kind": "pptx_directory_batch",
        "complete": False,
        "incomplete_reason_codes": ["pptx_batch_discovery_start_failure"],
        "results": [],
        "skipped": [{"path": ".", "reason": "pptx_batch_discovery_start_failure"}],
        "error": {
            "reason_code": "pptx_batch_discovery_start_failure",
            "details": {"supervisor_reason_code": "unsafe_worker_process_metadata"},
        },
    }
    assert captured.err == "ERROR: pptx_batch_discovery_start_failure\n"


def test_directory_identity_zero_is_reported_instead_of_silently_deduplicated(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    monkeypatch.setattr(
        pptx_extraction,
        "_usable_directory_identity",
        lambda _value: None,
    )

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(root, [])

    assert discovered == []
    assert skipped == [
        {"path": ".", "reason": "pptx_batch_directory_identity_unavailable"}
    ]


def test_directory_identity_collision_is_reported_not_silently_omitted(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    nested = root / "nested"
    nested.mkdir(parents=True)
    root_identity = (root.lstat().st_dev, root.lstat().st_ino)
    original = pptx_extraction._usable_directory_identity

    def colliding_identity(value):
        identity = original(value)
        if identity == (nested.lstat().st_dev, nested.lstat().st_ino):
            return root_identity
        return identity

    monkeypatch.setattr(
        pptx_extraction,
        "_usable_directory_identity",
        colliding_identity,
    )

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(root, [])

    assert discovered == []
    assert skipped == [
        {
            "path": "nested",
            "reason": "pptx_batch_directory_identity_collision",
        }
    ]


def test_batch_continues_after_one_supervised_failure_with_relative_paths(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    (root / "a-bad.pptx").write_bytes(b"bad")
    (root / "b-good.pptx").write_bytes(b"good")
    _use_in_process_directory_discovery(pptx_extraction, monkeypatch)

    def fake_extract(
        path,
        *,
        trusted_root,
        ocr,
        source_size_limit_bytes,
        deadline_monotonic,
    ):
        assert trusted_root == root
        assert ocr is False
        assert source_size_limit_bytes > 0
        assert deadline_monotonic > 0
        if path.name == "a-bad.pptx":
            raise pptx_extraction.PptxEvidenceError(
                "synthetic path must stay private",
                reason_code="pptx_probe_timeout",
            )
        return {
            "pptx_path": str(path),
            "slide_count": 1,
            "input_fingerprint": {"size_bytes": path.stat().st_size},
        }

    monkeypatch.setattr(pptx_extraction, "extract_pptx", fake_extract)
    batch = pptx_extraction.batch_extract(root, [], ocr=False)

    assert batch["complete"] is False
    assert batch["incomplete_reason_codes"] == ["pptx_probe_timeout"]
    assert batch["results"] == [
        {
            "pptx_path": "b-good.pptx",
            "slide_count": 1,
            "input_fingerprint": {"size_bytes": 4},
        }
    ]
    assert batch["skipped"] == [
        {
            "path": "a-bad.pptx",
            "reason": "pptx_probe_timeout",
        }
    ]
    assert str(root) not in json.dumps(batch)


@pytest.mark.skipif(
    sys.platform == "win32",
    reason="this deterministic swap uses POSIX directory symlinks",
)
def test_batch_rejects_intermediate_directory_swapped_after_discovery(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    nested = root / "nested"
    nested.mkdir(parents=True)
    make_deck(1).save(nested / "deck.pptx")
    outside = tmp_path / "outside"
    outside.mkdir()
    make_deck(2).save(outside / "deck.pptx")
    parked = tmp_path / "original-nested"

    def discover(directory, _patterns, _exclusions, *, deadline):
        assert directory == root
        assert deadline > 0
        nested.rename(parked)
        nested.symlink_to(outside, target_is_directory=True)
        return ["nested/deck.pptx"], [], True, []

    monkeypatch.setattr(
        pptx_extraction,
        "_run_supervised_directory_discovery",
        discover,
    )

    batch = pptx_extraction.batch_extract(root, [], ocr=False)

    assert batch["results"] == []
    assert batch["complete"] is False
    assert batch["skipped"] == [
        {
            "path": "nested/deck.pptx",
            "reason": "pptx_artifact_unavailable",
        }
    ]


def test_batch_file_budget_stops_before_extra_worker(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    (root / "a.pptx").write_bytes(b"a")
    (root / "b.pptx").write_bytes(b"b")
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_FILES", 1)
    _use_in_process_directory_discovery(pptx_extraction, monkeypatch)
    calls = []

    def fake_extract(
        path,
        *,
        trusted_root,
        ocr,
        source_size_limit_bytes,
        deadline_monotonic,
    ):
        assert trusted_root == root
        calls.append(path.name)
        return {
            "pptx_path": str(path),
            "slide_count": 1,
            "input_fingerprint": {"size_bytes": path.stat().st_size},
        }

    monkeypatch.setattr(pptx_extraction, "extract_pptx", fake_extract)
    batch = pptx_extraction.batch_extract(root, [], ocr=False)

    assert calls == ["a.pptx"]
    assert batch["results"][0]["pptx_path"] == "a.pptx"
    assert batch["complete"] is False
    assert batch["incomplete_reason_codes"] == [
        "pptx_batch_file_limit",
        "pptx_batch_scan_incomplete_file_limit",
    ]
    assert batch["skipped"] == [
        {
            "path": "b.pptx",
            "reason": "pptx_batch_file_limit",
        },
        {
            "path": ".",
            "reason": "pptx_batch_scan_incomplete_file_limit",
        },
    ]


def test_directory_entry_budget_stops_before_sorting_unbounded_listing(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    (root / "a.pptx").write_bytes(b"a")
    (root / "b.pptx").write_bytes(b"b")
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_ENTRIES", 1)

    discovered, skipped, _started = pptx_extraction._discover_pptx_files(root, [])

    assert discovered == []
    assert skipped == [{"path": ".", "reason": "pptx_batch_entry_limit"}]


def test_batch_input_budget_uses_launched_generation_not_discovery_size(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    for name in ("a.pptx", "b.pptx", "c.pptx"):
        (root / name).write_bytes(b"x")
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_INPUT_BYTES", 5)
    _use_in_process_directory_discovery(pptx_extraction, monkeypatch)
    calls = []

    def fake_extract(
        path,
        *,
        trusted_root,
        ocr,
        source_size_limit_bytes,
        deadline_monotonic,
    ):
        assert trusted_root == root
        calls.append((path.name, source_size_limit_bytes))
        if path.name == "a.pptx":
            # The admitted generation grew after discovery; charge the exact
            # worker-bound size rather than the stale one-byte directory stat.
            return {
                "pptx_path": str(path),
                "slide_count": 1,
                "input_fingerprint": {"size_bytes": 4},
            }
        raise pptx_extraction.PptxEvidenceError(
            "source no longer fits the exact remaining bytes",
            reason_code="pptx_batch_input_limit",
        )

    monkeypatch.setattr(pptx_extraction, "extract_pptx", fake_extract)
    batch = pptx_extraction.batch_extract(root, [], ocr=False)

    assert calls == [("a.pptx", 5), ("b.pptx", 1)]
    assert batch["results"][0]["input_fingerprint"]["size_bytes"] == 4
    assert batch["complete"] is False
    assert batch["skipped"] == [
        {"path": "b.pptx", "reason": "pptx_batch_input_limit"},
        {"path": "c.pptx", "reason": "pptx_batch_input_limit"},
    ]


def test_batch_charges_admitted_generation_reported_by_failed_launch(
    pptx_extraction,
    monkeypatch,
):
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_INPUT_BYTES", 5)
    monkeypatch.setattr(
        pptx_extraction,
        "_run_supervised_directory_discovery",
        lambda _root, _patterns, _exclusions, *, deadline: (
            ["a.pptx", "b.pptx", "c.pptx"],
            [],
            True,
            [],
        ),
    )
    calls = []

    def fake_extract(path, **options):
        calls.append((path.name, options["source_size_limit_bytes"]))
        if path.name == "a.pptx":
            raise pptx_extraction.PptxEvidenceError(
                "parse failed after admission",
                reason_code="pptx_evidence_invalid",
                details={"admitted_source_size_bytes": 4},
            )
        raise pptx_extraction.PptxEvidenceError(
            "remaining byte budget is smaller than the admitted generation",
            reason_code="pptx_batch_input_limit",
        )

    monkeypatch.setattr(pptx_extraction, "extract_pptx", fake_extract)

    batch = pptx_extraction.batch_extract("/lexical/root", [], ocr=False)

    assert batch["results"] == []
    assert calls == [("a.pptx", 5), ("b.pptx", 1)]
    assert batch["complete"] is False
    assert batch["skipped"] == [
        {"path": "a.pptx", "reason": "pptx_evidence_invalid"},
        {"path": "b.pptx", "reason": "pptx_batch_input_limit"},
        {"path": "c.pptx", "reason": "pptx_batch_input_limit"},
    ]


def test_batch_passes_one_absolute_deadline_and_stops_after_wall_exhaustion(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    first = root / "a.pptx"
    second = root / "b.pptx"
    first.write_bytes(b"a")
    second.write_bytes(b"b")
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_WALL_SECONDS", 10)
    deadlines = []

    def discover(_directory, _patterns, _exclusions, *, deadline):
        deadlines.append(deadline)
        return ["a.pptx", "b.pptx"], [], True, []

    monkeypatch.setattr(
        pptx_extraction, "_run_supervised_directory_discovery", discover
    )
    monkeypatch.setattr(pptx_extraction.time, "monotonic", lambda: 100.0)
    calls = []

    def fake_extract(path, **options):
        calls.append((path.name, options["deadline_monotonic"]))
        raise pptx_extraction.PptxEvidenceError(
            "deadline exhausted",
            reason_code="pptx_batch_wall_limit",
        )

    monkeypatch.setattr(pptx_extraction, "extract_pptx", fake_extract)
    batch = pptx_extraction.batch_extract(root, [], ocr=False)

    assert batch["results"] == []
    assert deadlines == [110.0]
    assert calls == [("a.pptx", 110.0)]
    assert batch["complete"] is False
    assert batch["skipped"] == [
        {"path": "a.pptx", "reason": "pptx_batch_wall_limit"},
        {"path": "b.pptx", "reason": "pptx_batch_wall_limit"},
    ]


def test_batch_compact_output_budget_includes_wrapper_and_skips(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    root = tmp_path / "decks"
    root.mkdir()
    lock = root / "~$locked.pptx"
    lock.write_bytes(b"lock")
    expected_skips = [
        {"path": "~$locked.pptx", "reason": "pptx_batch_office_lock_file"}
    ]
    expected_batch = pptx_extraction._build_batch_output([], expected_skips)
    exact = pptx_extraction._encode_batch_output(expected_batch) + b"\n"
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_OUTPUT_BYTES", len(exact))
    _use_in_process_directory_discovery(pptx_extraction, monkeypatch)

    batch = pptx_extraction.batch_extract(root, [], ocr=False)

    assert batch["results"] == []
    assert batch["complete"] is True
    assert batch["skipped"] == expected_skips
    assert len(pptx_extraction._encode_batch_output(batch)) + 1 == len(exact)


def test_directory_cli_emits_the_same_compact_bytes_budgeted_by_batch(
    pptx_extraction,
    monkeypatch,
    tmp_path,
    capsys,
):
    root = tmp_path / "decks"
    root.mkdir()
    (root / "~$locked.pptx").write_bytes(b"lock")
    expected = (
        pptx_extraction._encode_batch_output(
            pptx_extraction._build_batch_output(
                [],
                [
                    {
                        "path": "~$locked.pptx",
                        "reason": "pptx_batch_office_lock_file",
                    }
                ],
            )
        )
        + b"\n"
    )
    monkeypatch.setattr(pptx_extraction, "_BATCH_MAX_OUTPUT_BYTES", len(expected))
    _use_in_process_directory_discovery(pptx_extraction, monkeypatch)

    assert pptx_extraction.main(["--directory", str(root), "--no-ocr"]) == 0

    captured = capsys.readouterr()
    assert captured.err == ""
    assert captured.out.encode("utf-8") == expected
    assert len(captured.out.encode("utf-8")) <= pptx_extraction._BATCH_MAX_OUTPUT_BYTES


def test_directory_cli_partial_batch_exits_zero_with_exact_public_envelope(
    pptx_extraction,
    monkeypatch,
    capsys,
):
    monkeypatch.setattr(
        pptx_extraction,
        "_run_supervised_directory_discovery",
        lambda _root, _patterns, _exclusions, *, deadline: (
            ["broken.pptx"],
            [],
            True,
            [],
        ),
    )
    monkeypatch.setattr(
        pptx_extraction,
        "extract_pptx",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(
            pptx_extraction.PptxEvidenceError(
                "fixture parse failure",
                reason_code="pptx_parse_failure",
            )
        ),
    )

    exit_code = pptx_extraction.main(["--directory", "/lexical/root", "--no-ocr"])

    captured = capsys.readouterr()
    assert exit_code == 0
    assert captured.err == ""
    assert json.loads(captured.out) == {
        "schema_version": 1,
        "kind": "pptx_directory_batch",
        "complete": False,
        "incomplete_reason_codes": ["pptx_parse_failure"],
        "results": [],
        "skipped": [{"path": "broken.pptx", "reason": "pptx_parse_failure"}],
    }


def test_per_slide_visual_count(pptx_extraction, tmp_path):
    prs = make_deck(3)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    result = pptx_extraction._extract_pptx_in_process(path)
    assert len(result["per_slide_visual"]) == 3


def test_template_layouts_emitted(pptx_extraction, tmp_path):
    """extract_pptx must emit a top-level template_layouts key."""
    prs = make_deck(1)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    result = pptx_extraction._extract_pptx_in_process(path)
    assert "template_layouts" in result
    assert isinstance(result["template_layouts"], list)


def test_template_layouts_default_count(pptx_extraction):
    """A default python-pptx Presentation ships with 11 stock layouts under 1 master."""
    prs = Presentation()
    layouts = pptx_extraction.extract_template_layouts(prs)
    assert len(layouts) == 11
    assert all(layout["master_index"] == 0 for layout in layouts)


def test_template_layouts_entry_shape(pptx_extraction):
    """Each layout entry has the canonical keys with the documented types."""
    prs = Presentation()
    layouts = pptx_extraction.extract_template_layouts(prs)
    expected_keys = {"index", "master_index", "name", "placeholders"}
    for layout in layouts:
        assert set(layout.keys()) == expected_keys
        assert isinstance(layout["index"], int)
        assert isinstance(layout["master_index"], int)
        assert isinstance(layout["name"], str)
        assert isinstance(layout["placeholders"], list)


def test_template_layouts_index_is_global_and_sequential(pptx_extraction):
    """The `index` field is a global running counter, not per-master."""
    prs = Presentation()
    layouts = pptx_extraction.extract_template_layouts(prs)
    indices = [layout["index"] for layout in layouts]
    assert indices == list(range(len(layouts)))


def test_template_layouts_placeholder_shape(pptx_extraction):
    """Placeholder entries carry idx (int) + type (canonical name string)."""
    prs = Presentation()
    layouts = pptx_extraction.extract_template_layouts(prs)
    # python-pptx stock template has at least one layout with a TITLE placeholder.
    title_layouts = [
        layout
        for layout in layouts
        if any(p.get("type") == "TITLE" for p in layout["placeholders"])
    ]
    assert title_layouts, "expected at least one layout with a TITLE placeholder"
    for layout in layouts:
        for p in layout["placeholders"]:
            assert set(p.keys()) == {"idx", "type"}
            assert isinstance(p["idx"], int)
            assert isinstance(p["type"], str)
            # Type names should be canonical enum identifiers (no surrounding " (3)" digits)
            assert "(" not in p["type"]


def test_template_layouts_known_layout_name(pptx_extraction):
    """At least one layout from the python-pptx stock template carries a recognizable name."""
    prs = Presentation()
    layouts = pptx_extraction.extract_template_layouts(prs)
    names = [layout["name"] for layout in layouts]
    # python-pptx stock layouts include "Title Slide" and "Blank" as part of the default master.
    assert "Title Slide" in names
    assert "Blank" in names


# ── text_extraction_confidence (issue #116) ──────────────────────────
#
# `pptx-extraction.py` reads text out of PPTX shapes. Text rendered inside a
# picture — the norm for AI-generated illustration decks — is invisible to it.
# These tests pin the contract that the extractor reports that blindness
# instead of asserting the slide is wordless.


def _png(path, w=16, h=16):
    """Emit a minimal valid PNG from stdlib only (no image library needed)."""
    import struct
    import zlib

    def chunk(tag, data):
        body = tag + data
        return (
            struct.pack(">I", len(data))
            + body
            + struct.pack(">I", zlib.crc32(body) & 0xFFFFFFFF)
        )

    raw = b"".join(b"\x00" + b"\x7f\x7f\x7f" * w for _ in range(h))
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )
    return str(path)


def _first_slide(pptx_extraction, prs, tmp_path):
    path = str(tmp_path / "deck.pptx")
    prs.save(path)
    return pptx_extraction._extract_pptx_in_process(path)["per_slide_visual"][0]


def test_picture_area_ratio_full_bleed(pptx_extraction, tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png(tmp_path / "i.png"),
        Inches(0),
        Inches(0),
        width=deck_width(prs),
        height=deck_height(prs),
    )
    data = _first_slide(pptx_extraction, prs, tmp_path)
    assert data["image_area_ratio"] > 0.99


def test_picture_area_ratio_missing_geometry_is_zero(pptx_extraction):
    """Unknown size is not evidence of a large picture."""

    class _Shape:
        width = None
        height = None

    class _Prs:
        slide_width = 9144000
        slide_height = 6858000

    assert pptx_extraction.picture_area_ratio(_Shape(), _Prs()) == 0.0


def test_full_bleed_image_slide_does_not_assert_absence(
    pptx_extraction,
    tmp_path,
):
    """Issue #116: a full-bleed image slide must not read as 'no text'."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png(tmp_path / "i.png"),
        Inches(0),
        Inches(0),
        width=deck_width(prs),
        height=deck_height(prs),
    )
    data = _first_slide(pptx_extraction, prs, tmp_path)

    assert data["has_image"] is True
    assert data["has_text_frame_shapes"] is False
    assert data["text_content_preview"] == ""
    # The load-bearing assertion: unreadable is reported as low confidence,
    # never as evidence the slide is wordless.
    assert data["text_extraction_confidence"] == "low"


def test_text_slide_is_high_confidence(pptx_extraction, tmp_path):
    """No picture — extractable text is the whole story."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_title(slide).text = "A real title"
    data = _first_slide(pptx_extraction, prs, tmp_path)
    assert data["has_text_frame_shapes"] is True
    assert data["text_extraction_confidence"] == "high"


def test_small_decorative_image_stays_high_confidence(
    pptx_extraction,
    tmp_path,
):
    """A logo-sized picture cannot be hiding the slide's content."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_title(slide).text = "Title"
    slide.shapes.add_picture(
        _png(tmp_path / "i.png"),
        Inches(0),
        Inches(0),
        width=Emu(int(deck_width(prs) * 0.1)),
        height=Emu(int(deck_height(prs) * 0.1)),
    )
    data = _first_slide(pptx_extraction, prs, tmp_path)
    assert data["has_image"] is True
    assert data["image_area_ratio"] < 0.5
    assert data["text_extraction_confidence"] == "high"


def test_text_overlay_over_full_bleed_is_still_low_confidence(
    pptx_extraction,
    tmp_path,
):
    """Extracting *some* text is not evidence of extracting *all* of it."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_title(slide).text = "Overlay"
    slide.shapes.add_picture(
        _png(tmp_path / "i.png"),
        Inches(0),
        Inches(0),
        width=deck_width(prs),
        height=deck_height(prs),
    )
    data = _first_slide(pptx_extraction, prs, tmp_path)
    assert data["has_text_frame_shapes"] is True
    assert data["text_extraction_confidence"] == "low"


def test_retired_field_is_gone(pptx_extraction, tmp_path):
    """`has_text_placeholder` named a claim the extractor cannot make."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    data = _first_slide(pptx_extraction, prs, tmp_path)
    assert "has_text_placeholder" not in data
    assert "has_text_frame_shapes" in data


def test_image_background_slide_is_low_confidence(
    pptx_extraction,
    tmp_path,
    monkeypatch,
):
    """An image *background* covers the slide and can carry baked-in text.

    It is not a PICTURE shape, so the shape walk never sees it — the same
    blindness as issue #116, one layer down. python-pptx has no public authoring
    API for image backgrounds, so this test stubs only the classifier. XML/blob
    extraction is covered by the synthetic background fixture below.
    """
    monkeypatch.setattr(
        pptx_extraction,
        "get_background_color",
        lambda slide: (None, "image"),
    )
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # no pictures, no text
    data = _first_slide(pptx_extraction, prs, tmp_path)

    assert data["background_type"] == "image"
    assert data["has_image"] is False  # not a PICTURE shape
    assert data["image_area_ratio"] == 0.0  # no picture geometry at all
    assert data["text_extraction_confidence"] == "low"


def test_area_ratio_is_not_rounded_across_the_threshold(pptx_extraction):
    """Rounding must not decide classification.

    A picture at 0.4996 of the slide is below the threshold; rounding it to
    0.5 first would flip it and make the threshold depend on the rounding.
    """

    class _Prs:
        slide_width = 10000
        slide_height = 10000

    class _Shape:
        # 0.4996 of the slide area — just under the threshold.
        width = 4996
        height = 10000

    ratio = pptx_extraction.picture_area_ratio(_Shape(), _Prs())
    assert ratio < pptx_extraction.PPTX_TEXT_BEARING_IMAGE_AREA_RATIO
    assert ratio == pytest.approx(0.4996)


# ── OCR inventory on low-confidence slides (issue #129) ──────────────
#
# #116/#119 stopped asserting absence. #129 fills a word inventory from
# picture blobs so analysis can cite and cross-check, without replacing
# the vision pass for design dimensions.


def _png_with_text(path, text="VENUE PREPARATION", w=800, h=200):
    """Build a high-contrast PNG with clear text via Pillow (no binary fixture).

    Prefer a TrueType font when the OS has one; otherwise draw with the
    default bitmap font on a small canvas and nearest-neighbor upscale so OCR
    stays reliable without fixtures.
    """
    from PIL import Image, ImageDraw, ImageFont

    font = None
    for candidate in (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",  # CI (Ubuntu)
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",  # macOS
        "/Library/Fonts/Arial Bold.ttf",
    ):
        try:
            font = ImageFont.truetype(candidate, 48)
            break
        except OSError:
            continue

    if font is None:
        scale = 8
        small = Image.new("RGB", (max(w // scale, 100), max(h // scale, 40)), "white")
        ImageDraw.Draw(small).text(
            (2, 2),
            text,
            fill="black",
            font=ImageFont.load_default(),
        )
        img = small.resize((w, h), Image.Resampling.NEAREST)
    else:
        img = Image.new("RGB", (w, h), "white")
        ImageDraw.Draw(img).text((20, 60), text, fill="black", font=font)

    img.save(path, format="PNG")
    return str(path)


def test_normalize_ocr_text(pptx_extraction):
    assert pptx_extraction.normalize_ocr_text("  a \n\tb  ") == "a b"
    assert pptx_extraction.normalize_ocr_text("") == ""
    assert pptx_extraction.normalize_ocr_text(None) == ""


def test_ocr_confidence_is_paired_only_with_retained_token(pptx_extraction):
    text, confidence = pptx_extraction._ocr_text_and_confidence(
        {
            "text": ["", "VISIBLE"],
            "conf": [99.0, 12.0],
        }
    )
    assert text == "VISIBLE"
    assert confidence == 12.0


def test_high_confidence_slide_method_is_shapes_only(pptx_extraction, tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_title(slide).text = "A real title"
    data = _first_slide(pptx_extraction, prs, tmp_path)
    assert data["text_extraction_confidence"] == "high"
    assert data["text_extraction_method"] == "shapes"
    assert data["ocr_text"] == ""


def test_full_bleed_runs_ocr_fn_and_records_inventory(
    pptx_extraction,
    tmp_path,
):
    """Low-confidence picture slide gets ocr_text from the OCR path."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png_with_text(tmp_path / "labeled.png"),
        Inches(0),
        Inches(0),
        width=deck_width(prs),
        height=deck_height(prs),
    )
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    data = pptx_extraction._extract_pptx_in_process(
        path,
        ocr_fn=lambda blobs: "VENUE PREPARATION",
    )["per_slide_visual"][0]

    assert data["text_extraction_confidence"] == "low"
    assert data["text_content_preview"] == ""  # shapes still empty
    assert data["ocr_text"] == "VENUE PREPARATION"
    assert data["text_extraction_method"] == "shapes+ocr"


def test_multi_image_ocr_emits_one_identity_and_outcome_receipt_per_asset(
    pptx_extraction,
    tmp_path,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png_with_text(tmp_path / "left.png", text="LEFT LABEL"),
        Inches(0),
        Inches(0),
        width=Emu(int(deck_width(prs) / 2)),
        height=deck_height(prs),
    )
    slide.shapes.add_picture(
        _png_with_text(tmp_path / "right.png", text="RIGHT LABEL"),
        Emu(int(deck_width(prs) / 2)),
        Inches(0),
        width=Emu(int(deck_width(prs) / 2)),
        height=deck_height(prs),
    )
    path = tmp_path / "multi-image.pptx"
    prs.save(path)
    outcomes = iter(
        [
            {
                "engine": "fixture-ocr",
                "engine_version": "1.0",
                "result_status": "text_recovered",
                "result_confidence": 94.5,
                "recovered_text": "LEFT LABEL",
                "trustworthy_text": True,
                "error": None,
            },
            {
                "engine": "fixture-ocr",
                "engine_version": "1.0",
                "result_status": "low_confidence_text",
                "result_confidence": 31.0,
                "recovered_text": "RIGHT LAB3L",
                "trustworthy_text": False,
                "error": None,
            },
        ]
    )

    data = pptx_extraction._extract_pptx_in_process(
        str(path),
        ocr_fn=lambda _blobs: next(outcomes),
    )["per_slide_visual"][0]
    channel = next(
        item for item in data["text_channels"] if item["channel"] == "picture_ocr"
    )

    assert channel["status"] == "partial"
    assert channel["reason"] == "partial_ocr_results"
    assert channel["text"] == "LEFT LABEL | RIGHT LAB3L"
    assert channel["result_confidence"] == pytest.approx(62.75)
    receipts = channel["ocr_receipts"]
    assert len(receipts) == 2
    assert [receipt["result_status"] for receipt in receipts] == [
        "text_recovered",
        "low_confidence_text",
    ]
    assert [receipt["part_name"] for receipt in receipts] == [
        "ppt/media/image1.png",
        "ppt/media/image2.png",
    ]
    assert all(len(receipt["asset_sha256"]) == 64 for receipt in receipts)
    assert receipts[0]["asset_sha256"] != receipts[1]["asset_sha256"]
    assert all(receipt["shape_path"] for receipt in receipts)
    assert receipts[1]["recovered_text"] == "RIGHT LAB3L"
    assert receipts[1]["trustworthy_text"] is False


def test_genuine_empty_ocr_is_distinct_from_asset_failure(
    pptx_extraction,
    tmp_path,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for index in range(2):
        slide.shapes.add_picture(
            _png(tmp_path / f"asset-{index}.png", w=16 + index),
            Emu(int(deck_width(prs) * index / 2)),
            Inches(0),
            width=Emu(int(deck_width(prs) / 2)),
            height=deck_height(prs),
        )
    path = tmp_path / "empty-vs-failure.pptx"
    prs.save(path)
    outcomes = iter(
        [
            {
                "engine": "fixture-ocr",
                "engine_version": "1.0",
                "result_status": "genuine_empty",
                "result_confidence": None,
                "recovered_text": "",
                "trustworthy_text": False,
                "error": None,
            },
            {
                "engine": "fixture-ocr",
                "engine_version": "1.0",
                "result_status": "failed",
                "result_confidence": None,
                "recovered_text": "",
                "trustworthy_text": False,
                "error": "engine_error: synthetic failure",
            },
        ]
    )

    data = pptx_extraction._extract_pptx_in_process(
        str(path),
        ocr_fn=lambda _blobs: next(outcomes),
    )["per_slide_visual"][0]
    channel = next(
        item for item in data["text_channels"] if item["channel"] == "picture_ocr"
    )

    assert channel["text"] == ""
    assert channel["status"] == "failed"
    assert [item["result_status"] for item in channel["ocr_receipts"]] == [
        "genuine_empty",
        "failed",
    ]
    assert channel["ocr_receipts"][0]["error"] is None
    assert "synthetic failure" in channel["ocr_receipts"][1]["error"]


def test_truncated_picture_is_one_failed_receipt_not_a_deck_abort(
    pptx_extraction,
    monkeypatch,
    tmp_path,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png(tmp_path / "healthy.png", w=64),
        Inches(0),
        Inches(0),
        width=Emu(int(deck_width(prs) / 2)),
        height=deck_height(prs),
    )
    slide.shapes.add_picture(
        _png(tmp_path / "truncated.png", w=65),
        Emu(int(deck_width(prs) / 2)),
        Inches(0),
        width=Emu(int(deck_width(prs) / 2)),
        height=deck_height(prs),
    )
    path = tmp_path / "mixed-assets.pptx"
    prs.save(path)
    source = io.BytesIO(path.read_bytes())
    rewritten = io.BytesIO()
    with (
        zipfile.ZipFile(source) as archive,
        zipfile.ZipFile(rewritten, "w") as destination,
    ):
        for member in archive.infolist():
            payload = archive.read(member)
            if member.filename == "ppt/media/image2.png":
                payload = payload[:40]
            destination.writestr(member, payload)
    path.write_bytes(rewritten.getvalue())

    class FakeTesseractError(Exception):
        pass

    def image_to_data(image, *, output_type):
        assert output_type == "DICT"
        image.load()
        return {"text": ["HEALTHY"], "conf": ["93"]}

    fake_pytesseract = SimpleNamespace(
        Output=SimpleNamespace(DICT="DICT"),
        TesseractNotFoundError=FakeTesseractError,
        TesseractError=FakeTesseractError,
        image_to_data=image_to_data,
    )
    monkeypatch.setitem(sys.modules, "pytesseract", fake_pytesseract)
    monkeypatch.setattr(pptx_extraction, "_tesseract_available", True)
    monkeypatch.setattr(pptx_extraction, "_tesseract_version", "fixture-1")

    data = pptx_extraction._extract_pptx_in_process(path)["per_slide_visual"][0]
    channel = next(
        item for item in data["text_channels"] if item["channel"] == "picture_ocr"
    )

    assert channel["status"] == "partial"
    assert channel["text"] == "HEALTHY"
    assert [item["result_status"] for item in channel["ocr_receipts"]] == [
        "text_recovered",
        "failed",
    ]
    assert [item["part_name"] for item in channel["ocr_receipts"]] == [
        "ppt/media/image1.png",
        "ppt/media/image2.png",
    ]
    failed = channel["ocr_receipts"][1]
    assert failed["attempted"] is True
    assert failed["trustworthy_text"] is False
    assert failed["error"].startswith("image_decode_error:")
    assert "/" not in failed["error"]


def test_no_ocr_flag_skips_inventory(pptx_extraction, tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png_with_text(tmp_path / "labeled.png"),
        Inches(0),
        Inches(0),
        width=deck_width(prs),
        height=deck_height(prs),
    )
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    called = []

    def spy(blobs):
        called.append(blobs)
        return "SHOULD NOT APPEAR"

    data = pptx_extraction._extract_pptx_in_process(
        path,
        ocr=False,
        ocr_fn=spy,
    )["per_slide_visual"][0]

    assert called == []
    assert data["ocr_text"] == ""
    assert data["text_extraction_method"] == "shapes"
    assert data["text_extraction_confidence"] == "low"
    channel = next(
        item for item in data["text_channels"] if item["channel"] == "picture_ocr"
    )
    assert channel["status"] == "skipped"
    assert channel["ocr_receipts"][0]["attempted"] is False
    assert channel["ocr_receipts"][0]["result_status"] == "skipped"
    assert channel["ocr_receipts"][0]["part_name"] == "ppt/media/image1.png"


def test_small_decorative_image_does_not_ocr(pptx_extraction, tmp_path):
    """High-confidence slides must not pay for OCR."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_title(slide).text = "Title"
    slide.shapes.add_picture(
        _png_with_text(tmp_path / "logo.png", text="LOGO"),
        Inches(0),
        Inches(0),
        width=Emu(int(deck_width(prs) * 0.1)),
        height=Emu(int(deck_height(prs) * 0.1)),
    )
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    called = []
    data = pptx_extraction._extract_pptx_in_process(
        path,
        ocr_fn=lambda blobs: called.append(blobs) or "X",
    )["per_slide_visual"][0]

    assert data["text_extraction_confidence"] == "high"
    assert called == []
    assert data["ocr_text"] == ""
    assert data["text_extraction_method"] == "shapes"


def test_ocr_unavailable_records_method_not_crash(pptx_extraction, tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png(tmp_path / "i.png"),
        Inches(0),
        Inches(0),
        width=deck_width(prs),
        height=deck_height(prs),
    )
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    def boom(_blobs):
        raise pptx_extraction.OcrUnavailableError("no engine")

    data = pptx_extraction._extract_pptx_in_process(path, ocr_fn=boom)[
        "per_slide_visual"
    ][0]
    assert data["ocr_text"] == ""
    assert data["text_extraction_method"] == "shapes+ocr_unavailable"
    assert data["text_extraction_confidence"] == "low"
    channel = next(
        item for item in data["text_channels"] if item["channel"] == "picture_ocr"
    )
    assert channel["attempted"] is False
    assert channel["engine"] == "tesseract"
    assert channel["engine_version"] is None
    assert channel["status"] == "unavailable"
    assert channel["reason"] == "ocr_engine_unavailable"
    assert channel["ocr_receipts"][0]["attempted"] is False


def test_reported_image_background_without_blob_does_not_claim_ocr(
    pptx_extraction,
    tmp_path,
    monkeypatch,
):
    """A reported background without actual XML/blob never claims OCR ran."""
    monkeypatch.setattr(
        pptx_extraction,
        "get_background_color",
        lambda slide: (None, "image"),
    )
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    called = []
    data = pptx_extraction._extract_pptx_in_process(
        path,
        ocr_fn=lambda blobs: called.append(blobs) or "X",
    )["per_slide_visual"][0]

    assert data["text_extraction_confidence"] == "low"
    assert called == []
    assert data["ocr_text"] == ""
    assert data["text_extraction_method"] == "shapes"
    channel = next(
        item
        for item in data["text_channels"]
        if item["channel"] == "background_image_ocr"
    )
    assert channel["attempted"] is False
    assert channel["engine"] == "tesseract"
    assert channel["engine_version"] is None
    assert channel["status"] == "unavailable"
    assert channel["reason"] == "no_readable_asset"
    assert channel["ocr_receipts"] == []
    assert data["render_required"] is True


def test_ocr_text_capped(pptx_extraction, monkeypatch):
    monkeypatch.setattr(
        pptx_extraction,
        "ocr_image_bytes",
        lambda blob: "A" * 9000,
    )
    out = pptx_extraction.ocr_picture_blobs([b"x"])
    assert len(out) == pptx_extraction._OCR_TEXT_MAX_CHARS


@pytest.mark.parametrize(
    "injected",
    [
        "A" * 9000,
        {
            "engine": "fixture",
            "engine_version": "1",
            "result_status": "text_recovered",
            "result_confidence": 90.0,
            "recovered_text": "A" * 9000,
            "trustworthy_text": True,
            "error": None,
        },
    ],
)
def test_each_ocr_receipt_text_is_capped(pptx_extraction, injected):
    slide_data = {
        "text_channels": [],
        "ocr_text": "",
        "text_extraction_method": "shapes",
    }
    pptx_extraction._run_ocr_channel(
        slide_data,
        [
            {
                "blob": b"exact-asset",
                "shape_path": ["Picture 1"],
                "part_name": "ppt/media/image1.png",
            }
        ],
        channel="picture_ocr",
        provenance={"source": "embedded_picture_blobs"},
        ocr_fn=lambda _blobs: injected,
    )
    receipt = slide_data["text_channels"][0]["ocr_receipts"][0]
    assert len(receipt["recovered_text"]) == pptx_extraction._OCR_TEXT_MAX_CHARS
    assert len(slide_data["ocr_text"]) == pptx_extraction._OCR_TEXT_MAX_CHARS


def test_untrustworthy_recovered_receipt_keeps_channel_partial(pptx_extraction):
    slide_data = {
        "text_channels": [],
        "ocr_text": "",
        "text_extraction_method": "shapes",
    }
    pptx_extraction._run_ocr_channel(
        slide_data,
        [
            {
                "blob": b"exact-asset",
                "shape_path": ["Picture 1"],
                "part_name": "ppt/media/image1.png",
            }
        ],
        channel="picture_ocr",
        provenance={"source": "embedded_picture_blobs"},
        ocr_fn=lambda _blobs: {
            "engine": "fixture",
            "engine_version": "1",
            "result_status": "text_recovered",
            "result_confidence": 10.0,
            "recovered_text": "UNCERTAIN",
            "trustworthy_text": False,
            "error": None,
        },
    )
    channel = slide_data["text_channels"][0]
    assert channel["text"] == "UNCERTAIN"
    assert channel["status"] == "partial"
    assert channel["reason"] == "partial_ocr_results"


def test_ocr_image_bytes_reads_clear_text(pptx_extraction, tmp_path):
    """Integration: real tesseract on a programmatic text PNG."""
    pytesseract = pytest.importorskip("pytesseract")
    try:
        pytesseract.get_tesseract_version()
    except pytesseract.TesseractNotFoundError:
        pytest.skip("tesseract binary not installed")

    path = tmp_path / "clear.png"
    _png_with_text(path, text="VENUE PREPARATION")
    text = pptx_extraction.ocr_image_bytes(path.read_bytes())
    assert "VENUE" in text.upper()
    assert "PREPARATION" in text.upper()


def test_full_bleed_with_baked_text_end_to_end(pptx_extraction, tmp_path):
    """Integration: full-bleed labeled picture yields non-empty ocr_text."""
    pytesseract = pytest.importorskip("pytesseract")
    try:
        pytesseract.get_tesseract_version()
    except pytesseract.TesseractNotFoundError:
        pytest.skip("tesseract binary not installed")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        _png_with_text(tmp_path / "labeled.png", text="VENUE PREPARATION"),
        Inches(0),
        Inches(0),
        width=deck_width(prs),
        height=deck_height(prs),
    )
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    data = pptx_extraction._extract_pptx_in_process(path)["per_slide_visual"][0]
    assert data["text_extraction_confidence"] == "low"
    assert data["text_content_preview"] == ""
    assert data["text_extraction_method"] == "shapes+ocr"
    assert "VENUE" in data["ocr_text"].upper()
    assert "PREPARATION" in data["ocr_text"].upper()


# ── recursive/container fidelity and provenance ──────────────────────


def test_grouped_shapes_are_walked_recursively_and_lower_confidence(
    pptx_extraction,
    tmp_path,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    outer = slide.shapes.add_group_shape()
    outer_text = outer.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    )
    outer_text.text_frame.text = "Text inside outer group"
    inner = outer.shapes.add_group_shape()
    inner_text = inner.shapes.add_textbox(
        Inches(2),
        Inches(2),
        Inches(4),
        Inches(1),
    )
    inner_text.text_frame.text = "Text inside nested group"

    data = _first_slide(pptx_extraction, prs, tmp_path)

    assert "Text inside outer group" in data["text_content_preview"]
    assert "Text inside nested group" in data["text_content_preview"]
    assert data["shape_count"] == 1
    assert data["shape_count_recursive"] == 4
    assert data["text_extraction_confidence"] == "low"
    assert data["render_required"] is True
    assert "grouped_shapes" in data["render_required_reasons"]

    nested = next(
        channel
        for channel in data["text_channels"]
        if channel["text"] == "Text inside nested group"
    )
    assert nested["confidence"] == "medium"
    assert nested["provenance"]["source"] == "pptx_shape_text_frame"
    assert len(nested["provenance"]["shape_path"]) == 3


def test_table_cell_text_has_its_own_provenance_channel(
    pptx_extraction,
    tmp_path,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(2),
    )
    table_shape.table.cell(0, 0).text = "Name"
    table_shape.table.cell(0, 1).text = "Count"
    table_shape.table.cell(1, 0).text = "Hooks"
    table_shape.table.cell(1, 1).text = "12"

    data = _first_slide(pptx_extraction, prs, tmp_path)

    assert "Hooks" in data["text_content_preview"]
    channel = next(
        item for item in data["text_channels"] if item["channel"] == "table_cell_text"
    )
    assert channel["confidence"] == "medium"
    assert channel["status"] == "extracted"
    assert channel["provenance"]["source"] == "pptx_table_cells"
    assert channel["provenance"]["cells"] == ["R1C1", "R1C2", "R2C1", "R2C2"]
    assert data["text_extraction_confidence"] == "low"
    assert "table" in data["render_required_reasons"]
    summary = next(
        item
        for item in data["shapes_summary"]
        if item.get("graphic_frame_type") == "table"
    )
    assert summary["table_rows"] == 2
    assert summary["table_columns"] == 2


def test_smartart_and_unknown_graphic_frames_are_explicitly_unsupported(
    pptx_extraction,
    tmp_path,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    smartart = slide.shapes.add_table(
        1,
        1,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    )
    graphic_frame_element(
        smartart
    ).graphic.graphicData.uri = (
        "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    )
    other = slide.shapes.add_table(
        1,
        1,
        Inches(1),
        Inches(3),
        Inches(3),
        Inches(1),
    )
    graphic_frame_element(
        other
    ).graphic.graphicData.uri = "urn:example:unsupported-graphic"

    data = _first_slide(pptx_extraction, prs, tmp_path)

    assert data["has_unsupported_content"] is True
    assert data["text_extraction_confidence"] == "low"
    assert data["render_required"] is True
    kinds = {item["content_type"] for item in data["unsupported_content"]}
    assert kinds == {"smartart", "graphic_frame"}
    assert {
        item["channel"]
        for item in data["text_channels"]
        if item["status"] == "unsupported"
    } == {"smartart_text", "graphic_frame_text"}
    assert all(item["graphic_data_uri"] for item in data["unsupported_content"])


def _set_background_image(slide, image_path):
    """Author an image background directly in DrawingML for fixture coverage."""
    _, r_id = slide.part.get_or_add_image_part(str(image_path))
    bg_pr = slide.element.cSld.get_or_add_bgPr()
    existing_fill = bg_pr.eg_fillProperties
    if existing_fill is not None:
        bg_pr.remove(existing_fill)
    blip_fill = parse_xml(
        f"<a:blipFill {nsdecls('a', 'r')}>"
        f'<a:blip r:embed="{r_id}"/>'
        "<a:stretch><a:fillRect/></a:stretch>"
        "</a:blipFill>"
    )
    bg_pr.insert(0, blip_fill)
    return r_id


def test_background_image_blob_is_ocrd_with_distinct_provenance(
    pptx_extraction,
    tmp_path,
):
    image_path = tmp_path / "background.png"
    _png_with_text(image_path, text="BACKGROUND LABEL")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background_image(slide, image_path)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    seen = []
    data = pptx_extraction._extract_pptx_in_process(
        path,
        ocr_fn=lambda blobs: seen.append(blobs) or "BACKGROUND LABEL",
    )["per_slide_visual"][0]

    assert data["background_type"] == "image"
    assert data["has_image"] is False
    assert data["text_extraction_confidence"] == "low"
    assert data["render_required"] is True
    assert "background_image" in data["render_required_reasons"]
    assert len(seen) == 1 and seen[0][0].startswith(b"\x89PNG")
    channel = next(
        item
        for item in data["text_channels"]
        if item["channel"] == "background_image_ocr"
    )
    assert channel["status"] == "extracted"
    assert channel["text"] == "BACKGROUND LABEL"
    assert channel["confidence"] == "low"
    assert channel["provenance"]["source"] == "pptx_background_image"
    assert channel["provenance"]["part_name"].startswith("ppt/media/")
    assert data["ocr_text"] == "BACKGROUND LABEL"


def test_background_inspection_does_not_break_inheritance(pptx_extraction):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert slide.element.cSld.bg is None

    assert pptx_extraction.get_background_color(slide) == (None, "unknown")

    # Access through python-pptx's public fill property would author a noFill
    # node here and sever inheritance. Extraction must stay read-only.
    assert slide.element.cSld.bg is None


def test_missing_background_blob_requires_render_without_claiming_ocr(
    pptx_extraction,
    tmp_path,
):
    image_path = tmp_path / "background.png"
    _png(image_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background_image(slide, image_path)
    blip = next(
        iter(
            slide.element.cSld.bg.bgPr.eg_fillProperties.iter(
                pptx_extraction.qn("a:blip")
            )
        )
    )
    blip.set(pptx_extraction.qn("r:embed"), "rIdMissing")

    data = _first_slide(pptx_extraction, prs, tmp_path)

    channel = next(
        item
        for item in data["text_channels"]
        if item["channel"] == "background_image_ocr"
    )
    assert channel["status"] == "unavailable"
    assert channel["text"] == ""
    assert data["text_extraction_method"] == "shapes"
    assert data["text_extraction_confidence"] == "low"
    assert data["render_required"] is True


def _damage_first_media_member(path):
    """Flip a compressed payload byte without updating its recorded CRC."""
    with zipfile.ZipFile(path) as archive:
        member = next(
            item
            for item in archive.infolist()
            if item.filename.startswith("ppt/media/") and item.file_size
        )
    package = bytearray(path.read_bytes())
    name_size, extra_size = struct.unpack_from(
        "<HH",
        package,
        member.header_offset + 26,
    )
    payload_offset = member.header_offset + 30 + name_size + extra_size
    package[payload_offset + (member.compress_size // 2)] ^= 0xFF
    path.write_bytes(package)
    return member.filename


def test_bad_crc_media_is_recovered_without_losing_the_deck(
    pptx_extraction,
    tmp_path,
):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "Healthy native text"
    slide.shapes.add_picture(
        _png(tmp_path / "asset.png"),
        Inches(1),
        Inches(2),
        Inches(2),
        Inches(2),
    )
    _append_timing_xml(slide, "<p:animEffect/>")
    path = tmp_path / "corrupt-media.pptx"
    prs.save(path)
    damaged_name = _damage_first_media_member(path)

    result = pptx_extraction._extract_pptx_in_process(str(path), ocr=False)
    data = result["per_slide_visual"][0]

    assert result["slide_count"] == 1
    assert result["corrupt_assets"] == [
        {
            "part_name": damaged_name,
            "error_type": "crc_mismatch",
            "status": "recovered_with_placeholder",
        }
    ]
    assert "Healthy native text" in data["text_content_preview"]
    assert data["native_timing"]["animation_behavior_counts"]["effect"] == 1
    assert data["text_extraction_confidence"] == "low"
    assert "corrupt_embedded_asset" in data["render_required_reasons"]
    assert any(
        item["content_type"] == "corrupt_embedded_asset"
        for item in data["unsupported_content"]
    )
    channel = next(
        item for item in data["text_channels"] if item["channel"] == "picture_ocr"
    )
    assert channel["attempted"] is False
    assert channel["status"] == "unavailable"
    assert channel["reason"] == "no_readable_asset"
    assert channel["ocr_receipts"] == []
    assert data["render_required"] is True


# ── native timing/build structure (issue #151) ──────────────────────


def _append_timing_xml(slide, behavior_xml):
    """Append a synthetic but package-native PresentationML timing tree."""
    timing = parse_xml(
        f"<p:timing {nsdecls('p')}>"
        '<p:tnLst><p:par><p:cTn id="1"><p:childTnLst>'
        f"{behavior_xml}"
        "</p:childTnLst></p:cTn></p:par></p:tnLst>"
        "</p:timing>"
    )
    slide.element.append(timing)


def _append_transition_xml(slide):
    slide.element.append(
        parse_xml(f"<p:transition {nsdecls('p')}><p:fade/></p:transition>")
    )


def _append_build_list_xml(slide):
    """Append real PresentationML build entries without inferring playback."""
    timing = slide.element.xpath("./p:timing")[0]
    timing.append(
        parse_xml(
            f"<p:bldLst {nsdecls('p')}>"
            '<p:bldP spid="2" grpId="0"/>'
            '<p:bldDgm spid="3" grpId="1"/>'
            '<p:bldOleChart spid="4" grpId="2"/>'
            '<p:bldGraphic spid="5" grpId="3"/>'
            "</p:bldLst>"
        )
    )


def test_native_timing_categories_and_deck_totals_stay_distinct(
    pptx_extraction, tmp_path
):
    prs = Presentation()
    animated = prs.slides.add_slide(prs.slide_layouts[6])
    _append_transition_xml(animated)
    _append_timing_xml(
        animated,
        (
            '<p:set><p:cBhvr><p:cTn id="2"/><p:attrNameLst>'
            "<p:attrName>style.visibility</p:attrName>"
            '</p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            '<p:set><p:cBhvr><p:cTn id="3"/><p:attrNameLst>'
            "<p:attrName>style.opacity</p:attrName>"
            '</p:attrNameLst></p:cBhvr><p:to><p:fltVal val="1"/></p:to></p:set>'
            "<p:anim/><p:animClr/><p:animEffect/><p:animEffect/>"
            "<p:animMotion/><p:animRot/><p:animScale/>"
            '<p:audio><p:cMediaNode><p:cTn id="4"/></p:cMediaNode></p:audio>'
            '<p:video><p:cMediaNode><p:cTn id="5"/></p:cMediaNode></p:video>'
        ),
    )
    _append_build_list_xml(animated)

    media_only = prs.slides.add_slide(prs.slide_layouts[6])
    _append_timing_xml(
        media_only, ('<p:audio><p:cMediaNode><p:cTn id="6"/></p:cMediaNode></p:audio>')
    )

    transition_only = prs.slides.add_slide(prs.slide_layouts[6])
    _append_transition_xml(transition_only)

    path = tmp_path / "timing-structure.pptx"
    prs.save(path)
    result = pptx_extraction._extract_pptx_in_process(str(path), ocr=False)
    slides = result["per_slide_visual"]

    first = slides[0]["native_timing"]
    assert first["timing_element_present"] is True
    assert first["timing_element_count"] == 1
    assert first["transition_count"] == 1
    assert first["set_action_count"] == 2
    assert first["visibility_set_action_count"] == 1
    assert first["animation_behavior_counts"] == {
        "general": 1,
        "color": 1,
        "effect": 2,
        "motion": 1,
        "rotation": 1,
        "scale": 1,
        "total": 7,
    }
    assert first["media_timing_counts"] == {"audio": 1, "video": 1, "total": 2}
    assert first["build_list_present"] is True
    assert first["build_list_count"] == 1
    assert first["build_entry_counts"] == {
        "paragraph": 1,
        "diagram": 1,
        "ole_chart": 1,
        "graphic": 1,
        "total": 4,
    }
    assert first["has_build_entries"] is True
    assert first["has_animation_behaviors"] is True
    assert first["has_media_timing"] is True
    assert first["provenance"] == {
        "source": "pptx_package_xml",
        "measurement": "raw_ooxml_element_counts",
        "observed_playback": False,
        "part_name": "ppt/slides/slide1.xml",
    }

    # A timing container carrying only embedded-media timing is not generic
    # animation and especially not a motion observation.
    second = slides[1]["native_timing"]
    assert second["timing_element_present"] is True
    assert second["has_animation_behaviors"] is False
    assert second["animation_behavior_counts"]["motion"] == 0
    assert second["animation_behavior_counts"]["total"] == 0
    assert second["has_media_timing"] is True
    assert second["media_timing_counts"] == {"audio": 1, "video": 0, "total": 1}
    assert second["build_list_present"] is False
    assert second["build_list_count"] == 0
    assert second["build_entry_counts"] == {
        "paragraph": 0,
        "diagram": 0,
        "ole_chart": 0,
        "graphic": 0,
        "total": 0,
    }
    assert second["has_build_entries"] is False

    third = slides[2]["native_timing"]
    assert third["timing_element_present"] is False
    assert third["transition_count"] == 1
    assert third["has_animation_behaviors"] is False
    assert third["has_media_timing"] is False

    assert result["native_timing_summary"] == {
        "slides_with_timing_elements": 2,
        "slides_with_transitions": 2,
        "slides_with_animation_behaviors": 1,
        "slides_with_media_timing": 2,
        "slides_with_build_lists": 1,
        "slides_with_build_entries": 1,
        "timing_element_count": 2,
        "transition_count": 2,
        "set_action_count": 2,
        "visibility_set_action_count": 1,
        "build_list_count": 1,
        "animation_behavior_counts": {
            "general": 1,
            "color": 1,
            "effect": 2,
            "motion": 1,
            "rotation": 1,
            "scale": 1,
            "total": 7,
        },
        "media_timing_counts": {"audio": 2, "video": 1, "total": 3},
        "build_entry_counts": {
            "paragraph": 1,
            "diagram": 1,
            "ole_chart": 1,
            "graphic": 1,
            "total": 4,
        },
        "provenance": {
            "source": "pptx_package_xml",
            "measurement": "raw_ooxml_element_counts",
            "observed_playback": False,
        },
    }


def test_adjacent_static_progressive_builds_do_not_invent_native_timing(
    pptx_extraction, tmp_path
):
    """Duplicate-slide builds stay visible states, not inferred animation."""
    prs = Presentation()
    first = prs.slides.add_slide(prs.slide_layouts[6])
    first.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.text = "Base diagram"
    second = prs.slides.add_slide(prs.slide_layouts[6])
    second.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.text = "Base diagram"
    second.shapes.add_textbox(
        Inches(1), Inches(2), Inches(5), Inches(1)
    ).text_frame.text = "Step 2 annotation"
    path = tmp_path / "static-progressive-build.pptx"
    prs.save(path)

    result = pptx_extraction._extract_pptx_in_process(str(path), ocr=False)

    assert "Base diagram" in result["per_slide_visual"][0]["text_content_preview"]
    assert "Step 2 annotation" in result["per_slide_visual"][1]["text_content_preview"]
    assert all(
        not slide["native_timing"]["timing_element_present"]
        for slide in result["per_slide_visual"]
    )
    assert result["native_timing_summary"]["animation_behavior_counts"]["total"] == 0
    assert result["native_timing_summary"]["media_timing_counts"]["total"] == 0
    assert all(
        slide["native_timing"]["build_list_present"] is False
        and slide["native_timing"]["build_entry_counts"]["total"] == 0
        for slide in result["per_slide_visual"]
    )
    assert result["native_timing_summary"]["build_entry_counts"]["total"] == 0


def test_versions_and_input_fingerprint_are_stable_and_content_addressed(
    pptx_extraction,
    tmp_path,
):
    prs = make_deck(1)
    source = tmp_path / "source.pptx"
    prs.save(source)
    copy = tmp_path / "copy.pptx"
    copy.write_bytes(source.read_bytes())
    modified = tmp_path / "modified.pptx"
    modified.write_bytes(source.read_bytes() + b"\x00")

    first = pptx_extraction._extract_pptx_in_process(str(source), ocr=False)
    second = pptx_extraction._extract_pptx_in_process(str(source), ocr=False)
    copied = pptx_extraction._extract_pptx_in_process(str(copy), ocr=False)
    changed = pptx_extraction._extract_pptx_in_process(str(modified), ocr=False)

    assert pptx_extraction.SCHEMA_VERSION == 4
    assert pptx_extraction.PIPELINE_VERSION == "1.5.0"
    assert first["schema_version"] == pptx_extraction.SCHEMA_VERSION
    assert first["pipeline_version"] == pptx_extraction.PIPELINE_VERSION
    assert first["input_fingerprint"] == second["input_fingerprint"]
    assert first["input_fingerprint"] == copied["input_fingerprint"]
    assert first["input_fingerprint"] != changed["input_fingerprint"]
    assert first["input_fingerprint"]["algorithm"] == "sha256"
    assert len(first["input_fingerprint"]["digest"]) == 64
    assert first["input_fingerprint"]["size_bytes"] == source.stat().st_size


def test_version_flag_is_machine_readable(pptx_extraction):
    proc = subprocess.run(
        [sys.executable, pptx_extraction.__file__, "--version"],
        capture_output=True,
        text=True,
        check=True,
        timeout=30,
    )
    assert json.loads(proc.stdout) == {
        "schema_version": pptx_extraction.SCHEMA_VERSION,
        "pipeline_version": pptx_extraction.PIPELINE_VERSION,
    }


_PAGE_DECIMAL_ALPHABETS = (
    "0123456789",
    "٠١٢٣٤٥٦٧٨٩",
    "۰۱۲۳۴۵۶۷۸۹",
    "０１２３４５６７８９",
    "𝟘𝟙𝟚𝟛𝟜𝟝𝟞𝟟𝟠𝟡",
)
_PAGE_DECIMAL_ALPHABET_IDS = (
    "ascii",
    "arabic-indic",
    "extended-arabic-indic",
    "fullwidth",
    "mathematical",
)
_INSPECTED_PAGES_GRAMMAR_ERROR = "--inspected-pages values must be PAGE or START-END"


def _render_page_decimal(value, alphabet):
    return "".join(alphabet[int(digit)] for digit in str(value))


def _render_mixed_page_decimal(value):
    return "".join(
        _PAGE_DECIMAL_ALPHABETS[index % len(_PAGE_DECIMAL_ALPHABETS)][int(digit)]
        for index, digit in enumerate(str(value))
    )


class _NoSplitOrStripPageRange(str):
    def split(self, *_args, **_kwargs):
        raise AssertionError("page-range parsing must not split its input")

    def strip(self, *_args, **_kwargs):
        raise AssertionError("page-range parsing must not copy a stripped token")


@pytest.mark.parametrize(
    "alphabet",
    _PAGE_DECIMAL_ALPHABETS,
    ids=_PAGE_DECIMAL_ALPHABET_IDS,
)
def test_page_range_parser_supports_unicode_decimal_digits_at_global_ceiling(
    pptx_extraction,
    pptx_evidence,
    alphabet,
):
    limit = pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS
    limit_error = (
        "--inspected-pages page numbers must not exceed the bounded PPTX "
        f"page limit {limit}"
    )

    assert pptx_extraction.parse_page_range_arguments(
        [f"{_render_page_decimal(1, alphabet)}-{_render_page_decimal(3, alphabet)}"]
    ) == [[1, 3]]
    assert pptx_extraction.parse_page_range_arguments(
        [_render_page_decimal(limit, alphabet)]
    ) == [[limit, limit]]

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction.parse_page_range_arguments(
            [_render_page_decimal(limit + 1, alphabet)]
        )
    assert str(caught.value) == limit_error


def test_page_range_parser_supports_mixed_script_decimal_numbers(
    pptx_extraction,
    pptx_evidence,
):
    limit = pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS
    limit_error = (
        "--inspected-pages page numbers must not exceed the bounded PPTX "
        f"page limit {limit}"
    )

    assert pptx_extraction.parse_page_range_arguments(
        [f"{_render_mixed_page_decimal(12345)}-{_render_mixed_page_decimal(12347)}"]
    ) == [[12345, 12347]]
    assert pptx_extraction.parse_page_range_arguments(
        [_render_mixed_page_decimal(limit)]
    ) == [[limit, limit]]

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction.parse_page_range_arguments(
            [_render_mixed_page_decimal(limit + 1)]
        )
    assert str(caught.value) == limit_error


def test_page_range_parser_rejects_local_impossibilities_with_fixed_errors(
    pptx_extraction,
    pptx_evidence,
):
    limit = pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS
    cases = (
        ("0", "--inspected-pages page numbers must be at least 1"),
        ("0-1", "--inspected-pages page numbers must be at least 1"),
        (
            "5-3",
            "--inspected-pages range end must not be less than its start",
        ),
        (
            str(limit + 1),
            "--inspected-pages page numbers must not exceed the bounded PPTX "
            f"page limit {limit}",
        ),
        (
            "9" * 5_000,
            "--inspected-pages page numbers must not exceed the bounded PPTX "
            f"page limit {limit}",
        ),
    )

    for value, expected_message in cases:
        with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
            pptx_extraction.parse_page_range_arguments([value])
        assert str(caught.value) == expected_message
        assert value not in str(caught.value)


def test_page_range_parser_bounds_malformed_diagnostics_without_echoing_input(
    pptx_extraction,
):
    malformed = "x" + ("９" * 5_000)

    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction.parse_page_range_arguments([malformed])

    diagnostic = str(caught.value)
    assert diagnostic == _INSPECTED_PAGES_GRAMMAR_ERROR
    assert malformed not in diagnostic
    assert len(diagnostic.encode("utf-8")) < 128


def test_page_range_parser_enforces_range_count_without_split_amplification(
    pptx_extraction,
    pptx_evidence,
):
    limit = pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS
    exact_limit = _NoSplitOrStripPageRange(("1," * (limit - 1)) + "1")
    parsed = pptx_extraction.parse_page_range_arguments([exact_limit])

    assert len(parsed) == limit
    assert parsed[0] == parsed[-1] == [1, 1]

    malformed = "x" + ("９" * 5_000)
    limit_plus_one = _NoSplitOrStripPageRange(f"{exact_limit},{malformed}")
    with pytest.raises(pptx_extraction.PptxEvidenceError) as caught:
        pptx_extraction.parse_page_range_arguments([limit_plus_one])

    diagnostic = str(caught.value)
    assert diagnostic == (f"--inspected-pages must contain no more than {limit} ranges")
    assert malformed not in diagnostic
    assert len(diagnostic.encode("utf-8")) < 128


def test_page_range_parser_preserves_repeated_comma_and_leading_zero_inputs(
    pptx_extraction,
):
    leading_zero_page = ("0" * 5_000) + "1"
    leading_zero_end = ("０" * 5_000) + "３"
    canonical = pptx_extraction.parse_page_range_arguments(["1", "2-4,6", "8-9"])
    equivalent = pptx_extraction.parse_page_range_arguments(
        [leading_zero_page, "0002-０００４,٠٠٠٦", "０００８-0009"]
    )

    assert canonical == equivalent == [[1, 1], [2, 4], [6, 6], [8, 9]]
    assert pptx_extraction.parse_page_range_arguments(
        [f"{leading_zero_page}-{leading_zero_end}"]
    ) == [[1, 3]]


def test_unicode_ranges_have_canonical_normalized_json_byte_equivalence(
    pptx_extraction,
    pptx_evidence,
):
    limit = pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS
    canonical_ranges = pptx_evidence.normalize_page_ranges(
        pptx_extraction.parse_page_range_arguments(["1-3", "5", str(limit)]),
        page_count=limit,
        allow_empty=True,
    )
    unicode_ranges = pptx_evidence.normalize_page_ranges(
        pptx_extraction.parse_page_range_arguments(
            [
                " ١-۳, ５ ",
                _render_mixed_page_decimal(limit),
            ]
        ),
        page_count=limit,
        allow_empty=True,
    )

    canonical_bytes = json.dumps(
        canonical_ranges,
        ensure_ascii=False,
        separators=(",", ":"),
    ).encode("utf-8")
    unicode_bytes = json.dumps(
        unicode_ranges,
        ensure_ascii=False,
        separators=(",", ":"),
    ).encode("utf-8")
    assert unicode_bytes == canonical_bytes


def test_page_range_parser_leaves_deck_and_cross_range_checks_to_normalizer(
    pptx_extraction,
    pptx_evidence,
):
    deck_overflow = pptx_extraction.parse_page_range_arguments(["10"])
    overlapping = pptx_extraction.parse_page_range_arguments(["2-3", "3-4"])

    with pytest.raises(pptx_evidence.PptxEvidenceError):
        pptx_evidence.normalize_page_ranges(
            deck_overflow,
            page_count=9,
            allow_empty=True,
        )
    with pytest.raises(pptx_evidence.PptxEvidenceError):
        pptx_evidence.normalize_page_ranges(
            overlapping,
            page_count=4,
            allow_empty=True,
        )


def test_cli_rejects_impossible_page_ranges_before_artifact_work(
    pptx_extraction,
    pptx_evidence,
    monkeypatch,
    capsys,
):
    limit = pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS
    artifact_calls = []

    def forbidden_artifact_work(*args, **kwargs):
        artifact_calls.append((args, kwargs))
        pytest.fail("invalid page range reached PPTX artifact work")

    monkeypatch.setattr(pptx_extraction, "extract_pptx", forbidden_artifact_work)
    cases = (
        ("0", "--inspected-pages page numbers must be at least 1"),
        (
            "5-3",
            "--inspected-pages range end must not be less than its start",
        ),
        (
            _render_mixed_page_decimal(limit + 1),
            "--inspected-pages page numbers must not exceed the bounded PPTX "
            f"page limit {limit}",
        ),
    )

    for value, expected_message in cases:
        with pytest.raises(SystemExit) as raised:
            pptx_extraction.main(
                [
                    "/unread/deck.pptx",
                    "--rendered-pdf",
                    "/unread/deck.pdf",
                    "--inspected-pages",
                    value,
                    "--no-ocr",
                ]
            )
        assert raised.value.code == 2
        captured = capsys.readouterr()
        assert captured.out == ""
        assert expected_message in captured.err
        assert "Traceback" not in captured.err
        assert value not in captured.err

    assert artifact_calls == []


def test_cli_bounds_resource_sized_page_range_errors_before_artifact_work(
    pptx_extraction,
    pptx_evidence,
    monkeypatch,
    capsys,
):
    limit = pptx_evidence.PPTX_ARCHIVE_MAX_MEMBERS
    artifact_calls = []
    malformed = "x" + ("９" * 5_000)
    too_many = ("１," * limit) + malformed

    def forbidden_artifact_work(*args, **kwargs):
        artifact_calls.append((args, kwargs))
        pytest.fail("resource-sized page range reached PPTX artifact work")

    monkeypatch.setattr(pptx_extraction, "extract_pptx", forbidden_artifact_work)
    cases = (
        (malformed, _INSPECTED_PAGES_GRAMMAR_ERROR),
        (
            too_many,
            f"--inspected-pages must contain no more than {limit} ranges",
        ),
    )

    for value, expected_message in cases:
        with pytest.raises(SystemExit) as raised:
            pptx_extraction.main(
                [
                    "/unread/deck.pptx",
                    "--rendered-pdf",
                    "/unread/deck.pdf",
                    "--inspected-pages",
                    value,
                    "--no-ocr",
                ]
            )
        assert raised.value.code == 2
        captured = capsys.readouterr()
        assert captured.out == ""
        assert expected_message in captured.err
        assert malformed not in captured.err
        assert "Traceback" not in captured.err
        assert len(captured.err.encode("utf-8")) < 4_096

    assert artifact_calls == []


def test_cli_normalizes_long_leading_zero_page_range_before_artifact_work(
    pptx_extraction,
    monkeypatch,
    capsys,
):
    artifact_calls = []
    leading_zero_start = ("٠" * 5_000) + "１"
    leading_zero_end = ("۰" * 5_000) + "𝟛"

    def capture_artifact_work(path, **options):
        artifact_calls.append((path, options))
        return {"slide_count": 3}

    monkeypatch.setattr(pptx_extraction, "extract_pptx", capture_artifact_work)

    assert (
        pptx_extraction.main(
            [
                "/unread/deck.pptx",
                "--rendered-pdf",
                "/unread/deck.pdf",
                "--inspected-pages",
                f"{leading_zero_start}-{leading_zero_end}",
                "--no-ocr",
            ]
        )
        == 0
    )

    captured = capsys.readouterr()
    assert captured.err == ""
    assert json.loads(captured.out) == {"slide_count": 3}
    assert artifact_calls == [
        (
            "/unread/deck.pptx",
            {
                "ocr": False,
                "rendered_pdf_path": "/unread/deck.pdf",
                "inspected_page_ranges": [[1, 3]],
            },
        )
    ]


def test_directory_worker_main_reports_closed_supervisor_failure(
    pptx_extraction,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    def fail():
        raise pptx_extraction.SupervisorError("worker_output_limit_exceeded")

    monkeypatch.setattr(pptx_extraction, "_run_directory_worker_child", fail)
    monkeypatch.setattr(
        pptx_extraction.sys,
        "argv",
        [pptx_extraction.__file__, pptx_extraction._DIRECTORY_WORKER_FLAG],
    )

    assert pptx_extraction._main() == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert (
        captured.err == "pptx directory worker failed: worker_output_limit_exceeded\n"
    )
    assert "Traceback" not in captured.err


def test_directory_worker_main_closes_unexpected_failure_diagnostic(
    pptx_extraction,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    leaked_path = "/private/vault/source.pptx"

    def fail():
        raise RuntimeError(f"failure at {leaked_path}")

    monkeypatch.setattr(pptx_extraction, "_run_directory_worker_child", fail)
    monkeypatch.setattr(
        pptx_extraction.sys,
        "argv",
        [pptx_extraction.__file__, pptx_extraction._DIRECTORY_WORKER_FLAG],
    )

    assert pptx_extraction._main() == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == "pptx directory worker failed: unexpected_error\n"
    assert leaked_path not in captured.err
    assert "Traceback" not in captured.err


def test_directory_worker_main_preserves_success_output_contract(
    pptx_extraction,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    monkeypatch.setattr(pptx_extraction, "_run_directory_worker_child", lambda: 0)
    monkeypatch.setattr(
        pptx_extraction.sys,
        "argv",
        [pptx_extraction.__file__, pptx_extraction._DIRECTORY_WORKER_FLAG],
    )

    assert pptx_extraction._main() == 0
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == ""
