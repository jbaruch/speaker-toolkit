"""Tests for insert-placeholder-slides.py — placeholder slide insertion."""

import json
import os
import subprocess
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor

from conftest import make_deck, SCRIPTS_PC

SCRIPT = os.path.abspath(os.path.join(SCRIPTS_PC, "insert-placeholder-slides.py"))


def _slide_texts(slide):
    return [s.text_frame.text for s in slide.shapes if s.has_text_frame]


def test_yellow_background(insert_placeholder, tmp_path):
    prs = make_deck(3)
    slide = insert_placeholder.add_placeholder_slide(prs, "Test Title")
    # Check that a shape has the yellow fill
    found_yellow = False
    for shape in slide.shapes:
        try:
            if shape.fill.fore_color.rgb == RGBColor(0xFF, 0xF2, 0x9E):
                found_yellow = True
                break
        except Exception:
            continue
    assert found_yellow


def test_placeholder_title_text(insert_placeholder, tmp_path):
    prs = make_deck(3)
    slide = insert_placeholder.add_placeholder_slide(prs, "My Section")
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    assert any("[PLACEHOLDER] My Section" in t for t in texts)


def test_subtitle_present(insert_placeholder, tmp_path):
    prs = make_deck(3)
    slide = insert_placeholder.add_placeholder_slide(prs, "Title", "Subtitle text")
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    assert any("Subtitle text" in t for t in texts)


def test_multi_insert_positioning(insert_placeholder, tmp_path):
    prs = make_deck(3)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    # Insert at positions 1 and 3 (1-indexed, final positions)
    insert_placeholder.add_placeholder_slide(prs2, "Second insert")
    insert_placeholder.move_slide(prs2, len(prs2.slides) - 1, 2)  # pos 3 (0-indexed=2)
    insert_placeholder.add_placeholder_slide(prs2, "First insert")
    insert_placeholder.move_slide(prs2, len(prs2.slides) - 1, 0)  # pos 1 (0-indexed=0)

    assert len(prs2.slides) == 5


def test_slide_count_increases(insert_placeholder, tmp_path):
    prs = make_deck(3)
    insert_placeholder.add_placeholder_slide(prs, "New Slide")
    assert len(prs.slides) == 4


def test_move_slide(insert_placeholder, tmp_path):
    prs = make_deck(3)
    # Add a placeholder at the end, then move to front
    insert_placeholder.add_placeholder_slide(prs, "Should be first")
    insert_placeholder.move_slide(prs, 3, 0)

    # Verify the slide was moved — the new first slide should have our text
    first_slide = prs.slides[0]
    texts = []
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    assert any("[PLACEHOLDER] Should be first" in t for t in texts)


def test_output_flag(insert_placeholder, tmp_path):
    """Saving to a different output path preserves the original."""
    prs = make_deck(3)
    original = str(tmp_path / "original.pptx")
    prs.save(original)

    prs2 = Presentation(original)
    insert_placeholder.add_placeholder_slide(prs2, "Added")
    output = str(tmp_path / "modified.pptx")
    prs2.save(output)

    # Original unchanged
    assert len(Presentation(original).slides) == 3
    # Output has the new slide
    assert len(Presentation(output).slides) == 4


def test_title_not_double_prefixed(insert_placeholder):
    """A title that already starts with '[PLACEHOLDER] ' must not be re-prefixed."""
    prs = make_deck(1)
    slide = insert_placeholder.add_placeholder_slide(prs, "[PLACEHOLDER] Foo")
    texts = _slide_texts(slide)
    assert any(t == "[PLACEHOLDER] Foo" for t in texts)
    assert not any("[PLACEHOLDER] [PLACEHOLDER]" in t for t in texts)


def test_format_title_helper(insert_placeholder):
    """The title formatter prepends the prefix only when missing."""
    f = insert_placeholder._format_title
    assert f("Foo") == "[PLACEHOLDER] Foo"
    assert f("[PLACEHOLDER] Foo") == "[PLACEHOLDER] Foo"
    # Case-insensitive: already-prefixed in lowercase should not double-prefix.
    assert f("[placeholder] foo") == "[placeholder] foo"


def test_format_title_preserves_leading_whitespace(insert_placeholder):
    """Leading whitespace must be preserved consistently regardless of whether
    the prefix is already present — earlier behavior stripped only on the
    already-prefixed branch."""
    f = insert_placeholder._format_title
    # Already prefixed with leading whitespace — return as-is, don't strip.
    assert f("  [PLACEHOLDER] Foo") == "  [PLACEHOLDER] Foo"
    # Not prefixed with leading whitespace — prepend, don't strip the input.
    assert f("  Foo") == "[PLACEHOLDER]   Foo"


def _run_cli(deck_path, json_path):
    """Invoke the script as a subprocess so argparse and stderr behave realistically."""
    return subprocess.run(
        [sys.executable, SCRIPT, deck_path, json_path],
        capture_output=True,
        text=True,
    )


def test_cli_positions_land_at_declared_finals(tmp_path):
    """Regression for #18: high positions must land at their declared final positions,
    not silently append because list.insert clamped an out-of-range index."""
    prs = make_deck(8)
    deck = str(tmp_path / "deck.pptx")
    prs.save(deck)

    # Original length 8, 4 placeholders at final positions 2, 6, 11, 12 (total_after=12).
    # Pre-fix with high-to-low iteration, positions 11 and 12 landed at wrong indexes.
    spec = [
        {"position": 2,  "title": "A", "subtitle": ""},
        {"position": 6,  "title": "B", "subtitle": ""},
        {"position": 11, "title": "C", "subtitle": ""},
        {"position": 12, "title": "D", "subtitle": ""},
    ]
    spec_path = str(tmp_path / "spec.json")
    with open(spec_path, "w") as f:
        json.dump(spec, f)

    result = _run_cli(deck, spec_path)
    assert result.returncode == 0, result.stderr

    final = Presentation(deck)
    assert len(final.slides) == 12

    # Each placeholder's title should appear at its declared 1-indexed position.
    for entry in spec:
        idx = entry["position"] - 1
        texts = _slide_texts(final.slides[idx])
        assert any(f"[PLACEHOLDER] {entry['title']}" in t for t in texts), (
            f"expected '{entry['title']}' at position {entry['position']}, "
            f"got texts {texts}"
        )


def test_cli_rejects_duplicate_positions(tmp_path):
    prs = make_deck(3)
    deck = str(tmp_path / "deck.pptx")
    prs.save(deck)

    spec_path = str(tmp_path / "spec.json")
    with open(spec_path, "w") as f:
        json.dump([
            {"position": 2, "title": "A", "subtitle": ""},
            {"position": 2, "title": "B", "subtitle": ""},
        ], f)

    result = _run_cli(deck, spec_path)
    assert result.returncode != 0
    assert "duplicate position" in result.stderr


def test_cli_rejects_out_of_range(tmp_path):
    prs = make_deck(3)
    deck = str(tmp_path / "deck.pptx")
    prs.save(deck)

    spec_path = str(tmp_path / "spec.json")
    with open(spec_path, "w") as f:
        # Original 3 + 1 placeholder = total_after 4; position 99 is out of range.
        json.dump([{"position": 99, "title": "A", "subtitle": ""}], f)

    result = _run_cli(deck, spec_path)
    assert result.returncode != 0
    assert "out of range" in result.stderr
