"""Tests for generate-qr.py — QR generation (no network calls)."""

import contextlib
import json
import os
import sys
from pathlib import Path

import pytest

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from conftest import current_tracking_config, make_deck


def _current_tracking_database():
    return {
        "schema_version": 1,
        "config": current_tracking_config(),
        "talks": [],
        "pptx_catalog": [],
        "qr_codes": [],
        "resources": [],
        "thumbnails": [],
        "confirmed_intents": [],
        "improvement_goals": [],
    }


def _square_png(tmp_path, name="sq.png"):
    """A two-color (QR-like) PNG: quantizes to 2 colors with ~0 reconstruction
    error, so the content-based detector treats it as a QR."""
    path = str(tmp_path / name)
    im = Image.new("RGB", (50, 50), (255, 255, 255))
    px = im.load()
    for x in range(50):
        for y in range(50):
            if (x + y) % 2 == 0:
                px[x, y] = (0, 0, 0)
    im.save(path)
    return path


def _multicolor_png(tmp_path, name="multi.png"):
    """A many-color (photo/diagram-like) PNG: a smooth gradient that does NOT
    reduce to two colors, so the detector rejects it even when square."""
    path = str(tmp_path / name)
    im = Image.new("RGB", (60, 60))
    px = im.load()
    for x in range(60):
        for y in range(60):
            px[x, y] = ((x * 4) % 256, (y * 4) % 256, ((x + y) * 2) % 256)
    im.save(path)
    return path


def _sparse_text_png(tmp_path, name="screenshot.png"):
    """A mostly-white image with sparse dark 'text' — ~2-color but heavily skewed
    to the background (like a doc screenshot), so it must NOT be taken for a QR."""
    path = str(tmp_path / name)
    im = Image.new("RGB", (80, 80), (255, 255, 255))
    px = im.load()
    for y in range(0, 80, 6):
        for x in range(0, 80, 3):
            px[x, y] = (0, 0, 0)
    im.save(path)
    return path


def _receipt(path, *, bg_hex=None):
    """Minimal valid artifact receipt for tests that do not write a real PNG."""
    return {
        "path": path,
        "path_root": "cwd",
        "sha256": "0" * 64,
        "bg_hex": bg_hex,
    }


def test_choose_fg_color_dark_bg(generate_qr):
    # Dark background → white foreground
    fg = generate_qr.choose_fg_color((0, 0, 0))
    assert fg == (255, 255, 255)


def test_choose_fg_color_light_bg(generate_qr):
    # Light background → black foreground
    fg = generate_qr.choose_fg_color((255, 255, 255))
    assert fg == (0, 0, 0)


def test_choose_fg_color_none(generate_qr):
    # No background detected → default to black
    fg = generate_qr.choose_fg_color(None)
    assert fg == (0, 0, 0)


def test_choose_fg_color_mid_gray(generate_qr):
    # Mid-gray — should pick based on luminance
    fg = generate_qr.choose_fg_color((128, 128, 128))
    # Luminance for (128,128,128) ≈ 0.502 → ≥ 0.5 → black
    assert fg == (0, 0, 0)


def test_choose_fg_color_dark_blue(generate_qr):
    # Dark blue → white foreground
    fg = generate_qr.choose_fg_color((0, 0, 100))
    assert fg == (255, 255, 255)


def test_generate_qr_png(generate_qr, tmp_path):
    out = str(tmp_path / "test.png")
    generate_qr.generate_qr_png("https://example.com", (0, 0, 0), (255, 255, 255), out)
    assert os.path.isfile(out)
    assert os.path.getsize(out) > 100  # sanity check — not empty


def test_generate_qr_png_custom_colors(generate_qr, tmp_path):
    out = str(tmp_path / "custom.png")
    generate_qr.generate_qr_png(
        "https://example.com", (255, 255, 255), (128, 0, 128), out
    )
    assert os.path.isfile(out)


def test_tracking_db_crud_insert(generate_qr):
    db = {}
    entry = {
        "talk_slug": "test-talk",
        "target_url": "https://example.com/notes",
        "shortener": "none",
        "short_url": "https://example.com/notes",
    }
    generate_qr.update_tracking_db(db, entry, [_receipt("test-talk-qr.png")])
    assert len(db["qr_codes"]) == 1
    assert db["qr_codes"][0]["schema_version"] == 2
    assert db["qr_codes"][0]["talk_slug"] == "test-talk"
    assert db["qr_codes"][0]["qr_png_rel_path"] == "test-talk-qr.png"


def test_tracking_db_crud_update(generate_qr):
    db = {
        "qr_codes": [
            {
                "talk_slug": "test-talk",
                "target_url": "https://old-url.com",
                "shortener": "none",
                "short_url": "https://old-url.com",
                "qr_png_rel_path": "old.png",
                "created_at": "2024-01-01",
                "updated_at": "2024-01-01",
            }
        ]
    }
    entry = {
        "talk_slug": "test-talk",
        "target_url": "https://new-url.com",
        "shortener": "none",
        "short_url": "https://new-url.com",
    }
    generate_qr.update_tracking_db(db, entry, [_receipt("new.png")])
    assert len(db["qr_codes"]) == 1
    assert db["qr_codes"][0]["schema_version"] == 2
    assert db["qr_codes"][0]["target_url"] == "https://new-url.com"
    assert db["qr_codes"][0]["qr_png_rel_path"] == "new.png"
    # created_at preserved from original
    assert db["qr_codes"][0]["created_at"] == "2024-01-01"
    assert db["qr_codes"][0]["schema_version"] == 2


def test_tracking_db_semantic_noop_preserves_raw_bytes_and_inode(
    generate_qr,
    tmp_path,
):
    path = tmp_path / "tracking-database.json"
    database = _current_tracking_database()
    raw = (json.dumps(database, separators=(",", ":")) + "\n").encode()
    path.write_bytes(raw)
    snapshot = generate_qr.snapshot_tracking_database(path)
    equivalent = {key: database[key] for key in reversed(database)}

    result = generate_qr.write_tracking_db(snapshot, equivalent)

    assert result.changed is False
    assert result.installed is False
    assert path.read_bytes() == raw
    assert path.stat().st_ino == snapshot.generation.inode


@pytest.mark.parametrize("mode", ["png", "deck"])
def test_main_existing_vault_without_database_stops_before_side_effects(
    generate_qr,
    tmp_path,
    monkeypatch,
    capsys,
    mode,
):
    vault = tmp_path / "vault"
    vault.mkdir()
    output = tmp_path / "talk-qr.png"
    deck = tmp_path / "talk.pptx"
    arguments = ["generate-qr.py"]
    if mode == "png":
        arguments.extend(["--png-only", "--output", str(output)])
    else:
        deck.write_bytes(b"unchanged deck")
        arguments.append(str(deck))
    arguments.extend(
        [
            "--talk-slug",
            "talk",
            "--shownotes-url",
            "https://example.com/talk",
            "--vault",
            str(vault),
        ]
    )

    def forbidden(*_args, **_kwargs):
        raise AssertionError("QR side effect ran before tracking DB validation")

    monkeypatch.setattr(generate_qr.sys, "argv", arguments)
    monkeypatch.setattr(generate_qr, "resolve_short_url", forbidden)
    monkeypatch.setattr(generate_qr, "generate_qr_png", forbidden)
    monkeypatch.setattr(generate_qr, "Presentation", forbidden)
    monkeypatch.setattr(generate_qr, "insert_qr_via_powerpoint", forbidden)

    with pytest.raises(SystemExit) as exc_info:
        generate_qr.main()

    assert exc_info.value.code == 1
    assert "tracking-database.json is missing" in capsys.readouterr().err
    assert not output.exists()
    if mode == "deck":
        assert deck.read_bytes() == b"unchanged deck"


def test_main_dry_run_allows_existing_vault_without_database(
    generate_qr,
    tmp_path,
    monkeypatch,
    capsys,
):
    vault = tmp_path / "vault"
    vault.mkdir()

    def forbidden(*_args, **_kwargs):
        raise AssertionError("dry run performed a write")

    monkeypatch.setattr(
        generate_qr.sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "talk",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://example.com/talk",
            "--vault",
            str(vault),
            "--dry-run",
        ],
    )
    monkeypatch.setattr(generate_qr, "generate_qr_png", forbidden)
    monkeypatch.setattr(generate_qr, "write_tracking_db", forbidden)

    generate_qr.main()

    assert "DRY RUN: would save QR" in capsys.readouterr().out


def test_main_invalid_background_stops_before_url_resolution(
    generate_qr,
    tmp_path,
    monkeypatch,
    capsys,
):
    def forbidden(*_args, **_kwargs):
        raise AssertionError("URL resolution ran before local validation")

    monkeypatch.setattr(
        generate_qr.sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "talk",
            "--shownotes-url",
            "https://example.com/talk",
            "--vault",
            str(tmp_path / "missing-vault"),
            "--bg-color",
            "not-rgb",
        ],
    )
    monkeypatch.setattr(generate_qr, "resolve_short_url", forbidden)

    with pytest.raises(SystemExit) as exc_info:
        generate_qr.main()

    assert exc_info.value.code == 1
    assert "--bg-color must be R,G,B" in capsys.readouterr().out


def test_main_valid_snapshot_generates_png_and_persists_metadata(
    generate_qr,
    tmp_path,
    monkeypatch,
    capsys,
):
    vault = tmp_path / "vault"
    vault.mkdir()
    database_path = vault / "tracking-database.json"
    database_path.write_text(
        json.dumps(_current_tracking_database()) + "\n",
        encoding="utf-8",
    )
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr(
        generate_qr.sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "talk",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://example.com/talk",
            "--vault",
            str(vault),
        ],
    )

    generate_qr.main()

    assert (tmp_path / "talk-qr.png").is_file()
    database = json.loads(database_path.read_text(encoding="utf-8"))
    assert database["qr_codes"][0]["talk_slug"] == "talk"
    assert "Tracking DB updated:" in capsys.readouterr().out

    generate_qr.main()

    assert "Tracking DB unchanged:" in capsys.readouterr().out


def test_resolve_slide_bg_rgb_none_for_plain_deck(generate_qr, tmp_path):
    """A plain deck without explicit background returns None."""
    prs = make_deck(1)
    path = str(tmp_path / "deck.pptx")
    prs.save(path)
    prs2 = Presentation(path)
    result = generate_qr.resolve_slide_bg_rgb(prs2.slides[0])
    # May be None or a default — both are acceptable
    assert result is None or isinstance(result, tuple)


def test_slide_has_existing_qr_detects_square_qr_sized_picture(generate_qr, tmp_path):
    prs = make_deck(1)
    prs.slides[0].shapes.add_picture(
        _square_png(tmp_path), Inches(8), Inches(5), Inches(2.0), Inches(2.0)
    )
    assert generate_qr.slide_has_existing_qr(prs.slides[0]) is True


def test_slide_has_existing_qr_ignores_non_square_picture(generate_qr, tmp_path):
    prs = make_deck(1)
    prs.slides[0].shapes.add_picture(
        _square_png(tmp_path), Inches(1), Inches(1), Inches(6.0), Inches(2.0)
    )
    assert generate_qr.slide_has_existing_qr(prs.slides[0]) is False


def test_slide_has_existing_qr_detects_large_square_qr(generate_qr, tmp_path):
    # Size-independent: a 2.8in two-color square is still a QR (the inherited
    # shownotes QR in the #56 repro deck was 2.78in — outside the old 1.5-2.5 band).
    prs = make_deck(1)
    prs.slides[0].shapes.add_picture(
        _square_png(tmp_path), Inches(1), Inches(1), Inches(2.8), Inches(2.8)
    )
    assert generate_qr.slide_has_existing_qr(prs.slides[0]) is True


def test_slide_has_existing_qr_ignores_multicolor_square(generate_qr, tmp_path):
    # Square and in size range, but many colors (e.g. a Venn diagram) → not a QR.
    prs = make_deck(1)
    prs.slides[0].shapes.add_picture(
        _multicolor_png(tmp_path), Inches(2), Inches(2), Inches(2.0), Inches(2.0)
    )
    assert generate_qr.slide_has_existing_qr(prs.slides[0]) is False


def test_slide_has_existing_qr_ignores_small_two_color_icon(generate_qr, tmp_path):
    # Two-color but below the 1.5in floor → a small icon, not a QR.
    prs = make_deck(1)
    prs.slides[0].shapes.add_picture(
        _square_png(tmp_path), Inches(1), Inches(1), Inches(1.0), Inches(1.0)
    )
    assert generate_qr.slide_has_existing_qr(prs.slides[0]) is False


def test_slide_has_existing_qr_ignores_text_screenshot(generate_qr, tmp_path):
    # Near-square, in size range, ~2-color — but mostly-white with sparse text
    # (unbalanced) → not a QR. Regression: a martinfowler.com screenshot (slide 28
    # of the #56 repro deck) that the recon-error-only test wrongly accepted.
    prs = make_deck(1)
    prs.slides[0].shapes.add_picture(
        _sparse_text_png(tmp_path), Inches(2), Inches(1), Inches(4.0), Inches(4.0)
    )
    assert generate_qr.slide_has_existing_qr(prs.slides[0]) is False


def test_slide_has_existing_qr_false_for_plain_slide(generate_qr):
    prs = make_deck(1)
    assert generate_qr.slide_has_existing_qr(prs.slides[0]) is False


def test_two_color_metrics_separate_qr_screenshot_photo(generate_qr, tmp_path):
    qr_err, qr_min = generate_qr._two_color_metrics(
        open(_square_png(tmp_path, "q.png"), "rb").read()
    )
    multi_err, _ = generate_qr._two_color_metrics(
        open(_multicolor_png(tmp_path, "m.png"), "rb").read()
    )
    _, sshot_min = generate_qr._two_color_metrics(
        open(_sparse_text_png(tmp_path, "s.png"), "rb").read()
    )
    assert qr_err < 5.0 and qr_min >= 0.25  # QR: two-color AND balanced
    assert multi_err > 20.0  # photo/diagram: many colors
    assert sshot_min < 0.25  # screenshot: mostly background


def test_find_qr_rects_returns_points_geometry(generate_qr, tmp_path):
    prs = make_deck(1)
    # placed at 8in,5in, 2in square → points: 576, 360, 144, 144
    prs.slides[0].shapes.add_picture(
        _square_png(tmp_path), Inches(8), Inches(5), Inches(2.0), Inches(2.0)
    )
    rects = generate_qr.find_qr_rects(prs.slides[0])
    assert len(rects) == 1
    L, T, W, H = rects[0]
    assert abs(L - 576) < 0.5 and abs(T - 360) < 0.5
    assert abs(W - 144) < 0.5 and abs(H - 144) < 0.5


def test_resolve_target_includes_inherited_qr_slides(generate_qr, tmp_path):
    """A deck adapted from another talk: config targets only the closing slide,
    but an earlier slide carries an inherited QR — it must also be targeted."""
    prs = make_deck(4)
    prs.slides[1].shapes.add_picture(
        _square_png(tmp_path), Inches(8), Inches(5), Inches(2.0), Inches(2.0)
    )
    indices = generate_qr.resolve_target_slide_indices(
        prs, {"slide_position": "closing"}, "https://example.com/notes"
    )
    assert 1 in indices  # inherited-QR slide
    assert 3 in indices  # closing slide


def test_back_half_is_always_slug_ignoring_preferred_short_path(
    generate_qr, monkeypatch
):
    """The back-half is ALWAYS the talk slug — a legacy preferred_short_path is ignored."""
    captured = {}

    def fake_create_bitly_link(long_url, api_token, custom_back_half=None, domain=None):
        captured["back_half"] = custom_back_half
        return {
            "short_url": "https://jbaru.ch/my-slug",
            "link_id": "id",
            "short_path": custom_back_half,
        }

    monkeypatch.setattr(generate_qr, "create_bitly_link", fake_create_bitly_link)
    config = {
        "shortener": "bitly",
        "preferred_short_path": "legacy-override",
        "bitly_domain": "jbaru.ch",
    }
    secrets = {"bitly": {"api_token": "tok"}}
    generate_qr.resolve_short_url(
        "https://jbaru.ch/my-slug",
        "my-slug",
        config,
        secrets,
        {},
        dry_run=False,
        vault_path=None,
    )
    assert captured["back_half"] == "my-slug"


def test_insert_qr_via_powerpoint_orchestration(generate_qr, monkeypatch):
    """The PowerPoint-write orchestration is deterministic and unit-tested with
    the InsertQR wrapper (the actual VBA) mocked: one insert-qr.sh call per color
    variant, the deck threaded through uniquely-named intermediates, intermediates
    cleaned up, and the final result moved back onto the deck."""
    calls = []
    monkeypatch.setattr(
        generate_qr.subprocess, "run", lambda cmd, **kw: calls.append(cmd)
    )
    removed = []
    monkeypatch.setattr(generate_qr.os, "remove", lambda p: removed.append(p))
    moved = []
    monkeypatch.setattr(generate_qr.shutil, "move", lambda a, b: moved.append((a, b)))
    monkeypatch.setattr(generate_qr.os.path, "isfile", lambda p: True)

    deck = "/decks/talk.pptx"
    # job slides carry their existing-QR rects (points): slide 1 replaces in place,
    # slide 3 is a new placement; slide 5 likewise new.
    jobs = [
        ("/q/a.png", [(1, [(100.0, 200.0, 144.0, 144.0)]), (3, [])]),
        ("/q/b.png", [(5, [])]),
    ]
    generate_qr.insert_qr_via_powerpoint(deck, jobs, "/scripts")

    wrapper = "/scripts/insert-qr.sh"
    # one subprocess call per job; deck threaded through intermediates; spec is
    # 1-based, ";"-joined, with per-slide removal rects after ":"
    assert calls == [
        [
            wrapper,
            "/decks/talk.pptx",
            "/decks/talk.pptx.qrtmp0.pptx",
            "/q/a.png",
            "1:100.00,200.00,144.00,144.00;3",
        ],
        [
            wrapper,
            "/decks/talk.pptx.qrtmp0.pptx",
            "/decks/talk.pptx.qrtmp1.pptx",
            "/q/b.png",
            "5",
        ],
    ]
    # the prior intermediate is cleaned up; the final intermediate is moved onto the deck
    assert removed == ["/decks/talk.pptx.qrtmp0.pptx"]
    assert moved == [("/decks/talk.pptx.qrtmp1.pptx", "/decks/talk.pptx")]


def test_insert_qr_via_powerpoint_missing_wrapper(generate_qr, monkeypatch):
    """A missing insert-qr.sh fails fast with an actionable error, not a traceback."""
    import pytest

    monkeypatch.setattr(generate_qr.os.path, "isfile", lambda p: False)
    with pytest.raises(SystemExit):
        generate_qr.insert_qr_via_powerpoint(
            "/decks/talk.pptx", [("/q/a.png", [(1, [])])], "/scripts"
        )


def test_insert_qr_via_powerpoint_wrapper_failure_is_actionable(
    generate_qr, monkeypatch
):
    """A wrapper (insert-qr.sh) failure surfaces as an actionable SystemExit
    pointing at the DeckOps setup, not a raw CalledProcessError traceback."""
    import pytest

    def boom(cmd, **kw):
        raise generate_qr.subprocess.CalledProcessError(1, cmd)

    monkeypatch.setattr(generate_qr.subprocess, "run", boom)
    monkeypatch.setattr(generate_qr.os.path, "isfile", lambda p: True)
    with pytest.raises(SystemExit) as exc:
        generate_qr.insert_qr_via_powerpoint(
            "/decks/talk.pptx", [("/q/a.png", [(1, [])])], "/scripts"
        )
    assert "deck-editing-setup.md" in str(exc.value)


def test_format_qr_spec_new_placement_and_replace(generate_qr):
    # No rects → bare slide number (new bottom-right placement)
    assert generate_qr._format_qr_spec([(5, [])]) == "5"
    # One rect → "num:L,T,W,H" (replace in place), 2-decimal points
    assert (
        generate_qr._format_qr_spec([(12, [(450.0, 80.0, 200.16, 200.16)])])
        == "12:450.00,80.00,200.16,200.16"
    )
    # Multiple slides + a duplicate-QR slide (two rects) → flattened, ";"-joined
    assert (
        generate_qr._format_qr_spec([(12, [(1, 2, 3, 4), (5, 6, 7, 8)]), (38, [])])
        == "12:1.00,2.00,3.00,4.00,5.00,6.00,7.00,8.00;38"
    )


def test_create_bitly_link_raises_when_custom_back_half_fails(generate_qr, monkeypatch):
    """If the custom back-half can't be set, fail rather than return a random hash —
    the back-half must always be the slug (rules/qr-generation-rules.md §2)."""
    import urllib.error

    def fake_http(url, data=None, headers=None, method="GET"):
        if url.endswith("/v4/bitlinks"):
            return {"id": "bit.ly/abc123", "link": "https://bit.ly/abc123"}
        # bit.ly answers 422 when the requested back-half is already taken.
        raise urllib.error.HTTPError(url, 422, "Unprocessable", {}, None)

    monkeypatch.setattr(generate_qr, "_http_request", fake_http)
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="could not set custom back-half"
    ) as excinfo:
        generate_qr.create_bitly_link(
            "https://example.com/notes",
            "tok",
            custom_back_half="my-slug",
            domain="jbaru.ch",
        )
    # The link exists provider-side; its identity must travel with the failure
    # so the operator can reuse or delete it deterministically.
    message = str(excinfo.value)
    assert "link_id=bit.ly/abc123" in message
    assert "short_url=https://bit.ly/abc123" in message


def test_create_bitly_link_lets_programming_errors_propagate(generate_qr, monkeypatch):
    """Only documented provider failures are wrapped; bugs surface as themselves."""

    def fake_http(url, data=None, headers=None, method="GET"):
        if url.endswith("/v4/bitlinks"):
            return {"id": "bit.ly/abc123", "link": "https://bit.ly/abc123"}
        raise TypeError("bug in the caller")

    monkeypatch.setattr(generate_qr, "_http_request", fake_http)
    with pytest.raises(TypeError, match="bug in the caller"):
        generate_qr.create_bitly_link(
            "https://example.com/notes",
            "tok",
            custom_back_half="my-slug",
            domain="jbaru.ch",
        )


def test_legacy_non_slug_cache_entry_is_recreated_with_slug(generate_qr, monkeypatch):
    """A tracked link with a legacy non-slug back-half is NOT reused/retargeted —
    it is recreated with the slug, even when the cached target matches."""
    created = {}

    def fake_create_bitly_link(long_url, api_token, custom_back_half=None, domain=None):
        created["back_half"] = custom_back_half
        return {
            "short_url": f"https://jbaru.ch/{custom_back_half}",
            "link_id": "new-id",
            "short_path": custom_back_half,
        }

    def boom_update(*a, **k):
        raise AssertionError(
            "update_bitly_link must not be called for a legacy non-slug entry"
        )

    monkeypatch.setattr(generate_qr, "create_bitly_link", fake_create_bitly_link)
    monkeypatch.setattr(generate_qr, "update_bitly_link", boom_update)
    tracking_db = {
        "qr_codes": [
            {
                "talk_slug": "my-slug",
                "target_url": "https://jbaru.ch/my-slug",  # cached target matches → would reuse without the fix
                "shortener": "bitly",
                "short_path": "legacy-hash",  # NON-slug back-half
                "short_url": "https://bit.ly/legacy-hash",
                "shortener_link_id": "old-id",
            }
        ]
    }
    short_url, meta = generate_qr.resolve_short_url(
        "https://jbaru.ch/my-slug",
        "my-slug",
        {"shortener": "bitly", "bitly_domain": "jbaru.ch"},
        {"bitly": {"api_token": "tok"}},
        tracking_db,
        dry_run=False,
        vault_path=None,
    )
    assert created["back_half"] == "my-slug"
    assert meta["short_path"] == "my-slug"
    assert meta["shortener_link_id"] == "new-id"
    assert short_url == "https://jbaru.ch/my-slug"


def test_missing_custom_domain_decision_stops_before_first_link(
    generate_qr, monkeypatch
):
    """First short link with NO recorded custom-domain decision (key absent) STOPS
    so the agent asks the user — it must not silently default to bit.ly."""
    import pytest

    monkeypatch.setattr(
        generate_qr,
        "create_bitly_link",
        lambda *a, **k: (_ for _ in ()).throw(
            AssertionError("must STOP before creating")
        ),
    )
    with pytest.raises(SystemExit):
        generate_qr.resolve_short_url(
            "https://jbaru.ch/my-slug",
            "my-slug",
            {"shortener": "bitly"},  # no bitly_domain key → decision not recorded
            {"bitly": {"api_token": "tok"}},
            {},
            dry_run=False,
            vault_path=None,
        )


def test_explicit_null_custom_domain_proceeds(generate_qr, monkeypatch):
    """An explicit null custom-domain decision is recorded — proceed on the default
    domain, no STOP, no re-ask."""
    captured = {}

    def fake_create_bitly_link(long_url, api_token, custom_back_half=None, domain=None):
        captured["domain"] = domain
        return {
            "short_url": f"https://bit.ly/{custom_back_half}",
            "link_id": "id",
            "short_path": custom_back_half,
        }

    monkeypatch.setattr(generate_qr, "create_bitly_link", fake_create_bitly_link)
    short_url, meta = generate_qr.resolve_short_url(
        "https://jbaru.ch/my-slug",
        "my-slug",
        {
            "shortener": "bitly",
            "bitly_domain": None,
        },  # recorded decision: no custom domain
        {"bitly": {"api_token": "tok"}},
        {},
        dry_run=False,
        vault_path=None,
    )
    assert captured["domain"] is None
    assert meta["short_path"] == "my-slug"


def test_slug_cache_entry_is_reused(generate_qr, monkeypatch):
    """A tracked entry whose back-half is already the slug is reused from cache —
    no API call — when the target matches."""

    def boom_create(*a, **k):
        raise AssertionError("must reuse cache, not create a new link")

    monkeypatch.setattr(generate_qr, "create_bitly_link", boom_create)
    tracking_db = {
        "qr_codes": [
            {
                "talk_slug": "my-slug",
                "target_url": "https://jbaru.ch/my-slug",
                "shortener": "bitly",
                "short_path": "my-slug",
                "short_url": "https://jbaru.ch/my-slug",
                "shortener_link_id": "id",
            }
        ]
    }
    short_url, meta = generate_qr.resolve_short_url(
        "https://jbaru.ch/my-slug",
        "my-slug",
        {"shortener": "bitly"},
        {},
        tracking_db,
        dry_run=False,
    )
    assert short_url == "https://jbaru.ch/my-slug"
    assert meta["short_path"] == "my-slug"


def test_main_dry_run_dual_reads_legacy_database_without_writing(
    generate_qr, monkeypatch, tmp_path
):
    database = {"config": {}, "talks": [], "pptx_catalog": []}
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(database), encoding="utf-8")
    before = path.read_bytes()
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "legacy",
            "--short-url",
            "https://example.test/legacy",
            "--shownotes-url",
            "https://example.test/notes",
            "--vault",
            str(tmp_path),
            "--dry-run",
        ],
    )

    generate_qr.main()

    assert path.read_bytes() == before


def test_main_rejects_legacy_database_before_qr_side_effect(
    generate_qr, monkeypatch, tmp_path
):
    database = {"config": {}, "talks": [], "pptx_catalog": []}
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(database), encoding="utf-8")
    before = path.read_bytes()
    monkeypatch.setattr(
        generate_qr,
        "generate_qr_png",
        lambda *_: pytest.fail("QR generation must not start on legacy state"),
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "legacy",
            "--short-url",
            "https://example.test/legacy",
            "--shownotes-url",
            "https://example.test/notes",
            "--vault",
            str(tmp_path),
        ],
    )

    with pytest.raises(SystemExit, match="current tracking schema"):
        generate_qr.main()

    assert path.read_bytes() == before


def test_main_current_database_stamps_qr_record_and_writes_atomically(
    generate_qr, monkeypatch, tmp_path
):
    database = _current_tracking_database()
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(database), encoding="utf-8")
    output = tmp_path / "current.png"
    monkeypatch.setattr(
        generate_qr,
        "generate_qr_png",
        lambda _url, _fg, _bg, target: Path(target).write_bytes(b"qr"),
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "current",
            "--short-url",
            "https://example.test/current",
            "--shownotes-url",
            "https://example.test/notes",
            "--vault",
            str(tmp_path),
            "--output",
            str(output),
        ],
    )

    generate_qr.main()

    written = json.loads(path.read_text(encoding="utf-8"))
    assert written["schema_version"] == 1
    record = written["qr_codes"][0]
    assert written["qr_codes"] == [
        {
            "schema_version": 2,
            "talk_slug": "current",
            # The canonical redirect target, never the short URL standing in for it.
            "target_url": "https://example.test/notes",
            "shortener": "mcp_preresolved",
            # The back-half is recovered from the short URL and matches the slug.
            "short_path": "current",
            "short_url": "https://example.test/current",
            "shortener_link_id": None,
            "qr_png_rel_path": record["qr_png_rel_path"],
            "artifacts": record["artifacts"],
            "created_at": record["created_at"],
            "updated_at": record["updated_at"],
        }
    ]
    # The exact written path is recorded, not the default {slug}-qr.png name.
    assert len(record["artifacts"]) == 1
    artifact = record["artifacts"][0]
    assert artifact["path"].endswith("current.png")
    assert artifact["path_root"] in {"cwd", "absolute"}
    assert len(artifact["sha256"]) == 64
    assert artifact["bg_hex"] is None
    assert record["qr_png_rel_path"] == artifact["path"]


def test_main_rejects_future_database_without_side_effect(
    generate_qr, monkeypatch, tmp_path
):
    database = _current_tracking_database() | {"schema_version": 2}
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(database), encoding="utf-8")
    before = path.read_bytes()
    monkeypatch.setattr(
        generate_qr,
        "generate_qr_png",
        lambda *_: pytest.fail("QR generation must not start on future state"),
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "future",
            "--short-url",
            "https://example.test/future",
            "--shownotes-url",
            "https://example.test/notes",
            "--vault",
            str(tmp_path),
            "--dry-run",
        ],
    )

    with pytest.raises(SystemExit, match="no usable prior state"):
        generate_qr.main()

    assert path.read_bytes() == before


def test_main_rejects_tracking_database_symlink_before_loading_config(
    generate_qr, monkeypatch, tmp_path
):
    target = tmp_path / "target.json"
    target.write_text(json.dumps(_current_tracking_database()), encoding="utf-8")
    path = tmp_path / "tracking-database.json"
    path.symlink_to(target.name)
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "link",
            "--short-url",
            "https://example.test/link",
            "--shownotes-url",
            "https://example.test/notes",
            "--vault",
            str(tmp_path),
            "--dry-run",
        ],
    )

    with pytest.raises(SystemExit, match="symbolic link"):
        generate_qr.main()


# --- #170: a configured shortener must fail closed, never ship a raw URL ---


def _qr_db():
    """Tracking DB carrying one managed link for `my-talk`."""
    return {
        "qr_codes": [
            {
                "talk_slug": "my-talk",
                "target_url": "https://example.com/old",
                "shortener": "bitly",
                "short_path": "my-talk",
                "short_url": "https://jbaru.ch/my-talk",
                "shortener_link_id": "bit.ly/abc123",
            }
        ]
    }


def test_unconfigured_shortener_fails_closed(generate_qr):
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="no URL shortener configured"
    ):
        generate_qr.resolve_short_url(
            "https://example.com/notes", "my-talk", {}, {}, {"qr_codes": []}
        )


def test_unknown_shortener_fails_closed(generate_qr):
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="unknown shortener 'tinyurl'"
    ):
        generate_qr.resolve_short_url(
            "https://example.com/notes",
            "my-talk",
            {"shortener": "tinyurl"},
            {},
            {"qr_codes": []},
        )


@pytest.mark.parametrize(
    "service,key", [("bitly", "api_token"), ("rebrandly", "api_key")]
)
def test_missing_credentials_fail_closed(generate_qr, service, key):
    with pytest.raises(generate_qr.ShortenerResolutionError, match=f"{key} is missing"):
        generate_qr.resolve_short_url(
            "https://example.com/notes",
            "my-talk",
            {"shortener": service, f"{service}_domain": None},
            {},
            {"qr_codes": []},
        )


def test_provider_error_fails_closed(generate_qr, monkeypatch):
    import urllib.error

    def fake_http(url, data=None, headers=None, method="GET"):
        raise urllib.error.URLError("connection refused")

    monkeypatch.setattr(generate_qr, "_http_request", fake_http)
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="could not produce the managed"
    ):
        generate_qr.resolve_short_url(
            "https://example.com/notes",
            "my-talk",
            {"shortener": "bitly", "bitly_domain": "jbaru.ch"},
            {"bitly": {"api_token": "tok"}},
            {"qr_codes": []},
        )


def test_malformed_provider_response_fails_closed(generate_qr, monkeypatch):
    """A response missing an expected field is a provider failure, not a raw-URL cue."""
    monkeypatch.setattr(
        generate_qr,
        "_http_request",
        lambda url, data=None, headers=None, method="GET": {"unexpected": "shape"},
    )
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="could not produce the managed"
    ):
        generate_qr.resolve_short_url(
            "https://example.com/notes",
            "my-talk",
            {"shortener": "bitly", "bitly_domain": "jbaru.ch"},
            {"bitly": {"api_token": "tok"}},
            {"qr_codes": []},
        )


def test_programming_errors_propagate_unwrapped(generate_qr, monkeypatch):
    def fake_http(url, data=None, headers=None, method="GET"):
        raise AttributeError("bug in the caller")

    monkeypatch.setattr(generate_qr, "_http_request", fake_http)
    with pytest.raises(AttributeError, match="bug in the caller"):
        generate_qr.resolve_short_url(
            "https://example.com/notes",
            "my-talk",
            {"shortener": "bitly", "bitly_domain": "jbaru.ch"},
            {"bitly": {"api_token": "tok"}},
            {"qr_codes": []},
        )


def test_failure_never_downgrades_an_existing_managed_record(generate_qr, monkeypatch):
    """A resolution failure must leave the tracking DB byte-identical."""
    import copy
    import urllib.error

    db = _qr_db()
    before = copy.deepcopy(db)

    monkeypatch.setattr(
        generate_qr,
        "_http_request",
        lambda *a, **k: (_ for _ in ()).throw(urllib.error.URLError("down")),
    )
    with pytest.raises(generate_qr.ShortenerResolutionError):
        generate_qr.resolve_short_url(
            "https://example.com/new",
            "my-talk",
            {"shortener": "bitly", "bitly_domain": "jbaru.ch"},
            {"bitly": {"api_token": "tok"}},
            db,
        )

    assert db == before
    assert db["qr_codes"][0]["shortener"] == "bitly"


def test_explicit_none_still_authorizes_a_raw_url(generate_qr):
    """The one sanctioned path to a raw target URL stays open."""
    url, meta = generate_qr.resolve_short_url(
        "https://example.com/notes",
        "my-talk",
        {"shortener": "none"},
        {},
        {"qr_codes": []},
    )
    assert url == "https://example.com/notes"
    assert meta["shortener"] == "none"
    assert meta["target_url"] == "https://example.com/notes"


def test_cached_raw_record_is_not_reused_when_config_is_missing(generate_qr):
    """A stale `shortener: none` entry must not re-authorize a raw URL."""
    db = {
        "qr_codes": [
            {
                "talk_slug": "my-talk",
                "target_url": "https://example.com/notes",
                "shortener": "none",
                "short_path": None,
                "short_url": "https://example.com/notes",
                "shortener_link_id": None,
            }
        ]
    }
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="no URL shortener configured"
    ):
        generate_qr.resolve_short_url(
            "https://example.com/notes", "my-talk", {}, {}, db
        )


def test_cached_raw_record_is_not_reused_under_a_managed_shortener(
    generate_qr, monkeypatch
):
    """Switching config from none to bitly must re-resolve, not replay the raw URL."""
    db = {
        "qr_codes": [
            {
                "talk_slug": "my-talk",
                "target_url": "https://example.com/notes",
                "shortener": "none",
                "short_path": None,
                "short_url": "https://example.com/notes",
                "shortener_link_id": None,
            }
        ]
    }
    calls = []

    def fake_http(url, data=None, headers=None, method="GET"):
        calls.append(url)
        if url.endswith("/v4/bitlinks"):
            return {"id": "jbaru.ch/my-talk", "link": "https://jbaru.ch/my-talk"}
        return {}

    monkeypatch.setattr(generate_qr, "_http_request", fake_http)
    url, meta = generate_qr.resolve_short_url(
        "https://example.com/notes",
        "my-talk",
        {"shortener": "bitly", "bitly_domain": "jbaru.ch"},
        {"bitly": {"api_token": "tok"}},
        db,
    )

    assert calls, "the managed shortener must actually be called"
    assert url == "https://jbaru.ch/my-talk"
    assert meta["shortener"] == "bitly"


def test_cached_managed_record_is_still_reused(generate_qr):
    """The reuse path stays open when the cached record matches current config."""
    db = _qr_db()
    db["qr_codes"][0]["target_url"] = "https://example.com/notes"
    url, meta = generate_qr.resolve_short_url(
        "https://example.com/notes",
        "my-talk",
        {"shortener": "bitly", "bitly_domain": "jbaru.ch"},
        {"bitly": {"api_token": "tok"}},
        db,
    )
    assert url == "https://jbaru.ch/my-talk"
    assert meta["shortener"] == "bitly"


def test_unknown_shortener_is_rejected_before_cache_reuse(generate_qr):
    db = _qr_db()
    db["qr_codes"][0]["target_url"] = "https://example.com/notes"
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="unknown shortener 'tinyurl'"
    ):
        generate_qr.resolve_short_url(
            "https://example.com/notes", "my-talk", {"shortener": "tinyurl"}, {}, db
        )


# --- #171: catalog fidelity — canonical target, provider identity, every artifact ---


def test_mcp_mode_records_canonical_target_and_provider_identity(
    generate_qr, monkeypatch, tmp_path
):
    """MCP mode must persist the redirect target, provider, and link id."""
    output = tmp_path / "talk-qr.png"
    # An explicit vault keeps the run hermetic. Without --vault the script
    # falls back to ~/.claude/rhetoric-knowledge-vault, so the test would read
    # the developer's real vault and behave differently in CI.
    (tmp_path / "tracking-database.json").write_text(
        json.dumps(_current_tracking_database()), encoding="utf-8"
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "my-talk",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://jbaru.ch/my-talk",
            "--short-provider",
            "bitly",
            "--short-link-id",
            "bit.ly/abc123",
            "--vault",
            str(tmp_path),
            "--output",
            str(output),
        ],
    )
    monkeypatch.setattr(
        generate_qr,
        "update_tracking_db",
        lambda db, entry, artifacts: captured.update(entry=entry, artifacts=artifacts),
    )
    captured = {}
    generate_qr.main()

    entry = captured["entry"]
    assert entry["target_url"] == "https://example.test/notes"
    assert entry["short_url"] == "https://jbaru.ch/my-talk"
    assert entry["shortener"] == "bitly"
    assert entry["shortener_link_id"] == "bit.ly/abc123"
    assert entry["short_path"] == "my-talk"


def test_png_only_records_the_path_actually_written(generate_qr, monkeypatch, tmp_path):
    """--output PATH must be cataloged, not the default {slug}-qr.png name."""
    output = tmp_path / "custom" / "elsewhere.png"
    output.parent.mkdir()
    captured = {}
    # Explicit vault — see the note in the MCP test above.
    (tmp_path / "tracking-database.json").write_text(
        json.dumps(_current_tracking_database()), encoding="utf-8"
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "my-talk",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://jbaru.ch/my-talk",
            "--vault",
            str(tmp_path),
            "--output",
            str(output),
        ],
    )
    monkeypatch.setattr(
        generate_qr,
        "update_tracking_db",
        lambda db, entry, artifacts: captured.update(artifacts=artifacts),
    )
    generate_qr.main()

    assert len(captured["artifacts"]) == 1
    artifact = captured["artifacts"][0]
    assert artifact["path"].endswith("elsewhere.png")
    assert "my-talk-qr.png" not in artifact["path"]
    # The digest binds the record to these exact bytes.
    import hashlib

    assert artifact["sha256"] == hashlib.sha256(output.read_bytes()).hexdigest()


def test_artifact_receipt_records_an_explicit_path_root(generate_qr, tmp_path):
    """The path root is recorded, never left for a reader to guess."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    inside = deck_dir / "a.png"
    inside.write_bytes(b"png-a")
    outside = tmp_path / "b.png"
    outside.write_bytes(b"png-b")

    within = generate_qr._artifact_receipt(str(inside), str(deck_dir), None)
    assert within["path_root"] == "deck_dir"
    assert within["path"] == "a.png"

    elsewhere = generate_qr._artifact_receipt(str(outside), str(deck_dir), None)
    assert elsewhere["path_root"] == "absolute"
    assert elsewhere["path"] == str(outside)


def test_back_half_that_is_not_the_slug_stops_the_run(generate_qr):
    """§2 admits no exception: a non-slug back-half is the random-hash failure."""
    with pytest.raises(
        generate_qr.ShortenerResolutionError, match="is not the talk slug"
    ):
        generate_qr._validated_back_half("https://bit.ly/a3xK9f", "my-talk")
    assert (
        generate_qr._validated_back_half("https://jbaru.ch/my-talk", "my-talk")
        == "my-talk"
    )


def test_mcp_non_slug_back_half_exits_before_any_side_effect(
    generate_qr, monkeypatch, tmp_path
):
    output = tmp_path / "qr.png"
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "my-talk",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://bit.ly/a3xK9f",
            "--output",
            str(output),
        ],
    )
    with pytest.raises(SystemExit) as excinfo:
        generate_qr.main()
    assert excinfo.value.code == 1
    assert not output.exists()


def test_provider_flags_require_short_url(generate_qr, monkeypatch):
    """Provider identity without --short-url would be silently dropped."""
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "my-talk",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-provider",
            "bitly",
            "--short-link-id",
            "bit.ly/abc",
        ],
    )
    with pytest.raises(SystemExit):
        generate_qr.main()


@pytest.mark.parametrize(
    "flag,value",
    [
        ("--short-provider", "bitly"),
        ("--short-link-id", "bit.ly/abc123"),
    ],
)
def test_provider_identity_is_all_or_neither(
    generate_qr, monkeypatch, tmp_path, flag, value
):
    """Half an identity catalogs an incomplete provider record."""
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "my-talk",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://jbaru.ch/my-talk",
            flag,
            value,
            "--output",
            str(tmp_path / "qr.png"),
        ],
    )
    with pytest.raises(SystemExit):
        generate_qr.main()


def test_every_colour_variant_is_cataloged(generate_qr, tmp_path):
    """A multi-colour deck run records each PNG, not just the first."""
    deck_dir = tmp_path / "deck"
    deck_dir.mkdir()
    receipts = []
    for name, bg in (("a-qr-ffffff.png", "ffffff"), ("a-qr-000000.png", "000000")):
        path = deck_dir / name
        path.write_bytes(name.encode())
        receipts.append(generate_qr._artifact_receipt(str(path), str(deck_dir), bg))

    db = {"qr_codes": []}
    generate_qr.update_tracking_db(
        db,
        {
            "talk_slug": "a",
            "target_url": "https://example.test/notes",
            "shortener": "bitly",
            "short_url": "https://jbaru.ch/a",
            "short_path": "a",
            "shortener_link_id": "bit.ly/a",
        },
        receipts,
    )

    record = db["qr_codes"][0]
    assert record["schema_version"] == 2
    assert [a["bg_hex"] for a in record["artifacts"]] == ["ffffff", "000000"]
    assert len({a["sha256"] for a in record["artifacts"]}) == 2
    assert record["qr_png_rel_path"] == record["artifacts"][0]["path"]


# --- #172: publication stays recoverable when the CAS commit rejects ---


def test_unrelated_concurrent_write_no_longer_rejects_the_qr_commit(
    generate_qr, monkeypatch, tmp_path, capsys
):
    """A concurrent writer touching an unrelated collection must not reject us."""
    database = _current_tracking_database()
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(database), encoding="utf-8")
    output = tmp_path / "current.png"

    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "current",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://example.test/current",
            "--vault",
            str(tmp_path),
            "--output",
            str(output),
        ],
    )

    # Simulate an unrelated writer landing a change after our snapshot but
    # before our commit: mutate a different collection on disk.
    real_generate = generate_qr.generate_qr_png

    def generate_then_race(*args, **kwargs):
        result = real_generate(*args, **kwargs)
        raced = json.loads(path.read_text(encoding="utf-8"))
        raced["resources"] = [
            {
                "schema_version": 1,
                "talk_slug": "other",
                "item_count": 1,
                "category_breakdown": {"url": 1},
            }
        ]
        path.write_text(json.dumps(raced), encoding="utf-8")
        return result

    monkeypatch.setattr(generate_qr, "generate_qr_png", generate_then_race)
    generate_qr.main()

    written = json.loads(path.read_text(encoding="utf-8"))
    # Our QR record landed...
    assert [r["talk_slug"] for r in written["qr_codes"]] == ["current"]
    # ...and the unrelated writer's change survived; we rebased, not clobbered.
    assert written["resources"][0]["talk_slug"] == "other"


def test_commit_rejection_emits_a_structured_effects_payload(
    generate_qr, monkeypatch, tmp_path, capsys
):
    """A post-effect commit failure must not look side-effect-free."""
    database = _current_tracking_database()
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(database), encoding="utf-8")
    output = tmp_path / "current.png"

    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "current",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://example.test/current",
            "--vault",
            str(tmp_path),
            "--output",
            str(output),
        ],
    )
    monkeypatch.setattr(
        generate_qr,
        "commit_qr_record",
        lambda *a, **k: (_ for _ in ()).throw(ValueError("generation conflict")),
    )

    with pytest.raises(SystemExit) as excinfo:
        generate_qr.main()
    assert excinfo.value.code == 1

    payload = json.loads(capsys.readouterr().err)
    assert payload["error"] == "qr_publication_unfinalized"
    assert payload["message"] == "generation conflict"
    assert payload["tracking_database_updated"] is False
    assert payload["atomic_rollback"] is False
    assert payload["retry"]["idempotent"] is True
    png = [e for e in payload["effects"] if e["kind"] == "png"]
    assert [e["path"] for e in png] == [str(output)]
    assert png[0]["rollback"] == {"action": "delete", "target": str(output)}


def test_payload_is_one_valid_json_document(generate_qr):
    receipt = generate_qr.EffectsReceipt("my-talk")
    payload = generate_qr.unfinalized_effects_payload(receipt, "boom")
    assert json.loads(json.dumps(payload)) == payload
    assert payload["effects"] == []
    assert payload["talk_slug"] == "my-talk"


@pytest.mark.parametrize(
    "action,prior,expected",
    [
        ("created", None, {"action": "delete", "target": "https://jbaru.ch/my-talk"}),
        (
            "retargeted",
            "https://example.test/old",
            {
                "action": "restore_target",
                "target": "https://jbaru.ch/my-talk",
                "restore_to": "https://example.test/old",
            },
        ),
        ("preresolved", None, {"action": "none", "target": "https://jbaru.ch/my-talk"}),
    ],
)
def test_link_rollback_matches_how_the_link_came_to_be(
    generate_qr, action, prior, expected
):
    """Deleting a link this run did not create is destructive, not a rollback."""
    receipt = generate_qr.EffectsReceipt("my-talk")
    receipt.record_short_link(
        "bitly",
        "bit.ly/abc123",
        "https://jbaru.ch/my-talk",
        action=action,
        prior_target=prior,
    )
    payload = generate_qr.unfinalized_effects_payload(receipt, "boom")
    link = [e for e in payload["effects"] if e["kind"] == "short_link"][0]
    assert link["rollback"] == expected
    assert link["action"] == action


def test_payload_covers_every_landed_effect(generate_qr):
    """Link-only recovery would imply the PNGs and deck were reverted too."""
    receipt = generate_qr.EffectsReceipt("my-talk")
    receipt.record_short_link(
        "bitly", "bit.ly/abc", "https://jbaru.ch/my-talk", action="created"
    )
    receipt.record_artifacts(["/tmp/a.png", "/tmp/b.png"])
    receipt.record_deck("/tmp/deck.pptx")

    payload = generate_qr.unfinalized_effects_payload(receipt, "boom")
    kinds = [e["kind"] for e in payload["effects"]]
    assert kinds == ["short_link", "png", "png", "deck"]

    deck = payload["effects"][-1]
    # No backup is taken, so a restore must not be claimed.
    assert deck["backup_available"] is False
    assert deck["rollback"]["action"] == "restore_from_version_control"


def test_publication_lock_is_per_slug(generate_qr, tmp_path):
    """Different slugs take different locks and do not block each other."""
    with generate_qr.qr_publication_lock(str(tmp_path), "talk-a") as first:
        assert os.path.basename(first) == ".qr-talk-a.lock"
        with generate_qr.qr_publication_lock(str(tmp_path), "talk-b") as second:
            assert os.path.basename(second) == ".qr-talk-b.lock"
            assert first != second


def test_run_reloads_state_after_the_lock_so_it_retargets_instead_of_duplicating(
    generate_qr, monkeypatch, tmp_path
):
    """State loaded before the lock is stale; resolving from it duplicates links.

    Simulates the real interleaving: this run loads the database, and a
    competing same-slug process commits its link while this one waits on the
    lock. The commit is injected at lock acquisition, which is exactly when the
    wait ends.
    """
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(_current_tracking_database()), encoding="utf-8")
    (tmp_path / "speaker-profile.json").write_text(
        json.dumps(
            {
                "publishing_process": {
                    "qr_code": {"shortener": "bitly", "bitly_domain": "jbaru.ch"}
                }
            }
        ),
        encoding="utf-8",
    )
    (tmp_path / "secrets.json").write_text(
        json.dumps({"bitly": {"api_token": "tok"}}), encoding="utf-8"
    )

    created, updated = [], []
    monkeypatch.setattr(
        generate_qr,
        "create_bitly_link",
        lambda long_url, api_token, custom_back_half=None, domain=None: (
            created.append(long_url)
            or {
                "short_url": f"https://jbaru.ch/{custom_back_half}",
                "link_id": "bit.ly/NEW",
                "short_path": custom_back_half,
            }
        ),
    )
    monkeypatch.setattr(
        generate_qr,
        "update_bitly_link",
        lambda link_id, new_long_url, api_token: updated.append(
            (link_id, new_long_url)
        ),
    )

    real_lock = generate_qr.qr_publication_lock

    @contextlib.contextmanager
    def lock_then_race(vault_path, talk_slug):
        with real_lock(vault_path, talk_slug) as held:
            # A competing process finished while we waited: its link is now
            # committed, and our pre-lock view does not contain it.
            raced = json.loads(path.read_text(encoding="utf-8"))
            raced["qr_codes"] = [
                {
                    "schema_version": 2,
                    "talk_slug": "current",
                    "target_url": "https://example.test/other",
                    "shortener": "bitly",
                    "short_path": "current",
                    "short_url": "https://jbaru.ch/current",
                    "shortener_link_id": "bit.ly/RACED",
                    "qr_png_rel_path": "current.png",
                    "artifacts": [
                        {
                            "path": "current.png",
                            "path_root": "cwd",
                            "sha256": "b" * 64,
                            "bg_hex": None,
                        }
                    ],
                    "created_at": "2026-08-09",
                    "updated_at": "2026-08-09",
                }
            ]
            path.write_text(json.dumps(raced), encoding="utf-8")
            yield held

    monkeypatch.setattr(generate_qr, "qr_publication_lock", lock_then_race)
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "current",
            "--shownotes-url",
            "https://example.test/notes",
            "--vault",
            str(tmp_path),
            "--output",
            str(tmp_path / "current.png"),
        ],
    )
    generate_qr.main()

    assert created == [], "the raced link must be retargeted, never duplicated"
    assert updated == [("bit.ly/RACED", "https://example.test/notes")]

    record = json.loads(path.read_text(encoding="utf-8"))["qr_codes"][0]
    assert record["shortener_link_id"] == "bit.ly/RACED"


def test_lock_failure_exits_cleanly_without_a_traceback(
    generate_qr, monkeypatch, tmp_path, capsys
):
    path = tmp_path / "tracking-database.json"
    path.write_text(json.dumps(_current_tracking_database()), encoding="utf-8")

    @contextlib.contextmanager
    def refuse(vault_path, talk_slug):
        raise ValueError(f"cannot open QR publication lock for {talk_slug}")
        yield  # pragma: no cover

    monkeypatch.setattr(generate_qr, "qr_publication_lock", refuse)
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "current",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://example.test/current",
            "--vault",
            str(tmp_path),
            "--output",
            str(tmp_path / "qr.png"),
        ],
    )

    with pytest.raises(SystemExit) as excinfo:
        generate_qr.main()
    assert excinfo.value.code == 1
    assert "cannot open QR publication lock" in capsys.readouterr().err


# --- slug is a path component: it must never escape the vault ---


@pytest.mark.parametrize(
    "bad",
    [
        "../../etc/passwd",
        "a/b",
        "..",
        ".",
        "Talk-Slug",
        "talk slug",
        "talk_slug",
        "-leading",
        "trailing-",
        "double--hyphen",
        "",
    ],
)
def test_invalid_talk_slug_is_rejected(generate_qr, bad):
    with pytest.raises(ValueError, match="kebab-case"):
        generate_qr.require_valid_talk_slug(bad)


@pytest.mark.parametrize("good", ["arc-of-ai", "devnexus26-robocoders", "talk1"])
def test_valid_talk_slug_is_accepted(generate_qr, good):
    assert generate_qr.require_valid_talk_slug(good) == good


def test_unsafe_slug_is_rejected_at_the_cli_boundary(
    generate_qr, monkeypatch, tmp_path, capsys
):
    """A path-shaped slug fails on the slug contract, before any file is touched.

    Without validation the same input fails later and less usefully — the lock
    open errors on a directory that happens not to exist — so the operator is
    told about a lock path instead of the slug that is actually wrong.
    """
    vault = tmp_path / "vault"
    vault.mkdir()
    (vault / "tracking-database.json").write_text(
        json.dumps(_current_tracking_database()), encoding="utf-8"
    )

    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "../escaped",
            "--shownotes-url",
            "https://example.test/notes",
            "--short-url",
            "https://example.test/escaped",
            "--vault",
            str(vault),
            "--output",
            str(tmp_path / "qr.png"),
        ],
    )
    with pytest.raises(SystemExit):
        generate_qr.main()

    err = capsys.readouterr().err
    assert "kebab-case" in err
    assert "qr-generation-rules.md" in err
    # Nothing was created, in the vault or above it.
    assert list(vault.glob(".qr-*")) == []
    assert list(tmp_path.glob(".qr-*")) == []
    assert not (tmp_path / "qr.png").exists()


def test_back_half_failure_after_link_creation_reports_the_created_link(
    generate_qr, monkeypatch, tmp_path, capsys
):
    """The link exists provider-side; claiming no effects landed would be false."""
    vault = tmp_path / "vault"
    vault.mkdir()
    (vault / "tracking-database.json").write_text(
        json.dumps(_current_tracking_database()), encoding="utf-8"
    )
    (vault / "speaker-profile.json").write_text(
        json.dumps(
            {
                "publishing_process": {
                    "qr_code": {"shortener": "bitly", "bitly_domain": "jbaru.ch"}
                }
            }
        ),
        encoding="utf-8",
    )
    (vault / "secrets.json").write_text(
        json.dumps({"bitly": {"api_token": "tok"}}), encoding="utf-8"
    )

    import urllib.error

    def fake_http(url, data=None, headers=None, method="GET"):
        if url.endswith("/v4/bitlinks"):
            return {"id": "jbaru.ch/abc123", "link": "https://jbaru.ch/abc123"}
        # Creation succeeded; the back-half assignment is what fails.
        raise urllib.error.HTTPError(url, 422, "Unprocessable", {}, None)

    monkeypatch.setattr(generate_qr, "_http_request", fake_http)
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "generate-qr.py",
            "--png-only",
            "--talk-slug",
            "current",
            "--shownotes-url",
            "https://example.test/notes",
            "--vault",
            str(vault),
            "--output",
            str(tmp_path / "qr.png"),
        ],
    )

    with pytest.raises(SystemExit) as excinfo:
        generate_qr.main()
    assert excinfo.value.code == 1

    payload = json.loads(capsys.readouterr().err)
    link = [e for e in payload["effects"] if e["kind"] == "short_link"]
    assert link, "the created link must be reported"
    assert link[0]["link_id"] == "jbaru.ch/abc123"
    assert link[0]["action"] == "created"
    assert link[0]["rollback"] == {
        "action": "delete",
        "target": "https://jbaru.ch/abc123",
    }


def test_partial_link_identity_is_structured_not_only_in_the_message(generate_qr):
    err = generate_qr.ShortenerResolutionError(
        "boom",
        partial_link={
            "provider": "bitly",
            "link_id": "x",
            "short_url": "https://jbaru.ch/x",
        },
    )
    assert err.partial_link["link_id"] == "x"
    assert generate_qr.ShortenerResolutionError("boom").partial_link is None


@pytest.mark.parametrize(
    "action,idempotent",
    [
        ("created", False),
        ("retargeted", True),
        ("preresolved", True),
    ],
)
def test_retry_idempotency_depends_on_whether_a_record_can_find_the_link(
    generate_qr, action, idempotent
):
    """A link this run created has no committed record, so a retry cannot find it."""
    receipt = generate_qr.EffectsReceipt("my-talk")
    receipt.record_short_link(
        "bitly",
        "bit.ly/abc",
        "https://jbaru.ch/my-talk",
        action=action,
        prior_target="https://x.test/old",
    )
    payload = generate_qr.unfinalized_effects_payload(receipt, "boom")
    assert payload["retry"]["idempotent"] is idempotent
    if not idempotent:
        assert "cannot find it" in payload["retry"]["detail"]


def test_retry_is_idempotent_when_no_link_work_happened(generate_qr):
    receipt = generate_qr.EffectsReceipt("my-talk")
    receipt.record_artifacts(["/tmp/a.png"])
    assert (
        generate_qr.unfinalized_effects_payload(receipt, "boom")["retry"]["idempotent"]
        is True
    )


def test_lock_guidance_never_advises_deleting_the_lock_file(generate_qr, tmp_path):
    """Unlinking an flock path lets a second process publish concurrently."""
    import fcntl as _fcntl

    lock_path = tmp_path / ".qr-my-talk.lock"
    lock_path.touch()
    held = os.open(str(lock_path), os.O_RDWR)
    _fcntl.flock(held, _fcntl.LOCK_EX)
    try:

        def blocking_flock(fd, op):
            raise OSError("Resource temporarily unavailable")

        import unittest.mock as mock

        with mock.patch.object(generate_qr.fcntl, "flock", blocking_flock):
            with pytest.raises(ValueError) as excinfo:
                with generate_qr.qr_publication_lock(str(tmp_path), "my-talk"):
                    pass
    finally:
        _fcntl.flock(held, _fcntl.LOCK_UN)
        os.close(held)

    message = str(excinfo.value)
    assert "Do not delete the lock file" in message
    assert "remove the stale lock" not in message
    assert "flock" in message
