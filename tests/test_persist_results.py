"""Tests for persist-results.py — deterministic Step 4 merge of subagent returns.

Regression coverage for #97: structured_data computed by subagents must land in
the tracking DB, with the declared queryable scalars promoted to the talk top level.
"""

import copy
import importlib
import json
import pathlib
import shutil
import subprocess
import sys
from datetime import datetime, timedelta, timezone

import pytest
from pypdf import PdfWriter
from pptx import Presentation

from conftest import current_tracking_config


def test_atomic_json_write_cleans_stage_and_propagates_interrupt(
    persist_results, tracking_database_io, tmp_path, monkeypatch,
):
    target = tmp_path / "tracking-database.json"
    target.write_text('{"old": true}\n', encoding="utf-8")

    def interrupt(_source, _target):
        raise KeyboardInterrupt

    monkeypatch.setattr(tracking_database_io.os, "replace", interrupt)

    with pytest.raises(KeyboardInterrupt):
        persist_results.atomic_write_json(target, {"new": True})

    assert json.loads(target.read_text(encoding="utf-8")) == {"old": True}
    assert {path.name for path in tmp_path.iterdir()} == {
        target.name,
        ".tracking-database.json.lock",
    }


def _return(**overrides):
    ret = {
        "filename": "talk.md",
        "return_schema_version": 2,
        "queue_claim": {
            "run_id": "reparse",
            "batch_id": "25",
            "reprocess_generation": 1,
        },
        "status": "processed",
        "processed_date": "2026-06-18",
        "transcript_source": "youtube_auto",
        "slide_source": "pptx",
        "rhetoric_notes": "notes",
        "areas_for_improvement": "improve",
        "adherence_assessment": "above baseline",
        "new_patterns": "",
        "summary_updates": "",
        "structured_data": {
            "delivery_language": "en",
            "co_presenter": False,
            "slide_count": 62,
            "audience_interaction_count": 3,
            "opening_type": "demo_cold_open",
            "closing_type": "summary_cta",
            "narrative_arc_type": "problem_diagnosis_solution",
            "slide_design_style": "comic_book",
            "illustration_style": "comic_book",
        },
        "verbatim_examples": {"jokes": ["j1"]},
        "pattern_observations": {
            "patterns_detected": [
                {"pattern_id": "narrative-arc", "confidence": "strong",
                 "evidence_source": "transcript",
                 "evidence": "The talk follows a four-act argument."},
                {"pattern_id": "bookends", "confidence": "moderate",
                 "evidence_source": "static_slides",
                 "evidence": "Repeated dividers mark section boundaries."},
            ],
            "antipatterns_detected": [
                {"pattern_id": "shortchanged", "confidence": "weak",
                 "evidence_source": "transcript",
                 "evidence": "The close is compressed."},
            ],
            "evidence_sources": ["transcript", "native_deck", "static_slides",
                                 "delivery_video", "source_comparison"],
            "not_evaluable": [],
            "pattern_score": {"patterns_used": 2, "antipatterns_detected": 1, "score": 1},
        },
        "catalog_feedback": {
            "unmatched_observations": [],
            "confusable_pairs": [],
            "definition_problems": [],
            "scoring_problems": [],
            "tensions": [],
        },
    }
    ret.update(overrides)
    if (ret.get("return_schema_version") == 3 and
            "adherence_assessment" not in overrides and
            "adherence_comparison" not in overrides):
        ret["adherence_assessment"] = ""
    return ret


def _talk(**overrides):
    talk = {
        "schema_version": 5,
        "filename": "talk.md",
        "status": "reprocessing-inflight",
        "reprocess_generation": 1,
        "video_url": "https://youtu.be/AbCdEfGhI_1",
        "youtube_id": "AbCdEfGhI_1",
        "pptx_path": "Conference/Talk.pptx",
        "slides_url": "https://drive.google.com/file/d/slides-id/view",
        "_queue_claim": {
            "schema_version": 1,
            "run_id": "reparse",
            "batch_id": "25",
            "claimed_at": "2026-07-31T18:00:00+00:00",
            "previous_status": "needs-reprocessing",
            "reprocess_generation": 1,
            "state": "claimed",
        },
        "structured_data": {},
        "verbatim_examples": {},
        "pattern_observations": {"pattern_ids": [], "antipattern_ids": [], "pattern_score": 0},
    }
    talk.update(overrides)
    return talk


def _write_tiny_mp4(path):
    ffmpeg = shutil.which("ffmpeg")
    assert ffmpeg is not None, "source-video persistence test requires ffmpeg"
    assert shutil.which("ffprobe") is not None, (
        "source-video persistence test requires ffprobe"
    )
    created = subprocess.run(
        [
            ffmpeg,
            "-nostdin",
            "-hide_banner",
            "-loglevel",
            "error",
            "-f",
            "lavfi",
            "-i",
            "color=c=black:s=160x90:r=2",
            "-t",
            "0.25",
            "-an",
            "-c:v",
            "mpeg4",
            "-pix_fmt",
            "yuv420p",
            "-y",
            str(path),
        ],
        capture_output=True,
        check=False,
    )
    assert created.returncode == 0, created.stderr.decode(
        "utf-8", errors="replace"
    )


def _db_json(database):
    """Render a current tracking database around one test's talk fixtures."""
    database["schema_version"] = 1
    config = database.setdefault("config", {})
    if isinstance(config, dict):
        existing = dict(config)
        config.clear()
        config.update(current_tracking_config(**existing))
    database.setdefault("pptx_catalog", [])
    database.setdefault("qr_codes", [])
    database.setdefault("resources", [])
    database.setdefault("thumbnails", [])
    database.setdefault("confirmed_intents", [])
    database.setdefault("improvement_goals", [])
    return json.dumps(database)


def _adherence_baseline(persist_results, *, count=0, filenames=("talk.md",)):
    return {
        "schema_version": 1,
        "as_of": "2026-07-31T18:00:00+00:00",
        "scope": "global",
        "active_batch_excluded": True,
        "excluded_filenames": sorted(filenames),
        "eligible_statuses": ["processed", "processed_partial"],
        "pattern_scoring_generation_status": "current",
        "pattern_scoring_generation_reasons": [],
        "pattern_catalog_fingerprint": persist_results.load_catalog().fingerprint,
        "pattern_scoring_schema_version": (
            persist_results.PATTERN_SCORING_SCHEMA_VERSION),
        "scored_talk_count": count,
        "pattern_score_sum": count,
        "average_pattern_score": 1.0 if count else None,
    }


def _v3_talk_and_return(persist_results, *, filename="talk.md"):
    ret = _return(
        filename=filename,
        return_schema_version=3,
        adherence_assessment="",
    )
    talk = _talk(filename=filename)
    talk["_queue_claim"].update({
        "schema_version": 3,
        "required_return_schema_version": 3,
        "adherence_baseline": _adherence_baseline(
            persist_results,
            filenames=(filename,),
        ),
    })
    return talk, ret


def _skipped_return(**overrides):
    ret = {
        "filename": "talk.md",
        "return_schema_version": 2,
        "queue_claim": {
            "run_id": "reparse",
            "batch_id": "25",
            "reprocess_generation": 1,
        },
        "status": "skipped_no_sources",
    }
    ret.update(overrides)
    return ret


def _v4_transcript_batch(return_validation, tmp_path, *, language="en"):
    """Build one real-artifact v4 batch with exhaustive catalog outcomes."""
    transcript = tmp_path / "transcripts" / "manual-talk.txt"
    transcript.parent.mkdir()
    lines = [
        "A uniquely phrased production failure opens this synthetic talk clearly."
    ]
    lines.extend(
        f"Synthetic transcript line {number} provides substantive source evidence "
        "for deterministic catalog verification."
        for number in range(2, 82)
    )
    transcript.write_text("\n".join(lines), encoding="utf-8")
    transcript_timing = importlib.import_module("transcript_timing")
    transcript_text = transcript.read_text(encoding="utf-8")
    transcript_timing.write_quality_receipt(
        transcript,
        transcript_text,
        transcript_timing.build_quality_policy(400),
        {"kind": "fixed_default"},
    )
    filename = "talk.md"
    baseline = _adherence_baseline(
        return_validation, count=0, filenames=(filename,))
    claim = {
        "schema_version": 4,
        "run_id": "reparse-v4",
        "batch_id": "batch-1",
        "claimed_at": "2026-07-31T18:00:00+00:00",
        "previous_status": "needs-reprocessing",
        "reprocess_generation": 1,
        "state": "claimed",
        "required_return_schema_version": 4,
        "adherence_baseline": baseline,
    }
    talk = {
        "schema_version": 4,
        "filename": filename,
        "title": "Synthetic v4 Talk",
        "status": "reprocessing-inflight",
        "reprocess_generation": 1,
        "transcript_path": "transcripts/manual-talk.txt",
        "transcript_source": "manual",
        "slide_source": "none",
        "delivery_language": language,
        "_queue_claim": claim,
    }
    detected_id = "echo-chamber"
    catalog = return_validation.load_catalog()
    not_evaluable = []
    for pattern_id, entry in sorted(catalog.entries.items()):
        if not entry.observable or pattern_id == detected_id:
            continue
        gate = entry.absence_evaluable_from
        if gate is None:
            not_evaluable.append({
                "pattern_id": pattern_id,
                "reason_code": (
                    "absence_not_authorized_by_catalog"
                    if entry.evaluable_from is not None
                    else "source_gate_pending_owner_review"
                ),
            })
        elif not any(group == frozenset({"transcript"}) for group in gate):
            not_evaluable.append({
                "pattern_id": pattern_id,
                "reason_code": "missing_required_source_coverage",
            })
    ret = {
        "filename": filename,
        "return_schema_version": 4,
        "queue_claim": {
            "run_id": "reparse-v4",
            "batch_id": "batch-1",
            "reprocess_generation": 1,
        },
        "status": "processed_partial",
        "transcript_source": "manual",
        "slide_source": "none",
        "rhetoric_notes": "The talk uses a clear incident-led narrative.",
        "areas_for_improvement": "The close could return to the opening incident.",
        "adherence_assessment": "",
        "new_patterns": "",
        "summary_updates": "",
        "structured_data": {
            "delivery_language": language,
            "co_presenter": False,
        },
        "verbatim_examples": {},
        "pattern_observations": {
            "patterns_detected": [{
                "pattern_id": detected_id,
                "confidence": "moderate",
                "evidence_source": "transcript",
                "evidence": "A concrete production failure opens the narrative.",
                "evidence_citations": [{
                    "source": "transcript",
                    "channel": "transcript",
                    "quote": lines[0],
                }],
            }],
            "antipatterns_detected": [],
            "evidence_sources": ["transcript"],
            "source_inspection": [{
                "source": "transcript",
                "line_ranges": [[1, len(lines)]],
            }],
            "not_evaluable": not_evaluable,
            "pattern_score": {
                "patterns_used": 1,
                "antipatterns_detected": 0,
                "score": 1,
            },
        },
        "catalog_feedback": {
            "unmatched_observations": [],
            "confusable_pairs": [],
            "definition_problems": [],
            "scoring_problems": [],
            "tensions": [],
        },
    }
    return talk, ret


def _complete_unavailable_source_gates(return_validation, ret):
    available = set(ret["pattern_observations"]["evidence_sources"])
    catalog = return_validation.load_catalog()
    ret["pattern_observations"]["not_evaluable"] = [{
        "pattern_id": pattern_id,
        "evidence_source": sorted(available)[0],
        "reason": "The inspected fixture sources cannot evaluate this pattern.",
    } for pattern_id, entry in sorted(catalog.entries.items())
        if entry.observable and entry.evaluable_from is not None and
        not return_validation.qualifying_evidence_groups(
            entry.evaluable_from, available)]
    return ret


def _gradual_comparison_return(
        return_validation, *, version=2, include_delivery=False,
        evidence_sources_used=None):
    """Build one exact-pair legacy/current comparison fixture."""
    sources = ["transcript", "static_slides", "native_deck"]
    if include_delivery:
        sources.append("delivery_video")
    sources.append("source_comparison")
    detection = {
        "pattern_id": "gradual-consistency",
        "confidence": "moderate",
        "evidence_source": "source_comparison",
        "evidence": "The rendered and native deck views agree.",
    }
    if evidence_sources_used is not None:
        detection["evidence_sources_used"] = evidence_sources_used
    ret = _return(return_schema_version=version)
    ret["pattern_observations"].update({
        "patterns_detected": [detection],
        "antipatterns_detected": [],
        "evidence_sources": sources,
        "pattern_score": {
            "patterns_used": 1,
            "antipatterns_detected": 0,
            "score": 1,
        },
    })
    return _complete_unavailable_source_gates(return_validation, ret)


def _visual_row(*, background="current", content_type="title"):
    return {
        "slide_number": 1,
        "background_color_name": background,
        "content_type": content_type,
        "image_composition": "full_bleed_with_text",
        "has_speech_bubble": False,
        "has_starburst": False,
        "has_footer": True,
    }


def test_promotes_queryable_scalars(persist_results):
    talk = _talk()
    persist_results.merge_talk(talk, _return())
    assert talk["slide_count"] == 62
    assert talk["delivery_language"] == "en"
    assert talk["co_presenter"] is False  # boolean false is meaningful, not "empty"
    assert talk["opening_type"] == "demo_cold_open"
    assert talk["illustration_style"] == "comic_book"
    assert talk["pattern_score"] == 1
    assert talk["audience_interaction_count"] == 3


def test_full_structured_data_persisted(persist_results):
    talk = _talk()
    persist_results.merge_talk(talk, _return())
    # The whole block lands, not just the promoted scalars.
    assert talk["structured_data"]["narrative_arc_type"] == "problem_diagnosis_solution"
    assert talk["structured_data"]["slide_design_style"] == "comic_book"


def test_skipped_merge_cannot_mutate_prior_analysis_even_without_validator(
        persist_results):
    talk = _talk(
        processed_date="2026-06-18",
        transcript_source="youtube_auto",
        slide_source="pptx",
        rhetoric_notes="trusted prior analysis",
        structured_data={"slide_count": 42},
        verbatim_examples={"jokes": ["kept"]},
        pattern_scoring_generation_status="current",
        pattern_scoring_generation_reasons=[],
        pattern_scoring_schema_version=2,
        pattern_catalog_fingerprint="0" * 64,
        video_url=None,
        youtube_id=None,
        pptx_path=None,
        slides_url=None,
    )
    malicious = _return(
        status="skipped_no_sources",
        rhetoric_notes=["wrong type"],
        slide_source="not-an-enum",
        structured_data={"video_extraction": {"schema_version": "invalid"}},
        clear_fields=["structured_data.slide_count"],
    )

    persist_results.merge_talk(talk, malicious, run_date="2026-07-31")

    assert talk["status"] == "skipped_no_sources"
    assert talk["rhetoric_notes"] == "trusted prior analysis"
    assert talk["structured_data"] == {"slide_count": 42}
    assert talk["verbatim_examples"] == {"jokes": ["kept"]}
    assert talk["pattern_scoring_generation_status"] == "current"
    assert talk["pattern_scoring_schema_version"] == 2
    assert talk["pattern_catalog_fingerprint"] == "0" * 64
    assert talk["processed_date"] == "2026-06-18"
    assert talk["transcript_source"] == "youtube_auto"
    assert talk["slide_source"] == "pptx"


def test_legacy_deep_merge_is_additive(persist_results):
    talk = _talk(structured_data={"act_structure": {"act_count": 4}})
    ret = _return()
    ret["return_schema_version"] = 1
    ret["structured_data"]["act_structure"] = {"named_acts": True}
    persist_results.merge_talk(talk, ret)
    acts = talk["structured_data"]["act_structure"]
    assert acts["act_count"] == 4  # earlier-run data preserved
    assert acts["named_acts"] is True  # new data merged in


@pytest.mark.parametrize(
    "field", ["slides_local_path", "slides_pdf_path", "pdf_path"])
def test_legacy_return_persists_exact_local_pdf_preclaim_before_drive(
        persist_results, return_validation, tmp_path, field):
    local_path = "slides/descriptive-name.pdf"
    pdf_path = tmp_path / local_path
    pdf_path.parent.mkdir()
    writer = PdfWriter()
    writer.add_blank_page(width=96, height=54)
    with pdf_path.open("wb") as stream:
        writer.write(stream)
    talk = _talk(
        slide_source="pdf",
        google_drive_id="slides-id",
        **{field: local_path},
    )
    ret = _return(slide_source="pdf", slides_local_path=local_path)
    ret["pattern_observations"]["evidence_sources"] = [
        "transcript", "static_slides", "delivery_video"]
    _complete_unavailable_source_gates(return_validation, ret)
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert stored["slides_local_path"] == local_path


def test_nonvideo_return_cannot_persist_video_extraction_manifest(
        persist_results, tmp_path):
    ret = _return(slide_source="pptx")
    ret["structured_data"]["video_extraction"] = {
        "schema_version": 3,
        "artifacts": [{"path": "/outside/untrusted.pdf"}],
    }
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = _db_json({"talks": [_talk()]})
    db.write_text(original)
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "allowed only when slide_source is 'video_extracted'" in result.stderr
    assert db.read_text() == original


@pytest.mark.parametrize(
    ("fault", "expected_error"),
    [
        ("missing_context", "cannot admit returned PDF artifact"),
        ("outside_source", "return video extraction source is unavailable"),
        ("mismatched_promoted", "content digest disagrees"),
    ],
)
def test_legacy_video_return_bounds_every_manifest_artifact_before_merge(
        persist_results, return_validation, tmp_path, fault, expected_error):
    video_id = "AbCdEfGhI_1"
    rebuild = tmp_path / "slides-rebuild" / video_id
    rebuild.mkdir(parents=True)
    source_video = rebuild / f"{video_id}.mp4"
    if fault == "outside_source":
        source_video = (
            tmp_path.parent / f"{tmp_path.name}-outside" / f"{video_id}.mp4"
        )
        source_video.parent.mkdir()
    _write_tiny_mp4(source_video)
    slide_region = rebuild / f"{video_id}.slide-region.pdf"
    promoted = tmp_path / "slides" / f"{video_id}.pdf"
    promoted.parent.mkdir()
    for path in (slide_region, promoted):
        writer = PdfWriter()
        writer.add_blank_page(width=96, height=54)
        with path.open("wb") as stream:
            writer.write(stream)
    if fault == "mismatched_promoted":
        writer = PdfWriter()
        writer.add_blank_page(width=96, height=54)
        writer.add_metadata({"/Title": "Different valid one-page PDF"})
        with promoted.open("wb") as stream:
            writer.write(stream)
    missing_context = rebuild / f"{video_id}.context.pdf"
    if fault != "missing_context":
        writer = PdfWriter()
        writer.add_blank_page(width=96, height=54)
        with missing_context.open("wb") as stream:
            writer.write(stream)
    shared = {
        "page_count": 1,
        "source_video_id": video_id,
        "source_video_path": str(source_video),
    }
    manifest = {
        "slide_source": "video_extracted",
        "schema_version": 3,
        "pipeline_version": "0.10.0",
        "source_video_id": video_id,
        "source_video_path": str(source_video),
        "total_frames_extracted": 1,
        "unique_frame_count": 1,
        "authored_slide_count": None,
        "hash_threshold_used": 8,
        "slide_region_detected": False,
        "slide_region_applied": True,
        "slide_region_method": "manual",
        "slide_region_verified": True,
        "slide_region": [0.05, 0.02, 0.78, 0.98],
        "fps_used": 0.5,
        "retained_frames": [
            {"page_number": 1, "frame_index": 0, "timestamp_seconds": 0.0}
        ],
        "review_required": False,
        "review_reason": None,
        "artifacts": [
            {
                "path": str(slide_region),
                "artifact_scope": "slide_region",
                "crop_method": "manual",
                "crop_verified": True,
                "trusted_for_authored_slide_analysis": True,
                **shared,
            },
            {
                "path": str(missing_context),
                "artifact_scope": "full_frame_context",
                "crop_method": "none",
                "crop_verified": False,
                "trusted_for_authored_slide_analysis": False,
                **shared,
            },
        ],
    }
    talk = _talk()
    ret = _return(
        return_schema_version=2,
        slide_source="video_extracted",
        slides_local_path=f"slides/{video_id}.pdf",
    )
    ret["pattern_observations"]["evidence_sources"] = [
        "transcript",
        "static_slides",
        "delivery_video",
    ]
    _complete_unavailable_source_gates(return_validation, ret)
    ret["structured_data"]["video_extraction"] = manifest
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    original = db.read_bytes()
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode != 0
    assert expected_error in result.stderr
    assert db.read_bytes() == original


def test_v2_adds_only_inside_registered_extension_namespace(persist_results):
    talk = _talk(structured_data={
        "extensions": {"argument_map": {"act_count": 4}},
    })
    ret = _return()
    ret["structured_data"]["extensions"] = {
        "argument_map": {"named_acts": True},
    }

    persist_results.merge_talk(talk, ret)

    assert talk["structured_data"]["extensions"] == {
        "argument_map": {"act_count": 4, "named_acts": True},
    }


@pytest.mark.parametrize("field", [
    "color_coded_backgrounds",
    "typography_observations",
    "footer_observations",
    "shape_observations",
    "video_extraction",
    "key_data_points",
    "named_authorities",
    "time_bound_promotion",
    "native_deck_audit",
    "native_timing_audit",
    "source_comparison",
    "source_identity",
    "animation_observations",
    "pptx_pdf_reconciliation",
])
def test_v2_registered_maps_are_atomic_snapshots(persist_results, field):
    talk = _talk(structured_data={
        field: {"current": {"value": 1}, "stale_child": True},
    })
    ret = _return()
    ret["structured_data"][field] = {"current": {"value": 2}}

    persist_results.merge_talk(talk, ret)

    assert talk["structured_data"][field] == {"current": {"value": 2}}


@pytest.mark.parametrize(
    ("claim_version", "return_version", "closed_version"),
    [(1, 1, 2), (2, 2, 2), (3, 3, 3)],
)
def test_queue_claim_closure_preserves_the_version_matrix(
        persist_results, claim_version, return_version, closed_version):
    if claim_version == 3:
        talk, ret = _v3_talk_and_return(persist_results)
    else:
        talk = _talk()
        talk["_queue_claim"]["schema_version"] = claim_version
        ret = _return(return_schema_version=return_version)

    persist_results.merge_talk(
        talk,
        ret,
        run_date="2026-07-31T18:05:00+00:00",
        enforce_queue_claim=True,
    )

    claim = talk["_queue_claim"]
    assert claim["schema_version"] == closed_version
    assert claim["state"] == "completed"
    assert claim["result_payload_sha256"] == \
        persist_results.canonical_return_sha256(ret)


def test_v2_snapshot_clears_stale_authenticated_adherence_comparison(
        persist_results):
    stale = {
        "schema_version": 1,
        "baseline": {"untrusted": "stale"},
        "talk_pattern_score": 99,
    }
    talk = _talk(adherence_comparison=copy.deepcopy(stale))

    persist_results.merge_talk(talk, _return(return_schema_version=2))

    assert "adherence_comparison" not in talk

    legacy = _talk(adherence_comparison=copy.deepcopy(stale))
    persist_results.merge_talk(legacy, _return(return_schema_version=1))
    assert legacy["adherence_comparison"] == stale


def test_v2_atomic_snapshot_matches_rendered_analysis(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    analyses = tmp_path / "analyses"
    talk = _talk(structured_data={
        "typography_observations": {
            "current_marker": "legacy",
            "stale_child": "must disappear",
        },
    })
    ret = _return()
    ret["structured_data"]["typography_observations"] = {
        "current_marker": "current",
    }
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([ret]))

    persisted = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )
    assert persisted.returncode == 0, persisted.stderr
    analysis_script = persist_results.__file__.replace(
        "persist-results.py", "write-analysis.py")
    rendered = subprocess.run(
        [sys.executable, analysis_script, str(batch), str(analyses),
         "--talks", str(db)],
        capture_output=True,
        text=True,
    )

    assert rendered.returncode == 0, rendered.stderr
    stored_map = json.loads(db.read_text())["talks"][0]["structured_data"][
        "typography_observations"]
    markdown = (analyses / "talk.md").read_text()
    assert stored_map == ret["structured_data"]["typography_observations"]
    assert "current_marker" in markdown
    assert "current" in markdown
    assert "stale_child" not in markdown


def test_v2_image_distribution_and_basis_replace_as_one_snapshot(persist_results):
    talk = _talk(structured_data={
        "image_source_distribution": {"legacy": 9, "unknown": 1},
        "image_source_distribution_basis": "Unit: legacy asset.",
    })
    ret = _return()
    ret["structured_data"].update({
        "image_source_distribution": {"unknown": 2},
        "image_source_distribution_basis": "Unit: slide; unknown origins stay unknown.",
    })

    persist_results.merge_talk(talk, ret)

    assert talk["structured_data"]["image_source_distribution"] == {"unknown": 2}
    assert talk["structured_data"]["image_source_distribution_basis"] == (
        "Unit: slide; unknown origins stay unknown.")


def test_v2_per_slide_ledger_replaces_atomically(persist_results):
    old_row = _visual_row(background="legacy")
    old_row["legacy_note"] = "must disappear"
    talk = _talk(structured_data={"per_slide_visual": [old_row]})
    ret = _return()
    ret["structured_data"].update({
        "slide_count": 1,
        "per_slide_visual": [_visual_row()],
    })

    persist_results.merge_talk(talk, ret)

    assert talk["structured_data"]["per_slide_visual"] == [_visual_row()]


def test_v2_empty_verbatim_and_structured_lists_replace_prior_findings(
        persist_results):
    talk = _talk(
        structured_data={"visual_continuity_devices": ["legacy_device"]},
        verbatim_examples={"jokes": ["legacy line"]},
    )
    ret = _return(verbatim_examples={"jokes": []})
    ret["structured_data"]["visual_continuity_devices"] = []

    persist_results.merge_talk(talk, ret)

    assert talk["structured_data"]["visual_continuity_devices"] == []
    assert talk["verbatim_examples"]["jokes"] == []


def test_v2_empty_scalar_replaces_prior_assessment(persist_results):
    talk = _talk(adherence_assessment="legacy assessment")

    persist_results.merge_talk(talk, _return(adherence_assessment=""))

    assert talk["adherence_assessment"] == ""


def test_v2_empty_promoted_list_stays_synchronized(persist_results):
    talk = _talk(
        co_presenter=True,
        co_presenters=["Legacy Name"],
        structured_data={
            "co_presenter": True,
            "co_presenters": ["Legacy Name"],
        },
    )
    ret = _return()
    ret["structured_data"].update({
        "co_presenter": False,
        "co_presenters": [],
    })

    persist_results.merge_talk(talk, ret)

    assert talk["structured_data"]["co_presenter"] is False
    assert talk["structured_data"]["co_presenters"] == []
    assert talk["co_presenter"] is False
    assert talk["co_presenters"] == []


def test_v2_unknown_object_fails_without_mutating_talk(persist_results):
    talk = _talk(structured_data={"slide_count": 4})
    original = copy.deepcopy(talk)
    ret = _return()
    ret["structured_data"]["undeclared_snapshot"] = {"value": 1}

    with pytest.raises(ValueError, match="unregistered object"):
        persist_results.merge_talk(talk, ret)

    assert talk == original


@pytest.mark.parametrize(
    ("field", "old_value", "new_value"),
    [
        ("typography_observations", ["legacy"], {"current": True}),
        ("visual_continuity_devices", {"legacy": True}, []),
    ],
)
def test_v2_nonempty_type_changes_replace_exactly(
        persist_results, field, old_value, new_value):
    talk = _talk(structured_data={field: old_value})
    ret = _return()
    ret["structured_data"][field] = new_value

    persist_results.merge_talk(talk, ret)

    assert talk["structured_data"][field] == new_value


def test_v2_wrong_type_never_preserves_old_value_or_mutates_talk(persist_results):
    talk = _talk(structured_data={"typography_observations": {"legacy": True}})
    original = copy.deepcopy(talk)
    ret = _return()
    ret["structured_data"]["typography_observations"] = []

    with pytest.raises(ValueError, match="must be an object"):
        persist_results.merge_talk(talk, ret)

    assert talk == original


def test_v2_repairs_legacy_verbatim_and_pattern_array_containers(persist_results):
    talk = _talk(
        verbatim_examples=["legacy"],
        pattern_observations=["legacy"],
    )

    persist_results.merge_talk(talk, _return(verbatim_examples={}))

    assert talk["verbatim_examples"] == {}
    assert isinstance(talk["pattern_observations"], dict)
    assert talk["pattern_observations"]["pattern_ids"] == [
        "narrative-arc", "bookends"]


def test_v2_clear_can_remove_an_undeclared_legacy_verbatim_lane(persist_results):
    talk = _talk(verbatim_examples={
        "jokes": ["legacy line"],
        "legacy_lane": ["undeclared"],
    })
    ret = _return(
        clear_fields=["verbatim_examples.legacy_lane"],
        verbatim_examples={"jokes": []},
    )

    persist_results.merge_talk(talk, ret)

    assert talk["verbatim_examples"] == {"jokes": []}


def test_v2_malformed_structured_container_fails_before_mutation(persist_results):
    talk = _talk(structured_data=["legacy"])
    original = copy.deepcopy(talk)

    with pytest.raises(ValueError, match="stored structured_data"):
        persist_results.merge_talk(talk, _return())

    assert talk == original


def test_v2_pattern_snapshot_removes_stale_children_and_accepts_zero_score(
        persist_results):
    talk = _talk(
        pattern_score=7,
        pattern_observations={
            "pattern_score": 7,
            "pattern_ids": ["legacy"],
            "legacy_note": "must disappear",
        },
    )
    ret = _return()
    ret["pattern_observations"].update({
        "patterns_detected": [],
        "antipatterns_detected": [],
        "not_evaluable": [],
        "pattern_score": {
            "patterns_used": 0,
            "antipatterns_detected": 0,
            "score": 0,
        },
    })

    persist_results.merge_talk(talk, ret)

    assert talk["pattern_score"] == 0
    assert talk["pattern_observations"] == {
        "patterns_detected": [],
        "pattern_ids": [],
        "antipatterns_detected": [],
        "antipattern_ids": [],
        "not_evaluable": [],
        "not_evaluable_ids": [],
        "evidence_sources": [
            "transcript", "native_deck", "static_slides", "delivery_video",
            "source_comparison"],
        "pattern_score": 0,
    }


def test_legacy_omitted_verbatim_lane_preserves_prior_value(persist_results):
    talk = _talk(verbatim_examples={"jokes": ["legacy line"]})
    ret = _return(return_schema_version=1, verbatim_examples={})

    persist_results.merge_talk(talk, ret)

    assert talk["verbatim_examples"]["jokes"] == ["legacy line"]


def test_video_extraction_manifest_replaces_legacy_manifest_atomically(
        persist_results):
    talk = _talk(structured_data={
        "video_extraction": {
            "schema_version": 2,
            "output_pdf": "/legacy/full-frame.pdf",
            "unique_slides_count": 80,
        },
    })
    ret = _return()
    ret["structured_data"]["video_extraction"] = {
        "schema_version": 3,
        "artifacts": [{"artifact_scope": "slide_region"}],
    }
    persist_results.merge_talk(talk, ret)
    assert talk["structured_data"]["video_extraction"] == {
        "schema_version": 3,
        "artifacts": [{"artifact_scope": "slide_region"}],
    }


def test_legacy_empty_values_never_clobber(persist_results):
    talk = _talk(structured_data={"slide_count": 62})
    ret = _return()
    ret["return_schema_version"] = 1
    ret["structured_data"]["slide_count"] = None  # empty must not overwrite
    persist_results.merge_talk(talk, ret)
    assert talk["structured_data"]["slide_count"] == 62


def test_legacy_empty_substantive_prose_remains_a_noop(persist_results):
    talk = _talk(
        rhetoric_notes="trusted prior notes",
        areas_for_improvement="trusted prior improvements",
    )
    ret = _return(
        return_schema_version=1,
        rhetoric_notes="",
        areas_for_improvement="",
    )

    persist_results.merge_talk(talk, ret)

    assert talk["rhetoric_notes"] == "trusted prior notes"
    assert talk["areas_for_improvement"] == "trusted prior improvements"


def test_pattern_observations_normalized(persist_results):
    talk = _talk()
    persist_results.merge_talk(talk, _return())
    obs = talk["pattern_observations"]
    assert obs["pattern_ids"] == ["narrative-arc", "bookends"]
    assert obs["antipattern_ids"] == ["shortchanged"]
    assert obs["pattern_score"] == 1  # flattened from {"score": 1}
    assert len(obs["patterns_detected"]) == 2  # detailed arrays kept for Section 15
    assert obs["evidence_sources"] == [
        "transcript", "native_deck", "static_slides", "delivery_video",
        "source_comparison"]
    assert obs["not_evaluable"] == []
    assert obs["not_evaluable_ids"] == []


def test_transcript_quote_is_verified_and_locations_are_engine_owned(
        persist_results, return_validation, tmp_path):
    talk, ret = _v4_transcript_batch(return_validation, tmp_path)
    talk["schema_version"] = persist_results.TALK_SCHEMA_VERSION
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [
            sys.executable,
            persist_results.__file__,
            str(db),
            str(batch),
            "--run-date",
            "2026-07-31T18:05:00+00:00",
        ],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    observations = stored["pattern_observations"]
    citation = observations["patterns_detected"][0]["evidence_citations"][0]
    assert citation["line_start"] == 1
    assert citation["line_end"] == 1
    assert citation["artifact_root"] == "vault"
    assert citation["artifact_path"] == "transcripts/manual-talk.txt"
    assert len(citation["artifact_sha256"]) == 64
    assert observations["evidence_schema_version"] == 1
    assert observations["source_inspection"][0]["coverage_complete"] is True
    assert stored["pattern_scoring_generation_status"] == \
        "legacy_unbaselineable"
    assert "return_schema_precedes_exhaustive_outcomes" in stored[
        "pattern_scoring_generation_reasons"
    ]
    assert "pattern_scoring_schema_version" not in stored
    assert "pattern_catalog_fingerprint" not in stored
    # The worker receipt remains the exact raw claim; engine-owned locations
    # exist only in the canonical persisted projection.
    raw = json.loads(batch.read_text())[0]
    assert "line_start" not in raw["pattern_observations"][
        "patterns_detected"][0]["evidence_citations"][0]


@pytest.mark.parametrize("version", [1, 2])
def test_legacy_single_comparison_pair_is_inferred_but_remains_historical(
        persist_results, return_validation, version):
    talk = _talk()
    ret = _gradual_comparison_return(return_validation, version=version)

    persist_results.merge_talk(talk, ret)

    detection = talk["pattern_observations"]["patterns_detected"][0]
    assert detection["evidence_sources_used"] == [
        "static_slides", "native_deck"]
    assert talk["pattern_scoring_generation_status"] == \
        "legacy_unbaselineable"
    assert talk["pattern_scoring_generation_reasons"] == [
        "return_schema_precedes_source_locations"]
    assert "pattern_scoring_schema_version" not in talk
    assert "pattern_catalog_fingerprint" not in talk


@pytest.mark.parametrize("version", [1, 2])
def test_legacy_ambiguous_comparison_replays_but_is_excluded_from_baselines(
        persist_results, return_validation, version):
    talk = _talk(
        pattern_scoring_schema_version=2,
        pattern_catalog_fingerprint="0" * 64,
    )
    ret = _gradual_comparison_return(
        return_validation, version=version, include_delivery=True)

    persist_results.merge_talk(talk, ret)

    detection = talk["pattern_observations"]["patterns_detected"][0]
    assert "evidence_sources_used" not in detection
    assert talk["pattern_scoring_generation_status"] == \
        "legacy_unbaselineable"
    assert talk["pattern_scoring_generation_reasons"] == [
        "comparison_group_ambiguous:gradual-consistency",
        "return_schema_precedes_source_locations",
    ]
    assert "pattern_scoring_schema_version" not in talk
    assert "pattern_catalog_fingerprint" not in talk


@pytest.mark.parametrize("version", [1, 2])
def test_legacy_new_strong_gate_failure_replays_without_current_fingerprint(
        persist_results, version):
    talk = _talk(
        pattern_scoring_schema_version=2,
        pattern_catalog_fingerprint="0" * 64,
    )
    ret = _return(return_schema_version=version)
    ret["pattern_observations"].update({
        "patterns_detected": [{
            "pattern_id": "traveling-highlights",
            "confidence": "strong",
            "evidence_source": "static_slides",
            "evidence": "The highlight appears in the static rendering.",
        }],
        "antipatterns_detected": [],
        "pattern_score": {
            "patterns_used": 1,
            "antipatterns_detected": 0,
            "score": 1,
        },
    })

    persist_results.merge_talk(talk, ret)

    assert talk["pattern_scoring_generation_status"] == \
        "legacy_unbaselineable"
    assert talk["pattern_scoring_generation_reasons"] == [
        "return_schema_precedes_source_locations",
        "strong_gate_unsatisfied:traveling-highlights",
    ]
    assert "pattern_scoring_schema_version" not in talk
    assert "pattern_catalog_fingerprint" not in talk


def test_direct_v3_ineligible_merge_replays_as_historical(
        persist_results):
    talk = _talk()
    ret = _return(return_schema_version=3)
    ret["pattern_observations"].update({
        "patterns_detected": [{
            "pattern_id": "traveling-highlights",
            "confidence": "strong",
            "evidence_source": "static_slides",
            "evidence": "Static pages cannot establish a strong detection.",
        }],
        "antipatterns_detected": [],
        "pattern_score": {
            "patterns_used": 1,
            "antipatterns_detected": 0,
            "score": 1,
        },
    })

    persist_results.merge_talk(talk, ret)

    assert talk["pattern_scoring_generation_status"] == \
        "legacy_unbaselineable"
    assert talk["pattern_scoring_generation_reasons"] == [
        "return_schema_precedes_source_locations",
        "strong_gate_unsatisfied:traveling-highlights",
    ]


def test_direct_catalog_fingerprint_mismatch_fails_without_mutating_talk(
        persist_results, return_validation):
    talk = _talk()
    original = copy.deepcopy(talk)

    with pytest.raises(ValueError, match="does not match the catalog"):
        persist_results.merge_talk(
            talk,
            _return(),
            catalog_fingerprint="0" * 64,
            catalog=return_validation.load_catalog(),
        )

    assert talk == original


@pytest.mark.parametrize(
    ("return_schema_version", "expected"),
    [(1, "legacy assessment"), (2, ""), (3, "")],
)
def test_only_v2_and_v3_use_snapshot_scalar_semantics(
        persist_results, return_schema_version, expected):
    talk = _talk(adherence_assessment="legacy assessment")
    ret = _return(
        return_schema_version=return_schema_version,
        adherence_assessment="",
    )

    persist_results.merge_talk(talk, ret)

    assert talk["adherence_assessment"] == expected


def test_not_evaluable_patterns_are_persisted_but_never_scored(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["not_evaluable"] = [{
        "pattern_id": "composite-animation",
        "evidence_source": "static_slides",
        "reason": "No animation timing survives in the rendered pages.",
    }]
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    obs = json.loads(db.read_text())["talks"][0]["pattern_observations"]
    assert obs["not_evaluable_ids"] == ["composite-animation"]
    assert obs["pattern_score"] == 1


def test_scalar_result_fields_copied(persist_results):
    talk = _talk()
    persist_results.merge_talk(
        talk, _return(slides_local_path="slides/source.pdf"))
    assert talk["status"] == "processed"
    assert talk["processed_date"] == "2026-06-18"
    assert talk["rhetoric_notes"] == "notes"
    assert talk["transcript_source"] == "youtube_auto"
    assert talk["slides_local_path"] == "slides/source.pdf"


def test_run_date_stamped_when_return_omits_processed_date(persist_results):
    """The reparse regression: a return that reports status but no date left the
    previous run's date in place, so the DB could not say which talks it covered."""
    ret = _return()
    del ret["processed_date"]
    talk = _talk(processed_date="2026-04-09")
    _, stamped, _, _ = persist_results.merge_talk(talk, ret, run_date="2026-07-26")
    assert talk["processed_date"] == "2026-07-26"
    assert stamped is True


def test_batch_run_date_wins_over_legacy_return_date(persist_results):
    talk = _talk()
    _, stamped, _, _ = persist_results.merge_talk(
        talk, _return(), run_date="2026-07-26")
    assert talk["processed_date"] == "2026-07-26"
    assert stamped is True


def test_empty_processed_date_is_stamped(persist_results):
    talk = _talk(processed_date="2026-04-09")
    _, stamped, _, _ = persist_results.merge_talk(
        talk, _return(processed_date=""), run_date="2026-07-26")
    assert talk["processed_date"] == "2026-07-26"
    assert stamped is True


def test_no_run_date_leaves_processed_date_untouched(persist_results):
    ret = _return()
    del ret["processed_date"]
    talk = _talk(processed_date="2026-04-09")
    _, stamped, _, _ = persist_results.merge_talk(talk, ret)
    assert talk["processed_date"] == "2026-04-09"
    assert stamped is False


def test_cli_run_date_pins_the_stamp(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    del ret["processed_date"]
    db.write_text(_db_json({"talks": [_talk(processed_date="2026-04-09")]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch),
         "--run-date", "2026-07-26"],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    assert json.loads(db.read_text())["talks"][0]["processed_date"] == "2026-07-26"
    report = json.loads(result.stdout)
    assert report["run_date"] == "2026-07-26"
    assert report["talks"][0]["stamped_processed_date"] is True


def test_cli_batch_timestamp_overrides_date_only_return_stamp_everywhere(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    authoritative = "2026-07-27T14:03:22+00:00"
    db.write_text(_db_json({"talks": [_talk(processed_date="2026-04-09")]}))
    batch.write_text(json.dumps([_return(processed_date="2026-07-27")]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch),
         "--run-date", authoritative],
        capture_output=True, text=True,
    )

    assert result.returncode == 0, result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert stored["processed_date"] == authoritative
    assert stored["_queue_claim"]["released_at"] == authoritative
    assert json.loads(result.stdout)["talks"][0]["stamped_processed_date"] is True


def test_cli_rejects_conflicting_explicit_return_timestamp_before_write(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = {"talks": [_talk(processed_date="2026-04-09")]}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([
        _return(processed_date="2026-07-27T08:00:00+00:00")]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch),
         "--run-date", "2026-07-27T14:03:22+00:00"],
        capture_output=True, text=True,
    )

    assert result.returncode == 1
    assert "explicit timestamp" in result.stderr
    assert "conflicts with authoritative batch run_date" in result.stderr
    assert json.loads(db.read_text()) == original


def test_matching_explicit_return_timestamp_is_accepted_and_batch_stamped(
        persist_results):
    talk = _talk(processed_date="2026-04-09")
    _, stamped, _, _ = persist_results.merge_talk(
        talk,
        _return(processed_date="2026-07-27T16:03:22.987+02:00"),
        run_date="2026-07-27T14:03:22+00:00",
    )
    assert talk["processed_date"] == "2026-07-27T14:03:22+00:00"
    assert stamped is True


def test_cli_rejects_malformed_run_date(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch),
         "--run-date", "26-07-2026"],
        capture_output=True, text=True,
    )
    assert result.returncode != 0
    assert "YYYY-MM-DD" in result.stderr


def test_cli_writes_db_and_reports(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    script = persist_results.__file__
    result = subprocess.run(
        [sys.executable, script, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    out = json.loads(db.read_text())["talks"][0]
    assert out["slide_count"] == 62
    # Structured JSON summary on stdout (not prose), per script-delegation.
    report = json.loads(result.stdout)
    assert report["persisted"] == 1
    assert report["talks"][0]["filename"] == "talk.md"
    assert "slide_count" in report["talks"][0]["promoted"]


@pytest.mark.parametrize("config_mode", ["absent", "null", "exact"])
def test_cli_accepts_database_bound_vault_root_authority(
    persist_results,
    tmp_path,
    config_mode,
):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    database = {"talks": [_talk()]}
    _db_json(database)
    if config_mode == "null":
        database["config"]["vault_storage_path"] = None
    elif config_mode == "exact":
        database["config"]["vault_storage_path"] = str(tmp_path)
    db.write_text(json.dumps(database), encoding="utf-8")
    batch.write_text(json.dumps([_return()]), encoding="utf-8")

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    assert json.loads(db.read_text(encoding="utf-8"))["talks"][0]["status"] == (
        "processed"
    )


@pytest.mark.skipif(
    sys.platform == "win32",
    reason="directory symlink setup is privileged",
)
@pytest.mark.parametrize("configured_identity", ["alias", "physical"])
def test_cli_compares_symlinked_vault_roots_by_lexical_identity_before_persistence(
    persist_results,
    tmp_path,
    configured_identity,
):
    physical_root = tmp_path / "physical-vault"
    physical_root.mkdir()
    alias_root = tmp_path / "alias-vault"
    alias_root.symlink_to(physical_root, target_is_directory=True)
    db = physical_root / "tracking-database.json"
    alias_db = alias_root / db.name
    batch = tmp_path / "batch-returns.json"
    database = {"talks": [_talk()]}
    _db_json(database)
    configured_root = (
        alias_root if configured_identity == "alias" else physical_root
    )
    database["config"]["vault_storage_path"] = str(configured_root)
    db.write_text(json.dumps(database), encoding="utf-8")
    batch.write_text(json.dumps([_return()]), encoding="utf-8")
    before = db.read_bytes()

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(alias_db), str(batch)],
        capture_output=True,
        text=True,
    )

    if configured_identity == "alias":
        assert result.returncode == 0, result.stderr
        assert json.loads(db.read_text(encoding="utf-8"))["talks"][0][
            "status"
        ] == "processed"
    else:
        assert result.returncode == 1
        assert result.stderr == (
            "ERROR: vault_root_authority_mismatch:database_path:config_root\n"
        )
        assert db.read_bytes() == before


FOREIGN_ABSOLUTE_VAULT_ROOT = (
    "/foreign/vault" if sys.platform == "win32" else r"C:\foreign\vault"
)
NATIVE_DOT_VAULT_ROOT = (
    r"C:\trusted\other\..\vault"
    if sys.platform == "win32"
    else "/trusted/other/../vault"
)


@pytest.mark.parametrize(
    ("configured_root", "locator_reason"),
    [
        ("", "artifact_locator_empty_or_whitespace"),
        (" ", "artifact_locator_empty_or_whitespace"),
        ("relative/vault", "artifact_root_not_native_absolute"),
        ("C:vault", "artifact_locator_windows_drive_relative"),
        (r"\vault", "artifact_locator_windows_current_drive_rooted"),
        (NATIVE_DOT_VAULT_ROOT, "artifact_locator_dot_segment"),
        ("~/private-vault", "artifact_locator_home_expansion_unsupported"),
        (FOREIGN_ABSOLUTE_VAULT_ROOT, "artifact_locator_foreign_absolute"),
        (r"\\?\C:\private-vault", "artifact_locator_windows_device_namespace"),
    ],
)
def test_cli_rejects_invalid_configured_vault_root_before_persistence(
    persist_results,
    tmp_path,
    configured_root,
    locator_reason,
):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    database = {"talks": [_talk()]}
    _db_json(database)
    database["config"]["vault_storage_path"] = configured_root
    original = json.dumps(database)
    db.write_text(original, encoding="utf-8")
    batch.write_text(json.dumps([_return()]), encoding="utf-8")

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert result.stderr == (
        f"ERROR: vault_root_config_invalid:{locator_reason}\n"
    )
    assert "Traceback" not in result.stderr
    if configured_root.strip():
        assert configured_root not in result.stderr
    assert db.read_text(encoding="utf-8") == original


def test_cli_rejects_configured_vault_root_authority_mismatch(
    persist_results,
    tmp_path,
):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    database = {"talks": [_talk()]}
    _db_json(database)
    mismatched_root = tmp_path / "other-vault"
    database["config"]["vault_storage_path"] = str(mismatched_root)
    original = json.dumps(database)
    db.write_text(original, encoding="utf-8")
    batch.write_text(json.dumps([_return()]), encoding="utf-8")

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert result.stderr == (
        "ERROR: vault_root_authority_mismatch:database_path:config_root\n"
    )
    assert str(mismatched_root) not in result.stderr
    assert db.read_text(encoding="utf-8") == original


def test_cli_rejects_relative_database_authority_before_open(
    persist_results,
    tmp_path,
):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    database = {"talks": [_talk()]}
    original = _db_json(database)
    db.write_text(original, encoding="utf-8")
    batch.write_text(json.dumps([_return()]), encoding="utf-8")

    result = subprocess.run(
        [sys.executable, persist_results.__file__, db.name, str(batch)],
        cwd=tmp_path,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert result.stderr == (
        "ERROR: vault_root_database_path_invalid:"
        "artifact_root_not_native_absolute\n"
    )
    assert db.name not in result.stderr
    assert db.read_text(encoding="utf-8") == original


def test_root_authority_rejection_precedes_every_artifact_and_write_boundary(
    persist_results,
    tmp_path,
    monkeypatch,
    capsys,
):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    database = {"talks": [_talk()]}
    _db_json(database)
    database["config"]["vault_storage_path"] = "~/private-vault"
    original = json.dumps(database)
    db.write_text(original, encoding="utf-8")
    batch.write_text(json.dumps([_return()]), encoding="utf-8")
    called = []

    def forbidden(name):
        def invoke(*_args, **_kwargs):
            called.append(name)
            raise AssertionError(f"{name} ran after root-authority rejection")

        return invoke

    for name in (
        "load_json",
        "validate_batch",
        "assess_batch_artifact_capabilities",
        "validate_batch_claims_against_talks",
        "admit_return_artifacts",
        "canonicalize_return_evidence",
        "assess_current_persisted_pattern_evidence_freshness",
        "atomic_write_json",
    ):
        monkeypatch.setattr(persist_results, name, forbidden(name))
    monkeypatch.setattr(
        sys,
        "argv",
        [persist_results.__file__, str(db), str(batch)],
    )

    with pytest.raises(SystemExit) as caught:
        persist_results.main()

    assert caught.value.code == 1
    assert called == []
    assert capsys.readouterr().err == (
        "ERROR: vault_root_config_invalid:"
        "artifact_locator_home_expansion_unsupported\n"
    )
    assert db.read_text(encoding="utf-8") == original


def test_cli_fails_visibly_on_filename_mismatch(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk(filename="a.md")]}))
    batch.write_text(json.dumps([_return(filename="missing.md")]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "must exactly match" in result.stderr
    assert "unexpected ['missing.md']" in result.stderr


def test_cli_missing_input_file_is_actionable(persist_results, tmp_path):
    batch = tmp_path / "batch-returns.json"
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(tmp_path / "nope.json"), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "tracking database file not found" in result.stderr
    assert "Traceback" not in result.stderr  # no raw traceback


def test_cli_malformed_json_is_actionable(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text("{ not valid json ")
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "not valid JSON" in result.stderr
    assert "Traceback" not in result.stderr


def test_cli_non_array_batch_is_actionable(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps({"filename": "talk.md"}))  # object, not array
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "must be a JSON array" in result.stderr


def test_default_stamp_is_second_resolution(persist_results):
    """A day-granular stamp cannot order a talk against a fix that shipped the
    same day — during the 2026-07-26 reparse, 90 talks shared one date and the
    re-check backlog had to flag every one because ordering was unknowable.

    The clock is injected, not read: a test asserting a stamp produced from the
    live wall clock is non-deterministic by construction."""
    frozen = datetime(2026, 7, 27, 14, 3, 22, 456789, tzinfo=timezone.utc)
    assert persist_results.default_stamp(frozen) == "2026-07-27T14:03:22+00:00"


def test_default_stamp_normalizes_to_utc(persist_results):
    """Stamps written on machines in different zones must order against each other."""
    frozen = datetime(2026, 7, 27, 16, 3, 22, tzinfo=timezone(timedelta(hours=2)))
    assert persist_results.default_stamp(frozen) == "2026-07-27T14:03:22+00:00"


def test_explicit_offset_timestamp_is_normalized_to_utc(persist_results):
    assert persist_results.normalize_stamp(
        "2026-07-27T16:03:22.987654+02:00") == "2026-07-27T14:03:22+00:00"


def test_return_processed_timestamp_is_persisted_in_canonical_utc(persist_results):
    talk = _talk()

    persist_results.merge_talk(
        talk,
        _return(processed_date="2026-07-27T16:03:22.987654+02:00"),
    )

    assert talk["processed_date"] == "2026-07-27T14:03:22+00:00"


def test_naive_timestamp_is_rejected_with_an_actionable_message(persist_results):
    """A naive timestamp cannot be ordered against one from another machine."""
    with pytest.raises(ValueError) as excinfo:
        persist_results.normalize_stamp("2026-07-27T14:03:22")
    assert "+00:00" in str(excinfo.value)


def test_cli_rejects_a_naive_timestamp(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch),
         "--run-date", "2026-07-27T14:03:22"],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "timezone-aware" in result.stderr


def test_explicit_date_only_stamp_is_still_accepted(persist_results, tmp_path):
    """Callers pinning a stamp for reproducibility keep working, and records
    written before this change stay readable."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    del ret["processed_date"]
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch),
         "--run-date", "2026-07-26"],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    assert json.loads(db.read_text())["talks"][0]["processed_date"] == "2026-07-26"


def test_explicit_timestamp_is_accepted(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    del ret["processed_date"]
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch),
         "--run-date", "2026-07-27T14:03:22+00:00"],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    assert json.loads(db.read_text())["talks"][0]["processed_date"].startswith("2026-07-27T14:03:22")


def test_bare_int_pattern_score_is_coerced_and_promoted(persist_results, tmp_path):
    """The bare int is not harmless: PROMOTE digs
    `pattern_observations.pattern_score.score`, `dig` returns None on an int, and
    the queryable top-level scalar is silently dropped — the exact missing-scalar
    defect this script exists to fix, reintroduced through the input shape.

    Roughly a third of returns arrive this way; the schema invites it.
    """
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["pattern_score"] = 1  # 2 patterns - 1 antipattern
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["talks"][0]["coerced_pattern_score"] is True
    assert "pattern_score" in payload["talks"][0]["promoted"]
    talk = json.loads(db.read_text())["talks"][0]
    assert talk["pattern_score"] == 1


def test_dict_pattern_score_is_not_reported_as_coerced(persist_results, tmp_path):
    """Guard the guard: the flag must distinguish the two shapes."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["talks"][0]["coerced_pattern_score"] is False
    assert "pattern_score" in payload["talks"][0]["promoted"]


def test_bare_int_contradicting_the_arrays_fails_loudly(persist_results, tmp_path):
    """A coerced value that disagrees with the arrays is a real inconsistency.

    Recomputing silently would launder it; the script refuses to guess which
    number is right.
    """
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["pattern_score"] = 42  # arrays say 2 - 1 = 1
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "2 patterns minus 1 antipatterns is 1" in result.stderr
    assert "42" in result.stderr


@pytest.mark.parametrize("bad", [True, False, "19", ["19"], 1.5, 1.0])
def test_non_numeric_pattern_score_never_reaches_the_db(persist_results, tmp_path, bad):
    """Assert the persisted OUTCOME, not the helper.

    An earlier version of this test called `canonicalize_pattern_score` directly
    and passed while `merge_talk` still persisted `pattern_score: True` — because
    `isinstance(True, int)` holds in Python, so a bool sailed through
    `normalize_pattern_observations`'s numeric branch. Testing the helper in
    isolation is exactly what hid the bug.
    """
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["pattern_score"] = bad
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1, f"{bad!r} was accepted: {result.stdout}"
    assert "pattern_score" in result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert "pattern_score" not in stored, f"{bad!r} reached the DB as {stored.get('pattern_score')!r}"


@pytest.mark.parametrize("inner", [True, False, "19", ["19"], 1.5, 1.0])
def test_invalid_score_inside_the_declared_dict_is_rejected(
        persist_results, tmp_path, inner):
    """The dict is the declared shape, but its `score` is what PROMOTE writes.

    Type-checking only the bare form left `{"score": True}` reaching the DB
    unexamined — the same defect one level in.
    """
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["pattern_score"] = {
        "patterns_used": 2, "antipatterns_detected": 1, "score": inner}
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1, f"{inner!r} was accepted: {result.stdout}"
    assert "pattern_score.score" in result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert "pattern_score" not in stored


def test_merge_stamps_the_talk_schema_version(persist_results, tmp_path):
    """stateful-artifacts requires a schema_version on every record."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    talk = json.loads(db.read_text())["talks"][0]
    assert talk["schema_version"] == persist_results.TALK_SCHEMA_VERSION


def test_persistence_accepts_historical_talk_under_current_root(
        persist_results, tmp_path):
    """Current roots may retain v1-v4 talks until that exact talk is persisted."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    old = _talk()
    old["schema_version"] = 1
    db.write_text(_db_json({"talks": [old]}))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    assert (
        json.loads(db.read_text())["talks"][0]["schema_version"]
        == persist_results.TALK_SCHEMA_VERSION
    )


def test_persistence_refuses_legacy_database_without_rewriting(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(json.dumps({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    before = db.read_bytes()

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "migrate-tracking-database.py" in result.stderr
    assert db.read_bytes() == before


def test_root_migration_preserves_historical_pattern_evidence(tracking_database):
    talk = {
        "filename": "legacy.md",
        "schema_version": 2,
        "status": "processed",
        "pattern_observations": {
            "patterns_detected": [{
                "pattern_id": "narrative-arc",
                "confidence": "strong",
                "evidence": "Legacy prose without a source location.",
            }],
            "antipatterns_detected": [{
                "pattern_id": "shortchanged",
                "confidence": "moderate",
                "evidence": "A legacy model-supplied location is untrusted.",
                "evidence_citations": [{
                    "source": "transcript",
                    "channel": "transcript",
                    "line_start": 999,
                    "line_end": 999,
                }],
            }],
        },
    }
    database = {"talks": [talk]}

    before = copy.deepcopy(talk)

    migrated = tracking_database.migrate_tracking_database(database).database

    assert migrated["talks"][0] == before


def test_future_talk_schema_rejects_before_any_record_is_migrated(
        persist_results, tracking_database):
    database = {
        "talks": [
            {"filename": "legacy.md", "schema_version": 1},
            {
                "filename": "future.md",
                "schema_version": persist_results.TALK_SCHEMA_VERSION + 1,
            },
        ],
    }
    before = copy.deepcopy(database)

    with pytest.raises(
        tracking_database.TrackingDatabaseError,
        match="talks_schema_version_unsupported",
    ):
        tracking_database.migrate_tracking_database(database)

    assert database == before


def test_cli_rejects_future_talk_schema_without_rewriting(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = {"talks": [_talk(schema_version=99)]}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([_return()]))
    before = db.read_bytes()

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "talks_schema_version_unsupported" in result.stderr
    assert db.read_bytes() == before


def test_cli_assesses_future_root_before_old_talks_shape(
    persist_results,
    tmp_path,
):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = {
        "schema_version": 99,
        "future_inventory": {"records": "opaque"},
    }
    db.write_text(json.dumps(original))
    batch.write_text(json.dumps([_return()]))
    before = db.read_bytes()

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "tracking_database_schema_version_unsupported" in result.stderr
    assert "expected a JSON object with a `talks` array" not in result.stderr
    assert db.read_bytes() == before


def test_non_object_talk_member_is_actionable_and_does_not_rewrite(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk(), "not-a-talk"]}))
    batch.write_text(json.dumps([_return()]))
    before = db.read_bytes()

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "talks[1] must be a JSON object" in result.stderr
    assert "Traceback" not in result.stderr
    assert db.read_bytes() == before


def test_tracking_database_symlink_is_rejected_without_splitting_state(
        persist_results, tmp_path):
    target = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    target.write_text(json.dumps({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    before = target.read_bytes()
    link = tmp_path / "tracking-link.json"
    link.symlink_to(target.name)

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(link), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "symbolic link" in result.stderr
    assert link.is_symlink()
    assert target.read_bytes() == before


@pytest.mark.parametrize("block", ["structured_data", "verbatim_examples",
                                   "pattern_observations"])
def test_a_malformed_content_block_fails_loudly(persist_results, tmp_path, block):
    """A wrong-typed block used to be SKIPPED while the merge reported success.

    `structured_data` arriving as a list lost the entire analysis and still
    exited 0 — the silent-drop shape this script exists to eliminate.
    """
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret[block] = [{"slide_count": 60}]  # a list where an object is declared
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1, f"{block} was skipped silently: {result.stdout}"
    assert block in result.stderr


@pytest.mark.parametrize("bad", [
    ["narrative-arc", "bookends"],          # bare id strings
    "narrative-arc",                        # a string: len() counts characters
    {"pattern_id": "narrative-arc"},        # a single object, not an array
])
def test_malformed_detection_arrays_fail_loudly(persist_results, tmp_path, bad):
    """List-of-strings raised AttributeError mid-merge and killed the script
    before it printed its JSON; a plain string made `len()` count characters as
    detections and fed a silently wrong number into the score cross-check."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["patterns_detected"] = bad
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1, f"{bad!r} was accepted: {result.stdout}"
    assert "patterns_detected" in result.stderr
    assert "Traceback" not in result.stderr, "died instead of reporting"


def test_a_malformed_return_leaves_the_talk_untouched(persist_results, tmp_path):
    """Validation runs before any write, so a bad return is not half-merged."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["patterns_detected"] = ["narrative-arc"]
    original = _talk()
    db.write_text(_db_json({"talks": [original]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert json.loads(db.read_text())["talks"][0] == original


@pytest.mark.parametrize(
    ("field", "value", "message"),
    [
        ("rhetoric_notes", "", "must be a non-whitespace string"),
        ("areas_for_improvement", " \n\t", "must be a non-whitespace string"),
        ("adherence_assessment", " \n", "exact empty string sentinel"),
        ("transcript_source", None, "must be omitted when provenance is unknown"),
    ],
)
def test_v2_destructive_empty_scalars_leave_database_byte_stable(
        persist_results, tmp_path, field, value, message):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = {"talks": [_talk(
        rhetoric_notes="trusted prior analysis",
        areas_for_improvement="trusted prior improvement notes",
        transcript_source="manual",
    )]}
    ret = _return(**{field: value})
    db.write_text(_db_json(original))
    before = db.read_bytes()
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert message in result.stderr
    assert db.read_bytes() == before


def test_malformed_per_slide_visual_leaves_database_untouched(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["structured_data"].update({
        "slide_count": 1,
        # Synthetic legacy four-key shape: aliases cannot cross persistence.
        "per_slide_visual": [{
            "slide_number": 1,
            "content_type": "title",
            "background": "purple_halftone",
            "has_text_footer": True,
        }],
    })
    original = {"talks": [_talk()]}
    db.write_text(_db_json(original))
    before = db.read_bytes()
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "per_slide_visual[0] must contain exactly the canonical fields" in result.stderr
    assert db.read_bytes() == before


def test_effective_per_slide_dependencies_fail_before_database_write(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(structured_data={
        "slide_count": 1,
        "background_color_sequence": ["legacy"],
        "meme_count": 1,
    })
    ret = _return()
    ret["structured_data"].update({
        "slide_count": 1,
        "per_slide_visual": [_visual_row()],
    })
    db.write_text(_db_json({"talks": [talk]}))
    before = db.read_bytes()
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "effective merged analysis is invalid" in result.stderr
    assert "background_color_sequence" in result.stderr
    assert db.read_bytes() == before


def test_late_effective_failure_keeps_complete_database_byte_stable(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    first_talk = _talk(filename="first.md")
    second_talk = _talk(
        filename="second.md",
        structured_data={
            "slide_count": 1,
            "background_color_sequence": ["legacy"],
        },
    )
    first_return = _return(filename="first.md")
    second_return = _return(filename="second.md")
    second_return["structured_data"].update({
        "slide_count": 1,
        "per_slide_visual": [_visual_row()],
    })
    db.write_text(_db_json({"talks": [first_talk, second_talk]}))
    before = db.read_bytes()
    batch.write_text(json.dumps([first_return, second_return]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "second.md" in result.stderr
    assert db.read_bytes() == before


def test_malformed_stored_structured_container_keeps_database_byte_stable(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk(structured_data=["legacy"])]}))
    before = db.read_bytes()
    batch.write_text(json.dumps([_return()]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "stored structured_data" in result.stderr
    assert db.read_bytes() == before


def test_undeclared_legacy_verbatim_lane_keeps_database_byte_stable(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(verbatim_examples={
        "jokes": ["legacy line"],
        "legacy_lane": ["undeclared"],
    })
    db.write_text(_db_json({"talks": [talk]}))
    before = db.read_bytes()
    batch.write_text(json.dumps([_return(verbatim_examples={"jokes": []})]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "unknown verbatim_examples lanes" in result.stderr
    assert db.read_bytes() == before


def test_missing_image_source_distribution_basis_leaves_database_untouched(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["structured_data"]["image_source_distribution"] = {
        "speaker_created": 3,
        "unknown": 2,
    }
    original = {"talks": [_talk()]}
    db.write_text(_db_json(original))
    before = db.read_bytes()
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "image_source_distribution_basis is required" in result.stderr
    assert db.read_bytes() == before


def test_v2_image_source_basis_without_map_leaves_database_untouched(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["structured_data"]["image_source_distribution_basis"] = (
        "Unit: slide; unknown origins stay unknown.")
    db.write_text(_db_json({"talks": [_talk()]}))
    before = db.read_bytes()
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "to be supplied together" in result.stderr
    assert db.read_bytes() == before


def test_persistence_reports_that_schema_migration_was_not_run(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    merged, untouched = _talk(), _talk()
    untouched["filename"] = "other-talk.md"
    untouched["_queue_claim"]["batch_id"] = "unrelated-batch"
    db.write_text(_db_json({"talks": [merged, untouched]}))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    assert json.loads(result.stdout)["migrated_records"] == 0
    for talk in json.loads(db.read_text())["talks"]:
        assert talk["schema_version"] == persist_results.TALK_SCHEMA_VERSION


def test_score_object_without_a_score_key_fails_loudly(persist_results, tmp_path):
    """Present-but-incomplete is malformed, not absent — returning "no score"
    would drop the value exactly like the bare int used to."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    ret = _return()
    ret["pattern_observations"]["pattern_score"] = {
        "patterns_used": 2, "antipatterns_detected": 1}
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "pattern_score.score must be an integer" in result.stderr
    assert "pattern_score" not in json.loads(db.read_text())["talks"][0]


def test_clear_fields_removes_contaminated_values_and_promoted_scalars(
        persist_results, tmp_path):
    """A corrective reparse can explicitly remove claims disproved by artifacts."""
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(
        slide_count=999,
        slides_local_path="slides/stale.pdf",
        structured_data={"slide_count": 999, "stale_claim": "borrowed deck"},
        verbatim_examples={"jokes": ["quote from a sibling delivery"]},
    )
    ret = _return(clear_fields=[
        "slides_local_path",
        "structured_data.slide_count",
        "structured_data.stale_claim",
        "verbatim_examples.jokes",
    ])
    del ret["structured_data"]["slide_count"]
    ret["verbatim_examples"]["jokes"] = []
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert "slide_count" not in stored
    assert "slides_local_path" not in stored
    assert "slide_count" not in stored["structured_data"]
    assert "stale_claim" not in stored["structured_data"]
    # A v2 replacement follows the clear and is authoritative, even when empty.
    assert stored["verbatim_examples"]["jokes"] == []
    report = json.loads(result.stdout)
    assert "structured_data.slide_count" in report["talks"][0]["cleared"]
    assert "slide_count" in report["talks"][0]["cleared"]


def test_pattern_score_clear_removes_both_nested_and_promoted_value(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(
        pattern_score=99,
        pattern_observations={"pattern_score": 99, "pattern_ids": ["narrative-arc"]},
    )
    ret = _return(clear_fields=["pattern_observations.pattern_score"])
    # The valid replacement score is written after the clear; the important
    # invariant is that no old 99 survives on either surface.
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([ret]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert stored["pattern_score"] == 1
    assert stored["pattern_observations"]["pattern_score"] == 1


def test_historical_return_is_marked_unbaselineable_and_claim_is_closed(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(reprocess_reason="pattern_scoring_generation:legacy_generation")
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    report = json.loads(result.stdout)
    stored = json.loads(db.read_text())["talks"][0]
    assert stored["pattern_scoring_generation_status"] == \
        "legacy_unbaselineable"
    assert stored["pattern_scoring_generation_reasons"] == [
        "return_schema_precedes_source_locations"]
    assert "pattern_scoring_schema_version" not in stored
    assert "pattern_catalog_fingerprint" not in stored
    assert report["pattern_catalog_fingerprint"]
    assert stored["_queue_claim"]["state"] == "completed"
    assert stored["_queue_claim"]["schema_version"] == \
        persist_results.PREVIOUS_QUEUE_CLAIM_SCHEMA_VERSION
    assert stored["_queue_claim"]["result_status"] == "processed"
    assert stored["_queue_claim"]["release_reason"] == "return_persisted"
    assert stored["_queue_claim"]["result_payload_sha256"] == \
        persist_results.canonical_return_sha256(_return())
    assert "reprocess_reason" not in stored


@pytest.mark.parametrize("status", ["processed", "processed_partial"])
def test_analysis_terminal_outcome_clears_live_reprocess_reason(
        persist_results, status):
    talk = _talk(
        reprocess_reason="pattern_scoring_generation:missing_generation_status")
    ret = _return(status=status)

    persist_results.merge_talk(
        talk,
        ret,
        run_date="2026-07-31T18:30:00+00:00",
        enforce_queue_claim=True,
    )

    assert talk["status"] == status
    assert "reprocess_reason" not in talk
    assert talk["_queue_claim"]["state"] == "completed"


def test_v3_member_baseline_mismatch_rejects_whole_batch_without_write(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    filenames = ("a.md", "b.md")
    baseline = _adherence_baseline(
        persist_results, filenames=filenames)
    talks = []
    returns = []
    for filename in filenames:
        talk, ret = _v3_talk_and_return(
            persist_results, filename=filename)
        talk["_queue_claim"]["adherence_baseline"] = copy.deepcopy(baseline)
        talks.append(talk)
        returns.append(ret)
    talks[1]["_queue_claim"]["adherence_baseline"].update({
        "scored_talk_count": 1,
        "pattern_score_sum": 1,
        "average_pattern_score": 1.0,
    })
    original = {"talks": talks}
    db.write_text(_db_json(original))
    before = db.read_bytes()
    batch.write_text(json.dumps(returns))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "share one immutable adherence_baseline" in result.stderr
    assert db.read_bytes() == before


def test_post_batch_stdout_uses_complete_replacement_cohort_and_keeps_claim(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk, ret = _v3_talk_and_return(persist_results)
    claim_baseline = copy.deepcopy(talk["_queue_claim"]["adherence_baseline"])
    talk.update({
        "pattern_score": 99,
        "pattern_observations": {
            **talk["pattern_observations"],
            "pattern_score": 99,
        },
        "pattern_scoring_generation_status": "current",
        "pattern_scoring_generation_reasons": [],
        "pattern_scoring_schema_version": (
            persist_results.PATTERN_SCORING_SCHEMA_VERSION),
        "pattern_catalog_fingerprint": persist_results.load_catalog().fingerprint,
    })
    prior = _talk(filename="prior.md", status="processed")
    prior.pop("_queue_claim")
    prior.update({
        "pattern_score": 3,
        "pattern_observations": {
            **prior["pattern_observations"],
            "pattern_score": 3,
            "opportunity_coverage_identity": "0" * 64,
        },
        "pattern_scoring_generation_status": "current",
        "pattern_scoring_generation_reasons": [],
        "pattern_scoring_schema_version": (
            persist_results.PATTERN_SCORING_SCHEMA_VERSION),
        "pattern_catalog_fingerprint": persist_results.load_catalog().fingerprint,
    })
    db.write_text(_db_json({"talks": [talk, prior]}))
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [
            sys.executable,
            persist_results.__file__,
            str(db),
            str(batch),
            "--run-date",
            "2026-07-31T18:05:00+00:00",
        ],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    report = json.loads(result.stdout)
    cohort = report["current_adherence_baseline"]
    assert cohort["active_batch_excluded"] is False
    assert cohort["excluded_filenames"] == []
    assert cohort["as_of"] == "2026-07-31T18:05:00+00:00"
    assert cohort["scored_talk_count"] == 0
    assert cohort["pattern_score_sum"] == 0
    assert cohort["average_pattern_score"] is None
    stored = json.loads(db.read_text())["talks"][0]
    assert stored["pattern_score"] == 1
    assert stored["_queue_claim"]["adherence_baseline"] == claim_baseline


def test_missing_status_rejects_the_whole_batch_without_migrating_db(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = {"talks": [_talk()]}
    invalid = _return()
    del invalid["status"]
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([_return(filename="talk.md"), invalid]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "status is required" in result.stderr
    assert json.loads(db.read_text()) == original


def test_duplicate_db_filenames_are_rejected_before_write(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = {"talks": [_talk(), _talk()]}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([_return()]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "duplicate talk filename" in result.stderr
    assert json.loads(db.read_text()) == original


def test_partial_three_member_batch_is_rejected_before_any_db_write(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talks = [_talk(filename=f"{name}.md") for name in ("a", "b", "c")]
    returns = [_return(filename=f"{name}.md") for name in ("a", "b")]
    original = {"talks": talks}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps(returns))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )

    assert result.returncode == 1
    assert "missing ['c.md']" in result.stderr
    assert json.loads(db.read_text()) == original


def test_partially_closed_batch_cannot_be_finished_piecemeal(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talks = [_talk(filename=f"{name}.md") for name in ("a", "b", "c")]
    for talk in talks[:2]:
        talk["status"] = "processed"
        talk["_queue_claim"].update({
            "state": "completed",
            "released_at": "2026-07-31T18:05:00+00:00",
            "release_reason": "return_persisted",
            "result_status": "processed",
        })
    original = {"talks": talks}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([
        _return(filename=f"{name}.md") for name in ("a", "b", "c")]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )

    assert result.returncode == 1
    assert "closed or stranded member" in result.stderr
    assert json.loads(db.read_text()) == original


def test_stale_return_generation_cannot_roll_back_a_requeue(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(reprocess_generation=2)
    talk["_queue_claim"]["reprocess_generation"] = 2
    original = {"talks": [talk]}
    stale = _return()  # generation 1
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([stale]))
    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert result.returncode == 1
    assert "does not match active claim value 2" in result.stderr
    assert json.loads(db.read_text()) == original


def test_claim_generation_cannot_lag_top_level_generation(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(reprocess_generation=2)
    original = {"talks": [talk]}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([_return()]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )

    assert result.returncode == 1
    assert "current claim generation 1 disagrees with talk generation 2" in result.stderr
    assert json.loads(db.read_text()) == original


def test_active_claim_with_undeclared_fields_is_rejected_before_write(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk()
    talk["_queue_claim"]["released_at"] = None
    original = {"talks": [talk]}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([_return()]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )

    assert result.returncode == 1
    assert "must use exactly the schema fields" in result.stderr
    assert "released_at" in result.stderr
    assert json.loads(db.read_text()) == original


def test_skipped_return_preserves_prior_analysis_and_persists_terminal_metadata(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(
        reprocess_reason="source_identity_correction",
        processed_date="2026-06-18",
        transcript_source="youtube_auto",
        slide_source="pptx",
        rhetoric_notes="trusted prior analysis",
        structured_data={"slide_count": 42},
        verbatim_examples={"jokes": ["kept"]},
        pattern_scoring_generation_status="current",
        pattern_scoring_generation_reasons=[],
        pattern_scoring_schema_version=2,
        pattern_catalog_fingerprint="0" * 64,
        video_url=None,
        youtube_id=None,
        pptx_path=None,
        slides_url=None,
    )
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([_skipped_return()]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )

    assert result.returncode == 0, result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert stored["status"] == "skipped_no_sources"
    assert stored["processed_date"] == "2026-06-18"
    assert stored["transcript_source"] == "youtube_auto"
    assert stored["slide_source"] == "pptx"
    assert stored["rhetoric_notes"] == "trusted prior analysis"
    assert stored["structured_data"] == {"slide_count": 42}
    assert stored["verbatim_examples"] == {"jokes": ["kept"]}
    assert stored["pattern_scoring_generation_status"] == "current"
    assert stored["pattern_scoring_schema_version"] == 2
    assert stored["pattern_catalog_fingerprint"] == "0" * 64
    assert stored["_queue_claim"]["state"] == "completed"
    assert "reprocess_reason" not in stored
    report = json.loads(result.stdout)
    assert report["talks"][0]["pattern_scoring_generation_status"] == \
        "not_applicable"
    assert report["talks"][0]["pattern_scoring_generation_reasons"] == []


@pytest.mark.parametrize(
    "capability", ["remote_video", "transcript", "pptx", "pdf"])
def test_skipped_no_sources_rejects_each_live_capability(
        persist_results, tmp_path, capability):
    base = {
        "video_url": None,
        "youtube_id": None,
        "pptx_path": None,
        "slides_url": None,
        "google_drive_id": None,
        "transcript_path": None,
        "slides_local_path": None,
    }
    source_fields = {}
    if capability == "remote_video":
        source_fields = {"video_url": "https://youtu.be/AbCdEfGhI_1"}
    elif capability == "transcript":
        path = tmp_path / "transcripts" / "AbCdEfGhI_1.txt"
        path.parent.mkdir()
        path.write_text(" ".join(["substantive"] * 600), encoding="utf-8")
        source_fields = {
            "transcript_path": "transcripts/AbCdEfGhI_1.txt"}
    elif capability == "pptx":
        path = tmp_path / "decks" / "talk.pptx"
        path.parent.mkdir()
        deck = Presentation()
        deck.slides.add_slide(deck.slide_layouts[6])
        deck.save(str(path))
        source_fields = {"pptx_path": "decks/talk.pptx"}
    else:
        path = tmp_path / "slides" / "talk.pdf"
        path.parent.mkdir()
        writer = PdfWriter()
        writer.add_blank_page(width=640, height=480)
        with path.open("wb") as stream:
            writer.write(stream)
        source_fields = {"slides_local_path": "slides/talk.pdf"}
    talk = _talk(**{**base, **source_fields})
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([_skipped_return()]))
    before = db.read_bytes()

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "cannot finish skipped_no_sources" in result.stderr
    assert db.read_bytes() == before


@pytest.mark.parametrize("surviving_source", ["deck", "transcript"])
def test_invalid_source_cannot_hide_independent_terminal_capability(
        persist_results, tmp_path, surviving_source):
    fields = {
        "video_url": None,
        "youtube_id": None,
        "slides_url": None,
        "google_drive_id": None,
        "slides_local_path": None,
    }
    if surviving_source == "deck":
        deck = tmp_path / "decks" / "talk.pptx"
        deck.parent.mkdir()
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(deck)
        fields.update({
            "transcript_path": "../bad.txt",
            "pptx_path": "decks/talk.pptx",
        })
    else:
        transcript = tmp_path / "transcripts" / "manual-talk.txt"
        transcript.parent.mkdir()
        text = " ".join(["substantive transcript evidence"] * 200)
        transcript.write_text(text, encoding="utf-8")
        transcript_timing = importlib.import_module("transcript_timing")
        transcript_timing.write_quality_receipt(
            transcript,
            text,
            transcript_timing.build_quality_policy(400),
            {"kind": "fixed_default"},
        )
        fields.update({
            "transcript_path": "transcripts/manual-talk.txt",
            "transcript_source": "manual",
            "pptx_path": "../bad.pptx",
        })
    talk = _talk(**fields)
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([_skipped_return()]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "cannot finish skipped_no_sources" in result.stderr
    assert surviving_source in result.stderr or (
        surviving_source == "deck" and "slides" in result.stderr)


@pytest.mark.parametrize("artifact_kind", ["transcript", "pptx", "pdf"])
def test_download_failure_rejects_when_a_local_artifact_remains(
        persist_results, tmp_path, artifact_kind):
    local_source = {}
    if artifact_kind == "transcript":
        path = tmp_path / "transcripts" / "AbCdEfGhI_1.txt"
        path.parent.mkdir()
        path.write_text(" ".join(["substantive"] * 600), encoding="utf-8")
        local_source = {"transcript_path": "transcripts/AbCdEfGhI_1.txt"}
    elif artifact_kind == "pptx":
        path = tmp_path / "decks" / "talk.pptx"
        path.parent.mkdir()
        deck = Presentation()
        deck.slides.add_slide(deck.slide_layouts[6])
        deck.save(str(path))
        local_source = {"pptx_path": "decks/talk.pptx"}
    else:
        path = tmp_path / "slides" / "talk.pdf"
        path.parent.mkdir()
        writer = PdfWriter()
        writer.add_blank_page(width=640, height=480)
        with path.open("wb") as stream:
            writer.write(stream)
        local_source = {"slides_local_path": "slides/talk.pdf"}
    talk = _talk(**{
        "pptx_path": None,
        "slides_url": None,
        "transcript_path": None,
        "slides_local_path": None,
        **local_source,
    })
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([
        _skipped_return(status="skipped_download_failed")]))
    before = db.read_bytes()

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "source capabilities remain verified or repairable" in result.stderr
    assert db.read_bytes() == before


def test_broken_local_reference_allows_no_sources_terminal(
        persist_results, tmp_path):
    talk = _talk(
        video_url=None,
        youtube_id=None,
        slides_url=None,
        pptx_path="decks/missing.pptx",
    )
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([_skipped_return()]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    assert json.loads(db.read_text())["talks"][0]["status"] == \
        "skipped_no_sources"


def test_broken_local_reference_with_remote_allows_download_failed(
        persist_results, tmp_path):
    talk = _talk(
        slides_url=None,
        pptx_path="decks/missing.pptx",
    )
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([
        _skipped_return(status="skipped_download_failed")]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    assert json.loads(db.read_text())["talks"][0]["status"] == \
        "skipped_download_failed"


def test_download_failure_requires_a_remote_acquisition_path(
        persist_results, tmp_path):
    talk = _talk(
        video_url=None,
        youtube_id=None,
        pptx_path=None,
        slides_url=None,
    )
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([
        _skipped_return(status="skipped_download_failed")]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert (
        "without a usable remote video, transcript/YouTube identity, or slide "
        "acquisition path"
    ) in result.stderr


def test_malformed_remote_references_do_not_authorize_download_failed(
        persist_results, tmp_path):
    talk = _talk(
        video_url="https://youtu.be/too-short",
        youtube_id=None,
        pptx_path=None,
        slides_url="javascript:alert(1)",
        google_drive_id="../",
    )
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([
        _skipped_return(status="skipped_download_failed")]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "without a usable remote video" in result.stderr


def test_duplicate_skip_requires_a_bound_canonical_talk(
        persist_results, tmp_path):
    talk = _talk(video_url=None, pptx_path=None, slides_url=None)
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([_skipped_return(status="skipped_duplicate")]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "source_relation.type='duplicate'" in result.stderr


@pytest.mark.parametrize(
    ("status", "talk_overrides"),
    [
        (
            "skipped_download_failed",
            {"video_url": "https://youtu.be/AbCdEfGhI_1"},
        ),
        (
            "skipped_duplicate",
            {
                "source_relation": {
                    "type": "duplicate",
                    "target_filename": "canonical.md",
                },
            },
        ),
    ],
)
def test_mechanically_bound_skip_status_can_close_claim(
        persist_results, tmp_path, status, talk_overrides):
    talk = _talk(**{
        "reprocess_reason": "source_identity_correction",
        "video_url": None,
        "pptx_path": None,
        "slides_url": None,
        "transcript_path": None,
        "slides_local_path": None,
        **talk_overrides,
    })
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [talk]}))
    batch.write_text(json.dumps([_skipped_return(status=status)]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    stored = json.loads(db.read_text())["talks"][0]
    assert stored["status"] == status
    assert stored["_queue_claim"]["result_status"] == status
    assert "reprocess_reason" not in stored


def test_skipped_return_with_analysis_fields_aborts_without_write(
        persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    original = {"talks": [_talk(rhetoric_notes="trusted prior analysis")]}
    invalid = _skipped_return(rhetoric_notes="stale rejected analysis")
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([invalid]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )

    assert result.returncode == 1
    assert "cannot mutate or clear prior analysis fields" in result.stderr
    assert json.loads(db.read_text()) == original


def test_wrong_transcript_source_cannot_be_resurrected(
        persist_results, return_validation, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    talk = _talk(
        video_url=None,
        youtube_id=None,
        pptx_path=None,
        transcript_source="none",
        slide_source="pdf",
        google_drive_id="slides-id",
    )
    ret = _return(status="processed_partial", slide_source="pdf")
    ret["pattern_observations"]["evidence_sources"] = [
        "transcript", "static_slides"]
    _complete_unavailable_source_gates(return_validation, ret)
    original = {"talks": [talk]}
    db.write_text(_db_json(original))
    batch.write_text(json.dumps([ret]))

    result = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )

    assert result.returncode == 1
    assert "no transcript reference or active video source" in result.stderr
    assert json.loads(db.read_text()) == original


def test_completed_generation_cannot_be_replayed(persist_results, tmp_path):
    db = tmp_path / "tracking-database.json"
    batch = tmp_path / "batch-returns.json"
    db.write_text(_db_json({"talks": [_talk()]}))
    batch.write_text(json.dumps([_return()]))
    first = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert first.returncode == 0, first.stderr
    completed = db.read_bytes()
    replay = subprocess.run(
        [sys.executable, persist_results.__file__, str(db), str(batch)],
        capture_output=True, text=True,
    )
    assert replay.returncode == 1
    assert "closed or stranded member" in replay.stderr
    assert db.read_bytes() == completed


# --- #203: the CLI has a closed failure boundary that names commit position ---

@pytest.mark.parametrize("committed", [False, True])
def test_outer_boundary_reports_whether_the_commit_landed(
        persist_results, capsys, monkeypatch, committed):
    """A late failure must say whether the atomic write already happened.

    Without it an operator cannot tell a pre-commit abort from a post-commit
    reporting failure, and a blind retry could re-persist the batch.
    """
    def explode(*_args, **_kwargs):
        raise RuntimeError("injected failure at /private/vault/tracking.json")

    monkeypatch.setattr(persist_results, "main", explode)
    monkeypatch.setitem(persist_results._COMMIT_STATE, "database_written", committed)

    assert persist_results.run_cli() == 2

    captured = capsys.readouterr()
    assert captured.out == ""                     # stdout stays clean
    payload = json.loads(captured.err.splitlines()[0])
    assert payload["error"] == "persist_results_unexpected_failure"
    assert payload["error_type"] == "RuntimeError"
    assert payload["database_written"] is committed
    # Path-neutral: the exception text and its path never reach the output.
    assert "injected failure" not in captured.err
    assert "/private/vault/tracking.json" not in captured.err
    assert "Traceback" not in captured.err


def test_outer_boundary_lets_a_clean_run_report_success(
        persist_results, monkeypatch):
    """The boundary must not swallow or alter a normal run."""
    monkeypatch.setattr(persist_results, "main", lambda *a, **k: None)
    assert persist_results.run_cli() == 0


def test_outer_boundary_does_not_catch_sys_exit(persist_results, monkeypatch):
    """main()'s own sys.exit paths keep their exit codes."""
    def bail(*_args, **_kwargs):
        raise SystemExit(1)

    monkeypatch.setattr(persist_results, "main", bail)
    with pytest.raises(SystemExit) as excinfo:
        persist_results.run_cli()
    assert excinfo.value.code == 1
