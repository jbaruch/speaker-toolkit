"""Regression tests for the offline vault identity/source preflight."""

from copy import deepcopy
import hashlib
import importlib
import json
import os
from pathlib import Path
import pathlib
import shutil
import struct
import subprocess
import sys
from typing import Any
import zipfile

import pytest
from PIL import Image
from pypdf import PdfWriter
from pptx import Presentation
from pptx.util import Inches


VIDEO_ID = "AbCdEfGhI_1"
OTHER_VIDEO_ID = "ZyXwVuTsR_2"
DRIVE_ID = "drive-file-123"
DEFAULT_DIRECTORY_EXCLUSIONS = [
    ".venv",
    "venv",
    "node_modules",
    ".git",
    ".hg",
    ".svn",
    "__pycache__",
    ".pytest_cache",
    ".mypy_cache",
    ".ruff_cache",
    ".tox",
    ".nox",
    ".tessl",
]
QUEUE_STATE_SCRIPT = (
    Path(__file__).parents[1]
    / "skills"
    / "vault-ingress"
    / "scripts"
    / "queue-state.py"
)
_TINY_VIDEO_BYTES: bytes | None = None


def foreign_absolute_locator(name: str) -> str:
    if os.name == "nt":
        return f"/foreign/{name}"
    return rf"C:\foreign\{name}"


@pytest.fixture
def vault_fixture(tmp_path):
    vault = tmp_path / "vault"
    transcripts = vault / "transcripts"
    slides = vault / "slides"
    pptx_source = tmp_path / "presentations"
    transcripts.mkdir(parents=True)
    slides.mkdir()
    pptx_source.mkdir()
    return {
        "root": vault,
        "transcripts": transcripts,
        "slides": slides,
        "pptx_source": pptx_source,
        "database": vault / "tracking-database.json",
    }


def base_talk(**updates: Any) -> dict[str, Any]:
    talk = {
        "filename": "2026-07-30-perfect-ingress.md",
        "title": "Perfect Vault Ingress",
        "date": "2026-07-30",
        "video_url": f"https://www.youtube.com/watch?v={VIDEO_ID}",
        "youtube_id": VIDEO_ID,
        "transcript_source": "youtube_auto",
        "slide_source": "none",
        "status": "processed",
    }
    talk.update(updates)
    return talk


def current_v5_talk(preflight_vault, **updates: Any) -> dict[str, Any]:
    talk = base_talk()
    talk.update(
        {
            "schema_version": 5,
            "pattern_scoring_generation_status": "current",
            "pattern_scoring_generation_reasons": [],
            "pattern_scoring_schema_version": (
                preflight_vault.PATTERN_SCORING_SCHEMA_VERSION
            ),
            "pattern_observations": {
                "evidence_schema_version": (
                    preflight_vault.PATTERN_EVIDENCE_SCHEMA_VERSION
                ),
            },
        }
    )
    talk.update(updates)
    return talk


def source_identity(**updates):
    evidence = {
        "schema_version": 1,
        "provider": "youtube",
        "video_id": VIDEO_ID,
        "title": "Perfect Vault Ingress — Baruch Sadogursky",
        "uploader": "Conference Channel",
        "uploader_id": "@conference",
        "speakers": ["Baruch Sadogursky"],
        "recorded_date": "2026-07-30",
        "upload_date": "2026-07-31",
        "duration_seconds": 2700,
        "webpage_url": f"https://www.youtube.com/watch?v={VIDEO_ID}",
        "webpage_video_id": VIDEO_ID,
        "captured_at": "2026-07-31T12:00:00Z",
    }
    evidence.update(updates)
    return evidence


def write_database(fixture, talks, config=None, *, current=False):
    database = {
        "config": config
        or {
            "speaker_name": "Baruch Sadogursky",
            "pptx_source_dir": str(fixture["pptx_source"]),
        },
        "talks": talks,
    }
    if current:
        database.update(
            {
                "schema_version": 1,
                "pptx_catalog": [],
                "qr_codes": [],
                "resources": [],
                "thumbnails": [],
                "confirmed_intents": [],
                "improvement_goals": [],
            }
        )
        database["config"]["schema_version"] = 2
        database["config"].setdefault(
            "pptx_directory_exclusions",
            deepcopy(DEFAULT_DIRECTORY_EXCLUSIONS),
        )
        for talk in talks:
            if isinstance(talk, dict):
                talk["schema_version"] = 5
    fixture["database"].write_text(json.dumps(database, indent=2), encoding="utf-8")
    return fixture["database"]


def materialize_transcript(fixture, video_id=VIDEO_ID):
    path = fixture["transcripts"] / f"{video_id}.txt"
    # Long enough to remain plausible for the 45-minute provider-duration
    # fixtures. Tests targeting identity/provenance should not accidentally
    # exercise the shared partial-transcript guard.
    text = " ".join(["substantive transcript evidence"] * 600)
    path.write_text(text, encoding="utf-8")
    transcript_timing = importlib.import_module("transcript_timing")
    transcript_timing.write_quality_receipt(
        path,
        text,
        transcript_timing.build_quality_policy(400),
        {"kind": "fixed_default"},
    )
    return path


def write_pdf(path: Path, *, page_count: int = 1) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    writer = PdfWriter()
    for _index in range(page_count):
        writer.add_blank_page(width=640, height=480)
    with path.open("wb") as stream:
        writer.write(stream)
    return path


def write_tiny_video(path: Path) -> Path:
    """Materialize one valid MP4 while keeping ffmpeg local to video tests."""
    global _TINY_VIDEO_BYTES
    path.parent.mkdir(parents=True, exist_ok=True)
    if _TINY_VIDEO_BYTES is not None:
        path.write_bytes(_TINY_VIDEO_BYTES)
        return path

    ffmpeg = shutil.which("ffmpeg")
    assert ffmpeg is not None, "source-video manifest tests require ffmpeg"
    assert shutil.which("ffprobe") is not None, (
        "source-video manifest tests require ffprobe"
    )
    created = subprocess.run(
        [
            ffmpeg,
            "-nostdin",
            "-hide_banner",
            "-loglevel",
            "error",
            "-f",
            "lavfi",
            "-i",
            "color=c=black:s=160x90:r=1",
            "-t",
            "1",
            "-an",
            "-c:v",
            "mpeg4",
            "-pix_fmt",
            "yuv420p",
            "-movflags",
            "+faststart",
            "-y",
            str(path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    assert created.returncode == 0, created.stderr
    _TINY_VIDEO_BYTES = path.read_bytes()
    return path


def materialize_crc_damaged_pptx(fixture):
    """Create a deck with one CRC-damaged media member under the source root."""
    image_path = fixture["pptx_source"] / "asset.png"
    Image.new("RGB", (64, 64), "navy").save(image_path)
    path = fixture["pptx_source"] / "damaged.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    deck.save(str(path))
    with zipfile.ZipFile(path) as archive:
        member = next(
            item
            for item in archive.infolist()
            if item.filename.startswith("ppt/media/") and item.file_size
        )
    package = bytearray(path.read_bytes())
    name_size, extra_size = struct.unpack_from(
        "<HH", package, member.header_offset + 26
    )
    payload_offset = member.header_offset + 30 + name_size + extra_size
    package[payload_offset + (member.compress_size // 2)] ^= 0xFF
    path.write_bytes(package)
    return path, member.filename


def materialize_shared_crc_damaged_tiff_pptx(fixture, *, slide_count=73):
    """Model the stable multi-owner vault deck without retaining private bytes."""
    image_path = fixture["pptx_source"] / "shared-asset.tiff"
    Image.new("RGB", (64, 64), "navy").save(image_path, format="TIFF")
    path = fixture["pptx_source"] / "shared-damaged-73-slides.pptx"
    deck = Presentation()
    for index in range(slide_count):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        if index == 0:
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    deck.save(str(path))
    with zipfile.ZipFile(path) as archive:
        member = next(
            item
            for item in archive.infolist()
            if item.filename.startswith("ppt/media/") and item.file_size
        )
    assert member.filename.endswith(".tiff")
    package = bytearray(path.read_bytes())
    name_size, extra_size = struct.unpack_from(
        "<HH", package, member.header_offset + 26
    )
    payload_offset = member.header_offset + 30 + name_size + extra_size
    package[payload_offset + (member.compress_size // 2)] ^= 0xFF
    path.write_bytes(package)
    return path, member.filename


def trusted_video_manifest(fixture, page_count=1):
    rebuild = fixture["root"] / "slides-rebuild" / VIDEO_ID
    rebuild.mkdir(parents=True, exist_ok=True)
    source_video = rebuild / f"{VIDEO_ID}.mp4"
    write_tiny_video(source_video)
    candidate = rebuild / f"{VIDEO_ID}.slide-region.pdf"
    write_pdf(candidate, page_count=page_count)
    retained = [
        {
            "page_number": page,
            "frame_index": page - 1,
            "timestamp_seconds": float((page - 1) * 2),
        }
        for page in range(1, page_count + 1)
    ]
    return {
        "slide_source": "video_extracted",
        "schema_version": 3,
        "pipeline_version": "0.10.0",
        "source_video_id": VIDEO_ID,
        "source_video_path": str(source_video),
        "total_frames_extracted": page_count,
        "unique_frame_count": page_count,
        "authored_slide_count": None,
        "hash_threshold_used": 8,
        "slide_region_detected": False,
        "slide_region_applied": True,
        "slide_region_method": "manual",
        "slide_region_verified": True,
        "slide_region": [0.05, 0.02, 0.78, 0.98],
        "fps_used": 0.5,
        "retained_frames": retained,
        "review_required": False,
        "review_reason": None,
        "artifacts": [
            {
                "path": str(candidate),
                "artifact_scope": "slide_region",
                "page_count": page_count,
                "source_video_id": VIDEO_ID,
                "source_video_path": str(source_video),
                "crop_method": "manual",
                "crop_verified": True,
                "trusted_for_authored_slide_analysis": True,
            }
        ],
    }


def context_video_manifest(fixture, page_count=1):
    manifest = trusted_video_manifest(fixture, page_count)
    rebuild = fixture["root"] / "slides-rebuild" / VIDEO_ID
    context = rebuild / f"{VIDEO_ID}.context.pdf"
    write_pdf(context, page_count=page_count)
    source_video = manifest["source_video_path"]
    manifest.update(
        {
            "slide_region_detected": False,
            "slide_region_applied": False,
            "slide_region_method": "none",
            "slide_region_verified": False,
            "slide_region": None,
            "review_required": True,
            "review_reason": "No verified slide region is available.",
            "artifacts": [
                {
                    "path": str(context),
                    "artifact_scope": "full_frame_context",
                    "page_count": page_count,
                    "source_video_id": VIDEO_ID,
                    "source_video_path": source_video,
                    "crop_method": "none",
                    "crop_verified": False,
                    "trusted_for_authored_slide_analysis": False,
                }
            ],
        }
    )
    return manifest


def finding_codes(report, severity=None):
    return {
        finding["code"]
        for finding in report["findings"]
        if severity is None or finding["severity"] == severity
    }


@pytest.mark.parametrize(
    ("url", "expected"),
    [
        (f"https://www.youtube.com/watch?v={VIDEO_ID}&t=42", VIDEO_ID),
        (f"https://youtu.be/{VIDEO_ID}?feature=shared", VIDEO_ID),
        (f"https://www.youtube.com/shorts/{VIDEO_ID}?si=x", VIDEO_ID),
        (f"https://www.youtube.com/embed/{VIDEO_ID}", VIDEO_ID),
        (f"https://www.youtube-nocookie.com/embed/{VIDEO_ID}", VIDEO_ID),
    ],
)
def test_youtube_id_parser_covers_supported_source_forms(
    preflight_vault, url, expected
):
    assert preflight_vault.parse_youtube_id(url) == expected


@pytest.mark.parametrize(
    "url",
    [
        "https://www.youtube.com/watch?v=too-short",
        "https://www.youtube.com/live/AbCdEfGhI_1",
        "https://videos.example.com/watch?v=AbCdEfGhI_1",
    ],
)
def test_youtube_id_parser_rejects_unsupported_or_invalid_urls(preflight_vault, url):
    assert preflight_vault.parse_youtube_id(url) is None


def test_preflight_input_and_finding_paths_are_lexical(
    preflight_vault,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    database = vault / "tracking-database.json"
    monkeypatch.setattr(
        Path,
        "is_dir",
        lambda *_args, **_kwargs: pytest.fail("preflight probed input directory"),
    )
    monkeypatch.setattr(
        Path,
        "exists",
        lambda *_args, **_kwargs: pytest.fail("preflight probed input existence"),
    )
    monkeypatch.setattr(
        Path,
        "resolve",
        lambda *_args, **_kwargs: pytest.fail("preflight resolved a finding path"),
    )

    assert preflight_vault.resolve_input(vault) == (vault, database)
    assert preflight_vault.resolve_input(database) == (vault, database)
    json_named_vault = vault / "conference-catalog.json"
    assert preflight_vault.resolve_input(json_named_vault) == (
        json_named_vault,
        json_named_vault / "tracking-database.json",
    )
    case_variant_database = vault / "TRACKING-DATABASE.JSON"
    assert preflight_vault.resolve_input(case_variant_database) == (
        vault,
        case_variant_database,
    )
    validator = preflight_vault.VaultPreflight({}, vault, database)
    validator.add(
        "warning",
        "synthetic",
        "synthetic finding",
        artifact_path=vault / "slides" / "talk.pptx",
    )
    assert validator.findings[0]["artifact_path"] == str(vault / "slides" / "talk.pptx")


@pytest.mark.parametrize(
    ("raw_root", "locator_reason_code"),
    [
        ("relative-vault", "artifact_root_not_native_absolute"),
        ("~/vault", "artifact_locator_home_expansion_unsupported"),
        (
            foreign_absolute_locator("vault"),
            "artifact_locator_foreign_absolute",
        ),
        (r"\\?\C:\vault", "artifact_locator_windows_device_namespace"),
    ],
)
def test_preflight_rejects_direct_root_before_database_snapshot(
    preflight_vault,
    monkeypatch: pytest.MonkeyPatch,
    raw_root: str,
    locator_reason_code: str,
) -> None:
    monkeypatch.setattr(
        preflight_vault,
        "snapshot_tracking_database",
        lambda *_args, **_kwargs: pytest.fail(
            "invalid direct root reached the database snapshot boundary"
        ),
    )

    report = preflight_vault.run_preflight(raw_root)

    assert report["database"] is None
    assert report["vault_root"] is None
    assert report["blocking_count"] == 1
    finding = report["findings"][0]
    assert finding["code"] == "vault_root_cli_invalid"
    assert finding["field"] == "cli.vault_root"
    assert finding["actual"] == {
        "reason_code": "vault_root_cli_invalid",
        "locator_reason_code": locator_reason_code,
    }
    assert raw_root not in json.dumps(report, sort_keys=True)


def test_preflight_labels_relative_direct_database_before_snapshot(
    preflight_vault,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    raw_database = "private-vault/tracking-database.json"
    monkeypatch.setattr(
        preflight_vault,
        "snapshot_tracking_database",
        lambda *_args, **_kwargs: pytest.fail(
            "invalid direct database reached the snapshot boundary"
        ),
    )

    report = preflight_vault.run_preflight(raw_database)

    assert report["database"] is None
    assert report["vault_root"] is None
    assert report["blocking_count"] == 1
    finding = report["findings"][0]
    assert finding["code"] == "vault_root_database_path_invalid"
    assert finding["field"] == "database.path"
    assert finding["actual"] == {
        "reason_code": "vault_root_database_path_invalid",
        "locator_reason_code": "artifact_root_not_native_absolute",
    }
    assert raw_database not in json.dumps(report, sort_keys=True)


@pytest.mark.parametrize(
    ("configured_root", "expected_code", "expected_locator_reason"),
    [
        (
            "",
            "vault_root_config_invalid",
            "artifact_locator_empty_or_whitespace",
        ),
        (
            " ",
            "vault_root_config_invalid",
            "artifact_locator_empty_or_whitespace",
        ),
        (
            "relative-vault",
            "vault_root_config_invalid",
            "artifact_root_not_native_absolute",
        ),
        (
            "C:vault",
            "vault_root_config_invalid",
            "artifact_locator_windows_drive_relative",
        ),
        (
            r"\vault",
            "vault_root_config_invalid",
            "artifact_locator_windows_current_drive_rooted",
        ),
        (
            r"C:\trusted\other\..\vault"
            if os.name == "nt"
            else "/trusted/other/../vault",
            "vault_root_config_invalid",
            "artifact_locator_dot_segment",
        ),
        (
            "~/vault",
            "vault_root_config_invalid",
            "artifact_locator_home_expansion_unsupported",
        ),
        (
            foreign_absolute_locator("vault"),
            "vault_root_config_invalid",
            "artifact_locator_foreign_absolute",
        ),
        (
            r"\\?\C:\vault",
            "vault_root_config_invalid",
            "artifact_locator_windows_device_namespace",
        ),
    ],
)
def test_preflight_rejects_invalid_configured_vault_root_before_catalog_checks(
    preflight_vault,
    vault_fixture,
    monkeypatch: pytest.MonkeyPatch,
    configured_root: str,
    expected_code: str,
    expected_locator_reason: str,
) -> None:
    write_database(
        vault_fixture,
        [],
        config={
            "speaker_name": "Baruch Sadogursky",
            "vault_storage_path": configured_root,
        },
    )
    monkeypatch.setattr(
        preflight_vault.VaultPreflight,
        "_validate_filenames",
        lambda *_args, **_kwargs: pytest.fail(
            "invalid configured vault root reached catalog checks"
        ),
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 1
    finding = report["findings"][0]
    assert finding["code"] == expected_code
    assert finding["field"] == "config.vault_storage_path"
    assert finding["actual"] == {
        "reason_code": expected_code,
        "locator_reason_code": expected_locator_reason,
    }
    if configured_root.strip():
        assert configured_root not in json.dumps(finding, sort_keys=True)


def test_preflight_rejects_configured_vault_authority_mismatch_before_catalog_checks(
    preflight_vault,
    vault_fixture,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    configured_root = vault_fixture["root"].parent / "different-vault"
    write_database(
        vault_fixture,
        [],
        config={
            "speaker_name": "Baruch Sadogursky",
            "vault_storage_path": str(configured_root),
        },
    )
    monkeypatch.setattr(
        preflight_vault.VaultPreflight,
        "_validate_filenames",
        lambda *_args, **_kwargs: pytest.fail(
            "mismatched vault root reached catalog checks"
        ),
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 1
    finding = report["findings"][0]
    assert finding["code"] == "vault_root_authority_mismatch"
    assert finding["field"] == "config.vault_storage_path"
    assert finding["actual"] == {
        "reason_code": "vault_root_authority_mismatch",
        "authorities": ["database_path", "config_root"],
    }
    assert str(configured_root) not in json.dumps(finding, sort_keys=True)


def test_preflight_names_cli_authority_when_database_parent_disagrees(
    preflight_vault,
    tmp_path: Path,
) -> None:
    cli_root = tmp_path / "cli-vault"
    database_path = tmp_path / "database-vault" / "tracking-database.json"
    validator = preflight_vault.VaultPreflight(
        {"config": {}, "talks": []},
        cli_root,
        database_path,
    )

    report = validator.run()

    assert report["blocking_count"] == 1
    finding = report["findings"][0]
    assert finding["code"] == "vault_root_authority_mismatch"
    assert finding["field"] == "cli.vault_root"
    assert finding["actual"] == {
        "reason_code": "vault_root_authority_mismatch",
        "authorities": ["database_path", "cli_root"],
    }
    diagnostic = json.dumps(finding, sort_keys=True)
    assert str(cli_root) not in diagnostic
    assert str(database_path) not in diagnostic


@pytest.mark.parametrize("configured", [None, "same"])
def test_preflight_accepts_null_or_matching_configured_vault_authority(
    preflight_vault,
    vault_fixture,
    configured: str | None,
) -> None:
    configured_value = str(vault_fixture["root"]) if configured == "same" else None
    write_database(
        vault_fixture,
        [],
        config={
            "speaker_name": "Baruch Sadogursky",
            "vault_storage_path": configured_value,
        },
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert not any(
        finding["code"].startswith("vault_root_") for finding in report["findings"]
    )


@pytest.mark.skipif(
    sys.platform == "win32",
    reason="directory symlink setup is privileged",
)
@pytest.mark.parametrize("configured_identity", ["alias", "physical"])
def test_preflight_compares_symlinked_vault_roots_by_lexical_identity(
    preflight_vault,
    vault_fixture,
    tmp_path: Path,
    configured_identity: str,
) -> None:
    physical_root = vault_fixture["root"]
    alias_root = tmp_path / "alias-vault"
    alias_root.symlink_to(physical_root, target_is_directory=True)
    configured_root = alias_root if configured_identity == "alias" else physical_root
    write_database(
        vault_fixture,
        [],
        config={
            "speaker_name": "Baruch Sadogursky",
            "vault_storage_path": str(configured_root),
        },
    )

    report = preflight_vault.run_preflight(alias_root)

    if configured_identity == "alias":
        assert report["blocking_count"] == 0
        assert report["vault_root"] == str(alias_root)
        assert report["database"] == str(alias_root / "tracking-database.json")
    else:
        assert report["blocking_count"] == 1
        finding = report["findings"][0]
        assert finding["code"] == "vault_root_authority_mismatch"
        assert finding["actual"] == {
            "reason_code": "vault_root_authority_mismatch",
            "authorities": ["database_path", "config_root"],
        }
        diagnostic = json.dumps(finding, sort_keys=True)
        assert str(alias_root) not in diagnostic
        assert str(physical_root) not in diagnostic


def test_pptx_artifact_validation_never_calls_parent_is_file(
    preflight_vault,
    vault_fixture,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    talk = base_talk(
        transcript_source="none",
        slide_source="pptx",
        pptx_path="Conference/Talk.pptx",
    )
    write_database(vault_fixture, [talk])
    capabilities = {
        "verified_capabilities": ("slides",),
        "verified_evidence_sources": ("native_deck",),
        "acquisition_capabilities": (),
        "repair_capabilities": (),
        "source_reasons": {},
        "degraded_evidence_sources": {},
        "unavailable_evidence_sources": {},
    }
    monkeypatch.setattr(
        preflight_vault,
        "assess_talk_artifact_capabilities",
        lambda *_args, **_kwargs: capabilities,
    )
    original_is_file = Path.is_file

    def guarded_is_file(path: Path) -> bool:
        if path.suffix.casefold() == ".pptx":
            pytest.fail("preflight called is_file before the bounded PPTX probe")
        return original_is_file(path)

    monkeypatch.setattr(Path, "is_file", guarded_is_file)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert not {
        "slide_pptx_artifact_missing",
        "slide_pptx_artifact_unreadable",
    }.intersection(finding_codes(report))


@pytest.mark.parametrize(
    "manifest_factory",
    [trusted_video_manifest, context_video_manifest],
    ids=["slide-region", "full-frame-context"],
)
def test_video_pdf_validation_never_touches_pdf_leaf_in_parent(
    preflight_vault,
    vault_fixture,
    monkeypatch: pytest.MonkeyPatch,
    manifest_factory,
) -> None:
    manifest = manifest_factory(vault_fixture)
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])
    pattern_evidence = importlib.import_module("pattern_evidence")

    for method_name in (
        "is_file",
        "resolve",
        "stat",
        "lstat",
        "open",
        "read_bytes",
    ):
        original = getattr(Path, method_name)

        def guarded(
            path: Path,
            *args,
            _original=original,
            _method_name=method_name,
            **kwargs,
        ):
            if path.suffix.casefold() == ".pdf":
                pytest.fail(f"preflight called {_method_name} for a PDF artifact")
            return _original(path, *args, **kwargs)

        monkeypatch.setattr(Path, method_name, guarded)

    original_hash = pattern_evidence._sha256_file

    def guarded_hash(path: Path) -> str:
        if path.suffix.casefold() == ".pdf":
            pytest.fail("preflight hashed a PDF artifact in the owner")
        return original_hash(path)

    monkeypatch.setattr(pattern_evidence, "_sha256_file", guarded_hash)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "video_extraction_provenance_invalid" not in finding_codes(report)


def test_processed_video_extraction_requires_promoted_canonical_pdf(
    preflight_vault,
    vault_fixture,
) -> None:
    manifest = trusted_video_manifest(vault_fixture)
    talk = base_talk(
        status="processed",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "slide_video_artifact_missing" in finding_codes(report)


def test_promoted_video_pdf_must_match_trusted_slide_region_digest(
    preflight_vault,
    vault_fixture,
) -> None:
    manifest = trusted_video_manifest(vault_fixture)
    promoted = vault_fixture["slides"] / f"{VIDEO_ID}.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=640, height=480)
    writer.add_metadata({"/Title": "Different valid one-page PDF"})
    with promoted.open("wb") as stream:
        writer.write(stream)
    talk = base_talk(
        status="processed",
        transcript_source="none",
        slide_source="video_extracted",
        slides_local_path=promoted.relative_to(vault_fixture["root"]).as_posix(),
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "video_extraction_provenance_invalid"
    )
    assert finding["severity"] == "blocking"
    assert any("digest must match" in error for error in finding["actual"])


def test_video_extraction_manifest_is_rejected_outside_video_slide_lane(
    preflight_vault,
    vault_fixture,
    monkeypatch,
) -> None:
    def reject_probe(*_args, **_kwargs):
        raise AssertionError("wrong-lane video artifacts must not be probed")

    monkeypatch.setattr(preflight_vault, "probe_pdf_artifact", reject_probe)
    talk = base_talk(
        status="pending",
        video_url=None,
        youtube_id=None,
        transcript_source="none",
        slide_source="none",
        structured_data={
            "video_extraction": {
                "artifacts": [{"path": "/outside/untrusted.pdf"}],
            }
        },
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "video_extraction_slide_source_mismatch"
    )
    assert finding["severity"] == "blocking"
    assert finding["field"] == "structured_data.video_extraction"
    assert finding["actual"] == "none"


@pytest.mark.parametrize(
    "manifest_factory",
    [trusted_video_manifest, context_video_manifest],
    ids=["slide-region", "full-frame-context"],
)
def test_video_manifest_pdf_page_count_is_bounded_and_verified(
    preflight_vault,
    vault_fixture,
    manifest_factory,
) -> None:
    manifest = manifest_factory(vault_fixture, page_count=2)
    artifact = Path(manifest["artifacts"][0]["path"])
    write_pdf(artifact, page_count=1)
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "video_extraction_provenance_invalid"
    )
    assert any(
        "page_count must match the bounded PDF probe" in error
        for error in finding["actual"]
    )


def test_video_manifest_rejects_pdf_spoofed_as_source_video_before_is_file(
    preflight_vault,
    vault_fixture,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    manifest = trusted_video_manifest(vault_fixture)
    spoofed = manifest["artifacts"][0]["path"]
    manifest["source_video_path"] = spoofed
    manifest["artifacts"][0]["source_video_path"] = spoofed
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])
    original_is_file = Path.is_file

    def guarded_is_file(path: Path) -> bool:
        if path.suffix.casefold() == ".pdf":
            pytest.fail("preflight statted a PDF spoofed as source video")
        return original_is_file(path)

    monkeypatch.setattr(Path, "is_file", guarded_is_file)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "video_extraction_provenance_invalid" in finding_codes(report)


def test_video_manifest_rejects_source_video_outside_vault(
    preflight_vault,
    vault_fixture,
) -> None:
    manifest = trusted_video_manifest(vault_fixture)
    outside = vault_fixture["root"].parent / "outside" / f"{VIDEO_ID}.mp4"
    outside.parent.mkdir()
    outside.write_bytes(b"outside video")
    manifest["source_video_path"] = str(outside)
    for artifact in manifest["artifacts"]:
        artifact["source_video_path"] = str(outside)
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "video_extraction_provenance_invalid" in finding_codes(report)


def test_missing_source_video_reports_closed_path_neutral_error(
    preflight_vault,
    vault_fixture,
) -> None:
    manifest = trusted_video_manifest(vault_fixture)
    source = Path(manifest["source_video_path"])
    source.unlink()
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "source_video_artifact_missing"
    )
    diagnostic = json.dumps(
        {"message": finding["message"], "actual": finding["actual"]},
        sort_keys=True,
    )
    assert str(source) not in diagnostic
    assert finding["severity"] == "warning"
    assert finding["field"] == ("structured_data.video_extraction.source_video_path")
    assert finding["actual"]["reason_code"] == "video_artifact_unavailable"
    assert finding["actual"]["details"]["failure_kind"] == "missing"
    assert "video_extraction_provenance_invalid" not in finding_codes(report)


def test_video_manifest_rejects_symlinked_source_video(
    preflight_vault,
    vault_fixture,
) -> None:
    manifest = trusted_video_manifest(vault_fixture)
    source = Path(manifest["source_video_path"])
    outside = vault_fixture["root"].parent / "outside" / source.name
    outside.parent.mkdir()
    outside.write_bytes(b"outside video")
    source.unlink()
    try:
        source.symlink_to(outside)
    except OSError:
        pytest.skip("directory policy does not permit symlink creation")
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "video_extraction_provenance_invalid" in finding_codes(report)


def test_video_manifest_pdf_probes_accept_documented_symlink_vault_root(
    preflight_vault,
    vault_fixture,
    tmp_path: Path,
) -> None:
    manifest = context_video_manifest(vault_fixture)
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])
    locator = tmp_path / "canonical-vault"
    try:
        locator.symlink_to(vault_fixture["root"], target_is_directory=True)
    except OSError:
        pytest.skip("directory policy does not permit symlink creation")

    report = preflight_vault.run_preflight(locator)

    assert "video_extraction_provenance_invalid" not in finding_codes(report)


def test_slide_contract_taxonomy_includes_bounded_pdf_failures(
    preflight_vault,
) -> None:
    assert {
        "slide_pdf_artifact_unavailable",
        "slide_pdf_artifact_unreadable",
        "slide_video_artifact_unavailable",
        "slide_video_artifact_unreadable",
    } <= preflight_vault.SLIDE_CONTRACT_CODES


def test_source_video_contract_taxonomy_is_separate_and_closed(
    preflight_vault,
) -> None:
    assert preflight_vault.SOURCE_VIDEO_CONTRACT_CODES == {
        "source_video_artifact_missing",
        "source_video_artifact_unavailable",
        "source_video_artifact_unreadable",
    }
    assert not (
        preflight_vault.SOURCE_VIDEO_CONTRACT_CODES
        & preflight_vault.SLIDE_CONTRACT_CODES
    )


@pytest.mark.parametrize(
    ("reason_code", "details", "expected_code"),
    [
        (
            "video_artifact_unavailable",
            {"failure_kind": "missing"},
            "source_video_artifact_missing",
        ),
        (
            "video_cloud_placeholder_unavailable",
            {"availability": {"state": "offline"}},
            "source_video_artifact_unavailable",
        ),
        *[
            (reason, {}, "source_video_artifact_unreadable")
            for reason in (
                "video_artifact_changed",
                "video_artifact_too_large",
                "video_batch_wall_limit",
                "video_dependency_unavailable",
                "video_duration_unavailable",
                "video_evidence_invalid",
                "video_invalid_container",
                "video_no_video_stream",
                "video_parser_rejected",
                "video_parser_repair_required",
                "video_probe_containment_unavailable",
                "video_probe_crash",
                "video_probe_malformed_result",
                "video_probe_monitor_identity_changed",
                "video_probe_monitor_unavailable",
                "video_probe_request_oversized",
                "video_probe_resource_unavailable",
                "video_probe_result_oversized",
                "video_probe_start_failure",
                "video_probe_timeout",
                "video_stream_limit",
            )
        ],
        (
            "video_artifact_unavailable",
            {"failure_kind": "io"},
            "source_video_artifact_unreadable",
        ),
        (
            "video_artifact_unavailable",
            {},
            "source_video_artifact_unreadable",
        ),
        *[
            (
                "video_artifact_unavailable",
                {"failure_kind": failure_kind},
                "video_extraction_provenance_invalid",
            )
            for failure_kind in (
                "not_regular",
                "root_escape",
                "symlink_or_reparse",
            )
        ],
        (
            "video_evidence_invalid",
            {"locator_failure": "artifact_locator_dot_segment"},
            "video_extraction_provenance_invalid",
        ),
    ],
)
def test_source_video_failure_mapping_is_closed(
    preflight_vault,
    reason_code: str,
    details: dict[str, object],
    expected_code: str,
) -> None:
    video_evidence = importlib.import_module("video_evidence")
    error = video_evidence.VideoEvidenceError(
        "synthetic bounded video failure",
        reason_code=reason_code,
        details=details,
    )

    assert preflight_vault._source_video_failure_code(error) == expected_code


@pytest.mark.parametrize(
    ("reason_code", "details", "expected_reason"),
    [
        (
            "video_evidence_invalid",
            {"locator_failure": "artifact_locator_dot_segment"},
            "source_video_path: artifact_locator_dot_segment",
        ),
        (
            "video_artifact_unavailable",
            {"failure_kind": "root_escape"},
            "source_video_path: root_escape",
        ),
    ],
)
def test_source_video_provenance_failures_keep_string_reason_list(
    preflight_vault,
    vault_fixture,
    monkeypatch: pytest.MonkeyPatch,
    reason_code: str,
    details: dict[str, object],
    expected_reason: str,
) -> None:
    video_evidence = importlib.import_module("video_evidence")
    manifest = trusted_video_manifest(vault_fixture)

    def fail_probe(*_args, **_kwargs):
        raise video_evidence.VideoEvidenceError(
            "synthetic bounded video provenance failure",
            reason_code=reason_code,
            details=details,
        )

    monkeypatch.setattr(video_evidence.VideoEvidenceAssessment, "probe", fail_probe)
    write_database(
        vault_fixture,
        [
            base_talk(
                status="processed_partial",
                transcript_source="none",
                slide_source="video_extracted",
                structured_data={"video_extraction": manifest},
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "video_extraction_provenance_invalid"
    )
    assert finding["actual"] == [expected_reason]
    assert all(isinstance(reason, str) for reason in finding["actual"])


@pytest.mark.parametrize(
    ("reason_code", "details", "expected_code", "current", "severity"),
    [
        (
            "video_artifact_unavailable",
            {"failure_kind": "missing", "exception_type": "FileNotFoundError"},
            "source_video_artifact_missing",
            False,
            "warning",
        ),
        (
            "video_cloud_placeholder_unavailable",
            {"availability": {"state": "offline"}, "reparse_tag": None},
            "source_video_artifact_unavailable",
            False,
            "warning",
        ),
        (
            "video_invalid_container",
            {},
            "source_video_artifact_unreadable",
            False,
            "warning",
        ),
        (
            "video_invalid_container",
            {},
            "source_video_artifact_unreadable",
            True,
            "blocking",
        ),
    ],
)
def test_source_video_probe_failure_is_structured_current_aware_and_not_duplicated(
    preflight_vault,
    vault_fixture,
    monkeypatch: pytest.MonkeyPatch,
    reason_code: str,
    details: dict[str, object],
    expected_code: str,
    current: bool,
    severity: str,
) -> None:
    video_evidence = importlib.import_module("video_evidence")
    manifest = trusted_video_manifest(vault_fixture)

    def fail_probe(*_args, **_kwargs):
        raise video_evidence.VideoEvidenceError(
            "synthetic bounded video failure",
            reason_code=reason_code,
            details=details,
        )

    monkeypatch.setattr(video_evidence.VideoEvidenceAssessment, "probe", fail_probe)
    talk_updates = {
        "status": "processed_partial",
        "transcript_source": "none",
        "slide_source": "video_extracted",
        "structured_data": {"video_extraction": manifest},
    }
    talk = (
        current_v5_talk(preflight_vault, **talk_updates)
        if current
        else base_talk(**talk_updates)
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = [item for item in report["findings"] if item["code"] == expected_code]
    assert len(findings) == 1
    finding = findings[0]
    assert finding["severity"] == severity
    assert finding["field"] == ("structured_data.video_extraction.source_video_path")
    assert finding["actual"] == {
        "reason_code": reason_code,
        "details": details,
    }
    assert "video_extraction_provenance_invalid" not in finding_codes(report)


def test_clean_record_has_no_findings(preflight_vault, vault_fixture):
    materialize_transcript(vault_fixture)
    write_pdf(vault_fixture["slides"] / f"{DRIVE_ID}.pdf")
    pptx = vault_fixture["pptx_source"] / "Conf" / "Perfect.pptx"
    pptx.parent.mkdir()
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(str(pptx))
    talk = base_talk(
        slide_source="both",
        google_drive_id=DRIVE_ID,
        pptx_path="Conf/Perfect.pptx",
        duration_seconds=2700,
        source_identity=source_identity(),
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["ok"] is True
    assert report["blocking_count"] == 0
    assert report["warning_count"] == 0
    assert report["findings"] == []


def test_missing_quality_receipt_is_actionable_for_completed_evidence(
    preflight_vault,
    vault_fixture,
):
    transcript = materialize_transcript(vault_fixture)
    transcript.with_suffix(".quality.json").unlink()
    talk = base_talk(source_identity=source_identity())
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "transcript_quality_receipt_unverified" in finding_codes(report, "warning")
    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "transcript_quality_receipt_unverified"
    )
    assert finding["capability_fact"]["repair_capabilities"] == ["transcript"]


def test_current_v5_missing_quality_receipt_remains_blocking(
    preflight_vault,
    vault_fixture,
):
    transcript = materialize_transcript(vault_fixture)
    transcript.with_suffix(".quality.json").unlink()
    talk = current_v5_talk(
        preflight_vault,
        source_identity=source_identity(),
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "transcript_quality_receipt_unverified" in finding_codes(report, "blocking")


def test_current_v5_missing_declared_artifact_remains_blocking(
    preflight_vault,
    vault_fixture,
):
    talk = current_v5_talk(
        preflight_vault,
        video_url=None,
        youtube_id=None,
        transcript_source="manual",
        transcript_path="transcripts/missing.txt",
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "transcript_artifact_missing" in finding_codes(report, "blocking")


def test_legacy_quality_warning_allows_normalize_to_requeue(
    preflight_vault,
    vault_fixture,
):
    transcript = materialize_transcript(vault_fixture)
    transcript.with_suffix(".quality.json").unlink()
    write_database(vault_fixture, [base_talk()], current=True)

    report = preflight_vault.run_preflight(vault_fixture["root"])
    assert report["ok"] is True
    assert "transcript_quality_receipt_unverified" in finding_codes(report, "warning")

    completed = subprocess.run(
        [
            sys.executable,
            str(QUEUE_STATE_SCRIPT),
            str(vault_fixture["database"]),
            "normalize",
        ],
        cwd=QUEUE_STATE_SCRIPT.parents[3],
        text=True,
        capture_output=True,
        check=False,
        timeout=30,
    )

    assert completed.returncode == 0, completed.stderr
    normalized = json.loads(vault_fixture["database"].read_text(encoding="utf-8"))
    assert normalized["talks"][0]["status"] == "needs-reprocessing"
    assert normalized["talks"][0]["reprocess_reason"].endswith(
        ":missing_generation_status"
    )


def test_quality_receipt_duration_uses_acquisition_level_tolerance(
    preflight_vault,
    vault_fixture,
):
    transcript = materialize_transcript(vault_fixture)
    transcript_timing = importlib.import_module("transcript_timing")
    text = transcript.read_text(encoding="utf-8")
    transcript_timing.write_quality_receipt(
        transcript,
        text,
        transcript_timing.build_quality_policy(400, trusted_duration_seconds=2700.0),
        {
            "kind": "youtube_duration",
            "video_id": VIDEO_ID,
            "duration_seconds": 2700.0,
        },
    )
    # This would have passed the retired max(60 seconds, 5%) comparison.
    talk = base_talk(source_identity=source_identity(duration_seconds=2702.0))
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "transcript_quality_provenance_mismatch" in finding_codes(report, "blocking")


def test_h1b_upload_before_delivery_is_blocking(preflight_vault, vault_fixture):
    """H1b shape: source evidence predates the event it allegedly records."""
    materialize_transcript(vault_fixture)
    talk = base_talk(
        duration_seconds=2700,
        source_identity=source_identity(upload_date="2025-12-31"),
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["database"])

    assert "source_identity_upload_predates_talk" in finding_codes(report, "blocking")
    assert report["ok"] is False


def test_h4_wrong_recording_evidence_is_blocking(preflight_vault, vault_fixture):
    """H4 shape: valid metadata was captured, but for an unrelated recording."""
    materialize_transcript(vault_fixture)
    talk = base_talk(
        duration_seconds=2700,
        source_identity=source_identity(
            video_id=OTHER_VIDEO_ID,
            title="Kubernetes Security from First Principles",
            speakers=["Alice Example"],
            duration_seconds=300,
        ),
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert {
        "source_identity_video_id_mismatch",
        "source_identity_title_mismatch",
        "source_identity_speaker_mismatch",
        "source_identity_duration_mismatch",
    } <= finding_codes(report, "blocking")


def test_short_provider_title_with_matching_event_passes_identity_preflight(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    catalog_title = (
        "Coding Fast and Slow: Applying Kahneman's Insights to Improve "
        "Development Practices and Efficiency"
    )
    talk = base_talk(
        title=catalog_title,
        conference="Dev2Next 2026",
        source_identity=source_identity(
            title="Coding Fast and Slow at Dev2Next 2026",
        ),
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "source_identity_title_mismatch" not in finding_codes(report)
    assert "source_identity_event_mismatch" not in finding_codes(report)


def test_short_provider_title_from_wrong_event_is_blocking(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    catalog_title = (
        "Coding Fast and Slow: Applying Kahneman's Insights to Improve "
        "Development Practices and Efficiency"
    )
    active = base_talk(
        title=catalog_title,
        conference="Conference Alpha 2024",
        date="2024-06-26",
        source_identity=source_identity(
            title="Coding Fast and Slow at Devoxx Poland 2024",
            recorded_date="2024-06-26",
            upload_date="2024-07-01",
        ),
    )
    known_other_event = base_talk(
        filename="2024-06-19-devoxx-poland-catalog-only.md",
        title=catalog_title,
        conference="Devoxx Poland 2024",
        date="2024-06-19",
        video_url=None,
        youtube_id=None,
        transcript_source="none",
        status="skipped_no_sources",
    )
    write_database(vault_fixture, [active, known_other_event])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "source_identity_title_mismatch" not in finding_codes(report)
    assert "source_identity_event_mismatch" in finding_codes(
        report,
        "blocking",
    )


def test_speaker_match_accepts_surname_but_not_unrelated_given_name(preflight_vault):
    assert preflight_vault.names_agree("Baruch Sadogursky", "Sadogursky")
    assert preflight_vault.names_agree(
        "Baruch Sadogursky", "Sadogursky, Baruch (JFrog)"
    )
    assert not preflight_vault.names_agree("Baruch Sadogursky", "Baruch")


def test_url_and_stored_youtube_id_must_agree(preflight_vault, vault_fixture):
    materialize_transcript(vault_fixture, OTHER_VIDEO_ID)
    talk = base_talk(youtube_id=OTHER_VIDEO_ID)
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item for item in report["findings"] if item["code"] == "youtube_id_mismatch"
    )
    assert finding["expected"] == VIDEO_ID
    assert finding["actual"] == OTHER_VIDEO_ID
    assert finding["severity"] == "blocking"


def test_duplicate_video_requires_an_explicit_relation(preflight_vault, vault_fixture):
    materialize_transcript(vault_fixture)
    first = base_talk()
    second = deepcopy(first)
    second["filename"] = "2026-07-31-second-catalog-record.md"
    write_database(vault_fixture, [first, second])

    unqualified = preflight_vault.run_preflight(vault_fixture["root"])
    assert "duplicate_youtube_id" in finding_codes(unqualified, "blocking")

    second["source_relation"] = {
        "type": "borrowed_recording",
        "target_filename": first["filename"],
    }
    write_database(vault_fixture, [first, second])
    qualified = preflight_vault.run_preflight(vault_fixture["root"])
    assert "duplicate_youtube_id" not in finding_codes(qualified)
    assert qualified["blocking_count"] == 0


def test_known_bad_source_cannot_be_reactivated_in_another_url_form(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    talk = base_talk(
        source_rejections=[
            {
                "source_type": "video",
                "url": f"https://youtu.be/{VIDEO_ID}",
                "reason": "wrong_delivery",
                "evidence": "provider metadata names a different conference",
                "verified_at": "2026-07-31T14:00:00-05:00",
            }
        ]
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "rejected_source_reactivated" in finding_codes(report, "blocking")


def test_inactive_well_formed_source_rejection_is_valid(
    preflight_vault,
    vault_fixture,
):
    talk = base_talk(
        video_url=None,
        youtube_id=None,
        transcript_source="none",
        source_rejections=[
            {
                "source_type": "video",
                "url": f"https://youtu.be/{VIDEO_ID}",
                "reason": "non_delivery_clip",
                "evidence": "duration is only 226 seconds",
                "verified_at": "2026-07-31T14:00:00-05:00",
            }
        ],
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert not finding_codes(report)


def test_known_bad_slide_source_cannot_return_in_another_drive_url_form(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    talk = base_talk(
        slides_url=f"https://drive.google.com/open?id={DRIVE_ID}",
        source_rejections=[
            {
                "source_type": "slides",
                "url": f"https://drive.google.com/file/d/{DRIVE_ID}/view",
                "reason": "wrong_delivery",
                "evidence": "the footer names a different conference",
                "verified_at": "2026-07-31T14:00:00-05:00",
            }
        ],
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "rejected_source_reactivated" in finding_codes(report, "blocking")


@pytest.mark.parametrize(
    "bad_rejections",
    [
        {},
        ["not an object"],
        [{"source_type": "video", "url": "https://example.com"}],
        [
            {
                "source_type": "video",
                "url": "https://example.com",
                "reason": "wrong",
                "evidence": "verified",
                "verified_at": "2026-07-31T14:00:00",
            }
        ],
    ],
)
def test_malformed_source_rejection_is_blocking(
    preflight_vault,
    vault_fixture,
    bad_rejections,
):
    materialize_transcript(vault_fixture)
    talk = base_talk(source_rejections=bad_rejections)
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert finding_codes(report, "blocking") & {
        "source_rejections_shape_invalid",
        "source_rejection_invalid",
    }


def test_legacy_duplicate_relation_only_waives_the_same_recording(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    first = base_talk()
    second = deepcopy(first)
    second.update(
        {
            "filename": "legacy-duplicate.md",
            "_duplicate_of": first["filename"],
        }
    )
    write_database(vault_fixture, [first, second])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "duplicate_youtube_id" not in finding_codes(report)
    assert report["blocking_count"] == 0


@pytest.mark.parametrize(
    ("fault_code", "severity"),
    [
        ("slide_source_unsupported", "blocking"),
        ("slide_pptx_reference_missing", "warning"),
        ("slide_pptx_artifact_missing", "warning"),
        ("slide_pdf_reference_missing", "warning"),
        ("slide_pdf_artifact_missing", "warning"),
        ("slide_video_reference_missing", "blocking"),
        ("slide_video_artifact_missing", "warning"),
    ],
)
def test_all_seven_slide_contract_fault_classes_are_reported(
    preflight_vault,
    vault_fixture,
    fault_code,
    severity,
):
    materialize_transcript(vault_fixture)
    talk = base_talk()
    if fault_code == "slide_source_unsupported":
        talk["slide_source"] = "transcript_only"
    elif fault_code == "slide_pptx_reference_missing":
        talk["slide_source"] = "pptx"
    elif fault_code == "slide_pptx_artifact_missing":
        talk.update(slide_source="pptx", pptx_path="missing.pptx")
    elif fault_code == "slide_pdf_reference_missing":
        talk["slide_source"] = "pdf"
    elif fault_code == "slide_pdf_artifact_missing":
        talk.update(slide_source="pdf", google_drive_id=DRIVE_ID)
    elif fault_code == "slide_video_reference_missing":
        talk.update(
            slide_source="video_extracted",
            video_url=None,
            youtube_id=None,
            transcript_source="none",
        )
    elif fault_code == "slide_video_artifact_missing":
        talk["slide_source"] = "video_extracted"
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert fault_code in finding_codes(report, severity)
    assert fault_code in preflight_vault.SLIDE_CONTRACT_CODES


@pytest.mark.parametrize(
    "source_fields,expected_code,expected_reason",
    [
        (
            {
                "slide_source": "pptx",
                "pptx_path": foreign_absolute_locator("talk.pptx"),
            },
            "slide_pptx_artifact_unreadable",
            "pptx_path: artifact_locator_foreign_absolute",
        ),
        (
            {
                "slide_source": "pptx",
                "pptx_path": r"conference\talk.pptx",
            },
            "slide_pptx_artifact_unreadable",
            "pptx_path: artifact_locator_noncanonical_relative",
        ),
        (
            {
                "slide_source": "pptx",
                "pptx_path": "//server/share/talk.pptx",
            },
            "slide_pptx_artifact_unreadable",
            "pptx_path: artifact_locator_ambiguous_double_slash",
        ),
        (
            {
                "slide_source": "pdf",
                "slides_local_path": "~/slides/talk.pdf",
            },
            "slide_pdf_artifact_unreadable",
            "slides_local_path: artifact_locator_home_expansion_unsupported",
        ),
        (
            {
                "slide_source": "pptx",
                "pptx_path": " deck.pptx ",
            },
            "slide_pptx_artifact_unreadable",
            "pptx_path: artifact_locator_empty_or_whitespace",
        ),
        (
            {
                "slide_source": "pdf",
                "slides_local_path": " slides/talk.pdf ",
            },
            "slide_pdf_artifact_unreadable",
            "slides_local_path: artifact_locator_empty_or_whitespace",
        ),
    ],
)
def test_preflight_uses_context_locator_contract_without_rebased_artifact_path(
    preflight_vault,
    vault_fixture,
    source_fields,
    expected_code,
    expected_reason,
) -> None:
    talk = base_talk(
        status="pending",
        transcript_source="none",
        video_url=None,
        youtube_id=None,
        **source_fields,
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(item for item in report["findings"] if item["code"] == expected_code)
    assert finding["artifact_path"] is None
    capability = finding["capability_fact"]
    lane = "native_deck" if source_fields["slide_source"] == "pptx" else "static_slides"
    assert capability["source_reasons"][lane] == expected_reason
    raw_locator = next(
        source_fields[field]
        for field in ("pptx_path", "slides_local_path")
        if field in source_fields
    )
    assert raw_locator not in json.dumps(finding, sort_keys=True)


def test_preflight_rejects_invalid_configured_pptx_root_without_vault_fallback(
    preflight_vault,
    vault_fixture,
) -> None:
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(str(vault_fixture["root"] / "shadow.pptx"))
    talk = base_talk(
        status="pending",
        transcript_source="none",
        video_url=None,
        youtube_id=None,
        slide_source="pptx",
        pptx_path="shadow.pptx",
    )
    write_database(
        vault_fixture,
        [talk],
        config={
            "speaker_name": "Baruch Sadogursky",
            "pptx_source_dir": "relative-presentations",
        },
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "slide_pptx_artifact_unreadable"
    )
    assert finding["artifact_path"] is None
    assert finding["capability_fact"]["source_reasons"]["native_deck"] == (
        "pptx_source_dir: artifact_root_not_native_absolute"
    )


@pytest.mark.parametrize(
    "configured_root,reason_code",
    [
        ("", "artifact_locator_empty_or_whitespace"),
        ("relative-presentations", "artifact_root_not_native_absolute"),
        ("~/presentations", "artifact_locator_home_expansion_unsupported"),
        (
            foreign_absolute_locator("presentations"),
            "artifact_locator_foreign_absolute",
        ),
        (r"\\?\C:\presentations", "artifact_locator_windows_device_namespace"),
    ],
)
def test_preflight_reports_invalid_configured_pptx_root_without_talks(
    preflight_vault,
    vault_fixture,
    configured_root,
    reason_code,
) -> None:
    write_database(
        vault_fixture,
        [],
        config={
            "speaker_name": "Baruch Sadogursky",
            "pptx_source_dir": configured_root,
        },
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = [
        item for item in report["findings"] if item["code"] == "pptx_source_dir_invalid"
    ]
    assert len(findings) == 1
    finding = findings[0]
    assert finding["severity"] == "blocking"
    assert finding["actual"] == {"reason_code": reason_code}
    assert finding["artifact_path"] is None
    if configured_root:
        assert configured_root not in json.dumps(finding, sort_keys=True)


def test_preflight_reports_malformed_pptx_directory_exclusions(
    preflight_vault,
    vault_fixture,
):
    write_database(
        vault_fixture,
        [],
        config={
            "speaker_name": "Baruch Sadogursky",
            "pptx_source_dir": str(vault_fixture["pptx_source"]),
            "pptx_directory_exclusions": ["venv", "VENV"],
        },
        current=True,
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["ok"] is False
    findings = [
        item
        for item in report["findings"]
        if item["code"] == "pptx_directory_exclusions_invalid"
    ]
    assert len(findings) == 1
    assert findings[0]["field"] == "config.pptx_directory_exclusions"
    assert findings[0]["actual"] == ["venv", "VENV"]
    assert "tracking_database_schema_invalid" not in finding_codes(report)


def test_preflight_reports_missing_current_pptx_directory_exclusions_once(
    preflight_vault,
    vault_fixture,
):
    write_database(vault_fixture, [], current=True)
    database = json.loads(vault_fixture["database"].read_text(encoding="utf-8"))
    database["config"].pop("pptx_directory_exclusions")
    vault_fixture["database"].write_text(
        json.dumps(database, indent=2),
        encoding="utf-8",
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = [
        item
        for item in report["findings"]
        if item["code"] == "pptx_directory_exclusions_invalid"
    ]
    assert len(findings) == 1
    assert findings[0]["field"] == "config.pptx_directory_exclusions"
    assert findings[0]["actual"] == {"state": "missing"}
    assert "tracking_database_schema_invalid" not in finding_codes(report)


def test_preflight_preserves_unrelated_schema_fault_with_invalid_exclusions(
    preflight_vault,
    vault_fixture,
):
    write_database(vault_fixture, [], current=True)
    database = json.loads(vault_fixture["database"].read_text(encoding="utf-8"))
    database["config"]["pptx_directory_exclusions"] = ["venv", "VENV"]
    database["pptx_catalog"] = {}
    vault_fixture["database"].write_text(
        json.dumps(database, indent=2),
        encoding="utf-8",
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    blocking_codes = finding_codes(report, "blocking")
    assert "pptx_directory_exclusions_invalid" in blocking_codes
    assert "tracking_database_schema_invalid" in blocking_codes
    assert (
        sum(
            finding["code"] == "pptx_directory_exclusions_invalid"
            for finding in report["findings"]
        )
        == 1
    )
    assert (
        sum(
            finding["code"] == "tracking_database_schema_invalid"
            for finding in report["findings"]
        )
        == 1
    )


def test_preflight_accepts_historical_config_v1_for_owner_migration(
    preflight_vault,
    vault_fixture,
):
    write_database(vault_fixture, [], current=True)
    database = json.loads(vault_fixture["database"].read_text(encoding="utf-8"))
    database["config"]["schema_version"] = 1
    database["config"].pop("pptx_directory_exclusions")
    vault_fixture["database"].write_text(
        json.dumps(database, indent=2),
        encoding="utf-8",
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "pptx_directory_exclusions_invalid" not in finding_codes(report)
    assert "tracking_database_schema_unsupported" not in finding_codes(report)


def test_preflight_fails_closed_on_future_config_generation(
    preflight_vault,
    vault_fixture,
):
    write_database(vault_fixture, [], current=True)
    database = json.loads(vault_fixture["database"].read_text(encoding="utf-8"))
    database["config"]["schema_version"] = 3
    vault_fixture["database"].write_text(
        json.dumps(database, indent=2),
        encoding="utf-8",
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["ok"] is False
    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "tracking_database_schema_unsupported"
    )
    assert finding["severity"] == "blocking"
    assert finding["field"] == "config.schema_version"
    assert finding["expected"] == [1, 2]
    assert finding["actual"] == {
        "schema_version": 3,
        "reason_codes": ["config_schema_version_unsupported"],
    }


def test_preflight_null_configured_pptx_root_uses_vault_fallback(
    preflight_vault,
    vault_fixture,
) -> None:
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(str(vault_fixture["root"] / "deck.pptx"))
    talk = base_talk(
        status="pending",
        transcript_source="none",
        video_url=None,
        youtube_id=None,
        slide_source="pptx",
        pptx_path="deck.pptx",
    )
    write_database(
        vault_fixture,
        [talk],
        config={
            "speaker_name": "Baruch Sadogursky",
            "pptx_source_dir": None,
        },
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "pptx_source_dir_invalid" not in finding_codes(report)
    assert not any(
        code.startswith("slide_pptx_artifact_") for code in finding_codes(report)
    )


def test_preflight_absolute_transcript_locator_cannot_select_existing_bytes(
    preflight_vault,
    vault_fixture,
) -> None:
    external = vault_fixture["root"].parent / "external.txt"
    external.write_text("substantive transcript evidence " * 600, encoding="utf-8")
    talk = base_talk(
        status="pending",
        transcript_source="manual",
        transcript_path=str(external),
        video_url=None,
        youtube_id=None,
        slide_source="none",
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "transcript_reference_missing"
    )
    assert finding["artifact_path"] is None
    assert str(external) not in json.dumps(finding, sort_keys=True)


def test_preflight_never_strips_a_transcript_locator_before_artifact_io(
    preflight_vault,
    vault_fixture,
    monkeypatch,
) -> None:
    transcript = vault_fixture["transcripts"] / "manual.txt"
    transcript.write_text("substantive transcript evidence " * 600, encoding="utf-8")
    raw_locator = " transcripts/manual.txt "
    talk = base_talk(
        status="pending",
        transcript_source="manual",
        transcript_path=raw_locator,
        video_url=None,
        youtube_id=None,
        slide_source="none",
    )
    write_database(vault_fixture, [talk])
    original_is_file = Path.is_file
    original_read_text = Path.read_text

    def guarded_is_file(path: Path) -> bool:
        if path == transcript:
            pytest.fail("stripped transcript locator reached is_file")
        return original_is_file(path)

    def guarded_read_text(path: Path, *args, **kwargs) -> str:
        if path == transcript:
            pytest.fail("stripped transcript locator reached read_text")
        return original_read_text(path, *args, **kwargs)

    monkeypatch.setattr(Path, "is_file", guarded_is_file)
    monkeypatch.setattr(Path, "read_text", guarded_read_text)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "transcript_reference_missing"
    )
    assert finding["artifact_path"] is None
    assert raw_locator not in json.dumps(finding, sort_keys=True)


def test_preflight_never_reads_a_transcript_owned_by_another_youtube_id(
    preflight_vault,
    vault_fixture,
    monkeypatch,
) -> None:
    transcript = vault_fixture["transcripts"] / "other-owner.txt"
    transcript.write_text("substantive transcript evidence " * 600, encoding="utf-8")
    raw_locator = "transcripts/other-owner.txt"
    talk = base_talk(
        status="pending",
        transcript_source="manual",
        transcript_path=raw_locator,
        slide_source="none",
    )
    write_database(vault_fixture, [talk])
    original_is_file = Path.is_file
    original_read_text = Path.read_text

    def guarded_is_file(path: Path) -> bool:
        if path == transcript:
            pytest.fail("mismatched transcript owner reached is_file")
        return original_is_file(path)

    def guarded_read_text(path: Path, *args, **kwargs) -> str:
        if path == transcript:
            pytest.fail("mismatched transcript owner reached read_text")
        return original_read_text(path, *args, **kwargs)

    monkeypatch.setattr(Path, "is_file", guarded_is_file)
    monkeypatch.setattr(Path, "read_text", guarded_read_text)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "transcript_reference_missing"
    )
    assert finding["artifact_path"] is None
    assert raw_locator not in json.dumps(finding, sort_keys=True)


def test_preflight_path_helpers_never_strip_raw_locator_grammar(
    preflight_vault,
    vault_fixture,
) -> None:
    validator = preflight_vault.VaultPreflight(
        {},
        vault_fixture["root"],
        vault_fixture["database"],
    )

    assert validator._pptx_path({"pptx_path": " deck.pptx "}) is None
    assert validator._slide_pdf_path({"slides_local_path": " slides/deck.pdf "}) is None


def test_preflight_foreign_video_manifest_locator_is_path_neutral(
    preflight_vault,
    vault_fixture,
) -> None:
    manifest = context_video_manifest(vault_fixture)
    foreign_video = foreign_absolute_locator(f"{VIDEO_ID}.mp4")
    manifest["source_video_path"] = foreign_video
    for artifact in manifest["artifacts"]:
        artifact["source_video_path"] = foreign_video
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": manifest},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "video_extraction_provenance_invalid"
    )
    assert finding["artifact_path"] is None
    assert "artifact_locator_foreign_absolute" in finding["actual"]
    assert foreign_video not in json.dumps(finding, sort_keys=True)


def test_pending_artifact_gaps_are_warnings_not_blockers(
    preflight_vault,
    vault_fixture,
):
    talk = base_talk(status="pending", slide_source="video_extracted")
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert report["ok"] is True
    assert {
        "transcript_artifact_missing",
        "slide_video_artifact_missing",
    } <= finding_codes(report, "warning")


@pytest.mark.parametrize(
    ("source_fields", "expected_code"),
    [
        (
            {
                "transcript_source": "none",
                "slide_source": "pdf",
                "slides_local_path": "slides/missing.pdf",
            },
            "slide_pdf_artifact_missing",
        ),
        (
            {
                "transcript_source": "none",
                "slide_source": "pptx",
                "pptx_path": "missing.pptx",
            },
            "slide_pptx_artifact_missing",
        ),
        (
            {
                "transcript_source": "manual",
                "transcript_path": "transcripts/missing.txt",
                "slide_source": "none",
            },
            "transcript_artifact_missing",
        ),
    ],
)
def test_pending_nonvideo_sources_still_run_artifact_preflight(
    preflight_vault,
    vault_fixture,
    source_fields,
    expected_code,
):
    talk = base_talk(
        status="pending",
        video_url=None,
        youtube_id=None,
        **source_fields,
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert report["ok"] is True
    assert expected_code in finding_codes(report, "warning")


@pytest.mark.parametrize(
    ("source_fields", "expected_code"),
    [
        (
            {"transcript_source": [], "slide_source": "none"},
            "transcript_source_unsupported",
        ),
        (
            {"transcript_source": "none", "slide_source": {}},
            "slide_source_unsupported",
        ),
    ],
)
def test_pending_unhashable_source_enums_report_instead_of_crashing(
    preflight_vault,
    vault_fixture,
    source_fields,
    expected_code,
):
    talk = base_talk(
        status="pending",
        video_url=None,
        youtube_id=None,
        **source_fields,
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert expected_code in finding_codes(report, "blocking")


def test_recoverable_legacy_processed_artifact_gaps_are_warnings(
    preflight_vault,
    vault_fixture,
):
    talk = base_talk(slide_source="pdf", google_drive_id=DRIVE_ID)
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert {
        "transcript_artifact_missing",
        "slide_pdf_artifact_missing",
    } <= finding_codes(report, "warning")
    assert report["ok"] is True


def test_legacy_missing_remote_pdf_is_a_reacquisition_warning(
    preflight_vault,
    vault_fixture,
):
    talk = base_talk(
        video_url=None,
        youtube_id=None,
        transcript_source="none",
        slide_source="pdf",
        slides_url=f"https://drive.google.com/file/d/{DRIVE_ID}/view",
        google_drive_id=DRIVE_ID,
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["ok"] is True
    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "slide_pdf_artifact_missing"
    )
    assert finding["severity"] == "warning"
    assert finding["capability_fact"]["acquisition_capabilities"] == ["slides"]


def test_crc_damaged_pptx_preflight_is_structured_and_severity_is_current_aware(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    deck, damaged_part = materialize_crc_damaged_pptx(vault_fixture)
    legacy = base_talk(
        slide_source="pptx",
        pptx_path=deck.name,
    )
    write_database(vault_fixture, [legacy])

    warning_report = preflight_vault.run_preflight(vault_fixture["root"])

    warning = next(
        item
        for item in warning_report["findings"]
        if item["code"] == "slide_pptx_artifact_degraded"
    )
    assert warning["severity"] == "warning"
    assert warning["actual"]["archive_recovery"][0]["part_name"] == damaged_part

    current = current_v5_talk(
        preflight_vault,
        slide_source="pptx",
        pptx_path=deck.name,
    )
    write_database(vault_fixture, [current])
    cli_run = subprocess.run(
        [sys.executable, preflight_vault.__file__, str(vault_fixture["root"])],
        capture_output=True,
        text=True,
        check=False,
        timeout=30,
    )

    blocking_report = json.loads(cli_run.stdout)
    blocking = next(
        item
        for item in blocking_report["findings"]
        if item["code"] == "slide_pptx_artifact_degraded"
    )
    assert cli_run.returncode == 1
    assert "Traceback" not in cli_run.stderr
    assert blocking["severity"] == "blocking"
    assert blocking["actual"]["status"] == "degraded_recoverable"


def test_shared_73_slide_damaged_tiff_deck_reports_both_owners_without_promotion(
    preflight_vault,
    vault_fixture,
):
    deck, damaged_part = materialize_shared_crc_damaged_tiff_pptx(vault_fixture)
    talks = []
    for filename, video_id in (
        ("2018-01-01-shared-deck-a.md", VIDEO_ID),
        ("2018-01-02-shared-deck-b.md", OTHER_VIDEO_ID),
    ):
        talks.append(
            base_talk(
                filename=filename,
                title=f"Shared damaged deck {filename[-4]}",
                video_url=f"https://www.youtube.com/watch?v={video_id}",
                youtube_id=video_id,
                transcript_source="none",
                slide_source="pptx",
                pptx_path=deck.name,
            )
        )
    write_database(vault_fixture, talks)

    cli_run = subprocess.run(
        [sys.executable, preflight_vault.__file__, str(vault_fixture["root"])],
        capture_output=True,
        text=True,
        check=False,
        timeout=30,
    )

    report = json.loads(cli_run.stdout)
    findings = [
        item
        for item in report["findings"]
        if item["code"] == "slide_pptx_artifact_degraded"
    ]
    assert cli_run.returncode == 0
    assert "Traceback" not in cli_run.stderr
    assert [item["filename"] for item in findings] == [
        "2018-01-01-shared-deck-a.md",
        "2018-01-02-shared-deck-b.md",
    ]
    assert all(item["severity"] == "warning" for item in findings)
    assert all(
        item["actual"]["archive_recovery"][0]["part_name"] == damaged_part
        for item in findings
    )
    assert all(
        "native_deck" not in item["capability_fact"]["verified_evidence_sources"]
        for item in findings
    )


def test_unrecoverable_pptx_preflight_is_a_structured_unavailable_finding(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    deck = vault_fixture["pptx_source"] / "unrecoverable.pptx"
    deck.write_bytes(b"not an OOXML ZIP package")
    write_database(
        vault_fixture,
        [base_talk(slide_source="pptx", pptx_path=deck.name)],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "slide_pptx_artifact_unreadable"
    )
    assert finding["severity"] == "warning"
    assert finding["actual"]["reason_code"] == "pptx_invalid_container"
    assert "invalid PPTX ZIP container" in finding["actual"]["reason"]
    assert "native_deck" not in finding["capability_fact"]["verified_evidence_sources"]


def test_transcript_only_partial_record_can_repair_a_missing_slide_lane(
    preflight_vault,
    vault_fixture,
):
    transcript = vault_fixture["transcripts"] / "manual.txt"
    text = " ".join(["substantive transcript evidence"] * 600)
    transcript.write_text(text, encoding="utf-8")
    transcript_timing = importlib.import_module("transcript_timing")
    transcript_timing.write_quality_receipt(
        transcript,
        text,
        transcript_timing.build_quality_policy(400),
        {"kind": "fixed_default"},
    )
    talk = base_talk(
        status="processed_partial",
        video_url=None,
        youtube_id=None,
        transcript_source="manual",
        transcript_path="transcripts/manual.txt",
        slide_source="pdf",
        google_drive_id=None,
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["ok"] is True
    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "slide_pdf_reference_missing"
    )
    assert finding["severity"] == "warning"
    assert finding["capability_fact"]["verified_capabilities"] == ["transcript"]


@pytest.mark.parametrize("field", ["slides_local_path", "slides_pdf_path", "pdf_path"])
@pytest.mark.parametrize("slide_source", ["pdf", "video_extracted"])
def test_explicit_local_pdf_path_satisfies_legacy_artifact_contract(
    preflight_vault,
    vault_fixture,
    field,
    slide_source,
):
    materialize_transcript(vault_fixture)
    artifact = vault_fixture["slides"] / "descriptive-legacy-name.pdf"
    write_pdf(artifact)
    talk = base_talk(
        status=(
            "needs-reprocessing" if slide_source == "video_extracted" else "processed"
        ),
        slide_source=slide_source,
        google_drive_id=None,
        **{field: "slides/descriptive-legacy-name.pdf"},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    if slide_source == "video_extracted":
        assert finding_codes(report, "warning") == {
            "video_extraction_provenance_missing"
        }
    else:
        assert not finding_codes(report)


def test_malformed_pdf_is_reported_without_erasing_independent_sources(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    artifact = vault_fixture["slides"] / "malformed.pdf"
    artifact.write_bytes(b"not a PDF")
    deck_path = vault_fixture["pptx_source"] / "independent.pptx"
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(str(deck_path))
    write_database(
        vault_fixture,
        [
            base_talk(
                slide_source="both",
                slides_local_path="slides/malformed.pdf",
                pptx_path="independent.pptx",
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "slide_pdf_artifact_unreadable"
    )
    capability = finding["capability_fact"]
    assert set(capability["verified_evidence_sources"]) == {
        "native_deck",
        "transcript",
    }
    unavailable = capability["unavailable_evidence_sources"]["static_slides"]
    assert unavailable["reason_code"] == "pdf_invalid_container"


@pytest.mark.parametrize(
    ("reason_code", "expected_code"),
    [
        ("pdf_cloud_placeholder_unavailable", "slide_pdf_artifact_unavailable"),
        ("pdf_probe_timeout", "slide_pdf_artifact_unreadable"),
        ("pdf_dependency_unavailable", "slide_pdf_artifact_unreadable"),
        ("pdf_probe_monitor_unavailable", "slide_pdf_artifact_unreadable"),
        ("pdf_probe_monitor_identity_changed", "slide_pdf_artifact_unreadable"),
        ("pdf_probe_containment_unavailable", "slide_pdf_artifact_unreadable"),
        ("pdf_probe_resource_unavailable", "slide_pdf_artifact_unreadable"),
    ],
)
def test_pdf_probe_failures_are_lane_local_and_structured(
    preflight_vault,
    vault_fixture,
    monkeypatch,
    reason_code,
    expected_code,
):
    materialize_transcript(vault_fixture)
    artifact = write_pdf(vault_fixture["slides"] / "bounded.pdf")
    deck_path = vault_fixture["pptx_source"] / "independent.pptx"
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(str(deck_path))
    pattern_evidence = importlib.import_module("pattern_evidence")
    pdf_evidence = importlib.import_module("pdf_evidence")

    def fail_pdf_probe(*_args, **_kwargs):
        raise pdf_evidence.PdfEvidenceError(
            "synthetic bounded PDF failure",
            reason_code=reason_code,
        )

    monkeypatch.setattr(pattern_evidence, "probe_pdf_artifact", fail_pdf_probe)
    write_database(
        vault_fixture,
        [
            base_talk(
                slide_source="both",
                slides_local_path=artifact.relative_to(
                    vault_fixture["root"]
                ).as_posix(),
                pptx_path="independent.pptx",
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(item for item in report["findings"] if item["code"] == expected_code)
    if expected_code == "slide_pdf_artifact_unreadable":
        assert finding["message"] == (
            "declared slide PDF could not complete bounded evidence inspection"
        )
    capability = finding["capability_fact"]
    assert set(capability["verified_evidence_sources"]) == {
        "native_deck",
        "transcript",
    }
    unavailable = capability["unavailable_evidence_sources"]["static_slides"]
    assert unavailable["reason_code"] == reason_code


@pytest.mark.parametrize(
    ("slide_source", "expected_code"),
    [
        ("pdf", "slide_pdf_artifact_missing"),
        ("video_extracted", "slide_video_artifact_missing"),
    ],
)
def test_missing_explicit_local_pdf_does_not_fall_back_to_another_identity(
    preflight_vault,
    vault_fixture,
    slide_source,
    expected_code,
):
    materialize_transcript(vault_fixture)
    write_pdf(vault_fixture["slides"] / f"{VIDEO_ID}.pdf")
    write_pdf(vault_fixture["slides"] / f"{DRIVE_ID}.pdf")
    talk = base_talk(
        slide_source=slide_source,
        google_drive_id=DRIVE_ID,
        slides_local_path="slides/missing-explicit.pdf",
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert expected_code in finding_codes(report, "warning")


def test_video_pdf_page_count_is_never_treated_as_authored_slide_count(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    # The bounded probe verifies the real page tree, but that delivered-frame
    # count is never reinterpreted as authored slide_count metadata.
    write_pdf(vault_fixture["slides"] / f"{VIDEO_ID}.pdf", page_count=99)
    talk = base_talk(
        slide_source="video_extracted",
        slide_count=3,
        structured_data={
            "video_extraction": trusted_video_manifest(vault_fixture, 99),
        },
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert not any("slide_count" in item["code"] for item in report["findings"])


@pytest.mark.parametrize(
    "manifest_factory",
    [
        trusted_video_manifest,
        context_video_manifest,
    ],
)
def test_processed_partial_video_manifest_needs_no_promoted_deck(
    preflight_vault,
    vault_fixture,
    manifest_factory,
):
    materialize_transcript(vault_fixture)
    manifest = manifest_factory(vault_fixture)
    write_database(
        vault_fixture,
        [
            base_talk(
                status="processed_partial",
                slide_source="video_extracted",
                structured_data={"video_extraction": manifest},
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert not finding_codes(report)


def test_context_video_manifest_cannot_back_a_promoted_deck(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    promoted = vault_fixture["slides"] / f"{VIDEO_ID}.pdf"
    write_pdf(promoted)
    manifest = context_video_manifest(vault_fixture)
    write_database(
        vault_fixture,
        [
            base_talk(
                status="processed_partial",
                slide_source="video_extracted",
                slides_local_path=f"slides/{VIDEO_ID}.pdf",
                structured_data={"video_extraction": manifest},
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "video_extraction_untrusted" in finding_codes(report, "blocking")


def test_completed_legacy_video_pdf_without_provenance_is_repairable(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    write_pdf(vault_fixture["slides"] / f"{VIDEO_ID}.pdf")
    write_database(
        vault_fixture,
        [base_talk(slide_source="video_extracted")],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "video_extraction_provenance_missing" in finding_codes(report, "warning")


def test_requeued_legacy_video_pdf_without_provenance_is_a_warning(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    write_pdf(vault_fixture["slides"] / f"{VIDEO_ID}.pdf")
    write_database(
        vault_fixture,
        [base_talk(status="needs-reprocessing", slide_source="video_extracted")],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert "video_extraction_provenance_missing" in finding_codes(report, "warning")


def test_unverified_video_crop_cannot_support_completed_deck_analysis(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    write_pdf(vault_fixture["slides"] / f"{VIDEO_ID}.pdf")
    manifest = trusted_video_manifest(vault_fixture)
    rebuild = vault_fixture["root"] / "slides-rebuild" / VIDEO_ID
    context = rebuild / f"{VIDEO_ID}.context.pdf"
    write_pdf(context)
    manifest.update(
        {
            "slide_region_detected": True,
            "slide_region_method": "auto",
            "slide_region_verified": False,
            "review_required": True,
            "review_reason": "auto crop needs review",
        }
    )
    manifest["artifacts"][0].update(
        {
            "crop_method": "auto",
            "crop_verified": False,
            "trusted_for_authored_slide_analysis": False,
        }
    )
    manifest["artifacts"].append(
        {
            "path": str(context),
            "artifact_scope": "full_frame_context",
            "page_count": manifest["unique_frame_count"],
            "source_video_id": VIDEO_ID,
            "source_video_path": manifest["source_video_path"],
            "crop_method": "none",
            "crop_verified": False,
            "trusted_for_authored_slide_analysis": False,
        }
    )
    write_database(
        vault_fixture,
        [
            base_talk(
                slide_source="video_extracted",
                structured_data={"video_extraction": manifest},
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "video_extraction_untrusted" in finding_codes(report, "blocking")


def _observation_block(**lanes):
    """A current persisted block, as the canonical writer emits it."""
    block = {
        "patterns_detected": [],
        "antipatterns_detected": [],
        "not_evaluable": [],
    }
    block.update(lanes)
    return block


def _catalog_entry(preflight_vault, entry_type="pattern"):
    catalog = preflight_vault.load_catalog()
    return next(
        entry
        for entry in sorted(catalog.entries.values(), key=lambda item: item.pattern_id)
        if entry.observable and entry.entry_type == entry_type
    )


def test_a_corrupt_persisted_block_blocks_before_any_claim(
    preflight_vault,
    vault_fixture,
):
    """The swapped-field signature from the live vault (#167).

    The record-schema check reads the talk's own fields; nothing read the
    nested block, so this reached rendered analyses.
    """
    entry = _catalog_entry(preflight_vault)
    materialize_transcript(vault_fixture)
    write_database(
        vault_fixture,
        [
            base_talk(
                pattern_observations=_observation_block(
                    patterns_detected=[
                        {
                            "pattern_id": entry.pattern_id,
                            "confidence": "strong",
                            "evidence": list(entry.vault_dimensions),
                            "dimensions": "The speaker opened with the map.",
                        }
                    ]
                )
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = [
        item
        for item in report["findings"]
        if item["code"].startswith("persisted_observations_")
    ]
    assert [item["code"] for item in findings] == [
        "persisted_observations_detection_fields_swapped"
    ]
    assert findings[0]["severity"] == "blocking"
    assert report["ok"] is False


@pytest.mark.parametrize(
    ("observations", "expected_field"),
    [
        ([{"pattern_id": "x"}], "pattern_observations"),
        ("not a block", "pattern_observations"),
        (
            {"patterns_detected": [], "antipatterns_detected": []},
            "pattern_observations.not_evaluable",
        ),
    ],
)
def test_the_finding_names_the_field_a_repair_would_edit(
    preflight_vault,
    vault_fixture,
    observations,
    expected_field,
):
    """A root finding already names the block; a lane finding is relative to it.

    Prefixing both yields `pattern_observations.pattern_observations`, which
    points a repair at a field that does not exist.
    """
    materialize_transcript(vault_fixture)
    write_database(vault_fixture, [base_talk(pattern_observations=observations)])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    fields = [
        item["field"]
        for item in report["findings"]
        if item["code"].startswith("persisted_observations_")
    ]
    assert fields == [expected_field]


def test_an_unreadable_catalog_says_how_to_recover(
    preflight_vault,
    vault_fixture,
    monkeypatch,
):
    """`error-handling` Actionable Messages: name the fix, not just the fault."""
    materialize_transcript(vault_fixture)
    write_database(
        vault_fixture,
        [
            base_talk(
                pattern_observations={
                    "patterns_detected": [],
                    "antipatterns_detected": [],
                    "not_evaluable": [],
                }
            )
        ],
    )

    attempts = []

    def unreadable():
        attempts.append(1)
        raise preflight_vault.ReturnValidationError("no pattern entries found")

    monkeypatch.setattr(preflight_vault, "load_catalog", unreadable)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = [
        item
        for item in report["findings"]
        if item["code"] == "pattern_catalog_unreadable"
    ]
    assert len(findings) == 1
    assert "rerun preflight" in findings[0]["message"]
    assert "Restore it" in findings[0]["message"]
    assert attempts == [1]


def test_an_unreadable_catalog_is_one_finding_for_the_whole_run(
    preflight_vault,
    vault_fixture,
    monkeypatch,
):
    """One condition, one finding.

    Repeating it per analyzed talk inflates the blocking count and buries
    every other finding in the report.
    """
    block = {
        "patterns_detected": [],
        "antipatterns_detected": [],
        "not_evaluable": [],
    }
    materialize_transcript(vault_fixture)
    write_database(
        vault_fixture,
        [
            base_talk(pattern_observations=block),
            base_talk(
                filename="2026-07-31-second-talk.md",
                video_url="https://www.youtube.com/watch?v=SECONDVIDEO",
                youtube_id="SECONDVIDEO",
                pattern_observations=block,
            ),
        ],
    )
    monkeypatch.setattr(
        preflight_vault,
        "load_catalog",
        lambda: (_ for _ in ()).throw(
            preflight_vault.ReturnValidationError("no pattern entries found")
        ),
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert [
        item["code"]
        for item in report["findings"]
        if item["code"] == "pattern_catalog_unreadable"
    ] == ["pattern_catalog_unreadable"]


def test_an_unobservable_entry_warns_without_blocking(
    preflight_vault,
    vault_fixture,
):
    """The catalog moved; the record did not. 641 live detections are these."""
    catalog = preflight_vault.load_catalog()
    archival = next(
        entry
        for entry in sorted(catalog.entries.values(), key=lambda item: item.pattern_id)
        if not entry.observable
    )
    lane = (
        "antipatterns_detected"
        if archival.entry_type == "antipattern"
        else "patterns_detected"
    )
    materialize_transcript(vault_fixture)
    write_database(
        vault_fixture,
        [
            base_talk(
                pattern_observations=_observation_block(
                    **{
                        lane: [
                            {
                                "pattern_id": archival.pattern_id,
                                "confidence": "strong",
                                "evidence": "The speaker did the thing, at 12:03.",
                                "dimensions": list(archival.vault_dimensions),
                            }
                        ]
                    }
                )
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = [
        item
        for item in report["findings"]
        if item["code"].startswith("persisted_observations_")
    ]
    assert [item["severity"] for item in findings] == ["warning"]
    assert report["blocking_count"] == 0


@pytest.mark.parametrize("talk_updates", [{}, {"pattern_observations": None}])
def test_an_unscored_talk_is_not_reported_as_corrupt(
    preflight_vault,
    vault_fixture,
    talk_updates,
):
    """Absence is incompleteness, not corruption.

    A legacy record preserves the field as an explicit null rather than
    dropping it, and that is the same "never scored" state. Blocking — or even
    warning on — every talk that predates pattern scoring would flood a queue
    that is working.
    """
    materialize_transcript(vault_fixture)
    write_database(vault_fixture, [base_talk(**talk_updates)])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert [
        item
        for item in report["findings"]
        if item["code"].startswith("persisted_observations_")
    ] == []


def test_a_pending_talk_is_never_held_to_the_current_block_contract(
    preflight_vault,
    vault_fixture,
):
    """A talk not claiming analysis carries whatever the queue left it."""
    materialize_transcript(vault_fixture)
    write_database(
        vault_fixture,
        [base_talk(status="pending", pattern_observations={"patterns_detected": "x"})],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert [
        item
        for item in report["findings"]
        if item["code"].startswith("persisted_observations_")
    ] == []


def test_absent_legacy_identity_metadata_is_not_a_finding(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    write_database(vault_fixture, [base_talk()])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["findings"] == []


def test_partial_identity_metadata_gaps_are_warnings(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    write_database(
        vault_fixture,
        [base_talk(source_identity={"schema_version": 1, "provider": "youtube"})],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 0
    assert {
        "source_identity_video_id_missing",
        "source_identity_title_missing",
        "source_identity_speakers_missing",
        "source_identity_date_missing",
        "source_identity_duration_seconds_missing",
    } <= finding_codes(report, "warning")


@pytest.mark.parametrize(
    ("field", "value", "code"),
    [
        ("transcript_source", "captions_maybe", "transcript_source_unsupported"),
        ("transcript_source", ["youtube_auto"], "transcript_source_unsupported"),
        ("slide_source", "transcript_only", "slide_source_unsupported"),
        ("slide_source", ["pdf"], "slide_source_unsupported"),
    ],
)
def test_source_enums_are_closed(
    preflight_vault,
    vault_fixture,
    field,
    value,
    code,
):
    materialize_transcript(vault_fixture)
    write_database(vault_fixture, [base_talk(**{field: value})])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert code in finding_codes(report, "blocking")


@pytest.mark.parametrize("status", ["skipped_no_sources", "skipped_no_video"])
def test_source_less_skip_status_cannot_hide_pdf_source(
    preflight_vault,
    vault_fixture,
    status,
):
    talk = base_talk(
        status=status,
        video_url=None,
        youtube_id=None,
        transcript_source="none",
        slide_source="pdf",
        slides_url=f"https://drive.google.com/file/d/{DRIVE_ID}/view",
        google_drive_id=DRIVE_ID,
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(
        item
        for item in report["findings"]
        if item["code"] == "status_source_reachability_conflict"
    )
    assert finding["severity"] == "blocking"
    assert finding["actual"] == {
        "status": status,
        "independent_sources": ["pdf"],
    }


def test_source_less_skip_status_cannot_hide_pptx_source(
    preflight_vault,
    vault_fixture,
):
    talk = base_talk(
        status="skipped_no_video",
        video_url=None,
        youtube_id=None,
        transcript_source="none",
        slide_source="pptx",
        pptx_path="Conference/Talk.pptx",
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "status_source_reachability_conflict" in finding_codes(report, "blocking")


def test_video_extraction_without_independent_slides_is_not_reachable(
    preflight_vault,
    vault_fixture,
):
    talk = base_talk(
        status="skipped_no_sources",
        video_url=None,
        youtube_id=None,
        transcript_source="none",
        slide_source="video_extracted",
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "status_source_reachability_conflict" not in finding_codes(report)


def test_filenames_must_be_nonempty_and_unique(preflight_vault, vault_fixture):
    talks = [
        {"filename": " ", "status": "skipped_no_sources"},
        {"filename": "same.md", "status": "skipped_no_sources"},
        {"filename": "same.md", "status": "skipped_no_sources"},
    ]
    write_database(vault_fixture, talks)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert {"filename_missing", "duplicate_filename"} <= finding_codes(
        report, "blocking"
    )


def test_source_indexes_survive_a_malformed_talk(preflight_vault, vault_fixture):
    talks = [
        None,
        {"filename": "same.md", "status": "skipped_no_sources"},
        {"filename": "same.md", "status": "skipped_no_sources"},
    ]
    write_database(vault_fixture, talks)

    report = preflight_vault.run_preflight(vault_fixture["root"])

    duplicate = next(
        item for item in report["findings"] if item["code"] == "duplicate_filename"
    )
    assert duplicate["actual"] == [1, 2]


def test_unhashable_relation_type_is_reported_not_raised(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    talk = base_talk(
        source_relation={
            "type": ["duplicate"],
            "target_filename": "other.md",
        }
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "source_relation_invalid" in finding_codes(report, "blocking")


def test_identity_date_and_duration_types_are_validated(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    evidence = source_identity(
        recorded_date=20260730,
        upload_date=None,
        duration_seconds=float("inf"),
    )
    write_database(vault_fixture, [base_talk(source_identity=evidence)])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert finding_codes(report, "blocking") == {"database_json_invalid"}
    # The public message no longer echoes the rejected literal (#200).
    assert report["findings"][0]["message"] == (
        "tracking database contains a non-standard JSON number"
    )
    assert "Infinity" not in json.dumps(report["findings"], sort_keys=True)
    # The report remains strict JSON even when Python's input decoder accepted
    # a non-standard Infinity token from a legacy artifact.
    json.dumps(report, allow_nan=False)


def test_identity_date_requires_hyphenated_iso_form(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    evidence = source_identity(recorded_date="20260730", upload_date=None)
    write_database(vault_fixture, [base_talk(source_identity=evidence)])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "source_identity_date_invalid" in finding_codes(report, "blocking")


@pytest.mark.parametrize(
    ("field", "value", "code"),
    [
        ("uploader", [], "source_identity_provider_fact_invalid"),
        ("uploader_id", " ", "source_identity_provider_fact_invalid"),
        (
            "webpage_url",
            "https://videos.example.com/watch?v=AbCdEfGhI_1",
            "source_identity_webpage_url_invalid",
        ),
        (
            "webpage_url",
            f"https://www.youtube.com/watch?v={OTHER_VIDEO_ID}",
            "source_identity_webpage_identity_mismatch",
        ),
        ("webpage_video_id", "too-short", "source_identity_webpage_video_id_invalid"),
        (
            "webpage_video_id",
            OTHER_VIDEO_ID,
            "source_identity_webpage_identity_mismatch",
        ),
        ("captured_at", "2026-07-31T12:00:00", "source_identity_captured_at_invalid"),
    ],
)
def test_identity_provider_facts_are_validated_when_present(
    preflight_vault,
    vault_fixture,
    field,
    value,
    code,
):
    materialize_transcript(vault_fixture)
    evidence = source_identity(**{field: value})
    write_database(
        vault_fixture,
        [
            base_talk(
                duration_seconds=2700,
                source_identity=evidence,
            )
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    finding = next(item for item in report["findings"] if item["code"] == code)
    assert finding["severity"] == "blocking"
    assert finding["field"] == f"source_identity.{field}"


def test_live_audit_provider_proposal_satisfies_preflight_contract(
    preflight_vault,
    audit_source_identities,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    metadata = {
        "id": VIDEO_ID,
        "title": "Perfect Vault Ingress",
        "uploader": "Conference Channel",
        "uploader_id": "@conference",
        "upload_date": "20260731",
        "duration": 2700,
        "webpage_url": f"https://www.youtube.com/watch?v={VIDEO_ID}",
    }
    evidence, faults = audit_source_identities.provider_evidence(
        VIDEO_ID,
        metadata,
        "2026-07-31T12:00:00Z",
    )
    proposal = audit_source_identities.proposed_source_identity(evidence)
    assert faults == []

    talk = base_talk(duration_seconds=2700, source_identity=proposal)
    write_database(vault_fixture, [talk])
    clean = preflight_vault.run_preflight(vault_fixture["root"])

    assert clean["blocking_count"] == 0
    assert finding_codes(clean, "warning") == {"source_identity_speakers_missing"}

    corrupt = deepcopy(proposal)
    corrupt["webpage_video_id"] = OTHER_VIDEO_ID
    write_database(
        vault_fixture,
        [
            base_talk(
                duration_seconds=2700,
                source_identity=corrupt,
            )
        ],
    )
    rejected = preflight_vault.run_preflight(vault_fixture["root"])

    assert "source_identity_webpage_identity_mismatch" in finding_codes(
        rejected,
        "blocking",
    )


def test_boolean_identity_schema_version_is_not_version_one(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    evidence = source_identity(schema_version=True)
    write_database(vault_fixture, [base_talk(source_identity=evidence)])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert "source_identity_schema_unsupported" in finding_codes(report, "warning")


def test_report_is_deterministic_and_preflight_is_read_only(
    preflight_vault,
    vault_fixture,
):
    materialize_transcript(vault_fixture)
    database = write_database(vault_fixture, [base_talk()])
    before = database.read_bytes()

    first = preflight_vault.run_preflight(database)
    second = preflight_vault.run_preflight(vault_fixture["root"])

    assert first == second
    assert database.read_bytes() == before


def test_future_tracking_database_is_blocking_no_usable_state(
    preflight_vault,
    vault_fixture,
):
    database = write_database(vault_fixture, [base_talk()])
    payload = json.loads(database.read_text(encoding="utf-8"))
    payload["schema_version"] = 99
    database.write_text(json.dumps(payload), encoding="utf-8")
    before = database.read_bytes()

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert finding_codes(report, "blocking") == {"tracking_database_schema_unsupported"}
    assert database.read_bytes() == before


def test_cli_emits_json_and_only_blocks_on_integrity_errors(
    preflight_vault,
    vault_fixture,
):
    pending = base_talk(status="pending", slide_source="video_extracted")
    write_database(vault_fixture, [pending])
    warning_run = subprocess.run(
        [sys.executable, preflight_vault.__file__, str(vault_fixture["root"])],
        capture_output=True,
        text=True,
        check=False,
        timeout=30,
    )
    warning_report = json.loads(warning_run.stdout)
    assert warning_run.returncode == 0
    assert warning_report["warning_count"] > 0
    assert warning_report["blocking_count"] == 0

    write_database(vault_fixture, [base_talk(youtube_id=OTHER_VIDEO_ID)])
    blocking_run = subprocess.run(
        [sys.executable, preflight_vault.__file__, str(vault_fixture["database"])],
        capture_output=True,
        text=True,
        check=False,
        timeout=30,
    )
    blocking_report = json.loads(blocking_run.stdout)
    assert blocking_run.returncode == 1
    assert blocking_report["blocking_count"] > 0


def test_unreadable_database_is_a_structured_blocking_report(
    preflight_vault,
    vault_fixture,
):
    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 1
    assert report["findings"][0]["code"] == "database_unreadable"


def test_malformed_json_is_a_structured_blocking_report(
    preflight_vault,
    vault_fixture,
):
    vault_fixture["database"].write_text('{"talks": [}', encoding="utf-8")

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 1
    assert report["findings"][0]["code"] == "database_json_invalid"


def test_non_utf8_database_is_a_structured_blocking_report(
    preflight_vault,
    vault_fixture,
):
    vault_fixture["database"].write_bytes(b"\xff\xfe\x00")

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert report["blocking_count"] == 1
    assert report["findings"][0]["code"] == "database_encoding_invalid"


# --- #203: the CLI has a closed failure boundary ---


def test_outer_boundary_emits_one_blocking_report_without_a_traceback(
    preflight_vault, capsys, monkeypatch
):
    """An unexpected failure still produces the machine-readable blocking signal.

    Callers gate claiming on this report. A traceback with no JSON would read as
    "preflight did not run", and a partial document would not parse.
    """

    def explode(*_args, **_kwargs):
        raise RuntimeError("injected outer failure at /private/vault/secret.md")

    monkeypatch.setattr(preflight_vault, "main", explode)
    assert preflight_vault.run_cli() == 2

    captured = capsys.readouterr()
    report = json.loads(captured.out)  # exactly one valid document
    assert report["ok"] is False
    assert report["blocking_count"] == 1
    finding = report["findings"][0]
    assert finding["code"] == "preflight_unexpected_failure"
    assert finding["severity"] == "blocking"
    assert finding["actual"] == "RuntimeError"
    # Path-neutral: neither the exception text nor its path reaches the report.
    assert "injected outer failure" not in captured.out
    assert "/private/vault/secret.md" not in captured.out
    assert "Traceback" not in captured.out
    assert "do not begin claiming" in captured.err


def test_outer_boundary_finding_matches_the_normal_finding_shape(
    preflight_vault, vault_fixture, capsys, monkeypatch
):
    """The failure finding must parse like every other finding."""
    write_database(vault_fixture, [base_talk()], current=True)
    normal = preflight_vault.run_preflight(vault_fixture["root"])
    capsys.readouterr()

    monkeypatch.setattr(
        preflight_vault,
        "main",
        lambda *a, **k: (_ for _ in ()).throw(RuntimeError("boom")),
    )
    assert preflight_vault.run_cli() == 2

    failure = json.loads(capsys.readouterr().out)
    # Superset, not equality: the contract is that a consumer reading any
    # normal key cannot KeyError on the failure shape. The failure report adds
    # `origin`, which no normal finding carries and no consumer reads.
    assert set(normal["findings"][0]) <= set(failure["findings"][0])
    assert set(normal) <= set(failure)
    assert failure["schema_version"] == normal["schema_version"]


def test_outer_boundary_passes_a_clean_exit_code_through(preflight_vault, monkeypatch):
    """A normal run's exit code is not rewritten by the boundary."""
    monkeypatch.setattr(preflight_vault, "main", lambda *a, **k: 0)
    assert preflight_vault.run_cli() == 0
    monkeypatch.setattr(preflight_vault, "main", lambda *a, **k: 1)
    assert preflight_vault.run_cli() == 1


def test_failure_guidance_names_a_concrete_next_action(
    preflight_vault, capsys, monkeypatch
):
    """Telling the operator to repair whatever the exception type named is not
    an actionable message."""
    monkeypatch.setattr(
        preflight_vault,
        "main",
        lambda *a, **k: (_ for _ in ()).throw(RuntimeError("boom")),
    )
    assert preflight_vault.run_cli() == 2
    err = capsys.readouterr().err
    assert "code locations that failed" in err
    assert "check-runtime.py" in err
    assert "do not begin claiming" in err


def test_failure_report_names_code_locations_without_exception_text(
    preflight_vault, capsys, monkeypatch
):
    """no-secrets forbids exception text; the origin frames replace it."""

    def explode(*_args, **_kwargs):
        raise RuntimeError("token=SECRET at /private/vault/creds.json")

    monkeypatch.setattr(preflight_vault, "main", explode)
    assert preflight_vault.run_cli() == 2

    captured = capsys.readouterr()
    report = json.loads(captured.out)
    origin = report["findings"][0]["origin"]
    assert origin, "the failing code location must be reported"
    assert all(":" in frame and " in " in frame for frame in origin)
    # Neither the message, the credential, nor a full path may appear anywhere.
    for stream in (captured.out, captured.err):
        assert "SECRET" not in stream
        assert "/private/vault/creds.json" not in stream
        assert "Traceback" not in stream
    assert "code locations that failed" in captured.err


# --- #200: public diagnostics route on typed reasons, not exception prose ---


@pytest.mark.parametrize(
    "payload,expected_code",
    [
        (b"\xff\xfe not utf-8", "database_encoding_invalid"),
        (b"{not json", "database_json_invalid"),
        (b'{"a": 1, "a": 2}', "database_json_invalid"),
        (b'{"a": NaN}', "database_json_invalid"),
        (b'["not", "an", "object"]', "database_json_invalid"),
    ],
)
def test_database_read_finding_code_comes_from_the_typed_reason(
    preflight_vault, tmp_path, payload, expected_code
):
    """The public taxonomy must not depend on upstream message wording."""
    database = tmp_path / "tracking-database.json"
    database.write_bytes(payload)

    report = preflight_vault.run_preflight(database)
    codes = [item["code"] for item in report["findings"]]
    assert expected_code in codes, codes


def test_unmapped_decoder_reason_falls_back_to_unreadable(preflight_vault):
    """An unrecognised reason must not invent a code."""
    assert preflight_vault._DATABASE_READ_DIAGNOSTICS.get("brand_new_reason") is None
    assert preflight_vault._DATABASE_READ_FALLBACK[0] == "database_unreadable"


def test_transcript_decode_failure_reports_a_typed_reason_not_os_text(
    preflight_vault, vault_fixture
):
    """`actual` is public: a raw OSError string carries the host path."""
    # The transcript resolves from youtube_id, so the file must use that name.
    transcript = vault_fixture["root"] / "transcripts" / f"{VIDEO_ID}.txt"
    transcript.parent.mkdir(parents=True, exist_ok=True)
    transcript.write_bytes(b"\xff\xfe invalid utf-8 bytes")
    write_database(vault_fixture, [base_talk(status="processed")], current=True)

    report = preflight_vault.run_preflight(vault_fixture["root"])
    finding = next(
        (
            f
            for f in report["findings"]
            if f["code"] == "transcript_artifact_unreadable"
        ),
        None,
    )
    assert finding is not None, [f["code"] for f in report["findings"]]
    assert finding["actual"] == "not_utf8"
    assert finding["message"].startswith(
        "transcript artifact is not valid UTF-8 speech text"
    )
    assert "rerun preflight" in finding["message"], (
        "error-handling requires the message to say what to do next"
    )
    # `artifact_path` is a documented structured field and legitimately carries
    # an absolute path (#200). Every OTHER field must stay free of raw decoder
    # prose and host paths.
    rest = {k: v for k, v in finding.items() if k != "artifact_path"}
    serialized = json.dumps(rest, sort_keys=True)
    assert "codec" not in serialized
    assert "invalid start byte" not in serialized
    assert str(vault_fixture["root"]) not in serialized


def test_transcript_read_failure_does_not_claim_a_decode_failure(
    preflight_vault, vault_fixture, monkeypatch
):
    """A permission denial is a read failure — the message must say so (#253)."""
    transcript = vault_fixture["root"] / "transcripts" / f"{VIDEO_ID}.txt"
    transcript.parent.mkdir(parents=True, exist_ok=True)
    transcript.write_text("speech text", encoding="utf-8")
    write_database(vault_fixture, [base_talk(status="processed")], current=True)

    real_read_text = Path.read_text

    def denied(self, *args, **kwargs):
        if self == transcript:
            raise PermissionError(13, "Permission denied", str(transcript))
        return real_read_text(self, *args, **kwargs)

    monkeypatch.setattr(Path, "read_text", denied)

    report = preflight_vault.run_preflight(vault_fixture["root"])
    finding = next(
        (
            f
            for f in report["findings"]
            if f["code"] == "transcript_artifact_unreadable"
        ),
        None,
    )
    assert finding is not None, [f["code"] for f in report["findings"]]
    assert finding["actual"] == "unreadable:PermissionError"
    assert finding["message"].startswith("transcript artifact could not be read")
    assert "rerun preflight" in finding["message"], (
        "error-handling requires the message to say what to do next"
    )
    # The location lives in the documented `artifact_path` field (#200); neither
    # the message nor `actual` may carry the errno prose or the host path.
    for field in ("message", "actual"):
        assert "Permission denied" not in finding[field]
        assert str(vault_fixture["root"]) not in finding[field]


def test_reworded_decoder_message_keeps_its_finding_code(
    preflight_vault, tmp_path, monkeypatch
):
    """The taxonomy survives upstream rewording — that is the point of the fix.

    Substring matching on the message would fall through to the generic
    `database_unreadable` as soon as the wording changed.
    """
    io_module = importlib.import_module("tracking_database_io")
    database = tmp_path / "tracking-database.json"
    database.write_text("{}", encoding="utf-8")

    def reworded(*_args, **_kwargs):
        raise io_module.TrackingDatabaseIOError(
            "the file could not be interpreted as text",  # no legacy substring
            reason_code="encoding_invalid",
        )

    monkeypatch.setattr(preflight_vault, "decode_json_object", reworded)

    report = preflight_vault.run_preflight(database)
    codes = [item["code"] for item in report["findings"]]
    assert "database_encoding_invalid" in codes, codes
    assert "database_unreadable" not in codes


def test_manifest_rejection_reports_a_typed_reason_not_the_rejected_value(
    preflight_vault, vault_fixture
):
    """ReturnValidationError text can echo the malformed input it rejected."""
    poisoned = "SENSITIVE_VALUE_abc123"
    talk = base_talk(
        status="processed_partial",
        transcript_source="none",
        slide_source="video_extracted",
        structured_data={"video_extraction": {"schema_version": poisoned}},
    )
    write_database(vault_fixture, [talk])

    report = preflight_vault.run_preflight(vault_fixture["root"])
    finding = next(
        (
            f
            for f in report["findings"]
            if f["code"] == "video_extraction_provenance_invalid"
        ),
        None,
    )
    assert finding is not None, [f["code"] for f in report["findings"]]
    assert finding["actual"].startswith("video_extraction.")
    assert poisoned not in json.dumps(finding, sort_keys=True)


@pytest.mark.parametrize(
    "payload,leaked",
    [
        (
            b'{"a": 1, "SENSITIVE_KEY_abc": 2, "SENSITIVE_KEY_abc": 3}',
            "SENSITIVE_KEY_abc",
        ),
        (
            b'{"a": 123456789012345678901234567890123456789.5}',
            "123456789012345678901234567890",
        ),
    ],
)
def test_database_read_report_never_echoes_malformed_input(
    preflight_vault, tmp_path, payload, leaked
):
    """Decoder messages embed the rejected key or value; the report must not."""
    database = tmp_path / "tracking-database.json"
    database.write_bytes(payload)

    report = preflight_vault.run_preflight(database)
    serialized = json.dumps(
        {k: v for k, v in report.items() if k not in {"database", "vault_root"}},
        sort_keys=True,
    )
    assert leaked not in serialized, serialized
    assert "database_json_invalid" in serialized


def test_database_read_report_never_echoes_the_host_path(preflight_vault, tmp_path):
    """The decoder message embeds the artifact path; findings must not."""
    database = tmp_path / "secret-vault-name" / "tracking-database.json"
    database.parent.mkdir(parents=True)
    database.write_bytes(b"\xff\xfe not utf-8")

    report = preflight_vault.run_preflight(database)
    # `database` and `vault_root` are documented structured fields.
    findings = json.dumps(report["findings"], sort_keys=True)
    assert "secret-vault-name" not in findings, findings
    assert report["findings"][0]["code"] == "database_encoding_invalid"


@pytest.mark.parametrize(
    "payload,expected_message",
    [
        (
            # MAX_JSON_NESTING_DEPTH is 200; go comfortably past it.
            (
                '{"config": {}, "talks": [], "deep": ' + "[" * 260 + "]" * 260 + "}"
            ).encode(),
            "tracking database exceeds the maximum supported JSON nesting depth",
        ),
        (
            b'{"config": {}, "talks": [], "s": "\\ud800"}',
            "tracking database contains an unpaired UTF-16 surrogate in a JSON string",
        ),
    ],
)
def test_tree_validator_defects_route_to_their_specific_reason(
    preflight_vault, tmp_path, payload, expected_message
):
    """These raise after a successful decode; untyped they fall back to generic."""
    database = tmp_path / "tracking-database.json"
    database.write_bytes(payload)

    report = preflight_vault.run_preflight(database)
    finding = report["findings"][0]
    assert finding["code"] == "database_json_invalid", finding
    assert finding["message"] == expected_message


def test_every_emitted_decoder_reason_is_mapped(preflight_vault):
    """A reason with no mapping silently degrades to the generic code."""
    import re

    # `__file__` is Optional because a namespace or builtin module has none;
    # this one is imported from a real path.
    io_module_file = importlib.import_module("tracking_database_io").__file__
    assert io_module_file is not None
    io_source = pathlib.Path(io_module_file).read_text(encoding="utf-8")
    emitted = set(re.findall(r'reason_code="([a-z_]+)"', io_source))
    mapped = set(preflight_vault._DATABASE_READ_DIAGNOSTICS)
    assert emitted <= mapped, sorted(emitted - mapped)


# PPTX catalog visual-evidence selection (#229). A stored receipt is a hint;
# preflight fingerprints the deck on disk before letting one claim currency.


def _write_deck(fixture, name: str = "Talk.pptx", body: bytes = b"deck-bytes") -> Path:
    path = fixture["pptx_source"] / name
    path.write_bytes(body)
    return path


def _fingerprint(path: Path) -> dict[str, Any]:
    data = path.read_bytes()
    return {
        "algorithm": "sha256",
        "digest": hashlib.sha256(data).hexdigest(),
        "size_bytes": len(data),
    }


def _write_artifact(fixture, body: bytes = b'{"slides": []}') -> Path:
    """The extraction artifact a receipt points at, vault-root-relative."""
    path = fixture["root"] / "evidence" / "talk.json"
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(body)
    return path


def _catalog_record(
    preflight_vault, deck: Path, artifact: Path, **overrides: Any
) -> dict[str, Any]:
    """One schema-v2 catalog record whose receipt matches deck and artifact."""
    record = {
        "schema_version": 2,
        "pptx_path": deck.name,
        "talk_filename": None,
        "matched": False,
        "slide_count": 42,
        "visual_extracted": True,
        "visual_evidence": {
            "outcome": "succeeded",
            "extractor_schema_version": preflight_vault.PPTX_EXTRACTION_SCHEMA_VERSION,
            "pipeline_version": preflight_vault.PPTX_EXTRACTION_PIPELINE_VERSION,
            "source_fingerprint": _fingerprint(deck),
            "artifact": {
                "path": "evidence/talk.json",
                "sha256": _fingerprint(artifact)["digest"],
            },
        },
    }
    record.update(overrides)
    return record


def _write_catalog(fixture, records: list[dict[str, Any]]) -> Path:
    database = write_database(fixture, [], current=True)
    payload = json.loads(database.read_text(encoding="utf-8"))
    payload["pptx_catalog"] = records
    database.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return database


def _catalog_findings(report: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        finding
        for finding in report["findings"]
        if finding["code"].startswith("pptx_visual_evidence")
    ]


def test_preflight_is_quiet_when_the_receipt_matches_the_deck_on_disk(
    preflight_vault,
    vault_fixture,
) -> None:
    deck = _write_deck(vault_fixture)
    artifact = _write_artifact(vault_fixture)
    _write_catalog(vault_fixture, [_catalog_record(preflight_vault, deck, artifact)])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    assert _catalog_findings(report) == []


def test_preflight_warns_when_the_deck_changed_after_extraction(
    preflight_vault,
    vault_fixture,
) -> None:
    """The receipt still names the current extractor; the bytes moved on."""
    deck = _write_deck(vault_fixture)
    artifact = _write_artifact(vault_fixture)
    record = _catalog_record(preflight_vault, deck, artifact)
    _write_catalog(vault_fixture, [record])
    deck.write_bytes(b"deck-bytes-edited-since-extraction")

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = _catalog_findings(report)
    assert [finding["code"] for finding in findings] == [
        "pptx_visual_evidence_not_current"
    ]
    assert findings[0]["actual"]["classification"] == "stale"


def test_preflight_warns_when_the_deck_cannot_be_fingerprinted(
    preflight_vault,
    vault_fixture,
) -> None:
    """Stored metadata alone never proves currency."""
    deck = _write_deck(vault_fixture)
    artifact = _write_artifact(vault_fixture)
    record = _catalog_record(preflight_vault, deck, artifact)
    _write_catalog(vault_fixture, [record])
    deck.unlink()

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = _catalog_findings(report)
    assert findings[0]["actual"]["classification"] == "unverified"


def test_preflight_warns_about_a_legacy_record(
    preflight_vault,
    vault_fixture,
) -> None:
    """A v1 bare claim cannot name the extractor generation it came from."""
    deck = _write_deck(vault_fixture)
    _write_catalog(
        vault_fixture,
        [
            {
                "schema_version": 1,
                "pptx_path": deck.name,
                "talk_filename": None,
                "matched": False,
                "slide_count": 42,
                "visual_extracted": True,
            }
        ],
    )

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = _catalog_findings(report)
    assert findings[0]["code"] == "pptx_visual_evidence_not_current"
    assert findings[0]["severity"] == "warning"
    assert findings[0]["actual"]["classification"] == "unknown_legacy"


def test_preflight_reports_an_unreadable_receipt_instead_of_crashing(
    preflight_vault,
    vault_fixture,
) -> None:
    deck = _write_deck(vault_fixture)
    artifact = _write_artifact(vault_fixture)
    broken = _catalog_record(preflight_vault, deck, artifact)
    broken["visual_evidence"]["artifact"] = None
    _write_catalog(vault_fixture, [broken])

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = _catalog_findings(report)
    assert [finding["code"] for finding in findings] == [
        "pptx_visual_evidence_unreadable"
    ]
    # The closed code and its neutral prose, never the rejected persisted
    # value the exception message names (`no-secrets` -> Logging).
    assert findings[0]["message"] == (
        "visual evidence records a success but names no artifact"
    )
    # A bad receipt is per-record evidence trouble, so the vault is still
    # usable: it must not surface as unusable owner state.
    assert report["blocking_count"] == 0


def test_preflight_warns_when_the_extraction_artifact_is_gone(
    preflight_vault,
    vault_fixture,
) -> None:
    """A deleted artifact must not stay authoritative."""
    deck = _write_deck(vault_fixture)
    artifact = _write_artifact(vault_fixture)
    _write_catalog(vault_fixture, [_catalog_record(preflight_vault, deck, artifact)])
    artifact.unlink()

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = _catalog_findings(report)
    assert findings[0]["actual"]["classification"] == "unverified"
    assert findings[0]["actual"]["artifact_observed"] is False
    assert findings[0]["actual"]["source_observed"] is True


def test_preflight_warns_when_the_extraction_artifact_was_replaced(
    preflight_vault,
    vault_fixture,
) -> None:
    deck = _write_deck(vault_fixture)
    artifact = _write_artifact(vault_fixture)
    _write_catalog(vault_fixture, [_catalog_record(preflight_vault, deck, artifact)])
    artifact.write_bytes(b'{"slides": ["replaced"]}')

    report = preflight_vault.run_preflight(vault_fixture["root"])

    findings = _catalog_findings(report)
    assert findings[0]["actual"]["classification"] == "stale"
