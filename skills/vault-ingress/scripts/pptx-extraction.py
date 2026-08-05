#!/usr/bin/env python3
"""Extract visual design data from .pptx files using python-pptx.

Produces per-slide visual data and global design statistics as JSON.
Skips static exports, conflict copies, and configured filename patterns.

Text is collected recursively from shape frames and native table cells, with
source/confidence recorded per channel. Low-confidence picture and background
blobs are OCR'd for a word inventory; groups, tables, SmartArt, charts, and
other unsupported containers explicitly require a rendered pass. OCR remains
inventory only, not a substitute for visual design judgment.

Native slide XML is also inventoried for timing containers, exact animation
behavior elements, visibility set actions, transitions, and audio/video timing.
Those are raw package counts with explicit non-playback provenance, not claims
about observed motion, concurrency, or delivered behavior.

Usage:
    pptx-extraction.py <file.pptx> [--no-ocr]
        [--rendered-pdf <path>] [--inspected-pages <PAGE|START-END>]
    pptx-extraction.py --directory <directory> [--skip pattern ...]
        [--exclude-directory component ...] [--no-ocr]
    pptx-extraction.py --version

    <path>       Path to one .pptx, or a root when --directory is explicit
    --skip       One configured skip pattern; repeat for each pattern and omit for none
    --no-ocr     Skip OCR even on low-confidence slides (shape walk only)

Examples:
    pptx-extraction.py /path/to/talk.pptx
    pptx-extraction.py --directory /path/to/Presentations --skip template --skip draft
"""

import argparse
import hashlib
import io
import json
import math
import os
import re
import stat
import sys
import time
from collections import Counter
from dataclasses import replace
from math import gcd
from pathlib import Path
from typing import cast
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

import artifact_metadata
from artifact_locator import ArtifactLocatorError, materialize_native_root
from artifact_supervisor import (
    JsonValue,
    SupervisorError,
    SupervisorLimits,
    WorkerRequest,
    isolate_protocol_output,
    read_worker_request,
    run_authenticated_worker,
    write_worker_response,
)
from pptx_evidence import (
    PPTX_EXTRACTION_PIPELINE_VERSION,
    PPTX_EXTRACTION_SCHEMA_VERSION,
    PPTX_OCR_TRUST_CONFIDENCE,
    PPTX_TEXT_BEARING_IMAGE_AREA_RATIO,
    PptxEvidenceError,
    build_native_deck_audit,
    _build_rendered_page_inspection_in_process,
    finite_confidence,
    parse_page_range_arguments,
    presentation_with_media_recovery,
    run_supervised_pptx_extraction,
    sha256_bytes,
    snapshot_regular_file,
)
from pptx_discovery_contract import (
    PPTX_DIRECTORY_BATCH_KIND,
    PPTX_DIRECTORY_BATCH_SCHEMA_VERSION,
    PPTX_DIRECTORY_INCOMPLETE_REASON_CODES,
    PPTX_DIRECTORY_MANIFEST_KIND,
    PPTX_DIRECTORY_MANIFEST_SCHEMA_VERSION,
    PptxDiscoveryContractError,
    build_pptx_directory_batch,
    directory_component_is_excluded,
    directory_incomplete_reason_codes,
    validate_pptx_directory_exclusions,
)

# Field-shape and behavior versions are deliberately separate. A missing
# schema_version/pipeline_version identifies the legacy extractor output.
# v2 added fixed-shape native timing. V3 added a raw build-list lane, a closed
# native-deck audit, and an exact rendered-page inspection receipt. V4 makes
# shape, picture, and background capability/asset bindings non-optional.
SCHEMA_VERSION = PPTX_EXTRACTION_SCHEMA_VERSION
PIPELINE_VERSION = PPTX_EXTRACTION_PIPELINE_VERSION

# PresentationML timing elements are counted by exact qualified name and kept
# in separate semantic lanes. In particular, a <p:timing> container or a media
# node is not a generic animation/motion observation.
_ANIMATION_BEHAVIOR_ELEMENTS = {
    "general": "p:anim",
    "color": "p:animClr",
    "effect": "p:animEffect",
    "motion": "p:animMotion",
    "rotation": "p:animRot",
    "scale": "p:animScale",
}
_MEDIA_TIMING_ELEMENTS = {
    "audio": "p:audio",
    "video": "p:video",
}
_BUILD_ENTRY_ELEMENTS = {
    "paragraph": "p:bldP",
    "diagram": "p:bldDgm",
    "ole_chart": "p:bldOleChart",
    "graphic": "p:bldGraphic",
}
_TIMING_PROVENANCE = {
    "source": "pptx_package_xml",
    "measurement": "raw_ooxml_element_counts",
    "observed_playback": False,
}

# DrawingML graphic-frame URIs. python-pptx exposes chart/table helpers, but
# returns shape_type=None for SmartArt and other graphic frames, so the URI is
# the only reliable discriminator for those objects.
_GRAPHIC_DATA_URI_TABLE = "http://schemas.openxmlformats.org/drawingml/2006/table"
_GRAPHIC_DATA_URI_CHART = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_GRAPHIC_DATA_URI_OLE = "http://schemas.openxmlformats.org/presentationml/2006/ole"

# Cap so one dense manual page cannot blow out the JSON. Inventory is for
# cites and transcript cross-checks, not a full document dump.
_OCR_TEXT_MAX_CHARS = 8000
_NATIVE_TEXT_MAX_CHARS = 8000
_SHAPE_PATH_COMPONENT_MAX_CHARS = 4096
_PACKAGE_PART_NAME_MAX_CHARS = 2048


def _bounded_package_part_name(value):
    part_name = str(value).lstrip("/")
    if (
        not part_name
        or "\x00" in part_name
        or len(part_name) > _PACKAGE_PART_NAME_MAX_CHARS
    ):
        raise PptxEvidenceError(
            "PPTX package part name exceeds the bounded evidence contract",
            reason_code="pptx_probe_resource_unavailable",
        )
    return part_name


# Directory mode is intentionally a fixed, non-user-expandable batch contract.
# It discovers only local, non-symlink regular files and stops before launching
# more workers once any aggregate budget is exhausted.
_BATCH_MAX_DEPTH = 32
_BATCH_MAX_DIRECTORIES = 5_000
_BATCH_MAX_ENTRIES = 50_000
_BATCH_MAX_POLICY_EXCLUDED_ENTRIES = 5_000
_BATCH_MAX_FILES = 256
_BATCH_MAX_RELATIVE_PATH_CHARS = 4_096
_BATCH_MAX_INPUT_BYTES = 16 * 1024 * 1024 * 1024
_BATCH_MAX_OUTPUT_BYTES = 512 * 1024 * 1024
_BATCH_MAX_WALL_SECONDS = 3_600
_WINDOWS_REPARSE_POINT_ATTRIBUTE = artifact_metadata.WINDOWS_REPARSE_POINT_ATTRIBUTE
_WINDOWS_OFFLINE_ATTRIBUTE = artifact_metadata.WINDOWS_OFFLINE_ATTRIBUTE
_WINDOWS_RECALL_ON_OPEN_ATTRIBUTE = artifact_metadata.WINDOWS_RECALL_ON_OPEN_ATTRIBUTE
_WINDOWS_RECALL_ON_DATA_ACCESS_ATTRIBUTE = (
    artifact_metadata.WINDOWS_RECALL_ON_DATA_ACCESS_ATTRIBUTE
)
_WINDOWS_UNAVAILABLE_CLOUD_ATTRIBUTES = (
    artifact_metadata.WINDOWS_UNAVAILABLE_CLOUD_ATTRIBUTES
)
_WINDOWS_CLOUD_REPARSE_TAGS = artifact_metadata.WINDOWS_CLOUD_REPARSE_TAGS
_BATCH_MAX_ROOT_PATH_CHARS = 4_096
_BATCH_MAX_SKIP_PATTERNS = 64
_BATCH_MAX_SKIP_PATTERN_CHARS = 256
_DIRECTORY_MANIFEST_SCHEMA_VERSION = PPTX_DIRECTORY_MANIFEST_SCHEMA_VERSION
_DIRECTORY_WORKER_FLAG = "--directory-worker"
_DIRECTORY_OPERATION = "pptx_resolve_input"
_DIRECTORY_MANIFEST_SKIP_REASONS = frozenset(
    {
        "pptx_batch_cloud_placeholder_unavailable",
        "pptx_batch_conflict_copy",
        "pptx_batch_depth_limit",
        "pptx_batch_directory_changed",
        "pptx_batch_directory_identity_collision",
        "pptx_batch_directory_identity_unavailable",
        "pptx_batch_directory_excluded",
        "pptx_batch_directory_limit",
        "pptx_batch_directory_unavailable",
        "pptx_batch_entry_limit",
        "pptx_batch_entry_unavailable",
        "pptx_batch_file_limit",
        "pptx_batch_office_lock_file",
        "pptx_batch_path_invalid",
        "pptx_batch_path_limit",
        "pptx_batch_reparse_point_rejected",
        "pptx_batch_scan_incomplete_file_limit",
        "pptx_batch_skip_pattern",
        "pptx_batch_static_export",
        "pptx_batch_symlink_rejected",
        "pptx_batch_wall_limit",
    }
)
_BATCH_EXTRACTION_FAILURE_REASONS = frozenset(
    {
        "pptx_archive_recovery_required",
        "pptx_artifact_changed",
        "pptx_artifact_unavailable",
        "pptx_batch_input_limit",
        "pptx_batch_wall_limit",
        "pptx_cloud_placeholder_unavailable",
        "pptx_dependency_unavailable",
        "pptx_evidence_invalid",
        "pptx_invalid_container",
        "pptx_no_slides",
        "pptx_parse_failure",
        "pptx_probe_crash",
        "pptx_probe_exception",
        "pptx_probe_malformed_result",
        "pptx_probe_materialization_changed",
        "pptx_probe_containment_unavailable",
        "pptx_probe_monitor_identity_changed",
        "pptx_probe_monitor_unavailable",
        "pptx_probe_request_oversized",
        "pptx_probe_resource_unavailable",
        "pptx_probe_result_oversized",
        "pptx_probe_start_failure",
        "pptx_probe_timeout",
        "pptx_recovery_failure",
        "pptx_structural_damage",
    }
)
_DIRECTORY_RESOURCE_FAILURES = frozenset(
    {
        "worker_containment_unavailable",
        "worker_diagnostic_limit_exceeded",
        "worker_input_limit_exceeded",
        "worker_memory_limit_exceeded",
        "worker_monitor_identity_changed",
        "worker_monitor_unavailable",
        "worker_process_limit_exceeded",
    }
)
_DIRECTORY_START_FAILURES = frozenset(
    {
        "invalid_worker_command",
        "unsafe_worker_process_metadata",
        "worker_pipe_setup_failed",
        "worker_request_write_failed",
        "worker_start_failed",
    }
)
_DIRECTORY_WORKER_FAILURES = frozenset(
    {
        "worker_cleanup_failed",
        "worker_diagnostic_read_failed",
        "worker_exit",
        "worker_exit_before_barrier",
        "worker_output_read_failed",
        "worker_process_tree_leak",
    }
)
_DIRECTORY_CHILD_FAILURES = frozenset(
    {
        "pptx_batch_request_invalid",
        "pptx_batch_root_invalid",
        "pptx_batch_root_unavailable",
    }
)
_DIRECTORY_BATCH_FAILURE_REASONS = frozenset(
    {
        *_DIRECTORY_CHILD_FAILURES,
        "pptx_batch_discovery_output_limit",
        "pptx_batch_discovery_protocol_invalid",
        "pptx_batch_discovery_resource_unavailable",
        "pptx_batch_discovery_start_failure",
        "pptx_batch_discovery_timeout",
        "pptx_batch_discovery_worker_failure",
        "pptx_batch_manifest_invalid",
        "pptx_batch_wall_limit",
    }
)
_DIRECTORY_LIMITS = SupervisorLimits(
    profile_id="pptx-directory-discovery-v2",
    wall_seconds=60,
    max_memory_bytes=512 * 1024 * 1024,
    max_input_bytes=64 * 1024,
    max_output_bytes=16 * 1024 * 1024,
    max_processes=1,
)

# Mean Tesseract token confidence uses the engine's documented 0..100 scale.
# Text below this floor remains visible for spelling review but is not marked
# trustworthy evidence.
_OCR_TRUST_CONFIDENCE = PPTX_OCR_TRUST_CONFIDENCE

# One stderr warning per process when the OCR engine is missing — not once
# per slide on a 100-slide deck.
_ocr_unavailable_warned = False

# Cached tesseract availability for this process: None = not checked yet.
_tesseract_available = None

# Populated by the same availability probe, then repeated on every asset
# receipt so an OCR result is reproducible and distinguishable from a skip.
_tesseract_version = None


class OcrUnavailableError(Exception):
    """Tesseract (or its Python binding) is not available on this machine."""


def _count_in_containers(containers, qualified_name):
    """Count exact descendants across selected PresentationML containers."""
    tag = qn(qualified_name)
    return sum(1 for container in containers for _element in container.iter(tag))


def _is_visibility_set_action(set_element):
    """Return whether a <p:set> explicitly targets a visibility attribute."""
    for attribute_name in set_element.iter(qn("p:attrName")):
        value = (attribute_name.text or "").strip().lower()
        if value == "visibility" or value.endswith(".visibility"):
            return True
    return False


def extract_native_timing(slide):
    """Inventory raw timing/transition XML without claiming observed motion.

    Counts are structural package evidence only. Markup-compatibility Choice and
    Fallback branches are both retained and therefore both counted; resolving
    which branch a particular presenter executes requires a playback engine.
    """
    root = slide.element
    timing_elements = list(root.iter(qn("p:timing")))
    build_lists = list(root.iter(qn("p:bldLst")))
    set_actions = [
        element for timing in timing_elements for element in timing.iter(qn("p:set"))
    ]
    animation_counts = {
        name: _count_in_containers(timing_elements, qualified_name)
        for name, qualified_name in _ANIMATION_BEHAVIOR_ELEMENTS.items()
    }
    animation_counts["total"] = sum(animation_counts.values())
    media_counts = {
        name: _count_in_containers(timing_elements, qualified_name)
        for name, qualified_name in _MEDIA_TIMING_ELEMENTS.items()
    }
    media_counts["total"] = sum(media_counts.values())
    build_counts = {
        name: _count_in_containers(build_lists, qualified_name)
        for name, qualified_name in _BUILD_ENTRY_ELEMENTS.items()
    }
    build_counts["total"] = sum(build_counts.values())
    part_name = _bounded_package_part_name(slide.part.partname)
    return {
        "timing_element_present": bool(timing_elements),
        "timing_element_count": len(timing_elements),
        "transition_count": sum(1 for _ in root.iter(qn("p:transition"))),
        "set_action_count": len(set_actions),
        "visibility_set_action_count": sum(
            1 for element in set_actions if _is_visibility_set_action(element)
        ),
        "animation_behavior_counts": animation_counts,
        "media_timing_counts": media_counts,
        "build_list_present": bool(build_lists),
        "build_list_count": len(build_lists),
        "build_entry_counts": build_counts,
        "has_animation_behaviors": animation_counts["total"] > 0,
        "has_media_timing": media_counts["total"] > 0,
        "has_build_entries": build_counts["total"] > 0,
        "provenance": {
            **_TIMING_PROVENANCE,
            "part_name": part_name,
        },
    }


def summarize_native_timing(per_slide_visual):
    """Aggregate fixed-key deck totals from per-slide structural metadata."""
    animation_counts = {name: 0 for name in _ANIMATION_BEHAVIOR_ELEMENTS}
    media_counts = {name: 0 for name in _MEDIA_TIMING_ELEMENTS}
    build_counts = {name: 0 for name in _BUILD_ENTRY_ELEMENTS}
    summary = {
        "slides_with_timing_elements": 0,
        "slides_with_transitions": 0,
        "slides_with_animation_behaviors": 0,
        "slides_with_media_timing": 0,
        "slides_with_build_lists": 0,
        "slides_with_build_entries": 0,
        "timing_element_count": 0,
        "transition_count": 0,
        "set_action_count": 0,
        "visibility_set_action_count": 0,
        "build_list_count": 0,
    }
    for slide_data in per_slide_visual:
        timing = slide_data["native_timing"]
        summary["slides_with_timing_elements"] += int(timing["timing_element_present"])
        summary["slides_with_transitions"] += int(timing["transition_count"] > 0)
        summary["slides_with_animation_behaviors"] += int(
            timing["has_animation_behaviors"]
        )
        summary["slides_with_media_timing"] += int(timing["has_media_timing"])
        summary["slides_with_build_lists"] += int(timing["build_list_present"])
        summary["slides_with_build_entries"] += int(timing["has_build_entries"])
        for field in (
            "timing_element_count",
            "transition_count",
            "set_action_count",
            "visibility_set_action_count",
            "build_list_count",
        ):
            summary[field] += timing[field]
        for name in animation_counts:
            animation_counts[name] += timing["animation_behavior_counts"][name]
        for name in media_counts:
            media_counts[name] += timing["media_timing_counts"][name]
        for name in build_counts:
            build_counts[name] += timing["build_entry_counts"][name]
    animation_counts["total"] = sum(animation_counts.values())
    media_counts["total"] = sum(media_counts.values())
    build_counts["total"] = sum(build_counts.values())
    return {
        **summary,
        "animation_behavior_counts": animation_counts,
        "media_timing_counts": media_counts,
        "build_entry_counts": build_counts,
        "provenance": dict(_TIMING_PROVENANCE),
    }


def _input_fingerprint(blob):
    """Return a deterministic fingerprint for the exact source bytes."""
    return {
        "algorithm": "sha256",
        "digest": hashlib.sha256(blob).hexdigest(),
        "size_bytes": len(blob),
    }


def picture_area_ratio(shape, prs):
    """Return a picture shape's area as a fraction of the slide (0.0–1.0).

    Missing geometry (any of width/height/slide dimensions absent or zero)
    returns 0.0 — unknown size is not evidence of a large picture.
    """
    slide_area = (prs.slide_width or 0) * (prs.slide_height or 0)
    if not slide_area:
        return 0.0
    if not shape.width or not shape.height:
        return 0.0
    return min((shape.width * shape.height) / slide_area, 1.0)


def normalize_ocr_text(text):
    """Collapse whitespace and strip; empty input stays empty."""
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


def _ocr_text_and_confidence(data):
    """Pair each retained Tesseract token with only its own confidence."""
    tokens = []
    confidences = []
    for token, raw_confidence in zip(data.get("text", []), data.get("conf", [])):
        normalized = normalize_ocr_text(token)
        if not normalized:
            continue
        tokens.append(normalized)
        try:
            confidence = float(raw_confidence)
        except (TypeError, ValueError):
            continue
        normalized_confidence = finite_confidence(confidence)
        if normalized_confidence is not None:
            confidences.append(normalized_confidence)
    text = normalize_ocr_text(" ".join(tokens))[:_OCR_TEXT_MAX_CHARS]
    confidence = round(sum(confidences) / len(confidences), 3) if confidences else None
    return text, confidence


def _require_tesseract():
    """Ensure tesseract + bindings are available; cache the result per process.

    Raises OcrUnavailableError when missing. Subsequent calls in the same
    process do not re-spawn the version check.
    """
    global _tesseract_available, _tesseract_version
    if _tesseract_available is True:
        return _tesseract_version
    if _tesseract_available is False:
        raise OcrUnavailableError(
            "tesseract binary not found; install tesseract-ocr (apt) or "
            "tesseract (brew)"
        )

    try:
        import pytesseract  # pyright: ignore[reportMissingImports] - optional OCR dependency checked at runtime
    except ImportError as e:
        _tesseract_available = False
        raise OcrUnavailableError(
            "OCR requires Pillow and pytesseract; install project dependencies"
        ) from e

    try:
        version = pytesseract.get_tesseract_version()
    except pytesseract.TesseractNotFoundError as e:
        _tesseract_available = False
        raise OcrUnavailableError(
            "tesseract binary not found; install tesseract-ocr (apt) or "
            "tesseract (brew)"
        ) from e

    _tesseract_available = True
    _tesseract_version = str(version)[:128]
    return _tesseract_version


def _ocr_image_result(blob):
    """OCR one image and return an outcome distinct from asset provenance."""
    global _tesseract_available
    engine_version = _require_tesseract()

    try:
        import pytesseract  # pyright: ignore[reportMissingImports] - optional OCR dependency checked at runtime
        from PIL import Image, UnidentifiedImageError
    except ImportError as e:
        raise OcrUnavailableError(
            "OCR requires Pillow and pytesseract; install project dependencies"
        ) from e
    pillow_errors = (
        UnidentifiedImageError,
        Image.DecompressionBombError,
        OSError,
    )

    try:
        img = Image.open(io.BytesIO(blob))
    except pillow_errors as e:
        error_code = f"image_decode_error:{type(e).__name__}"[:512]
        sys.stderr.write(f"WARN: OCR skipped image blob ({error_code})\n")
        return {
            "attempted": True,
            "engine": "tesseract",
            "engine_version": engine_version,
            "result_status": "failed",
            "result_confidence": None,
            "error": error_code,
            "recovered_text": "",
            "trustworthy_text": False,
        }

    try:
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        data = pytesseract.image_to_data(
            img,
            output_type=pytesseract.Output.DICT,
        )
    except pytesseract.TesseractNotFoundError as e:
        _tesseract_available = False
        raise OcrUnavailableError(
            "tesseract binary not found; install tesseract-ocr (apt) or "
            "tesseract (brew)"
        ) from e
    except pillow_errors as e:
        error_code = f"image_decode_error:{type(e).__name__}"[:512]
        sys.stderr.write(f"WARN: OCR failed on image blob ({error_code})\n")
        return {
            "attempted": True,
            "engine": "tesseract",
            "engine_version": engine_version,
            "result_status": "failed",
            "result_confidence": None,
            "error": error_code,
            "recovered_text": "",
            "trustworthy_text": False,
        }
    except pytesseract.TesseractError as e:
        error_code = f"engine_error:{type(e).__name__}"[:512]
        sys.stderr.write(f"WARN: OCR failed on image blob ({error_code})\n")
        return {
            "attempted": True,
            "engine": "tesseract",
            "engine_version": engine_version,
            "result_status": "failed",
            "result_confidence": None,
            "error": error_code,
            "recovered_text": "",
            "trustworthy_text": False,
        }

    text, confidence = _ocr_text_and_confidence(data)
    trustworthy = bool(
        text and confidence is not None and confidence >= _OCR_TRUST_CONFIDENCE
    )
    if not text:
        status = "genuine_empty"
    elif trustworthy:
        status = "text_recovered"
    else:
        status = "low_confidence_text"
    return {
        "attempted": True,
        "engine": "tesseract",
        "engine_version": engine_version,
        "result_status": status,
        "result_confidence": confidence,
        "error": None,
        "recovered_text": text,
        "trustworthy_text": trustworthy,
    }


def ocr_image_bytes(blob):
    """OCR one image and return recovered text for legacy callers."""
    return _ocr_image_result(blob)["recovered_text"]


def ocr_picture_blobs(blobs):
    """OCR picture blobs; join non-empty results with ' | '.

    Raises OcrUnavailableError if the engine is missing. Returns "" when every
    blob is empty or unreadable. Caller is responsible for sort order.
    """
    parts = []
    for blob in blobs:
        text = ocr_image_bytes(blob)
        if text:
            parts.append(text)
    joined = " | ".join(parts)
    if len(joined) > _OCR_TEXT_MAX_CHARS:
        return joined[:_OCR_TEXT_MAX_CHARS]
    return joined


def _append_ocr_text(slide_data, text):
    """Append one OCR channel's text to the backward-compatible aggregate."""
    if not text:
        return
    current = slide_data.get("ocr_text", "")
    combined = f"{current} | {text}" if current else text
    slide_data["ocr_text"] = combined[:_OCR_TEXT_MAX_CHARS]


def _run_ocr_channel(
    slide_data,
    assets,
    *,
    channel,
    provenance,
    ocr=True,
    ocr_fn=None,
):
    """OCR image assets and emit one outcome receipt for every exact blob."""
    record = {
        "channel": channel,
        "text": "",
        "confidence": "low",
        "result_confidence": None,
        "status": "unavailable" if not assets else "pending",
        "attempted": False,
        "engine": "tesseract",
        "engine_version": None,
        "reason": "no_readable_asset" if not assets else None,
        "provenance": provenance,
        "ocr_receipts": [],
    }
    slide_data["text_channels"].append(record)

    if not assets:
        return
    if not ocr:
        record["status"] = "skipped"
        record["reason"] = "ocr_disabled"
        for asset in assets:
            record["ocr_receipts"].append(
                {
                    "attempted": False,
                    "engine": "tesseract",
                    "engine_version": None,
                    "result_status": "skipped",
                    "result_confidence": None,
                    "error": "ocr_disabled",
                    "part_name": asset.get("part_name"),
                    "asset_sha256": sha256_bytes(asset["blob"]),
                    "shape_path": list(asset.get("shape_path") or []),
                    "recovered_text": "",
                    "trustworthy_text": False,
                }
            )
        return

    global _ocr_unavailable_warned
    for asset in assets:
        blob = asset["blob"]
        try:
            if ocr_fn is None:
                outcome = _ocr_image_result(blob)
            else:
                injected = ocr_fn([blob])
                if isinstance(injected, dict):
                    recovered_text = normalize_ocr_text(
                        injected.get("recovered_text", "")
                    )[:_OCR_TEXT_MAX_CHARS]
                    injected_engine = injected.get("engine", "injected")
                    injected_version = injected.get("engine_version")
                    injected_error = injected.get("error")
                    outcome = {
                        "attempted": injected.get("attempted", True) is True,
                        "engine": (
                            injected_engine[:128]
                            if isinstance(injected_engine, str)
                            else "injected"
                        ),
                        "engine_version": (
                            injected_version[:128]
                            if isinstance(injected_version, str)
                            else None
                        ),
                        "result_status": injected.get("result_status")
                        or ("text_recovered" if recovered_text else "genuine_empty"),
                        "result_confidence": finite_confidence(
                            injected.get("result_confidence")
                        ),
                        "error": (
                            injected_error[:512]
                            if isinstance(injected_error, str)
                            else None
                        ),
                        "recovered_text": recovered_text,
                        "trustworthy_text": (
                            injected.get("trustworthy_text", bool(recovered_text))
                            is True
                        ),
                    }
                else:
                    text = normalize_ocr_text(injected)[:_OCR_TEXT_MAX_CHARS]
                    outcome = {
                        "attempted": True,
                        "engine": "injected",
                        "engine_version": None,
                        "result_status": (
                            "text_recovered" if text else "genuine_empty"
                        ),
                        "result_confidence": None,
                        "error": None,
                        "recovered_text": text,
                        "trustworthy_text": bool(text),
                    }
        except OcrUnavailableError as e:
            outcome = {
                "attempted": False,
                "engine": "tesseract",
                "engine_version": None,
                "result_status": "unavailable",
                "result_confidence": None,
                "error": str(e)[:512],
                "recovered_text": "",
                "trustworthy_text": False,
            }
            if not _ocr_unavailable_warned:
                sys.stderr.write(
                    f"WARN: OCR unavailable ({e}); low-confidence slides will "
                    "retain per-asset unavailable receipts. Install tesseract "
                    "to enable OCR.\n"
                )
                _ocr_unavailable_warned = True
        record["ocr_receipts"].append(
            {
                **outcome,
                "part_name": asset.get("part_name"),
                "asset_sha256": sha256_bytes(blob),
                "shape_path": list(asset.get("shape_path") or []),
            }
        )

    recovered = [
        receipt["recovered_text"]
        for receipt in record["ocr_receipts"]
        if receipt["recovered_text"]
    ]
    text = normalize_ocr_text(" | ".join(recovered))[:_OCR_TEXT_MAX_CHARS]
    record["text"] = text
    confidences = [
        receipt["result_confidence"]
        for receipt in record["ocr_receipts"]
        if receipt["result_confidence"] is not None
    ]
    if confidences:
        record["result_confidence"] = round(sum(confidences) / len(confidences), 3)
    statuses = {receipt["result_status"] for receipt in record["ocr_receipts"]}
    record["attempted"] = any(
        receipt["attempted"] for receipt in record["ocr_receipts"]
    )
    engines = {
        receipt["engine"]
        for receipt in record["ocr_receipts"]
        if isinstance(receipt.get("engine"), str) and receipt["engine"]
    }
    record["engine"] = next(iter(engines)) if len(engines) == 1 else None
    versions = {
        receipt["engine_version"]
        for receipt in record["ocr_receipts"]
        if isinstance(receipt.get("engine_version"), str) and receipt["engine_version"]
    }
    record["engine_version"] = next(iter(versions)) if len(versions) == 1 else None
    all_recovered_text_trustworthy = all(
        receipt["trustworthy_text"] is True
        for receipt in record["ocr_receipts"]
        if receipt["recovered_text"]
    )
    if text and statuses == {"text_recovered"} and all_recovered_text_trustworthy:
        record["status"] = "extracted"
    elif text:
        record["status"] = "partial"
    elif statuses == {"genuine_empty"}:
        record["status"] = "empty"
    elif "failed" in statuses:
        record["status"] = "failed"
    elif statuses == {"unavailable"}:
        record["status"] = "unavailable"
        record["reason"] = "ocr_engine_unavailable"
    else:
        record["status"] = "partial"
    if record["status"] == "failed":
        record["reason"] = "ocr_failed"
    elif record["status"] == "partial":
        record["reason"] = "partial_ocr_results"
    _append_ocr_text(slide_data, text)
    if any(receipt["attempted"] for receipt in record["ocr_receipts"]):
        slide_data["text_extraction_method"] = "shapes+ocr"
    elif statuses == {"unavailable"}:
        slide_data["text_extraction_method"] = "shapes+ocr_unavailable"


def apply_ocr_to_slide(
    slide_data,
    picture_blobs,
    *,
    ocr=True,
    ocr_fn=None,
    shape_paths=None,
    part_names=None,
):
    """Fill picture OCR fields on a per-slide dict.

    OCR runs only when confidence is low and at least one picture blob is
    available. Image-background slides with no PICTURE shapes have no blob to
    OCR here (rendering the page is out of this script's scope).

    ocr_fn is injectable for tests (signature: list[bytes] -> str|dict). The
    production path runs Tesseract independently for each asset.
    """
    if slide_data.get("text_extraction_confidence") != "low":
        return slide_data
    if not picture_blobs:
        if slide_data.get("has_image") is True:
            _run_ocr_channel(
                slide_data,
                [],
                channel="picture_ocr",
                provenance={
                    "source": "embedded_picture_blobs",
                    "shape_paths": shape_paths or [],
                },
                ocr=ocr,
                ocr_fn=ocr_fn,
            )
        return slide_data

    paths = shape_paths or [[] for _blob in picture_blobs]
    names = part_names or [None for _blob in picture_blobs]
    assets = [
        {"blob": blob, "shape_path": path, "part_name": part_name}
        for blob, path, part_name in zip(picture_blobs, paths, names)
    ]
    _run_ocr_channel(
        slide_data,
        assets,
        channel="picture_ocr",
        provenance={
            "source": "embedded_picture_blobs",
            "shape_paths": shape_paths or [],
        },
        ocr=ocr,
        ocr_fn=ocr_fn,
    )
    return slide_data


def apply_background_ocr(slide_data, background_image, *, ocr=True, ocr_fn=None):
    """OCR an actual background fill blob, or record why rendering is needed."""
    if background_image is None:
        background_image = {
            "blob": None,
            "status": "unavailable",
            "provenance": {"source": "pptx_background_image"},
        }
    blob = background_image.get("blob")
    provenance = background_image["provenance"]
    _run_ocr_channel(
        slide_data,
        (
            [
                {
                    "blob": blob,
                    "shape_path": [],
                    "part_name": provenance.get("part_name"),
                }
            ]
            if blob
            else []
        ),
        channel="background_image_ocr",
        provenance=provenance,
        ocr=ocr,
        ocr_fn=ocr_fn,
    )
    return slide_data


def rgb_to_hex(rgb):
    """Convert RGBColor to hex string."""
    if rgb is None:
        return None
    return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"


def _local_name(element):
    """Return the namespace-free local name for an lxml element."""
    return element.tag.rsplit("}", 1)[-1]


def _background_candidates(slide):
    """Yield slide/layout/master objects in background inheritance order."""
    yield "slide", slide
    layout = slide.slide_layout
    if layout is not None:
        yield "layout", layout
        master = layout.slide_master
        if master is not None:
            yield "master", master


def _effective_background(slide):
    """Return the first explicit background without mutating PPTX XML.

    Accessing ``slide.background.fill`` creates a no-fill background when the
    slide inherits one. Reading the raw cSld tree avoids that destructive
    behavior and lets us retain the relationship owner for image fills.
    """
    for source, owner in _background_candidates(slide):
        bg = owner.element.cSld.bg
        if bg is None:
            continue
        bg_pr = bg.bgPr
        if bg_pr is None:
            # A theme background reference explicitly wins over inheritance,
            # but it does not expose a concrete color/image blob here.
            return {
                "source": source,
                "owner": owner,
                "fill": None,
                "type": "unknown",
                "color": None,
            }

        fill = bg_pr.eg_fillProperties
        if fill is None:
            fill_type = "unknown"
        else:
            fill_type = {
                "solidFill": "solid",
                "pattFill": "pattern",
                "gradFill": "gradient",
                "blipFill": "image",
                "noFill": "unknown",
                "grpFill": "unknown",
            }.get(_local_name(fill), "unknown")

        color = None
        if fill is not None:
            srgb = next(iter(fill.iter(qn("a:srgbClr"))), None)
            if srgb is not None and srgb.get("val"):
                color = f"#{srgb.get('val').upper()}"
            else:
                system = next(iter(fill.iter(qn("a:sysClr"))), None)
                if system is not None and system.get("lastClr"):
                    color = f"#{system.get('lastClr').upper()}"

        reported_type = fill_type
        if fill_type == "solid" and source != "slide":
            reported_type = f"solid_from_{source}"
        return {
            "source": source,
            "owner": owner,
            "fill": fill,
            "type": reported_type,
            "color": color,
        }

    return {
        "source": "none",
        "owner": None,
        "fill": None,
        "type": "unknown",
        "color": None,
    }


def get_background_color(slide):
    """Extract background color/type without altering inherited backgrounds."""
    background = _effective_background(slide)
    return background["color"], background["type"]


def get_background_image(slide):
    """Return an explicit background-image blob and its provenance.

    The result is ``None`` for non-image backgrounds. For an image background
    whose relationship/blob is missing, the returned record has ``blob=None``
    and ``status=unavailable`` so callers can require a rendered-page pass.
    """
    background = _effective_background(slide)
    background_type = background["type"]
    if background_type != "image":
        return None

    fill = background["fill"]
    owner = background["owner"]
    source = background["source"]
    provenance = {
        "source": "pptx_background_image",
        "background_owner": source,
    }
    if fill is None or owner is None:
        return {"blob": None, "status": "unavailable", "provenance": provenance}

    blip = next(iter(fill.iter(qn("a:blip"))), None)
    r_id = blip.get(qn("r:embed")) if blip is not None else None
    if not r_id:
        return {"blob": None, "status": "unavailable", "provenance": provenance}

    provenance["relationship_id"] = r_id[:_SHAPE_PATH_COMPONENT_MAX_CHARS]
    try:
        image_part = owner.part.related_part(r_id)
        blob = image_part.blob
    except (KeyError, ValueError, AttributeError):
        return {"blob": None, "status": "unavailable", "provenance": provenance}
    if not blob:
        return {"blob": None, "status": "unavailable", "provenance": provenance}
    provenance["part_name"] = _bounded_package_part_name(image_part.partname)
    return {"blob": blob, "status": "available", "provenance": provenance}


def _picture_payload(shape):
    """Return a picture blob plus the package part that supplied it."""
    r_id = shape.element.blip_rId
    if not r_id:
        raise ValueError("picture has no embedded image relationship")
    image_part = shape.part.related_part(r_id)
    blob = image_part.blob
    if not blob:
        raise ValueError("picture has an empty embedded image part")
    return blob, _bounded_package_part_name(image_part.partname)


def _is_picture_shape(shape):
    """Include inserted picture placeholders, but not movie poster frames."""
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
        shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
        and _local_name(shape.element) == "pic"
    )


def extract_shape_info(shape):
    """Extract visual properties from a shape."""
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "has_text_frame": bool(shape.has_text_frame),
        "is_picture": _is_picture_shape(shape),
        "is_graphic_frame": _local_name(shape.element) == "graphicFrame",
        "graphic_frame_type": None,
        "graphic_data_uri": None,
        "left": round(shape.left / 914400, 2) if shape.left else None,
        "top": round(shape.top / 914400, 2) if shape.top else None,
        "width": round(shape.width / 914400, 2) if shape.width else None,
        "height": round(shape.height / 914400, 2) if shape.height else None,
    }

    # Text properties
    if shape.has_text_frame:
        tf = shape.text_frame
        info["text_preview"] = tf.text[:_NATIVE_TEXT_MAX_CHARS] if tf.text else ""
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font:
                    info["font_name"] = run.font.name
                    if info["font_name"]:
                        info["font_name"] = info["font_name"][
                            :_SHAPE_PATH_COMPONENT_MAX_CHARS
                        ]
                    info["font_size"] = run.font.size.pt if run.font.size else None
                    try:
                        info["font_color"] = (
                            rgb_to_hex(run.font.color.rgb) if run.font.color else None
                        )
                    except AttributeError:
                        info["font_color"] = None
                    info["bold"] = run.font.bold
                    info["italic"] = run.font.italic
                    break
            if "font_name" in info:
                break

    # Fill properties
    if hasattr(shape, "fill"):
        try:
            fill = shape.fill
            if fill.type == 1:  # solid
                info["fill_color"] = rgb_to_hex(fill.fore_color.rgb)
        except (AttributeError, TypeError, ValueError, NotImplementedError):
            pass

    # Line/outline properties
    if hasattr(shape, "line"):
        try:
            line = shape.line
            if line.fill.type == 1:
                info["line_color"] = rgb_to_hex(line.color.rgb)
                info["line_width"] = line.width.pt if line.width else None
        except (AttributeError, TypeError, ValueError, NotImplementedError):
            pass

    # Auto-shape type (for speech bubbles, starbursts, etc.)
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            info["auto_shape_type"] = str(shape.auto_shape_type)
        except (AttributeError, ValueError):
            info["auto_shape_type"] = "UNKNOWN"

    return info


def walk_shapes(shapes, parent_path=()):
    """Yield every shape recursively with a stable, human-readable path."""
    for index, shape in enumerate(shapes):
        name = (shape.name or f"shape_{index + 1}")[:_SHAPE_PATH_COMPONENT_MAX_CHARS]
        path = parent_path + (name,)
        yield shape, path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if len(path) >= 64 and len(shape.shapes):
                raise PptxEvidenceError(
                    "PPTX group nesting exceeds the bounded evidence contract",
                    reason_code="pptx_probe_resource_unavailable",
                )
            yield from walk_shapes(shape.shapes, path)


def _graphic_data_uri(shape):
    """Return a graphic-frame URI, including those unknown to python-pptx."""
    element = shape.element
    if _local_name(element) != "graphicFrame":
        return None
    uri = element.graphicData_uri
    return uri[:_SHAPE_PATH_COMPONENT_MAX_CHARS] if uri else None


def _classify_graphic_frame(shape):
    """Classify a DrawingML graphic frame from its URI."""
    if _local_name(shape.element) != "graphicFrame":
        return None
    uri = _graphic_data_uri(shape)
    if uri is None:
        return "graphic_frame"
    if uri == _GRAPHIC_DATA_URI_TABLE:
        return "table"
    if uri == _GRAPHIC_DATA_URI_CHART:
        return "chart"
    if uri == _GRAPHIC_DATA_URI_OLE:
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            return "embedded_ole_object"
        return "linked_ole_object"
    if "diagram" in uri.lower():
        return "smartart"
    return "graphic_frame"


def _extract_table_channel(shape, shape_path):
    """Extract visible native table-cell text, skipping merged spill cells."""
    table = shape.table
    cell_text = []
    fonts = Counter()
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            if cell.is_spanned:
                continue
            text = cell.text.strip()
            if text:
                cell_text.append(
                    {
                        "cell": f"R{row_index + 1}C{column_index + 1}",
                        "text": text,
                    }
                )
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        fonts[run.font.name[:_SHAPE_PATH_COMPONENT_MAX_CHARS]] += 1

    combined = " | ".join(item["text"] for item in cell_text)[:_NATIVE_TEXT_MAX_CHARS]
    channel = {
        "channel": "table_cell_text",
        "text": combined,
        "confidence": "medium",
        "status": "extracted" if combined else "empty",
        "provenance": {
            "source": "pptx_table_cells",
            "shape_path": list(shape_path),
            "cells": [item["cell"] for item in cell_text],
        },
    }
    details = {
        "table_rows": len(table.rows),
        "table_columns": len(table.columns),
        "table_text_preview": combined,
        "table_fonts": dict(fonts),
    }
    return channel, details


def _shape_text_channel(shape, shape_path, *, in_group):
    """Return a provenance-bearing channel for one shape text frame."""
    text = shape.text_frame.text[:_NATIVE_TEXT_MAX_CHARS]
    return {
        "channel": "shape_text",
        "text": text,
        "confidence": "medium" if in_group else "high",
        "status": "extracted" if text.strip() else "empty",
        "provenance": {
            "source": "pptx_shape_text_frame",
            "shape_path": list(shape_path),
        },
    }


def _mark_render_required(slide_data, reason):
    """Conservatively lower confidence and record a stable render reason."""
    slide_data["text_extraction_confidence"] = "low"
    slide_data["render_required"] = True
    if reason not in slide_data["render_required_reasons"]:
        slide_data["render_required_reasons"].append(reason)


def _record_unsupported(slide_data, kind, shape, shape_path, *, uri=None):
    """Record content that requires rendering or a specialized extractor."""
    provenance = {
        "source": "pptx_unsupported_visual_container",
        "shape_path": list(shape_path),
    }
    if uri:
        provenance["graphic_data_uri"] = uri
    slide_data["text_channels"].append(
        {
            "channel": f"{kind}_text",
            "text": "",
            "confidence": "low",
            "status": "unsupported",
            "provenance": provenance,
        }
    )
    entry = {
        "content_type": kind,
        "shape_name": shape_path[-1],
        "shape_path": list(shape_path),
        "reason": "visible text or labels may not be represented in PPTX text frames",
        "render_required": True,
    }
    if uri:
        entry["graphic_data_uri"] = uri
    slide_data["unsupported_content"].append(entry)
    _mark_render_required(slide_data, kind)


def _has_speaker_notes(slide):
    """Return whether a slide has non-empty speaker notes."""
    if not slide.has_notes_slide:
        return False
    notes_frame = slide.notes_slide.notes_text_frame
    return bool(notes_frame is not None and notes_frame.text.strip())


def extract_template_layouts(prs):
    """Enumerate slide layouts defined by the presentation's masters.

    Returns a list of {index, master_index, name, placeholders: [{idx, type}]}
    entries. `master_index` distinguishes layouts that share a name across
    different slide masters (PowerPoint allows reuse of layout names like
    "Title and Content" across masters) — the vault-profile aggregator keys
    on the (master_index, name) pair when preserving curated `use_for`
    values across regenerations.

    The `use_for` field documented in the speaker-profile schema is
    intentionally curated by the speaker and is not emitted here.
    """
    layouts = []
    index = 0
    for master_index, master in enumerate(prs.slide_masters):
        for layout in master.slide_layouts:
            placeholders = []
            for ph in layout.placeholders:
                try:
                    pf = ph.placeholder_format
                    pt = pf.type
                    type_name = getattr(pt, "name", None) or str(pt).split(" ", 1)[0]
                    placeholders.append(
                        {
                            "idx": pf.idx,
                            "type": type_name,
                        }
                    )
                except AttributeError as e:
                    # Malformed placeholder — record skip with context, continue.
                    sys.stderr.write(
                        f"WARN: skipping placeholder in layout "
                        f"master={master_index} '{layout.name}': {e}\n"
                    )
            layouts.append(
                {
                    "index": index,
                    "master_index": master_index,
                    "name": layout.name[:_SHAPE_PATH_COMPONENT_MAX_CHARS],
                    "placeholders": placeholders,
                }
            )
            index += 1
    return layouts


def _extract_pptx_in_process(
    pptx_path,
    *,
    ocr=True,
    ocr_fn=None,
    rendered_pdf_path=None,
    inspected_page_ranges=None,
    rendered_pdf_generation=None,
):
    """Extract one deck inside an already-contained worker process.

    ocr: when True (default), low-confidence image channels get an OCR review
         inventory in ocr_text/text_channels; affirmative text is gated by each
         receipt's trustworthy_text field.
    ocr_fn: optional callable(list[bytes]) -> str for tests; default uses
            tesseract via ocr_picture_blobs.
    rendered_pdf_path/inspected_page_ranges: optional exact rendered artifact
            plus the page ranges actually inspected for the native-deck audit.
    rendered_pdf_generation: supervisor-owned identity for that rendered PDF.
    """
    package_blob = snapshot_regular_file(pptx_path, label="PPTX artifact")
    source_fingerprint = _input_fingerprint(package_blob)
    prs, archive_recovery = presentation_with_media_recovery(package_blob)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    if slide_width is None or slide_height is None:
        raise ValueError("PPTX presentation has no slide dimensions")
    slide_width_value = int(slide_width)
    slide_height_value = int(slide_height)
    if slide_width_value <= 0 or slide_height_value <= 0:
        raise ValueError("PPTX presentation has invalid slide dimensions")
    ratio_divisor = gcd(slide_width_value, slide_height_value)
    corrupt_part_names = {str(record["part_name"]) for record in archive_recovery}
    result = {
        "schema_version": SCHEMA_VERSION,
        "pipeline_version": PIPELINE_VERSION,
        "input_fingerprint": source_fingerprint,
        "pptx_path": os.fspath(pptx_path),
        "slide_count": len(prs.slides),
        "slide_width_inches": round(slide_width_value / 914400, 2),
        "slide_height_inches": round(slide_height_value / 914400, 2),
        "aspect_ratio": (
            f"{slide_width_value // ratio_divisor}:"
            f"{slide_height_value // ratio_divisor}"
        ),
        "corrupt_assets": [
            {
                "part_name": record["part_name"],
                "error_type": record["error_type"],
                "status": "recovered_with_placeholder",
            }
            for record in archive_recovery
        ],
        "archive_recovery": archive_recovery,
        "template_layouts": extract_template_layouts(prs),
        "per_slide_visual": [],
        "global_design": {
            "fonts_used": Counter(),
            "background_colors": Counter(),
            "shape_types_used": Counter(),
            "color_sequence": [],
        },
    }

    for i, slide in enumerate(prs.slides):
        bg_hex, bg_type = get_background_color(slide)

        slide_data = {
            "slide_number": i + 1,
            "slide_part_name": _bounded_package_part_name(slide.part.partname),
            "background_color_hex": bg_hex,
            "background_type": bg_type,
            "background_asset_status": (
                "unavailable" if bg_type == "image" else "not_applicable"
            ),
            "background_part_name": None,
            "background_asset_sha256": None,
            "layout_name": (
                slide.slide_layout.name[:_SHAPE_PATH_COMPONENT_MAX_CHARS]
                if slide.slide_layout
                else None
            ),
            "shape_count": len(slide.shapes),
            # True when at least one shape carries a text frame. Names what it
            # measures — shapes the extractor can read text out of. It is NOT
            # a claim about whether the slide shows text; text rendered inside
            # a picture is invisible here (see text_extraction_confidence).
            "has_text_frame_shapes": False,
            "has_image": False,
            # Largest picture's area as a fraction of the slide (0.0–1.0).
            "image_area_ratio": 0.0,
            # Low is also used for groups, tables, unsupported graphic frames,
            # and image backgrounds: extracted text in one channel never proves
            # another visual channel is wordless.
            "text_extraction_confidence": "high",
            "text_channels": [],
            "unsupported_content": [],
            "has_unsupported_content": False,
            "render_required": False,
            "render_required_reasons": [],
            "has_speaker_notes": _has_speaker_notes(slide),
            # Raw package structure only: separate behavior/media/transition
            # lanes deliberately make no claim about playback or delivery.
            "native_timing": extract_native_timing(slide),
            # Shape-frame text only. Baked-in picture text lands in ocr_text.
            "text_content_preview": "",
            # OCR inventory when confidence is low and picture blobs exist.
            "ocr_text": "",
            # shapes | shapes+ocr | shapes+ocr_unavailable
            "text_extraction_method": "shapes",
            "footer_text": "",
            "shapes_summary": [],
        }

        text_parts = []
        # Retain the maximum raw geometry until the closed, reported value is
        # derived below. The render decision uses that same reported value so
        # a near-threshold deck cannot disagree with its own evidence record.
        max_image_ratio = 0.0
        # (ratio, blob, path, package-part) entries — sorted largest-first
        # before OCR so the primary full-bleed image is inventoried first.
        picture_entries = []
        recursive_shape_count = 0
        for shape, shape_path in walk_shapes(slide.shapes):
            recursive_shape_count += 1
            shape_info = extract_shape_info(shape)
            shape_info["name"] = shape_path[-1]
            shape_info["shape_path"] = list(shape_path)
            shape_info["group_depth"] = len(shape_path) - 1
            slide_data["shapes_summary"].append(shape_info)

            # Track fonts
            if "font_name" in shape_info and shape_info["font_name"]:
                result["global_design"]["fonts_used"][shape_info["font_name"]] += 1

            # Track shape types
            if "auto_shape_type" in shape_info:
                result["global_design"]["shape_types_used"][
                    shape_info["auto_shape_type"]
                ] += 1

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _mark_render_required(slide_data, "grouped_shapes")
                slide_data["text_channels"].append(
                    {
                        "channel": "group_container_text",
                        "text": "",
                        "confidence": "low",
                        "status": "requires_render",
                        "provenance": {
                            "source": "pptx_group_container",
                            "shape_path": list(shape_path),
                        },
                    }
                )

            # Check for images
            if _is_picture_shape(shape):
                slide_data["has_image"] = True
                shape_info["picture_asset_status"] = "unavailable"
                shape_info["picture_part_name"] = None
                shape_info["picture_asset_sha256"] = None
                ratio = picture_area_ratio(shape, prs)
                if ratio > max_image_ratio:
                    max_image_ratio = ratio
                try:
                    blob, part_name = _picture_payload(shape)
                    shape_info["picture_part_name"] = part_name
                    shape_info["picture_asset_sha256"] = hashlib.sha256(
                        blob
                    ).hexdigest()
                    if part_name in corrupt_part_names:
                        shape_info["picture_asset_status"] = "corrupt"
                        _record_unsupported(
                            slide_data,
                            "corrupt_embedded_asset",
                            shape,
                            shape_path,
                        )
                    else:
                        shape_info["picture_asset_status"] = "available"
                        picture_entries.append(
                            (ratio, blob, list(shape_path), part_name)
                        )
                except (KeyError, ValueError, AttributeError) as e:
                    _record_unsupported(
                        slide_data,
                        "unreadable_picture",
                        shape,
                        shape_path,
                    )
                    sys.stderr.write(
                        f"WARN: could not read picture blob on slide {i + 1}: {e}\n"
                    )

            # Check for text-frame shapes
            if shape.has_text_frame:
                slide_data["has_text_frame_shapes"] = True
                channel = _shape_text_channel(
                    shape,
                    shape_path,
                    in_group=len(shape_path) > 1,
                )
                slide_data["text_channels"].append(channel)
                text_parts.append(channel["text"])

                # Detect footer by position (bottom 15% of slide) and small font
                if (
                    len(shape_path) == 1
                    and shape_info["top"] is not None
                    and shape_info["top"] > result["slide_height_inches"] * 0.85
                ):
                    slide_data["footer_text"] = channel["text"]

            graphic_kind = _classify_graphic_frame(shape)
            if graphic_kind is not None:
                uri = _graphic_data_uri(shape)
                shape_info["graphic_frame_type"] = graphic_kind
                shape_info["graphic_data_uri"] = uri
                if graphic_kind == "table":
                    table_channel, table_details = _extract_table_channel(
                        shape,
                        shape_path,
                    )
                    shape_info.update(table_details)
                    slide_data["text_channels"].append(table_channel)
                    text_parts.append(table_channel["text"])
                    for font_name, count in table_details["table_fonts"].items():
                        result["global_design"]["fonts_used"][font_name] += count
                    _mark_render_required(slide_data, "table")
                else:
                    _record_unsupported(
                        slide_data,
                        graphic_kind,
                        shape,
                        shape_path,
                        uri=uri,
                    )
            elif shape.shape_type in {
                MSO_SHAPE_TYPE.DIAGRAM,
                MSO_SHAPE_TYPE.IGX_GRAPHIC,
                MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
                MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
                MSO_SHAPE_TYPE.MEDIA,
                MSO_SHAPE_TYPE.WEB_VIDEO,
            }:
                _record_unsupported(
                    slide_data,
                    str(shape.shape_type).split(" ", 1)[0].lower(),
                    shape,
                    shape_path,
                )

        slide_data["shape_count_recursive"] = recursive_shape_count

        slide_data["text_content_preview"] = " | ".join(
            t[:50] for t in text_parts if t.strip()
        )[:200]

        slide_data["image_area_ratio"] = round(max_image_ratio, 3)

        # A picture large enough to carry text, or an image *background* (which
        # covers the whole slide by definition), can both be hiding rendered
        # text the shape walk never sees.
        if (
            slide_data["image_area_ratio"] >= PPTX_TEXT_BEARING_IMAGE_AREA_RATIO
            or bg_type == "image"
        ):
            reason = "background_image" if bg_type == "image" else "large_picture"
            _mark_render_required(slide_data, reason)

        picture_entries.sort(key=lambda item: item[0], reverse=True)
        picture_blobs = [entry[1] for entry in picture_entries]
        picture_paths = [entry[2] for entry in picture_entries]
        picture_part_names = [entry[3] for entry in picture_entries]
        apply_ocr_to_slide(
            slide_data,
            picture_blobs,
            ocr=ocr,
            ocr_fn=ocr_fn,
            shape_paths=picture_paths,
            part_names=picture_part_names,
        )

        if bg_type == "image":
            background_image = get_background_image(slide)
            if background_image is not None and background_image.get("blob"):
                slide_data["background_asset_status"] = "available"
                slide_data["background_part_name"] = background_image["provenance"].get(
                    "part_name"
                )
                slide_data["background_asset_sha256"] = hashlib.sha256(
                    background_image["blob"]
                ).hexdigest()
            if (
                background_image is not None
                and background_image["provenance"].get("part_name")
                in corrupt_part_names
            ):
                _mark_render_required(slide_data, "corrupt_embedded_asset")
                slide_data["background_asset_status"] = "corrupt"
                background_image["blob"] = None
                background_image["status"] = "recovered_with_placeholder"
                background_image["provenance"]["asset_status"] = (
                    "recovered_with_placeholder"
                )
                slide_data["unsupported_content"].append(
                    {
                        "content_type": "corrupt_embedded_asset",
                        "shape_name": None,
                        "shape_path": [],
                        "reason": "background image bytes failed package CRC validation",
                        "render_required": True,
                    }
                )
            apply_background_ocr(
                slide_data,
                background_image,
                ocr=ocr,
                ocr_fn=ocr_fn,
            )

        slide_data["has_extracted_text"] = any(
            channel["text"].strip() for channel in slide_data["text_channels"]
        )
        slide_data["has_unsupported_content"] = bool(slide_data["unsupported_content"])

        # Track background colors
        if bg_hex:
            result["global_design"]["background_colors"][bg_hex] += 1
        result["global_design"]["color_sequence"].append(bg_hex or "unknown")

        slide_data["render_required_reasons"].sort()
        result["per_slide_visual"].append(slide_data)

    required_reasons = {
        slide_data["slide_number"]: slide_data["render_required_reasons"]
        for slide_data in result["per_slide_visual"]
        if slide_data["render_required"]
    }
    if inspected_page_ranges and rendered_pdf_path is None:
        raise PptxEvidenceError(
            "inspected_page_ranges requires rendered_pdf_path so the assertion "
            "can be identity-bound"
        )
    render_receipt = None
    if rendered_pdf_path is not None:
        render_receipt = _build_rendered_page_inspection_in_process(
            source_pptx_sha256=source_fingerprint["digest"],
            rendered_pdf_path=rendered_pdf_path,
            inspected_page_ranges=inspected_page_ranges or [],
            required_slide_numbers=sorted(required_reasons),
            slide_count=len(prs.slides),
            rendered_pdf_generation=rendered_pdf_generation,
        )
    result["native_deck_audit"] = build_native_deck_audit(
        source_pptx_sha256=source_fingerprint["digest"],
        source_pptx_size_bytes=source_fingerprint["size_bytes"],
        slide_count=len(prs.slides),
        render_required_reasons=required_reasons,
        rendered_page_inspection=render_receipt,
    )

    # Convert Counters to dicts for JSON serialization
    result["global_design"]["fonts_used"] = dict(result["global_design"]["fonts_used"])
    result["global_design"]["background_colors"] = dict(
        result["global_design"]["background_colors"]
    )
    result["global_design"]["shape_types_used"] = dict(
        result["global_design"]["shape_types_used"]
    )
    result["native_timing_summary"] = summarize_native_timing(
        result["per_slide_visual"]
    )

    return result


def extract_pptx(
    pptx_path,
    *,
    trusted_root=None,
    ocr=True,
    ocr_fn=None,
    rendered_pdf_path=None,
    inspected_page_ranges=None,
    source_size_limit_bytes=None,
    deadline_monotonic=None,
):
    """Extract one deck through the authenticated resource supervisor.

    ``ocr_fn`` is deliberately confined to ``_extract_pptx_in_process`` as a
    unit-test seam; callables cannot cross the private worker protocol.
    """
    if ocr_fn is not None:
        raise PptxEvidenceError(
            "ocr_fn is available only inside an already-contained worker",
            reason_code="pptx_evidence_invalid",
        )
    options = {
        "trusted_root": trusted_root,
        "ocr": ocr,
        "rendered_pdf_path": rendered_pdf_path,
        "inspected_page_ranges": inspected_page_ranges,
    }
    if source_size_limit_bytes is not None:
        options["source_size_limit_bytes"] = source_size_limit_bytes
    if deadline_monotonic is not None:
        options["deadline_monotonic"] = deadline_monotonic
    if trusted_root is None:
        options.pop("trusted_root")
    return run_supervised_pptx_extraction(pptx_path, **options)


def should_skip(basename, skip_patterns):
    """Check if a .pptx file should be skipped."""
    lower = basename.lower()
    if basename.startswith("~$"):
        return True, "Office lock file"
    # Skip static exports
    if "static" in lower:
        return True, "static export"
    # Skip Google Drive conflict copies: (N).pptx
    if re.search(r"\(\d+\)\.pptx$", basename):
        return True, "conflict copy"
    # Skip files matching user-provided skip patterns (case-insensitive)
    for pat in skip_patterns:
        if pat.lower() in lower:
            return True, f"matches skip pattern '{pat}'"
    return False, None


def _batch_skip_reason(reason):
    """Collapse human-facing skip prose to a bounded stable reason code."""
    if reason == "Office lock file":
        return "pptx_batch_office_lock_file"
    if reason == "static export":
        return "pptx_batch_static_export"
    if reason == "conflict copy":
        return "pptx_batch_conflict_copy"
    return "pptx_batch_skip_pattern"


def _batch_error_reason(exc):
    """Return a path-free reason code for one failed supervised extraction."""
    reason_code = getattr(exc, "reason_code", None)
    if (
        isinstance(reason_code, str)
        and reason_code in _BATCH_EXTRACTION_FAILURE_REASONS
    ):
        return reason_code
    return "pptx_extraction_failed"


def _is_windows_reparse_point(stat_result):
    attributes = getattr(stat_result, "st_file_attributes", 0)
    return bool(
        isinstance(attributes, int)
        and not isinstance(attributes, bool)
        and attributes & _WINDOWS_REPARSE_POINT_ATTRIBUTE
    )


def _windows_leaf_rejection_reason(stat_result):
    """Reject unavailable Cloud Files and every unsupported redirecting leaf."""
    attributes = getattr(stat_result, "st_file_attributes", 0)
    if (
        not isinstance(attributes, int)
        or isinstance(attributes, bool)
        or attributes < 0
    ):
        return "pptx_batch_reparse_point_rejected"
    if attributes & _WINDOWS_UNAVAILABLE_CLOUD_ATTRIBUTES:
        return "pptx_batch_cloud_placeholder_unavailable"
    if not _is_windows_reparse_point(stat_result):
        return None
    tag = getattr(stat_result, "st_reparse_tag", None)
    if (
        isinstance(tag, int)
        and not isinstance(tag, bool)
        and tag in _WINDOWS_CLOUD_REPARSE_TAGS
    ):
        # A hydrated Cloud Files leaf can retain its supported Cloud tag after
        # its data is local. Offline/recall attributes above remain fail-closed.
        return None
    return "pptx_batch_reparse_point_rejected"


def _usable_directory_identity(stat_result):
    device = getattr(stat_result, "st_dev", None)
    inode = getattr(stat_result, "st_ino", None)
    if (
        isinstance(device, bool)
        or not isinstance(device, int)
        or device <= 0
        or isinstance(inode, bool)
        or not isinstance(inode, int)
        or inode <= 0
    ):
        return None
    return (device, inode)


def _entry_is_real_excluded_directory(entry, exclusions):
    """Identify one policy-prunable dirent without trusting its name alone."""
    if not directory_component_is_excluded(entry.name, exclusions):
        return False
    try:
        entry_stat = entry.stat(follow_symlinks=False)
        return bool(
            entry.is_dir(follow_symlinks=False)
            and not _is_windows_reparse_point(entry_stat)
            and not entry.is_symlink()
        )
    except OSError:
        return False


def _discover_pptx_files(directory, skip_patterns, directory_exclusions=()):
    """Discover a deterministic, bounded set inside the contained worker."""
    root = Path(_validated_directory_root(directory))
    try:
        exclusions = validate_pptx_directory_exclusions(directory_exclusions)
    except PptxDiscoveryContractError as exc:
        raise PptxEvidenceError(
            "PPTX directory exclusions are invalid",
            reason_code="pptx_batch_request_invalid",
        ) from exc
    try:
        root_stat = root.lstat()
    except OSError as exc:
        raise PptxEvidenceError(
            "PPTX batch root is unavailable",
            reason_code="pptx_batch_root_unavailable",
            details={"exception_type": type(exc).__name__},
        ) from exc
    if (
        stat.S_ISLNK(root_stat.st_mode)
        or _is_windows_reparse_point(root_stat)
        or not stat.S_ISDIR(root_stat.st_mode)
    ):
        raise PptxEvidenceError(
            "PPTX batch root must be a non-symlink directory",
            reason_code="pptx_batch_root_invalid",
        )

    started = time.monotonic()
    discovered = []
    skipped = []
    skip_receipts = set()

    def record_skip(path, reason):
        """Keep the worker's own closed manifest self-compatible."""
        receipt = (path, reason)
        if receipt not in skip_receipts:
            skip_receipts.add(receipt)
            skipped.append({"path": path, "reason": reason})

    stack = [(root, "", 0)]
    visited = set()
    directory_count = 0
    entry_count = 0
    policy_excluded_entry_count = 0

    while stack:
        if time.monotonic() - started > _BATCH_MAX_WALL_SECONDS:
            record_skip(".", "pptx_batch_wall_limit")
            break
        current, relative_directory, depth = stack.pop()
        try:
            current_stat = current.lstat()
        except OSError:
            record_skip(
                relative_directory or ".",
                "pptx_batch_directory_unavailable",
            )
            continue
        if (
            stat.S_ISLNK(current_stat.st_mode)
            or _is_windows_reparse_point(current_stat)
            or not stat.S_ISDIR(current_stat.st_mode)
        ):
            record_skip(
                relative_directory or ".",
                "pptx_batch_directory_changed",
            )
            continue
        identity = _usable_directory_identity(current_stat)
        if identity is None:
            record_skip(
                relative_directory or ".",
                "pptx_batch_directory_identity_unavailable",
            )
            continue
        if identity in visited:
            record_skip(
                relative_directory or ".",
                "pptx_batch_directory_identity_collision",
            )
            continue
        visited.add(identity)
        directory_count += 1
        if directory_count > _BATCH_MAX_DIRECTORIES:
            record_skip(".", "pptx_batch_directory_limit")
            break
        try:
            entries = []
            entry_limit_hit = False
            with os.scandir(current) as iterator:
                for entry in iterator:
                    if _entry_is_real_excluded_directory(entry, exclusions):
                        if (
                            policy_excluded_entry_count
                            >= _BATCH_MAX_POLICY_EXCLUDED_ENTRIES
                        ):
                            entry_limit_hit = True
                            break
                        entries.append(entry)
                        policy_excluded_entry_count += 1
                        continue
                    if entry_count >= _BATCH_MAX_ENTRIES:
                        entry_limit_hit = True
                        break
                    entries.append(entry)
                    entry_count += 1
        except OSError:
            record_skip(
                relative_directory or ".",
                "pptx_batch_directory_unavailable",
            )
            continue
        if entry_limit_hit:
            record_skip(".", "pptx_batch_entry_limit")
            break
        entries.sort(key=lambda item: item.name)

        child_directories = []
        limit_hit = False
        for entry in entries:
            relative = (
                f"{relative_directory}/{entry.name}"
                if relative_directory
                else entry.name
            )
            if entry.name in {"", ".", ".."} or "\\" in entry.name:
                record_skip(".", "pptx_batch_path_invalid")
                continue
            if len(relative) > _BATCH_MAX_RELATIVE_PATH_CHARS:
                record_skip(
                    relative[:_BATCH_MAX_RELATIVE_PATH_CHARS],
                    "pptx_batch_path_limit",
                )
                continue
            try:
                entry_stat = entry.stat(follow_symlinks=False)
                entry_is_directory = entry.is_dir(follow_symlinks=False)
                leaf_rejection = _windows_leaf_rejection_reason(entry_stat)
                if leaf_rejection is not None and not entry_is_directory:
                    record_skip(relative, leaf_rejection)
                    continue
                if entry_is_directory and _is_windows_reparse_point(entry_stat):
                    record_skip(
                        relative,
                        "pptx_batch_reparse_point_rejected",
                    )
                    continue
                if entry.is_symlink():
                    record_skip(
                        relative,
                        "pptx_batch_symlink_rejected",
                    )
                    continue
                if entry_is_directory:
                    if directory_component_is_excluded(entry.name, exclusions):
                        record_skip(relative, "pptx_batch_directory_excluded")
                    elif depth >= _BATCH_MAX_DEPTH:
                        record_skip(relative, "pptx_batch_depth_limit")
                    else:
                        child_directories.append(
                            (Path(entry.path), relative, depth + 1)
                        )
                    continue
                if not entry.is_file(
                    follow_symlinks=False
                ) or not entry.name.lower().endswith(".pptx"):
                    continue
                skip, reason = should_skip(entry.name, skip_patterns)
                if skip:
                    record_skip(relative, _batch_skip_reason(reason))
                    continue
            except OSError:
                record_skip(relative, "pptx_batch_entry_unavailable")
                continue
            if len(discovered) >= _BATCH_MAX_FILES:
                record_skip(
                    relative,
                    "pptx_batch_file_limit",
                )
                record_skip(
                    ".",
                    "pptx_batch_scan_incomplete_file_limit",
                )
                stack.clear()
                limit_hit = True
                break
            # Do not retain discovery-time size as budget authority. Cloud
            # hydration or replacement can change it before worker admission;
            # the supervisor snapshots and enforces the exact launched generation.
            discovered.append((Path(entry.path), relative))
        if not limit_hit:
            stack.extend(reversed(child_directories))

    return discovered, skipped, started


def _validated_directory_root(value):
    try:
        root = materialize_native_root(value)
    except ArtifactLocatorError as exc:
        raise PptxEvidenceError(
            "directory root must be a bounded native absolute path",
            reason_code="pptx_batch_root_invalid",
            details={"locator_failure": exc.reason_code},
        ) from exc
    rendered = os.fspath(root)
    if len(rendered) > _BATCH_MAX_ROOT_PATH_CHARS:
        raise PptxEvidenceError(
            "directory root must be a bounded native absolute path",
            reason_code="pptx_batch_root_invalid",
            details={"locator_failure": "directory_root_path_limit"},
        )
    return rendered


def _validated_skip_patterns(value):
    if not isinstance(value, (list, tuple)) or len(value) > _BATCH_MAX_SKIP_PATTERNS:
        raise PptxEvidenceError(
            "directory skip patterns exceed their bounded contract",
            reason_code="pptx_batch_request_invalid",
        )
    normalized = []
    for pattern in value:
        if (
            not isinstance(pattern, str)
            or not pattern
            or "\x00" in pattern
            or len(pattern) > _BATCH_MAX_SKIP_PATTERN_CHARS
        ):
            raise PptxEvidenceError(
                "directory skip pattern is invalid",
                reason_code="pptx_batch_request_invalid",
            )
        normalized.append(pattern)
    return normalized


def _validated_directory_exclusions(value):
    try:
        return validate_pptx_directory_exclusions(
            value,
            label="directory_exclusions",
        )
    except PptxDiscoveryContractError as exc:
        raise PptxEvidenceError(
            "directory exclusions exceed their bounded exact-component contract",
            reason_code="pptx_batch_request_invalid",
        ) from exc


def _validated_relative_manifest_path(value, *, allow_root=False):
    if allow_root and value == ".":
        return value
    if (
        not isinstance(value, str)
        or not value
        or len(value) > _BATCH_MAX_RELATIVE_PATH_CHARS
        or "\x00" in value
        or "\\" in value
        or value.startswith("/")
        or re.match(r"^[A-Za-z]:", value)
    ):
        raise PptxEvidenceError(
            "directory worker returned a noncanonical relative path",
            reason_code="pptx_batch_manifest_invalid",
        )
    components = value.split("/")
    if any(component in {"", ".", ".."} for component in components):
        raise PptxEvidenceError(
            "directory worker returned a noncanonical relative path",
            reason_code="pptx_batch_manifest_invalid",
        )
    return value


def _decode_directory_manifest(value, *, expected_directory_exclusions=()):
    """Decode only the closed, authenticated directory-worker body."""
    try:
        expected_exclusions = validate_pptx_directory_exclusions(
            expected_directory_exclusions,
            label="expected_directory_exclusions",
        )
    except PptxDiscoveryContractError as exc:
        raise PptxEvidenceError(
            "expected directory exclusions violate their contract",
            reason_code="pptx_batch_manifest_invalid",
        ) from exc
    if not isinstance(value, dict) or set(value) != {
        "schema_version",
        "kind",
        "complete",
        "directory_exclusions",
        "incomplete_reason_codes",
        "files",
        "skipped",
    }:
        raise PptxEvidenceError(
            "directory worker returned an invalid manifest",
            reason_code="pptx_batch_manifest_invalid",
        )
    if (
        type(value.get("schema_version")) is not int
        or value.get("schema_version") != _DIRECTORY_MANIFEST_SCHEMA_VERSION
        or value.get("kind") != PPTX_DIRECTORY_MANIFEST_KIND
        or type(value.get("complete")) is not bool
        or not isinstance(value.get("incomplete_reason_codes"), list)
    ):
        raise PptxEvidenceError(
            "directory worker returned an unsupported manifest",
            reason_code="pptx_batch_manifest_invalid",
        )
    try:
        returned_exclusions = validate_pptx_directory_exclusions(
            value.get("directory_exclusions"),
            label="manifest.directory_exclusions",
        )
    except PptxDiscoveryContractError as exc:
        raise PptxEvidenceError(
            "directory worker returned invalid exclusion policy",
            reason_code="pptx_batch_manifest_invalid",
        ) from exc
    if returned_exclusions != expected_exclusions:
        raise PptxEvidenceError(
            "directory worker exclusion policy does not match its request",
            reason_code="pptx_batch_manifest_invalid",
        )
    exclusion_identities = {
        component.casefold() for component in expected_exclusions
    }
    raw_files = value.get("files")
    raw_skipped = value.get("skipped")
    if (
        not isinstance(raw_files, list)
        or len(raw_files) > _BATCH_MAX_FILES
        or not isinstance(raw_skipped, list)
        or len(raw_skipped) > _BATCH_MAX_ENTRIES + _BATCH_MAX_DIRECTORIES + 1
    ):
        raise PptxEvidenceError(
            "directory worker manifest exceeds its collection bounds",
            reason_code="pptx_batch_manifest_invalid",
        )

    files = []
    seen_paths = set()
    for item in raw_files:
        if not isinstance(item, dict) or set(item) != {"path"}:
            raise PptxEvidenceError(
                "directory worker file binding is invalid",
                reason_code="pptx_batch_manifest_invalid",
            )
        relative = _validated_relative_manifest_path(item.get("path"))
        basename = relative.rsplit("/", 1)[-1]
        if (
            relative in seen_paths
            or not relative.casefold().endswith(".pptx")
            or basename.startswith("~$")
        ):
            raise PptxEvidenceError(
                "directory worker returned duplicate file bindings",
                reason_code="pptx_batch_manifest_invalid",
            )
        seen_paths.add(relative)
        files.append(relative)

    skipped = []
    seen_skip_receipts = set()
    seen_nonroot_skip_paths = set()
    for item in raw_skipped:
        if not isinstance(item, dict) or set(item) != {"path", "reason"}:
            raise PptxEvidenceError(
                "directory worker skip receipt is invalid",
                reason_code="pptx_batch_manifest_invalid",
            )
        path = _validated_relative_manifest_path(item.get("path"), allow_root=True)
        reason = item.get("reason")
        if (
            not isinstance(reason, str)
            or reason not in _DIRECTORY_MANIFEST_SKIP_REASONS
        ):
            raise PptxEvidenceError(
                "directory worker skip reason is invalid",
                reason_code="pptx_batch_manifest_invalid",
            )
        receipt = (path, reason)
        if (
            receipt in seen_skip_receipts
            or path in seen_paths
            or (path != "." and path in seen_nonroot_skip_paths)
        ):
            raise PptxEvidenceError(
                "directory worker returned duplicate skip receipts",
                reason_code="pptx_batch_manifest_invalid",
            )
        seen_skip_receipts.add(receipt)
        if path != ".":
            seen_nonroot_skip_paths.add(path)
        skipped.append({"path": path, "reason": reason})
    skipped_nonroot_paths = {
        item["path"] for item in skipped if item["path"] != "."
    }
    for item in skipped:
        path = item["path"]
        if path == ".":
            continue
        components = path.split("/")
        if any(
            component.casefold() in exclusion_identities
            for component in components[:-1]
        ):
            raise PptxEvidenceError(
                "directory worker returned a skip beneath an excluded directory",
                reason_code="pptx_batch_manifest_invalid",
            )
        if (
            item["reason"] == "pptx_batch_directory_excluded"
            and components[-1].casefold() not in exclusion_identities
        ):
            raise PptxEvidenceError(
                "directory worker fabricated an exclusion receipt",
                reason_code="pptx_batch_manifest_invalid",
            )
    for relative in files:
        components = relative.split("/")
        if any(
            component.casefold() in exclusion_identities
            for component in components[:-1]
        ) or any(
            relative.startswith(f"{skipped_path}/")
            for skipped_path in skipped_nonroot_paths
        ):
            raise PptxEvidenceError(
                "directory worker returned a file beneath a skipped directory",
                reason_code="pptx_batch_manifest_invalid",
            )
    try:
        incomplete_reason_codes = directory_incomplete_reason_codes(skipped)
    except PptxDiscoveryContractError as exc:
        raise PptxEvidenceError(
            "directory worker completeness receipts are invalid",
            reason_code="pptx_batch_manifest_invalid",
        ) from exc
    if (
        value.get("complete") != (not incomplete_reason_codes)
        or value.get("incomplete_reason_codes") != incomplete_reason_codes
    ):
        raise PptxEvidenceError(
            "directory worker completeness does not match its skip receipts",
            reason_code="pptx_batch_manifest_invalid",
        )
    return files, skipped, not incomplete_reason_codes, incomplete_reason_codes


def _directory_limits_before_deadline(deadline_monotonic):
    if (
        isinstance(deadline_monotonic, bool)
        or not isinstance(deadline_monotonic, (int, float))
        or not math.isfinite(float(deadline_monotonic))
    ):
        raise PptxEvidenceError(
            "directory deadline must be a finite monotonic timestamp",
            reason_code="pptx_batch_request_invalid",
        )
    remaining = (
        float(deadline_monotonic) - time.monotonic() - _DIRECTORY_LIMITS.cleanup_seconds
    )
    if remaining <= 0:
        raise PptxEvidenceError(
            "batch deadline expired before directory discovery",
            reason_code="pptx_batch_wall_limit",
        )
    if remaining >= _DIRECTORY_LIMITS.wall_seconds:
        return _DIRECTORY_LIMITS, False
    return replace(_DIRECTORY_LIMITS, wall_seconds=remaining), True


def _run_supervised_directory_discovery(
    directory,
    skip_patterns,
    directory_exclusions=(),
    *,
    deadline,
):
    """Resolve one directory only through the authenticated bounded worker."""
    root = _validated_directory_root(directory)
    patterns = _validated_skip_patterns(skip_patterns)
    exclusions = _validated_directory_exclusions(directory_exclusions)
    limits, deadline_limited = _directory_limits_before_deadline(deadline)
    command = [sys.executable, os.path.abspath(__file__), _DIRECTORY_WORKER_FLAG]
    payload: dict[str, JsonValue] = {
        "root_path": root,
        "skip_patterns": cast(JsonValue, patterns),
        "directory_exclusions": cast(JsonValue, exclusions),
    }
    try:
        result = run_authenticated_worker(
            command,
            _DIRECTORY_OPERATION,
            {},
            payload,
            limits,
            immutable_process_identity=command[:2],
            sensitive_values=(root,),
            schema_generation=SCHEMA_VERSION,
            pipeline_generation=PIPELINE_VERSION,
        )
    except SupervisorError as exc:
        if exc.reason_code == "worker_timeout":
            reason = (
                "pptx_batch_wall_limit"
                if deadline_limited
                else "pptx_batch_discovery_timeout"
            )
        elif exc.reason_code == "worker_output_limit_exceeded":
            reason = "pptx_batch_discovery_output_limit"
        elif exc.reason_code in _DIRECTORY_CHILD_FAILURES:
            reason = exc.reason_code
        elif exc.reason_code in _DIRECTORY_RESOURCE_FAILURES:
            reason = "pptx_batch_discovery_resource_unavailable"
        elif exc.reason_code in _DIRECTORY_START_FAILURES:
            reason = "pptx_batch_discovery_start_failure"
        elif exc.reason_code in _DIRECTORY_WORKER_FAILURES:
            reason = "pptx_batch_discovery_worker_failure"
        else:
            reason = "pptx_batch_discovery_protocol_invalid"
        raise PptxEvidenceError(
            "bounded directory discovery failed",
            reason_code=reason,
            details={"supervisor_reason_code": exc.reason_code},
        ) from exc
    return _decode_directory_manifest(
        result.payload,
        expected_directory_exclusions=exclusions,
    )


def _dispatch_directory_worker(request: WorkerRequest):
    if (
        request.operation != _DIRECTORY_OPERATION
        or request.limit_profile_id != _DIRECTORY_LIMITS.profile_id
        or request.schema_generation != SCHEMA_VERSION
        or request.pipeline_generation != PIPELINE_VERSION
        or request.expected_generations
        or not isinstance(request.payload, dict)
        or set(request.payload)
        != {"root_path", "skip_patterns", "directory_exclusions"}
    ):
        raise SupervisorError("invalid_worker_request")
    try:
        root = _validated_directory_root(request.payload.get("root_path"))
        patterns = _validated_skip_patterns(request.payload.get("skip_patterns"))
        exclusions = _validated_directory_exclusions(
            request.payload.get("directory_exclusions")
        )
        discovered, skipped, _started = _discover_pptx_files(
            root,
            patterns,
            exclusions,
        )
    except PptxEvidenceError as exc:
        locator_failure = exc.details.get("locator_failure")
        details = (
            {"locator_failure": locator_failure}
            if isinstance(locator_failure, str)
            else {}
        )
        raise SupervisorError(exc.reason_code, details) from exc
    incomplete_reason_codes = directory_incomplete_reason_codes(skipped)
    return {
        "schema_version": _DIRECTORY_MANIFEST_SCHEMA_VERSION,
        "kind": PPTX_DIRECTORY_MANIFEST_KIND,
        "complete": not incomplete_reason_codes,
        "directory_exclusions": exclusions,
        "incomplete_reason_codes": incomplete_reason_codes,
        "files": [{"path": relative} for _path, relative in discovered],
        "skipped": skipped,
    }


def _run_directory_worker_child():
    request = read_worker_request(max_input_bytes=_DIRECTORY_LIMITS.max_input_bytes)
    protocol_output = isolate_protocol_output()
    try:
        try:
            payload = _dispatch_directory_worker(request)
            write_worker_response(
                request,
                payload=payload,
                observed_generations={},
                stream=protocol_output,
                max_output_bytes=_DIRECTORY_LIMITS.max_output_bytes,
            )
        except SupervisorError as exc:
            write_worker_response(
                request,
                error=SupervisorError(exc.reason_code, exc.details),
                observed_generations={},
                stream=protocol_output,
                max_output_bytes=_DIRECTORY_LIMITS.max_output_bytes,
            )
    finally:
        protocol_output.close()
    return 0


def _build_batch_output(results, skipped):
    """Build the public completeness envelope from all retained receipts."""
    try:
        return build_pptx_directory_batch(results, skipped)
    except PptxDiscoveryContractError as exc:
        raise PptxEvidenceError(
            "PPTX directory batch output violates its closed contract",
            reason_code="pptx_batch_manifest_invalid",
        ) from exc


def _encode_batch_output(batch):
    """Encode the exact compact public JSON emitted by directory mode."""
    return json.dumps(
        batch,
        ensure_ascii=False,
        separators=(",", ":"),
    ).encode("utf-8")


def _encode_batch_failure(error):
    """Encode one path-neutral whole-root failure for machine callers."""
    details = {}
    supervisor_reason = error.details.get("supervisor_reason_code")
    if isinstance(supervisor_reason, str):
        details["supervisor_reason_code"] = supervisor_reason
    try:
        batch = build_pptx_directory_batch(
            [],
            [{"path": ".", "reason": error.reason_code}],
            error={"reason_code": error.reason_code, "details": details},
        )
    except PptxDiscoveryContractError as exc:
        raise PptxEvidenceError(
            "whole-root failure is outside the public batch contract",
            reason_code="pptx_batch_manifest_invalid",
        ) from exc
    return _encode_batch_output(batch)


def _encoded_json_size(value):
    return len(
        json.dumps(
            value,
            ensure_ascii=False,
            separators=(",", ":"),
        ).encode("utf-8")
    )


def _batch_output_size(result_sizes, skip_sizes, incomplete_reason_codes):
    """Return exact compact envelope size from pre-encoded element sizes."""
    metadata = {
        "schema_version": PPTX_DIRECTORY_BATCH_SCHEMA_VERSION,
        "kind": PPTX_DIRECTORY_BATCH_KIND,
        "complete": not incomplete_reason_codes,
        "incomplete_reason_codes": list(incomplete_reason_codes),
        "results": [],
        "skipped": [],
    }
    empty_size = _encoded_json_size(metadata)
    return (
        empty_size
        + sum(result_sizes)
        + max(0, len(result_sizes) - 1)
        + sum(skip_sizes)
        + max(0, len(skip_sizes) - 1)
    )


def batch_extract(
    directory,
    skip_patterns,
    directory_exclusions=(),
    *,
    ocr=True,
):
    """Extract a bounded deterministic directory batch through the supervisor."""
    results = []
    started = time.monotonic()
    deadline = started + _BATCH_MAX_WALL_SECONDS
    root = Path(_validated_directory_root(directory))
    (
        relative_files,
        skipped,
        _discovery_complete,
        _discovery_incomplete_reasons,
    ) = _run_supervised_directory_discovery(
        root,
        skip_patterns,
        directory_exclusions,
        deadline=deadline,
    )
    discovered = [
        (root.joinpath(*relative.split("/")), relative) for relative in relative_files
    ]
    input_bytes = 0
    result_sizes = []
    skip_sizes = [_encoded_json_size(item) for item in skipped]

    def append_skip(path, reason):
        record = {"path": path, "reason": reason}
        skipped.append(record)
        skip_sizes.append(_encoded_json_size(record))

    def skip_remaining(index, reason):
        for remaining in discovered[index:]:
            append_skip(remaining[1], reason)

    for index, (pptx_path, relative) in enumerate(discovered):
        if time.monotonic() >= deadline:
            skip_remaining(index, "pptx_batch_wall_limit")
            break
        remaining_input_bytes = _BATCH_MAX_INPUT_BYTES - input_bytes
        if remaining_input_bytes <= 0:
            skip_remaining(index, "pptx_batch_input_limit")
            break
        try:
            data = extract_pptx(
                pptx_path,
                trusted_root=root,
                ocr=ocr,
                source_size_limit_bytes=remaining_input_bytes,
                deadline_monotonic=deadline,
            )
            if time.monotonic() >= deadline:
                skip_remaining(index, "pptx_batch_wall_limit")
                break
            fingerprint = data.get("input_fingerprint")
            exact_size = (
                fingerprint.get("size_bytes") if isinstance(fingerprint, dict) else None
            )
            if (
                isinstance(exact_size, bool)
                or not isinstance(exact_size, int)
                or exact_size < 1
                or exact_size > remaining_input_bytes
            ):
                raise PptxEvidenceError(
                    "supervised extraction returned an invalid source size",
                    reason_code="pptx_evidence_invalid",
                )
            # Directory output is relocatable and never exposes the root path.
            data["pptx_path"] = relative
            encoded_size = _encoded_json_size(data)
            # Reserve a maximum-length stable reason for every deck not yet
            # launched. This guarantees later per-file failures still fit the
            # same compact wrapper without retroactively dropping results.
            future_skip_sizes = [
                _encoded_json_size({"path": remaining[1], "reason": "x" * 96})
                for remaining in discovered[index + 1 :]
            ]
            if (
                _batch_output_size(
                    [*result_sizes, encoded_size],
                    [*skip_sizes, *future_skip_sizes],
                    sorted(PPTX_DIRECTORY_INCOMPLETE_REASON_CODES),
                )
                + 1
                > _BATCH_MAX_OUTPUT_BYTES
            ):
                skip_remaining(index, "pptx_batch_output_limit")
                break
            results.append(data)
            result_sizes.append(encoded_size)
            input_bytes += exact_size
        except (
            PptxEvidenceError,
            OSError,
            ValueError,
            KeyError,
            AttributeError,
        ) as exc:
            # The evidence boundary reports the exact generation admitted for
            # a launched parse. Discovery metadata is intentionally not a
            # substitute: hydration or replacement can change the generation
            # before admission.
            admitted_size = getattr(exc, "details", {}).get(
                "admitted_source_size_bytes"
            )
            if (
                isinstance(admitted_size, int)
                and not isinstance(admitted_size, bool)
                and 0 < admitted_size <= remaining_input_bytes
            ):
                input_bytes += admitted_size
            reason = _batch_error_reason(exc)
            if reason in {"pptx_batch_input_limit", "pptx_batch_wall_limit"}:
                skip_remaining(index, reason)
                break
            append_skip(relative, reason)
            if (
                _batch_output_size(
                    result_sizes,
                    skip_sizes,
                    sorted(PPTX_DIRECTORY_INCOMPLETE_REASON_CODES),
                )
                + 1
                > _BATCH_MAX_OUTPUT_BYTES
            ):
                return _build_batch_output(
                    results,
                    [{"path": ".", "reason": "pptx_batch_output_limit"}],
                )

    batch = _build_batch_output(results, skipped)
    encoded_batch = _encode_batch_output(batch)
    if len(encoded_batch) + 1 > _BATCH_MAX_OUTPUT_BYTES:
        return _build_batch_output(
            results,
            [{"path": ".", "reason": "pptx_batch_output_limit"}],
        )
    return batch


def main(argv=None):
    parser = argparse.ArgumentParser(
        description="Extract visual design data from .pptx files."
    )
    parser.add_argument(
        "path",
        nargs="?",
        help="Single .pptx file, or a directory root with --directory",
    )
    parser.add_argument(
        "--directory",
        action="store_true",
        help="Treat path as a bounded recursive directory root",
    )
    parser.add_argument(
        "--skip",
        action="append",
        default=None,
        help=(
            "Configured skip pattern (case-insensitive); repeat for each pattern "
            "and omit all --skip flags for an empty set"
        ),
    )
    parser.add_argument(
        "--exclude-directory",
        action="append",
        default=None,
        metavar="COMPONENT",
        help=(
            "Exact directory-name component to prune; repeat for each configured "
            "exclusion and omit all flags for an empty set"
        ),
    )
    parser.add_argument(
        "--no-ocr",
        action="store_true",
        help="Skip OCR on low-confidence slides (shape walk only)",
    )
    parser.add_argument(
        "--rendered-pdf",
        help="PDF rendered from the exact input PPTX for identity-bound review",
    )
    parser.add_argument(
        "--inspected-pages",
        action="append",
        default=[],
        metavar="PAGE|START-END",
        help="Rendered pages actually inspected; repeat or comma-separate values",
    )
    parser.add_argument(
        "--version",
        action="store_true",
        help="Print extractor schema and pipeline versions as JSON",
    )
    args = parser.parse_args(argv)
    if args.version:
        print(
            json.dumps(
                {
                    "schema_version": SCHEMA_VERSION,
                    "pipeline_version": PIPELINE_VERSION,
                }
            )
        )
        return
    if args.path is None:
        parser.error("path is required unless --version is used")
    ocr = not args.no_ocr
    try:
        inspected_page_ranges = parse_page_range_arguments(args.inspected_pages)
    except PptxEvidenceError as exc:
        parser.error(str(exc))

    if args.directory:
        if args.rendered_pdf or inspected_page_ranges:
            parser.error("--rendered-pdf/--inspected-pages require a single PPTX input")
        try:
            batch = batch_extract(
                args.path,
                args.skip or [],
                args.exclude_directory or [],
                ocr=ocr,
            )
        except PptxEvidenceError as exc:
            if exc.reason_code in _DIRECTORY_BATCH_FAILURE_REASONS:
                sys.stdout.write(_encode_batch_failure(exc).decode("utf-8") + "\n")
            print(f"ERROR: {exc.reason_code}", file=sys.stderr)
            return 1
        sys.stdout.write(_encode_batch_output(batch).decode("utf-8") + "\n")
    else:
        try:
            result = extract_pptx(
                args.path,
                ocr=ocr,
                rendered_pdf_path=args.rendered_pdf,
                inspected_page_ranges=inspected_page_ranges,
            )
        except PptxEvidenceError as exc:
            print(f"ERROR: {exc.reason_code}", file=sys.stderr)
            return 1
        # outer-boundary-process-contract: the preflight caller treats an
        # invalid/nonzero process result as its silent-failure shape; this catch
        # emits one concise diagnostic because propagation would replace that
        # contract with an unstructured traceback.
        except Exception:  # noqa: BLE001 - outer-boundary-process-contract
            print("ERROR: pptx_extraction_failed", file=sys.stderr)
            return 1
        print(json.dumps(result, indent=2))
    return 0


def _main():
    if sys.argv[1:] == [_DIRECTORY_WORKER_FLAG]:
        try:
            return _run_directory_worker_child()
        except SupervisorError as exc:
            print(
                f"pptx directory worker failed: {exc.reason_code}",
                file=sys.stderr,
                flush=True,
            )
            return 2
        # The supervisor treats a nonzero child without an authenticated
        # response as a bounded crash. Emit a path-neutral stderr diagnostic
        # plus exit 2 because propagation would leak a traceback and violate
        # the one-frame response contract. outer-boundary-process-contract.
        except Exception:  # noqa: BLE001
            print(
                "pptx directory worker failed: unexpected_error",
                file=sys.stderr,
                flush=True,
            )
            return 2
    return main()


if __name__ == "__main__":
    raise SystemExit(_main())
