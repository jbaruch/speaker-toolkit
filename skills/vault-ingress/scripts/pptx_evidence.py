#!/usr/bin/env python3
"""Shared native-deck recovery and exact render-inspection receipts."""

from __future__ import annotations

import base64
import copy
import hashlib
import importlib.util
import io
import json
import math
import os
import re
import signal
import stat as stat_module
import sys
import tempfile
import time
import xml.etree.ElementTree as ElementTree
import zipfile
from bisect import bisect_left
from collections import Counter
from collections.abc import Mapping
from dataclasses import dataclass, replace
from pathlib import Path
from typing import Any, BinaryIO, Callable, cast
from zlib import error as ZlibError

from artifact_supervisor import (
    FileGeneration,
    SupervisorError,
    SupervisorLimits,
    WorkerRequest,
    isolate_protocol_output,
    read_worker_request,
    run_authenticated_worker,
    write_worker_response,
)


PPTX_EXTRACTION_SCHEMA_VERSION = 4
PPTX_EXTRACTION_PIPELINE_VERSION = "1.4.0"
ARCHIVE_RECOVERY_SCHEMA_VERSION = 1
NATIVE_DECK_AUDIT_SCHEMA_VERSION = 1
RENDER_INSPECTION_SCHEMA_VERSION = 1
PPTX_ARTIFACT_PROBE_SCHEMA_VERSION = 1
PPTX_ARTIFACT_PROBE_TIMEOUT_SECONDS = 15
PPTX_OCR_TRUST_CONFIDENCE = 50.0
# A reported picture-area ratio at or above this boundary requires rendered
# inspection because the picture can carry text invisible to the shape walk.
# Producer and validator both consume this one authority so they cannot drift.
PPTX_TEXT_BEARING_IMAGE_AREA_RATIO = 0.5
PPTX_ARTIFACT_PROBE_MAX_RESULT_BYTES = 64 * 1024
PPTX_ARTIFACT_PROBE_MAX_RECOVERY_RECORDS = 64
PPTX_ARCHIVE_MAX_MEMBERS = 65_536
PPTX_ARCHIVE_MAX_EXPANDED_BYTES = 4 * 1024 * 1024 * 1024
PPTX_ARCHIVE_MAX_MEMBER_BYTES = 512 * 1024 * 1024
PPTX_ARTIFACT_PROBE_CHILD_FLAG = "--artifact-probe-child"
PPTX_NATIVE_AUDIT_CHILD_FLAG = "--native-audit-child"
PPTX_SUPERVISED_WORKER_FLAG = "--supervised-worker"
PPTX_METADATA_OPERATION = "pptx_metadata"
PPTX_PROBE_OPERATION = "pptx_probe"
PPTX_NATIVE_AUDIT_OPERATION = "pptx_native_audit"
PPTX_EXTRACT_OPERATION = "pptx_extract"
# A hydrated source deck may legitimately be much larger than its cloud
# placeholder metadata suggested.  Keep a hard per-artifact ceiling, but leave
# enough room for the ~1.045 GB authored decks present in real vaults.
PPTX_MAX_INPUT_BYTES = 2 * 1024 * 1024 * 1024
PPTX_METADATA_LIMITS = SupervisorLimits(
    profile_id="pptx-metadata-v1",
    wall_seconds=15,
    max_memory_bytes=256 * 1024 * 1024,
    max_input_bytes=64 * 1024,
    max_output_bytes=64 * 1024,
    max_processes=1,
)
PPTX_PROBE_LIMITS = SupervisorLimits(
    profile_id="pptx-probe-v1",
    wall_seconds=15,
    max_memory_bytes=4 * 1024 * 1024 * 1024,
    max_input_bytes=64 * 1024,
    max_output_bytes=64 * 1024,
)
PPTX_NATIVE_AUDIT_LIMITS = SupervisorLimits(
    profile_id="pptx-native-audit-v1",
    wall_seconds=180,
    max_memory_bytes=4 * 1024 * 1024 * 1024,
    max_input_bytes=64 * 1024,
    max_output_bytes=2 * 1024 * 1024,
)
PPTX_EXTRACT_NO_OCR_LIMITS = SupervisorLimits(
    profile_id="pptx-extract-no-ocr-v1",
    wall_seconds=300,
    max_memory_bytes=4 * 1024 * 1024 * 1024,
    max_input_bytes=64 * 1024,
    max_output_bytes=128 * 1024 * 1024,
)
PPTX_EXTRACT_OCR_LIMITS = SupervisorLimits(
    profile_id="pptx-extract-ocr-v1",
    wall_seconds=900,
    max_memory_bytes=4 * 1024 * 1024 * 1024,
    max_input_bytes=64 * 1024,
    max_output_bytes=128 * 1024 * 1024,
)
PPTX_MACOS_DATALESS_FLAG = int(
    getattr(
        stat_module,
        "SF_DATALESS",
        0x40000000 if sys.platform == "darwin" else 0,
    )
)
_WINDOWS_OFFLINE_FILE_ATTRIBUTES = 0x001000 | 0x040000 | 0x400000
PPTX_WINDOWS_CLOUD_FILE_ATTRIBUTES = (
    _WINDOWS_OFFLINE_FILE_ATTRIBUTES if os.name == "nt" else 0
)
PPTX_WINDOWS_REPARSE_POINT_ATTRIBUTE = 0x000400
PPTX_WINDOWS_CLOUD_REPARSE_TAGS = frozenset(
    0x9000001A + (suffix << 12) for suffix in range(16)
)

_METADATA_SCHEMA_VERSION = 1
_METADATA_FAILURE_KINDS = frozenset(
    {"io", "missing", "not_regular", "root_escape", "symlink_or_reparse"}
)

_SHA256_RE = re.compile(r"^[0-9a-f]{64}$")
_OPC_ASCII_CASE_FOLD = str.maketrans(
    "ABCDEFGHIJKLMNOPQRSTUVWXYZ",
    "abcdefghijklmnopqrstuvwxyz",
)
_OPC_UNRESERVED_BYTES = frozenset(
    b"ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789-._~"
)
_OPC_HEX_DIGITS = frozenset("0123456789ABCDEFabcdef")
_OPC_RELATIONSHIP_TAG = (
    "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
)
_OPC_CONTENT_TYPE_DEFAULT_TAG = (
    "{http://schemas.openxmlformats.org/package/2006/content-types}Default"
)
_OPC_CONTENT_TYPE_OVERRIDE_TAG = (
    "{http://schemas.openxmlformats.org/package/2006/content-types}Override"
)
_PRESENTATION_SLIDE_ID_TAG = (
    "{http://schemas.openxmlformats.org/presentationml/2006/main}sldId"
)
_GRAPHIC_DATA_URI_TABLE = "http://schemas.openxmlformats.org/drawingml/2006/table"
_GRAPHIC_DATA_URI_CHART = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_GRAPHIC_DATA_URI_OLE = "http://schemas.openxmlformats.org/presentationml/2006/ole"
_AUDIT_FIELDS = frozenset(
    {
        "schema_version",
        "extraction_schema_version",
        "extraction_pipeline_version",
        "source_pptx_sha256",
        "source_pptx_size_bytes",
        "slide_count",
        "render_required_slide_numbers",
        "render_required_reasons",
        "extraction_receipt_sha256",
        "rendered_page_inspection",
    }
)
_EXTRACTION_FIELDS = frozenset(
    {
        "schema_version",
        "pipeline_version",
        "input_fingerprint",
        "pptx_path",
        "slide_count",
        "slide_width_inches",
        "slide_height_inches",
        "aspect_ratio",
        "corrupt_assets",
        "archive_recovery",
        "template_layouts",
        "per_slide_visual",
        "global_design",
        "native_deck_audit",
        "native_timing_summary",
    }
)
_TEMPLATE_LAYOUT_FIELDS = frozenset({"index", "master_index", "name", "placeholders"})
_TEMPLATE_PLACEHOLDER_FIELDS = frozenset({"idx", "type"})
_SLIDE_VISUAL_FIELDS = frozenset(
    {
        "slide_number",
        "slide_part_name",
        "background_color_hex",
        "background_type",
        "background_asset_status",
        "background_part_name",
        "background_asset_sha256",
        "layout_name",
        "shape_count",
        "shape_count_recursive",
        "has_text_frame_shapes",
        "has_extracted_text",
        "has_image",
        "image_area_ratio",
        "text_extraction_confidence",
        "text_content_preview",
        "ocr_text",
        "text_extraction_method",
        "text_channels",
        "unsupported_content",
        "has_unsupported_content",
        "render_required",
        "render_required_reasons",
        "footer_text",
        "has_speaker_notes",
        "native_timing",
        "shapes_summary",
    }
)
_SHAPE_REQUIRED_FIELDS = frozenset(
    {
        "name",
        "shape_type",
        "has_text_frame",
        "is_picture",
        "is_graphic_frame",
        "graphic_frame_type",
        "graphic_data_uri",
        "left",
        "top",
        "width",
        "height",
        "shape_path",
        "group_depth",
    }
)
_SHAPE_OPTIONAL_FIELDS = frozenset(
    {
        "text_preview",
        "font_name",
        "font_size",
        "font_color",
        "bold",
        "italic",
        "fill_color",
        "line_color",
        "line_width",
        "auto_shape_type",
        "picture_asset_status",
        "picture_part_name",
        "picture_asset_sha256",
        "table_rows",
        "table_columns",
        "table_text_preview",
        "table_fonts",
    }
)
_TEXT_CHANNEL_BASE_FIELDS = frozenset(
    {"channel", "text", "confidence", "status", "provenance"}
)
_TEXT_CHANNEL_OCR_FIELDS = frozenset(
    {
        *_TEXT_CHANNEL_BASE_FIELDS,
        "result_confidence",
        "attempted",
        "engine",
        "engine_version",
        "reason",
        "ocr_receipts",
    }
)
_OCR_RECEIPT_FIELDS = frozenset(
    {
        "attempted",
        "engine",
        "engine_version",
        "result_status",
        "result_confidence",
        "error",
        "part_name",
        "asset_sha256",
        "shape_path",
        "recovered_text",
        "trustworthy_text",
    }
)
_UNSUPPORTED_CONTENT_FIELDS = frozenset(
    {"content_type", "shape_name", "shape_path", "reason", "render_required"}
)
_TIMING_ANIMATION_FIELDS = (
    "general",
    "color",
    "effect",
    "motion",
    "rotation",
    "scale",
)
_TIMING_MEDIA_FIELDS = ("audio", "video")
_TIMING_BUILD_FIELDS = ("paragraph", "diagram", "ole_chart", "graphic")
_TIMING_PROVENANCE = {
    "source": "pptx_package_xml",
    "measurement": "raw_ooxml_element_counts",
    "observed_playback": False,
}
_NATIVE_TIMING_FIELDS = frozenset(
    {
        "timing_element_present",
        "timing_element_count",
        "transition_count",
        "set_action_count",
        "visibility_set_action_count",
        "animation_behavior_counts",
        "media_timing_counts",
        "build_list_present",
        "build_list_count",
        "build_entry_counts",
        "has_animation_behaviors",
        "has_media_timing",
        "has_build_entries",
        "provenance",
    }
)
_NATIVE_TIMING_SUMMARY_COUNT_FIELDS = (
    "slides_with_timing_elements",
    "slides_with_transitions",
    "slides_with_animation_behaviors",
    "slides_with_media_timing",
    "slides_with_build_lists",
    "slides_with_build_entries",
    "timing_element_count",
    "transition_count",
    "set_action_count",
    "visibility_set_action_count",
    "build_list_count",
)
_NATIVE_TIMING_SUMMARY_FIELDS = frozenset(
    {
        *_NATIVE_TIMING_SUMMARY_COUNT_FIELDS,
        "animation_behavior_counts",
        "media_timing_counts",
        "build_entry_counts",
        "provenance",
    }
)
_RENDER_RECEIPT_FIELDS = frozenset(
    {
        "schema_version",
        "source_pptx_sha256",
        "rendered_pdf_sha256",
        "rendered_pdf_size_bytes",
        "rendered_page_count",
        "inspected_page_ranges",
        "inspected_required_slide_numbers",
        "complete",
        "binding_sha256",
    }
)
_RECOVERY_IMAGE_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)
_PptxCacheKey = tuple[str, int, int, int, int, int, int, int, int, str]


class PptxEvidenceError(ValueError):
    """A PPTX artifact or evidence receipt violates its closed contract."""

    def __init__(
        self,
        message: str,
        *,
        reason_code: str = "pptx_evidence_invalid",
        details: Mapping[str, object] | None = None,
    ) -> None:
        super().__init__(message)
        self.reason_code = reason_code
        self.details = dict(details or {})


@dataclass(frozen=True)
class PptxArtifactProbe:
    """One exact readable deck generation plus any loss-recovery records."""

    slide_count: int
    source_sha256: str
    source_size_bytes: int
    archive_recovery: tuple[dict[str, object], ...]


@dataclass(frozen=True)
class _MetadataReceipt:
    generation: FileGeneration
    root_generation: FileGeneration | None
    reparse_tag: int | None


_PPTX_ARTIFACT_PROBE_CACHE: dict[
    _PptxCacheKey,
    PptxArtifactProbe | tuple[str, str, dict[str, object]],
] = {}
_PPTX_NATIVE_AUDIT_CACHE: dict[
    _PptxCacheKey,
    dict[str, object] | tuple[str, str, dict[str, object]],
] = {}


def sha256_bytes(blob: bytes) -> str:
    """Return the lowercase SHA-256 digest for exact bytes."""
    return hashlib.sha256(blob).hexdigest()


def _canonical_sha256(value: Mapping[str, object]) -> str:
    payload = json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    return sha256_bytes(payload)


def _require_sha256(value: object, label: str) -> str:
    if not isinstance(value, str) or _SHA256_RE.fullmatch(value) is None:
        raise PptxEvidenceError(
            f"{label} must be a lowercase 64-character SHA-256 digest"
        )
    return value


def _file_generation(
    stat_result: os.stat_result,
) -> tuple[int, int, int, int, int, int, int]:
    """Return fields that identify one stable regular-file generation."""
    return (
        stat_result.st_dev,
        stat_result.st_ino,
        stat_result.st_size,
        stat_result.st_mtime_ns,
        stat_result.st_ctime_ns,
        int(getattr(stat_result, "st_flags", 0)),
        int(getattr(stat_result, "st_file_attributes", 0)),
    )


def snapshot_regular_file(path: str | Path, *, label: str) -> bytes:
    """Read one exact non-symlink regular-file generation into memory."""
    artifact = Path(path)
    try:
        initial = artifact.lstat()
    except OSError as exc:
        raise PptxEvidenceError(
            f"{label} is unavailable at {artifact}: {exc}",
            reason_code="pptx_artifact_unavailable",
            details={"exception_type": type(exc).__name__},
        ) from exc
    if stat_module.S_ISLNK(initial.st_mode) or not stat_module.S_ISREG(initial.st_mode):
        raise PptxEvidenceError(
            f"{label} must be a non-symlink regular file: {artifact}",
            reason_code="pptx_artifact_unavailable",
        )
    generation = _file_generation(initial)
    try:
        with artifact.open("rb") as stream:
            opened = os.fstat(stream.fileno())
            if _file_generation(opened) != generation:
                raise PptxEvidenceError(
                    f"{label} changed while opening: {artifact}",
                    reason_code="pptx_artifact_changed",
                )
            blob = stream.read()
            after_read = os.fstat(stream.fileno())
    except PptxEvidenceError:
        raise
    except OSError as exc:
        raise PptxEvidenceError(
            f"cannot read {label} at {artifact}: {exc}",
            reason_code="pptx_artifact_unavailable",
            details={"exception_type": type(exc).__name__},
        ) from exc
    try:
        current = artifact.lstat()
    except OSError as exc:
        raise PptxEvidenceError(
            f"{label} changed while it was read at {artifact}: {exc}",
            reason_code="pptx_artifact_changed",
            details={"exception_type": type(exc).__name__},
        ) from exc
    if (
        _file_generation(after_read) != generation
        or _file_generation(current) != generation
        or len(blob) != initial.st_size
    ):
        raise PptxEvidenceError(
            f"{label} changed while reading: {artifact}",
            reason_code="pptx_artifact_changed",
        )
    return blob


def _is_embedded_media_member(member_name: str) -> bool:
    normalized = member_name.replace("\\", "/").lstrip("/")
    return normalized.startswith("ppt/media/") and not normalized.endswith("/")


def _opc_member_name_key(member_name: str) -> str:
    """Return the ASCII-case-insensitive OPC identity for one ZIP member."""
    is_directory = member_name.endswith("/")
    part_name = member_name[:-1] if is_directory else member_name
    if (
        not part_name
        or part_name.startswith("/")
        or "\\" in part_name
        or "?" in part_name
        or "#" in part_name
    ):
        raise PptxEvidenceError(
            "PPTX archive contains an invalid package part name; "
            "restore or re-export the source deck",
            reason_code="pptx_invalid_container",
        )
    segments = part_name.split("/")
    if any(not segment or segment.endswith(".") for segment in segments):
        raise PptxEvidenceError(
            "PPTX archive contains an invalid package part-name segment; "
            "restore or re-export the source deck",
            reason_code="pptx_invalid_container",
        )

    canonical: list[str] = []
    cursor = 0
    while cursor < len(part_name):
        character = part_name[cursor]
        if character != "%":
            canonical.append(character)
            cursor += 1
            continue
        if (
            cursor + 2 >= len(part_name)
            or part_name[cursor + 1] not in _OPC_HEX_DIGITS
            or part_name[cursor + 2] not in _OPC_HEX_DIGITS
        ):
            raise PptxEvidenceError(
                "PPTX archive contains an invalid package part-name escape; "
                "restore or re-export the source deck",
                reason_code="pptx_invalid_container",
            )
        escaped_byte = int(part_name[cursor + 1 : cursor + 3], 16)
        if (
            escaped_byte in _OPC_UNRESERVED_BYTES
            or escaped_byte in {0x2F, 0x5C}
            or escaped_byte < 0x20
            or escaped_byte == 0x7F
        ):
            raise PptxEvidenceError(
                "PPTX archive contains a forbidden package part-name escape; "
                "restore or re-export the source deck",
                reason_code="pptx_invalid_container",
            )
        canonical.append(f"%{escaped_byte:02X}")
        cursor += 3

    key = "".join(canonical).translate(_OPC_ASCII_CASE_FOLD)
    return f"{key}/" if is_directory else key


def _validate_relationship_part(stream) -> None:
    """Reject missing or duplicate relationship IDs before package parsing."""
    relationship_ids: set[str] = set()
    try:
        elements = ElementTree.iterparse(stream, events=("end",))
        for _event, element in elements:
            if element.tag != _OPC_RELATIONSHIP_TAG:
                element.clear()
                continue
            relationship_id = element.attrib.get("Id")
            if not relationship_id or relationship_id in relationship_ids:
                raise PptxEvidenceError(
                    "PPTX archive contains a relationship part with missing or "
                    "duplicate IDs; restore or re-export the source deck",
                    reason_code="pptx_invalid_container",
                )
            relationship_ids.add(relationship_id)
            if len(relationship_ids) > PPTX_ARCHIVE_MAX_MEMBERS:
                raise PptxEvidenceError(
                    "PPTX relationship part exceeds the bounded relationship contract",
                    reason_code="pptx_probe_resource_unavailable",
                )
            element.clear()
    except ElementTree.ParseError as exc:
        raise PptxEvidenceError(
            "PPTX archive contains a malformed relationship part; "
            "restore or re-export the source deck",
            reason_code="pptx_invalid_container",
        ) from exc


def _validate_content_types_part(stream) -> None:
    """Reject ambiguous defaults/overrides before content-type dispatch."""
    default_extensions: set[str] = set()
    override_part_names: set[str] = set()
    record_count = 0
    try:
        elements = ElementTree.iterparse(stream, events=("end",))
        for _event, element in elements:
            if element.tag == _OPC_CONTENT_TYPE_DEFAULT_TAG:
                extension = element.attrib.get("Extension")
                content_type = element.attrib.get("ContentType")
                if (
                    not extension
                    or not content_type
                    or any(character in extension for character in "/\\.%?#")
                ):
                    raise PptxEvidenceError(
                        "PPTX archive contains an invalid content-type default; "
                        "restore or re-export the source deck",
                        reason_code="pptx_invalid_container",
                    )
                key = extension.translate(_OPC_ASCII_CASE_FOLD)
                if key in default_extensions:
                    raise PptxEvidenceError(
                        "PPTX archive contains duplicate case-equivalent "
                        "content-type defaults; restore or re-export the source deck",
                        reason_code="pptx_invalid_container",
                    )
                default_extensions.add(key)
                record_count += 1
            elif element.tag == _OPC_CONTENT_TYPE_OVERRIDE_TAG:
                part_name = element.attrib.get("PartName")
                content_type = element.attrib.get("ContentType")
                if (
                    not part_name
                    or not part_name.startswith("/")
                    or part_name.startswith("//")
                    or part_name.endswith("/")
                    or not content_type
                ):
                    raise PptxEvidenceError(
                        "PPTX archive contains an invalid content-type override; "
                        "restore or re-export the source deck",
                        reason_code="pptx_invalid_container",
                    )
                key = _opc_member_name_key(part_name[1:])
                if key in override_part_names:
                    raise PptxEvidenceError(
                        "PPTX archive contains duplicate or equivalent content-type "
                        "overrides; restore or re-export the source deck",
                        reason_code="pptx_invalid_container",
                    )
                override_part_names.add(key)
                record_count += 1
            element.clear()
            if record_count > PPTX_ARCHIVE_MAX_MEMBERS:
                raise PptxEvidenceError(
                    "PPTX content-type manifest exceeds the bounded record contract",
                    reason_code="pptx_probe_resource_unavailable",
                )
    except ElementTree.ParseError as exc:
        raise PptxEvidenceError(
            "PPTX archive contains a malformed content-type manifest; "
            "restore or re-export the source deck",
            reason_code="pptx_invalid_container",
        ) from exc


def _validate_presentation_slide_ids(stream) -> None:
    """Reject duplicate slide identities/references before slide enumeration."""
    numeric_ids: set[str] = set()
    relationship_ids: set[str] = set()
    relationship_attribute = (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    try:
        elements = ElementTree.iterparse(stream, events=("end",))
        for _event, element in elements:
            if element.tag != _PRESENTATION_SLIDE_ID_TAG:
                element.clear()
                continue
            numeric_id = element.attrib.get("id")
            relationship_id = element.attrib.get(relationship_attribute)
            if (
                not numeric_id
                or not relationship_id
                or numeric_id in numeric_ids
                or relationship_id in relationship_ids
            ):
                raise PptxEvidenceError(
                    "PPTX presentation contains missing or duplicate slide "
                    "identities; restore or re-export the source deck",
                    reason_code="pptx_invalid_container",
                )
            numeric_ids.add(numeric_id)
            relationship_ids.add(relationship_id)
            if len(numeric_ids) > PPTX_ARCHIVE_MAX_MEMBERS:
                raise PptxEvidenceError(
                    "PPTX presentation exceeds the bounded slide-identity contract",
                    reason_code="pptx_probe_resource_unavailable",
                )
            element.clear()
    except ElementTree.ParseError as exc:
        raise PptxEvidenceError(
            "PPTX archive contains a malformed presentation part; "
            "restore or re-export the source deck",
            reason_code="pptx_invalid_container",
        ) from exc


def _corrupt_zip_members(package_blob: bytes) -> list[str]:
    """Validate every ZIP member and return those whose payload is corrupt."""
    corrupt: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(package_blob)) as archive:
            members = archive.infolist()
            if len(members) > PPTX_ARCHIVE_MAX_MEMBERS:
                raise PptxEvidenceError(
                    "PPTX archive exceeds the bounded member-count contract",
                    reason_code="pptx_probe_resource_unavailable",
                )
            expanded_bytes = 0
            member_names: set[str] = set()
            part_names: list[str] = []
            for member in members:
                if (
                    not member.filename
                    or "\x00" in member.filename
                    or len(member.filename) > 2048
                ):
                    raise PptxEvidenceError(
                        "PPTX archive member name exceeds the bounded evidence contract",
                        reason_code="pptx_probe_resource_unavailable",
                    )
                member_name_key = _opc_member_name_key(member.filename)
                if member_name_key in member_names:
                    raise PptxEvidenceError(
                        "PPTX archive contains duplicate or case-equivalent "
                        "package part names; "
                        "restore or re-export the source deck",
                        reason_code="pptx_invalid_container",
                    )
                member_names.add(member_name_key)
                if not member.filename.endswith("/"):
                    part_names.append(member_name_key)
                if member.file_size > PPTX_ARCHIVE_MAX_MEMBER_BYTES:
                    raise PptxEvidenceError(
                        "PPTX archive member exceeds the bounded expansion contract",
                        reason_code="pptx_probe_resource_unavailable",
                    )
                expanded_bytes += member.file_size
                if expanded_bytes > PPTX_ARCHIVE_MAX_EXPANDED_BYTES:
                    raise PptxEvidenceError(
                        "PPTX archive exceeds the bounded expansion contract",
                        reason_code="pptx_probe_resource_unavailable",
                    )

            sorted_part_names = sorted(part_names)
            for part_name in sorted_part_names:
                segment_prefix = f"{part_name}/"
                candidate_index = bisect_left(sorted_part_names, segment_prefix)
                if candidate_index < len(sorted_part_names) and sorted_part_names[
                    candidate_index
                ].startswith(segment_prefix):
                    raise PptxEvidenceError(
                        "PPTX archive contains a package part whose name is a "
                        "segment prefix of another part; restore or re-export "
                        "the source deck",
                        reason_code="pptx_invalid_container",
                    )

            for member in members:
                try:
                    with archive.open(member) as stream:
                        member_name_key = _opc_member_name_key(member.filename)
                        if member_name_key == "[content_types].xml":
                            _validate_content_types_part(stream)
                        elif member_name_key == "ppt/presentation.xml":
                            _validate_presentation_slide_ids(stream)
                        elif member_name_key.endswith(".rels"):
                            _validate_relationship_part(stream)
                        else:
                            while stream.read(1024 * 1024):
                                pass
                except (zipfile.BadZipFile, ZlibError):
                    corrupt.append(member.filename)
                    if len(corrupt) > PPTX_ARTIFACT_PROBE_MAX_RECOVERY_RECORDS:
                        raise PptxEvidenceError(
                            "PPTX archive has more corrupt members than the "
                            "bounded recovery contract permits",
                            reason_code="pptx_probe_result_oversized",
                        )
    except zipfile.BadZipFile as exc:
        raise PptxEvidenceError(
            "invalid PPTX ZIP container; restore or re-export the source deck",
            reason_code="pptx_invalid_container",
        ) from exc
    return corrupt


def _parse_presentation(package_blob: bytes, *, recovered: bool) -> Any:
    try:
        from lxml.etree import XMLSyntaxError
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError as exc:
        raise PptxEvidenceError(
            "PPTX evidence requires the declared python-pptx and lxml runtime "
            "dependencies; install the speaker-toolkit project dependencies",
            reason_code="pptx_dependency_unavailable",
        ) from exc
    try:
        return Presentation(io.BytesIO(package_blob))
    except (
        zipfile.BadZipFile,
        ZlibError,
        PackageNotFoundError,
        XMLSyntaxError,
        OSError,
        ValueError,
        KeyError,
    ) as exc:
        prefix = "recovered PPTX package" if recovered else "PPTX package"
        raise PptxEvidenceError(
            f"{prefix} cannot be parsed; restore or re-export the source deck: {exc}",
            reason_code="pptx_parse_failure",
        ) from exc


def presentation_with_media_recovery(
    package_blob: bytes,
) -> tuple[Any, list[dict[str, object]]]:
    """Open a deck, replacing only CRC-damaged embedded media in memory."""
    corrupt_names = _corrupt_zip_members(package_blob)
    if not corrupt_names:
        return _parse_presentation(package_blob, recovered=False), []

    structural = [name for name in corrupt_names if not _is_embedded_media_member(name)]
    if structural:
        raise PptxEvidenceError(
            "corrupt structural PPTX member(s) are not recoverable: "
            f"{', '.join(sorted(structural))}; restore or re-export the source deck",
            reason_code="pptx_structural_damage",
            details={"part_names": sorted(structural)},
        )

    recovered_package = io.BytesIO()
    corrupt_set = set(corrupt_names)
    try:
        with (
            zipfile.ZipFile(io.BytesIO(package_blob)) as source,
            zipfile.ZipFile(recovered_package, "w") as destination,
        ):
            for member in source.infolist():
                # ZipFile.open(..., "w") mutates the supplied ZipInfo with the
                # destination header offset. Keep the source archive's object
                # intact so subsequent source reads remain correctly bound.
                destination_member = copy.copy(member)
                with destination.open(destination_member, "w") as output:
                    if member.filename in corrupt_set:
                        output.write(_RECOVERY_IMAGE_BYTES)
                        continue
                    with source.open(member) as input_stream:
                        while True:
                            chunk = input_stream.read(1024 * 1024)
                            if not chunk:
                                break
                            output.write(chunk)
    except (zipfile.BadZipFile, ZlibError, OSError) as exc:
        raise PptxEvidenceError(
            "could not recover corrupt PPTX media; restore or re-export the source deck",
            reason_code="pptx_recovery_failure",
        ) from exc

    recovery = [
        {
            "schema_version": ARCHIVE_RECOVERY_SCHEMA_VERSION,
            "part_name": name,
            "member_kind": "embedded_media",
            "error_type": "crc_mismatch",
            "status": "recovered_with_placeholder_asset",
            "content_replaced": True,
            "replacement_sha256": sha256_bytes(_RECOVERY_IMAGE_BYTES),
        }
        for name in sorted(corrupt_names)
    ]
    return _parse_presentation(recovered_package.getvalue(), recovered=True), recovery


def _probe_pptx_artifact_in_process(path: str | Path) -> PptxArtifactProbe:
    """Perform the expensive deck probe inside the bounded worker only."""
    package_blob = snapshot_regular_file(path, label="PPTX artifact")
    presentation, recovery = presentation_with_media_recovery(package_blob)
    slide_count = len(presentation.slides)
    if slide_count < 1:
        raise PptxEvidenceError(
            f"PPTX artifact has no slides: {path}",
            reason_code="pptx_no_slides",
        )
    return PptxArtifactProbe(
        slide_count=slide_count,
        source_sha256=sha256_bytes(package_blob),
        source_size_bytes=len(package_blob),
        archive_recovery=tuple(dict(item) for item in recovery),
    )


_CHILD_PROBE_REASON_CODES = frozenset(
    {
        "pptx_artifact_changed",
        "pptx_artifact_unavailable",
        "pptx_archive_recovery_required",
        "pptx_dependency_unavailable",
        "pptx_evidence_invalid",
        "pptx_invalid_container",
        "pptx_no_slides",
        "pptx_parse_failure",
        "pptx_probe_exception",
        "pptx_probe_resource_unavailable",
        "pptx_probe_result_oversized",
        "pptx_recovery_failure",
        "pptx_structural_damage",
    }
)

_ARCHIVE_INTEGRITY_CONFIRMATION_REASON_CODES = frozenset(
    {
        "pptx_artifact_unavailable",
        "pptx_archive_recovery_required",
        "pptx_invalid_container",
        "pptx_parse_failure",
        "pptx_probe_result_oversized",
        "pptx_recovery_failure",
        "pptx_structural_damage",
    }
)


def _probe_failure(
    reason_code: str,
    *,
    details: Mapping[str, object] | None = None,
) -> PptxEvidenceError:
    """Create one bounded, actionable parent-side probe error."""
    normalized = dict(details or {})
    timeout_seconds = normalized.get(
        "timeout_seconds", PPTX_ARTIFACT_PROBE_TIMEOUT_SECONDS
    )
    part_names = normalized.get("part_names")
    joined_parts = (
        ", ".join(str(name) for name in part_names)
        if isinstance(part_names, list)
        else ""
    )
    failure_kind = normalized.get("failure_kind")
    unavailable_message = {
        "root_escape": "PPTX artifact is outside its trusted root",
        "symlink_or_reparse": (
            "PPTX artifact must not traverse a symbolic link or reparse point"
        ),
    }.get(
        failure_kind if isinstance(failure_kind, str) else "",
        "PPTX artifact is unavailable; restore the file and retry",
    )
    messages = {
        "pptx_artifact_unavailable": unavailable_message,
        "pptx_archive_recovery_required": (
            "PPTX artifact required placeholder archive recovery and cannot "
            "authorize a fresh native-deck audit"
            + (f": {joined_parts}" if joined_parts else "")
            + "; restore or re-export the source deck"
        ),
        "pptx_dependency_unavailable": (
            "PPTX evidence requires its declared runtime dependencies; install "
            "the speaker-toolkit project dependencies"
        ),
        "pptx_evidence_invalid": (
            "PPTX artifact violates the evidence contract; restore or re-export "
            "the source deck"
        ),
        "pptx_invalid_container": (
            "invalid PPTX ZIP container; restore or re-export the source deck"
        ),
        "pptx_no_slides": "PPTX artifact has no slides; re-export the source deck",
        "pptx_parse_failure": (
            "PPTX package cannot be parsed; restore or re-export the source deck"
        ),
        "pptx_probe_exception": (
            "PPTX artifact probe failed unexpectedly inside its bounded worker"
        ),
        "pptx_probe_resource_unavailable": (
            "PPTX evidence exceeded a configured worker resource limit"
        ),
        "pptx_probe_request_oversized": (
            "bounded PPTX evidence worker request exceeded its input contract"
        ),
        "pptx_probe_result_oversized": (
            "bounded PPTX evidence worker result exceeded its output contract"
        ),
        "pptx_probe_monitor_unavailable": (
            "bounded PPTX evidence worker could not inspect its process tree"
        ),
        "pptx_probe_monitor_identity_changed": (
            "bounded PPTX evidence worker process identity changed during inspection"
        ),
        "pptx_probe_containment_unavailable": (
            "bounded PPTX evidence worker could not establish or preserve "
            "process-tree containment"
        ),
        "pptx_recovery_failure": (
            "could not recover corrupt PPTX media; restore or re-export the source deck"
        ),
        "pptx_structural_damage": (
            "corrupt structural PPTX member(s) are not recoverable"
            + (f": {joined_parts}" if joined_parts else "")
            + "; restore or re-export the source deck"
        ),
        "pptx_probe_timeout": (
            "bounded PPTX evidence operation timed out after "
            f"{timeout_seconds} seconds; use an independent "
            "healthy evidence lane or repair/re-export the deck"
        ),
        "pptx_probe_start_failure": "could not start the bounded PPTX evidence worker",
        "pptx_probe_crash": "bounded PPTX evidence worker terminated unexpectedly",
        "pptx_probe_malformed_result": (
            "bounded PPTX evidence worker returned an invalid protocol result"
        ),
        "pptx_probe_materialization_changed": (
            "PPTX artifact produced inconsistent bounded reads while cloud "
            "materialization was changing; retry after the file is fully local"
        ),
        "pptx_artifact_changed": (
            "PPTX artifact changed during bounded evidence inspection"
        ),
        "pptx_cloud_placeholder_unavailable": (
            "PPTX artifact is an offline cloud placeholder; download the file "
            "locally before using native-deck evidence"
        ),
    }
    return PptxEvidenceError(
        messages.get(reason_code, "PPTX artifact is unavailable"),
        reason_code=reason_code,
        details=normalized,
    )


def _probe_child_failure_details(exc: PptxEvidenceError) -> dict[str, object]:
    """Copy only closed, bounded diagnostic fields into the child result."""
    details: dict[str, object] = {}
    raw_names = exc.details.get("part_names")
    if isinstance(raw_names, list):
        names = [
            name
            for name in raw_names
            if isinstance(name, str) and 0 < len(name) <= 2048
        ]
        if len(names) == len(raw_names) and len(names) <= 64:
            details["part_names"] = names
    exception_type = exc.details.get("exception_type")
    if (
        isinstance(exception_type, str)
        and exception_type
        and len(exception_type) <= 128
    ):
        details["exception_type"] = exception_type
    return details


def _pptx_probe_child(path: str | Path) -> dict[str, object]:
    """Return the closed payload written by the isolated probe worker."""
    try:
        probe = _probe_pptx_artifact_in_process(path)
    except MemoryError:
        return {
            "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": "pptx_probe_resource_unavailable",
            "details": {},
        }
    except PptxEvidenceError as exc:
        reason_code = (
            exc.reason_code
            if exc.reason_code in _CHILD_PROBE_REASON_CODES
            else "pptx_evidence_invalid"
        )
        return {
            "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": reason_code,
            "details": _probe_child_failure_details(exc),
        }
    return {
        "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
        "status": "available",
        "slide_count": probe.slide_count,
        "source_sha256": probe.source_sha256,
        "source_size_bytes": probe.source_size_bytes,
        "archive_recovery": [dict(item) for item in probe.archive_recovery],
    }


def _write_pptx_probe_result(
    result_file: BinaryIO,
    payload: Mapping[str, object],
) -> None:
    """Replace the private result through the child's retained descriptor."""
    rendered = (json.dumps(payload, sort_keys=True) + "\n").encode("utf-8")
    if len(rendered) > PPTX_ARTIFACT_PROBE_MAX_RESULT_BYTES:
        rendered = (
            json.dumps(
                {
                    "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
                    "status": "unavailable",
                    "reason_code": "pptx_probe_result_oversized",
                    "details": {},
                },
                sort_keys=True,
            )
            + "\n"
        ).encode("utf-8")
    result_file.seek(0)
    result_file.truncate()
    remaining = memoryview(rendered)
    while remaining:
        written = result_file.write(remaining)
        if written is None or written <= 0:
            raise OSError("PPTX probe result write made no progress")
        remaining = remaining[written:]
    result_file.flush()


def _unique_json_object(pairs: list[tuple[str, object]]) -> dict[str, object]:
    result: dict[str, object] = {}
    for key, value in pairs:
        if key in result:
            raise ValueError("duplicate JSON object key")
        result[key] = value
    return result


def _reject_nonstandard_json_constant(_value: str) -> object:
    raise ValueError("non-standard JSON number")


def _validated_recovery_records(value: object) -> tuple[dict[str, object], ...]:
    if (
        not isinstance(value, list)
        or len(value) > PPTX_ARTIFACT_PROBE_MAX_RECOVERY_RECORDS
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    expected_fields = {
        "schema_version",
        "part_name",
        "member_kind",
        "error_type",
        "status",
        "content_replaced",
        "replacement_sha256",
    }
    records: list[dict[str, object]] = []
    for raw_record in value:
        if not isinstance(raw_record, dict) or set(raw_record) != expected_fields:
            raise _probe_failure("pptx_probe_malformed_result")
        part_name = raw_record.get("part_name")
        if (
            raw_record.get("schema_version") != ARCHIVE_RECOVERY_SCHEMA_VERSION
            or not isinstance(part_name, str)
            or not _is_embedded_media_member(part_name)
            or len(part_name) > 2048
            or not _valid_opc_part_name(part_name)
            or raw_record.get("member_kind") != "embedded_media"
            or raw_record.get("error_type") != "crc_mismatch"
            or raw_record.get("status") != "recovered_with_placeholder_asset"
            or raw_record.get("content_replaced") is not True
            or raw_record.get("replacement_sha256")
            != sha256_bytes(_RECOVERY_IMAGE_BYTES)
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        records.append(dict(raw_record))
    part_names = [str(record["part_name"]) for record in records]
    if part_names != sorted(set(part_names)):
        raise _probe_failure("pptx_probe_malformed_result")
    return tuple(records)


def _decode_pptx_probe_result(raw: bytes) -> PptxArtifactProbe:
    if len(raw) > PPTX_ARTIFACT_PROBE_MAX_RESULT_BYTES or not raw.strip():
        raise _probe_failure("pptx_probe_malformed_result")
    try:
        decoded = raw.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    if len(decoded.splitlines()) != 1:
        raise _probe_failure("pptx_probe_malformed_result")
    try:
        payload = json.loads(
            decoded,
            object_pairs_hook=_unique_json_object,
            parse_constant=_reject_nonstandard_json_constant,
        )
    except (json.JSONDecodeError, RecursionError, ValueError) as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    if (
        not isinstance(payload, dict)
        or payload.get("schema_version") != PPTX_ARTIFACT_PROBE_SCHEMA_VERSION
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    if payload.get("status") == "unavailable":
        if set(payload) != {"schema_version", "status", "reason_code", "details"}:
            raise _probe_failure("pptx_probe_malformed_result")
        reason_code = payload.get("reason_code")
        details = payload.get("details")
        if (
            not isinstance(reason_code, str)
            or reason_code not in _CHILD_PROBE_REASON_CODES
            or not isinstance(details, dict)
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        if set(details) - {"part_names", "exception_type"}:
            raise _probe_failure("pptx_probe_malformed_result")
        raw_names = details.get("part_names")
        if raw_names is not None and (
            not isinstance(raw_names, list)
            or len(raw_names) > 64
            or any(
                not isinstance(name, str) or not name or len(name) > 2048
                for name in raw_names
            )
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        exception_type = details.get("exception_type")
        if exception_type is not None and (
            not isinstance(exception_type, str)
            or not exception_type
            or len(exception_type) > 128
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        raise _probe_failure(str(reason_code), details=details)
    expected_fields = {
        "schema_version",
        "status",
        "slide_count",
        "source_sha256",
        "source_size_bytes",
        "archive_recovery",
    }
    slide_count = payload.get("slide_count")
    source_size = payload.get("source_size_bytes")
    source_sha = payload.get("source_sha256")
    if (
        payload.get("status") != "available"
        or set(payload) != expected_fields
        or isinstance(slide_count, bool)
        or not isinstance(slide_count, int)
        or not 1 <= slide_count <= PPTX_ARCHIVE_MAX_MEMBERS
        or isinstance(source_size, bool)
        or not isinstance(source_size, int)
        or source_size < 1
        or not isinstance(source_sha, str)
        or _SHA256_RE.fullmatch(source_sha) is None
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    return PptxArtifactProbe(
        slide_count=slide_count,
        source_sha256=source_sha,
        source_size_bytes=source_size,
        archive_recovery=_validated_recovery_records(payload.get("archive_recovery")),
    )


def _read_pptx_probe_result(result_file: BinaryIO) -> PptxArtifactProbe:
    try:
        result_file.seek(0)
        raw = result_file.read(PPTX_ARTIFACT_PROBE_MAX_RESULT_BYTES + 1)
    except OSError as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    return _decode_pptx_probe_result(raw)


def _metadata_failure(
    failure_kind: str,
    *,
    exception_type: str | None = None,
) -> PptxEvidenceError:
    details: dict[str, object] = {"failure_kind": failure_kind}
    if exception_type is not None:
        details["exception_type"] = exception_type
    return _probe_failure("pptx_artifact_unavailable", details=details)


def _reparse_tag(snapshot: os.stat_result) -> int | None:
    attributes = int(getattr(snapshot, "st_file_attributes", 0))
    if not attributes & PPTX_WINDOWS_REPARSE_POINT_ATTRIBUTE:
        return None
    raw_tag = getattr(snapshot, "st_reparse_tag", None)
    return (
        int(raw_tag)
        if isinstance(raw_tag, int) and not isinstance(raw_tag, bool)
        else -1
    )


def _is_unsupported_reparse(
    snapshot: os.stat_result,
    *,
    allow_hydrated_cloud_file: bool,
) -> bool:
    tag = _reparse_tag(snapshot)
    if tag is None:
        return False
    return not (allow_hydrated_cloud_file and tag in PPTX_WINDOWS_CLOUD_REPARSE_TAGS)


def _metadata_generation_in_worker(
    path: Path,
    *,
    trusted_root: Path | None,
) -> tuple[FileGeneration, FileGeneration | None, int | None]:
    """Inspect one file only inside the bounded metadata worker."""
    if not path.is_absolute() or Path(os.path.abspath(path)) != path:
        raise SupervisorError("invalid_worker_request")

    target = path
    snapshot: os.stat_result | None = None
    root_generation: FileGeneration | None = None
    if trusted_root is not None:
        if (
            not trusted_root.is_absolute()
            or Path(os.path.abspath(trusted_root)) != trusted_root
        ):
            raise SupervisorError("invalid_worker_request")
        try:
            relative = path.relative_to(trusted_root)
        except ValueError as exc:
            raise _metadata_failure("root_escape") from exc
        if not relative.parts or any(
            part in {"", ".", ".."} for part in relative.parts
        ):
            raise _metadata_failure("root_escape")
        try:
            root_snapshot = trusted_root.lstat()
        except FileNotFoundError as exc:
            raise _metadata_failure(
                "missing", exception_type=type(exc).__name__
            ) from exc
        except (OSError, RuntimeError) as exc:
            raise _metadata_failure("io", exception_type=type(exc).__name__) from exc
        if (
            stat_module.S_ISLNK(root_snapshot.st_mode)
            or _is_unsupported_reparse(
                root_snapshot,
                allow_hydrated_cloud_file=False,
            )
            or not stat_module.S_ISDIR(root_snapshot.st_mode)
        ):
            raise _metadata_failure("root_escape")
        root_generation = FileGeneration.from_stat(root_snapshot)
        target = trusted_root
        for index, component in enumerate(relative.parts):
            target = target / component
            try:
                snapshot = target.lstat()
            except FileNotFoundError as exc:
                raise _metadata_failure(
                    "missing", exception_type=type(exc).__name__
                ) from exc
            except OSError as exc:
                raise _metadata_failure(
                    "io", exception_type=type(exc).__name__
                ) from exc
            is_leaf = index == len(relative.parts) - 1
            if stat_module.S_ISLNK(snapshot.st_mode) or _is_unsupported_reparse(
                snapshot,
                allow_hydrated_cloud_file=is_leaf,
            ):
                raise _metadata_failure("symlink_or_reparse")
            if not is_leaf and not stat_module.S_ISDIR(snapshot.st_mode):
                raise _metadata_failure("not_regular")
    else:
        try:
            snapshot = target.lstat()
        except FileNotFoundError as exc:
            raise _metadata_failure(
                "missing", exception_type=type(exc).__name__
            ) from exc
        except OSError as exc:
            raise _metadata_failure("io", exception_type=type(exc).__name__) from exc

    if snapshot is None:  # pragma: no cover - guarded by non-empty relative parts
        raise SupervisorError("invalid_worker_request")
    if stat_module.S_ISLNK(snapshot.st_mode) or _is_unsupported_reparse(
        snapshot,
        allow_hydrated_cloud_file=True,
    ):
        raise _metadata_failure("symlink_or_reparse")
    if not stat_module.S_ISREG(snapshot.st_mode):
        raise _metadata_failure("not_regular")
    return FileGeneration.from_stat(snapshot), root_generation, _reparse_tag(snapshot)


def _metadata_child(payload: Mapping[str, object]) -> dict[str, object]:
    if set(payload) != {"pptx_path", "trusted_root"}:
        raise SupervisorError("invalid_worker_request")
    path = _worker_bound_path(payload.get("pptx_path"))
    root_value = payload.get("trusted_root")
    if root_value is not None and not isinstance(root_value, str):
        raise SupervisorError("invalid_worker_request")
    trusted_root = _worker_bound_path(root_value) if root_value is not None else None
    try:
        generation, root_generation, reparse_tag = _metadata_generation_in_worker(
            path,
            trusted_root=trusted_root,
        )
    except PptxEvidenceError as exc:
        return {
            "schema_version": _METADATA_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": exc.reason_code,
            "details": dict(exc.details),
        }
    return {
        "schema_version": _METADATA_SCHEMA_VERSION,
        "status": "available",
        "generation": generation.to_dict(),
        "root_generation": (
            root_generation.to_dict() if root_generation is not None else None
        ),
        "reparse_tag": reparse_tag,
    }


def _decode_metadata_payload(payload: object) -> _MetadataReceipt:
    if not isinstance(payload, Mapping):
        raise _probe_failure("pptx_probe_malformed_result")
    status = payload.get("status")
    if payload.get("schema_version") != _METADATA_SCHEMA_VERSION:
        raise _probe_failure("pptx_probe_malformed_result")
    if status == "unavailable":
        if set(payload) != {"schema_version", "status", "reason_code", "details"}:
            raise _probe_failure("pptx_probe_malformed_result")
        details = payload.get("details")
        if (
            payload.get("reason_code") != "pptx_artifact_unavailable"
            or not isinstance(details, Mapping)
            or set(details) - {"failure_kind", "exception_type"}
            or not isinstance(details.get("failure_kind"), str)
            or details.get("failure_kind") not in _METADATA_FAILURE_KINDS
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        exception_type = details.get("exception_type")
        if exception_type is not None and (
            not isinstance(exception_type, str)
            or not exception_type
            or len(exception_type) > 128
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        raise _probe_failure(
            "pptx_artifact_unavailable",
            details=dict(details),
        )
    if status != "available" or set(payload) != {
        "schema_version",
        "status",
        "generation",
        "root_generation",
        "reparse_tag",
    }:
        raise _probe_failure("pptx_probe_malformed_result")
    raw_generation = payload.get("generation")
    if not isinstance(raw_generation, Mapping):
        raise _probe_failure("pptx_probe_malformed_result")
    try:
        generation = FileGeneration.from_dict(raw_generation)
    except (TypeError, ValueError) as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    raw_root_generation = payload.get("root_generation")
    if raw_root_generation is None:
        root_generation = None
    elif isinstance(raw_root_generation, Mapping):
        try:
            root_generation = FileGeneration.from_dict(raw_root_generation)
        except (TypeError, ValueError) as exc:
            raise _probe_failure("pptx_probe_malformed_result") from exc
        if not stat_module.S_ISDIR(root_generation.mode):
            raise _probe_failure("pptx_probe_malformed_result")
    else:
        raise _probe_failure("pptx_probe_malformed_result")
    reparse_tag = payload.get("reparse_tag")
    if reparse_tag is not None and (
        isinstance(reparse_tag, bool) or not isinstance(reparse_tag, int)
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    attributes = generation.file_attributes or 0
    has_reparse_attribute = bool(attributes & PPTX_WINDOWS_REPARSE_POINT_ATTRIBUTE)
    if (
        generation.size < 0
        or not stat_module.S_ISREG(generation.mode)
        or has_reparse_attribute != (reparse_tag is not None)
        or (
            reparse_tag is not None
            and (reparse_tag not in PPTX_WINDOWS_CLOUD_REPARSE_TAGS)
        )
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    return _MetadataReceipt(generation, root_generation, reparse_tag)


def _invoke_metadata_worker(
    command: list[str],
    payload: dict[str, object],
    sensitive_values: tuple[Path, ...],
    limits: SupervisorLimits,
) -> Any:
    """Narrow injection seam for the authenticated metadata invocation."""
    return run_authenticated_worker(
        command,
        PPTX_METADATA_OPERATION,
        {},
        cast(Any, payload),
        limits,
        sensitive_values=sensitive_values,
        schema_generation=PPTX_EXTRACTION_SCHEMA_VERSION,
        pipeline_generation=PPTX_EXTRACTION_PIPELINE_VERSION,
    )


def _run_bounded_metadata_worker(
    path: Path,
    *,
    trusted_root: Path | None = None,
    deadline_monotonic: float | None = None,
) -> _MetadataReceipt:
    artifact = Path(os.path.abspath(os.fspath(path)))
    root = (
        Path(os.path.abspath(os.fspath(trusted_root)))
        if trusted_root is not None
        else None
    )
    command = [
        sys.executable,
        os.fspath(Path(__file__).absolute()),
        PPTX_SUPERVISED_WORKER_FLAG,
    ]
    payload: dict[str, object] = {
        "pptx_path": os.fspath(artifact),
        "trusted_root": os.fspath(root) if root is not None else None,
    }
    sensitive = (artifact,) if root is None else (artifact, root)
    limits = _limits_before_deadline(PPTX_METADATA_LIMITS, deadline_monotonic)
    deadline_limited = limits.wall_seconds < PPTX_METADATA_LIMITS.wall_seconds
    try:
        result = _invoke_metadata_worker(command, payload, sensitive, limits)
    except SupervisorError as exc:
        if deadline_limited and exc.reason_code == "worker_timeout":
            raise PptxEvidenceError(
                "PPTX batch wall deadline expired during bounded metadata inspection",
                reason_code="pptx_batch_wall_limit",
            ) from exc
        raise _supervisor_probe_failure(
            exc, timeout_seconds=limits.wall_seconds
        ) from exc
    receipt = _decode_metadata_payload(result.payload)
    if (root is None) != (receipt.root_generation is None):
        raise _probe_failure("pptx_probe_malformed_result")
    return receipt


def _signal_name(signal_number: int) -> str:
    try:
        return signal.Signals(signal_number).name
    except ValueError:
        return f"SIG{signal_number}"


def _run_bounded_private_worker(
    path: Path,
    *,
    child_flag: str,
    decoder: Callable[[bytes], Any],
    trusted_root: Path | None = None,
    deadline_monotonic: float | None = None,
) -> Any:
    """Run one authenticated, process-tree-bounded PPTX worker."""
    operation_limits = {
        PPTX_ARTIFACT_PROBE_CHILD_FLAG: (
            PPTX_PROBE_OPERATION,
            PPTX_PROBE_LIMITS,
        ),
        PPTX_NATIVE_AUDIT_CHILD_FLAG: (
            PPTX_NATIVE_AUDIT_OPERATION,
            PPTX_NATIVE_AUDIT_LIMITS,
        ),
    }
    try:
        operation, base_limits = operation_limits[child_flag]
    except KeyError as exc:
        raise _probe_failure("pptx_probe_start_failure") from exc
    metadata = _run_bounded_metadata_worker(
        path,
        trusted_root=trusted_root,
        deadline_monotonic=deadline_monotonic,
    )
    generation = metadata.generation
    if trusted_root is not None and metadata.root_generation is None:
        raise _probe_failure("pptx_probe_malformed_result")
    expected_generations = {"pptx": generation}
    if metadata.root_generation is not None:
        expected_generations["pptx_root"] = metadata.root_generation
    if generation.size > PPTX_MAX_INPUT_BYTES:
        raise _probe_failure(
            "pptx_probe_resource_unavailable",
            details={"limit_bytes": PPTX_MAX_INPUT_BYTES},
        )
    limits = _limits_before_deadline(base_limits, deadline_monotonic)
    deadline_limited = limits.wall_seconds < base_limits.wall_seconds
    command = [
        sys.executable,
        os.fspath(Path(__file__).absolute()),
        PPTX_SUPERVISED_WORKER_FLAG,
    ]
    try:
        result = run_authenticated_worker(
            command,
            operation,
            expected_generations,
            {
                "pptx_path": os.fspath(path),
                "trusted_root": (
                    os.fspath(trusted_root) if trusted_root is not None else None
                ),
            },
            limits,
            sensitive_values=(path,),
            schema_generation=PPTX_EXTRACTION_SCHEMA_VERSION,
            pipeline_generation=PPTX_EXTRACTION_PIPELINE_VERSION,
        )
    except SupervisorError as exc:
        if deadline_limited and exc.reason_code == "worker_timeout":
            raise PptxEvidenceError(
                "PPTX batch wall deadline expired inside the bounded worker",
                reason_code="pptx_batch_wall_limit",
            ) from exc
        raise _supervisor_probe_failure(
            exc, timeout_seconds=limits.wall_seconds
        ) from exc
    current = _run_bounded_metadata_worker(
        path,
        trusted_root=trusted_root,
        deadline_monotonic=deadline_monotonic,
    )
    if current != metadata:
        raise _probe_failure("pptx_artifact_changed")
    try:
        raw = (
            json.dumps(
                result.payload,
                ensure_ascii=False,
                allow_nan=False,
                sort_keys=True,
                separators=(",", ":"),
            )
            + "\n"
        ).encode("utf-8")
    except (TypeError, ValueError, UnicodeError) as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    return decoder(raw)


def _supervised_file_generation(
    path: Path,
    *,
    label: str,
    trusted_root: Path | None = None,
    deadline_monotonic: float | None = None,
) -> FileGeneration:
    del label  # Failure text is deliberately path- and artifact-label-free.
    return _run_bounded_metadata_worker(
        path,
        trusted_root=trusted_root,
        deadline_monotonic=deadline_monotonic,
    ).generation


def _supervisor_probe_failure(
    exc: SupervisorError,
    *,
    timeout_seconds: float,
) -> PptxEvidenceError:
    reason = exc.reason_code
    if reason in {"worker_generation_changed", "worker_generation_binding_mismatch"}:
        return _probe_failure("pptx_artifact_changed")
    if reason == "worker_timeout":
        return _probe_failure(
            "pptx_probe_timeout",
            details={"timeout_seconds": timeout_seconds},
        )
    closed_details = {"supervisor_reason_code": reason}
    if reason == "worker_input_limit_exceeded":
        return _probe_failure(
            "pptx_probe_request_oversized",
            details=closed_details,
        )
    if reason == "worker_output_limit_exceeded":
        return _probe_failure(
            "pptx_probe_result_oversized",
            details=closed_details,
        )
    if reason == "worker_monitor_unavailable":
        public_reason = (
            "pptx_dependency_unavailable"
            if exc.details.get("dependency") == "psutil"
            else "pptx_probe_monitor_unavailable"
        )
        return _probe_failure(public_reason, details=closed_details)
    if reason == "worker_monitor_identity_changed":
        return _probe_failure(
            "pptx_probe_monitor_identity_changed",
            details=closed_details,
        )
    if reason in {
        "worker_containment_unavailable",
        "worker_process_tree_leak",
        "worker_cleanup_failed",
    }:
        return _probe_failure(
            "pptx_probe_containment_unavailable",
            details=closed_details,
        )
    if reason in {
        "worker_memory_limit_exceeded",
        "worker_process_limit_exceeded",
        "worker_diagnostic_limit_exceeded",
    }:
        return _probe_failure(
            "pptx_probe_resource_unavailable",
            details=closed_details,
        )
    if reason in {
        "worker_start_failed",
        "worker_pipe_setup_failed",
        "worker_exit_before_barrier",
        "worker_request_write_failed",
        "invalid_worker_command",
        "unsafe_worker_process_metadata",
    }:
        return _probe_failure(
            "pptx_probe_start_failure",
            details=closed_details,
        )
    if reason in {
        "worker_exit",
        "worker_diagnostic_read_failed",
        "worker_output_read_failed",
    }:
        return _probe_failure(
            "pptx_probe_crash",
            details=closed_details,
        )
    return _probe_failure(
        "pptx_probe_malformed_result",
        details=closed_details,
    )


def _confirm_archive_integrity_failure(
    artifact: Path,
    key: _PptxCacheKey,
    first_error: PptxEvidenceError,
    *,
    child_flag: str,
    decoder: Callable[[bytes], Any],
    trusted_root: Path | None = None,
    deadline_monotonic: float | None = None,
) -> None:
    """Require two identical bounded failures before treating damage as stable."""
    try:
        _run_bounded_private_worker(
            artifact,
            child_flag=child_flag,
            decoder=decoder,
            trusted_root=trusted_root,
            deadline_monotonic=deadline_monotonic,
        )
    except PptxEvidenceError as confirmation_error:
        if confirmation_error.reason_code == "pptx_batch_wall_limit":
            raise
        _confirmed_path, confirmed_key = _probe_file_identity(
            artifact,
            trusted_root=trusted_root,
            deadline_monotonic=deadline_monotonic,
        )
        if confirmed_key != key:
            raise _probe_failure("pptx_artifact_changed") from confirmation_error
        if (
            confirmation_error.reason_code == first_error.reason_code
            and confirmation_error.details == first_error.details
        ):
            return
        raise _probe_failure(
            "pptx_probe_materialization_changed"
        ) from confirmation_error
    raise _probe_failure("pptx_probe_materialization_changed") from first_error


def _run_bounded_pptx_probe(
    path: Path,
    *,
    trusted_root: Path | None = None,
    deadline_monotonic: float | None = None,
) -> PptxArtifactProbe:
    """Probe one deck in an output-silent, time-bounded worker."""
    result = _run_bounded_private_worker(
        path,
        child_flag=PPTX_ARTIFACT_PROBE_CHILD_FLAG,
        decoder=_decode_pptx_probe_result,
        trusted_root=trusted_root,
        deadline_monotonic=deadline_monotonic,
    )
    if not isinstance(result, PptxArtifactProbe):
        raise _probe_failure("pptx_probe_malformed_result")
    return result


def _probe_file_identity(
    path: str | Path,
    *,
    trusted_root: str | Path | None = None,
    deadline_monotonic: float | None = None,
) -> tuple[Path, _PptxCacheKey]:
    # Keep the lexical absolute path. Resolving here would follow a path that
    # was swapped to a symlink after lstat and could move worker I/O outside
    # the admitted artifact boundary.
    canonical = Path(os.path.abspath(os.fspath(path)))
    root = (
        Path(os.path.abspath(os.fspath(trusted_root)))
        if trusted_root is not None
        else None
    )
    generation = _supervised_file_generation(
        canonical,
        label="PPTX artifact",
        trusted_root=root,
        deadline_monotonic=deadline_monotonic,
    )
    return canonical, (
        os.fspath(canonical),
        generation.device,
        generation.inode,
        generation.size,
        generation.mtime_ns,
        generation.ctime_ns,
        generation.flags or 0,
        generation.file_attributes or 0,
        PPTX_EXTRACTION_SCHEMA_VERSION,
        PPTX_EXTRACTION_PIPELINE_VERSION,
    )


def _cloud_placeholder_failure(
    key: _PptxCacheKey,
) -> PptxEvidenceError | None:
    macos_flags = key[-4]
    file_attributes = key[-3]
    if PPTX_MACOS_DATALESS_FLAG and macos_flags & PPTX_MACOS_DATALESS_FLAG:
        return _probe_failure(
            "pptx_cloud_placeholder_unavailable",
            details={"st_flags": macos_flags},
        )
    if (
        PPTX_WINDOWS_CLOUD_FILE_ATTRIBUTES
        and file_attributes & PPTX_WINDOWS_CLOUD_FILE_ATTRIBUTES
    ):
        return _probe_failure(
            "pptx_cloud_placeholder_unavailable",
            details={"file_attributes": file_attributes},
        )
    return None


def _copy_probe(probe: PptxArtifactProbe) -> PptxArtifactProbe:
    return PptxArtifactProbe(
        slide_count=probe.slide_count,
        source_sha256=probe.source_sha256,
        source_size_bytes=probe.source_size_bytes,
        archive_recovery=tuple(dict(item) for item in probe.archive_recovery),
    )


def _cache_probe_result(
    key: _PptxCacheKey,
    value: PptxArtifactProbe | tuple[str, str, dict[str, object]],
) -> None:
    for stale_key in [
        candidate
        for candidate in _PPTX_ARTIFACT_PROBE_CACHE
        if candidate[0] == key[0] and candidate != key
    ]:
        _PPTX_ARTIFACT_PROBE_CACHE.pop(stale_key, None)
    _PPTX_ARTIFACT_PROBE_CACHE[key] = value


def clear_pptx_artifact_probe_cache() -> None:
    """Clear process-local probe memoization for tests and explicit refreshes."""
    _PPTX_ARTIFACT_PROBE_CACHE.clear()
    _PPTX_NATIVE_AUDIT_CACHE.clear()


def probe_pptx_artifact(
    path: str | Path,
    *,
    trusted_root: str | Path | None = None,
    deadline_monotonic: float | None = None,
) -> PptxArtifactProbe:
    """Return exact deck evidence through a bounded, generation-cached worker."""
    root = (
        Path(os.path.abspath(os.fspath(trusted_root)))
        if trusted_root is not None
        else None
    )
    artifact, key = _probe_file_identity(
        path,
        trusted_root=root,
        deadline_monotonic=deadline_monotonic,
    )
    cached = _PPTX_ARTIFACT_PROBE_CACHE.get(key)
    if isinstance(cached, PptxArtifactProbe):
        return _copy_probe(cached)
    if isinstance(cached, tuple):
        message, reason_code, details = cached
        raise PptxEvidenceError(
            message,
            reason_code=reason_code,
            details=details,
        )
    error = _cloud_placeholder_failure(key)
    if error is not None:
        cached_error = (str(error), error.reason_code, dict(error.details))
        _cache_probe_result(key, cached_error)
        raise error
    try:
        probe = _run_bounded_pptx_probe(
            artifact,
            trusted_root=root,
            deadline_monotonic=deadline_monotonic,
        )
    except PptxEvidenceError as exc:
        if exc.reason_code == "pptx_batch_wall_limit":
            raise
        try:
            _current_path, current_key = _probe_file_identity(
                artifact,
                trusted_root=root,
                deadline_monotonic=deadline_monotonic,
            )
        except PptxEvidenceError as identity_error:
            if identity_error.reason_code == "pptx_batch_wall_limit":
                raise
            raise _probe_failure("pptx_artifact_changed") from exc
        if current_key != key:
            raise _probe_failure("pptx_artifact_changed") from exc
        if exc.reason_code in {
            "pptx_artifact_changed",
            "pptx_probe_materialization_changed",
        }:
            raise _probe_failure("pptx_probe_materialization_changed") from exc
        if exc.reason_code in _ARCHIVE_INTEGRITY_CONFIRMATION_REASON_CODES:
            _confirm_archive_integrity_failure(
                artifact,
                key,
                exc,
                child_flag=PPTX_ARTIFACT_PROBE_CHILD_FLAG,
                decoder=_decode_pptx_probe_result,
                trusted_root=root,
                deadline_monotonic=deadline_monotonic,
            )
            cached_error = (str(exc), exc.reason_code, dict(exc.details))
            _cache_probe_result(key, cached_error)
            raise PptxEvidenceError(
                cached_error[0],
                reason_code=cached_error[1],
                details=cached_error[2],
            ) from exc
        # A timeout, crash, resource/monitor fault, or malformed protocol is an
        # invocation failure, not an artifact fact. Never make it sticky.
        raise PptxEvidenceError(
            str(exc),
            reason_code=exc.reason_code,
            details=dict(exc.details),
        ) from exc
    _between_path, between_key = _probe_file_identity(
        artifact,
        trusted_root=root,
        deadline_monotonic=deadline_monotonic,
    )
    if between_key != key:
        raise _probe_failure("pptx_artifact_changed")
    if probe.archive_recovery:
        try:
            confirmation = _run_bounded_pptx_probe(
                artifact,
                trusted_root=root,
                deadline_monotonic=deadline_monotonic,
            )
        except PptxEvidenceError as exc:
            if exc.reason_code == "pptx_batch_wall_limit":
                raise
            raise _probe_failure("pptx_probe_materialization_changed") from exc
        _confirmed_path, confirmed_key = _probe_file_identity(
            artifact,
            trusted_root=root,
            deadline_monotonic=deadline_monotonic,
        )
        if confirmed_key != key:
            raise _probe_failure("pptx_artifact_changed")
        if confirmation != probe:
            raise _probe_failure("pptx_probe_materialization_changed")
    _current_path, current_key = _probe_file_identity(
        artifact,
        trusted_root=root,
        deadline_monotonic=deadline_monotonic,
    )
    if current_key != key or probe.source_size_bytes != key[3]:
        raise _probe_failure("pptx_artifact_changed")
    cached_probe = _copy_probe(probe)
    _cache_probe_result(key, cached_probe)
    return _copy_probe(cached_probe)


def normalize_page_ranges(
    ranges: object,
    *,
    page_count: int,
    allow_empty: bool,
    label: str = "inspected_page_ranges",
) -> list[list[int]]:
    """Validate inclusive page ranges and return a normalized copy."""
    if (
        isinstance(page_count, bool)
        or not isinstance(page_count, int)
        or not 1 <= page_count <= PPTX_ARCHIVE_MAX_MEMBERS
    ):
        raise PptxEvidenceError(
            f"{label} page_count exceeds the bounded PPTX member contract"
        )
    if (
        not isinstance(ranges, list)
        or len(ranges) > PPTX_ARCHIVE_MAX_MEMBERS
        or (not ranges and not allow_empty)
    ):
        qualifier = "an array" if allow_empty else "a non-empty array"
        raise PptxEvidenceError(f"{label} must be {qualifier}")
    normalized: list[list[int]] = []
    prior_end = 0
    for index, raw_range in enumerate(ranges):
        range_label = f"{label}[{index}]"
        if not isinstance(raw_range, list) or len(raw_range) != 2:
            raise PptxEvidenceError(
                f"{range_label} must be a two-item [start, end] array"
            )
        start, end = raw_range
        if (
            isinstance(start, bool)
            or isinstance(end, bool)
            or not isinstance(start, int)
            or not isinstance(end, int)
            or start < 1
            or end < start
            or end > page_count
            or start <= prior_end
        ):
            raise PptxEvidenceError(
                f"{range_label} must be ascending, non-overlapping, and inside "
                f"the verified 1..{page_count} page bound"
            )
        normalized.append([start, end])
        prior_end = end
    return normalized


def _pages_covered_by_ranges(
    ranges: list[list[int]],
    pages: list[int],
) -> list[int]:
    """Return covered candidates without expanding any declared interval."""
    covered: list[int] = []
    range_index = 0
    for page in pages:
        while range_index < len(ranges) and ranges[range_index][1] < page:
            range_index += 1
        if (
            range_index < len(ranges)
            and ranges[range_index][0] <= page <= ranges[range_index][1]
        ):
            covered.append(page)
    return covered


def _normalize_required_slides(value: object, *, slide_count: int) -> list[int]:
    if (
        isinstance(slide_count, bool)
        or not isinstance(slide_count, int)
        or not 1 <= slide_count <= PPTX_ARCHIVE_MAX_MEMBERS
    ):
        raise PptxEvidenceError("slide_count exceeds the bounded PPTX contract")
    if not isinstance(value, list) or len(value) > PPTX_ARCHIVE_MAX_MEMBERS:
        raise PptxEvidenceError("render_required_slide_numbers must be an array")
    if any(
        isinstance(number, bool)
        or not isinstance(number, int)
        or number < 1
        or number > slide_count
        for number in value
    ):
        raise PptxEvidenceError(
            "render_required_slide_numbers must contain integers inside the "
            f"verified 1..{slide_count} slide bound"
        )
    if value != sorted(set(value)):
        raise PptxEvidenceError(
            "render_required_slide_numbers must be sorted and duplicate-free"
        )
    return list(value)


def build_native_deck_audit(
    *,
    source_pptx_sha256: str,
    source_pptx_size_bytes: int,
    slide_count: int,
    render_required_reasons: Mapping[int, list[str]],
    rendered_page_inspection: Mapping[str, object] | None = None,
) -> dict[str, object]:
    """Build the deterministic native-deck audit emitted by the extractor."""
    _require_sha256(source_pptx_sha256, "source_pptx_sha256")
    if (
        isinstance(source_pptx_size_bytes, bool)
        or not isinstance(source_pptx_size_bytes, int)
        or source_pptx_size_bytes < 1
    ):
        raise PptxEvidenceError("source_pptx_size_bytes must be positive")
    if (
        isinstance(slide_count, bool)
        or not isinstance(slide_count, int)
        or not 1 <= slide_count <= PPTX_ARCHIVE_MAX_MEMBERS
    ):
        raise PptxEvidenceError("slide_count must be positive")
    normalized_reasons: dict[str, list[str]] = {}
    for slide_number in sorted(render_required_reasons):
        reasons = render_required_reasons[slide_number]
        if (
            isinstance(slide_number, bool)
            or not isinstance(slide_number, int)
            or slide_number < 1
            or slide_number > slide_count
            or not isinstance(reasons, list)
            or not reasons
            or any(not isinstance(reason, str) or not reason for reason in reasons)
            or reasons != sorted(set(reasons))
        ):
            raise PptxEvidenceError(
                "render_required_reasons must map valid slide numbers to sorted, "
                "duplicate-free non-empty reason arrays"
            )
        normalized_reasons[str(slide_number)] = list(reasons)
    required = sorted(render_required_reasons)
    identity: dict[str, object] = {
        "schema_version": NATIVE_DECK_AUDIT_SCHEMA_VERSION,
        "extraction_schema_version": PPTX_EXTRACTION_SCHEMA_VERSION,
        "extraction_pipeline_version": PPTX_EXTRACTION_PIPELINE_VERSION,
        "source_pptx_sha256": source_pptx_sha256,
        "source_pptx_size_bytes": source_pptx_size_bytes,
        "slide_count": slide_count,
        "render_required_slide_numbers": required,
        "render_required_reasons": normalized_reasons,
    }
    return {
        **identity,
        "extraction_receipt_sha256": _canonical_sha256(identity),
        "rendered_page_inspection": (
            dict(rendered_page_inspection)
            if rendered_page_inspection is not None
            else None
        ),
    }


def _pdf_page_count(path: Path) -> int:
    pdf_read_error: type[Exception] = ValueError
    try:
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError as ImportedPdfReadError

        pdf_read_error = ImportedPdfReadError

        count = len(PdfReader(os.fspath(path), strict=True).pages)
    except ImportError as exc:
        raise PptxEvidenceError(
            "render receipt validation requires the declared pypdf dependency; "
            "install the speaker-toolkit project dependencies"
        ) from exc
    except (pdf_read_error, OSError, ValueError, KeyError, EOFError) as exc:
        raise PptxEvidenceError(
            f"rendered PDF is unreadable at {path}: {type(exc).__name__}"
        ) from exc
    if not 1 <= count <= PPTX_ARCHIVE_MAX_MEMBERS:
        raise PptxEvidenceError(f"rendered PDF has no pages at {path}")
    return count


def snapshot_rendered_pdf(path: str | Path) -> tuple[str, int, int]:
    """Copy, hash, and page-count one exact rendered-PDF generation."""
    artifact = Path(path)
    try:
        initial = artifact.lstat()
    except OSError as exc:
        raise PptxEvidenceError(
            f"rendered PDF is unavailable at {artifact}: {exc}"
        ) from exc
    if stat_module.S_ISLNK(initial.st_mode) or not stat_module.S_ISREG(initial.st_mode):
        raise PptxEvidenceError(
            f"rendered PDF must be a non-symlink regular file: {artifact}"
        )
    generation = _file_generation(initial)
    digest = hashlib.sha256()
    copied_size = 0
    with tempfile.TemporaryDirectory(prefix="speaker-toolkit-render-") as temp_dir:
        snapshot = Path(temp_dir) / "rendered.pdf"
        try:
            with artifact.open("rb") as source, snapshot.open("xb") as target:
                opened = os.fstat(source.fileno())
                if _file_generation(opened) != generation:
                    raise PptxEvidenceError(
                        f"rendered PDF changed while opening: {artifact}"
                    )
                while True:
                    chunk = source.read(1024 * 1024)
                    if not chunk:
                        break
                    digest.update(chunk)
                    target.write(chunk)
                    copied_size += len(chunk)
                after_read = os.fstat(source.fileno())
        except PptxEvidenceError:
            raise
        except OSError as exc:
            raise PptxEvidenceError(
                f"cannot snapshot rendered PDF at {artifact}: {exc}"
            ) from exc
        try:
            current = artifact.lstat()
        except OSError as exc:
            raise PptxEvidenceError(
                f"rendered PDF changed while it was read at {artifact}: {exc}"
            ) from exc
        if (
            _file_generation(after_read) != generation
            or _file_generation(current) != generation
            or copied_size != initial.st_size
        ):
            raise PptxEvidenceError(f"rendered PDF changed while reading: {artifact}")
        page_count = _pdf_page_count(snapshot)
    return digest.hexdigest(), copied_size, page_count


def build_rendered_page_inspection(
    *,
    source_pptx_sha256: str,
    rendered_pdf_path: str | Path,
    inspected_page_ranges: object,
    required_slide_numbers: list[int],
    slide_count: int,
) -> dict[str, object]:
    """Bind asserted page inspection to exact PPTX and rendered-PDF identities."""
    _require_sha256(source_pptx_sha256, "source_pptx_sha256")
    if (
        isinstance(slide_count, bool)
        or not isinstance(slide_count, int)
        or not 1 <= slide_count <= PPTX_ARCHIVE_MAX_MEMBERS
    ):
        raise PptxEvidenceError("slide_count must be a positive integer")
    required = _normalize_required_slides(
        required_slide_numbers, slide_count=slide_count
    )
    pdf_sha256, pdf_size, page_count = snapshot_rendered_pdf(Path(rendered_pdf_path))
    if page_count != slide_count:
        raise PptxEvidenceError(
            "rendered PDF page count must equal the source deck slide count; "
            f"expected {slide_count}, got {page_count}"
        )
    ranges = normalize_page_ranges(
        inspected_page_ranges,
        page_count=page_count,
        allow_empty=True,
    )
    inspected_required = _pages_covered_by_ranges(ranges, required)
    identity: dict[str, object] = {
        "schema_version": RENDER_INSPECTION_SCHEMA_VERSION,
        "source_pptx_sha256": source_pptx_sha256,
        "rendered_pdf_sha256": pdf_sha256,
        "rendered_pdf_size_bytes": pdf_size,
        "rendered_page_count": page_count,
        "inspected_page_ranges": ranges,
        "inspected_required_slide_numbers": inspected_required,
        "complete": inspected_required == required,
    }
    return {**identity, "binding_sha256": _canonical_sha256(identity)}


def validate_native_deck_audit(
    value: object,
    *,
    slide_count: int | None = None,
) -> dict[str, object]:
    """Validate a closed native-deck audit and its optional render receipt."""
    if not isinstance(value, Mapping):
        raise PptxEvidenceError("native_deck_audit must be an object")
    unknown = sorted(set(value) - _AUDIT_FIELDS)
    missing = sorted(_AUDIT_FIELDS - set(value))
    if unknown or missing:
        raise PptxEvidenceError(
            "native_deck_audit must contain exactly the schema fields; "
            f"missing={missing}, unknown={unknown}"
        )
    if value.get("schema_version") != NATIVE_DECK_AUDIT_SCHEMA_VERSION:
        raise PptxEvidenceError(
            "native_deck_audit.schema_version must be "
            f"{NATIVE_DECK_AUDIT_SCHEMA_VERSION}"
        )
    if value.get("extraction_schema_version") != PPTX_EXTRACTION_SCHEMA_VERSION:
        raise PptxEvidenceError(
            "native_deck_audit.extraction_schema_version must match the current "
            f"extractor schema {PPTX_EXTRACTION_SCHEMA_VERSION}"
        )
    if value.get("extraction_pipeline_version") != PPTX_EXTRACTION_PIPELINE_VERSION:
        raise PptxEvidenceError(
            "native_deck_audit.extraction_pipeline_version must match the current "
            f"extractor pipeline {PPTX_EXTRACTION_PIPELINE_VERSION}"
        )
    source_sha = _require_sha256(
        value.get("source_pptx_sha256"),
        "native_deck_audit.source_pptx_sha256",
    )
    size = value.get("source_pptx_size_bytes")
    if isinstance(size, bool) or not isinstance(size, int) or size < 1:
        raise PptxEvidenceError(
            "native_deck_audit.source_pptx_size_bytes must be a positive integer"
        )
    recorded_slide_count = value.get("slide_count")
    if (
        isinstance(recorded_slide_count, bool)
        or not isinstance(recorded_slide_count, int)
        or not 1 <= recorded_slide_count <= PPTX_ARCHIVE_MAX_MEMBERS
    ):
        raise PptxEvidenceError(
            "native_deck_audit.slide_count must be a positive integer"
        )
    if slide_count is not None and recorded_slide_count != slide_count:
        raise PptxEvidenceError(
            "native_deck_audit.slide_count must equal structured_data.slide_count"
        )
    required = _normalize_required_slides(
        value.get("render_required_slide_numbers"),
        slide_count=recorded_slide_count,
    )
    reasons = value.get("render_required_reasons")
    if not isinstance(reasons, Mapping):
        raise PptxEvidenceError(
            "native_deck_audit.render_required_reasons must be an object"
        )
    normalized_reasons: dict[str, list[str]] = {}
    for raw_slide, raw_reasons in reasons.items():
        if (
            not isinstance(raw_slide, str)
            or len(raw_slide) > 20
            or re.fullmatch(r"[1-9][0-9]*", raw_slide) is None
        ):
            raise PptxEvidenceError(
                "native_deck_audit.render_required_reasons keys must be bounded "
                "canonical decimal slide-number strings"
            )
        slide_number = int(raw_slide)
        if slide_number not in required:
            raise PptxEvidenceError(
                "native_deck_audit.render_required_reasons keys must exactly "
                "match render_required_slide_numbers"
            )
        if (
            not isinstance(raw_reasons, list)
            or not raw_reasons
            or any(not isinstance(reason, str) or not reason for reason in raw_reasons)
            or raw_reasons != sorted(set(raw_reasons))
        ):
            raise PptxEvidenceError(
                "native_deck_audit.render_required_reasons values must be sorted, "
                "duplicate-free non-empty string arrays"
            )
        normalized_reasons[raw_slide] = list(raw_reasons)
    if sorted(int(number) for number in normalized_reasons) != required:
        raise PptxEvidenceError(
            "native_deck_audit.render_required_reasons keys must exactly match "
            "render_required_slide_numbers"
        )
    identity: dict[str, object] = {
        "schema_version": NATIVE_DECK_AUDIT_SCHEMA_VERSION,
        "extraction_schema_version": PPTX_EXTRACTION_SCHEMA_VERSION,
        "extraction_pipeline_version": PPTX_EXTRACTION_PIPELINE_VERSION,
        "source_pptx_sha256": source_sha,
        "source_pptx_size_bytes": size,
        "slide_count": recorded_slide_count,
        "render_required_slide_numbers": required,
        "render_required_reasons": normalized_reasons,
    }
    if value.get("extraction_receipt_sha256") != _canonical_sha256(identity):
        raise PptxEvidenceError(
            "native_deck_audit.extraction_receipt_sha256 does not bind the "
            "declared extraction identity and render requirements"
        )
    raw_receipt = value.get("rendered_page_inspection")
    receipt = (
        validate_rendered_page_inspection(
            raw_receipt,
            required_slide_numbers=required,
            slide_count=recorded_slide_count,
        )
        if raw_receipt is not None
        else None
    )
    if receipt is not None and receipt["source_pptx_sha256"] != source_sha:
        raise PptxEvidenceError(
            "rendered_page_inspection.source_pptx_sha256 must match the "
            "native-deck audit source identity"
        )
    return {
        **identity,
        "extraction_receipt_sha256": value["extraction_receipt_sha256"],
        "rendered_page_inspection": receipt,
    }


def validate_rendered_page_inspection(
    value: object,
    *,
    required_slide_numbers: list[int],
    slide_count: int,
) -> dict[str, object]:
    """Validate a closed rendered-page receipt and recompute its binding."""
    if not isinstance(value, Mapping):
        raise PptxEvidenceError("rendered_page_inspection must be an object")
    unknown = sorted(set(value) - _RENDER_RECEIPT_FIELDS)
    missing = sorted(_RENDER_RECEIPT_FIELDS - set(value))
    if unknown or missing:
        raise PptxEvidenceError(
            "rendered_page_inspection must contain exactly the schema fields; "
            f"missing={missing}, unknown={unknown}"
        )
    if value.get("schema_version") != RENDER_INSPECTION_SCHEMA_VERSION:
        raise PptxEvidenceError(
            "rendered_page_inspection.schema_version must be "
            f"{RENDER_INSPECTION_SCHEMA_VERSION}"
        )
    source_sha = _require_sha256(
        value.get("source_pptx_sha256"),
        "rendered_page_inspection.source_pptx_sha256",
    )
    pdf_sha = _require_sha256(
        value.get("rendered_pdf_sha256"),
        "rendered_page_inspection.rendered_pdf_sha256",
    )
    pdf_size = value.get("rendered_pdf_size_bytes")
    page_count = value.get("rendered_page_count")
    if isinstance(pdf_size, bool) or not isinstance(pdf_size, int) or pdf_size < 1:
        raise PptxEvidenceError(
            "rendered_page_inspection.rendered_pdf_size_bytes must be positive"
        )
    if (
        isinstance(page_count, bool)
        or not isinstance(page_count, int)
        or page_count != slide_count
    ):
        raise PptxEvidenceError(
            "rendered_page_inspection.rendered_page_count must equal the source "
            f"deck slide count {slide_count}"
        )
    ranges = normalize_page_ranges(
        value.get("inspected_page_ranges"),
        page_count=page_count,
        allow_empty=True,
    )
    required = _normalize_required_slides(
        required_slide_numbers, slide_count=slide_count
    )
    inspected_required = _pages_covered_by_ranges(ranges, required)
    if value.get("inspected_required_slide_numbers") != inspected_required:
        raise PptxEvidenceError(
            "rendered_page_inspection.inspected_required_slide_numbers must be "
            "derived from inspected_page_ranges"
        )
    complete = inspected_required == required
    if value.get("complete") is not complete:
        raise PptxEvidenceError(
            "rendered_page_inspection.complete must reflect coverage of every "
            "render-required slide"
        )
    identity: dict[str, object] = {
        "schema_version": RENDER_INSPECTION_SCHEMA_VERSION,
        "source_pptx_sha256": source_sha,
        "rendered_pdf_sha256": pdf_sha,
        "rendered_pdf_size_bytes": pdf_size,
        "rendered_page_count": page_count,
        "inspected_page_ranges": ranges,
        "inspected_required_slide_numbers": inspected_required,
        "complete": complete,
    }
    if value.get("binding_sha256") != _canonical_sha256(identity):
        raise PptxEvidenceError(
            "rendered_page_inspection.binding_sha256 does not bind the exact "
            "source, render, and inspected ranges"
        )
    return {**identity, "binding_sha256": value["binding_sha256"]}


def _extract_native_deck_audit_in_process(path: Path) -> dict[str, object]:
    """Recompute one audit inside the bounded native-audit worker."""
    sys.modules.setdefault("pptx_evidence", sys.modules[__name__])
    extractor_path = Path(__file__).with_name("pptx-extraction.py")
    spec = importlib.util.spec_from_file_location(
        "_speaker_toolkit_bounded_pptx_extraction",
        extractor_path,
    )
    if spec is None or spec.loader is None:
        raise PptxEvidenceError(
            "cannot load the current PPTX extractor",
            reason_code="pptx_dependency_unavailable",
        )
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    extract = getattr(module, "_extract_pptx_in_process", None)
    if not callable(extract):
        raise PptxEvidenceError(
            "current PPTX extractor has no extraction entrypoint",
            reason_code="pptx_dependency_unavailable",
        )
    payload = extract(path, ocr=False)
    if not isinstance(payload, Mapping):
        raise PptxEvidenceError(
            "current PPTX extractor returned a non-object",
            reason_code="pptx_evidence_invalid",
        )
    raw_recovery = payload.get("archive_recovery")
    if isinstance(raw_recovery, list) and raw_recovery:
        part_names = sorted(
            str(item.get("part_name", "<unknown>"))
            for item in raw_recovery
            if isinstance(item, Mapping)
        )
        raise PptxEvidenceError(
            "PPTX extraction required placeholder archive recovery",
            reason_code="pptx_archive_recovery_required",
            details={"part_names": part_names},
        )
    return validate_native_deck_audit(payload.get("native_deck_audit"))


def _native_audit_child(path: str | Path) -> dict[str, object]:
    """Return a closed native-audit payload from the isolated worker."""
    try:
        audit = _extract_native_deck_audit_in_process(Path(path))
    except MemoryError:
        return {
            "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": "pptx_probe_resource_unavailable",
            "details": {},
        }
    except PptxEvidenceError as exc:
        reason_code = (
            exc.reason_code
            if exc.reason_code in _CHILD_PROBE_REASON_CODES
            else "pptx_evidence_invalid"
        )
        return {
            "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": reason_code,
            "details": _probe_child_failure_details(exc),
        }
    return {
        "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
        "status": "available",
        "native_deck_audit": audit,
    }


def _decode_native_audit_result(raw: bytes) -> dict[str, object]:
    if len(raw) > PPTX_ARTIFACT_PROBE_MAX_RESULT_BYTES or not raw.strip():
        raise _probe_failure("pptx_probe_malformed_result")
    try:
        decoded = raw.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    if len(decoded.splitlines()) != 1:
        raise _probe_failure("pptx_probe_malformed_result")
    try:
        payload = json.loads(
            decoded,
            object_pairs_hook=_unique_json_object,
            parse_constant=_reject_nonstandard_json_constant,
        )
    except (json.JSONDecodeError, RecursionError, ValueError) as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    if (
        not isinstance(payload, dict)
        or payload.get("schema_version") != PPTX_ARTIFACT_PROBE_SCHEMA_VERSION
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    if payload.get("status") == "unavailable":
        if set(payload) != {"schema_version", "status", "reason_code", "details"}:
            raise _probe_failure("pptx_probe_malformed_result")
        reason_code = payload.get("reason_code")
        details = payload.get("details")
        if (
            not isinstance(reason_code, str)
            or reason_code not in _CHILD_PROBE_REASON_CODES
            or not isinstance(details, dict)
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        if set(details) - {"part_names", "exception_type"}:
            raise _probe_failure("pptx_probe_malformed_result")
        part_names = details.get("part_names")
        if part_names is not None and (
            not isinstance(part_names, list)
            or len(part_names) > 64
            or any(
                not isinstance(name, str) or not name or len(name) > 2048
                for name in part_names
            )
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        exception_type = details.get("exception_type")
        if exception_type is not None and (
            not isinstance(exception_type, str)
            or not exception_type
            or len(exception_type) > 128
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        raise _probe_failure(str(reason_code), details=details)
    if (
        set(payload) != {"schema_version", "status", "native_deck_audit"}
        or payload.get("status") != "available"
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    try:
        return validate_native_deck_audit(payload.get("native_deck_audit"))
    except PptxEvidenceError as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc


def _cache_native_audit(
    key: _PptxCacheKey,
    value: dict[str, object] | tuple[str, str, dict[str, object]],
) -> None:
    stale = [
        candidate
        for candidate in _PPTX_NATIVE_AUDIT_CACHE
        if candidate[0] == key[0] and candidate != key
    ]
    for stale_key in stale:
        _PPTX_NATIVE_AUDIT_CACHE.pop(stale_key, None)
    _PPTX_NATIVE_AUDIT_CACHE[key] = value


def recompute_native_deck_audit(
    path: str | Path,
    *,
    trusted_root: str | Path | None = None,
) -> dict[str, object]:
    """Recompute an exact audit through a bounded, generation-cached worker."""
    root = (
        Path(os.path.abspath(os.fspath(trusted_root)))
        if trusted_root is not None
        else None
    )
    artifact, key = _probe_file_identity(path, trusted_root=root)
    cached = _PPTX_NATIVE_AUDIT_CACHE.get(key)
    if isinstance(cached, dict):
        return copy.deepcopy(cached)
    if isinstance(cached, tuple):
        message, reason_code, details = cached
        raise PptxEvidenceError(
            message,
            reason_code=reason_code,
            details=details,
        )
    error = _cloud_placeholder_failure(key)
    if error is not None:
        _cache_native_audit(
            key,
            (str(error), error.reason_code, dict(error.details)),
        )
        raise error
    try:
        result = _run_bounded_private_worker(
            artifact,
            child_flag=PPTX_NATIVE_AUDIT_CHILD_FLAG,
            decoder=_decode_native_audit_result,
            trusted_root=root,
        )
    except PptxEvidenceError as exc:
        _current_path, current_key = _probe_file_identity(
            artifact,
            trusted_root=root,
        )
        if current_key != key:
            raise _probe_failure("pptx_artifact_changed") from exc
        if exc.reason_code in {
            "pptx_artifact_changed",
            "pptx_probe_materialization_changed",
        }:
            raise _probe_failure("pptx_probe_materialization_changed") from exc
        if exc.reason_code in _ARCHIVE_INTEGRITY_CONFIRMATION_REASON_CODES:
            _confirm_archive_integrity_failure(
                artifact,
                key,
                exc,
                child_flag=PPTX_NATIVE_AUDIT_CHILD_FLAG,
                decoder=_decode_native_audit_result,
                trusted_root=root,
            )
            cached_error = (str(exc), exc.reason_code, dict(exc.details))
            _cache_native_audit(key, cached_error)
            raise PptxEvidenceError(
                cached_error[0],
                reason_code=cached_error[1],
                details=cached_error[2],
            ) from exc
        raise PptxEvidenceError(
            str(exc),
            reason_code=exc.reason_code,
            details=dict(exc.details),
        ) from exc
    if not isinstance(result, dict):
        raise _probe_failure("pptx_probe_malformed_result")
    _current_path, current_key = _probe_file_identity(
        artifact,
        trusted_root=root,
    )
    if current_key != key or result.get("source_pptx_size_bytes") != key[3]:
        raise _probe_failure("pptx_artifact_changed")
    cached_result = copy.deepcopy(result)
    _cache_native_audit(key, cached_result)
    return copy.deepcopy(cached_result)


def _generation_cloud_placeholder_details(
    generation: FileGeneration,
) -> dict[str, object] | None:
    if (
        PPTX_MACOS_DATALESS_FLAG
        and generation.flags is not None
        and generation.flags & PPTX_MACOS_DATALESS_FLAG
    ):
        return {"st_flags": generation.flags}
    if (
        PPTX_WINDOWS_CLOUD_FILE_ATTRIBUTES
        and generation.file_attributes is not None
        and generation.file_attributes & PPTX_WINDOWS_CLOUD_FILE_ATTRIBUTES
    ):
        return {"file_attributes": generation.file_attributes}
    return None


def _admit_supervised_input(
    path: str | Path,
    *,
    label: str,
    trusted_root: str | Path | None = None,
    deadline_monotonic: float | None = None,
) -> tuple[Path, FileGeneration, FileGeneration | None]:
    artifact = Path(os.path.abspath(os.fspath(path)))
    root = (
        Path(os.path.abspath(os.fspath(trusted_root)))
        if trusted_root is not None
        else None
    )
    root_generation: FileGeneration | None = None
    if root is None:
        generation = _supervised_file_generation(
            artifact,
            label=label,
            deadline_monotonic=deadline_monotonic,
        )
    else:
        receipt = _run_bounded_metadata_worker(
            artifact,
            trusted_root=root,
            deadline_monotonic=deadline_monotonic,
        )
        generation = receipt.generation
        root_generation = receipt.root_generation
        if root_generation is None:
            raise _probe_failure("pptx_probe_malformed_result")
    placeholder = _generation_cloud_placeholder_details(generation)
    if placeholder is not None:
        raise _probe_failure(
            "pptx_cloud_placeholder_unavailable",
            details=placeholder,
        )
    if generation.size > PPTX_MAX_INPUT_BYTES:
        raise _probe_failure(
            "pptx_probe_resource_unavailable",
            details={"limit_bytes": PPTX_MAX_INPUT_BYTES},
        )
    return artifact, generation, root_generation


def _limits_before_deadline(
    limits: SupervisorLimits,
    deadline_monotonic: float | None,
) -> SupervisorLimits:
    """Clamp one worker wall limit to an enclosing batch deadline.

    Reserve the configured cleanup window so timeout containment itself remains
    inside the caller's deadline instead of extending the aggregate batch wall.
    """
    if deadline_monotonic is None:
        return limits
    if (
        isinstance(deadline_monotonic, bool)
        or not isinstance(deadline_monotonic, (int, float))
        or not math.isfinite(float(deadline_monotonic))
    ):
        raise PptxEvidenceError(
            "deadline_monotonic must be a finite monotonic timestamp",
            reason_code="pptx_evidence_invalid",
        )
    remaining = float(deadline_monotonic) - time.monotonic() - limits.cleanup_seconds
    if remaining <= 0:
        raise PptxEvidenceError(
            "PPTX batch wall deadline was exhausted before worker launch",
            reason_code="pptx_batch_wall_limit",
        )
    if remaining >= limits.wall_seconds:
        return limits
    return replace(limits, wall_seconds=remaining)


def _validated_requested_page_ranges(value: object) -> list[list[int]]:
    if not isinstance(value, list):
        raise PptxEvidenceError(
            "inspected_page_ranges must be an array",
            reason_code="pptx_evidence_invalid",
        )
    normalized: list[list[int]] = []
    prior_end = 0
    for raw in value:
        if (
            not isinstance(raw, list)
            or len(raw) != 2
            or any(isinstance(item, bool) or not isinstance(item, int) for item in raw)
        ):
            raise PptxEvidenceError(
                "inspected_page_ranges entries must be integer pairs",
                reason_code="pptx_evidence_invalid",
            )
        start, end = raw
        if start < 1 or end < start or start <= prior_end:
            raise PptxEvidenceError(
                "inspected_page_ranges must be positive, ordered, and disjoint",
                reason_code="pptx_evidence_invalid",
            )
        normalized.append([start, end])
        prior_end = end
    return normalized


def _plain_nonnegative_int(value: object) -> bool:
    return type(value) is int and value >= 0


def _finite_number_between(
    value: object,
    *,
    minimum: float | None = None,
    maximum: float | None = None,
) -> bool:
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        return False
    try:
        numeric = float(value)
    except (OverflowError, ValueError):
        return False
    return (
        math.isfinite(numeric)
        and (minimum is None or numeric >= minimum)
        and (maximum is None or numeric <= maximum)
    )


def _bounded_text(
    value: object,
    *,
    maximum: int = 16 * 1024 * 1024,
    allow_none: bool = False,
    allow_empty: bool = True,
) -> bool:
    if value is None:
        return allow_none
    return (
        isinstance(value, str)
        and "\x00" not in value
        and len(value) <= maximum
        and (allow_empty or bool(value))
    )


def _valid_opc_part_name(value: object) -> bool:
    if not _bounded_text(value, maximum=2048, allow_empty=False):
        return False
    try:
        key = _opc_member_name_key(cast(str, value))
    except PptxEvidenceError:
        return False
    return bool(key) and not cast(str, value).endswith("/")


def _valid_string_path(value: object, *, allow_empty: bool = False) -> bool:
    return (
        isinstance(value, list)
        and (allow_empty or bool(value))
        and len(value) <= 64
        and all(_bounded_text(item, maximum=4096, allow_empty=False) for item in value)
    )


def _valid_count_map(
    value: object,
    *,
    required_keys: tuple[str, ...] | None = None,
    total_key: bool = False,
) -> bool:
    if not isinstance(value, dict) or len(value) > 65_536:
        return False
    if required_keys is not None:
        expected = {*required_keys, *(("total",) if total_key else ())}
        if set(value) != expected:
            return False
    if any(
        not _bounded_text(key, maximum=4096, allow_empty=False)
        or not _plain_nonnegative_int(count)
        for key, count in value.items()
    ):
        return False
    return not total_key or value["total"] == sum(
        value[key] for key in required_keys or ()
    )


def _valid_template_layouts(value: object) -> bool:
    if not isinstance(value, list) or len(value) > PPTX_ARCHIVE_MAX_MEMBERS:
        return False
    prior_master = 0
    for expected_index, layout in enumerate(value):
        if not isinstance(layout, dict) or set(layout) != _TEMPLATE_LAYOUT_FIELDS:
            return False
        placeholders = layout.get("placeholders")
        if (
            layout.get("index") != expected_index
            or not _plain_nonnegative_int(layout.get("master_index"))
            or not _bounded_text(layout.get("name"), maximum=4096)
            or not isinstance(placeholders, list)
            or len(placeholders) > PPTX_ARCHIVE_MAX_MEMBERS
        ):
            return False
        master_index = cast(int, layout["master_index"])
        if master_index < prior_master:
            return False
        prior_master = master_index
        for placeholder in placeholders:
            if (
                not isinstance(placeholder, dict)
                or set(placeholder) != _TEMPLATE_PLACEHOLDER_FIELDS
                or not _plain_nonnegative_int(placeholder.get("idx"))
                or not _bounded_text(
                    placeholder.get("type"), maximum=256, allow_empty=False
                )
            ):
                return False
    return True


def _parsed_reduced_aspect_ratio(value: object) -> tuple[int, int] | None:
    if not _bounded_text(value, maximum=64, allow_empty=False):
        return None
    ratio_match = re.fullmatch(r"([1-9][0-9]*):([1-9][0-9]*)", cast(str, value))
    if ratio_match is None:
        return None
    try:
        width = int(ratio_match.group(1))
        height = int(ratio_match.group(2))
    except ValueError:
        return None
    if math.gcd(width, height) != 1:
        return None
    return width, height


def _aspect_ratio_matches_dimensions(
    value: object,
    *,
    width_inches: object,
    height_inches: object,
) -> bool:
    """Bind the exact reduced ratio to dimensions rounded to two decimals."""
    ratio = _parsed_reduced_aspect_ratio(value)
    if (
        ratio is None
        or not _finite_number_between(width_inches, minimum=0.01)
        or not _finite_number_between(height_inches, minimum=0.01)
    ):
        return False
    width = float(cast(float, width_inches))
    height = float(cast(float, height_inches))
    width_low = max(0.0, width - 0.005)
    width_high = width + 0.005
    height_low = height - 0.005
    height_high = height + 0.005
    if height_low <= 0:
        return False
    claimed = ratio[0] / ratio[1]
    return width_low / height_high <= claimed <= width_high / height_low


def _picture_ratio_matches_geometry(
    value: object,
    *,
    shapes: list[dict[str, object]],
    slide_width_inches: object,
    slide_height_inches: object,
) -> bool:
    """Cross-bind a three-decimal image ratio to two-decimal shape geometry."""
    if (
        not _finite_number_between(value, minimum=0, maximum=1)
        or not _finite_number_between(slide_width_inches, minimum=0.01)
        or not _finite_number_between(slide_height_inches, minimum=0.01)
    ):
        return False
    slide_width = float(cast(float, slide_width_inches))
    slide_height = float(cast(float, slide_height_inches))
    slide_width_low = slide_width - 0.005
    slide_width_high = slide_width + 0.005
    slide_height_low = slide_height - 0.005
    slide_height_high = slide_height + 0.005
    if slide_width_low <= 0 or slide_height_low <= 0:
        return False

    lower = 0.0
    upper = 0.0
    pictures = [shape for shape in shapes if shape["is_picture"] is True]
    for picture in pictures:
        raw_width = picture.get("width")
        raw_height = picture.get("height")
        if raw_width is None or raw_height is None:
            picture_lower = 0.0
            picture_upper = 0.0
        else:
            if not _finite_number_between(
                raw_width, minimum=0
            ) or not _finite_number_between(raw_height, minimum=0):
                return False
            width = float(cast(float, raw_width))
            height = float(cast(float, raw_height))
            width_low = max(0.0, width - 0.005)
            width_high = width + 0.005
            height_low = max(0.0, height - 0.005)
            height_high = height + 0.005
            picture_lower = min(
                (width_low * height_low) / (slide_width_high * slide_height_high),
                1.0,
            )
            picture_upper = min(
                (width_high * height_high) / (slide_width_low * slide_height_low),
                1.0,
            )
        lower = max(lower, picture_lower)
        upper = max(upper, picture_upper)

    reported = float(cast(float, value))
    reported_low = max(0.0, reported - 0.0005)
    reported_high = min(1.0, reported + 0.0005)
    return reported_low <= upper and lower <= reported_high


def _valid_shape_summary(value: object) -> bool:
    if not isinstance(value, dict):
        return False
    fields = set(value)
    if not _SHAPE_REQUIRED_FIELDS.issubset(fields) or fields - (
        _SHAPE_REQUIRED_FIELDS | _SHAPE_OPTIONAL_FIELDS
    ):
        return False
    path = value.get("shape_path")
    if (
        not _bounded_text(value.get("name"), maximum=4096, allow_empty=False)
        or not _bounded_text(value.get("shape_type"), maximum=4096, allow_empty=False)
        or type(value.get("has_text_frame")) is not bool
        or type(value.get("is_picture")) is not bool
        or type(value.get("is_graphic_frame")) is not bool
        or not isinstance(path, list)
        or not _valid_string_path(path)
        or value.get("name") != path[-1]
        or not _plain_nonnegative_int(value.get("group_depth"))
        or value.get("group_depth") != len(path) - 1
    ):
        return False
    for field in ("left", "top"):
        raw = value.get(field)
        if raw is not None and not _finite_number_between(raw):
            return False
    for field in ("width", "height"):
        raw = value.get(field)
        if raw is not None and not _finite_number_between(raw, minimum=0):
            return False
    for field in (
        "text_preview",
        "font_name",
        "font_color",
        "fill_color",
        "line_color",
        "auto_shape_type",
        "graphic_frame_type",
        "graphic_data_uri",
        "table_text_preview",
        "picture_part_name",
    ):
        if field in value and not _bounded_text(
            value[field], maximum=8192, allow_none=True
        ):
            return False
    for field in ("font_size", "line_width"):
        if (
            field in value
            and value[field] is not None
            and not _finite_number_between(value[field], minimum=0)
        ):
            return False
    for field in ("bold", "italic"):
        if (
            field in value
            and value[field] is not None
            and type(value[field]) is not bool
        ):
            return False
    for field in ("table_rows", "table_columns"):
        if field in value and not _plain_nonnegative_int(value[field]):
            return False
    graphic_frame_type = value.get("graphic_frame_type")
    graphic_data_uri = value.get("graphic_data_uri")
    if graphic_frame_type is not None and (
        not isinstance(graphic_frame_type, str)
        or graphic_frame_type
        not in {
            "table",
            "chart",
            "embedded_ole_object",
            "linked_ole_object",
            "smartart",
            "graphic_frame",
        }
    ):
        return False
    if graphic_data_uri is not None and not _bounded_text(
        graphic_data_uri,
        maximum=8192,
        allow_empty=False,
    ):
        return False
    is_graphic_frame = cast(bool, value["is_graphic_frame"])
    is_picture = cast(bool, value["is_picture"])
    shape_prefix = cast(str, value["shape_type"]).split(" ", 1)[0]
    graphic_shape_prefixes = {
        "TABLE",
        "CHART",
        "DIAGRAM",
        "IGX_GRAPHIC",
        "EMBEDDED_OLE_OBJECT",
        "LINKED_OLE_OBJECT",
        "None",
    }
    expected_graphic_frame_type: str | None = (
        "graphic_frame" if is_graphic_frame and graphic_data_uri is None else None
    )
    if isinstance(graphic_data_uri, str):
        if graphic_data_uri == _GRAPHIC_DATA_URI_TABLE:
            expected_graphic_frame_type = "table"
        elif graphic_data_uri == _GRAPHIC_DATA_URI_CHART:
            expected_graphic_frame_type = "chart"
        elif graphic_data_uri == _GRAPHIC_DATA_URI_OLE:
            expected_graphic_frame_type = (
                "embedded_ole_object"
                if shape_prefix == "EMBEDDED_OLE_OBJECT"
                else "linked_ole_object"
            )
        elif "diagram" in graphic_data_uri.lower():
            expected_graphic_frame_type = "smartart"
        else:
            expected_graphic_frame_type = "graphic_frame"
    known_graphic_bindings = {
        "TABLE": ("table", _GRAPHIC_DATA_URI_TABLE),
        "CHART": ("chart", _GRAPHIC_DATA_URI_CHART),
        "EMBEDDED_OLE_OBJECT": ("embedded_ole_object", _GRAPHIC_DATA_URI_OLE),
        "LINKED_OLE_OBJECT": ("linked_ole_object", _GRAPHIC_DATA_URI_OLE),
    }
    known_graphic_binding = known_graphic_bindings.get(shape_prefix)
    picture_fields = {
        "picture_asset_status",
        "picture_part_name",
        "picture_asset_sha256",
    }
    text_format_fields = {
        "font_name",
        "font_size",
        "font_color",
        "bold",
        "italic",
    }
    if (
        is_graphic_frame is not (graphic_frame_type is not None)
        or (not is_graphic_frame and graphic_data_uri is not None)
        or graphic_frame_type != expected_graphic_frame_type
        or (is_graphic_frame and shape_prefix not in graphic_shape_prefixes)
        or (shape_prefix in {"TABLE", "CHART"} and not is_graphic_frame)
        or (
            known_graphic_binding is not None
            and (graphic_frame_type, graphic_data_uri) != known_graphic_binding
        )
        or (
            graphic_frame_type in {"table", "chart"}
            and shape_prefix != graphic_frame_type.upper()
        )
        or (
            graphic_frame_type in {"embedded_ole_object", "linked_ole_object"}
            and known_graphic_binding is None
        )
        or (
            (
                is_picture
                or shape_prefix
                in {"GROUP", "LINE", "LINKED_PICTURE", "MEDIA", "WEB_VIDEO"}
            )
            and (cast(bool, value["has_text_frame"]) or is_graphic_frame)
        )
        or (shape_prefix == "PICTURE" and not is_picture)
        or (is_picture and shape_prefix not in {"PICTURE", "PLACEHOLDER"})
        or (
            shape_prefix in {"AUTO_SHAPE", "FREEFORM", "PLACEHOLDER", "TEXT_BOX"}
            and not is_picture
            and not cast(bool, value["has_text_frame"])
        )
        or (is_graphic_frame and cast(bool, value["has_text_frame"]))
        or ("text_preview" in value) is not cast(bool, value["has_text_frame"])
        or (picture_fields.issubset(fields)) is not is_picture
        or bool(fields & picture_fields) is not is_picture
        or ("auto_shape_type" in fields) is not (shape_prefix == "AUTO_SHAPE")
        or (
            not cast(bool, value["has_text_frame"])
            and bool(fields & text_format_fields)
        )
        or (
            bool(fields & text_format_fields)
            and not text_format_fields.issubset(fields)
        )
    ):
        return False
    if is_picture:
        picture_status = value.get("picture_asset_status")
        part_name = value.get("picture_part_name")
        asset_sha256 = value.get("picture_asset_sha256")
        if picture_status not in {"available", "corrupt", "unavailable"}:
            return False
        if picture_status in {"available", "corrupt"}:
            if (
                not _valid_opc_part_name(part_name)
                or not isinstance(asset_sha256, str)
                or _SHA256_RE.fullmatch(asset_sha256) is None
            ):
                return False
        elif part_name is not None or asset_sha256 is not None:
            return False
    table_fields = {
        "table_rows",
        "table_columns",
        "table_text_preview",
        "table_fonts",
    }
    if graphic_frame_type == "table":
        if (
            not table_fields.issubset(fields)
            or cast(int, value["table_rows"]) < 1
            or cast(int, value["table_columns"]) < 1
        ):
            return False
    elif fields & table_fields:
        return False
    return "table_fonts" not in value or _valid_count_map(value["table_fonts"])


def _valid_shape_topology(shapes: list[dict[str, object]]) -> bool:
    """Require preorder nested paths whose prior parent is a real group."""
    prior_groups: list[tuple[str, ...]] = []
    for shape in shapes:
        path = tuple(cast(list[str], shape["shape_path"]))
        if len(path) > 1 and path[:-1] not in prior_groups:
            return False
        if shape["shape_type"] == "GROUP (6)":
            prior_groups.append(path)
    return True


def _valid_channel_provenance(value: object) -> bool:
    if not isinstance(value, dict):
        return False
    allowed = {
        "source",
        "shape_path",
        "shape_paths",
        "cells",
        "graphic_data_uri",
        "background_owner",
        "relationship_id",
        "part_name",
        "asset_status",
    }
    if not value or set(value) - allowed:
        return False
    if not _bounded_text(value.get("source"), maximum=256, allow_empty=False):
        return False
    if "shape_path" in value and not _valid_string_path(
        value["shape_path"], allow_empty=True
    ):
        return False
    if "shape_paths" in value:
        paths = value["shape_paths"]
        if (
            not isinstance(paths, list)
            or len(paths) > 65_536
            or any(not _valid_string_path(path, allow_empty=True) for path in paths)
        ):
            return False
    if "cells" in value:
        cells = value["cells"]
        if (
            not isinstance(cells, list)
            or len(cells) > 65_536
            or any(
                not _bounded_text(cell, maximum=64, allow_empty=False) for cell in cells
            )
        ):
            return False
    if "part_name" in value and value["part_name"] is not None:
        if not _valid_opc_part_name(value["part_name"]):
            return False
    return all(
        field not in value or _bounded_text(value[field], maximum=4096, allow_none=True)
        for field in (
            "graphic_data_uri",
            "background_owner",
            "relationship_id",
            "part_name",
            "asset_status",
        )
    )


def _valid_ocr_receipt(value: object) -> bool:
    if not isinstance(value, dict) or set(value) != _OCR_RECEIPT_FIELDS:
        return False
    confidence = value.get("result_confidence")
    recovered_text = value.get("recovered_text")
    status = value.get("result_status")
    engine_version = value.get("engine_version")
    error = value.get("error")
    if not isinstance(recovered_text, str) or not isinstance(status, str):
        return False
    normalized_text = re.sub(r"\s+", " ", recovered_text).strip()
    confidence_valid = confidence is None or _finite_number_between(
        confidence,
        minimum=0,
        maximum=100,
    )
    attempted_with_version = value.get("attempted") is True and _bounded_text(
        engine_version, maximum=128, allow_empty=False
    )
    status_semantics = {
        "text_recovered": (
            attempted_with_version
            and bool(recovered_text)
            and value.get("trustworthy_text") is True
            and confidence is not None
            and _finite_number_between(
                confidence,
                minimum=PPTX_OCR_TRUST_CONFIDENCE,
                maximum=100,
            )
            and error is None
        ),
        "low_confidence_text": (
            attempted_with_version
            and bool(recovered_text)
            and value.get("trustworthy_text") is False
            and (
                confidence is None
                or (
                    _finite_number_between(confidence, minimum=0, maximum=100)
                    and float(cast(float, confidence)) < PPTX_OCR_TRUST_CONFIDENCE
                )
            )
            and error is None
        ),
        "genuine_empty": (
            attempted_with_version
            and not recovered_text
            and value.get("trustworthy_text") is False
            and confidence is None
            and error is None
        ),
        "failed": (
            attempted_with_version
            and not recovered_text
            and value.get("trustworthy_text") is False
            and confidence is None
            and _bounded_text(error, maximum=512, allow_empty=False)
        ),
        "unavailable": (
            value.get("attempted") is False
            and not recovered_text
            and value.get("trustworthy_text") is False
            and confidence is None
            and engine_version is None
            and _bounded_text(error, maximum=512, allow_empty=False)
        ),
        "skipped": (
            value.get("attempted") is False
            and not recovered_text
            and value.get("trustworthy_text") is False
            and confidence is None
            and value.get("engine") == "tesseract"
            and value.get("engine_version") is None
            and value.get("error") == "ocr_disabled"
        ),
    }
    return (
        type(value.get("attempted")) is bool
        and value.get("engine") == "tesseract"
        and _bounded_text(engine_version, maximum=128, allow_none=True)
        and status in status_semantics
        and status_semantics[cast(str, status)]
        and confidence_valid
        and _bounded_text(error, maximum=512, allow_none=True)
        and _valid_opc_part_name(value.get("part_name"))
        and isinstance(value.get("asset_sha256"), str)
        and _SHA256_RE.fullmatch(value["asset_sha256"]) is not None
        and _valid_string_path(value.get("shape_path"), allow_empty=True)
        and _bounded_text(recovered_text, maximum=8000)
        and recovered_text == normalized_text
        and type(value.get("trustworthy_text")) is bool
    )


def _valid_ocr_channel_aggregate(value: dict[str, object]) -> bool:
    receipts = cast(list[dict[str, object]], value["ocr_receipts"])
    channel = value.get("channel")
    provenance = cast(dict[str, object], value["provenance"])
    if channel == "picture_ocr":
        if provenance.get("source") != "embedded_picture_blobs":
            return False
        expected_paths = [receipt["shape_path"] for receipt in receipts]
        if provenance.get("shape_paths") != expected_paths:
            return False
    elif channel == "background_image_ocr":
        if (
            provenance.get("source") != "pptx_background_image"
            or len(receipts) > 1
            or any(receipt["shape_path"] != [] for receipt in receipts)
            or (
                bool(receipts)
                and receipts[0].get("part_name") != provenance.get("part_name")
            )
        ):
            return False
    else:
        return False

    if not receipts:
        return (
            value.get("text") == ""
            and value.get("result_confidence") is None
            and value.get("status") == "unavailable"
            and value.get("attempted") is False
            and value.get("engine") == "tesseract"
            and value.get("engine_version") is None
            and value.get("reason") == "no_readable_asset"
        )

    statuses = {cast(str, receipt["result_status"]) for receipt in receipts}
    if "skipped" in statuses:
        return statuses == {"skipped"} and (
            value.get("text") == ""
            and value.get("result_confidence") is None
            and value.get("status") == "skipped"
            and value.get("attempted") is False
            and value.get("engine") == "tesseract"
            and value.get("engine_version") is None
            and value.get("reason") == "ocr_disabled"
        )

    recovered = [
        cast(str, receipt["recovered_text"])
        for receipt in receipts
        if receipt["recovered_text"]
    ]
    text = re.sub(r"\s+", " ", " | ".join(recovered)).strip()[:8000]
    confidences = [
        cast(float, receipt["result_confidence"])
        for receipt in receipts
        if receipt["result_confidence"] is not None
    ]
    expected_confidence = (
        round(sum(confidences) / len(confidences), 3) if confidences else None
    )
    attempted = any(cast(bool, receipt["attempted"]) for receipt in receipts)
    engines = {
        cast(str, receipt["engine"])
        for receipt in receipts
        if isinstance(receipt.get("engine"), str) and receipt["engine"]
    }
    engine = next(iter(engines)) if len(engines) == 1 else None
    versions = {
        cast(str, receipt["engine_version"])
        for receipt in receipts
        if isinstance(receipt.get("engine_version"), str) and receipt["engine_version"]
    }
    engine_version = next(iter(versions)) if len(versions) == 1 else None
    all_text_trustworthy = all(
        receipt["trustworthy_text"] is True
        for receipt in receipts
        if receipt["recovered_text"]
    )
    if text and statuses == {"text_recovered"} and all_text_trustworthy:
        status = "extracted"
        reason = None
    elif text:
        status = "partial"
        reason = "partial_ocr_results"
    elif statuses == {"genuine_empty"}:
        status = "empty"
        reason = None
    elif "failed" in statuses:
        status = "failed"
        reason = "ocr_failed"
    elif statuses == {"unavailable"}:
        status = "unavailable"
        reason = "ocr_engine_unavailable"
    else:
        status = "partial"
        reason = "partial_ocr_results"
    return (
        value.get("text") == text
        and value.get("result_confidence") == expected_confidence
        and value.get("status") == status
        and value.get("attempted") is attempted
        and value.get("engine") == engine
        and value.get("engine_version") == engine_version
        and value.get("reason") == reason
    )


def _valid_text_channel(value: object) -> bool:
    if not isinstance(value, dict):
        return False
    is_ocr = "ocr_receipts" in value
    if set(value) != (
        _TEXT_CHANNEL_OCR_FIELDS if is_ocr else _TEXT_CHANNEL_BASE_FIELDS
    ):
        return False
    channel = value.get("channel")
    confidence_label = value.get("confidence")
    status = value.get("status")
    if (
        not _bounded_text(channel, maximum=256, allow_empty=False)
        or not _bounded_text(value.get("text"))
        or not isinstance(confidence_label, str)
        or confidence_label not in {"high", "medium", "low"}
        or not isinstance(status, str)
        or status
        not in {
            "extracted",
            "empty",
            "partial",
            "failed",
            "skipped",
            "unavailable",
            "unsupported",
            "requires_render",
        }
        or not _valid_channel_provenance(value.get("provenance"))
    ):
        return False
    if not is_ocr:
        return True
    receipts = value.get("ocr_receipts")
    confidence = value.get("result_confidence")
    return (
        isinstance(channel, str)
        and channel in {"picture_ocr", "background_image_ocr"}
        and value.get("confidence") == "low"
        and type(value.get("attempted")) is bool
        and _bounded_text(value.get("engine"), maximum=128, allow_none=True)
        and _bounded_text(value.get("engine_version"), maximum=128, allow_none=True)
        and _bounded_text(value.get("reason"), maximum=512, allow_none=True)
        and (
            confidence is None
            or _finite_number_between(confidence, minimum=0, maximum=100)
        )
        and isinstance(receipts, list)
        and len(receipts) <= 65_536
        and all(_valid_ocr_receipt(receipt) for receipt in receipts)
        and _valid_ocr_channel_aggregate(cast(dict[str, object], value))
    )


def _valid_unsupported_content(value: object) -> bool:
    if not isinstance(value, dict):
        return False
    fields = set(value)
    if fields not in {
        _UNSUPPORTED_CONTENT_FIELDS,
        _UNSUPPORTED_CONTENT_FIELDS | {"graphic_data_uri"},
    }:
        return False
    shape_name = value.get("shape_name")
    expected_reason = (
        "visible text or labels may not be represented in PPTX text frames"
        if isinstance(shape_name, str)
        else "background image bytes failed package CRC validation"
    )
    return (
        _bounded_text(value.get("content_type"), maximum=256, allow_empty=False)
        and _bounded_text(shape_name, maximum=4096, allow_none=True)
        and _valid_string_path(value.get("shape_path"), allow_empty=True)
        and value.get("reason") == expected_reason
        and value.get("render_required") is True
        and (
            "graphic_data_uri" not in value
            or _bounded_text(value["graphic_data_uri"], maximum=4096, allow_empty=False)
        )
    )


def _valid_native_timing(value: object) -> bool:
    if not isinstance(value, dict) or set(value) != _NATIVE_TIMING_FIELDS:
        return False
    count_fields = (
        "timing_element_count",
        "transition_count",
        "set_action_count",
        "visibility_set_action_count",
        "build_list_count",
    )
    if any(not _plain_nonnegative_int(value.get(field)) for field in count_fields):
        return False
    if value["visibility_set_action_count"] > value["set_action_count"]:
        return False
    animation = value.get("animation_behavior_counts")
    media = value.get("media_timing_counts")
    builds = value.get("build_entry_counts")
    provenance = value.get("provenance")
    if (
        not isinstance(animation, dict)
        or not _valid_count_map(
            animation,
            required_keys=_TIMING_ANIMATION_FIELDS,
            total_key=True,
        )
        or not isinstance(media, dict)
        or not _valid_count_map(
            media,
            required_keys=_TIMING_MEDIA_FIELDS,
            total_key=True,
        )
        or not isinstance(builds, dict)
        or not _valid_count_map(
            builds,
            required_keys=_TIMING_BUILD_FIELDS,
            total_key=True,
        )
        or not isinstance(provenance, dict)
        or set(provenance) != {*_TIMING_PROVENANCE, "part_name"}
        or any(provenance.get(key) != item for key, item in _TIMING_PROVENANCE.items())
        or not _bounded_text(
            provenance.get("part_name"), maximum=2048, allow_empty=False
        )
    ):
        return False
    if (
        (
            value["set_action_count"] > 0
            or animation["total"] > 0
            or media["total"] > 0
            or value["build_list_count"] > 0
        )
        and value["timing_element_count"] == 0
    ) or (builds["total"] > 0 and value["build_list_count"] == 0):
        return False
    boolean_expectations = {
        "timing_element_present": value["timing_element_count"] > 0,
        "build_list_present": value["build_list_count"] > 0,
        "has_animation_behaviors": animation["total"] > 0,
        "has_media_timing": media["total"] > 0,
        "has_build_entries": builds["total"] > 0,
    }
    return all(
        value.get(field) is expected for field, expected in boolean_expectations.items()
    )


def _valid_native_text_channel_bindings(
    *,
    shapes: list[dict[str, object]],
    channels: list[dict[str, object]],
    unsupported: list[dict[str, object]],
    slide_height_inches: object,
) -> tuple[bool, str, str]:
    """Bind every native text channel and compatibility aggregate to shapes."""
    shape_bindings = [shape for shape in shapes if shape["has_text_frame"] is True]
    table_bindings = [
        shape for shape in shapes if shape.get("graphic_frame_type") == "table"
    ]
    group_bindings = [shape for shape in shapes if shape["shape_type"] == "GROUP (6)"]
    unsupported_bindings = [
        item for item in unsupported if isinstance(item.get("shape_name"), str)
    ]
    shape_channels = [
        channel for channel in channels if channel["channel"] == "shape_text"
    ]
    table_channels = [
        channel for channel in channels if channel["channel"] == "table_cell_text"
    ]
    group_channels = [
        channel for channel in channels if channel["channel"] == "group_container_text"
    ]
    reserved = {
        "shape_text",
        "table_cell_text",
        "group_container_text",
        "picture_ocr",
        "background_image_ocr",
    }
    unsupported_channels = [
        channel for channel in channels if channel["channel"] not in reserved
    ]
    if (
        len(shape_channels) != len(shape_bindings)
        or len(table_channels) != len(table_bindings)
        or len(group_channels) != len(group_bindings)
        or len(unsupported_channels) != len(unsupported_bindings)
    ):
        return False, "", ""

    for channel, shape in zip(shape_channels, shape_bindings):
        text = cast(str, shape["text_preview"])
        provenance = cast(dict[str, object], channel["provenance"])
        if (
            channel["text"] != text
            or channel["confidence"]
            != ("medium" if cast(int, shape["group_depth"]) > 0 else "high")
            or channel["status"] != ("extracted" if text.strip() else "empty")
            or set(provenance) != {"source", "shape_path"}
            or provenance["source"] != "pptx_shape_text_frame"
            or provenance["shape_path"] != shape["shape_path"]
        ):
            return False, "", ""

    for channel, shape in zip(table_channels, table_bindings):
        text = cast(str, shape.get("table_text_preview", ""))
        provenance = cast(dict[str, object], channel["provenance"])
        cells = provenance.get("cells")
        rows = cast(int, shape.get("table_rows", 0))
        columns = cast(int, shape.get("table_columns", 0))
        positions: list[tuple[int, int]] = []
        if isinstance(cells, list):
            for cell in cells:
                match = re.fullmatch(r"R([1-9][0-9]*)C([1-9][0-9]*)", cast(str, cell))
                if match is None:
                    return False, "", ""
                position = (int(match.group(1)), int(match.group(2)))
                if position[0] > rows or position[1] > columns:
                    return False, "", ""
                positions.append(position)
        if (
            channel["text"] != text
            or channel["confidence"] != "medium"
            or channel["status"] != ("extracted" if text else "empty")
            or set(provenance) != {"source", "shape_path", "cells"}
            or provenance["source"] != "pptx_table_cells"
            or provenance["shape_path"] != shape["shape_path"]
            or positions != sorted(set(positions))
            or bool(positions) is not bool(text)
        ):
            return False, "", ""

    for channel, shape in zip(group_channels, group_bindings):
        provenance = cast(dict[str, object], channel["provenance"])
        if (
            channel["text"] != ""
            or channel["confidence"] != "low"
            or channel["status"] != "requires_render"
            or set(provenance) != {"source", "shape_path"}
            or provenance["source"] != "pptx_group_container"
            or provenance["shape_path"] != shape["shape_path"]
        ):
            return False, "", ""

    for channel, item in zip(unsupported_channels, unsupported_bindings):
        provenance = cast(dict[str, object], channel["provenance"])
        expected_provenance: dict[str, object] = {
            "source": "pptx_unsupported_visual_container",
            "shape_path": item["shape_path"],
        }
        if "graphic_data_uri" in item:
            expected_provenance["graphic_data_uri"] = item["graphic_data_uri"]
        if (
            channel["channel"] != f"{item['content_type']}_text"
            or channel["text"] != ""
            or channel["confidence"] != "low"
            or channel["status"] != "unsupported"
            or provenance != expected_provenance
        ):
            return False, "", ""

    native_text_channels = [
        channel
        for channel in channels
        if channel["channel"] in {"shape_text", "table_cell_text"}
    ]
    preview = " | ".join(
        cast(str, channel["text"])[:50]
        for channel in native_text_channels
        if cast(str, channel["text"]).strip()
    )[:200]
    footer = ""
    slide_height = float(cast(float, slide_height_inches))
    for channel, shape in zip(shape_channels, shape_bindings):
        top = shape.get("top")
        if (
            shape["group_depth"] == 0
            and top is not None
            and float(cast(float, top)) > slide_height * 0.85
        ):
            footer = cast(str, channel["text"])
    return True, preview, footer


def _valid_slide_visual(
    value: object,
    *,
    slide_number: int,
    slide_width_inches: object,
    slide_height_inches: object,
    expected_ocr: bool | None,
) -> bool:
    if not isinstance(value, dict) or set(value) != _SLIDE_VISUAL_FIELDS:
        return False
    shape_count = value.get("shape_count")
    recursive_count = value.get("shape_count_recursive")
    shapes = value.get("shapes_summary")
    channels = value.get("text_channels")
    unsupported = value.get("unsupported_content")
    reasons = value.get("render_required_reasons")
    background = value.get("background_color_hex")
    background_type = value.get("background_type")
    background_asset_status = value.get("background_asset_status")
    background_part_name = value.get("background_part_name")
    background_asset_sha256 = value.get("background_asset_sha256")
    confidence_label = value.get("text_extraction_confidence")
    extraction_method = value.get("text_extraction_method")
    if not _plain_nonnegative_int(shape_count) or not _plain_nonnegative_int(
        recursive_count
    ):
        return False
    shape_count_int = cast(int, shape_count)
    recursive_count_int = cast(int, recursive_count)
    if (
        value.get("slide_number") != slide_number
        or not _bounded_text(
            value.get("slide_part_name"), maximum=2048, allow_empty=False
        )
        or value.get("slide_part_name") != f"ppt/slides/slide{slide_number}.xml"
        or recursive_count_int < shape_count_int
        or not isinstance(shapes, list)
        or len(shapes) != recursive_count_int
        or any(not _valid_shape_summary(shape) for shape in shapes)
        or not _valid_shape_topology(cast(list[dict[str, object]], shapes))
        or sum(shape["group_depth"] == 0 for shape in shapes) != shape_count_int
        or not isinstance(channels, list)
        or len(channels) > (4 * recursive_count_int) + 4
        or any(not _valid_text_channel(channel) for channel in channels)
        or not isinstance(unsupported, list)
        or len(unsupported) > recursive_count_int + 1
        or any(not _valid_unsupported_content(item) for item in unsupported)
        or not isinstance(reasons, list)
        or any(
            not _bounded_text(reason, maximum=256, allow_empty=False)
            for reason in reasons
        )
        or reasons != sorted(set(reasons))
        or (
            background is not None
            and (
                not isinstance(background, str)
                or re.fullmatch(r"#[0-9A-F]{6}", background) is None
            )
        )
        or not isinstance(background_type, str)
        or background_type
        not in {
            "solid",
            "pattern",
            "image",
            "gradient",
            "solid_from_layout",
            "solid_from_master",
            "unknown",
        }
        or not _bounded_text(value.get("layout_name"), maximum=4096, allow_none=True)
        or not _finite_number_between(
            value.get("image_area_ratio"), minimum=0, maximum=1
        )
        or not isinstance(confidence_label, str)
        or confidence_label not in {"high", "low"}
        or not isinstance(extraction_method, str)
        or extraction_method not in {"shapes", "shapes+ocr", "shapes+ocr_unavailable"}
        or not _bounded_text(value.get("text_content_preview"), maximum=200)
        or not _bounded_text(value.get("ocr_text"), maximum=8000)
        or not _bounded_text(value.get("footer_text"))
        or not _valid_native_timing(value.get("native_timing"))
    ):
        return False
    native_timing = cast(dict[str, object], value["native_timing"])
    timing_provenance = cast(dict[str, object], native_timing["provenance"])
    if timing_provenance["part_name"] != value["slide_part_name"]:
        return False
    if background_type == "image":
        if background_asset_status not in {"available", "corrupt", "unavailable"}:
            return False
        if background_asset_status in {"available", "corrupt"}:
            if (
                not _valid_opc_part_name(background_part_name)
                or not isinstance(background_asset_sha256, str)
                or _SHA256_RE.fullmatch(background_asset_sha256) is None
            ):
                return False
        elif background_part_name is not None or background_asset_sha256 is not None:
            return False
    elif (
        background_asset_status != "not_applicable"
        or background_part_name is not None
        or background_asset_sha256 is not None
    ):
        return False
    for field in (
        "has_text_frame_shapes",
        "has_extracted_text",
        "has_image",
        "has_unsupported_content",
        "render_required",
        "has_speaker_notes",
    ):
        if type(value.get(field)) is not bool:
            return False
    typed_shapes = cast(list[dict[str, object]], shapes)
    typed_channels = cast(list[dict[str, object]], channels)
    typed_unsupported = cast(list[dict[str, object]], unsupported)
    typed_reasons = cast(list[str], reasons)
    native_text_valid, expected_preview, expected_footer = (
        _valid_native_text_channel_bindings(
            shapes=typed_shapes,
            channels=typed_channels,
            unsupported=typed_unsupported,
            slide_height_inches=slide_height_inches,
        )
    )
    if not native_text_valid:
        return False
    if not _picture_ratio_matches_geometry(
        value["image_area_ratio"],
        shapes=typed_shapes,
        slide_width_inches=slide_width_inches,
        slide_height_inches=slide_height_inches,
    ):
        return False
    expected_reasons = {cast(str, item["content_type"]) for item in typed_unsupported}
    if value["background_type"] == "image":
        expected_reasons.add("background_image")
    if cast(float, value["image_area_ratio"]) >= PPTX_TEXT_BEARING_IMAGE_AREA_RATIO:
        expected_reasons.add("large_picture")
    if any(shape["shape_type"] == "GROUP (6)" for shape in typed_shapes):
        expected_reasons.add("grouped_shapes")
    if any(shape.get("graphic_frame_type") == "table" for shape in typed_shapes):
        expected_reasons.add("table")

    derived_unsupported: list[tuple[str, tuple[str, ...], str, str | None]] = []
    unsupported_shape_prefixes = {
        "DIAGRAM",
        "IGX_GRAPHIC",
        "EMBEDDED_OLE_OBJECT",
        "LINKED_OLE_OBJECT",
        "MEDIA",
        "WEB_VIDEO",
    }
    for shape in typed_shapes:
        graphic_kind = shape.get("graphic_frame_type")
        picture_status = shape.get("picture_asset_status")
        if picture_status == "corrupt":
            content_type = "corrupt_embedded_asset"
        elif picture_status == "unavailable":
            content_type = "unreadable_picture"
        elif isinstance(graphic_kind, str) and graphic_kind != "table":
            content_type = graphic_kind
        else:
            shape_prefix = cast(str, shape["shape_type"]).split(" ", 1)[0]
            if shape_prefix not in unsupported_shape_prefixes:
                continue
            content_type = shape_prefix.lower()
        derived_unsupported.append(
            (
                content_type,
                tuple(cast(list[str], shape["shape_path"])),
                cast(str, shape["name"]),
                (
                    cast(str, shape["graphic_data_uri"])
                    if isinstance(shape.get("graphic_data_uri"), str)
                    else None
                ),
            )
        )
    actual_unsupported = [
        (
            cast(str, item["content_type"]),
            tuple(cast(list[str], item["shape_path"])),
            cast(str, item["shape_name"]),
            (
                cast(str, item["graphic_data_uri"])
                if isinstance(item.get("graphic_data_uri"), str)
                else None
            ),
        )
        for item in typed_unsupported
        if isinstance(item.get("shape_name"), str)
    ]
    unnamed_unsupported = [
        item for item in typed_unsupported if item.get("shape_name") is None
    ]
    expected_unnamed_unsupported = 1 if background_asset_status == "corrupt" else 0
    if (
        Counter(derived_unsupported) != Counter(actual_unsupported)
        or len(unnamed_unsupported) != expected_unnamed_unsupported
        or any(
            item.get("content_type") != "corrupt_embedded_asset"
            or item.get("shape_path") != []
            or "graphic_data_uri" in item
            for item in unnamed_unsupported
        )
    ):
        return False

    ocr_channels = [
        channel
        for channel in typed_channels
        if channel["channel"] in {"picture_ocr", "background_image_ocr"}
    ]
    if expected_ocr is False and any(
        channel["attempted"] is True
        or any(
            receipt["result_status"] != "skipped"
            for receipt in cast(list[dict[str, object]], channel["ocr_receipts"])
        )
        for channel in ocr_channels
    ):
        return False
    if expected_ocr is True and any(
        channel["status"] == "skipped"
        or any(
            receipt["result_status"] == "skipped"
            for receipt in cast(list[dict[str, object]], channel["ocr_receipts"])
        )
        for channel in ocr_channels
    ):
        return False
    ocr_channel_names = [cast(str, channel["channel"]) for channel in ocr_channels]
    if ocr_channel_names != [
        name
        for name in ("picture_ocr", "background_image_ocr")
        if name in ocr_channel_names
    ]:
        return False
    expected_has_image = any(shape["is_picture"] is True for shape in typed_shapes)
    picture_channel = next(
        (channel for channel in ocr_channels if channel["channel"] == "picture_ocr"),
        None,
    )
    background_channel = next(
        (
            channel
            for channel in ocr_channels
            if channel["channel"] == "background_image_ocr"
        ),
        None,
    )
    if (picture_channel is not None) is not (
        expected_has_image and bool(expected_reasons)
    ):
        return False
    if (background_channel is not None) is not (background_type == "image"):
        return False
    if background_channel is not None:
        background_receipts = cast(
            list[dict[str, object]],
            background_channel["ocr_receipts"],
        )
        expected_background_receipts = (
            Counter(
                {
                    (
                        (),
                        cast(str, background_part_name),
                        cast(str, background_asset_sha256),
                    ): 1
                }
            )
            if background_asset_status == "available"
            else Counter()
        )
        actual_background_receipts = Counter(
            (
                tuple(cast(list[str], receipt["shape_path"])),
                cast(str, receipt["part_name"]),
                cast(str, receipt["asset_sha256"]),
            )
            for receipt in background_receipts
        )
        provenance = cast(dict[str, object], background_channel["provenance"])
        expected_provenance_fields = {"source", "background_owner"}
        if background_asset_status in {"available", "corrupt"}:
            expected_provenance_fields.update({"relationship_id", "part_name"})
        elif "relationship_id" in provenance:
            expected_provenance_fields.add("relationship_id")
        if background_asset_status == "corrupt":
            expected_provenance_fields.add("asset_status")
        if (
            actual_background_receipts != expected_background_receipts
            or set(provenance) != expected_provenance_fields
            or provenance.get("source") != "pptx_background_image"
            or provenance.get("background_owner") not in {"slide", "layout", "master"}
            or (
                background_asset_status in {"available", "corrupt"}
                and provenance.get("part_name") != background_part_name
            )
            or (
                background_asset_status == "corrupt"
                and provenance.get("asset_status") != "recovered_with_placeholder"
            )
        ):
            return False
    if picture_channel is not None:
        expected_picture_receipts = Counter(
            (
                tuple(cast(list[str], shape["shape_path"])),
                cast(str, shape["picture_part_name"]),
                cast(str, shape["picture_asset_sha256"]),
            )
            for shape in typed_shapes
            if shape["is_picture"] is True
            and shape.get("picture_asset_status") == "available"
        )
        actual_picture_receipts = Counter(
            (
                tuple(cast(list[str], receipt["shape_path"])),
                cast(str, receipt["part_name"]),
                cast(str, receipt["asset_sha256"]),
            )
            for receipt in cast(
                list[dict[str, object]], picture_channel["ocr_receipts"]
            )
        )
        if actual_picture_receipts != expected_picture_receipts:
            return False
    expected_ocr_text = ""
    expected_method = "shapes"
    for channel in ocr_channels:
        channel_text = cast(str, channel["text"])
        if channel_text:
            expected_ocr_text = (
                f"{expected_ocr_text} | {channel_text}"
                if expected_ocr_text
                else channel_text
            )[:8000]
        receipts = cast(list[dict[str, object]], channel["ocr_receipts"])
        if receipts:
            statuses = {receipt["result_status"] for receipt in receipts}
            if any(receipt["attempted"] is True for receipt in receipts):
                expected_method = "shapes+ocr"
            elif statuses == {"unavailable"}:
                expected_method = "shapes+ocr_unavailable"

    expected_has_text_frames = any(
        shape["has_text_frame"] is True for shape in typed_shapes
    )
    return (
        value["has_extracted_text"]
        is any(cast(str, channel["text"]).strip() for channel in typed_channels)
        and value["has_text_frame_shapes"] is expected_has_text_frames
        and value["has_image"] is expected_has_image
        and (expected_has_image or value["image_area_ratio"] == 0)
        and value["has_unsupported_content"] is bool(typed_unsupported)
        and typed_reasons == sorted(expected_reasons)
        and value["render_required"] is bool(expected_reasons)
        and value["text_extraction_confidence"]
        == ("low" if expected_reasons else "high")
        and value["ocr_text"] == expected_ocr_text
        and value["text_extraction_method"] == expected_method
        and value["text_content_preview"] == expected_preview
        and value["footer_text"] == expected_footer
    )


def _expected_native_timing_summary(
    slides: list[dict[str, object]],
) -> dict[str, object]:
    counts = {field: 0 for field in _NATIVE_TIMING_SUMMARY_COUNT_FIELDS}
    animation = {field: 0 for field in _TIMING_ANIMATION_FIELDS}
    media = {field: 0 for field in _TIMING_MEDIA_FIELDS}
    builds = {field: 0 for field in _TIMING_BUILD_FIELDS}
    for slide in slides:
        timing = cast(dict[str, object], slide["native_timing"])
        counts["slides_with_timing_elements"] += int(
            cast(bool, timing["timing_element_present"])
        )
        counts["slides_with_transitions"] += int(
            cast(int, timing["transition_count"]) > 0
        )
        counts["slides_with_animation_behaviors"] += int(
            cast(bool, timing["has_animation_behaviors"])
        )
        counts["slides_with_media_timing"] += int(
            cast(bool, timing["has_media_timing"])
        )
        counts["slides_with_build_lists"] += int(
            cast(bool, timing["build_list_present"])
        )
        counts["slides_with_build_entries"] += int(
            cast(bool, timing["has_build_entries"])
        )
        for field in (
            "timing_element_count",
            "transition_count",
            "set_action_count",
            "visibility_set_action_count",
            "build_list_count",
        ):
            counts[field] += cast(int, timing[field])
        for target, source_name, fields in (
            (animation, "animation_behavior_counts", _TIMING_ANIMATION_FIELDS),
            (media, "media_timing_counts", _TIMING_MEDIA_FIELDS),
            (builds, "build_entry_counts", _TIMING_BUILD_FIELDS),
        ):
            source = cast(dict[str, int], timing[source_name])
            for field in fields:
                target[field] += source[field]
    for target in (animation, media, builds):
        target["total"] = sum(target.values())
    return {
        **counts,
        "animation_behavior_counts": animation,
        "media_timing_counts": media,
        "build_entry_counts": builds,
        "provenance": dict(_TIMING_PROVENANCE),
    }


def _valid_recovery_asset_bindings(
    slides: list[dict[str, object]],
    recovery: tuple[dict[str, object], ...],
) -> bool:
    """Bind corrupt/available visual assets to exact archive recovery parts."""
    recovered_assets = {
        (
            _opc_member_name_key(cast(str, record["part_name"])),
            cast(str, record["replacement_sha256"]),
        )
        for record in recovery
    }
    recovered_part_names = {part_name for part_name, _digest in recovered_assets}
    digest_by_part = {
        _opc_member_name_key(part_name): digest
        for part_name, digest in recovered_assets
    }
    for slide in slides:
        bound_assets = [
            (
                shape.get("picture_asset_status"),
                shape.get("picture_part_name"),
                shape.get("picture_asset_sha256"),
            )
            for shape in cast(list[dict[str, object]], slide["shapes_summary"])
            if shape["is_picture"] is True
        ]
        if slide["background_type"] == "image":
            bound_assets.append(
                (
                    slide["background_asset_status"],
                    slide["background_part_name"],
                    slide["background_asset_sha256"],
                )
            )
        for status, part_name, digest in bound_assets:
            if isinstance(part_name, str) and isinstance(digest, str):
                part_key = _opc_member_name_key(part_name)
                existing_digest = digest_by_part.setdefault(part_key, digest)
                if existing_digest != digest:
                    return False
            else:
                part_key = None
            if status == "corrupt" and (part_key, digest) not in recovered_assets:
                return False
            if status == "available" and part_key in recovered_part_names:
                return False
    return True


def _validate_nested_extraction(
    extraction: dict[str, object],
    *,
    slide_count: int,
    recovery: tuple[dict[str, object], ...],
    audit: dict[str, object],
    expected_ocr: bool | None,
) -> None:
    dimensions = (
        extraction.get("slide_width_inches"),
        extraction.get("slide_height_inches"),
    )
    per_slide = extraction.get("per_slide_visual")
    if (
        any(not _finite_number_between(value, minimum=0.01) for value in dimensions)
        or not _aspect_ratio_matches_dimensions(
            extraction.get("aspect_ratio"),
            width_inches=dimensions[0],
            height_inches=dimensions[1],
        )
        or not _valid_template_layouts(extraction.get("template_layouts"))
        or not isinstance(per_slide, list)
        or any(
            not _valid_slide_visual(
                slide,
                slide_number=index,
                slide_width_inches=dimensions[0],
                slide_height_inches=dimensions[1],
                expected_ocr=expected_ocr,
            )
            for index, slide in enumerate(per_slide, start=1)
        )
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    slides = cast(list[dict[str, object]], per_slide)
    slide_part_names = [cast(str, slide["slide_part_name"]) for slide in slides]
    if len(slide_part_names) != len(set(slide_part_names)):
        raise _probe_failure("pptx_probe_malformed_result")
    layouts = cast(list[dict[str, object]], extraction["template_layouts"])
    layout_names = {layout["name"] for layout in layouts}
    if any(
        slide["layout_name"] is not None and slide["layout_name"] not in layout_names
        for slide in slides
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    if not _valid_recovery_asset_bindings(slides, recovery):
        raise _probe_failure("pptx_probe_malformed_result")
    corrupt_assets = [
        {
            "part_name": record["part_name"],
            "error_type": "crc_mismatch",
            "status": "recovered_with_placeholder",
        }
        for record in recovery
    ]
    if extraction.get("corrupt_assets") != corrupt_assets:
        raise _probe_failure("pptx_probe_malformed_result")
    summary = extraction.get("native_timing_summary")
    if (
        not isinstance(summary, dict)
        or set(summary) != _NATIVE_TIMING_SUMMARY_FIELDS
        or summary != _expected_native_timing_summary(slides)
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    expected_reasons = {
        str(slide["slide_number"]): list(
            cast(list[str], slide["render_required_reasons"])
        )
        for slide in slides
        if slide["render_required"] is True
    }
    if audit["render_required_reasons"] != expected_reasons or audit[
        "render_required_slide_numbers"
    ] != [int(number) for number in expected_reasons]:
        raise _probe_failure("pptx_probe_malformed_result")
    expected_fonts: dict[str, int] = {}
    expected_shape_types: dict[str, int] = {}
    expected_backgrounds: dict[str, int] = {}
    expected_color_sequence: list[str] = []
    for slide in slides:
        background = slide["background_color_hex"]
        expected_color_sequence.append(cast(str, background or "unknown"))
        if isinstance(background, str):
            expected_backgrounds[background] = (
                expected_backgrounds.get(background, 0) + 1
            )
        for shape in cast(list[dict[str, object]], slide["shapes_summary"]):
            font_name = shape.get("font_name")
            if isinstance(font_name, str) and font_name:
                expected_fonts[font_name] = expected_fonts.get(font_name, 0) + 1
            table_fonts = shape.get("table_fonts")
            if isinstance(table_fonts, dict):
                for name, count in table_fonts.items():
                    expected_fonts[name] = expected_fonts.get(name, 0) + cast(
                        int, count
                    )
            shape_type = shape.get("auto_shape_type")
            if isinstance(shape_type, str):
                expected_shape_types[shape_type] = (
                    expected_shape_types.get(shape_type, 0) + 1
                )
    global_design = extraction.get("global_design")
    if (
        not isinstance(global_design, dict)
        or set(global_design)
        != {"fonts_used", "background_colors", "shape_types_used", "color_sequence"}
        or not _valid_count_map(global_design.get("fonts_used"))
        or not _valid_count_map(global_design.get("background_colors"))
        or not _valid_count_map(global_design.get("shape_types_used"))
        or global_design.get("fonts_used") != expected_fonts
        or global_design.get("background_colors") != expected_backgrounds
        or global_design.get("shape_types_used") != expected_shape_types
        or global_design.get("color_sequence") != expected_color_sequence
    ):
        raise _probe_failure("pptx_probe_malformed_result")


def _decode_extraction_worker_payload(
    value: object,
    *,
    expected_ocr: bool | None = None,
) -> dict[str, object]:
    if not isinstance(value, dict):
        raise _probe_failure("pptx_probe_malformed_result")
    if value.get("schema_version") != PPTX_ARTIFACT_PROBE_SCHEMA_VERSION:
        raise _probe_failure("pptx_probe_malformed_result")
    if value.get("status") == "unavailable":
        if set(value) != {"schema_version", "status", "reason_code", "details"}:
            raise _probe_failure("pptx_probe_malformed_result")
        reason = value.get("reason_code")
        details = value.get("details")
        if (
            not isinstance(reason, str)
            or reason not in _CHILD_PROBE_REASON_CODES
            or not isinstance(details, dict)
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        if set(details) - {"part_names", "exception_type"}:
            raise _probe_failure("pptx_probe_malformed_result")
        part_names = details.get("part_names")
        if part_names is not None and (
            not isinstance(part_names, list)
            or len(part_names) > PPTX_ARTIFACT_PROBE_MAX_RECOVERY_RECORDS
            or any(
                not isinstance(name, str) or not name or len(name) > 2048
                for name in part_names
            )
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        exception_type = details.get("exception_type")
        if exception_type is not None and (
            not isinstance(exception_type, str)
            or not exception_type
            or len(exception_type) > 128
        ):
            raise _probe_failure("pptx_probe_malformed_result")
        raise _probe_failure(str(reason), details=details)
    if (
        set(value) != {"schema_version", "status", "extraction"}
        or value.get("status") != "available"
        or not isinstance(value.get("extraction"), dict)
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    extraction = copy.deepcopy(value["extraction"])
    slide_count = extraction.get("slide_count")
    fingerprint = extraction.get("input_fingerprint")
    per_slide = extraction.get("per_slide_visual")
    if (
        set(extraction) != _EXTRACTION_FIELDS
        or not _bounded_text(
            extraction.get("pptx_path"), maximum=32_768, allow_empty=False
        )
        or extraction.get("schema_version") != PPTX_EXTRACTION_SCHEMA_VERSION
        or extraction.get("pipeline_version") != PPTX_EXTRACTION_PIPELINE_VERSION
        or isinstance(slide_count, bool)
        or not isinstance(slide_count, int)
        or not 1 <= slide_count <= PPTX_ARCHIVE_MAX_MEMBERS
        or not isinstance(fingerprint, dict)
        or set(fingerprint) != {"algorithm", "digest", "size_bytes"}
        or fingerprint.get("algorithm") != "sha256"
        or not isinstance(fingerprint.get("digest"), str)
        or _SHA256_RE.fullmatch(str(fingerprint.get("digest"))) is None
        or isinstance(fingerprint.get("size_bytes"), bool)
        or not isinstance(fingerprint.get("size_bytes"), int)
        or cast(int, fingerprint.get("size_bytes")) < 1
        or not isinstance(per_slide, list)
        or len(per_slide) != slide_count
        or any(
            not isinstance(slide, dict) or slide.get("slide_number") != index
            for index, slide in enumerate(per_slide, start=1)
        )
        or not isinstance(extraction.get("corrupt_assets"), list)
        or not isinstance(extraction.get("template_layouts"), list)
        or not isinstance(extraction.get("global_design"), dict)
        or not isinstance(extraction.get("native_timing_summary"), dict)
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    recovery = _validated_recovery_records(extraction.get("archive_recovery"))
    try:
        audit = validate_native_deck_audit(
            extraction.get("native_deck_audit"),
            slide_count=slide_count,
        )
    except PptxEvidenceError as exc:
        raise _probe_failure("pptx_probe_malformed_result") from exc
    if (
        audit["source_pptx_sha256"] != fingerprint["digest"]
        or audit["source_pptx_size_bytes"] != fingerprint["size_bytes"]
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    _validate_nested_extraction(
        extraction,
        slide_count=slide_count,
        recovery=recovery,
        audit=audit,
        expected_ocr=expected_ocr,
    )
    extraction["archive_recovery"] = [dict(item) for item in recovery]
    extraction["native_deck_audit"] = audit
    return extraction


def _validate_extraction_render_binding(
    extraction: Mapping[str, object],
    *,
    rendered_generation: FileGeneration | None,
    requested_ranges: list[list[int]],
) -> None:
    audit = extraction.get("native_deck_audit")
    if not isinstance(audit, Mapping):
        raise _probe_failure("pptx_probe_malformed_result")
    receipt = audit.get("rendered_page_inspection")
    if rendered_generation is None:
        if receipt is not None:
            raise _probe_failure("pptx_probe_malformed_result")
        return
    if (
        not isinstance(receipt, Mapping)
        or receipt.get("rendered_pdf_size_bytes") != rendered_generation.size
        or receipt.get("inspected_page_ranges") != requested_ranges
    ):
        raise _probe_failure("pptx_probe_malformed_result")


def run_supervised_pptx_extraction(
    pptx_path: str | Path,
    *,
    trusted_root: str | Path | None = None,
    ocr: bool = True,
    rendered_pdf_path: str | Path | None = None,
    inspected_page_ranges: list[list[int]] | None = None,
    source_size_limit_bytes: int | None = None,
    deadline_monotonic: float | None = None,
) -> dict[str, object]:
    """Extract one deck and bind failures to its exact admitted generation."""
    admitted_source_sizes: list[int] = []
    try:
        return _run_supervised_pptx_extraction_impl(
            pptx_path,
            trusted_root=trusted_root,
            ocr=ocr,
            rendered_pdf_path=rendered_pdf_path,
            inspected_page_ranges=inspected_page_ranges,
            source_size_limit_bytes=source_size_limit_bytes,
            deadline_monotonic=deadline_monotonic,
            admitted_source_sizes=admitted_source_sizes,
        )
    except PptxEvidenceError as exc:
        if admitted_source_sizes:
            exc.details["admitted_source_size_bytes"] = admitted_source_sizes[0]
        raise


def _run_supervised_pptx_extraction_impl(
    pptx_path: str | Path,
    *,
    trusted_root: str | Path | None,
    ocr: bool,
    rendered_pdf_path: str | Path | None,
    inspected_page_ranges: list[list[int]] | None,
    source_size_limit_bytes: int | None,
    deadline_monotonic: float | None,
    admitted_source_sizes: list[int],
) -> dict[str, object]:
    """Extract one deck only through the authenticated worker boundary."""
    if type(ocr) is not bool:
        raise PptxEvidenceError(
            "ocr must be a boolean",
            reason_code="pptx_evidence_invalid",
        )
    root = (
        Path(os.path.abspath(os.fspath(trusted_root)))
        if trusted_root is not None
        else None
    )
    artifact, generation, root_generation = _admit_supervised_input(
        pptx_path,
        label="PPTX artifact",
        trusted_root=root,
        deadline_monotonic=deadline_monotonic,
    )
    admitted_source_sizes.append(generation.size)
    if source_size_limit_bytes is not None:
        if (
            isinstance(source_size_limit_bytes, bool)
            or not isinstance(source_size_limit_bytes, int)
            or source_size_limit_bytes < 1
        ):
            raise PptxEvidenceError(
                "source_size_limit_bytes must be a positive integer",
                reason_code="pptx_evidence_invalid",
            )
        if generation.size > source_size_limit_bytes:
            raise PptxEvidenceError(
                "PPTX source exceeds the remaining aggregate input budget",
                reason_code="pptx_batch_input_limit",
                details={"limit_bytes": source_size_limit_bytes},
            )
    expected = {"pptx": generation}
    if root is not None:
        if root_generation is None:
            raise _probe_failure("pptx_probe_malformed_result")
        expected["pptx_root"] = root_generation
    rendered_artifact: Path | None = None
    rendered_generation: FileGeneration | None = None
    if rendered_pdf_path is not None:
        rendered_artifact, rendered_generation, rendered_root_generation = (
            _admit_supervised_input(
                rendered_pdf_path,
                label="rendered PDF",
                deadline_monotonic=deadline_monotonic,
            )
        )
        if rendered_root_generation is not None:  # pragma: no cover - no root given
            raise _probe_failure("pptx_probe_malformed_result")
        expected["rendered_pdf"] = rendered_generation
    ranges = _validated_requested_page_ranges(inspected_page_ranges or [])
    if ranges and rendered_artifact is None:
        raise PptxEvidenceError(
            "inspected_page_ranges requires rendered_pdf_path",
            reason_code="pptx_evidence_invalid",
        )
    payload: dict[str, object] = {
        "pptx_path": os.fspath(artifact),
        "trusted_root": os.fspath(root) if root is not None else None,
        "ocr": ocr,
        "rendered_pdf_path": (
            os.fspath(rendered_artifact) if rendered_artifact is not None else None
        ),
        "inspected_page_ranges": ranges,
        "extraction_schema_version": PPTX_EXTRACTION_SCHEMA_VERSION,
        "extraction_pipeline_version": PPTX_EXTRACTION_PIPELINE_VERSION,
    }
    base_limits = PPTX_EXTRACT_OCR_LIMITS if ocr else PPTX_EXTRACT_NO_OCR_LIMITS
    limits = _limits_before_deadline(base_limits, deadline_monotonic)
    deadline_limited = limits.wall_seconds < base_limits.wall_seconds
    command = [
        sys.executable,
        os.fspath(Path(__file__).absolute()),
        PPTX_SUPERVISED_WORKER_FLAG,
    ]
    sensitive: list[Path] = [artifact]
    if root is not None:
        sensitive.append(root)
    if rendered_artifact is not None:
        sensitive.append(rendered_artifact)
    try:
        worker_result = run_authenticated_worker(
            command,
            PPTX_EXTRACT_OPERATION,
            expected,
            cast(Any, payload),
            limits,
            sensitive_values=sensitive,
            schema_generation=PPTX_EXTRACTION_SCHEMA_VERSION,
            pipeline_generation=PPTX_EXTRACTION_PIPELINE_VERSION,
        )
        extraction = _decode_extraction_worker_payload(
            worker_result.payload,
            expected_ocr=ocr,
        )
        _validate_extraction_render_binding(
            extraction,
            rendered_generation=rendered_generation,
            requested_ranges=ranges,
        )
    except SupervisorError as exc:
        if deadline_limited and exc.reason_code == "worker_timeout":
            raise PptxEvidenceError(
                "PPTX batch wall deadline expired inside the bounded worker",
                reason_code="pptx_batch_wall_limit",
            ) from exc
        raise _supervisor_probe_failure(
            exc, timeout_seconds=limits.wall_seconds
        ) from exc
    except PptxEvidenceError as exc:
        if exc.reason_code in _ARCHIVE_INTEGRITY_CONFIRMATION_REASON_CODES:
            try:
                probe_pptx_artifact(
                    artifact,
                    trusted_root=root,
                    deadline_monotonic=deadline_monotonic,
                )
            except PptxEvidenceError as confirmation:
                if confirmation.reason_code == "pptx_batch_wall_limit":
                    raise
                if confirmation.reason_code in {
                    "pptx_artifact_changed",
                    "pptx_probe_materialization_changed",
                }:
                    raise confirmation from exc
                if confirmation.reason_code == exc.reason_code:
                    raise exc
                raise _probe_failure(
                    "pptx_probe_materialization_changed"
                ) from confirmation
            raise _probe_failure("pptx_probe_materialization_changed") from exc
        raise

    current_receipt = _run_bounded_metadata_worker(
        artifact,
        trusted_root=root,
        deadline_monotonic=deadline_monotonic,
    )
    if (
        current_receipt.generation != generation
        or current_receipt.root_generation != root_generation
    ):
        raise _probe_failure("pptx_artifact_changed")
    if rendered_artifact is not None and rendered_generation is not None:
        rendered_current = _supervised_file_generation(
            rendered_artifact,
            label="rendered PDF",
            deadline_monotonic=deadline_monotonic,
        )
        if rendered_current != rendered_generation:
            raise _probe_failure("pptx_artifact_changed")
    fingerprint = extraction["input_fingerprint"]
    if (
        not isinstance(fingerprint, dict)
        or fingerprint.get("size_bytes") != generation.size
        or extraction.get("pptx_path") != os.fspath(artifact)
    ):
        raise _probe_failure("pptx_probe_malformed_result")
    recovery = extraction.get("archive_recovery")
    if isinstance(recovery, list) and recovery:
        stable_probe = probe_pptx_artifact(
            artifact,
            trusted_root=root,
            deadline_monotonic=deadline_monotonic,
        )
        if (
            stable_probe.source_sha256 != fingerprint.get("digest")
            or stable_probe.source_size_bytes != generation.size
            or [dict(item) for item in stable_probe.archive_recovery] != recovery
        ):
            raise _probe_failure("pptx_probe_materialization_changed")
    # Preserve the caller's public path spelling; directory mode deliberately
    # replaces this with a root-relative record before emitting aggregate JSON.
    extraction["pptx_path"] = os.fspath(pptx_path)
    return extraction


def parse_page_range_arguments(values: list[str] | None) -> list[list[int]]:
    """Parse repeated CLI PAGE or START-END values."""
    if not values:
        return []
    parsed: list[list[int]] = []
    for value in values:
        for token in value.split(","):
            candidate = token.strip()
            match = re.fullmatch(r"(\d+)(?:-(\d+))?", candidate)
            if match is None:
                raise PptxEvidenceError(
                    "--inspected-pages values must be PAGE or START-END, "
                    f"got {candidate!r}"
                )
            start = int(match.group(1))
            end = int(match.group(2) or match.group(1))
            parsed.append([start, end])
    return parsed


def ranges_cover_pages(ranges: object, pages: list[int], *, page_count: int) -> bool:
    normalized = normalize_page_ranges(ranges, page_count=page_count, allow_empty=True)
    if (
        not isinstance(pages, list)
        or len(pages) > PPTX_ARCHIVE_MAX_MEMBERS
        or any(
            isinstance(page, bool)
            or not isinstance(page, int)
            or not 1 <= page <= page_count
            for page in pages
        )
    ):
        return False
    required = sorted(set(pages))
    return _pages_covered_by_ranges(normalized, required) == required


def finite_confidence(value: object) -> float | None:
    """Normalize an OCR confidence to a finite 0..100 float or ``None``."""
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        return None
    numeric = float(value)
    if not math.isfinite(numeric) or numeric < 0 or numeric > 100:
        return None
    return round(numeric, 3)


def _worker_request_payload(request: WorkerRequest) -> dict[str, object]:
    if not isinstance(request.payload, dict):
        raise SupervisorError("invalid_worker_request")
    return dict(request.payload)


def _worker_bound_path(value: object) -> Path:
    if not isinstance(value, str) or not value or "\x00" in value:
        raise SupervisorError("invalid_worker_request")
    return Path(value)


def _worker_generation(
    path: Path,
    *,
    trusted_root: Path | None = None,
) -> FileGeneration:
    try:
        generation, _root_generation, _reparse_tag_value = (
            _metadata_generation_in_worker(
                path,
                trusted_root=trusted_root,
            )
        )
        return generation
    except PptxEvidenceError as exc:
        raise SupervisorError("worker_generation_changed") from exc


def _worker_root_generation(path: Path) -> FileGeneration:
    try:
        snapshot = path.lstat()
    except (OSError, RuntimeError) as exc:
        raise SupervisorError("worker_generation_changed") from exc
    if (
        stat_module.S_ISLNK(snapshot.st_mode)
        or _is_unsupported_reparse(
            snapshot,
            allow_hydrated_cloud_file=False,
        )
        or not stat_module.S_ISDIR(snapshot.st_mode)
    ):
        raise SupervisorError("worker_generation_changed")
    return FileGeneration.from_stat(snapshot)


def _load_pptx_extractor() -> Any:
    sys.modules.setdefault("pptx_evidence", sys.modules[__name__])
    extractor_path = Path(__file__).with_name("pptx-extraction.py")
    spec = importlib.util.spec_from_file_location(
        "_speaker_toolkit_supervised_pptx_extraction",
        extractor_path,
    )
    if spec is None or spec.loader is None:
        raise PptxEvidenceError(
            "cannot load the current PPTX extractor",
            reason_code="pptx_dependency_unavailable",
        )
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _validated_extract_worker_payload(
    payload: Mapping[str, object],
) -> tuple[Path, Path | None, str | None, list[list[int]]]:
    """Validate the complete extraction request before any artifact I/O."""
    expected_fields = {
        "pptx_path",
        "trusted_root",
        "ocr",
        "rendered_pdf_path",
        "inspected_page_ranges",
        "extraction_schema_version",
        "extraction_pipeline_version",
    }
    if set(payload) != expected_fields:
        raise SupervisorError("invalid_worker_request")
    if (
        type(payload.get("ocr")) is not bool
        or payload.get("extraction_schema_version") != PPTX_EXTRACTION_SCHEMA_VERSION
        or payload.get("extraction_pipeline_version")
        != PPTX_EXTRACTION_PIPELINE_VERSION
    ):
        raise SupervisorError("invalid_worker_request")
    pptx_path = _worker_bound_path(payload.get("pptx_path"))
    trusted_root_value = payload.get("trusted_root")
    if trusted_root_value is not None and not isinstance(trusted_root_value, str):
        raise SupervisorError("invalid_worker_request")
    trusted_root = (
        _worker_bound_path(trusted_root_value)
        if trusted_root_value is not None
        else None
    )
    rendered_value = payload.get("rendered_pdf_path")
    if rendered_value is not None and not isinstance(rendered_value, str):
        raise SupervisorError("invalid_worker_request")
    ranges = payload.get("inspected_page_ranges")
    try:
        ranges = _validated_requested_page_ranges(ranges)
    except PptxEvidenceError as exc:
        raise SupervisorError("invalid_worker_request") from exc
    return pptx_path, trusted_root, rendered_value, ranges


def _extract_child(payload: Mapping[str, object]) -> dict[str, object]:
    pptx_path, _trusted_root, rendered_value, ranges = (
        _validated_extract_worker_payload(payload)
    )
    try:
        extractor = _load_pptx_extractor()
        extract = getattr(extractor, "_extract_pptx_in_process", None)
        if not callable(extract):
            raise PptxEvidenceError(
                "current PPTX extractor has no contained entrypoint",
                reason_code="pptx_dependency_unavailable",
            )
        result = extract(
            pptx_path,
            ocr=payload["ocr"],
            rendered_pdf_path=(
                Path(rendered_value) if rendered_value is not None else None
            ),
            inspected_page_ranges=ranges,
        )
    except MemoryError:
        return {
            "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": "pptx_probe_resource_unavailable",
            "details": {},
        }
    except PptxEvidenceError as exc:
        reason_code = (
            exc.reason_code
            if exc.reason_code in _CHILD_PROBE_REASON_CODES
            else "pptx_evidence_invalid"
        )
        return {
            "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": reason_code,
            "details": _probe_child_failure_details(exc),
        }
    except (OSError, ValueError, KeyError, AttributeError) as exc:
        return {
            "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
            "status": "unavailable",
            "reason_code": "pptx_evidence_invalid",
            "details": {"exception_type": type(exc).__name__},
        }
    if not isinstance(result, dict):
        raise SupervisorError("worker_operation_failed")
    return {
        "schema_version": PPTX_ARTIFACT_PROBE_SCHEMA_VERSION,
        "status": "available",
        "extraction": result,
    }


def _dispatch_supervised_worker(
    request: WorkerRequest,
) -> tuple[dict[str, object], dict[str, FileGeneration]]:
    payload = _worker_request_payload(request)
    expected_profile = {
        PPTX_METADATA_OPERATION: PPTX_METADATA_LIMITS.profile_id,
        PPTX_PROBE_OPERATION: PPTX_PROBE_LIMITS.profile_id,
        PPTX_NATIVE_AUDIT_OPERATION: PPTX_NATIVE_AUDIT_LIMITS.profile_id,
        PPTX_EXTRACT_OPERATION: (
            PPTX_EXTRACT_OCR_LIMITS.profile_id
            if payload.get("ocr") is True
            else PPTX_EXTRACT_NO_OCR_LIMITS.profile_id
        ),
    }.get(request.operation)
    if (
        expected_profile is None
        or request.limit_profile_id != expected_profile
        or request.schema_generation != PPTX_EXTRACTION_SCHEMA_VERSION
        or request.pipeline_generation != PPTX_EXTRACTION_PIPELINE_VERSION
    ):
        raise SupervisorError("invalid_worker_request")
    if request.operation == PPTX_METADATA_OPERATION:
        if request.expected_generations:
            raise SupervisorError("invalid_worker_request")
        return _metadata_child(payload), {}

    if request.operation in {PPTX_PROBE_OPERATION, PPTX_NATIVE_AUDIT_OPERATION}:
        pptx_path = _worker_bound_path(payload.get("pptx_path"))
        trusted_root_value = payload.get("trusted_root")
        if set(payload) != {"pptx_path", "trusted_root"}:
            raise SupervisorError("invalid_worker_request")
        if trusted_root_value is not None and not isinstance(trusted_root_value, str):
            raise SupervisorError("invalid_worker_request")
        trusted_root = (
            _worker_bound_path(trusted_root_value)
            if trusted_root_value is not None
            else None
        )
    elif request.operation == PPTX_EXTRACT_OPERATION:
        pptx_path, trusted_root, _rendered_value, _ranges = (
            _validated_extract_worker_payload(payload)
        )
    else:
        raise SupervisorError("invalid_worker_operation")
    paths: dict[str, tuple[Path, Path | None]] = {"pptx": (pptx_path, trusted_root)}
    if request.operation == PPTX_EXTRACT_OPERATION:
        rendered_value = payload.get("rendered_pdf_path")
        if rendered_value is not None:
            paths["rendered_pdf"] = (_worker_bound_path(rendered_value), None)
    expected_names = set(paths)
    if trusted_root is not None:
        expected_names.add("pptx_root")
    if expected_names != set(request.expected_generations):
        raise SupervisorError("invalid_worker_request")
    before = {
        name: _worker_generation(path, trusted_root=root)
        for name, (path, root) in paths.items()
    }
    if trusted_root is not None:
        before["pptx_root"] = _worker_root_generation(trusted_root)
    if before != dict(request.expected_generations):
        raise SupervisorError("worker_generation_changed")

    if request.operation == PPTX_PROBE_OPERATION:
        result = _pptx_probe_child(pptx_path)
    elif request.operation == PPTX_NATIVE_AUDIT_OPERATION:
        result = _native_audit_child(pptx_path)
    elif request.operation == PPTX_EXTRACT_OPERATION:
        result = _extract_child(payload)
    else:
        raise SupervisorError("invalid_worker_operation")
    after = {
        name: _worker_generation(path, trusted_root=root)
        for name, (path, root) in paths.items()
    }
    if trusted_root is not None:
        after["pptx_root"] = _worker_root_generation(trusted_root)
    if after != before:
        raise SupervisorError("worker_generation_changed")
    return result, after


def _run_supervised_worker_child() -> int:
    request = read_worker_request(max_input_bytes=64 * 1024)
    protocol_output = isolate_protocol_output()
    try:
        try:
            payload, observed = _dispatch_supervised_worker(request)
            write_worker_response(
                request,
                payload=payload,
                observed_generations=observed,
                stream=protocol_output,
                max_output_bytes=128 * 1024 * 1024,
            )
        except SupervisorError as exc:
            # A pre-operation generation mismatch has no trustworthy observed
            # generation. Echoing the expected binding lets the authenticated
            # error arrive; the parent independently re-snapshots before use.
            write_worker_response(
                request,
                error=SupervisorError(exc.reason_code, exc.details),
                observed_generations=request.expected_generations,
                stream=protocol_output,
                max_output_bytes=128 * 1024 * 1024,
            )
    finally:
        protocol_output.close()
    return 0


def _main() -> int:
    if sys.argv[1:] != [PPTX_SUPERVISED_WORKER_FLAG]:
        raise SystemExit("pptx_evidence.py is a library; run pptx-extraction.py")
    try:
        return _run_supervised_worker_child()
    except SupervisorError as exc:
        print(
            f"pptx supervised worker failed: {exc.reason_code}",
            file=sys.stderr,
            flush=True,
        )
        return 2
    # The supervisor treats a nonzero child without an authenticated response
    # as a bounded crash. Emit a path-neutral stderr diagnostic plus exit 2
    # because propagation would leak a traceback and violate the one-frame
    # response contract. outer-boundary-process-contract.
    except Exception:  # noqa: BLE001
        print(
            "pptx supervised worker failed: unexpected_error",
            file=sys.stderr,
            flush=True,
        )
        return 2


if __name__ == "__main__":
    raise SystemExit(_main())
