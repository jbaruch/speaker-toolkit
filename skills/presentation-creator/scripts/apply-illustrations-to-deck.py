#!/usr/bin/env python3
"""Apply generated illustrations to a deck with designed title placement.

For each slide whose outline block has a `Safe zone:` field:
  1. Replace the background picture with the matching illustration.
  2. Add a zone-sized semi-transparent scrim between the picture and
     the text (if not already present).
  3. Reposition title text boxes into the designed safe zone:
     upper_third/middle_third/lower_third -> full-width band at the
     matching Y; left_half/right_half -> narrower column on that side.

The outline is the single source of truth for zone assignments — the
same `Safe zone:` lines that `generate-illustrations.py` reads.
See `rules/title-overlay-rules.md` for the policy behind this.

Usage:
    apply-illustrations-to-deck.py DECK ILLUSTRATIONS_DIR OUTLINE_MD \\
        [--out OUT_DECK] [--image-ext jpg|jpeg|png] \\
        [--scrim-color RRGGBB] [--scrim-alpha 0-100000]
"""
import argparse
import re
import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

# 16:9 slide geometry (inches)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Layout per zone: top/left of the title column, and column width.
# Horizontal bands use the full content width centered horizontally.
# Half-frame zones use a narrower column on the chosen side, with the
# title vertically centered.
TEXT_W_FULL_IN = 10.0
TEXT_W_HALF_IN = 5.5
HALF_MARGIN_IN = 0.4
_BAND_H_IN = 1.9    # horizontal-band title height (title + subtitle)
_HALF_H_IN = 4.5    # half-frame title column height
_HALF_TOP_IN = (SLIDE_H_IN - _HALF_H_IN) / 2

ZONE_LAYOUT = {
    "upper_third":  {"top_in": 0.4,         "left_in": (SLIDE_W_IN - TEXT_W_FULL_IN) / 2,            "width_in": TEXT_W_FULL_IN, "height_in": _BAND_H_IN},
    "middle_third": {"top_in": (SLIDE_H_IN - _BAND_H_IN) / 2, "left_in": (SLIDE_W_IN - TEXT_W_FULL_IN) / 2, "width_in": TEXT_W_FULL_IN, "height_in": _BAND_H_IN},
    "lower_third":  {"top_in": SLIDE_H_IN - _BAND_H_IN - 0.4, "left_in": (SLIDE_W_IN - TEXT_W_FULL_IN) / 2, "width_in": TEXT_W_FULL_IN, "height_in": _BAND_H_IN},
    "left_half":    {"top_in": _HALF_TOP_IN, "left_in": HALF_MARGIN_IN,                               "width_in": TEXT_W_HALF_IN, "height_in": _HALF_H_IN},
    "right_half":   {"top_in": _HALF_TOP_IN, "left_in": SLIDE_W_IN - HALF_MARGIN_IN - TEXT_W_HALF_IN, "width_in": TEXT_W_HALF_IN, "height_in": _HALF_H_IN},
}

SUBTITLE_OFFSET_IN = 1.2

# Default scrim: 45% black. Decks with a strong tonal style (warm sepia,
# cool night, etc.) should pass a sampled color via --scrim-color.
# See suggest-scrim-color.py (same directory) and rules/title-overlay-rules.md §5.
DEFAULT_SCRIM_HEX = "000000"
DEFAULT_SCRIM_ALPHA = 45000

SCRIM_SHAPE_NAME = "_title_scrim"


def parse_zones(outline_path: Path) -> dict:
    """Read `Safe zone:` lines from the outline.

    Parses per-slide blocks (### Slide N: ... up to the next ### or ##)
    so a Safe zone line is never matched against the wrong slide.

    Returns {slide_num: zone_name} where zone_name is a key of ZONE_LAYOUT.
    """
    text = outline_path.read_text()
    zones = {}
    # Split into per-slide blocks bounded by ### headers
    slide_block_re = re.compile(
        r"###\s+Slide\s+(\d+):(.*?)(?=\n###\s|\n##\s|\Z)", re.DOTALL
    )
    zone_re = re.compile(
        r"-\s*Safe zone:\s*(upper_third|middle_third|lower_third|left_half|right_half)"
    )
    for block_match in slide_block_re.finditer(text):
        slide_num = int(block_match.group(1))
        block_text = block_match.group(2)
        zone_match = zone_re.search(block_text)
        if zone_match:
            zones[slide_num] = zone_match.group(1)
    return zones


def replace_picture_blob(picture_shape, new_image_path: Path) -> None:
    """Swap the picture's embedded image to the one at new_image_path.

    Uses rel re-pointing rather than mutating ``image_part._blob`` in
    place. A template may seed every slide's picture from a single
    placeholder file; python-pptx dedupes identical source paths into
    one image part shared across slides. Mutating that shared blob
    would clobber every other slide referencing it — last swap wins
    and every slide ends up with the same final image.
    """
    slide_part = picture_shape.part
    _, new_rId = slide_part.get_or_add_image_part(str(new_image_path))
    picture_shape._element.blipFill.blip.set(qn("r:embed"), new_rId)


def ensure_scrim(slide, zone: str, scrim_hex: str, scrim_alpha: int) -> int:
    """Add a zone-sized semi-transparent rectangle between picture and text.

    Zone-scoped (not full-slide) — a full-slide scrim flattens the whole
    illustration; scoping to the title box keeps the rest at full brightness.
    OOXML spPr child order must be xfrm -> prstGeom -> solidFill -> ln.
    """
    # Skip if a scrim was already added (identified by name)
    if any(s.name == SCRIM_SHAPE_NAME for s in slide.shapes):
        return 0
    layout = ZONE_LAYOUT[zone]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(layout["left_in"]), top=Inches(layout["top_in"]),
        width=Inches(layout["width_in"]), height=Inches(layout["height_in"]),
    )
    shape.name = SCRIM_SHAPE_NAME
    sp = shape._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)

    spPr = sp.find(qn("p:spPr"))
    keep = {qn("a:xfrm"), qn("a:prstGeom"), qn("a:custGeom")}
    for child in list(spPr):
        if child.tag not in keep:
            spPr.remove(child)
    solid = etree.SubElement(spPr, qn("a:solidFill"))
    clr = etree.SubElement(solid, qn("a:srgbClr"))
    clr.set("val", scrim_hex.upper())
    alpha = etree.SubElement(clr, qn("a:alpha"))
    alpha.set("val", str(scrim_alpha))
    ln = etree.SubElement(spPr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))

    # Insert scrim just before the first shape with a text frame
    spTree = slide.shapes._spTree
    spTree.remove(sp)
    first_text_idx = None
    last_pic_idx = None
    for i, child in enumerate(spTree):
        # Track last picture for fallback insertion point
        if child.tag == qn("p:pic"):
            last_pic_idx = i
        # Detect any shape with text: textboxes (txBox="1") and placeholders with txBody
        nvSpPr = child.find(qn("p:nvSpPr"))
        if nvSpPr is not None:
            cNvSpPr = nvSpPr.find(qn("p:cNvSpPr"))
            has_txbox = cNvSpPr is not None and cNvSpPr.get("txBox") == "1"
            has_placeholder = nvSpPr.find(qn("p:nvPr")) is not None and \
                nvSpPr.find(qn("p:nvPr")).find(qn("p:ph")) is not None
            has_txbody = child.find(qn("p:txBody")) is not None
            if has_txbox or (has_placeholder and has_txbody):
                first_text_idx = i
                break
    if first_text_idx is not None:
        spTree.insert(first_text_idx, sp)
    elif last_pic_idx is not None:
        spTree.insert(last_pic_idx + 1, sp)
    else:
        spTree.append(sp)
    return 1


def reposition_title(slide, zone: str) -> int:
    """Reposition title and subtitle shapes into the designed zone.

    Only moves title/subtitle placeholders and text boxes — leaves body
    text, callouts, and other content shapes in their original positions.
    """
    # Prefer the slide's explicit title/subtitle placeholders
    title_shapes = []
    if slide.shapes.title is not None:
        title_shapes.append(slide.shapes.title)
    for s in slide.placeholders:
        if s.placeholder_format.idx == 1 and s not in title_shapes:  # subtitle
            title_shapes.append(s)
    # Fall back to text boxes if no placeholders found
    if not title_shapes:
        title_shapes = [
            s for s in slide.shapes
            if s.has_text_frame and s.name != SCRIM_SHAPE_NAME
            and s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        ]
    title_shapes.sort(key=lambda s: s.top)
    if not title_shapes:
        return 0

    layout = ZONE_LAYOUT[zone]
    title_top = Inches(layout["top_in"])
    text_left = Inches(layout["left_in"])
    text_width = Inches(layout["width_in"])

    for j, shape in enumerate(title_shapes):
        shape.left = int(text_left)
        shape.width = int(text_width)
        shape.top = int(title_top + Inches(SUBTITLE_OFFSET_IN * j))
    return len(title_shapes)


def apply(
    deck: Path, illust_dir: Path, zones: dict, out_deck: Path, ext: str,
    scrim_hex: str, scrim_alpha: int,
) -> list[dict]:
    if out_deck.exists():
        out_deck.unlink()
    shutil.copy2(deck, out_deck)

    prs = Presentation(str(out_deck))
    results = []

    for n, zone in sorted(zones.items()):
        illust = illust_dir / f"slide-{n:02d}.{ext}"
        if not illust.exists():
            print(f"  [{n:02d}] SKIP: missing {illust.name}")
            continue
        if n > len(prs.slides):
            print(f"  [{n:02d}] SKIP: out of deck range")
            continue

        slide = prs.slides[n - 1]
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pictures:
            print(f"  [{n:02d}] SKIP: no picture shape")
            continue
        bg = max(pictures, key=lambda s: (s.width or 0) * (s.height or 0))

        try:
            replace_picture_blob(bg, illust)
        except Exception as e:
            print(f"  [{n:02d}] FAILED image swap: {e}")
            continue

        scrim_added = ensure_scrim(slide, zone, scrim_hex, scrim_alpha)
        moved = reposition_title(slide, zone)
        print(f"  [{n:02d}] zone={zone}  moved={moved} text  scrim+{scrim_added}")
        results.append({
            "slide": n, "zone": zone,
            "text_moved": moved, "scrim_added": scrim_added,
        })

    prs.save(str(out_deck))
    return results


def main():
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("deck", type=Path, help="Path to source .pptx")
    ap.add_argument("illustrations", type=Path, help="Directory with slide-NN.<ext> files")
    ap.add_argument("outline", type=Path, help="Path to presentation-outline.md")
    ap.add_argument("--out", type=Path, default=None, help="Output deck (default: <stem>-with-titles.pptx)")
    ap.add_argument("--image-ext", default="jpg", choices=["jpg", "jpeg", "png"])
    ap.add_argument("--scrim-color", default=DEFAULT_SCRIM_HEX,
                    help="Scrim color as 6-digit hex (default: %(default)s). "
                         "Run suggest-scrim-color.py to sample one from the deck's illustrations.")
    ap.add_argument("--scrim-alpha", type=int, default=DEFAULT_SCRIM_ALPHA,
                    help="Scrim opacity in OOXML thousandths (0-100000, default: %(default)s = 45%%).")
    args = ap.parse_args()

    out_deck = args.out or args.deck.with_name(args.deck.stem + "-with-titles.pptx")
    zones = parse_zones(args.outline)
    if not zones:
        print(f"No `Safe zone:` lines found in {args.outline.name}. Nothing to do.")
        return

    scrim_hex = args.scrim_color.lstrip("#").upper()
    if len(scrim_hex) != 6 or any(c not in "0123456789ABCDEF" for c in scrim_hex):
        raise SystemExit(f"--scrim-color must be a 6-digit hex, got {args.scrim_color!r}")
    if not (0 <= args.scrim_alpha <= 100000):
        raise SystemExit(f"--scrim-alpha must be 0..100000, got {args.scrim_alpha}")

    results = apply(
        args.deck, args.illustrations, zones, out_deck, args.image_ext,
        scrim_hex, args.scrim_alpha,
    )
    print(f"\nSaved {out_deck}")
    print(f"Updated {len(results)}/{len(zones)} slides")


if __name__ == "__main__":
    main()
