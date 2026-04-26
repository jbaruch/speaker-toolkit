#!/usr/bin/env python3
"""Insert bright-yellow placeholder slides into a PowerPoint deck at specified positions.

Placeholders are visually loud (yellow background, bold [PLACEHOLDER] title) so they
stand out in thumbnail view when editing large decks. Each placeholder includes a
subtitle with context about what the slide needs to become.

Usage:
    insert-placeholder-slides.py <deck.pptx> <positions.json>
    insert-placeholder-slides.py <deck.pptx> --at 5 --title "New Slide" --subtitle "Description"
    insert-placeholder-slides.py <deck.pptx> <positions.json> --output adapted-deck.pptx

The JSON file format (multiple slides):
    [
        {"position": 5,  "title": "The Great Siloing", "subtitle": "Yegge 20/60/20 adoption data"},
        {"position": 16, "title": "Cost Curves",       "subtitle": "MIT SSRN paper visualization"}
    ]

Positions are 1-indexed (final slide number after all insertions) and must be distinct.
Slides are inserted from lowest position to highest so every move target stays within
the current deck bounds (Python's list.insert silently clamps out-of-range indices).

The title is auto-prefixed with "[PLACEHOLDER] "; if the JSON title already starts
with that prefix (any case), the script does not double-prefix.

If the target template has no layout named "Blank", the script falls back to the last
layout and warns. Use --blank-layout-name to override when the last layout carries
decorative shapes that would appear behind the yellow placeholder.

Requires:
    - python-pptx  (pip install python-pptx)
"""

import argparse
import json
import sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

YELLOW = RGBColor(0xFF, 0xF2, 0x9E)
TITLE_COLOR = RGBColor(0x20, 0x20, 0x20)
SUBTITLE_COLOR = RGBColor(0x40, 0x40, 0x40)

PLACEHOLDER_PREFIX = "[PLACEHOLDER] "


def _format_title(title):
    """Prepend '[PLACEHOLDER] ' unless the caller already included it."""
    stripped = title.lstrip()
    if stripped.lower().startswith(PLACEHOLDER_PREFIX.lower()):
        return stripped
    return PLACEHOLDER_PREFIX + title


def _find_blank_layout(prs, preferred_name="Blank"):
    """Return (layout, used_fallback) — warns once in the caller if fallback was used."""
    for layout in prs.slide_layouts:
        if layout.name == preferred_name:
            return layout, False
    return prs.slide_layouts[-1], True


def add_placeholder_slide(prs, title, subtitle="", blank_layout_name="Blank"):
    """Append a yellow placeholder slide with title and optional subtitle."""
    blank_layout, used_fallback = _find_blank_layout(prs, blank_layout_name)
    if used_fallback:
        print(
            f"  WARNING: No '{blank_layout_name}' layout found — using last layout "
            f"'{blank_layout.name}'. Any decorative shapes on that layout will appear "
            f"behind the placeholder. Pass --blank-layout-name to choose a different layout.",
            file=sys.stderr,
        )

    slide = prs.slides.add_slide(blank_layout)
    sw, sh = prs.slide_width, prs.slide_height

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = YELLOW
    bg.line.fill.background()

    margin = Emu(457200)
    title_h = Emu(1600000)
    box = slide.shapes.add_textbox(margin, Emu(sh // 3), sw - 2 * margin, title_h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = _format_title(title)
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            margin, Emu(sh // 3) + title_h + Emu(200000),
            sw - 2 * margin, Emu(2500000)
        )
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(20)
        srun.font.color.rgb = SUBTITLE_COLOR

    return slide


def move_slide(prs, old_index, new_index):
    """Reorder slide by manipulating sldIdLst XML."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def main():
    parser = argparse.ArgumentParser(
        description="Insert yellow placeholder slides into a PowerPoint deck."
    )
    parser.add_argument("deck", help="Path to .pptx file")
    parser.add_argument("json_file", nargs="?", help="JSON file with placeholder definitions")
    parser.add_argument("--at", type=int, action="append", help="Position (1-indexed, repeatable)")
    parser.add_argument("--title", nargs="*", help="Title(s) for --at mode")
    parser.add_argument("--subtitle", nargs="*", help="Subtitle(s) for --at mode")
    parser.add_argument("--output", "-o", help="Output path (default: overwrite input)")
    parser.add_argument(
        "--blank-layout-name",
        default="Blank",
        help="Slide layout name to use for placeholders (default: 'Blank'). "
             "Override when the template's last layout carries decorative shapes.",
    )
    args = parser.parse_args()

    if args.json_file:
        with open(args.json_file) as f:
            placeholders = json.load(f)
    elif args.at:
        titles = args.title or ["New Slide"] * len(args.at)
        subtitles = args.subtitle or [""] * len(args.at)
        if args.title and len(args.title) != len(args.at):
            print(f"ERROR: {len(args.at)} positions but {len(args.title)} titles",
                  file=sys.stderr)
            sys.exit(1)
        if args.subtitle and len(args.subtitle) != len(args.at):
            print(f"ERROR: {len(args.at)} positions but {len(args.subtitle)} subtitles",
                  file=sys.stderr)
            sys.exit(1)
        placeholders = [
            {"position": pos, "title": t, "subtitle": s}
            for pos, t, s in zip(args.at, titles, subtitles)
        ]
    else:
        print("ERROR: provide a JSON file or use --at/--title", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(args.deck)
    original_count = len(prs.slides)
    total_after = original_count + len(placeholders)
    print(f"Opened {args.deck}: {original_count} slides")

    # Validate positions: in range and distinct
    seen = set()
    for ph in placeholders:
        pos = ph["position"]
        if pos < 1 or pos > total_after:
            print(f"ERROR: position {pos} out of range (1..{total_after})", file=sys.stderr)
            sys.exit(1)
        if pos in seen:
            print(
                f"ERROR: duplicate position {pos} — 'final position after all insertions' "
                f"must be unique per placeholder",
                file=sys.stderr,
            )
            sys.exit(1)
        seen.add(pos)

    # Insert from lowest position to highest. Each iteration appends one placeholder,
    # growing the deck by one, then moves it to position-1. With distinct positions
    # sorted ascending, position[i] <= original_count + i + 1 is guaranteed, so the
    # move target is always within bounds (avoiding Python's list.insert clamp).
    for ph in sorted(placeholders, key=lambda x: x["position"]):
        title = ph["title"]
        subtitle = ph.get("subtitle", "")
        position = ph["position"]

        add_placeholder_slide(prs, title, subtitle, blank_layout_name=args.blank_layout_name)
        last_idx = len(prs.slides) - 1
        move_slide(prs, last_idx, position - 1)
        print(f"  Inserted {_format_title(title)} at position {position}")

    output_path = args.output or args.deck
    prs.save(output_path)
    print(f"Saved {output_path}: {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
