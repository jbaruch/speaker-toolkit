#!/usr/bin/env python3
"""Generate a QR code and optionally insert it into a PowerPoint deck.

Generates an unbranded QR code encoding a shownotes URL (optionally shortened
via bit.ly or rebrand.ly), matches the QR background color to the target slide,
auto-selects foreground color for contrast, and inserts the QR image into the
deck at the configured slide position.

Two URL-resolution modes:
  --shownotes-url URL   Script resolves shortening itself via {vault}/secrets.json
  --short-url URL       Agent pre-resolved the short URL (via MCP); requires
                        --shownotes-url as the canonical redirect target. Script just
                        generates the QR and inserts it

PNG-only mode (no deck required):
  --png-only            Generate the QR PNG without opening or modifying a deck.
                        Use --bg-color R,G,B and --output PATH to control colors
                        and output location. URL shortening and tracking DB
                        updates still run normally.

Usage:
    python3 generate-qr.py <deck.pptx> --talk-slug SLUG --shownotes-url URL
    python3 generate-qr.py <deck.pptx> --talk-slug SLUG --shownotes-url URL \\
        --short-url SHORT_URL --short-provider bitly --short-link-id LINK_ID
    python3 generate-qr.py <deck.pptx> --talk-slug SLUG --shownotes-url URL --dry-run
    python3 generate-qr.py <deck.pptx> --talk-slug SLUG --shownotes-url URL --profile PATH --vault PATH
    python3 generate-qr.py --png-only --talk-slug SLUG --shownotes-url URL --output qr.png --bg-color 128,0,128

Requires:
    - python-pptx  (pip install python-pptx)
    - qrcode       (pip install qrcode)
    - Pillow       (pip install Pillow — transitive dep of python-pptx)
"""

import argparse
import contextlib
import copy
import datetime
import fcntl
import hashlib
import io
import json
import os
from contextlib import contextmanager
from pathlib import Path
import shutil
import subprocess
import sys
import urllib.error
import urllib.parse
import urllib.request

try:
    import qrcode
    from qrcode.constants import ERROR_CORRECT_M
except ImportError:
    print("ERROR: 'qrcode' package not installed. Run: pip install qrcode")
    sys.exit(1)

try:
    from PIL import Image, ImageChops, ImageStat
except ImportError:
    print("ERROR: 'Pillow' package not installed. Run: pip install Pillow")
    sys.exit(1)

from pptx import Presentation  # read-only: background-color match + slide-finding
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

VAULT_INGRESS_SCRIPTS = (
    Path(__file__).resolve().parents[2] / "vault-ingress" / "scripts"
)
if str(VAULT_INGRESS_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(VAULT_INGRESS_SCRIPTS))

# Pyright cannot resolve this sibling script module added to sys.path at runtime.
from tracking_database import (  # noqa: E402  # pyright: ignore[reportMissingImports]
    QR_CODE_RECORD_SCHEMA_VERSION,
    TrackingDatabaseError,
    assess_tracking_database,
    require_current_tracking_database,
)
# Pyright cannot resolve this sibling script module added to sys.path at runtime.
from tracking_database_io import (  # noqa: E402  # pyright: ignore[reportMissingImports]
    TrackingDatabaseIOError,
    TrackingDatabaseSnapshot,
    decode_json_object,
    snapshot_tracking_database,
    write_json_object,
)

# --- Constants ---

QR_BOX_SIZE = 10       # pixels per QR module
QR_BORDER = 4          # quiet-zone modules
QR_ERROR_CORRECTION = ERROR_CORRECT_M

# QR placement: bottom-right, 2 inches wide, 0.3 inch margin from edges
QR_WIDTH_INCHES = 2.0
QR_MARGIN_INCHES = 0.3

# Existing-QR detection (content-based, size-independent). A QR is a SQUARE
# picture that is BOTH essentially two colors AND roughly balanced between them:
#   - it reconstructs near-perfectly when quantized to 2 colors (reconErr ≈ 0),
#     unlike photos/diagrams (many colors → high error); and
#   - its minority color covers a large fraction (~⅓+), unlike a text screenshot
#     that is mostly one background color with sparse text.
# Both axes are needed: a mostly-white doc screenshot is ~2-color but unbalanced.
# Size is NOT a signal — a deck adapted from another talk carries inherited QRs at
# arbitrary sizes (the same QR appeared at 1.8" and 2.8"); only a floor guards
# against tiny 2-color icons. Validated on a real inherited deck: QRs scored
# reconErr 0.0 / minority 0.34; a colored Venn 68 / 0.50; a martinfowler.com
# screenshot 6.4 / 0.18. The VBA writer can't run PIL, so detection lives here and
# the resulting geometry is handed to InsertQR.
QR_DETECT_MIN_INCHES = 1.5      # side floor — excludes small 2-color icons
QR_SQUARE_TOL_INCHES = 0.1      # |width − height| tolerance
QR_RECON_ERR_MAX = 5.0          # max mean 2-color reconstruction error
QR_MIN_MINORITY = 0.25          # min minority-color fraction (QRs are balanced)


# --- Vault / Config Loading ---

def load_vault_config(vault_path, profile_path=None):
    """Load speaker profile, secrets, and tracking database from the vault.

    Returns:
        tuple (speaker_profile, secrets, tracking_db, tracking_db_snapshot)
        Any of these may be empty dicts if the file doesn't exist.
    """
    speaker_profile = {}
    secrets = {}
    tracking_db = {}
    tracking_db_snapshot = None

    # Speaker profile
    sp_path = profile_path or os.path.join(vault_path, "speaker-profile.json")
    if os.path.isfile(sp_path):
        with open(sp_path, "r", encoding="utf-8") as f:
            speaker_profile = json.load(f)

    # Secrets
    secrets_path = os.path.join(vault_path, "secrets.json")
    if os.path.isfile(secrets_path):
        with open(secrets_path, "r", encoding="utf-8") as f:
            secrets = json.load(f)

    # Tracking database
    tdb_path = os.path.join(vault_path, "tracking-database.json")
    if os.path.isfile(tdb_path):
        try:
            tracking_db_snapshot = snapshot_tracking_database(tdb_path)
            tracking_db = decode_json_object(tracking_db_snapshot)
        except TrackingDatabaseIOError as exc:
            raise SystemExit(f"ERROR: cannot load tracking database: {exc}") from exc

    return speaker_profile, secrets, tracking_db, tracking_db_snapshot


def _require_tracking_db_snapshot(
    tracking_db_snapshot: object,
) -> TrackingDatabaseSnapshot:
    """Return the loaded snapshot or reject a vault that cannot be persisted."""
    if not isinstance(tracking_db_snapshot, TrackingDatabaseSnapshot):
        raise ValueError(
            "tracking-database.json is missing; initialize it through "
            "vault-ingress mutate-tracking-database.py before generating a QR"
        )
    return tracking_db_snapshot


def write_tracking_db(tracking_db_snapshot, tracking_db):
    """Commit QR metadata against the generation loaded before QR work."""
    snapshot = _require_tracking_db_snapshot(tracking_db_snapshot)
    try:
        require_current_tracking_database(tracking_db)
        return write_json_object(snapshot, tracking_db)
    except (TrackingDatabaseError, TrackingDatabaseIOError) as exc:
        raise ValueError(f"cannot safely update tracking database: {exc}") from exc


@contextmanager
def qr_publication_lock(vault_path, talk_slug):
    """Serialize QR publication for one talk slug across its external effects.

    The tracking-database lock is held only for the duration of a write. QR
    publication creates a remote link, writes PNGs, and mutates a deck before it
    commits, so two concurrent runs for the same slug can interleave those
    effects. This lock spans all of them.

    It is keyed per slug rather than per database so unrelated talks publish
    concurrently, and it is never held across an unrelated writer's commit.
    """
    lock_path = os.path.join(vault_path, f".qr-{talk_slug}.lock")
    flags = os.O_RDWR | os.O_CREAT | getattr(os, "O_CLOEXEC", 0)
    flags |= getattr(os, "O_NOFOLLOW", 0)
    try:
        descriptor = os.open(lock_path, flags, 0o600)
    except OSError as exc:
        raise ValueError(
            f"cannot open QR publication lock {lock_path}: {exc}"
        ) from exc
    try:
        try:
            fcntl.flock(descriptor, fcntl.LOCK_EX)
        except OSError as exc:
            raise ValueError(
                f"cannot acquire QR publication lock {lock_path}: {exc}"
            ) from exc
        yield lock_path
    finally:
        try:
            fcntl.flock(descriptor, fcntl.LOCK_UN)
        except OSError as exc:
            print(f"WARNING: could not unlock {lock_path}: {exc}", file=sys.stderr)
        try:
            os.close(descriptor)
        except OSError as exc:
            print(f"WARNING: could not close {lock_path}: {exc}", file=sys.stderr)


class EffectsReceipt:
    """What this publication changed outside the tracking database.

    A commit that rejects after these effects landed must not look
    side-effect-free. The receipt names each effect so the operator knows
    whether to finalize, retry idempotently, or roll back.
    """

    def __init__(self, talk_slug):
        self.talk_slug = talk_slug
        self.short_link = None      # dict, or None when no link work happened
        self.artifacts = []         # written PNG paths
        self.deck = None            # deck path, or None if untouched

    def record_short_link(self, provider, link_id, short_url, *, action,
                          prior_target=None):
        """Record link work. ``action`` is created, retargeted, or preresolved.

        Rollback differs by action: a link this run created can be deleted, a
        retargeted link must be restored to ``prior_target`` instead, and a
        preresolved link was never ours to remove.
        """
        self.short_link = {
            "provider": provider,
            "link_id": link_id,
            "short_url": short_url,
            "action": action,
            "prior_target": prior_target,
        }

    def record_artifacts(self, paths):
        self.artifacts = list(paths)

    def record_deck(self, deck_path):
        self.deck = deck_path

    def any_effects(self):
        return bool(self.short_link or self.artifacts or self.deck)


def _report_unfinalized_effects(receipt):
    """Print what already landed, and how a retry behaves against it."""
    if receipt is None or not receipt.any_effects():
        print(
            "  No external effects were made; this run is safe to retry as-is.",
            file=sys.stderr,
        )
        return

    print("", file=sys.stderr)
    print(
        "  The tracking database was NOT updated, but these effects already "
        "landed:",
        file=sys.stderr,
    )
    link = receipt.short_link
    if link:
        print(
            f"    short link {link['action']}: {link['short_url']} "
            f"(provider={link['provider'] or 'unknown'}, "
            f"link_id={link['link_id'] or 'unknown'})",
            file=sys.stderr,
        )
        if link["action"] == "retargeted":
            print(
                f"      previous target: {link['prior_target']}",
                file=sys.stderr,
            )
    for path in receipt.artifacts:
        print(f"    PNG written: {path}", file=sys.stderr)
    if receipt.deck:
        print(f"    deck mutated: {receipt.deck}", file=sys.stderr)
    print("", file=sys.stderr)
    print(
        "  Re-running the same command is idempotent for the short link: an "
        "existing managed link with the slug back-half is retargeted, never "
        "duplicated. The PNGs and deck are overwritten in place.",
        file=sys.stderr,
    )
    if link and link["action"] == "created":
        rollback = "delete the short link above"
    elif link and link["action"] == "retargeted":
        rollback = (
            f"point {link['short_url']} back at {link['prior_target']} — do NOT "
            "delete it, it predates this run"
        )
    elif link:
        rollback = (
            "leave the short link alone; it was supplied pre-resolved and this "
            "run did not create it"
        )
    else:
        rollback = "remove the PNGs above"
    print(f"  Re-run to finalize, or {rollback} to roll back.", file=sys.stderr)


def _reload_tracking_db(vault_path):
    """Re-read the database under the slug lock, before short-link resolution."""
    tdb_path = os.path.join(vault_path, "tracking-database.json")
    try:
        snapshot = snapshot_tracking_database(tdb_path)
        fresh = decode_json_object(snapshot)
        require_current_tracking_database(fresh)
    except (TrackingDatabaseError, TrackingDatabaseIOError) as exc:
        raise ValueError(
            f"cannot re-read tracking database under the publication lock: {exc}"
        ) from exc
    return fresh


def _qr_record_for(database, talk_slug):
    """Return this talk's qr_codes record, or None."""
    for record in database.get("qr_codes", []):
        if isinstance(record, dict) and record.get("talk_slug") == talk_slug:
            return record
    return None


def commit_qr_record(tdb_path, meta, artifacts, prior_record):
    """Rebase this run's single qr_codes upsert onto the current generation.

    Committing the snapshot loaded before QR work would reject whenever ANY
    writer touched the database meanwhile — after the short link, PNGs, and
    deck had already changed. Re-reading here narrows that to a conflict on
    this talk's own record.

    ``prior_record`` is this talk's record as it stood when publication began,
    under the slug lock. A same-talk change landing since then is another
    owner's decision about the same record; rebasing over it would silently
    discard that decision, so it rejects instead. The slug lock keeps competing
    QR runs out, but non-QR writers do not take it.
    """
    try:
        fresh_snapshot = snapshot_tracking_database(tdb_path)
        fresh = decode_json_object(fresh_snapshot)
        require_current_tracking_database(fresh)
    except (TrackingDatabaseError, TrackingDatabaseIOError) as exc:
        raise ValueError(
            f"cannot re-read tracking database before commit: {exc}"
        ) from exc

    current_record = _qr_record_for(fresh, meta["talk_slug"])
    if current_record != prior_record:
        raise ValueError(
            f"the qr_codes record for {meta['talk_slug']!r} changed during "
            "publication; another writer owns that change. Re-run to rebuild "
            "against it rather than discarding it"
        )

    update_tracking_db(fresh, meta, artifacts)
    return write_tracking_db(fresh_snapshot, fresh)


# --- URL Shortening ---

def _http_request(url, data=None, headers=None, method="GET"):
    """Make an HTTP request using stdlib urllib. Returns parsed JSON or raises."""
    if headers is None:
        headers = {}
    if data is not None:
        data = json.dumps(data).encode("utf-8")
        if "Content-Type" not in headers:
            headers["Content-Type"] = "application/json"

    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    with urllib.request.urlopen(req, timeout=30) as resp:
        return json.loads(resp.read().decode("utf-8"))


def create_bitly_link(long_url, api_token, custom_back_half=None, domain=None):
    """Create a new bit.ly short link.

    Args:
        long_url: The URL to shorten
        api_token: Bitly API token
        custom_back_half: Custom back-half for the short URL (e.g., talk slug).
            If provided, creates {domain}/{custom_back_half} instead of a random hash.
        domain: Custom Bitly domain (e.g., "jbaru.ch"). Defaults to "bit.ly".

    Returns:
        dict with keys: short_url, link_id, short_path
    """
    bitly_domain = domain or "bit.ly"
    payload = {"long_url": long_url, "domain": bitly_domain}
    if custom_back_half:
        payload["title"] = custom_back_half  # for tracking
    headers = {
        "Authorization": f"Bearer {api_token}",
        "Content-Type": "application/json",
    }

    # Create the link first
    result = _http_request(
        "https://api-ssl.bitly.com/v4/bitlinks",
        data=payload,
        headers=headers,
        method="POST",
    )
    link_id = result["id"]
    short_url = result["link"]
    short_path = link_id.split("/", 1)[-1] if "/" in link_id else link_id

    # Set custom back-half if requested
    if custom_back_half:
        try:
            _http_request(
                "https://api-ssl.bitly.com/v4/custom_bitlinks",
                data={
                    "bitlink_id": link_id,
                    "custom_bitlink": f"{bitly_domain}/{custom_back_half}",
                },
                headers=headers,
                method="POST",
            )
            short_url = f"https://{bitly_domain}/{custom_back_half}"
            short_path = custom_back_half
            print(f"  Custom back-half set: {bitly_domain}/{custom_back_half}")
        except (
            urllib.error.HTTPError,
            urllib.error.URLError,
            OSError,
            json.JSONDecodeError,
            KeyError,
        ) as e:
            # A random hash would silently break the slug=back-half contract
            # (rules/qr-generation-rules.md §2), so surface the failure instead.
            # The link already exists provider-side; carry its identity so the
            # operator can reuse or delete it deterministically.
            raise ShortenerResolutionError(
                f"could not set custom back-half '{custom_back_half}' on "
                f"{bitly_domain}: {e}. A provider-side link was already created "
                f"and must be reused or deleted: link_id={link_id} "
                f"short_url={short_url}"
            ) from e

    return {
        "short_url": short_url,
        "link_id": link_id,
        "short_path": short_path,
    }


def update_bitly_link(link_id, new_long_url, api_token):
    """Update an existing bit.ly link's target URL via PATCH."""
    headers = {
        "Authorization": f"Bearer {api_token}",
        "Content-Type": "application/json",
    }
    _http_request(
        f"https://api-ssl.bitly.com/v4/bitlinks/{link_id}",
        data={"long_url": new_long_url},
        headers=headers,
        method="PATCH",
    )


def create_rebrandly_link(long_url, api_key, domain=None, slashtag=None):
    """Create a new rebrand.ly short link.

    Returns:
        dict with keys: short_url, link_id, short_path
    """
    payload = {"destination": long_url}
    if domain:
        payload["domain"] = {"fullName": domain}
    if slashtag:
        payload["slashtag"] = slashtag

    headers = {
        "apikey": api_key,
        "Content-Type": "application/json",
    }
    result = _http_request(
        "https://api.rebrandly.com/v1/links",
        data=payload,
        headers=headers,
        method="POST",
    )
    return {
        "short_url": f"https://{result['shortUrl']}",
        "link_id": result["id"],
        "short_path": result["slashtag"],
    }


def update_rebrandly_link(link_id, new_long_url, api_key):
    """Update an existing rebrand.ly link's target URL."""
    headers = {
        "apikey": api_key,
        "Content-Type": "application/json",
    }
    _http_request(
        f"https://api.rebrandly.com/v1/links/{link_id}",
        data={"destination": new_long_url},
        headers=headers,
        method="POST",
    )


_SUPPORTED_SHORTENERS = frozenset({"bitly", "rebrandly", "none"})


class ShortenerResolutionError(RuntimeError):
    """A configured shortener could not produce its managed short link.

    Only an explicit ``shortener: none`` authorizes encoding a raw target URL.
    Every other resolution failure raises this so the run stops before a PNG,
    deck, or tracking-database write, rather than silently shipping a QR
    without the managed redirect layer and cataloging it as ``none``.
    """


def _print_missing_key_help(service, key_name, vault_path):
    """Print actionable help when an API key is missing from secrets.json."""
    secrets_path = os.path.join(vault_path, "secrets.json") if vault_path else "secrets.json"
    if vault_path and not os.path.isfile(secrets_path):
        print(f"  WARNING: No {service}.{key_name} found — secrets.json does not exist. Falling back to raw URL.")
        print("  Create it:")
        print(f'    echo \'{{\"{service}\": {{\"{key_name}\": \"YOUR_KEY\"}}}}\' > {secrets_path}')
        print(f"    chmod 600 {secrets_path}")
    else:
        print(f"  WARNING: No {service}.{key_name} in secrets.json, falling back to raw URL.")
        print(f"  Add to {secrets_path}:")
        print(f'    \"{service}\": {{\"{key_name}\": \"YOUR_KEY\"}}')


def _require_domain_decision(config, shortener, vault_path):
    """A custom-domain decision must be recorded before the FIRST short link is
    created. An ABSENT `{shortener}_domain` key means the user was never asked —
    STOP so the agent asks and saves the answer. A present value (a domain string,
    or null for "no custom domain") is a recorded decision and proceeds.
    """
    key = f"{shortener}_domain"
    if key in config:
        return
    default_domain = "bit.ly" if shortener == "bitly" else "the shortener default"
    profile = os.path.join(vault_path, "speaker-profile.json") if vault_path else "the speaker profile"
    print(f"ERROR: No custom-domain decision recorded for shortener '{shortener}'.")
    print("  Before creating the first short link, ask the user whether they have a")
    print("  custom domain (e.g. jbaru.ch), then save the answer under")
    print(f"  publishing_process.qr_code.{key} in {profile}:")
    print(f'    a domain string (e.g. "jbaru.ch"), or null for no custom domain ({default_domain}).')
    sys.exit(1)


def resolve_short_url(shownotes_url, talk_slug, config, secrets, tracking_db,
                      dry_run=False, vault_path=None, effects_receipt=None):
    """Resolve the short URL for a talk, using cache or API as needed.

    Args:
        shownotes_url: The full shownotes URL to shorten
        talk_slug: Unique talk identifier
        config: QR config from speaker profile (publishing_process.qr_code)
        secrets: Parsed secrets.json
        tracking_db: Parsed tracking-database.json
        dry_run: If True, skip API calls
        vault_path: Path to vault (for actionable error messages)

    Returns:
        tuple (short_url, metadata_dict)
        metadata_dict has keys: shortener, short_path, link_id, short_url, target_url
    """
    shortener = config.get("shortener")
    qr_codes = tracking_db.get("qr_codes", [])

    # Look up existing entry
    existing = None
    for entry in qr_codes:
        if entry.get("talk_slug") == talk_slug:
            existing = entry
            break

    # Enforce slug-only back-half on EXISTING links too: a tracked shortened link
    # whose back-half isn't the slug is legacy — don't reuse it from cache or
    # retarget it in place; drop it so a slug-based link is recreated below.
    if existing and existing.get("shortener_link_id") and existing.get("short_path") != talk_slug:
        print(f"  Legacy non-slug back-half '{existing.get('short_path')}' for '{talk_slug}' — recreating with the slug")
        existing = None

    # Configuration is validated BEFORE any cache reuse. A cached record proves
    # what was authorized on some earlier run, never what is authorized now, so
    # reusing one ahead of this check would let a stale `shortener: none` entry
    # re-authorize a raw URL under a missing or managed configuration.
    if shortener is None:
        raise ShortenerResolutionError(
            "no URL shortener configured at publishing_process.qr_code.shortener. "
            "Set 'shortener: bitly' or 'shortener: rebrandly', or set "
            "'shortener: none' to explicitly authorize an unmanaged raw URL."
        )

    if shortener not in _SUPPORTED_SHORTENERS:
        raise ShortenerResolutionError(
            f"unknown shortener '{shortener}' at "
            "publishing_process.qr_code.shortener; supported values are "
            f"{', '.join(repr(s) for s in sorted(_SUPPORTED_SHORTENERS))}"
        )

    # A cached record is reusable only when it was produced by the shortener now
    # configured. A mismatch means the configuration changed since it was
    # written, so it must be re-resolved rather than replayed.
    if existing and existing.get("shortener") != shortener:
        print(
            f"  Cached record for '{talk_slug}' was produced by "
            f"'{existing.get('shortener')}' but '{shortener}' is configured now "
            "— re-resolving"
        )
        existing = None

    # If cached under the current shortener and the target matches, reuse
    if existing and existing.get("target_url") == shownotes_url:
        print(f"  Reusing cached short URL for '{talk_slug}': {existing['short_url']}")
        return existing["short_url"], existing

    if shortener == "none":
        meta = {
            "talk_slug": talk_slug,
            "target_url": shownotes_url,
            "shortener": "none",
            "short_path": None,
            "short_url": shownotes_url,
            "shortener_link_id": None,
        }
        return shownotes_url, meta

    if dry_run:
        print(f"  DRY RUN: would call {shortener} API for {shownotes_url}")
        return shownotes_url, {
            "talk_slug": talk_slug,
            "target_url": shownotes_url,
            "shortener": shortener,
            "short_path": None,
            "short_url": shownotes_url,
            "shortener_link_id": None,
        }

    # The custom back-half (bit.ly) / slashtag (rebrand.ly) is ALWAYS the talk
    # slug — no override (see rules/qr-generation-rules.md). It is the decoupling
    # layer and the human-readable, traceable short path.
    custom_back_half = talk_slug

    try:
        if shortener == "bitly":
            api_token = secrets.get("bitly", {}).get("api_token")
            if not api_token:
                _print_missing_key_help("bitly", "api_token", vault_path)
                raise ShortenerResolutionError(
                    "bitly is configured but its api_token is missing from "
                    "secrets.json; see the guidance above"
                )

            bitly_domain = config.get("bitly_domain")  # e.g., "jbaru.ch"

            if existing and existing.get("shortener_link_id"):
                # Update existing link target
                print(f"  Updating bit.ly link {existing['shortener_link_id']} → {shownotes_url}")
                update_bitly_link(existing["shortener_link_id"], shownotes_url, api_token)
                meta = dict(existing)
                prior_target = existing.get("target_url")
                meta["target_url"] = shownotes_url
                meta["updated_at"] = datetime.date.today().isoformat()
                if effects_receipt is not None:
                    effects_receipt.record_short_link(
                        meta["shortener"], meta["shortener_link_id"],
                        meta["short_url"], action="retargeted",
                        prior_target=prior_target,
                    )
                return existing["short_url"], meta
            else:
                _require_domain_decision(config, "bitly", vault_path)
                # Create new link with talk slug as custom back-half
                domain_label = bitly_domain or "bit.ly"
                print(f"  Creating {domain_label} link for {shownotes_url} (back-half: {custom_back_half})")
                result = create_bitly_link(shownotes_url, api_token, custom_back_half, domain=bitly_domain)
                meta = {
                    "talk_slug": talk_slug,
                    "target_url": shownotes_url,
                    "shortener": "bitly",
                    "short_path": result["short_path"],
                    "short_url": result["short_url"],
                    "shortener_link_id": result["link_id"],
                }
                if effects_receipt is not None:
                    effects_receipt.record_short_link(
                        meta["shortener"], meta["shortener_link_id"],
                        meta["short_url"], action="created",
                    )
                return result["short_url"], meta

        elif shortener == "rebrandly":
            api_key = secrets.get("rebrandly", {}).get("api_key")
            if not api_key:
                _print_missing_key_help("rebrandly", "api_key", vault_path)
                raise ShortenerResolutionError(
                    "rebrandly is configured but its api_key is missing from "
                    "secrets.json; see the guidance above"
                )

            domain = config.get("rebrandly_domain")

            if existing and existing.get("shortener_link_id"):
                print(f"  Updating rebrand.ly link {existing['shortener_link_id']} → {shownotes_url}")
                update_rebrandly_link(existing["shortener_link_id"], shownotes_url, api_key)
                meta = dict(existing)
                prior_target = existing.get("target_url")
                meta["target_url"] = shownotes_url
                meta["updated_at"] = datetime.date.today().isoformat()
                if effects_receipt is not None:
                    effects_receipt.record_short_link(
                        meta["shortener"], meta["shortener_link_id"],
                        meta["short_url"], action="retargeted",
                        prior_target=prior_target,
                    )
                return existing["short_url"], meta
            else:
                _require_domain_decision(config, "rebrandly", vault_path)
                print(f"  Creating rebrand.ly link for {shownotes_url} (slashtag: {custom_back_half})")
                result = create_rebrandly_link(shownotes_url, api_key, domain, custom_back_half)
                meta = {
                    "talk_slug": talk_slug,
                    "target_url": shownotes_url,
                    "shortener": "rebrandly",
                    "short_path": result["short_path"],
                    "short_url": result["short_url"],
                    "shortener_link_id": result["link_id"],
                }
                if effects_receipt is not None:
                    effects_receipt.record_short_link(
                        meta["shortener"], meta["shortener_link_id"],
                        meta["short_url"], action="created",
                    )
                return result["short_url"], meta

        else:  # pragma: no cover - _SUPPORTED_SHORTENERS is checked above
            raise ShortenerResolutionError(
                f"unknown shortener '{shortener}' at "
                "publishing_process.qr_code.shortener"
            )

    except ShortenerResolutionError:
        raise
    except (
        urllib.error.HTTPError,
        urllib.error.URLError,
        OSError,
        json.JSONDecodeError,
        # A provider response missing an expected field is a malformed-response
        # failure, not a bug in this script.
        KeyError,
    ) as e:
        raise ShortenerResolutionError(
            f"{shortener} could not produce the managed short link for "
            f"'{talk_slug}': {e}"
        ) from e


# --- Slide Background Color Resolution ---

def resolve_slide_bg_rgb(slide):
    """Walk slide → layout → master to find an explicit solid fill background.

    Returns:
        tuple (R, G, B) as ints 0-255, or None if no solid fill found.
    """
    for obj in [slide, slide.slide_layout, slide.slide_layout.slide_master]:
        bg = obj.background
        fill = bg.fill
        if fill.type is not None:
            # Check for solid fill
            try:
                color = fill.fore_color
                if color.type is not None:
                    rgb = color.rgb
                    return (rgb[0], rgb[1], rgb[2])
            except (AttributeError, TypeError):
                continue
    return None


def choose_fg_color(bg_rgb):
    """Auto-contrast: pick black or white foreground based on background luminance.

    Uses WCAG relative luminance formula:
        L = 0.2126*R/255 + 0.7152*G/255 + 0.0722*B/255

    Returns:
        tuple (R, G, B) — white (255,255,255) for dark backgrounds,
        black (0,0,0) for light backgrounds.
    """
    if bg_rgb is None:
        return (0, 0, 0)  # default to black on unknown background

    r, g, b = bg_rgb
    luminance = 0.2126 * r / 255 + 0.7152 * g / 255 + 0.0722 * b / 255

    if luminance < 0.5:
        return (255, 255, 255)  # white foreground on dark background
    else:
        return (0, 0, 0)  # black foreground on light background


# --- QR Code Generation ---

def generate_qr_png(url, fg_rgb, bg_rgb, out_path):
    """Generate a QR code PNG file.

    Args:
        url: The URL to encode
        fg_rgb: Foreground color as (R, G, B) tuple
        bg_rgb: Background color as (R, G, B) tuple, or None for white
        out_path: Path to save the PNG file
    """
    if bg_rgb is None:
        bg_rgb = (255, 255, 255)

    qr = qrcode.QRCode(
        error_correction=QR_ERROR_CORRECTION,
        box_size=QR_BOX_SIZE,
        border=QR_BORDER,
    )
    qr.add_data(url)
    qr.make(fit=True)

    img = qr.make_image(fill_color=fg_rgb, back_color=bg_rgb)
    img.save(out_path)
    return out_path


# --- Slide Detection ---

def find_shownotes_slide(prs, shownotes_url):
    """Find the slide index containing the shownotes URL text.

    Returns:
        int (0-based slide index) or None if not found.
    """
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = "".join(run.text for run in paragraph.runs)
                    if shownotes_url in full_text:
                        return idx
    return None


def _two_color_metrics(blob):
    """(reconstruction error, minority-color fraction) for the image quantized to
    two colors.

    A QR is BOTH ~2-color (low reconstruction error) AND ~balanced (large minority
    fraction). Photos/diagrams have many colors → high error; text screenshots are
    ~2-color but mostly one background → tiny minority fraction. Color-agnostic, so
    black-on-white and white-on-purple QRs both qualify. Returns (255, 0) if the
    image can't be read.
    """
    try:
        im = Image.open(io.BytesIO(blob)).convert("RGB")
    except Exception:
        return 255.0, 0.0
    q = im.quantize(colors=2)
    means = ImageStat.Stat(ImageChops.difference(im, q.convert("RGB"))).mean
    err = sum(means) / len(means) if means else 255.0
    counts = sorted((c for c in q.histogram() if c > 0), reverse=True)
    total = sum(counts)
    minority = counts[1] / total if len(counts) > 1 and total else 0.0
    return err, minority


def _picture_is_qr(shape):
    """A picture is an (inherited) QR when it is square, at least
    QR_DETECT_MIN_INCHES on a side, ~2-color, and roughly balanced between those
    colors.

    Size-independent by design: a deck adapted from another talk carries that
    talk's QR at whatever size it used (the same QR can appear at 1.8" and 2.8"),
    so a size band can't identify it. The two-color + balance test does, while the
    floor rejects tiny icons and the square test rejects banners/photos.
    """
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    w, h = shape.width, shape.height
    if w is None or h is None:
        return False
    if abs(w - h) >= Inches(QR_SQUARE_TOL_INCHES) or min(w, h) < Inches(QR_DETECT_MIN_INCHES):
        return False
    try:
        blob = shape.image.blob
    except Exception:
        return False
    err, minority = _two_color_metrics(blob)
    return err < QR_RECON_ERR_MAX and minority >= QR_MIN_MINORITY


def find_qr_rects(slide):
    """QR pictures on the slide as (left, top, width, height) tuples in POINTS,
    in slide z-order (first = lowest-index shape).

    The VBA writer can't run PIL, so QR identification happens here and the
    geometry is handed to InsertQR for exact in-place replacement.
    """
    rects = []
    for shape in slide.shapes:
        if _picture_is_qr(shape) and shape.left is not None and shape.top is not None:
            rects.append((shape.left.pt, shape.top.pt, shape.width.pt, shape.height.pt))
    return rects


def slide_has_existing_qr(slide):
    """True if the slide carries an (inherited) QR picture (see find_qr_rects)."""
    return bool(find_qr_rects(slide))


def resolve_target_slide_indices(prs, config, shownotes_url):
    """Determine which slides should receive the QR code.

    Args:
        prs: Presentation object
        config: QR config from speaker profile
        shownotes_url: The shownotes URL (for content detection)

    Returns:
        list of 0-based slide indices
    """
    position = config.get("slide_position", "closing")
    slide_count = len(prs.slides)

    indices = []

    if position in ("shownotes_slide", "both"):
        sn_idx = find_shownotes_slide(prs, shownotes_url)
        if sn_idx is not None:
            indices.append(sn_idx)
        else:
            # Fallback: slide index 3 (common shownotes position)
            fallback = min(3, slide_count - 1)
            print(f"  WARNING: Shownotes URL not found in slide text, falling back to slide {fallback + 1}")
            indices.append(fallback)

    if position in ("closing", "both"):
        closing_idx = slide_count - 1
        if closing_idx not in indices:
            indices.append(closing_idx)

    # Always target every slide that already carries a QR — a deck adapted from
    # another talk inherits stale QRs (e.g. an early shownotes slide + closing)
    # that must be replaced in place, not left beside a freshly-added one.
    for idx, slide in enumerate(prs.slides):
        if idx not in indices and slide_has_existing_qr(slide):
            indices.append(idx)

    if not indices:
        # Default: last slide
        indices.append(slide_count - 1)

    return sorted(indices)


# --- QR Insertion ---

def _format_qr_spec(slide_specs):
    """Encode per-slide insertion targets for the InsertQR macro.

    Each entry is "<1-based num>", optionally followed by
    ":<rL,rT,rW,rH>[,<rL,rT,rW,rH>...]" giving the rects (POINTS) of existing QR
    pictures to remove on that slide. The macro deletes those shapes and places
    the new QR at the FIRST rect (exact in-place replacement, preserving the
    inherited size); with no rects it places a new QR bottom-right. Entries are
    joined by ";". Detection is content-based (see find_qr_rects) because the VBA
    writer can't run PIL — so the chosen geometry travels in this spec.
    """
    parts = []
    for num, rects in slide_specs:
        if rects:
            flat = ",".join(f"{v:.2f}" for rect in rects for v in rect)
            parts.append(f"{num}:{flat}")
        else:
            parts.append(str(num))
    return ";".join(parts)


def insert_qr_via_powerpoint(deck_path, jobs, scripts_dir):
    """Insert the QR PNG(s) into the deck via the real PowerPoint app.

    Replaces the old python-pptx write (`insert-qr.sh` → InsertQR VBA macro), so
    PowerPoint serializes a valid .pptx (see rules/deck-editing-rules.md). Python
    has already identified each slide's existing QR(s) by content; the macro just
    removes those exact shapes and places the new QR there (or bottom-right when a
    slide has none). macOS + Microsoft PowerPoint only.

    jobs: list of (png_path, slide_specs) — one per QR color variant, where
    slide_specs is a list of (1-based slide number, [removal rects in points]).
    Each variant is a separate InsertQR pass; the deck is threaded through
    uniquely-named intermediates (PowerPoint keys open decks by filename) and the
    final result is moved back to deck_path.
    """
    wrapper = os.path.join(scripts_dir, "insert-qr.sh")
    if not os.path.isfile(wrapper):
        raise SystemExit(f"insert-qr.sh not found at {wrapper} — reinstall the plugin")
    current = deck_path
    for n, (png_path, slide_specs) in enumerate(jobs):
        out = f"{deck_path}.qrtmp{n}.pptx"  # distinct basename — no open-deck collision
        spec = _format_qr_spec(slide_specs)
        try:
            subprocess.run([wrapper, current, out, png_path, spec], check=True)
        except subprocess.CalledProcessError:
            raise SystemExit(
                f"ERROR: QR insertion failed in PowerPoint (insert-qr.sh, spec '{spec}'). "
                "Confirm DeckOps.pptm is open with macros enabled and Automation consent "
                "granted — see skills/presentation-creator/references/deck-editing-setup.md"
            )
        if current != deck_path:
            os.remove(current)  # drop the prior intermediate
        current = out
    shutil.move(current, deck_path)


def _validated_back_half(short_url, talk_slug):
    """Return the short link's back-half, which MUST be the talk slug.

    rules/qr-generation-rules.md §2 admits no exception: a link whose back-half
    is not the slug is the random-hash failure mode, and it stops the run
    whether the script created the link or an agent pre-resolved it.
    """
    segment = urllib.parse.urlparse(short_url).path.strip("/").split("/")[-1]
    if segment != talk_slug:
        raise ShortenerResolutionError(
            f"pre-resolved short URL back-half '{segment}' is not the talk slug "
            f"'{talk_slug}'. The back-half must be the slug "
            "(rules/qr-generation-rules.md §2); recreate the short link with "
            "the slug as its back-half, then re-run."
        )
    return segment


def _artifact_receipt(path, deck_dir, bg_hex):
    """Bind one generated PNG to the exact path written plus its digest."""
    absolute = os.path.abspath(path)
    if deck_dir and os.path.commonpath([absolute, os.path.abspath(deck_dir)]) == \
            os.path.abspath(deck_dir):
        path_root = "deck_dir"
        recorded = os.path.relpath(absolute, os.path.abspath(deck_dir))
    elif not os.path.isabs(path):
        path_root = "cwd"
        recorded = os.path.relpath(absolute, os.path.abspath("."))
    else:
        path_root = "absolute"
        recorded = absolute
    digest = hashlib.sha256(Path(absolute).read_bytes()).hexdigest()
    return {
        "path": recorded,
        "path_root": path_root,
        "sha256": digest,
        "bg_hex": bg_hex,
    }


# --- Tracking Database ---

def update_tracking_db(tracking_db, entry, artifacts):
    """Append or replace a qr_codes entry in the tracking database.

    Args:
        tracking_db: The full tracking database dict (mutated in place)
        entry: Metadata dict from resolve_short_url
        artifacts: Non-empty list of receipts from _artifact_receipt(), one per
            generated PNG. `qr_png_rel_path` mirrors the first for schema-v1
            readers; `artifacts` is the authoritative record.
    """
    if "qr_codes" not in tracking_db:
        tracking_db["qr_codes"] = []

    today = datetime.date.today().isoformat()
    talk_slug = entry["talk_slug"]

    new_entry = {
        "schema_version": QR_CODE_RECORD_SCHEMA_VERSION,
        "talk_slug": talk_slug,
        "target_url": entry["target_url"],
        "shortener": entry["shortener"],
        "short_path": entry.get("short_path"),
        "short_url": entry["short_url"],
        "shortener_link_id": entry.get("shortener_link_id"),
        "qr_png_rel_path": artifacts[0]["path"],
        "artifacts": artifacts,
        "created_at": today,
        "updated_at": today,
    }

    # Replace existing entry for same talk_slug, or append
    replaced = False
    for i, existing in enumerate(tracking_db["qr_codes"]):
        if existing.get("talk_slug") == talk_slug:
            new_entry["created_at"] = existing.get("created_at", today)
            tracking_db["qr_codes"][i] = new_entry
            replaced = True
            break

    if not replaced:
        tracking_db["qr_codes"].append(new_entry)


# --- Main ---

def main():
    parser = argparse.ArgumentParser(
        description="Generate a QR code and insert it into a PowerPoint deck.",
        epilog="Examples:\n"
               "  %(prog)s deck.pptx --talk-slug arc-of-ai --shownotes-url https://example.com/arc-of-ai\n"
               "  %(prog)s deck.pptx --talk-slug arc-of-ai --shownotes-url https://example.com/arc-of-ai"
               " --short-url https://jbaru.ch/arc-of-ai --short-provider bitly --short-link-id bit.ly/abc\n"
               "  %(prog)s deck.pptx --talk-slug arc-of-ai --shownotes-url https://example.com/arc-of-ai --dry-run\n"
               "  %(prog)s --png-only --talk-slug SLUG --shownotes-url https://example.com/arc-of-ai --output qr.png\n"
               "  %(prog)s --png-only --talk-slug SLUG --shownotes-url https://example.com/arc-of-ai"
               " --short-url https://jbaru.ch/arc-of-ai --bg-color 128,0,128\n",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("deck", nargs="?", default=None, help="Path to the .pptx deck file (not required with --png-only)")
    parser.add_argument("--talk-slug", required=True, help="Unique talk identifier (e.g., arc-of-ai)")

    # The canonical redirect target is always required. --short-url supplies an
    # agent-preresolved managed link (MCP mode) and does NOT replace the target:
    # recording the short URL as its own target loses the redirect relationship
    # the catalog exists to describe.
    parser.add_argument("--shownotes-url", required=True,
                        help="Canonical shownotes URL — the short link's redirect target")
    parser.add_argument("--short-url",
                        help="Pre-resolved short URL (MCP mode; skips shortening)")
    parser.add_argument("--short-provider", metavar="NAME",
                        help="Provider that issued --short-url (e.g. bitly, rebrandly)")
    parser.add_argument("--short-link-id", metavar="ID",
                        help="Provider-side link id for --short-url")

    parser.add_argument("--png-only", action="store_true",
                        help="Generate QR PNG only, without opening or modifying a deck")
    parser.add_argument("--output", metavar="PATH",
                        help="Output path for QR PNG (default: {deck_dir}/{talk-slug}-qr.png, or ./{talk-slug}-qr.png with --png-only)")
    parser.add_argument("--bg-color", metavar="R,G,B",
                        help="QR background color as R,G,B (e.g., 128,0,128). Default: detected from deck, or white with --png-only")
    parser.add_argument("--profile", help="Path to speaker-profile.json (default: {vault}/speaker-profile.json)")
    parser.add_argument("--vault", help="Path to vault root directory")
    parser.add_argument("--dry-run", action="store_true", help="Skip API calls and deck modification")

    args = parser.parse_args()

    if not args.png_only and not args.deck:
        parser.error("deck is required unless --png-only is specified")

    # Provider identity describes an agent-preresolved link. Accepting it
    # without --short-url would silently drop it, and half of it would catalog
    # an incomplete identity.
    if (args.short_provider or args.short_link_id) and not args.short_url:
        parser.error("--short-provider and --short-link-id require --short-url")
    if bool(args.short_provider) != bool(args.short_link_id):
        parser.error(
            "--short-provider and --short-link-id must be given together; "
            "a provider without its link id catalogs an incomplete identity"
        )

    if args.deck and not os.path.isfile(args.deck):
        print(f"ERROR: Deck file not found: {args.deck}")
        sys.exit(1)

    # Validate local inputs before URL shortening can create or retarget a link.
    explicit_bg = None
    if args.bg_color:
        try:
            parts = [int(x.strip()) for x in args.bg_color.split(",")]
            if len(parts) != 3 or not all(0 <= x <= 255 for x in parts):
                raise ValueError
            explicit_bg = tuple(parts)
        except ValueError:
            print(
                "ERROR: --bg-color must be R,G,B with values 0-255 "
                f"(got: {args.bg_color})"
            )
            sys.exit(1)

    # Determine vault path
    vault_path = args.vault
    if not vault_path:
        vault_path = os.path.expanduser("~/.claude/rhetoric-knowledge-vault")
    vault_present_at_start = os.path.isdir(vault_path)

    # Load config
    (
        speaker_profile,
        secrets,
        tracking_db,
        tracking_db_snapshot,
    ) = load_vault_config(vault_path, args.profile)
    if tracking_db_snapshot is not None:
        try:
            database_assessment = assess_tracking_database(tracking_db)
        except TrackingDatabaseError as exc:
            raise SystemExit(f"ERROR: cannot use tracking database: {exc}") from exc
        if not database_assessment.usable:
            raise SystemExit(
                "ERROR: tracking database has no usable prior state for this reader: "
                + ", ".join(database_assessment.reason_codes)
            )
        if not args.dry_run and database_assessment.state != "current":
            raise SystemExit(
                "ERROR: QR persistence requires the current tracking schema; run "
                "vault-ingress migration before generating the short link"
            )
    qr_config = speaker_profile.get("publishing_process", {}).get("qr_code", {})
    if not args.dry_run and vault_present_at_start:
        try:
            _require_tracking_db_snapshot(tracking_db_snapshot)
        except ValueError as exc:
            print(f"ERROR: {exc}", file=sys.stderr)
            sys.exit(1)

    effects_receipt = EffectsReceipt(args.talk_slug)

    # Serialize this slug's publication across its external effects. A dry run
    # makes none, and a missing vault has nowhere to place the lock.
    serialized = not args.dry_run and vault_present_at_start
    try:
        lock_scope = (
            qr_publication_lock(vault_path, args.talk_slug)
            if serialized
            else contextlib.nullcontext()
        )
        with lock_scope:
            prior_record = None
            if serialized:
                # State loaded before the lock is stale by definition: another
                # same-slug process may have committed a link while we waited.
                # Resolving from that view would create a duplicate instead of
                # retargeting the link it just made.
                try:
                    tracking_db = _reload_tracking_db(vault_path)
                except ValueError as exc:
                    print(f"ERROR: {exc}", file=sys.stderr)
                    sys.exit(1)
                prior_record = copy.deepcopy(
                    _qr_record_for(tracking_db, args.talk_slug)
                )
            _publish(args, effects_receipt, vault_path, vault_present_at_start,
                     speaker_profile, secrets, tracking_db, qr_config,
                     explicit_bg, prior_record)
    except ValueError as exc:
        # Lock acquisition failures are a CLI error path, not a traceback.
        print(f"ERROR: {exc}", file=sys.stderr)
        sys.exit(1)


def _publish(args, effects_receipt, vault_path, vault_present_at_start,
             speaker_profile, secrets, tracking_db, qr_config, explicit_bg,
             prior_record):
    # Determine the URL to encode in the QR
    if args.short_url:
        # MCP-preresolved mode
        qr_url = args.short_url
        shownotes_url = args.shownotes_url
        try:
            short_path = _validated_back_half(qr_url, args.talk_slug)
        except ShortenerResolutionError as e:
            print(f"ERROR: {e}", file=sys.stderr)
            print(
                "  No QR PNG, deck, or tracking-database change was made.",
                file=sys.stderr,
            )
            sys.exit(1)
        meta = {
            "talk_slug": args.talk_slug,
            "target_url": shownotes_url,
            "shortener": args.short_provider or "mcp_preresolved",
            "short_path": short_path,
            "short_url": qr_url,
            "shortener_link_id": args.short_link_id,
        }
        print(f"Using pre-resolved short URL: {qr_url} -> {shownotes_url}")
        # The agent created this link, so a failed commit still leaves it behind.
        effects_receipt.record_short_link(
            meta["shortener"], meta["shortener_link_id"], qr_url,
            action="preresolved",
        )
    else:
        # Direct resolution mode
        shownotes_url = args.shownotes_url
        print(f"Resolving short URL for: {shownotes_url}")
        try:
            qr_url, meta = resolve_short_url(
                shownotes_url, args.talk_slug, qr_config, secrets, tracking_db,
                args.dry_run, vault_path=vault_path,
                effects_receipt=effects_receipt,
            )
        except ShortenerResolutionError as e:
            print(f"ERROR: {e}", file=sys.stderr)
            print(
                "  No QR PNG, deck, or tracking-database change was made. Only "
                "an explicit 'shortener: none' may encode a raw target URL.",
                file=sys.stderr,
            )
            sys.exit(1)


    print(f"QR will encode: {qr_url}")

    # --- PNG-only mode: no deck needed ---
    if args.png_only:
        qr_bg = explicit_bg or (255, 255, 255)
        qr_fg = choose_fg_color(qr_bg)
        print(f"QR colors: fg=RGB{qr_fg}, bg=RGB{qr_bg}")

        qr_filename = f"{args.talk_slug}-qr.png"
        qr_path = args.output or os.path.join(".", qr_filename)

        if not args.dry_run:
            generate_qr_png(qr_url, qr_fg, qr_bg, qr_path)
            size_kb = os.path.getsize(qr_path) / 1024
            print(f"QR PNG saved: {qr_path} ({size_kb:.1f} KB)")
            artifacts = [_artifact_receipt(qr_path, None, None)]
            effects_receipt.record_artifacts([qr_path])
        else:
            print(f"DRY RUN: would save QR to {qr_path}")
            artifacts = None

    # --- Deck mode: open deck, detect colors, insert ---
    else:
        prs = Presentation(args.deck)

        # Determine target slides
        target_url_for_detection = shownotes_url if args.shownotes_url else qr_url
        slide_indices = resolve_target_slide_indices(prs, qr_config, target_url_for_detection)
        print(f"Target slides: {[i + 1 for i in slide_indices]}")

        bg_match = qr_config.get("bg_color_match", True)
        deck_dir = os.path.dirname(os.path.abspath(args.deck))

        # Resolve background color per slide — different slides may have
        # different backgrounds (e.g., shownotes vs closing/thank-you).
        # Group slides by their QR color scheme to avoid redundant PNGs.
        slide_colors = {}  # idx -> (qr_bg, qr_fg)
        # Existing-QR geometry per target slide (points), captured while the deck
        # is open — the PowerPoint writer can't run PIL, so it gets these rects.
        qr_rects_by_idx = {}  # idx -> [(L, T, W, H), ...]
        for idx in slide_indices:
            if explicit_bg:
                bg_rgb = explicit_bg
            else:
                slide_bg = resolve_slide_bg_rgb(prs.slides[idx])
                if slide_bg:
                    bg_rgb = slide_bg
                else:
                    print(f"  WARNING: Could not detect background for slide {idx + 1}, defaulting to white")
                    bg_rgb = (255, 255, 255)
            qr_bg = bg_rgb if bg_match else (255, 255, 255)
            qr_fg = choose_fg_color(qr_bg)
            slide_colors[idx] = (qr_bg, qr_fg)
            qr_rects_by_idx[idx] = find_qr_rects(prs.slides[idx])
            placement = "replace in place" if qr_rects_by_idx[idx] else "new (bottom-right)"
            print(f"  Slide {idx + 1}: bg=RGB{qr_bg}, fg=RGB{qr_fg} — {placement}")

        # Group slides by color scheme to generate minimal QR PNGs
        color_groups = {}  # (qr_bg, qr_fg) -> [slide_indices]
        for idx, colors in slide_colors.items():
            color_groups.setdefault(colors, []).append(idx)

        if not args.dry_run:
            qr_paths_generated = []
            insert_jobs = []  # (png_path, [(1-based num, [removal rects])]) per color variant
            for (qr_bg, qr_fg), indices in color_groups.items():
                if len(color_groups) == 1:
                    qr_filename = f"{args.talk_slug}-qr.png"
                else:
                    # Multiple color variants — suffix with bg hex
                    bg_hex = "{:02x}{:02x}{:02x}".format(*qr_bg)
                    qr_filename = f"{args.talk_slug}-qr-{bg_hex}.png"

                qr_path = args.output if (args.output and len(color_groups) == 1) else os.path.join(deck_dir, qr_filename)
                generate_qr_png(qr_url, qr_fg, qr_bg, qr_path)
                size_kb = os.path.getsize(qr_path) / 1024
                print(f"  QR PNG saved: {qr_filename} ({size_kb:.1f} KB) — for slide(s) {[i + 1 for i in indices]}")
                qr_paths_generated.append(
                    (qr_path, None if len(color_groups) == 1 else bg_hex)
                )
                # VBA is 1-based; pair each slide with its existing-QR rects
                insert_jobs.append((qr_path, [(i + 1, qr_rects_by_idx[i]) for i in indices]))

            # Release the read-only deck handle, then write via the real
            # PowerPoint app (valid OOXML; see rules/deck-editing-rules.md).
            prs = None
            here = os.path.dirname(os.path.abspath(__file__))
            effects_receipt.record_artifacts(
                [path for path, _bg in qr_paths_generated]
            )
            insert_qr_via_powerpoint(args.deck, insert_jobs, here)
            print(f"Deck updated via PowerPoint: {args.deck}")
            effects_receipt.record_deck(args.deck)
            # Every generated variant is cataloged, not just the first.
            artifacts = [
                _artifact_receipt(path, deck_dir, bg_hex)
                for path, bg_hex in qr_paths_generated
            ]
        else:
            # Dry run — just report what would happen
            for (qr_bg, qr_fg), indices in color_groups.items():
                print(f"  DRY RUN: would generate QR bg=RGB{qr_bg} fg=RGB{qr_fg} for slides {[i + 1 for i in indices]}")
            artifacts = None

    # A dry run generated nothing, so there is no artifact to bind and the
    # catalog is left untouched. The upsert itself happens inside
    # commit_qr_record(), against a freshly read generation.
    meta["target_url"] = shownotes_url

    if not args.dry_run:
        tdb_path = os.path.join(vault_path, "tracking-database.json")
        if vault_present_at_start:
            try:
                # Rebase onto the current generation rather than committing the
                # snapshot taken before the link, PNGs, and deck were changed.
                write_result = commit_qr_record(
                    tdb_path, meta, artifacts, prior_record
                )
            except ValueError as exc:
                print(f"ERROR: {exc}", file=sys.stderr)
                _report_unfinalized_effects(effects_receipt)
                sys.exit(1)
            if write_result.installed:
                print(f"Tracking DB updated: {tdb_path}")
            else:
                print(f"Tracking DB unchanged: {tdb_path}")
            print(
                "Tracking DB SHA-256: "
                f"{write_result.input_sha256} -> {write_result.output_sha256} "
                f"({write_result.durability_state})",
                file=sys.stderr,
            )
            for warning in write_result.warnings:
                print(f"WARNING: {warning}", file=sys.stderr)
        else:
            print(f"  NOTE: Vault path {vault_path} not found, tracking DB not persisted")

    print("\nDone.")


if __name__ == "__main__":
    main()
