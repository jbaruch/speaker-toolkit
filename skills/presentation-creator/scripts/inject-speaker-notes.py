#!/usr/bin/env python3
"""Inject speaker notes into a PowerPoint deck from a JSON map.

The JSON file maps slide indices (0-based) to notes text:
  {"0": "", "1": "Brief intro.", "2": "Core argument starts here."}

Usage:
    inject-speaker-notes.py <deck.pptx> <notes.json>

Examples:
    inject-speaker-notes.py presentation.pptx speaker-notes.json

Keynote compatibility:
    After injection, this script post-processes the .pptx to add a
    <p:notesMasterIdLst> element to ppt/presentation.xml. python-pptx adds
    the notesMaster relationship but omits this element; PowerPoint tolerates
    the omission, but Keynote rejects the file as "invalid format". The patch
    is idempotent and only runs when a notesMaster relationship is present.
"""

import json
import re
import shutil
import sys
import zipfile
from pptx import Presentation


def patch_notes_master_idlst(pptx_path):
    """Add <p:notesMasterIdLst> to presentation.xml if a notesMaster rel exists.

    python-pptx writes the notesMaster relationship into
    ppt/_rels/presentation.xml.rels but never writes the corresponding
    <p:notesMasterIdLst> element into ppt/presentation.xml. The OOXML spec
    requires it, PowerPoint ignores its absence, Keynote does not.

    This function is idempotent: returns immediately if the element is
    already present or if no notesMaster relationship exists.
    """
    with zipfile.ZipFile(pptx_path, "r") as z:
        rels_xml = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        pres_xml = z.read("ppt/presentation.xml").decode("utf-8")

    m = re.search(
        r'<Relationship Id="(rId\d+)"[^>]*notesMaster[^>]*/>',
        rels_xml,
    )
    if not m:
        return False  # no notes master = nothing to patch
    if "<p:notesMasterIdLst>" in pres_xml:
        return False  # already patched

    notes_rid = m.group(1)
    notes_master_elem = (
        f'<p:notesMasterIdLst>'
        f'<p:notesMasterId r:id="{notes_rid}"/>'
        f'</p:notesMasterIdLst>'
    )
    patched = pres_xml.replace(
        "</p:sldIdLst>",
        f"</p:sldIdLst>{notes_master_elem}",
        1,
    )

    tmp_path = pptx_path + ".tmp"
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(
        tmp_path, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/presentation.xml":
                data = patched.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp_path, pptx_path)
    return True


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <deck.pptx> <notes.json>", file=sys.stderr)
        sys.exit(1)

    deck_path, notes_path = sys.argv[1], sys.argv[2]

    with open(notes_path) as f:
        notes_map = json.load(f)

    prs = Presentation(deck_path)
    injected = 0

    for idx_str, notes_text in notes_map.items():
        idx = int(idx_str)
        if notes_text and idx < len(prs.slides):
            slide = prs.slides[idx]
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
            injected += 1

    prs.save(deck_path)
    patched = patch_notes_master_idlst(deck_path)
    suffix = " (Keynote-compat patch applied)" if patched else ""
    print(f"Injected speaker notes into {injected} slides. Saved to {deck_path}{suffix}")
